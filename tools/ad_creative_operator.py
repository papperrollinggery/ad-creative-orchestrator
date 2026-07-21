#!/usr/bin/env python3
"""Non-developer operation surface for Ad Creative Orchestrator projects."""

from __future__ import annotations

import argparse
import csv
import fcntl
import hashlib
import html
import io
import importlib.metadata as metadata
import importlib.util
import json
import os
import re
import selectors
import shutil
import stat
import subprocess
import sys
import tempfile
import threading
import time
import unicodedata
import urllib.error
import urllib.parse
import urllib.request
import webbrowser
import zipfile
from contextvars import ContextVar
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import BinaryIO, Callable, Iterable
import xml.etree.ElementTree as ET

sys.dont_write_bytecode = True

from adco_core.facts import (
    export_intake_analysis_request,
    import_intake_analysis,
    load_fact_inventory,
    run_evidence_intake,
)
from adco_core.creative_contract import (
    BRIEF_CONTRACT_REL,
    BRIEF_SNAPSHOT_REL,
    CANDIDATE_IMPORT_RECEIPT_REL,
    CANDIDATE_SCHEMA_REL,
    CREATIVE_DIRECTIONS_REL,
    CRITIC_RECEIPT_REL,
    CURRENT_CANDIDATE_REL,
    GENERATION_REQUEST_REL,
    OPEN_GAPS_REL,
    OPTION_MATRIX_REL,
    create_creative_brief,
    import_creative_candidate,
    review_creative_candidate,
)
from adco_core.commands.run import RunPreflightError, execute_lightweight_run
from adco_core.ingestion import (
    LOCAL_SOURCE_PREFIX,
    ingest_source_rows,
    load_local_source_paths,
    load_local_source_paths_from_project_fd,
    register_local_source_path,
    source_path_label,
    source_row_files,
    source_row_material_roots,
)
from adco_core.incremental_validation import run_incremental_validation
from adco_core.specialist_exchange import (
    V2_CONTRACT_VERSION,
    build_v2_handoff,
    contained_project_path as contained_v2_project_path,
    current_scope_manifest as v2_scope_manifest,
    manifest_digest as v2_manifest_digest,
    negotiate_contract_version,
    validate_v2_exchange_row,
    validate_v2_receipt_outputs,
    v2_boundary_errors,
)
from init_project import agents_policy_status, copy_content_template, copy_template
from runtime_paths import (
    CONTENT_SURFACE,
    DELIVERY_SURFACE,
    project_surface,
    published_docs_root,
    repo_or_module_root,
    skill_draft_dir,
    source_root,
    template_root,
)
from specialist_schema_validation import (
    specialist_control_plane_errors,
    specialist_generation_authorization_errors,
    validate_specialist_payload,
)
from validate_project import (
    THREADOPS_RECEIPT_REQUIRED_PROOF,
    current_truth_value,
    legacy_baseline_message_allowed,
    receipt_proof_values,
    string_validation_errors,
    validate,
    validate_issues,
    validate_client_delivery_readiness,
)


REPO_ROOT = repo_or_module_root()
TEMPLATE_ROOT = template_root()
SKILL_DRAFT_DIR = skill_draft_dir()
PACKAGE_NAME = "ad-creative-orchestrator"
FALLBACK_VERSION = "0.3.2"
CONTROL_PLANE_SCHEMA_VERSION = "2.0"
CONTROL_PLANE_SCHEMA_REL = Path("AD-creative/orchestrator/control_plane_schema.json")
CONTROL_PLANE_MIGRATION_MANIFEST_REL = Path(
    "AD-creative/orchestrator/migrations/control_plane_v2_manifest.json"
)
FINAL_DELIVERY_CONFIRMATION_PROTOCOL = (
    "adco.final-delivery-reconciliation-confirmation"
)
FINAL_DELIVERY_CONFIRMATION_VERSION = "1.0"
DELIVERY_COMMAND_ACTIVE: ContextVar[bool] = ContextVar(
    "adco_delivery_command_active",
    default=False,
)
MIGRATION_CWD_LOCK = threading.RLock()
FINAL_DELIVERY_HOST_ATTESTATION_PROTOCOL = "adco.host-readback-attestation"
FINAL_DELIVERY_HOST_ATTESTATION_VERSION = "1.0"
FINAL_DELIVERY_HOST_ATTESTATION_ROOT = Path(
    "AD-creative/orchestrator/host_attestations"
)
SPECIALIST_EXCHANGE_PROTOCOL = "adco.specialist-exchange"
SPECIALIST_EXCHANGE_VERSION = "1.0"
DIRCREATIVE_PROFILE_ID = "dircreative.film-preproduction"
CLIENT_OUTLINE_CONFIRMATION_REL = Path(
    "AD-creative/client_review/client_outline_confirmation.json"
)
SPECIALIST_RESERVED_CLAIMS = (
    "client_ready",
    "ppt_ready",
    "final_delivery_ready",
    "send_ready",
    "project_complete",
    "control_plane_updated",
)
CURRENT_VERSION_TRUTH_KEYS = (
    "current_version_id",
    "current_pptx_artifact_id",
    "current_pdf_artifact_id",
    "current_preview_artifact_id",
    "current_text_extract_artifact_id",
    "current_ppt_editability_artifact_id",
    "version_map_status",
    "last_archive_before_edit",
)
CURRENT_VIEW_VERSION_STATUSES = {
    "draft",
    "internal_review",
    "ready",
    "active",
    "current",
}
SPECIALIST_ID_PATTERN = re.compile(r"^[A-Za-z0-9][A-Za-z0-9._-]{0,127}$")
DEFAULT_SKILL_INSTALL_DIR = Path.home() / ".codex/skills/ad-creative-orchestrator"
DEFAULT_SKILLSHUB_INSTALL_DIR = Path.home() / ".skillshub/ad-creative-orchestrator"
SKILL_INSTALL_MANIFEST_NAME = ".adco-skill-install-manifest.json"
DASHBOARD_REL = Path("AD-creative/handoff/操作台.html")
COUNCIL_REPORT_REL = Path("AD-creative/gates/THREE-COUNCIL-READINESS_report.md")
GOAL_PLAN_TEMPLATE_REL = Path("AD-creative/orchestrator/goal_iteration_plan_template.md")
GOAL_ITERATIONS_REL = Path("AD-creative/orchestrator/goal_iterations")
SUPPORT_BUNDLE_REL = Path("AD-creative/handoff/support_bundle.md")
PRIVATE_LOCAL_STATE_REL = Path(".adco-local")
PRIVATE_MARKER_SCAN_CHUNK_BYTES = 1024 * 1024
PRIVATE_MARKER_SCAN_MAX_FILE_BYTES = 256 * 1024 * 1024
PRIVATE_MARKER_SCAN_MAX_ARCHIVE_BYTES = 256 * 1024 * 1024
PRIVATE_MARKER_SCAN_MAX_MEMBER_BYTES = 64 * 1024 * 1024
PRIVATE_MARKER_SCAN_MAX_ARCHIVE_MEMBERS = 10_000
PRIVATE_MARKER_SCAN_MAX_PDF_PAGES = 2_000
PRIVATE_MARKER_SCAN_PDF_TIMEOUT_SECONDS = 30
PRIVATE_MARKER_SCAN_MAX_PDF_MEMORY_BYTES = 768 * 1024 * 1024
SAMPLE_MATERIAL_REL = Path("00_项目资料_ProjectMaterials/01_客户资料_ClientMaterials/sample_brief.md")
DEFAULT_DEMO_PROJECT = Path(tempfile.gettempdir()) / "adco-demo"
SAMPLE_GOAL = "基于内置 brief 提炼品牌策略，并提出三条机制不同的内部创意方向。"
SAMPLE_BRIEF = """# Sample Creative Brief

项目：NOVA Trail 户外功能饮料新品广告创意样例

客户希望输出一版广告创意提案，用于内部评审。
品牌主张：轻负担补给，目标人群是周末轻户外和城市通勤人群，竞品参考只借鉴场景节奏，不复制画面、包装或口号。
视觉资产计划需要规划关键视觉、产品露出、生成图边界、asset slot 和 visual QA。
本轮交付需要包含可编辑 PPT 结构、参考证据链、图片资产清单和内部确认清单。
客户明确不要使用未经授权 logo、真实品牌包装、不可追溯参考截图或未批准生成图。
样例已提供产品高清图、包装方向、字体方向和官方视觉规范摘要。
参考方向希望偏真实户外、清爽补给、清晨山路、手持产品、轻运动人群。
PPT 需要保留文本可编辑，客户稿发送前必须经过 Gate。
"""
ID_SUFFIX_PATTERN = re.compile(r"(\d+)$")
CLIENT_OWNER_PATTERN = re.compile(r"客户|不要")
CONFIRMATION_HEADER_CELLS = {"ID", "---"}
CLIENT_VISIBLE_VALUES = {"client_visible", "client_visible_ready", "sent"}
PASS_GATE_VALUES = {"pass", "passed"}
CLIENT_REVIEW_ASSET_STATUSES = {"selected", "approved", "done"}
GENERATED_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp"}
ACTIVE_ASSET_STATUSES = {"registered", "selected", "approved", "done"}
SELECTED_ASSET_STATUSES = {"selected", "approved", "done"}
UNAVAILABLE_ASSET_STATUSES = {
    "superseded",
    "deprecated",
    "archived",
    "rejected",
    "old",
    "obsolete",
    "removed",
}
CLIENT_ASSET_APPROVAL_VALUES = {
    "pass",
    "approved",
    "client_approved",
    "approved_for_client",
    "licensed",
    "cleared",
}
BROWSER_HELD_ASSET_MARKERS = (
    "browser",
    "浏览器",
    "grok",
    "chatgpt",
    "imagegen",
    "image_gen",
)
CLIENT_LANGUAGE_RELIABLE_NAME_PATTERN = re.compile(
    r"client[-_\s]?(copy|draft|visible|review)|客户稿|客户文案|客户审阅|审阅稿|final[-_\s]?(copy|delivery)|最终交付",
    re.IGNORECASE,
)
CLIENT_OUTLINE_BODY_MAX_CHARS = 420
CLIENT_OUTLINE_VISUAL_STATUSES = {
    "existing_image",
    "existing_asset",
    "to_generate",
    "pending_generation",
    "placeholder",
    "text_only",
    "no_visual",
}
EXISTING_IMAGE_STATUSES = {"existing_image", "existing_asset"}
CLIENT_OUTLINE_FIELDS = [
    "slide_id",
    "page_title",
    "body_copy",
    "client_confirmation_point",
    "material_role",
    "visual_slot",
    "visual_asset_status",
    "asset_ids",
    "visibility",
    "status",
    "notes",
]
ASSET_CURRENT_FIELDS = [
    "asset_id",
    "source",
    "platform",
    "conversation",
    "local_file",
    "path",
    "sha256",
    "original_or_processed",
    "approval",
    "direct_client_use",
    "used_in_slide",
    "qa_flags",
    "protected",
    "status",
    "notes",
]
ASSET_AUTHORIZATION_FIELDS = [
    "authorization_id",
    "asset_id",
    "asset_sha256",
    "approval_scope",
    "approved_by",
    "approved_at",
    "evidence_ref",
    "evidence_sha256",
    "status",
    "revoked_at",
    "notes",
]
FINAL_DELIVERY_LOCK_FIELDS = [
    "lock_id",
    "path",
    "sha256",
    "size_bytes",
    "mtime",
    "protected",
    "registered_at",
    "notes",
    "inventory_state",
    "reconciliation_state",
    "reconciliation_kind",
    "reconciles_lock_id",
    "supersedes_lock_id",
    "confirmed_by",
    "confirmed_at",
    "evidence_ref",
    "evidence_sha256",
    "host_attestation_ref",
    "host_attestation_sha256",
    "version_id",
    "supersedes_version_id",
    "status_reason",
]
ARTIFACT_INDEX_FIELDS = [
    "artifact_id",
    "artifact_type",
    "path",
    "stage",
    "version",
    "status",
    "visibility",
    "source_event_ids",
    "linked_requirements",
    "linked_work_items",
    "linked_references",
    "linked_assets",
    "gate_status",
    "supersedes_artifact_id",
    "created_at",
    "updated_at",
    "sha256",
    "size_bytes",
    "derived_from_artifact_id",
    "derived_from_sha256",
    "lifecycle_state",
    "original_path",
    "cleanup_ref",
    "removed_at",
    "removal_reason",
    "superseded_by",
    "status_reason",
]
ARTIFACT_LIFECYCLE_VALUES = {
    "active",
    "pending",
    "superseded",
    "withdrawn",
    "archived",
    "deprecated",
    "rejected",
    "removed",
    "legacy_unresolved_tombstone",
    "legacy_unknown",
}
ARTIFACT_INACTIVE_LIFECYCLE_VALUES = {
    "superseded",
    "withdrawn",
    "archived",
    "deprecated",
    "rejected",
    "removed",
    "legacy_unresolved_tombstone",
}
LEGACY_BASELINE_SOURCE_RELS = (
    "AD-creative/orchestrator/source_events.csv",
    "AD-creative/orchestrator/requirements.csv",
    "AD-creative/orchestrator/work_items.csv",
    "AD-creative/orchestrator/artifact_index.csv",
    "AD-creative/orchestrator/gate_log.csv",
    "AD-creative/orchestrator/version_map.csv",
    "AD-creative/references/reference_cards.csv",
    "AD-creative/visual_assets/asset_manifest.csv",
    "AD-creative/orchestrator/profile_knowledge/profile_subjects.csv",
    "AD-creative/orchestrator/profile_knowledge/meeting_voice_map.csv",
    "AD-creative/orchestrator/profile_knowledge/profile_insights.csv",
    "AD-creative/orchestrator/profile_knowledge/profile_conflicts.csv",
)
GATE_LOG_FIELDS = [
    "gate_id",
    "gate_run_id",
    "stage",
    "status",
    "score",
    "checked_artifacts",
    "target_ref",
    "target_sha256",
    "evidence_snapshot_ref",
    "evidence_snapshot_sha256",
    "blocking_issues",
    "revision_items",
    "questions",
    "next_state",
    "created_at",
    "owner",
    "supersedes_gate_run_id",
]
SPECIALIST_EXCHANGE_INDEX_FIELDS = [
    "exchange_id",
    "handoff_id",
    "attempt",
    "work_id",
    "provider_id",
    "profile_id",
    "contract_version",
    "descriptor_sha256",
    "handoff_sha256",
    "baseline_path",
    "baseline_sha256",
    "compatibility_status",
    "execution_mode",
    "lane_id",
    "thread_id",
    "handoff_path",
    "receipt_path",
    "receipt_sha256",
    "outcome",
    "adoption_path",
    "adoption_sha256",
    "adoption_decision",
    "thread_reconciliation_ref",
    "created_at",
    "updated_at",
]
CLIENT_LANGUAGE_BLOCKLIST = [
    "prompt",
    "thread",
    "worker",
    "lane",
    "receipt",
    "adoption",
    "cleanup",
    "gate",
    "AI",
    "ChatGPT",
    "Grok",
    "ImageGen",
    "制作表",
    "客户稿里标成",
    "可授权",
    "需确认",
    "内部",
    "执行过程",
    "待确认",
    "TBD",
    "TODO",
]
CLIENT_LANGUAGE_ASCII_PATTERN = re.compile(
    r"\b(prompt|thread|worker|lane|receipt|adoption|cleanup|gate|ai|chatgpt|grok|imagegen|tbd|todo)\b",
    re.IGNORECASE,
)
CLIENT_LANGUAGE_CJK_PATTERN = re.compile(
    "|".join(
        re.escape(pattern)
        for pattern in [
            "制作表",
            "客户稿里标成",
            "可授权",
            "需确认",
            "内部",
            "执行过程",
            "待确认",
            "提示词",
            "线程",
            "工作流",
            "执行记录",
        ]
    )
)
VISUAL_LAYOUT_RISK_PATTERNS = [
    "stretch",
    "stretched",
    "distorted",
    "crop",
    "cropped",
    "too small",
    "low-res",
    "card inside card",
    "nested card",
    "report-like",
    "too short",
    "over-simplified",
    "拉伸",
    "变形",
    "裁切错误",
    "图太小",
    "卡片套卡片",
    "报告感",
    "文字过短",
    "过度简化",
    "拥挤",
    "阅读顺序",
    "不匹配",
    "mismatch",
    "wrong image",
    "重复误用",
    "duplicate misuse",
    "比例不当",
    "aspect mismatch",
    "竖屏误用",
    "横屏误用",
]
VISUAL_LAYOUT_RISK_PATTERN = re.compile(
    "|".join(re.escape(pattern.lower()) for pattern in sorted(VISUAL_LAYOUT_RISK_PATTERNS, key=len, reverse=True))
)
TEXT_CLIENT_SCAN_SUFFIXES = {".md", ".txt", ".csv", ".json", ".yaml", ".yml"}
FINAL_DELIVERY_SUFFIXES = {".pptx", ".pdf", ".png", ".jpg", ".jpeg", ".webp", ".txt", ".md"}
CREATIVE_PRODUCTION_KINDS = {"moodboard", "ads", "shots"}
GOAL_PHASES = ("P0", "P1", "P2", "P3", "P4", "P5", "P6", "P7", "P8")
GOAL_PHASE_NAMES = {
    "P0": "Intake 与事实基线",
    "P1": "客户可读文本框架",
    "P2": "文本框架人工确认",
    "P3": "创意、参考与按需 Specialist",
    "P4": "Immutable PPT 版本导出",
    "P5": "语言、视觉、授权与可编辑性 Gate",
    "P6": "Fresh Client Pack Binding",
    "P7": "独立审阅与发送准备 Gate",
    "P8": "反馈合并与下一版本",
}
GOAL_PHASE_GATE_HINTS = {
    "P0": ("intake", "brief", "project_readiness"),
    "P1": ("client_outline", "client_review"),
    "P2": ("client_outline_confirmation",),
    "P3": ("creative", "reference_research", "specialist_handoff"),
    "P4": ("ppt_gate",),
    "P5": ("client_review", "visual_review", "ppt_gate"),
    "P6": ("client_pack", "final_delivery"),
    "P7": ("client_send_readiness", "final_delivery"),
    "P8": ("feedback", "next_version"),
}
VISUAL_RISK_PATTERNS = {
    "contact sheet",
    "low-quality collage",
    "collage",
    "placeholder-only",
    "fake logo",
    "假 logo",
    "假logo",
}
VISUAL_RISK_PATTERN = re.compile(
    "|".join(re.escape(pattern) for pattern in sorted(VISUAL_RISK_PATTERNS, key=len, reverse=True))
)
CREATIVE_PROPOSAL_ARTIFACTS = [
    (
        "ART-AUTO-CREATIVE-DIRECTIONS",
        "creative_directions",
        Path("AD-creative/creative/creative_directions.md"),
    ),
    (
        "ART-AUTO-CREATIVE-OPTION-MATRIX",
        "creative_option_matrix",
        Path("AD-creative/creative/option_matrix.csv"),
    ),
    (
        "ART-AUTO-PROPOSAL-STRUCTURE",
        "proposal_structure",
        Path("AD-creative/proposal_architecture/proposal_structure.md"),
    ),
    (
        "ART-AUTO-CLIENT-OUTLINE",
        "client_outline",
        Path("AD-creative/client_review/client_outline.csv"),
    ),
    (
        "ART-AUTO-SLIDE-SPEC",
        "slide_spec",
        Path("AD-creative/client_review/slide_spec.md"),
    ),
]
CREATIVE_PROPOSAL_REQUIRED_LABELS = [
    "business problem",
    "client real objective",
    "target audience",
    "behavior barrier",
    "consumer insight",
    "feature to benefit",
    "brand/category/competitor notes",
    "strategy path",
    "creative proposition",
    "core message",
    "key visual/action",
    "title/use case",
    "risk",
    "why choose",
    "proposal outline",
]
GENERIC_CREATIVE_PATTERNS = [
    "unlock",
    "elevate",
    "game changer",
    "next level",
    "seamless",
    "innovative",
    "empower",
    "reimagine",
    "breakthrough",
    "bold new",
    "打造全新体验",
    "重新定义",
    "引爆",
    "破圈",
    "赋能",
    "焕新",
    "不止于",
    "美好生活",
    "无限可能",
]
GENERIC_CREATIVE_PATTERN = re.compile(
    "|".join(re.escape(pattern) for pattern in sorted(GENERIC_CREATIVE_PATTERNS, key=len, reverse=True)),
    re.IGNORECASE,
)
CASE_CLAIM_PATTERN = re.compile(
    r"case study|案例|campaign proved|proven by|according to|数据显示|行业报告|竞品证明|参考案例",
    re.IGNORECASE,
)
INTERNAL_LANGUAGE_PATTERN = re.compile(
    r"prompt|thread|worker|lane|subagent|codex|worktree|执行线程|工作线程|提示词|子代理|泳道|lane plan",
    re.IGNORECASE,
)
CHATBOT_RESIDUE_PATTERN = re.compile(
    r"hope this helps|let me know|would you like|want me to|of course[!,]|certainly[!,]|"
    r"\bhere (?:is|are) (?:a|an|the)\b|"
    r"希望这对.{0,8}有帮助|请告诉我|您想要|当然[！!]|一定[！!]",
    re.IGNORECASE,
)
VAGUE_AUTHORITY_PATTERN = re.compile(
    r"industry reports?|experts? (?:argue|believe|say|suggest)|observers? (?:have )?(?:cited|noted|say)|"
    r"some critics argue|studies show|data shows|leading experts?|"
    r"行业报告|专家(?:认为|指出|表示)|观察者(?:指出|认为)|一些批评者|数据显示|多方(?:认为|指出)",
    re.IGNORECASE,
)
EXAGGERATED_SIGNIFICANCE_PATTERN = re.compile(
    r"stands as|serves as|testament to|pivotal moment|crucial role|key role|marks a shift|"
    r"broader trend|evolving landscape|lasting impact|transformative|game-changing|"
    r"标志着|见证了|是.+证明|关键(?:性)?(?:时刻|转折|作用)|至关重要|深远影响|"
    r"不断演变的格局|变革性|重塑(?:行业|市场|格局)",
    re.IGNORECASE,
)
NOT_ONLY_BUT_PATTERN = re.compile(
    r"not only .{0,120} but|not just .{0,120} (?:but|it'?s)|not merely .{0,120} but|"
    r"不仅.{0,80}而且|不止于.{0,80}(?:更是|而是)|不只是.{0,80}(?:更是|而是)",
    re.IGNORECASE,
)
GENERIC_AI_VOCABULARY_PATTERN = re.compile(
    r"additionally|align with|crucial|delve|enduring|enhance|foster(?:ing)?|highlight(?:ing)?|"
    r"interplay|intricate|landscape|pivotal|showcase|tapestry|testament|underscore|valuable|vibrant|"
    r"此外|至关重要|深入探讨|彰显|凸显|赋能|焕新|无缝|沉浸式|多维|格局|"
    r"重要抓手|有力抓手|强势助推|生态闭环",
    re.IGNORECASE,
)
CLIENT_WRITING_RISK_CODES = {
    "CHATBOT_RESIDUE": CHATBOT_RESIDUE_PATTERN,
    "VAGUE_AUTHORITY_CLAIM": VAGUE_AUTHORITY_PATTERN,
    "EXAGGERATED_SIGNIFICANCE": EXAGGERATED_SIGNIFICANCE_PATTERN,
    "FORMULAIC_NOT_ONLY_BUT": NOT_ONLY_BUT_PATTERN,
}
PROFILE_SUBJECT_FIELDS = [
    "subject_id",
    "subject_type",
    "name",
    "role_or_title",
    "organization",
    "source_event_ids",
    "first_seen_at",
    "last_seen_at",
    "profile_status",
    "influence_level",
    "decision_power",
    "traits",
    "needs",
    "preferences",
    "concerns",
    "notes",
]
PROFILE_VOICE_FIELDS = [
    "voice_id",
    "source_event_id",
    "file_path",
    "speaker",
    "utterance",
    "need_signal",
    "preference_signal",
    "concern_signal",
    "decision_signal",
    "influence_level",
    "decision_power",
    "evidence_quote",
    "confidence",
    "status",
]
PROFILE_INSIGHT_FIELDS = [
    "insight_id",
    "subject_id",
    "subject_type",
    "source_event_id",
    "file_path",
    "insight_type",
    "statement",
    "evidence_quote",
    "confidence",
    "status",
    "priority",
    "linked_requirement_ids",
    "supersedes_insight_id",
    "created_at",
    "updated_at",
]
PROFILE_CONFLICT_FIELDS = [
    "conflict_id",
    "topic",
    "source_event_ids",
    "subject_ids",
    "conflict_summary",
    "recommended_resolution",
    "status",
    "confidence",
    "evidence_quotes",
    "created_at",
    "updated_at",
]
PROFILE_STATUS_VALUES = {"candidate", "confirmed", "conflicted", "deprecated"}
PROFILE_SUBJECT_TYPES = {"participant", "brand", "company", "client_group"}
PROFILE_DECISION_LEVELS = {"high", "medium", "low", "unknown"}
PROFILE_SOURCE_LABELS = {
    "项目",
    "品牌",
    "产品",
    "参考方向",
    "本轮交付",
    "客户希望",
    "客户明确",
    "背景",
    "目标",
    "交付",
}
SPEAKER_LINE_PATTERN = re.compile(
    r"^(?:\[(?P<bracket>[^\]]{1,28})\]|(?P<label>[\w\u4e00-\u9fff·（）() /-]{1,28}))\s*[：:]\s*(?P<body>.+)$"
)
NEED_SIGNAL_PATTERN = re.compile(r"希望|想要|需要|目标|必须|本轮|交付|要做|要体现|诉求")
PREFERENCE_SIGNAL_PATTERN = re.compile(r"喜欢|偏好|更想|风格|调性|参考方向|年轻|高端|真实|清爽|专业")
CONCERN_SIGNAL_PATTERN = re.compile(r"担心|不要|不能|风险|怕|预算|时间|禁区|不希望|避免|限制")
DECISION_SIGNAL_PATTERN = re.compile(r"拍板|决定|确认|最终|老板|领导|负责人|决策|定下来|可以定|我来定")
BRAND_SIGNAL_PATTERN = re.compile(r"品牌|调性|主张|定位|人群|logo|Logo|包装|视觉规范|产品")
COMPANY_SIGNAL_PATTERN = re.compile(r"公司|团队|集团|部门|内部|统一意见|甲方|客户内部|汇报链路")
CONFLICT_SIGNAL_PATTERN = re.compile(r"分歧|不同意|冲突|但是|不过|有人认为|另一个方向|还没统一|需要统一")


def now_iso() -> str:
    return datetime.now().astimezone().isoformat(timespec="seconds")




def skill_tree_files(root: Path) -> dict[str, Path]:
    """Return the regular, non-hidden files managed by the packaged Skill tree."""
    if not root.is_dir():
        raise FileNotFoundError(f"skill draft directory not found: {root}")
    resolved_root = root.resolve()
    files: dict[str, Path] = {}
    for path in sorted(root.rglob("*")):
        rel = path.relative_to(root)
        if any(part.startswith(".") or part == "__pycache__" for part in rel.parts):
            continue
        if path.is_symlink():
            raise RuntimeError(f"skill source tree must not contain symlinks: {path}")
        if not path.is_file():
            continue
        resolved = path.resolve()
        try:
            resolved.relative_to(resolved_root)
        except ValueError as exc:
            raise RuntimeError(f"skill source file escapes source tree: {path}") from exc
        files[rel.as_posix()] = path
    if "SKILL.md" not in files:
        raise FileNotFoundError(f"skill draft not found: {root / 'SKILL.md'}")
    return files


def skill_tree_hash(files: dict[str, Path]) -> str:
    return skill_hash_map_digest(
        {rel_path: file_sha256(path) for rel_path, path in files.items()}
    )


def skill_hash_map_digest(hashes: dict[str, str]) -> str:
    digest = hashlib.sha256()
    for rel_path, file_hash in sorted(hashes.items()):
        digest.update(rel_path.encode("utf-8"))
        digest.update(b"\0")
        digest.update(file_hash.encode("ascii"))
        digest.update(b"\n")
    return digest.hexdigest()


def safe_skill_target_path(target: Path, rel_path: str) -> Path:
    relative = Path(rel_path)
    if relative.is_absolute() or not relative.parts or ".." in relative.parts:
        raise RuntimeError(f"unsafe managed Skill path: {rel_path}")
    current = target
    for part in relative.parts[:-1]:
        current = current / part
        if current.is_symlink():
            raise RuntimeError(f"Skill target contains a symlinked directory: {current}")
        if current.exists() and not current.is_dir():
            raise RuntimeError(f"Skill target parent is not a directory: {current}")
    destination = target / relative
    if destination.is_symlink():
        raise RuntimeError(f"Skill target file is a symlink: {destination}")
    if destination.exists() and not destination.is_file():
        raise RuntimeError(f"Skill target path is not a regular file: {destination}")
    return destination


def read_skill_install_manifest(target: Path) -> dict[str, object] | None:
    path = target / SKILL_INSTALL_MANIFEST_NAME
    if not path.exists():
        return None
    if path.is_symlink() or not path.is_file():
        raise RuntimeError(f"Skill install manifest is not a regular file: {path}")
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as exc:
        raise RuntimeError(f"Skill install manifest is malformed: {path}") from exc
    if not isinstance(payload, dict) or payload.get("skill") != "ad-creative-orchestrator":
        raise RuntimeError(f"Skill install manifest has the wrong owner: {path}")
    managed = payload.get("managed_files")
    if not isinstance(managed, list) or not all(isinstance(item, str) for item in managed):
        raise RuntimeError(f"Skill install manifest has invalid managed_files: {path}")
    managed_hashes = payload.get("managed_file_sha256")
    if (
        payload.get("schema_version") != 2
        or not isinstance(managed_hashes, dict)
        or set(managed_hashes) != set(managed)
        or not all(
            isinstance(rel_path, str)
            and isinstance(file_hash, str)
            and re.fullmatch(r"[0-9a-f]{64}", file_hash)
            for rel_path, file_hash in managed_hashes.items()
        )
        or skill_hash_map_digest(
            {str(key): str(value) for key, value in managed_hashes.items()}
        )
        != payload.get("source_tree_sha256")
    ):
        raise RuntimeError(f"Skill install manifest is legacy or tampered: {path}")
    for rel_path in managed:
        safe_skill_target_path(target, rel_path)
    return payload


def install_global_skill(target: Path = DEFAULT_SKILL_INSTALL_DIR) -> dict[str, object]:
    source_files = skill_tree_files(SKILL_DRAFT_DIR)
    requested_target = target.expanduser()
    if requested_target.is_symlink():
        resolved_link = requested_target.resolve()
        allowed_compatibility_link = (
            requested_target.absolute() == DEFAULT_SKILL_INSTALL_DIR.absolute()
            and resolved_link == DEFAULT_SKILLSHUB_INSTALL_DIR.resolve()
        )
        if not allowed_compatibility_link:
            raise RuntimeError(
                f"Skill target root symlink is not an approved compatibility link: {requested_target}"
            )
    target = requested_target.resolve()
    target.mkdir(parents=True, exist_ok=True)
    if not target.is_dir():
        raise RuntimeError(f"Skill target must be a real directory: {target}")

    manifest_warning = ""
    try:
        previous_manifest = read_skill_install_manifest(target)
    except RuntimeError as exc:
        previous_manifest = None
        manifest_warning = str(exc)
    previous_files = set(previous_manifest.get("managed_files", [])) if previous_manifest else set()
    current_files = set(source_files)
    removed_stale_files: list[str] = []
    preserved_stale_files: list[str] = []
    for rel_path in sorted(previous_files - current_files):
        try:
            stale = safe_skill_target_path(target, rel_path)
        except RuntimeError:
            preserved_stale_files.append(rel_path)
            continue
        # A manifest stored inside the writable target cannot prove historical
        # ownership, even when its hashes are internally consistent. Never
        # delete a stale path automatically; report and preserve it for an
        # explicit human cleanup decision.
        if stale.exists():
            preserved_stale_files.append(rel_path)

    for rel_path, source in sorted(source_files.items()):
        destination = safe_skill_target_path(target, rel_path)
        destination.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(source, destination)

    target_files = {
        rel_path: safe_skill_target_path(target, rel_path)
        for rel_path in sorted(source_files)
    }
    source_tree_hash = skill_tree_hash(source_files)
    target_tree_hash = skill_tree_hash(target_files)
    source_file_hashes = {
        rel_path: file_sha256(path) for rel_path, path in sorted(source_files.items())
    }
    manifest = {
        "schema_version": 2,
        "skill": "ad-creative-orchestrator",
        "managed_files": sorted(source_files),
        "managed_file_sha256": source_file_hashes,
        "source_tree_sha256": source_tree_hash,
    }
    write_json_object(target / SKILL_INSTALL_MANIFEST_NAME, manifest)

    source_skill = source_files["SKILL.md"]
    target_skill = target_files["SKILL.md"]
    source_hash = file_sha256(source_skill)
    target_hash = file_sha256(target_skill)
    return {
        "source": str(source_skill),
        "target": str(target_skill),
        "source_hash": source_hash,
        "target_hash": target_hash,
        "source_tree_hash": source_tree_hash,
        "target_tree_hash": target_tree_hash,
        "managed_files": sorted(source_files),
        "manifest": str(target / SKILL_INSTALL_MANIFEST_NAME),
        "manifest_warning": manifest_warning,
        "removed_stale_files": removed_stale_files,
        "preserved_stale_files": preserved_stale_files,
        "match": source_hash == target_hash and source_tree_hash == target_tree_hash,
    }


def package_version() -> str:
    root = source_root()
    pyproject = root / "pyproject.toml" if root else None
    if pyproject and pyproject.exists():
        try:
            import tomllib

            data = tomllib.loads(pyproject.read_text(encoding="utf-8"))
            version = data.get("project", {}).get("version")
            if isinstance(version, str) and version:
                return version
        except Exception:
            pass
    try:
        return metadata.version(PACKAGE_NAME)
    except metadata.PackageNotFoundError:
        pass
    return FALLBACK_VERSION


def module_available(name: str) -> tuple[bool, str]:
    try:
        imported = __import__(name)
        return True, str(getattr(imported, "__version__", "available") or "available")
    except Exception as exc:  # noqa: BLE001 - doctor should report exact import problem
        return False, str(exc)


def git_remote_summary() -> tuple[bool, str]:
    root = source_root()
    if not root or not (root / ".git").exists():
        return False, "not_source_git_checkout"
    try:
        result = subprocess.run(
            ["git", "remote", "-v"],
            cwd=root,
            text=True,
            capture_output=True,
            check=False,
        )
    except Exception as exc:  # noqa: BLE001 - actionable local diagnostic
        return False, f"git_remote_failed={exc}"
    remote_text = result.stdout.strip()
    if not remote_text:
        return False, "empty"
    return True, remote_text.replace("\n", " | ")


def doctor_report() -> tuple[str, list[str], list[str], list[str]]:
    issues: list[str] = []
    warnings: list[str] = []
    evidence: list[str] = [
        f"version={package_version()}",
        f"python={sys.version.split()[0]}",
        f"mode={'source' if source_root() else 'installed'}",
        f"runtime_root={REPO_ROOT}",
        f"template_root={TEMPLATE_ROOT}",
        f"skill_draft_dir={SKILL_DRAFT_DIR}",
    ]

    required_template_files = [
        "AD-creative/orchestrator/project.yml",
        "AD-creative/AGENTS.md",
        "AD-creative/orchestrator/requirements.csv",
        "AD-creative/orchestrator/thread_registry.csv",
        "AD-creative/orchestrator/thread_lane_plan_template.md",
        "AD-creative/orchestrator/agency_staff_selection_template.md",
        "AD-creative/agents/role_briefs/README.md",
        "AD-creative/handoff/项目看板.md",
        "AD-creative/gates/adversarial_council_gate_template.md",
    ]
    if not TEMPLATE_ROOT.exists():
        issues.append(f"template root missing: {TEMPLATE_ROOT}")
    else:
        evidence.append("template_root_exists=true")
        for rel_path in required_template_files:
            if not (TEMPLATE_ROOT / rel_path).exists():
                issues.append(f"template file missing: {rel_path}")

    skill_path = SKILL_DRAFT_DIR / "SKILL.md"
    if not skill_path.exists():
        issues.append(f"skill draft missing: {skill_path}")
    else:
        evidence.append(f"skill_draft={skill_path}")

    for module_name in ["PIL", "pptx"]:
        ok, note = module_available(module_name)
        evidence.append(f"module_{module_name}={note if ok else 'missing'}")
        if not ok:
            warnings.append(f"optional dependency unavailable: {module_name}: {note}")

    has_remote, remote_note = git_remote_summary()
    evidence.append(f"git_remote={remote_note}")
    if not has_remote and remote_note == "empty":
        warnings.append("git remote is not configured; push and GitHub Actions cannot run.")

    status = "PASS" if not issues else "CHECK"
    return status, issues, warnings, evidence


def evidence_map(evidence: list[str]) -> dict[str, str]:
    values: dict[str, str] = {}
    for item in evidence:
        if "=" in item:
            key, value = item.split("=", 1)
            values[key] = value
    return values


def release_status_payload() -> dict[str, object]:
    doctor_status, issues, warnings, evidence = doctor_report()
    values = evidence_map(evidence)
    mode = values.get("mode", "unknown")
    git_remote = values.get("git_remote", "unknown")
    remote_status = "PASS"
    if git_remote == "empty" or git_remote.startswith("git_remote_failed="):
        remote_status = "CHECK"
    elif git_remote == "not_source_git_checkout":
        remote_status = "NOT_APPLICABLE"

    if issues:
        release_status = "CHECK"
        next_action = "Fix doctor issues, then run make release-check."
    elif mode != "source":
        release_status = "INSTALLED_PACKAGE_OK"
        next_action = "Run adco release-status inside the source checkout before publishing."
    elif git_remote == "empty":
        release_status = "BLOCKED_REMOTE_MISSING"
        next_action = "Configure git remote, push the branch, and verify GitHub Actions."
    elif git_remote.startswith("git_remote_failed="):
        release_status = "CHECK"
        next_action = "Fix git remote diagnostics, then run make release-check."
    else:
        release_status = "READY_FOR_REMOTE_CHECKS"
        next_action = "Run make release-check, push the branch, and verify GitHub Actions."

    return {
        "release_status": release_status,
        "doctor_status": doctor_status,
        "remote_status": remote_status,
        "mode": mode,
        "git_remote": git_remote,
        "verify_command": "make release-check",
        "next_action": next_action,
        "issues": issues,
        "warnings": warnings,
        "evidence": evidence,
    }


def docs_payload() -> dict[str, object]:
    root = source_root()
    docs_root = published_docs_root()
    docs: list[dict[str, object]] = []
    for label, rel_path in [
        ("readme", "README.md"),
        ("changelog", "CHANGELOG.md"),
        ("install", "docs/operating/install.md"),
        ("adoption_patterns", "docs/operating/adoption_patterns.md"),
        ("release_plan", "docs/operating/open_source_release_plan.md"),
        ("first_run_transcript", "docs/assets/first-run-transcript.md"),
    ]:
        path = docs_root / rel_path
        docs.append({"label": label, "path": str(path), "exists": path.exists()})
    return {
        "mode": "source" if root else "installed",
        "source_root": str(root) if root else None,
        "docs_root": str(docs_root),
        "template_root": str(TEMPLATE_ROOT),
        "skill_draft": str(SKILL_DRAFT_DIR / "SKILL.md"),
        "docs": docs,
        "quickstart": [
            "adco --version",
            "adco doctor",
            "adco release-status",
            "adco quickstart",
            "adco next /tmp/adco-demo",
            "adco profile-analyze /tmp/adco-demo --brand <brand> --company <company>",
            "adco thread-plan /tmp/adco-demo --title ThreadOps --objective 'Coordinate Codex worker threads'",
            "adco hygiene /tmp/adco-demo",
            "adco open-dashboard /tmp/adco-demo --no-open",
            "adco check",
        ],
    }


def sanitized_doctor_evidence(evidence: list[str]) -> list[str]:
    sanitized: list[str] = []
    for item in evidence:
        if item.startswith("git_remote="):
            value = item.split("=", 1)[1]
            if value not in {"empty", "not_source_git_checkout"} and not value.startswith("git_remote_failed="):
                sanitized.append("git_remote=configured")
            else:
                sanitized.append(item)
        elif item.startswith(("runtime_root=", "template_root=", "skill_draft_dir=", "skill_draft=")):
            key, value = item.split("=", 1)
            sanitized.append(f"{key}={Path(value).name if value else value}")
        else:
            sanitized.append(item)
    return sanitized


def support_table(project: Path, rel_path: str, columns: list[str], limit: int = 5) -> list[str]:
    _, rows = read_csv_rows(project / rel_path)
    lines: list[str] = []
    for row in rows[-limit:]:
        values = [f"{column}={row.get(column, '')}" for column in columns if row.get(column, "")]
        lines.append("- " + " | ".join(values) if values else "- row_present")
    return lines or ["- none"]


def render_support_bundle(project: Path) -> Path:
    counts = read_counts(project)
    errors, validate_stats = validate(project)
    doctor_status, doctor_issues, doctor_warnings, doctor_evidence = doctor_report()
    dashboard = project / DASHBOARD_REL
    council_report = project / COUNCIL_REPORT_REL
    report = project / SUPPORT_BUNDLE_REL
    lines = [
        "# Support Bundle",
        "",
        "content_policy: sanitized diagnostics only; no client brief text, material body, prompt body, or image content included.",
        f"created_at: {now_iso()}",
        f"project_name: {project.name}",
        f"stage: {project_stage(project)}",
        f"validation: {'PASS' if not errors else 'CHECK'}",
        f"doctor: {doctor_status}",
        "",
        "## Environment",
        "",
        *[f"- {item}" for item in sanitized_doctor_evidence(doctor_evidence)],
        "",
        "## Doctor Warnings",
        "",
        *[f"- {warning}" for warning in doctor_warnings],
        *(["- none"] if not doctor_warnings else []),
        "",
        "## Doctor Issues",
        "",
        *[f"- {issue}" for issue in doctor_issues],
        *(["- none"] if not doctor_issues else []),
        "",
        "## Project Counts",
        "",
        *[f"- {key}: {value}" for key, value in counts.items()],
        "",
        "## Validation Counts",
        "",
        *[f"- {key}: {value}" for key, value in validate_stats.items()],
        "",
        "## Validation Errors",
        "",
        *[f"- {error}" for error in errors],
        *(["- none"] if not errors else []),
        "",
        "## Key Files",
        "",
        f"- dashboard: {'exists' if dashboard.exists() else 'missing'}",
        f"- council_report: {'exists' if council_report.exists() else 'missing'}",
        "",
        "## Latest Gates",
        "",
        *support_table(
            project,
            "AD-creative/orchestrator/gate_log.csv",
            ["gate_id", "stage", "status", "next_state", "created_at", "owner"],
        ),
        "",
        "## Latest Work Items",
        "",
        *support_table(
            project,
            "AD-creative/orchestrator/work_items.csv",
            ["work_id", "stage", "status", "owner", "gate_required", "client_visibility"],
        ),
    ]
    write_text(report, sanitize_private_source_markers(project, "\n".join(lines)))
    return report


def check_global_skill(target: Path = DEFAULT_SKILL_INSTALL_DIR) -> dict[str, object]:
    source_skill = SKILL_DRAFT_DIR / "SKILL.md"
    target_root = target.expanduser().resolve()
    target_skill = target_root / "SKILL.md"
    empty = {
        "source": str(source_skill),
        "target": str(target_skill),
        "source_hash": "",
        "target_hash": "",
        "source_tree_hash": "",
        "target_tree_hash": "",
        "managed_files": [],
        "manifest_match": False,
        "match": False,
    }
    try:
        source_files = skill_tree_files(SKILL_DRAFT_DIR)
        target_files = {
            rel_path: safe_skill_target_path(target_root, rel_path)
            for rel_path in sorted(source_files)
        }
    except (FileNotFoundError, RuntimeError):
        return empty
    if any(not path.is_file() for path in target_files.values()):
        return empty
    source_hash = file_sha256(source_files["SKILL.md"])
    target_hash = file_sha256(target_files["SKILL.md"])
    source_tree_hash = skill_tree_hash(source_files)
    target_tree_hash = skill_tree_hash(target_files)
    try:
        manifest = read_skill_install_manifest(target_root)
    except RuntimeError:
        manifest = None
    manifest_match = bool(
        manifest
        and manifest.get("schema_version") == 2
        and manifest.get("managed_files") == sorted(source_files)
        and manifest.get("managed_file_sha256")
        == {
            rel_path: file_sha256(path)
            for rel_path, path in sorted(source_files.items())
        }
        and manifest.get("source_tree_sha256") == source_tree_hash
    )
    return {
        "source": str(source_files["SKILL.md"]),
        "target": str(target_files["SKILL.md"]),
        "source_hash": source_hash,
        "target_hash": target_hash,
        "source_tree_hash": source_tree_hash,
        "target_tree_hash": target_tree_hash,
        "managed_files": sorted(source_files),
        "manifest_match": manifest_match,
        "match": (
            source_hash == target_hash
            and source_tree_hash == target_tree_hash
            and manifest_match
        ),
    }


def ensure_project(project: Path) -> tuple[int, int]:
    if DELIVERY_COMMAND_ACTIVE.get() or project_surface(project) == DELIVERY_SURFACE:
        return ensure_delivery_project(project)
    return copy_content_template(TEMPLATE_ROOT, project)


def ensure_delivery_project(project: Path) -> tuple[int, int]:
    """Materialize the full governance surface only for a delivery-risk command."""
    return copy_template(TEMPLATE_ROOT, project)


def read_csv_rows(path: Path) -> tuple[list[str], list[dict[str, str]]]:
    if not path.exists():
        return [], []
    with path.open(newline="", encoding="utf-8") as handle:
        reader = csv.DictReader(handle)
        return list(reader.fieldnames or []), list(reader)


def write_csv_rows(path: Path, fieldnames: list[str], rows: Iterable[dict[str, str]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temp_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            "w",
            newline="",
            encoding="utf-8",
            dir=path.parent,
            prefix=f".{path.name}.",
            suffix=".tmp",
            delete=False,
        ) as handle:
            temp_path = Path(handle.name)
            writer = csv.DictWriter(handle, fieldnames=fieldnames, lineterminator="\n")
            writer.writeheader()
            for row in rows:
                writer.writerow({key: row.get(key, "") for key in fieldnames})
        os.replace(temp_path, path)
    finally:
        if temp_path is not None:
            temp_path.unlink(missing_ok=True)


def append_csv_row(path: Path, row: dict[str, str]) -> None:
    fieldnames, rows = read_csv_rows(path)
    if not fieldnames:
        raise FileNotFoundError(f"CSV header not found: {path}")
    rows.append(row)
    write_csv_rows(path, fieldnames, rows)


def csv_rows_need_normalization(
    fieldnames: list[str], rows: Iterable[dict[str, str]]
) -> bool:
    return any(
        None in row or any(row.get(field) is None for field in fieldnames)
        for row in rows
    )


def ensure_csv_fields(path: Path, required_fields: list[str]) -> list[str]:
    fieldnames, rows = read_csv_rows(path)
    if not fieldnames:
        raise FileNotFoundError(f"CSV header not found: {path}")
    missing = [field for field in required_fields if field not in fieldnames]
    if missing or csv_rows_need_normalization(fieldnames, rows):
        fieldnames = [*fieldnames, *missing]
        write_csv_rows(path, fieldnames, rows)
    return fieldnames


def normalized_gate_log_data(
    fieldnames: list[str],
    rows: list[dict[str, str]],
) -> tuple[list[str], list[dict[str, str]]]:
    if not fieldnames:
        return GATE_LOG_FIELDS, []
    normalized_rows: list[dict[str, str]] = []
    previous_by_gate: dict[str, str] = {}
    used_run_ids: set[str] = set()
    for index, row in enumerate(rows, start=1):
        gate_id = row.get("gate_id", "").strip()
        run_id = row.get("gate_run_id", "").strip()
        if not run_id or run_id in used_run_ids:
            run_id = f"GATE-RUN-{index:04d}"
            while run_id in used_run_ids:
                index += 1
                run_id = f"GATE-RUN-{index:04d}"
        used_run_ids.add(run_id)
        normalized = {field: row.get(field, "") or "" for field in GATE_LOG_FIELDS}
        normalized["gate_run_id"] = run_id
        normalized["supersedes_gate_run_id"] = (
            row.get("supersedes_gate_run_id", "").strip()
            or previous_by_gate.get(gate_id, "")
        )
        normalized_rows.append(normalized)
        if gate_id:
            previous_by_gate[gate_id] = run_id
    return GATE_LOG_FIELDS, normalized_rows


def normalize_gate_log_schema(path: Path) -> tuple[list[str], list[dict[str, str]]]:
    fieldnames, rows = read_csv_rows(path)
    normalized_fields, normalized_rows = normalized_gate_log_data(fieldnames, rows)
    if fieldnames != GATE_LOG_FIELDS or rows != normalized_rows:
        write_csv_rows(path, normalized_fields, normalized_rows)
    return normalized_fields, normalized_rows


def update_or_append_csv_row(
    path: Path, key: str, row: dict[str, str], *, replace: bool = True
) -> None:
    fieldnames, rows = read_csv_rows(path)
    if not fieldnames:
        raise FileNotFoundError(f"CSV header not found: {path}")
    for index, existing in enumerate(rows):
        if existing.get(key) == row.get(key):
            if replace:
                rows[index] = row
            break
    else:
        rows.append(row)
    write_csv_rows(path, fieldnames, rows)


def next_id(rows: list[dict[str, str]], key: str, prefix: str) -> str:
    highest = 0
    for row in rows:
        match = ID_SUFFIX_PATTERN.search(row.get(key, ""))
        if match:
            highest = max(highest, int(match.group(1)))
    return f"{prefix}-{highest + 1:03d}"


def id_allocator(rows: list[dict[str, str]], key: str, prefix: str) -> Callable[[], str]:
    highest = 0
    for row in rows:
        match = ID_SUFFIX_PATTERN.search(row.get(key, ""))
        if match:
            highest = max(highest, int(match.group(1)))

    def allocate() -> str:
        nonlocal highest
        highest += 1
        return f"{prefix}-{highest:03d}"

    return allocate


def safe_rel(project: Path, path: Path) -> str:
    try:
        return str(path.resolve().relative_to(project.resolve()))
    except ValueError:
        return str(path.resolve())


def safe_rel_to(root: Path, path: Path) -> str:
    try:
        return str(path.resolve().relative_to(root.resolve()))
    except ValueError:
        return str(path.resolve())


def is_private_local_state_path(project: Path, path: Path) -> bool:
    try:
        relative = path.resolve().relative_to(project.resolve())
    except ValueError:
        return False
    return bool(relative.parts) and relative.parts[0] == PRIVATE_LOCAL_STATE_REL.name


def _safe_project_relative_parts(relative_path: str | Path) -> tuple[str, ...]:
    candidate = Path(relative_path)
    if (
        candidate.is_absolute()
        or not candidate.parts
        or any(part in {"", ".", ".."} for part in candidate.parts)
    ):
        raise ValueError(f"unsafe project-relative path: {relative_path}")
    return tuple(candidate.parts)


def _close_project_fd_chain(directory_fds: Iterable[int]) -> None:
    for fd in reversed(list(directory_fds)):
        os.close(fd)


def _open_project_root_fd(project: Path) -> int:
    no_follow = getattr(os, "O_NOFOLLOW", None)
    directory_flag = getattr(os, "O_DIRECTORY", None)
    if no_follow is None or directory_flag is None:
        raise ValueError("safe project access requires O_NOFOLLOW and O_DIRECTORY")
    try:
        fd = os.open(project, os.O_RDONLY | directory_flag | no_follow)
    except FileNotFoundError:
        raise
    except OSError as exc:
        raise ValueError(f"cannot safely open project root: {project}") from exc
    try:
        opened = os.fstat(fd)
        visible = os.stat(project, follow_symlinks=False)
        if (
            not stat.S_ISDIR(visible.st_mode)
            or (opened.st_dev, opened.st_ino)
            != (visible.st_dev, visible.st_ino)
        ):
            raise ValueError("project root changed while opening")
        return fd
    except BaseException:
        os.close(fd)
        raise


def _project_root_fd_is_current(project: Path, root_fd: int) -> bool:
    try:
        visible = os.stat(project, follow_symlinks=False)
        opened = os.fstat(root_fd)
    except OSError:
        return False
    return bool(
        stat.S_ISDIR(visible.st_mode)
        and (visible.st_dev, visible.st_ino) == (opened.st_dev, opened.st_ino)
    )


def _open_project_parent_chain(
    project: Path,
    relative_path: str | Path,
    *,
    create: bool = False,
    project_root_fd: int | None = None,
) -> tuple[list[int], tuple[str, ...]]:
    parts = _safe_project_relative_parts(relative_path)
    no_follow = getattr(os, "O_NOFOLLOW", None)
    directory_flag = getattr(os, "O_DIRECTORY", None)
    if no_follow is None or directory_flag is None:
        raise ValueError("safe project access requires O_NOFOLLOW and O_DIRECTORY")
    if project_root_fd is None:
        root_fd = _open_project_root_fd(project)
    else:
        root_fd = os.dup(project_root_fd)
        if not _project_root_fd_is_current(project, root_fd):
            os.close(root_fd)
            raise ValueError("project root changed during bound operation")
    directory_fds = [root_fd]
    try:
        for part in parts[:-1]:
            parent_fd = directory_fds[-1]
            if create:
                try:
                    os.mkdir(part, mode=0o755, dir_fd=parent_fd)
                except FileExistsError:
                    pass
            try:
                visible = os.stat(
                    part,
                    dir_fd=parent_fd,
                    follow_symlinks=False,
                )
                child_fd = os.open(
                    part,
                    os.O_RDONLY | directory_flag | no_follow,
                    dir_fd=parent_fd,
                )
            except FileNotFoundError:
                raise
            except OSError as exc:
                raise ValueError(
                    f"unsafe project path component: {'/'.join(parts[:-1])}"
                ) from exc
            opened = os.fstat(child_fd)
            if (
                not stat.S_ISDIR(visible.st_mode)
                or (visible.st_dev, visible.st_ino)
                != (opened.st_dev, opened.st_ino)
            ):
                os.close(child_fd)
                raise ValueError(
                    f"project path component changed: {'/'.join(parts[:-1])}"
                )
            directory_fds.append(child_fd)
        return directory_fds, parts
    except BaseException:
        _close_project_fd_chain(directory_fds)
        raise


def _project_parent_chain_is_current(
    project: Path,
    directory_fds: list[int],
    parts: tuple[str, ...],
) -> bool:
    try:
        root_visible = os.stat(project, follow_symlinks=False)
        root_opened = os.fstat(directory_fds[0])
        if (
            not stat.S_ISDIR(root_visible.st_mode)
            or (root_visible.st_dev, root_visible.st_ino)
            != (root_opened.st_dev, root_opened.st_ino)
        ):
            return False
        for index, part in enumerate(parts[:-1]):
            visible = os.stat(
                part,
                dir_fd=directory_fds[index],
                follow_symlinks=False,
            )
            opened = os.fstat(directory_fds[index + 1])
            if (
                not stat.S_ISDIR(visible.st_mode)
                or (visible.st_dev, visible.st_ino)
                != (opened.st_dev, opened.st_ino)
            ):
                return False
    except OSError:
        return False
    return True


def _open_project_relative_regular_file(
    project: Path,
    relative_path: str | Path,
    *,
    project_root_fd: int | None = None,
) -> tuple[int, os.stat_result, list[int], tuple[str, ...]]:
    directory_fds, parts = _open_project_parent_chain(
        project,
        relative_path,
        project_root_fd=project_root_fd,
    )
    file_fd: int | None = None
    try:
        parent_fd = directory_fds[-1]
        visible = os.stat(
            parts[-1],
            dir_fd=parent_fd,
            follow_symlinks=False,
        )
        if not stat.S_ISREG(visible.st_mode):
            raise ValueError(f"project file is not regular: {'/'.join(parts)}")
        file_fd = os.open(
            parts[-1],
            os.O_RDONLY | getattr(os, "O_NOFOLLOW", 0),
            dir_fd=parent_fd,
        )
        opened = os.fstat(file_fd)
        if (
            not stat.S_ISREG(opened.st_mode)
            or (visible.st_dev, visible.st_ino)
            != (opened.st_dev, opened.st_ino)
        ):
            raise ValueError(f"project file changed while opening: {'/'.join(parts)}")
        result = file_fd
        file_fd = None
        return result, opened, directory_fds, parts
    except BaseException:
        if file_fd is not None:
            os.close(file_fd)
        _close_project_fd_chain(directory_fds)
        raise


def _project_file_binding_is_current(
    project: Path,
    file_fd: int,
    opened: os.stat_result,
    directory_fds: list[int],
    parts: tuple[str, ...],
) -> bool:
    if not _project_parent_chain_is_current(project, directory_fds, parts):
        return False
    try:
        final_opened = os.fstat(file_fd)
        final_visible = os.stat(
            parts[-1],
            dir_fd=directory_fds[-1],
            follow_symlinks=False,
        )
    except OSError:
        return False
    return bool(
        stat.S_ISREG(final_visible.st_mode)
        and (final_visible.st_dev, final_visible.st_ino)
        == (opened.st_dev, opened.st_ino)
        and (final_opened.st_dev, final_opened.st_ino)
        == (opened.st_dev, opened.st_ino)
        and (
            final_opened.st_size,
            final_opened.st_mtime_ns,
            final_opened.st_ctime_ns,
        )
        == (opened.st_size, opened.st_mtime_ns, opened.st_ctime_ns)
    )


def private_source_path_markers(
    project: Path,
    *,
    project_root_fd: int | None = None,
) -> list[bytes]:
    try:
        sources = (
            load_local_source_paths(project)
            if project_root_fd is None
            else load_local_source_paths_from_project_fd(project_root_fd)
        )
    except ValueError as exc:
        raise ValueError(
            "private local source map is invalid or unreadable"
        ) from exc
    referenced_ids: set[str] = set()
    for rel_path in [
        "AD-creative/orchestrator/source_events.csv",
        "AD-creative/orchestrator/evidence_chunks.jsonl",
    ]:
        try:
            text = _read_project_text(
                project,
                rel_path,
                project_root_fd=project_root_fd,
            )
        except FileNotFoundError:
            continue
        except (OSError, UnicodeError, ValueError) as exc:
            raise ValueError(
                "private local source references are invalid or unreadable"
            ) from exc
        for line in text.splitlines():
            referenced_ids.update(
                match.group(1)
                for match in re.finditer(
                    rf"{re.escape(LOCAL_SOURCE_PREFIX)}([A-Za-z0-9._-]+)",
                    line,
                )
            )
    if referenced_ids - set(sources):
        raise ValueError("private local source map is missing registered aliases")
    return sorted(
        {
            value.encode("utf-8")
            for value in sources.values()
            if isinstance(value, str) and value
        },
        key=len,
        reverse=True,
    )


def _scan_marker_stream(
    handle: BinaryIO,
    markers: list[bytes],
    *,
    limit: int,
) -> tuple[bool, int, bool]:
    overlap = max(len(marker) for marker in markers) - 1
    tail = b""
    total = 0
    while True:
        remaining = limit - total
        if remaining < 0:
            return False, total, True
        chunk = handle.read(min(PRIVATE_MARKER_SCAN_CHUNK_BYTES, remaining + 1))
        if not chunk:
            return False, total, False
        total += len(chunk)
        if total > limit:
            return False, total, True
        data = tail + chunk
        if any(marker in data for marker in markers):
            return True, total, False
        tail = data[-overlap:] if overlap > 0 else b""


def _scan_text_marker_stream(
    handle: BinaryIO,
    markers: list[bytes],
    *,
    limit: int,
) -> tuple[bool, int, bool]:
    normalized_markers = [
        re.sub(rb"\s+", b"", marker) for marker in markers if marker
    ]
    overlap = max((len(marker) for marker in normalized_markers), default=1) - 1
    tail = b""
    total = 0
    while True:
        remaining = limit - total
        if remaining < 0:
            return False, total, True
        chunk = handle.read(min(PRIVATE_MARKER_SCAN_CHUNK_BYTES, remaining + 1))
        if not chunk:
            return False, total, False
        total += len(chunk)
        if total > limit:
            return False, total, True
        normalized = re.sub(rb"\s+", b"", chunk)
        data = tail + normalized
        if any(marker in data for marker in normalized_markers):
            return True, total, False
        tail = data[-overlap:] if overlap > 0 else b""


def _hash_open_fd(fd: int, *, limit: int) -> tuple[str, int, bool]:
    os.lseek(fd, 0, os.SEEK_SET)
    digest = hashlib.sha256()
    total = 0
    while True:
        chunk = os.read(fd, PRIVATE_MARKER_SCAN_CHUNK_BYTES)
        if not chunk:
            return digest.hexdigest(), total, False
        total += len(chunk)
        if total > limit:
            return "", total, True
        digest.update(chunk)


def _scan_and_hash_regular_fd(
    fd: int,
    markers: list[bytes],
    *,
    limit: int,
) -> tuple[str, int, bool, bool]:
    os.lseek(fd, 0, os.SEEK_SET)
    digest = hashlib.sha256()
    overlap = max((len(marker) for marker in markers), default=1) - 1
    tail = b""
    total = 0
    found = False
    while True:
        chunk = os.read(fd, PRIVATE_MARKER_SCAN_CHUNK_BYTES)
        if not chunk:
            return digest.hexdigest(), total, found, False
        total += len(chunk)
        if total > limit:
            return "", total, found, True
        digest.update(chunk)
        if markers:
            data = tail + chunk
            if any(marker in data for marker in markers):
                found = True
            tail = data[-overlap:] if overlap > 0 else b""


def _scan_zip_fd(fd: int, markers: list[bytes]) -> tuple[bool, str]:
    duplicate = os.dup(fd)
    try:
        os.lseek(duplicate, 0, os.SEEK_SET)
        with os.fdopen(duplicate, "rb", closefd=True) as handle:
            duplicate = -1
            if not zipfile.is_zipfile(handle):
                return False, ""
            handle.seek(0)
            with zipfile.ZipFile(handle) as archive:
                infos = archive.infolist()
                if len(infos) > PRIVATE_MARKER_SCAN_MAX_ARCHIVE_MEMBERS:
                    return True, "privacy scan archive member limit exceeded"
                archive_total = 0
                for info in infos:
                    if info.is_dir():
                        continue
                    if (
                        info.file_size < 0
                        or info.file_size > PRIVATE_MARKER_SCAN_MAX_MEMBER_BYTES
                        or archive_total + info.file_size
                        > PRIVATE_MARKER_SCAN_MAX_ARCHIVE_BYTES
                    ):
                        return True, "privacy scan archive limit exceeded"
                    remaining_archive = (
                        PRIVATE_MARKER_SCAN_MAX_ARCHIVE_BYTES - archive_total
                    )
                    member_limit = min(
                        PRIVATE_MARKER_SCAN_MAX_MEMBER_BYTES,
                        remaining_archive,
                    )
                    with archive.open(info) as member:
                        found, scanned, exceeded = _scan_marker_stream(
                            member,
                            markers,
                            limit=member_limit,
                        )
                    archive_total += scanned
                    if exceeded:
                        return True, "privacy scan archive limit exceeded"
                    if found:
                        return True, "private local source path marker detected"
            return True, ""
    finally:
        if duplicate >= 0:
            os.close(duplicate)


def _scan_pdf_command_fd(
    fd: int,
    markers: list[bytes],
    command: list[str],
) -> tuple[str, bool]:
    duplicate = os.dup(fd)
    process: subprocess.Popen[bytes] | None = None
    selector: selectors.BaseSelector | None = None
    try:
        os.lseek(duplicate, 0, os.SEEK_SET)
        with os.fdopen(duplicate, "rb", closefd=True) as pdf_input:
            duplicate = -1
            try:
                process = subprocess.Popen(
                    command,
                    stdin=pdf_input,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.DEVNULL,
                )
            except OSError:
                return "privacy scan PDF parser failed", False
            if process.stdout is None:
                return "privacy scan PDF parser failed", False
            output_fd = process.stdout.fileno()
            os.set_blocking(output_fd, False)
            selector = selectors.DefaultSelector()
            selector.register(output_fd, selectors.EVENT_READ)
            normalized_markers = [
                re.sub(rb"\s+", b"", marker) for marker in markers if marker
            ]
            overlap = max(
                (len(marker) for marker in normalized_markers),
                default=1,
            ) - 1
            tail = b""
            total = 0
            deadline = time.monotonic() + PRIVATE_MARKER_SCAN_PDF_TIMEOUT_SECONDS
            reached_eof = False
            while not reached_eof:
                remaining_time = deadline - time.monotonic()
                if remaining_time <= 0:
                    return "privacy scan PDF parser timed out", False
                events = selector.select(timeout=min(remaining_time, 0.1))
                if not events:
                    continue
                for _key, _mask in events:
                    try:
                        chunk = os.read(
                            output_fd,
                            min(
                                PRIVATE_MARKER_SCAN_CHUNK_BYTES,
                                PRIVATE_MARKER_SCAN_MAX_ARCHIVE_BYTES - total + 1,
                            ),
                        )
                    except BlockingIOError:
                        continue
                    if not chunk:
                        reached_eof = True
                        break
                    total += len(chunk)
                    if total > PRIVATE_MARKER_SCAN_MAX_ARCHIVE_BYTES:
                        return "privacy scan PDF text limit exceeded", False
                    normalized = re.sub(rb"\s+", b"", chunk)
                    data = tail + normalized
                    if any(marker in data for marker in normalized_markers):
                        return "", True
                    tail = data[-overlap:] if overlap > 0 else b""
            try:
                return_code = process.wait(timeout=1)
            except subprocess.TimeoutExpired:
                return "privacy scan PDF parser timed out", False
            if return_code != 0:
                return "privacy scan PDF parser failed", False
            return "", False
    finally:
        if selector is not None:
            selector.close()
        if process is not None and process.poll() is None:
            process.terminate()
            try:
                process.wait(timeout=1)
            except subprocess.TimeoutExpired:
                process.kill()
                process.wait()
        if process is not None and process.stdout is not None:
            process.stdout.close()
        if duplicate >= 0:
            os.close(duplicate)


def _scan_pdf_with_pypdf_fd(fd: int, markers: list[bytes]) -> tuple[str, bool]:
    if importlib.util.find_spec("pypdf") is None:
        return "privacy scan PDF parser unavailable", False
    script = "\n".join(
        [
            "import io, sys",
            "try:",
            "    import resource",
            "except ImportError:",
            "    raise SystemExit(78)",
            f"memory_limit = {PRIVATE_MARKER_SCAN_MAX_PDF_MEMORY_BYTES}",
            "limit_kind = getattr(resource, 'RLIMIT_AS', getattr(resource, 'RLIMIT_DATA', None))",
            "if limit_kind is None:",
            "    raise SystemExit(78)",
            "soft, hard = resource.getrlimit(limit_kind)",
            "bounded = memory_limit if hard == resource.RLIM_INFINITY else min(memory_limit, hard)",
            "resource.setrlimit(limit_kind, (bounded, hard))",
            "from pypdf import PdfReader",
            f"data = sys.stdin.buffer.read({PRIVATE_MARKER_SCAN_MAX_FILE_BYTES + 1})",
            f"if len(data) > {PRIVATE_MARKER_SCAN_MAX_FILE_BYTES}:",
            "    raise SystemExit(79)",
            "reader = PdfReader(io.BytesIO(data), strict=False)",
            f"if len(reader.pages) > {PRIVATE_MARKER_SCAN_MAX_PDF_PAGES}:",
            "    raise SystemExit(80)",
            "output = sys.stdout.buffer",
            "def emit(value):",
            "    output.write(str(value).encode('utf-8', errors='replace'))",
            "    output.write(b'\\n')",
            "metadata = reader.metadata",
            "if metadata:",
            "    for value in metadata.values():",
            "        emit(value)",
            "for page in reader.pages:",
            "    page.extract_text(visitor_text=lambda text, *_: emit(text))",
        ]
    )
    return _scan_pdf_command_fd(
        fd,
        markers,
        [sys.executable, "-I", "-c", script],
    )


def _scan_pdf_fd(fd: int, markers: list[bytes]) -> tuple[str, bool]:
    extractor = shutil.which("pdftotext")
    if extractor:
        issue, found = _scan_pdf_command_fd(
            fd,
            markers,
            [extractor, "-layout", "-", "-"],
        )
    else:
        issue, found = _scan_pdf_with_pypdf_fd(fd, markers)
    if issue or found:
        return issue, found

    metadata_extractor = shutil.which("pdfinfo")
    if metadata_extractor:
        return _scan_pdf_command_fd(
            fd,
            markers,
            [metadata_extractor, "-custom", "-"],
        )
    try:
        from pypdf import PdfReader  # noqa: F401
    except ImportError:
        return "privacy scan PDF metadata parser unavailable", False
    return _scan_pdf_with_pypdf_fd(fd, markers)


def private_source_candidate_evidence(
    project: Path,
    relative_path: str,
    markers: list[bytes],
    *,
    project_root_fd: int | None = None,
) -> tuple[str, str, int]:
    no_follow = getattr(os, "O_NOFOLLOW", None)
    if no_follow is None:
        return "privacy scan requires O_NOFOLLOW", "", 0
    fd: int | None = None
    directory_fds: list[int] = []
    try:
        fd, opened, directory_fds, parts = _open_project_relative_regular_file(
            project,
            relative_path,
            project_root_fd=project_root_fd,
        )
        if opened.st_size > PRIVATE_MARKER_SCAN_MAX_FILE_BYTES:
            return "privacy scan file limit exceeded", "", 0

        is_zip, zip_issue = _scan_zip_fd(fd, markers) if markers else (False, "")
        if zip_issue:
            return zip_issue, "", 0
        if is_zip:
            digest, bytes_read, exceeded = _hash_open_fd(
                fd,
                limit=PRIVATE_MARKER_SCAN_MAX_FILE_BYTES,
            )
            found = False
        else:
            digest, bytes_read, found, exceeded = _scan_and_hash_regular_fd(
                fd,
                markers,
                limit=PRIVATE_MARKER_SCAN_MAX_FILE_BYTES,
            )
        if exceeded:
            return "privacy scan file limit exceeded", "", 0
        if found:
            return "private local source path marker detected", "", 0
        if Path(relative_path).suffix.lower() == ".pdf" and markers:
            pdf_issue, pdf_found = _scan_pdf_fd(fd, markers)
            if pdf_issue:
                return pdf_issue, "", 0
            if pdf_found:
                return "private local source path marker detected in PDF", "", 0
        if bytes_read != opened.st_size or not _project_file_binding_is_current(
            project,
            fd,
            opened,
            directory_fds,
            parts,
        ):
            return "privacy scan target changed during inspection", "", 0
        return "", digest, opened.st_size
    except (
        EOFError,
        NotImplementedError,
        OSError,
        RuntimeError,
        ValueError,
        zipfile.BadZipFile,
        zipfile.LargeZipFile,
    ):
        return "privacy scan failed", "", 0
    finally:
        if fd is not None:
            os.close(fd)
        _close_project_fd_chain(directory_fds)


def private_source_marker_scan_issue(
    project: Path,
    relative_path: str,
    markers: list[bytes],
) -> str:
    issue, _, _ = private_source_candidate_evidence(
        project,
        relative_path,
        markers,
    )
    return issue


def file_contains_private_source_marker(
    project: Path,
    relative_path: str,
    markers: list[bytes],
) -> bool:
    return bool(private_source_marker_scan_issue(project, relative_path, markers))


def sanitize_private_source_markers(project: Path, text: str) -> str:
    sanitized = text
    for marker in private_source_path_markers(project):
        sanitized = sanitized.replace(marker.decode("utf-8"), "[private-local-source]")
    return sanitized


def codex_generated_images_root() -> Path:
    codex_home = Path(os.environ.get("CODEX_HOME", str(Path.home() / ".codex"))).expanduser()
    return codex_home / "generated_images"


def classify_material(path: Path) -> str:
    marker = str(path).lower()
    if "feedback" in marker or "客户反馈" in marker:
        return "feedback"
    if "director" in marker or "导演" in marker:
        return "director_note"
    if "meeting" in marker or "会议" in marker:
        return "meeting_note"
    if "approval" in marker or "确认" in marker:
        return "approval"
    if "rejection" in marker or "拒绝" in marker:
        return "rejection"
    if "change" in marker or "变更" in marker:
        return "change"
    return "initial"


def append_event(project: Path, payload: dict[str, object]) -> None:
    path = project / "AD-creative/orchestrator/events.jsonl"
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("a", encoding="utf-8") as handle:
        handle.write(json.dumps(payload, ensure_ascii=False, sort_keys=True) + "\n")


def register_materials(project: Path, material_paths: list[Path], goal: str) -> list[str]:
    source_path = project / "AD-creative/orchestrator/source_events.csv"
    _, rows = read_csv_rows(source_path)
    source_ids: list[str] = []
    allocate_source_id = id_allocator(rows, "source_event_id", "SRC")
    delivery_surface = project_surface(project) == DELIVERY_SURFACE

    for material in material_paths:
        if not material.exists():
            raise FileNotFoundError(f"material not found: {material}")
        source_id = allocate_source_id()
        try:
            material.resolve().relative_to(project.resolve())
            public_material_path = safe_rel(project, material)
        except ValueError:
            public_material_path = register_local_source_path(
                project,
                source_id,
                material,
            )
        row = {
            "source_event_id": source_id,
            "received_at": now_iso(),
            "source_owner": "operator",
            "source_type": "folder" if material.is_dir() else "file",
            "declared_semantics": classify_material(material),
            "file_paths": public_material_path,
            "raw_summary": f"待整理资料：{material.name}",
            "trust_level": "unreviewed",
            "affects_requirements": "unknown",
            "affects_artifacts": "",
            "supersedes_event_ids": "",
            "notes": goal,
        }
        rows.append(row)
        source_ids.append(source_id)
        if delivery_surface:
            append_event(
                project,
                {
                    "event_id": f"EVT-{source_id}",
                    "event_type": "material_registered",
                    "created_at": row["received_at"],
                    "source_event_id": source_id,
                    "material": row["file_paths"],
                    "goal": goal,
                },
            )

    if source_ids:
        fieldnames, _ = read_csv_rows(source_path)
        write_csv_rows(source_path, fieldnames, rows)
    return source_ids


def existing_source_ids_for_material(project: Path, material: Path) -> list[str]:
    _, rows = read_csv_rows(project / "AD-creative/orchestrator/source_events.csv")
    material_resolved = material.resolve()
    matches: list[str] = []
    for row in rows:
        try:
            roots = source_row_material_roots(project, row)
        except ValueError:
            continue
        if any(path.resolve() == material_resolved for path in roots):
            source_id = row.get("source_event_id", "")
            if source_id:
                matches.append(source_id)
    return matches


def write_sample_brief(project: Path, force: bool = False) -> tuple[Path, str]:
    material = project / SAMPLE_MATERIAL_REL
    if material.exists() and not force:
        return material, "existing"
    write_text(material, SAMPLE_BRIEF)
    return material, "overwritten" if force else "created"


def ensure_intake_work(project: Path, source_ids: list[str], goal: str) -> str:
    """Preserve the historical intake work item on the Delivery Surface."""
    work_path = project / "AD-creative/orchestrator/work_items.csv"
    fieldnames, rows = read_csv_rows(work_path)
    for row in rows:
        if row.get("stage") == "intake" and row.get("title") == "需求整理与缺口判断":
            return row.get("work_id", "")

    work_id = next_id(rows, "work_id", "WORK")
    rows.append(
        {
            "work_id": work_id,
            "stage": "intake",
            "title": "需求整理与缺口判断",
            "objective": goal or "整理客户资料，输出需求、缺口、追问和下一步建议。",
            "owner_agent": "Codex",
            "status": "ready",
            "priority": "high",
            "input_refs": ";".join(source_ids),
            "output_artifacts": "",
            "linked_requirements": "",
            "linked_source_events": ";".join(source_ids),
            "linked_references": "",
            "linked_assets": "",
            "linked_slides": "",
            "blocked_by": "",
            "gate_required": "Brief Gate",
            "client_visibility": "internal_only",
            "created_at": now_iso(),
            "updated_at": now_iso(),
            "supersedes_work_id": "",
        }
    )
    write_csv_rows(work_path, fieldnames, rows)
    return work_id


def ensure_profile_work(project: Path, source_ids: list[str], goal: str) -> str:
    if project_surface(project) != DELIVERY_SURFACE:
        return ""
    work_path = project / "AD-creative/orchestrator/work_items.csv"
    fieldnames, rows = read_csv_rows(work_path)
    for row in rows:
        if row.get("stage") == "intake" and row.get("title") == "会议画像与品牌画像分析":
            row["linked_source_events"] = join_unique_values(row.get("linked_source_events", ""), ";".join(source_ids))
            row["updated_at"] = now_iso()
            write_csv_rows(work_path, fieldnames, rows)
            return row.get("work_id", "")

    work_id = next_id(rows, "work_id", "WORK")
    rows.append(
        {
            "work_id": work_id,
            "stage": "intake",
            "title": "会议画像与品牌画像分析",
            "objective": goal or "分析会议资料中的人物画像、品牌画像、需求权重、决策权和分歧融合路径。",
            "owner_agent": "Codex",
            "status": "ready",
            "priority": "high",
            "input_refs": ";".join(source_ids),
            "output_artifacts": "",
            "linked_requirements": "",
            "linked_source_events": ";".join(source_ids),
            "linked_references": "",
            "linked_assets": "",
            "linked_slides": "",
            "blocked_by": "",
            "gate_required": "Profile Gate",
            "client_visibility": "internal_only",
            "created_at": now_iso(),
            "updated_at": now_iso(),
            "supersedes_work_id": "",
        }
    )
    write_csv_rows(work_path, fieldnames, rows)
    return work_id


def ensure_creative_proposal_work(
    project: Path,
    work_id: str,
    source_ids: str,
    objective: str,
) -> str:
    """Maintain the legacy creative work binding on the Delivery Surface."""
    brief_outputs = ";".join(
        [
            "ART-AUTO-CREATIVE-BRIEF-SNAPSHOT",
            "ART-AUTO-CREATIVE-BRIEF-CONTRACT",
            "ART-AUTO-CREATIVE-CANDIDATE-SCHEMA",
            "ART-AUTO-CREATIVE-GENERATION-REQUEST",
            "ART-AUTO-CREATIVE-OPEN-GAPS",
        ]
    )
    work_path = project / "AD-creative/orchestrator/work_items.csv"
    fieldnames, rows = read_csv_rows(work_path)
    if not fieldnames:
        raise FileNotFoundError(f"CSV header not found: {work_path}")
    for row in rows:
        if work_id and row.get("work_id") == work_id:
            row["linked_source_events"] = join_unique_values(
                row.get("linked_source_events", ""), source_ids
            )
            row["output_artifacts"] = join_unique_values(
                row.get("output_artifacts", ""), brief_outputs
            )
            row["updated_at"] = now_iso()
            write_csv_rows(work_path, fieldnames, rows)
            return work_id
        if not work_id and row.get("stage") == "creative" and row.get("title") in {
            "内部创意提案草案",
            "证据化创意 Brief",
        }:
            row["linked_source_events"] = join_unique_values(
                row.get("linked_source_events", ""), source_ids
            )
            row["title"] = "证据化创意 Brief"
            row["output_artifacts"] = brief_outputs
            row["updated_at"] = now_iso()
            write_csv_rows(work_path, fieldnames, rows)
            return row.get("work_id", "")
    new_work_id = work_id or next_id(rows, "work_id", "WORK")
    rows.append(
        {
            "work_id": new_work_id,
            "stage": "creative",
            "title": "证据化创意 Brief",
            "objective": objective or "生成证据快照、创意 Brief 合同与候选生成请求。",
            "owner_agent": "Codex",
            "status": "ready",
            "priority": "high",
            "input_refs": source_ids,
            "output_artifacts": brief_outputs,
            "linked_requirements": "",
            "linked_source_events": source_ids,
            "linked_references": "",
            "linked_assets": "",
            "linked_slides": "",
            "blocked_by": "",
            "gate_required": "Creative Quality Gate",
            "client_visibility": "internal_only",
            "created_at": now_iso(),
            "updated_at": now_iso(),
            "supersedes_work_id": "",
        }
    )
    write_csv_rows(work_path, fieldnames, rows)
    return new_work_id


def read_counts(project: Path) -> dict[str, int]:
    tables = {
        "source_events": "AD-creative/orchestrator/source_events.csv",
        "requirements": "AD-creative/orchestrator/requirements.csv",
        "gaps": "AD-creative/orchestrator/gaps.csv",
        "work_items": "AD-creative/orchestrator/work_items.csv",
        "agent_runs": "AD-creative/orchestrator/agent_runs.csv",
        "artifacts": "AD-creative/orchestrator/artifact_index.csv",
        "gates": "AD-creative/orchestrator/gate_log.csv",
        "references": "AD-creative/references/reference_cards.csv",
        "assets": "AD-creative/visual_assets/asset_manifest.csv",
        "profile_subjects": "AD-creative/orchestrator/profile_knowledge/profile_subjects.csv",
        "profile_insights": "AD-creative/orchestrator/profile_knowledge/profile_insights.csv",
        "profile_conflicts": "AD-creative/orchestrator/profile_knowledge/profile_conflicts.csv",
    }
    counts: dict[str, int] = {}
    for key, rel_path in tables.items():
        _, rows = read_csv_rows(project / rel_path)
        counts[key] = len(rows)
    return counts


def write_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temp_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            "w",
            encoding="utf-8",
            dir=path.parent,
            prefix=f".{path.name}.",
            suffix=".tmp",
            delete=False,
        ) as handle:
            temp_path = Path(handle.name)
            handle.write(content.rstrip() + "\n")
        os.replace(temp_path, path)
    finally:
        if temp_path is not None:
            temp_path.unlink(missing_ok=True)


def update_markdown_sections(path: Path, sections: dict[str, str]) -> None:
    """Update owned sections while preserving version truth and user-added sections."""
    text = path.read_text(encoding="utf-8") if path.exists() else "# Current Truth\n"
    for heading, body in sections.items():
        section = f"## {heading}\n{body.strip()}\n"
        pattern = re.compile(
            rf"(?ms)^## {re.escape(heading)}[ \t]*\n.*?(?=^## [^\n]+\n|\Z)"
        )
        if pattern.search(text):
            text = pattern.sub(section, text, count=1)
        else:
            text = text.rstrip() + "\n\n" + section
    write_text(path, text)


def normalize_current_version_truth_section(text: str) -> tuple[str, list[str]]:
    pattern = re.compile(
        r"(?ims)^##[ \t]+Current Version Truth[ \t]*\n(.*?)(?=^##[ \t]+|\Z)"
    )
    matches = list(pattern.finditer(text))
    if len(matches) != 1:
        return text, []
    match = matches[0]
    body = match.group(1)
    missing = [
        key
        for key in CURRENT_VERSION_TRUTH_KEYS
        if not re.search(rf"(?m)^[ \t]*{re.escape(key)}[ \t]*:", body)
    ]
    if not missing:
        return text, []
    insertion = "".join(f"{key}:\n" for key in missing)
    fence = body.rfind("```")
    if fence >= 0:
        normalized_body = body[:fence].rstrip() + "\n" + insertion + body[fence:]
    else:
        normalized_body = body.rstrip() + "\n" + insertion
    normalized = text[: match.start(1)] + normalized_body + text[match.end(1) :]
    return normalized, missing


def md_cell(value: str | None) -> str:
    text = re.sub(r"\s+", " ", value or "").strip()
    return (text or "-").replace("|", "\\|")


def markdown_rows(rows: list[tuple[str, ...]], empty: str) -> str:
    if not rows:
        return f"| - | {md_cell(empty)} | - |\n"
    return "".join(f"| {md_cell(a)} | {md_cell(b)} | {md_cell(c)} |\n" for a, b, c in rows)




def normalized_bool(value: str | None) -> bool:
    return (value or "").strip().lower() in {"true", "yes", "y", "1", "client_visible", "approved"}


def non_placeholder(value: str | None, *, min_chars: int = 1) -> bool:
    text = re.sub(r"\s+", "", value or "")
    lowered = (value or "").strip().lower()
    if lowered in {"", "-", "tbd", "todo", "template", "pending", "open question", "n/a", "na"}:
        return False
    if "tbd" in lowered or "open question" in lowered or "待确认" in lowered:
        return False
    return len(text) >= min_chars


def split_asset_refs(value: str | None) -> list[str]:
    if not value:
        return []
    refs: list[str] = []
    for item in re.split(r"[;,，、\s]+", value):
        cleaned = item.strip().strip("`")
        if not cleaned or cleaned.lower() in {"none", "placeholder", "n/a", "na", "-"}:
            continue
        refs.append(cleaned)
    return refs


def indicates_browser_held_assets(*values: str | None) -> bool:
    text = " ".join(value or "" for value in values).lower()
    return any(marker in text for marker in BROWSER_HELD_ASSET_MARKERS)


def row_is_client_visible(row: dict[str, str]) -> bool:
    return row.get("visibility", "").strip().lower() in CLIENT_VISIBLE_VALUES


def text_visibility_value(text: str) -> str:
    match = re.search(r"(?im)^\s*visibility\s*:\s*([^\n#]+)", text)
    return match.group(1).strip().lower() if match else ""


def client_language_text_for_path(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return pptx_text_content(path)
    if suffix == ".pdf":
        extraction_errors: list[str] = []
        extractor = shutil.which("pdftotext")
        if extractor:
            result = subprocess.run(
                [extractor, str(path), "-"],
                capture_output=True,
                text=True,
                check=False,
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout
            extraction_errors.append(
                "pdftotext: " + (result.stderr.strip() or "empty output")
            )
        try:
            import fitz  # type: ignore[import-not-found]

            with fitz.open(path) as document:
                text = "\n".join(page.get_text() for page in document)
            if text.strip():
                return text
        except Exception as exc:  # noqa: BLE001 - fail closed below
            extraction_errors.append(f"fitz: {exc}")
        else:
            extraction_errors.append("fitz: empty output")
        raise RuntimeError(
            "PDF text extraction unavailable or empty: "
            + "; ".join(extraction_errors or ["no extractor available"])
        )
    if suffix == ".csv":
        fields, rows = read_csv_rows(path)
        if "visibility" in fields:
            return "\n".join(row_text(row) for row in rows if row_is_client_visible(row))
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="utf-8", errors="ignore")


def is_reliably_client_language_file(project: Path, path: Path) -> bool:
    suffix = path.suffix.lower()
    if suffix not in TEXT_CLIENT_SCAN_SUFFIXES | {".pptx", ".pdf"}:
        return False
    rel_path = safe_rel(project, path)
    if path.name in {"README.md", "目录索引.md"}:
        return False
    if suffix in TEXT_CLIENT_SCAN_SUFFIXES:
        text = client_language_text_for_path(path)
        visibility = text_visibility_value(text)
        if visibility == "internal_only":
            return False
        if visibility in CLIENT_VISIBLE_VALUES:
            return True
        if suffix == ".csv":
            fields, rows = read_csv_rows(path)
            if "visibility" in fields:
                return any(row_is_client_visible(row) for row in rows)
    if "05_最终交付_FinalDelivery/" in rel_path or "04_客户审阅_ClientReview/" in rel_path:
        return True
    return bool(CLIENT_LANGUAGE_RELIABLE_NAME_PATTERN.search(path.name))


def markdown_table_rows(text: str) -> list[dict[str, str]]:
    lines = text.splitlines()
    rows: list[dict[str, str]] = []
    for index, line in enumerate(lines):
        if not line.strip().startswith("|") or index + 1 >= len(lines):
            continue
        next_line = lines[index + 1].strip()
        if not next_line.startswith("|") or "---" not in next_line:
            continue
        headers = [cell.strip() for cell in line.strip().strip("|").split("|")]
        for raw_row in lines[index + 2 :]:
            stripped = raw_row.strip()
            if not stripped.startswith("|"):
                break
            values = [cell.strip() for cell in stripped.strip("|").split("|")]
            values.extend([""] * max(0, len(headers) - len(values)))
            rows.append(dict(zip(headers, values)))
        if rows:
            break
    return rows


def infer_slide_asset_usage(project: Path) -> dict[str, list[str]]:
    usage: dict[str, list[str]] = {}
    client_outline = project / "AD-creative/client_review/client_outline.csv"
    _, outline_rows = read_csv_rows(client_outline)
    for row in outline_rows:
        slide_id = row.get("slide_id", "").strip()
        for asset_id in split_asset_refs(row.get("asset_ids", "")):
            usage.setdefault(asset_id, []).append(slide_id)

    slide_spec = project / "AD-creative/client_review/slide_spec.md"
    if slide_spec.exists():
        for row in markdown_table_rows(slide_spec.read_text(encoding="utf-8")):
            slide_id = first_nonempty(row.get("Slide"), row.get("slide"), row.get("Page"), default="")
            asset_cell = first_nonempty(
                row.get("Asset Slot"),
                row.get("Asset"),
                row.get("asset_ids"),
                row.get("Image"),
                default="",
            )
            for asset_id in split_asset_refs(asset_cell):
                usage.setdefault(asset_id, []).append(slide_id)
    return usage


def file_stat_mtime(path: Path) -> str:
    try:
        return datetime.fromtimestamp(path.stat().st_mtime).astimezone().isoformat(timespec="seconds")
    except OSError:
        return ""


def infer_asset_platform(source: str, rel_path: str, asset_type: str) -> str:
    text = " ".join([source, rel_path, asset_type]).lower()
    for platform in ["grok", "chatgpt", "imagegen", "image_gen", "creative production"]:
        if platform in text:
            return "ImageGen" if platform == "image_gen" else platform
    if "download" in text:
        return "download"
    if rel_path:
        return "local_file"
    return "unknown"


def sync_asset_current_manifest(project: Path) -> Path:
    ensure_csv_file(project / "AD-creative/visual_assets/asset_current_manifest.csv", ASSET_CURRENT_FIELDS)
    ensure_csv_file(project / "AD-creative/orchestrator/final_delivery_lock.csv", FINAL_DELIVERY_LOCK_FIELDS)
    _, assets = read_csv_rows(project / "AD-creative/visual_assets/asset_manifest.csv")
    fields, current_rows = read_csv_rows(project / "AD-creative/visual_assets/asset_current_manifest.csv")
    by_id = {row.get("asset_id", ""): row for row in current_rows if row.get("asset_id", "")}
    usage = infer_slide_asset_usage(project)
    now = now_iso()
    output_rows: list[dict[str, str]] = []
    seen: set[str] = set()
    for asset in assets:
        asset_id = asset.get("asset_id", "").strip()
        if not asset_id:
            continue
        seen.add(asset_id)
        existing = by_id.get(asset_id, {})
        rel_path = asset.get("path", "").strip()
        path = project / rel_path if rel_path else Path()
        sha = ""
        if rel_path and path.exists() and path.is_file():
            sha = file_sha256(path)
        direct_client_use = existing.get("direct_client_use") or (
            "yes" if asset.get("visibility", "").strip().lower() in CLIENT_VISIBLE_VALUES else "no"
        )
        source_value = existing.get("source") or first_nonempty(
            asset.get("prompt_or_edit_ref"),
            asset.get("reference_id"),
            asset.get("asset_type"),
            default="unrecorded_source",
        )
        platform_value = existing.get("platform") or infer_asset_platform(source_value, rel_path, asset.get("asset_type", ""))
        qa_items = []
        if asset.get("qa_status", "").strip():
            qa_items.append(f"qa_status={asset.get('qa_status', '').strip()}")
        if asset.get("risk_level", "").strip():
            qa_items.append(f"risk_level={asset.get('risk_level', '').strip()}")
        qa_flags = existing.get("qa_flags") or ";".join(qa_items)
        output_rows.append(
            {
                "asset_id": asset_id,
                "source": source_value,
                "platform": platform_value,
                "conversation": existing.get("conversation", ""),
                "local_file": existing.get("local_file") or rel_path,
                "path": rel_path,
                "sha256": sha or existing.get("sha256", ""),
                "original_or_processed": existing.get("original_or_processed")
                or ("processed" if "/selected/" in rel_path or "/processed/" in rel_path else "original"),
                "approval": existing.get("approval", ""),
                "direct_client_use": direct_client_use,
                "used_in_slide": existing.get("used_in_slide") or ";".join(sorted(set(usage.get(asset_id, [])))),
                "qa_flags": qa_flags,
                "protected": existing.get("protected") or "false",
                "status": asset.get("status", "") or existing.get("status", "registered"),
                "notes": existing.get("notes") or f"synced_from_asset_manifest_at={now}",
            }
        )
    for asset_id, row in by_id.items():
        if asset_id not in seen:
            output_rows.append(row)
    write_csv_rows(project / "AD-creative/visual_assets/asset_current_manifest.csv", fields or ASSET_CURRENT_FIELDS, output_rows)
    return project / "AD-creative/visual_assets/asset_current_manifest.csv"


def final_delivery_project_path(
    project: Path, raw_path: str, *, require_file: bool = False
) -> tuple[str, Path]:
    raw = (raw_path or "").strip().strip("`")
    candidate = Path(raw)
    if not raw or candidate.is_absolute() or ".." in candidate.parts:
        raise ValueError(f"FinalDelivery path must be project-relative: {raw_path}")
    rel = unicodedata.normalize("NFC", candidate.as_posix())
    relative = Path(rel)
    if not relative.parts or relative.parts[0] != "05_最终交付_FinalDelivery":
        raise ValueError(f"FinalDelivery path must stay inside 05_最终交付_FinalDelivery: {raw_path}")
    if project_relative_path_has_symlink_component(project, rel):
        raise ValueError(f"FinalDelivery path contains a symlink component: {raw_path}")
    lexical = project / relative
    resolved = lexical.resolve()
    final_root = (project / "05_最终交付_FinalDelivery").resolve()
    try:
        resolved.relative_to(final_root)
        resolved.relative_to(project.resolve())
    except ValueError as exc:
        raise ValueError(f"FinalDelivery path escapes the protected root: {raw_path}") from exc
    if require_file and (not lexical.is_file() or lexical.is_symlink()):
        raise ValueError(f"FinalDelivery path is not a regular file: {raw_path}")
    return rel, lexical


def final_delivery_human_identity_valid(value: str) -> bool:
    identity = value.strip()
    if not non_placeholder(identity, min_chars=2):
        return False
    normalized = re.sub(r"[^a-z0-9]+", "_", identity.lower()).strip("_")
    if "main_controller" in normalized or normalized == "maincontroller":
        return False
    tokens = {token for token in normalized.split("_") if token}
    forbidden_tokens = {
        "adco",
        "assistant",
        "automation",
        "agent",
        "bot",
        "chatgpt",
        "claude",
        "codex",
        "gemini",
        "model",
        "system",
        "worker",
        "ai",
    }
    if tokens & forbidden_tokens:
        return False
    return not any(token in identity for token in ("自动化", "机器人", "系统代理", "执行代理"))


def final_delivery_lock_id(rel_path: str) -> str:
    """Return a stable path identity without collapsing Unicode filenames."""
    canonical = unicodedata.normalize("NFC", rel_path.strip())
    digest = hashlib.sha256(canonical.encode("utf-8")).hexdigest().upper()
    return f"LOCK-{digest}"


def final_delivery_evidence_binding(
    project: Path, evidence_ref: str
) -> tuple[str, Path, str]:
    raw = (evidence_ref or "").strip().strip("`")
    candidate = Path(raw)
    if not raw or candidate.is_absolute() or ".." in candidate.parts:
        raise ValueError("FinalDelivery evidence_ref must be a project-relative file")
    rel = unicodedata.normalize("NFC", candidate.as_posix())
    if project_relative_path_has_symlink_component(project, rel):
        raise ValueError("FinalDelivery evidence_ref must not traverse symlinks")
    path = project / rel
    resolved = path.resolve()
    try:
        resolved.relative_to(project.resolve())
    except ValueError as exc:
        raise ValueError("FinalDelivery evidence_ref escapes the project") from exc
    if not path.is_file() or path.is_symlink() or path.stat().st_nlink != 1:
        raise ValueError(
            "FinalDelivery evidence_ref must be an existing non-symlink, non-hardlinked file"
        )
    return rel, path, file_sha256(path)


def final_delivery_evidence_binding_valid(
    project: Path, row: dict[str, str]
) -> bool:
    try:
        _, _, actual_sha = final_delivery_evidence_binding(
            project, row.get("evidence_ref", "")
        )
    except ValueError:
        return False
    expected_sha = row.get("evidence_sha256", "").strip()
    return bool(re.fullmatch(r"[0-9a-f]{64}", expected_sha)) and actual_sha == expected_sha


def final_delivery_artifact_for_path(
    project: Path, rel_path: str, *, require_active: bool
) -> dict[str, str] | None:
    _, artifacts = read_csv_rows(
        project / "AD-creative/orchestrator/artifact_index.csv"
    )
    matches: list[dict[str, str]] = []
    for row in artifacts:
        try:
            artifact_rel, _ = final_delivery_project_path(
                project, row.get("path") or row.get("original_path") or ""
            )
        except ValueError:
            continue
        if artifact_rel != rel_path:
            continue
        if require_active and normalized_artifact_lifecycle(row) != "active":
            continue
        matches.append(row)
    return matches[0] if len(matches) == 1 else None


def final_delivery_old_version_id(
    project: Path, old_row: dict[str, str]
) -> str:
    try:
        old_rel, _ = final_delivery_project_path(project, old_row.get("path", ""))
    except ValueError:
        return ""
    artifact = final_delivery_artifact_for_path(
        project, old_rel, require_active=False
    )
    if not artifact:
        return ""
    artifact_id = artifact.get("artifact_id", "").strip()
    if not artifact_id:
        return ""
    _, versions = read_csv_rows(
        project / "AD-creative/orchestrator/version_map.csv"
    )
    explicit = old_row.get("version_id", "").strip()
    candidates = [
        row
        for row in versions
        if row.get("artifact_id", "").strip() == artifact_id
        and (not explicit or row.get("version_id", "").strip() == explicit)
    ]
    if len(candidates) != 1:
        return ""
    return candidates[0].get("version_id", "").strip()


def final_delivery_confirmation_source_valid(
    project: Path,
    *,
    source_event_id: str,
    evidence_rel: str,
    confirmed_by: str,
    new_artifact_id: str,
) -> bool:
    _, source_events = read_csv_rows(
        project / "AD-creative/orchestrator/source_events.csv"
    )
    matches = [
        row
        for row in source_events
        if row.get("source_event_id", "").strip() == source_event_id.strip()
    ]
    if len(matches) != 1:
        return False
    source = matches[0]
    if source.get("source_type", "").strip().lower() != "file":
        return False
    if source.get("source_owner", "").strip().casefold() != confirmed_by.strip().casefold():
        return False
    if source.get("declared_semantics", "").strip().lower() not in {
        "approval",
        "confirmation",
        "final_delivery_reconciliation",
    }:
        return False
    if source.get("trust_level", "").strip().lower() not in {
        "confirmed",
        "user_confirmed",
        "client_confirmed",
    }:
        return False
    source_paths = {
        canonical_project_relative_path(project, value)
        for value in split_registered_paths(source.get("file_paths", ""))
    }
    if source_paths != {evidence_rel}:
        return False
    if new_artifact_id:
        affected = set(split_asset_refs(source.get("affects_artifacts", "")))
        if new_artifact_id not in affected:
            return False
    return True


def final_delivery_host_attestation_binding(
    project: Path,
    *,
    attestation_ref: str,
    confirmation_receipt_ref: str,
    confirmation_receipt_sha256: str,
) -> tuple[str, str]:
    raw = (attestation_ref or "").strip().strip("`")
    candidate = Path(raw)
    if candidate.is_absolute() or ".." in candidate.parts:
        raise ValueError("FinalDelivery host attestation must be project-relative")
    rel = unicodedata.normalize("NFC", candidate.as_posix())
    root_prefix = FINAL_DELIVERY_HOST_ATTESTATION_ROOT.as_posix() + "/"
    if not rel.startswith(root_prefix):
        raise ValueError(
            "FinalDelivery host attestation must stay in the host-only attestation root"
        )
    if project_relative_path_has_symlink_component(project, rel):
        raise ValueError("FinalDelivery host attestation must not traverse symlinks")
    path = project / rel
    try:
        path.resolve().relative_to(project.resolve())
    except ValueError as exc:
        raise ValueError("FinalDelivery host attestation escapes the project") from exc
    if not path.is_file() or path.is_symlink() or path.stat().st_nlink != 1:
        raise ValueError(
            "FinalDelivery host attestation must be an existing non-symlink, non-hardlinked file"
        )
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (UnicodeDecodeError, json.JSONDecodeError) as exc:
        raise ValueError("FinalDelivery host attestation must be structured JSON") from exc
    if not isinstance(payload, dict):
        raise ValueError("FinalDelivery host attestation must be a JSON object")
    required_exact = {
        "protocol_id": FINAL_DELIVERY_HOST_ATTESTATION_PROTOCOL,
        "schema_version": FINAL_DELIVERY_HOST_ATTESTATION_VERSION,
        "attestation_scope": "final_delivery_reconciliation",
        "attestation_role": "host_main_thread",
        "verified_by": "main_controller",
        "readback_status": "verified",
        "readback_tool": "codex_app.read_thread",
        "confirmation_receipt_ref": confirmation_receipt_ref,
        "confirmation_receipt_sha256": confirmation_receipt_sha256,
    }
    if any(
        str(payload.get(key, "")).strip() != expected
        for key, expected in required_exact.items()
    ):
        raise ValueError("FinalDelivery host attestation binding mismatch")
    if str(payload.get("authority", "")).strip() not in {
        "user",
        "client",
        "project_owner",
    }:
        raise ValueError("FinalDelivery host attestation authority is invalid")
    if not non_placeholder(str(payload.get("attestation_id", "")).strip(), min_chars=8):
        raise ValueError("FinalDelivery host attestation_id is invalid")
    try:
        parse_thread_timestamp(str(payload.get("verified_at", "")), "verified_at")
        validate_real_thread_id(
            str(payload.get("thread_id", "")).strip(),
            "host_attestation.thread_id",
        )
    except ValueError as exc:
        raise ValueError("FinalDelivery host attestation identity/time is invalid") from exc
    if not non_placeholder(
        str(payload.get("user_message_id", "")).strip(), min_chars=8
    ) or not re.fullmatch(
        r"[0-9a-f]{64}", str(payload.get("user_message_sha256", "")).strip()
    ):
        raise ValueError("FinalDelivery host attestation user-message binding is invalid")
    return rel, file_sha256(path)


def final_delivery_confirmation_receipt_binding(
    project: Path,
    *,
    evidence_ref: str,
    old_row: dict[str, str],
    new_row: dict[str, str],
    kind: str,
    confirmed_by: str,
    confirmed_at: str,
    version_id: str,
) -> tuple[str, str, str, str, str]:
    evidence_rel, evidence_path, evidence_sha = final_delivery_evidence_binding(
        project, evidence_ref
    )
    try:
        payload = json.loads(evidence_path.read_text(encoding="utf-8"))
    except (UnicodeDecodeError, json.JSONDecodeError) as exc:
        raise ValueError(
            "FinalDelivery evidence_ref must contain a structured confirmation receipt"
        ) from exc
    if not isinstance(payload, dict):
        raise ValueError("FinalDelivery confirmation receipt must be a JSON object")
    old_rel, _ = final_delivery_project_path(project, old_row.get("path", ""))
    new_rel, new_path = final_delivery_project_path(
        project, new_row.get("path", ""), require_file=True
    )
    new_artifact = final_delivery_artifact_for_path(
        project, new_rel, require_active=True
    )
    new_artifact_id = (
        new_artifact.get("artifact_id", "").strip() if new_artifact else ""
    )
    old_version_id = (
        final_delivery_old_version_id(project, old_row)
        if kind == "supersession"
        else ""
    )
    if kind == "supersession" and not old_version_id:
        raise ValueError(
            "FinalDelivery supersession requires an unambiguous old artifact/version binding"
        )
    expected = {
        "protocol_id": FINAL_DELIVERY_CONFIRMATION_PROTOCOL,
        "schema_version": FINAL_DELIVERY_CONFIRMATION_VERSION,
        "confirmation_scope": "final_delivery_reconciliation",
        "decision": "approve_reconciliation",
        "confirmed_by": confirmed_by.strip(),
        "confirmed_at": confirmed_at.strip(),
        "reconciliation_kind": kind,
        "old_lock_id": old_row.get("lock_id", "").strip(),
        "old_path": old_rel,
        "old_sha256": old_row.get("sha256", "").strip(),
        "new_path": new_rel,
        "new_sha256": file_sha256(new_path),
        "new_artifact_id": new_artifact_id,
        "version_id": version_id.strip(),
        "supersedes_version_id": old_version_id,
    }
    mismatched = [
        key
        for key, expected_value in expected.items()
        if str(payload.get(key, "")).strip() != expected_value
    ]
    if mismatched:
        raise ValueError(
            "FinalDelivery confirmation receipt binding mismatch: "
            + ", ".join(mismatched)
        )
    confirmation_id = str(payload.get("confirmation_id", "")).strip()
    source_event_id = str(payload.get("source_event_id", "")).strip()
    if not non_placeholder(confirmation_id, min_chars=8) or not source_event_id:
        raise ValueError(
            "FinalDelivery confirmation receipt requires confirmation_id and source_event_id"
        )
    if not final_delivery_confirmation_source_valid(
        project,
        source_event_id=source_event_id,
        evidence_rel=evidence_rel,
        confirmed_by=confirmed_by,
        new_artifact_id=new_artifact_id,
    ):
        raise ValueError(
            "FinalDelivery confirmation receipt is not bound to one confirmed user/client source event"
        )
    host_attestation_ref = str(payload.get("host_attestation_ref", "")).strip()
    host_rel, host_sha = final_delivery_host_attestation_binding(
        project,
        attestation_ref=host_attestation_ref,
        confirmation_receipt_ref=evidence_rel,
        confirmation_receipt_sha256=evidence_sha,
    )
    return evidence_rel, evidence_sha, host_rel, host_sha, old_version_id


def final_delivery_version_binding_valid(
    project: Path, rel_path: str, version_id: str
) -> bool:
    version_id = version_id.strip()
    if not version_id:
        return False
    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    artifact_matches: list[dict[str, str]] = []
    for row in artifacts:
        try:
            artifact_rel, _ = final_delivery_project_path(
                project, row.get("path") or row.get("original_path") or ""
            )
        except ValueError:
            continue
        if (
            artifact_rel == rel_path
            and normalized_artifact_lifecycle(row) == "active"
        ):
            artifact_matches.append(row)
    if len(artifact_matches) != 1:
        return False
    artifact_id = artifact_matches[0].get("artifact_id", "").strip()
    if not artifact_id:
        return False
    _, versions = read_csv_rows(project / "AD-creative/orchestrator/version_map.csv")
    version_matches = [
        row for row in versions if row.get("version_id", "").strip() == version_id
    ]
    if len(version_matches) != 1:
        return False
    version = version_matches[0]
    truth_path = project / "AD-creative/orchestrator/current_truth.md"
    truth_text = truth_path.read_text(encoding="utf-8") if truth_path.is_file() else ""
    truth_version_matches = current_truth_value(truth_text, "current_version_id") == version_id
    version_artifact_id = version.get("artifact_id", "").strip()
    if version.get("status", "").strip().lower() != "current":
        return False
    if not truth_version_matches:
        return False
    if version_artifact_id != artifact_id:
        return False
    artifact_type = artifact_matches[0].get("artifact_type", "").strip().lower()
    truth_key_by_type = {
        "pptx": "current_pptx_artifact_id",
        "pdf": "current_pdf_artifact_id",
        "preview": "current_preview_artifact_id",
        "deck_preview": "current_preview_artifact_id",
        "png_preview": "current_preview_artifact_id",
        "jpg_preview": "current_preview_artifact_id",
        "text_extract": "current_text_extract_artifact_id",
        "ppt_text_extract": "current_text_extract_artifact_id",
        "ppt_editability_check": "current_ppt_editability_artifact_id",
    }
    truth_key = truth_key_by_type.get(artifact_type, "")
    if not truth_key or current_truth_value(truth_text, truth_key) != artifact_id:
        return False
    return True


def final_delivery_supersession_chain_valid(
    project: Path, old_row: dict[str, str], new_version_id: str
) -> bool:
    old_version_id = final_delivery_old_version_id(project, old_row)
    if not old_version_id:
        return False
    _, versions = read_csv_rows(project / "AD-creative/orchestrator/version_map.csv")
    matches = [
        row
        for row in versions
        if row.get("version_id", "").strip() == new_version_id.strip()
    ]
    return bool(
        len(matches) == 1
        and matches[0].get("supersedes_version_id", "").strip()
        == old_version_id
    )


def final_delivery_lock_snapshot(
    project: Path,
    *,
    protected_value: str = "true",
) -> tuple[list[dict[str, str]], Path]:
    """Inventory first; never refresh an existing baseline or mutate user files."""
    lock_path = project / "AD-creative/orchestrator/final_delivery_lock.csv"
    ensure_csv_file(lock_path, FINAL_DELIVERY_LOCK_FIELDS)
    fields, rows = read_csv_rows(lock_path)
    issues: list[str] = []
    by_path: dict[str, dict[str, str]] = {}
    output_rows: list[dict[str, str]] = []
    duplicate_raw_paths: set[str] = set()
    lock_ids: set[str] = set()
    duplicate_lock_ids: set[str] = set()
    for row in rows:
        copied = dict(row)
        output_rows.append(copied)
        lock_id = row.get("lock_id", "").strip()
        if not lock_id:
            issues.append(
                "FinalDelivery lock row missing lock_id: "
                + (row.get("path", "").strip() or "<blank path>")
            )
        elif lock_id in lock_ids:
            duplicate_lock_ids.add(lock_id)
        else:
            lock_ids.add(lock_id)
        raw_path = row.get("path", "").strip()
        if not raw_path:
            continue
        if raw_path in by_path:
            duplicate_raw_paths.add(raw_path)
            continue
        by_path[raw_path] = copied
    for raw_path in sorted(duplicate_raw_paths):
        issues.append(f"duplicate FinalDelivery lock path: {raw_path}")
    for lock_id in sorted(duplicate_lock_ids):
        issues.append(f"duplicate FinalDelivery lock_id: {lock_id}")

    canonical_to_raw: dict[str, str] = {}
    pending_paths: list[str] = []

    for rel_path, row in by_path.items():
        try:
            canonical_rel, path = final_delivery_project_path(project, rel_path)
        except ValueError as exc:
            issues.append(f"unsafe lock path: {rel_path}: {exc}")
            continue
        if canonical_rel in canonical_to_raw:
            issues.append(
                "duplicate canonical FinalDelivery lock path: "
                f"{canonical_to_raw[canonical_rel]} and {rel_path}"
            )
            continue
        canonical_to_raw[canonical_rel] = rel_path
        if (
            row.get("inventory_state", "").strip() == "pending_reconciliation"
            and not normalized_bool(row.get("protected"))
            and path.is_file()
            and not is_final_delivery_metadata("", rel_path)
        ):
            pending_paths.append(rel_path)
        if not normalized_bool(row.get("protected")):
            continue
        row.setdefault("inventory_state", "protected_baseline")
        if not path.exists() or not path.is_file():
            if not final_delivery_reconciliation_valid(project, by_path, row):
                issues.append(f"protected file missing: {rel_path}")
            continue
        current_sha = file_sha256(path)
        current_size = str(path.stat().st_size)
        expected_sha = row.get("sha256", "").strip()
        expected_size = row.get("size_bytes", "").strip()
        if expected_sha and current_sha != expected_sha:
            issues.append(f"protected file changed: {rel_path}")
        elif expected_size and current_size != expected_size:
            issues.append(f"protected file size changed: {rel_path}")
        elif not expected_sha or not expected_size:
            issues.append(f"protected baseline incomplete: {rel_path}")

    final_root = project / "05_最终交付_FinalDelivery"
    if final_root.exists() and project_relative_path_has_symlink_component(
        project, "05_最终交付_FinalDelivery"
    ):
        issues.append("unsafe FinalDelivery root: symlink components are forbidden")
        inventory_paths: list[Path] = []
    else:
        inventory_paths = sorted(final_root.rglob("*")) if final_root.exists() else []
    for path in inventory_paths:
        if not path.is_file() or path.name in {"README.md", "目录索引.md", ".DS_Store"}:
            continue
        lexical_rel = path.relative_to(project).as_posix()
        try:
            rel_path, safe_path = final_delivery_project_path(
                project, lexical_rel, require_file=True
            )
        except ValueError as exc:
            issues.append(f"unsafe FinalDelivery inventory path: {lexical_rel}: {exc}")
            continue
        if rel_path in canonical_to_raw:
            continue
        metadata_only = is_final_delivery_metadata("", rel_path)
        new_row = {
            "lock_id": final_delivery_lock_id(rel_path),
            "path": rel_path,
            "sha256": file_sha256(safe_path),
            "size_bytes": str(safe_path.stat().st_size),
            "mtime": file_stat_mtime(safe_path),
            "protected": "no",
            "registered_at": now_iso(),
            "notes": (
                "generated_metadata_inventory_not_user_final"
                if metadata_only
                else "pending_final_delivery_inventory_do_not_move_or_overwrite"
            ),
            "inventory_state": "metadata_excluded" if metadata_only else "pending_reconciliation",
            "reconciliation_state": "not_applicable" if metadata_only else "pending",
            "reconciliation_kind": "",
            "reconciles_lock_id": "",
            "supersedes_lock_id": "",
            "confirmed_by": "",
            "confirmed_at": "",
            "evidence_ref": "",
            "evidence_sha256": "",
            "host_attestation_ref": "",
            "host_attestation_sha256": "",
            "version_id": "",
            "supersedes_version_id": "",
            "status_reason": "generated metadata is not a user final" if metadata_only else "new physical file inventoried safely",
        }
        by_path[rel_path] = new_row
        output_rows.append(new_row)
        if not metadata_only:
            pending_paths.append(rel_path)
        canonical_to_raw[rel_path] = rel_path

    # Safe inventory is durable even if an older protected path is missing or changed.
    write_csv_rows(lock_path, fields or FINAL_DELIVERY_LOCK_FIELDS, output_rows)

    if issues:
        raise RuntimeError("FinalDelivery lock integrity violation: " + "; ".join(issues))
    for rel_path in pending_paths:
        row = by_path[rel_path]
        row["protected"] = protected_value
        row["inventory_state"] = "protected_baseline"
        row["reconciliation_state"] = "not_required"
        row["status_reason"] = "initial immutable baseline registered"
    locked = output_rows
    write_csv_rows(lock_path, fields or FINAL_DELIVERY_LOCK_FIELDS, locked)
    return locked, lock_path


def final_delivery_reconciliation_valid(
    project: Path,
    rows_by_path: dict[str, dict[str, str]],
    old_row: dict[str, str],
) -> bool:
    old_lock_id = old_row.get("lock_id", "").strip()
    if not old_lock_id:
        return False
    candidates = [
        row
        for row in rows_by_path.values()
        if row.get("reconciles_lock_id", "").strip() == old_lock_id
    ]
    if len(candidates) != 1:
        return False
    new_row = candidates[0]
    if not (
        normalized_bool(new_row.get("protected"))
        and new_row.get("reconciliation_state", "").strip().lower() == "reconciled"
        and final_delivery_human_identity_valid(new_row.get("confirmed_by", ""))
        and new_row.get("confirmed_at", "").strip()
    ):
        return False
    try:
        parse_thread_timestamp(new_row.get("confirmed_at", ""), "confirmed_at")
    except ValueError:
        return False
    try:
        new_rel, new_path = final_delivery_project_path(
            project, new_row.get("path", ""), require_file=True
        )
        final_delivery_project_path(project, old_row.get("path", ""))
    except ValueError:
        return False
    actual_sha = file_sha256(new_path)
    if actual_sha != new_row.get("sha256", "").strip():
        return False
    kind = new_row.get("reconciliation_kind", "").strip().lower()
    try:
        evidence_rel, evidence_sha, host_rel, host_sha, old_version_id = (
            final_delivery_confirmation_receipt_binding(
                project,
                evidence_ref=new_row.get("evidence_ref", ""),
                old_row=old_row,
                new_row=new_row,
                kind=kind,
                confirmed_by=new_row.get("confirmed_by", ""),
                confirmed_at=new_row.get("confirmed_at", ""),
                version_id=new_row.get("version_id", ""),
            )
        )
    except ValueError:
        return False
    if (
        new_row.get("evidence_ref", "").strip() != evidence_rel
        or new_row.get("evidence_sha256", "").strip() != evidence_sha
        or new_row.get("host_attestation_ref", "").strip() != host_rel
        or new_row.get("host_attestation_sha256", "").strip() != host_sha
        or new_row.get("supersedes_version_id", "").strip() != old_version_id
    ):
        return False
    if kind == "rename":
        return actual_sha == old_row.get("sha256", "").strip()
    if kind == "supersession":
        return (
            actual_sha != old_row.get("sha256", "").strip()
            and bool(new_row.get("version_id", "").strip())
            and new_row.get("supersedes_lock_id", "").strip() == old_lock_id
            and final_delivery_version_binding_valid(
                project, new_rel, new_row.get("version_id", "")
            )
            and final_delivery_supersession_chain_valid(
                project, old_row, new_row.get("version_id", "")
            )
        )
    return False


def reconcile_final_delivery(
    project: Path,
    *,
    old_path: str,
    new_path: str,
    kind: str,
    confirmed_by: str,
    confirmed_at: str,
    evidence_ref: str,
    version_id: str = "",
) -> tuple[dict[str, str], Path]:
    """Record an evidence-bound rename or supersession without moving either file."""
    if kind not in {"rename", "supersession"}:
        raise ValueError("FinalDelivery reconciliation kind must be rename or supersession")
    if not final_delivery_human_identity_valid(confirmed_by):
        raise ValueError(
            "FinalDelivery reconciliation confirmed_by must identify a human, not ADCO/automation/agent"
        )
    if not confirmed_at.strip() or not evidence_ref.strip():
        raise ValueError(
            "FinalDelivery reconciliation requires confirmed_by, confirmed_at, and evidence_ref"
        )
    parse_thread_timestamp(confirmed_at, "confirmed_at")
    try:
        final_delivery_lock_snapshot(project)
    except RuntimeError:
        # Expected for a missing old baseline; the inventory-first write already persisted.
        pass
    lock_path = project / "AD-creative/orchestrator/final_delivery_lock.csv"
    fields, rows = read_csv_rows(lock_path)
    old_rel, old_file = final_delivery_project_path(project, old_path)
    new_rel, new_file = final_delivery_project_path(
        project, new_path, require_file=True
    )
    old_matches = [
        row
        for row in rows
        if canonical_project_relative_path(project, row.get("path")) == old_rel
    ]
    new_matches = [
        row
        for row in rows
        if canonical_project_relative_path(project, row.get("path")) == new_rel
    ]
    if len(old_matches) != 1 or len(new_matches) != 1:
        raise ValueError("FinalDelivery reconciliation paths must both be inventoried")
    old_row = old_matches[0]
    new_row = new_matches[0]
    if old_file.exists():
        raise ValueError("old FinalDelivery baseline still exists; reconciliation is only for a missing original path")
    old_sha = old_row.get("sha256", "").strip()
    new_sha = file_sha256(new_file)
    if kind == "rename" and old_sha != new_sha:
        raise ValueError("rename reconciliation requires the same sha256")
    if kind == "supersession":
        if old_sha == new_sha:
            raise ValueError("same-hash replacement should be reconciled as rename")
        if not version_id.strip():
            raise ValueError("different-hash supersession requires a new version_id")
        if not final_delivery_version_binding_valid(project, new_rel, version_id):
            raise ValueError(
                "different-hash supersession version_id must bind the new artifact in version_map/current truth"
            )
        if not final_delivery_supersession_chain_valid(
            project, old_row, version_id
        ):
            raise ValueError(
                "different-hash supersession version chain does not supersede the old baseline version"
            )
    evidence_rel, evidence_sha, host_rel, host_sha, old_version_id = (
        final_delivery_confirmation_receipt_binding(
            project,
            evidence_ref=evidence_ref,
            old_row=old_row,
            new_row=new_row,
            kind=kind,
            confirmed_by=confirmed_by,
            confirmed_at=confirmed_at,
            version_id=version_id,
        )
    )
    new_row.update(
        {
            "protected": "yes",
            "inventory_state": "protected_baseline",
            "reconciliation_state": "reconciled",
            "reconciliation_kind": kind,
            "reconciles_lock_id": old_row.get("lock_id", ""),
            "supersedes_lock_id": old_row.get("lock_id", "") if kind == "supersession" else "",
            "confirmed_by": confirmed_by.strip(),
            "confirmed_at": confirmed_at.strip(),
            "evidence_ref": evidence_rel,
            "evidence_sha256": evidence_sha,
            "host_attestation_ref": host_rel,
            "host_attestation_sha256": host_sha,
            "version_id": version_id.strip(),
            "supersedes_version_id": old_version_id,
            "status_reason": f"explicit {kind} reconciliation; no file move or overwrite performed",
        }
    )
    by_path = {
        row.get("path", "").strip(): row
        for row in rows
        if row.get("path", "").strip()
    }
    if not final_delivery_reconciliation_valid(project, by_path, old_row):
        raise RuntimeError("FinalDelivery reconciliation did not produce a valid evidence-bound state")
    write_csv_rows(lock_path, fields or FINAL_DELIVERY_LOCK_FIELDS, rows)
    return new_row, lock_path


def register_final_delivery_locks(project: Path) -> Path:
    _, lock_path = final_delivery_lock_snapshot(project)
    return lock_path


HUMAN_WORKSPACE_SPECS = (
    ("00_项目资料_ProjectMaterials", "项目资料 Project Materials", "登记客户资料、会议记录、导演组资料和客户反馈。"),
    ("01_参考资料_References", "参考资料 References", "登记官方参考、竞品案例、风格参考和视频参考。"),
    ("02_重要素材_KeyAssets", "重要素材 Key Assets", "登记品牌、产品、人物和有授权依据的关键素材。"),
    ("03_阶段成果_WorkInProgress", "阶段成果 Work In Progress", "登记内部方向草案、文案草案、视觉探索和方案结构。"),
    ("04_客户审阅_ClientReview", "客户审阅 Client Review", "登记准备给客户审阅的版本；内部执行资料不得混入。"),
    ("05_最终交付_FinalDelivery", "最终交付 Final Delivery", "登记最终确认交付物；这里只显示交付物，不显示 Gate/检查单/预览/文本抽取元数据。"),
)
HUMAN_INDEX_EXCLUDED_NAMES = {"README.md", "目录索引.md", ".DS_Store"}
CURRENT_ARTIFACT_TRUTH_KEYS = (
    "current_pptx_artifact_id",
    "current_pdf_artifact_id",
    "current_preview_artifact_id",
    "current_text_extract_artifact_id",
    "current_ppt_editability_artifact_id",
)
FINAL_DELIVERY_METADATA_ARTIFACT_TYPES = {
    "gate_report",
    "gate_checklist",
    "checklist",
    "preview",
    "deck_preview",
    "text_extract",
    "ppt_text_extract",
    "ppt_editability_check",
    "editability_report",
    "delivery_manifest",
    "final_delivery_index",
    "final_delivery_lock",
}
FINAL_DELIVERY_METADATA_MARKERS = {
    "gate",
    "checklist",
    "preview",
    "editability",
    "manifest",
    "lock",
}
FINAL_DELIVERY_ALWAYS_DELIVERABLE_SUFFIXES = {
    ".pdf",
    ".pptx",
    ".docx",
    ".xlsx",
    ".key",
    ".mov",
    ".mp4",
    ".zip",
}


def normalized_artifact_lifecycle(
    row: dict[str, str], *, superseded_ids: set[str] | None = None
) -> str:
    explicit = (row.get("lifecycle_state") or "").strip().lower().replace("-", "_")
    if explicit in ARTIFACT_LIFECYCLE_VALUES:
        return explicit
    artifact_id = (row.get("artifact_id") or "").strip()
    if artifact_id and artifact_id in (superseded_ids or set()):
        return "superseded"
    status = (row.get("status") or "").strip().lower().replace("-", "_").replace(" ", "_")
    if "removed" in status and any(token in status for token in ("clean", "cleanup")):
        return "legacy_unresolved_tombstone"
    status_map = {
        "superseded": "superseded",
        "withdrawn": "withdrawn",
        "archived": "archived",
        "deprecated": "deprecated",
        "rejected": "rejected",
        "removed": "removed",
        "deleted": "removed",
        "pending": "pending",
        "planned": "pending",
        "draft": "pending",
        "blocked": "pending",
        "not_run": "pending",
    }
    if status in status_map:
        return status_map[status]
    if status in {"", "active", "current", "done", "complete", "completed", "approved", "registered", "pass", "passed", "internal_review", "ready"}:
        return "active"
    return "legacy_unknown"


def canonical_project_relative_path(project: Path, raw_path: str | None) -> str:
    raw = (raw_path or "").strip().strip("`")
    if not raw or re.match(r"^[a-z][a-z0-9+.-]*://", raw, re.IGNORECASE):
        return ""
    candidate = Path(raw)
    try:
        resolved = candidate.resolve() if candidate.is_absolute() else (project / candidate).resolve()
        rel = resolved.relative_to(project.resolve())
    except (OSError, ValueError):
        return ""
    return unicodedata.normalize("NFC", rel.as_posix())


def split_registered_paths(value: str | None) -> list[str]:
    return [item.strip().strip("`") for item in re.split(r"[;\n]+", value or "") if item.strip().strip("`")]


def human_index_link(project: Path, index_path: Path, raw_target: str) -> str:
    target = (raw_target or "").strip()
    if re.match(r"^https://", target, re.IGNORECASE):
        return f"[打开链接]({target})"
    if project_relative_path_has_symlink_component(project, target):
        return f"`{md_cell(target)}` (UNSAFE SYMLINK; NOT LINKED)"
    rel = canonical_project_relative_path(project, target)
    if not rel:
        return md_cell(target)
    absolute = project / rel
    if not absolute.is_file():
        return f"`{md_cell(rel)}` (MISSING)"
    destination = os.path.relpath(absolute, index_path.parent).replace(os.sep, "/")
    destination = urllib.parse.quote(destination, safe="/._-~")
    label = md_cell(rel)
    return f"[{label}](<{destination}>)"


def is_final_delivery_metadata(artifact_type: str, raw_path: str) -> bool:
    normalized_type = re.sub(r"[^a-z0-9]+", "_", artifact_type.lower()).strip("_")
    if normalized_type in FINAL_DELIVERY_METADATA_ARTIFACT_TYPES:
        return True
    path = Path(raw_path)
    if path.suffix.lower() in FINAL_DELIVERY_ALWAYS_DELIVERABLE_SUFFIXES:
        return False
    stem = re.sub(r"[^a-z0-9]+", "_", path.stem.lower()).strip("_")
    tokens = {token for token in stem.split("_") if token}
    strong_phrases = {
        "text_extract",
        "gate_report",
        "gate_checklist",
        "delivery_gate",
        "final_delivery_index",
        "delivery_index",
        "lock_snapshot",
        "delivery_manifest",
    }
    return bool(tokens & FINAL_DELIVERY_METADATA_MARKERS) or any(
        phrase in stem for phrase in strong_phrases
    )


def physical_human_files(project: Path, folder: str) -> list[str]:
    root = project / folder
    if not root.exists():
        return []
    files: list[str] = []
    for path in root.rglob("*"):
        lexical_rel = path.relative_to(project).as_posix()
        if (
            path.name in HUMAN_INDEX_EXCLUDED_NAMES
            or any(part.startswith(".") for part in path.relative_to(root).parts)
            or path.is_symlink()
            or project_relative_path_has_symlink_component(project, lexical_rel)
            or not path.is_file()
        ):
            continue
        try:
            path.resolve().relative_to(project.resolve())
            path.resolve().relative_to(root.resolve())
        except ValueError:
            continue
        files.append(unicodedata.normalize("NFC", lexical_rel))
    return sorted(files, key=lambda value: (value.casefold(), value))


def physical_human_unsafe_files(project: Path, folder: str) -> list[str]:
    root = project / folder
    if not root.exists():
        return []
    unsafe: list[str] = []
    for path in root.rglob("*"):
        lexical_rel = path.relative_to(project).as_posix()
        if (
            path.name in HUMAN_INDEX_EXCLUDED_NAMES
            or any(part.startswith(".") for part in path.relative_to(root).parts)
        ):
            continue
        if path.is_symlink() or project_relative_path_has_symlink_component(
            project, lexical_rel
        ):
            unsafe.append(unicodedata.normalize("NFC", lexical_rel))
    return sorted(set(unsafe), key=lambda value: (value.casefold(), value))


def artifact_workspace_kinds(row: dict[str, str]) -> set[str]:
    path = (row.get("path") or row.get("original_path") or "").strip()
    stage = (row.get("stage") or "").strip().lower()
    visibility = (row.get("visibility") or "").strip().lower()
    artifact_type = (row.get("artifact_type") or "").strip().lower()
    kinds: set[str] = set()
    if (
        visibility not in CLIENT_VISIBLE_VALUES
        and stage not in {"client_review", "ppt_gate", "final_delivery"}
        and "/gates/" not in path
        and "/orchestrator/" not in path
    ):
        kinds.add("wip")
    if (
        visibility in CLIENT_VISIBLE_VALUES
        or "client_visible" in visibility
        or stage in {"client_review", "ppt_gate"}
        or "/client_review/" in path
        or "/ppt/previews/" in path
    ):
        kinds.add("client")
    if (
        stage == "final_delivery"
        or path.startswith("05_最终交付_FinalDelivery/")
        or "/delivery/" in path
    ) and not is_final_delivery_metadata(artifact_type, path):
        kinds.add("final")
    return kinds


def render_human_workspace_indexes(project: Path) -> list[Path]:
    """Render current-first, path-deduped indexes without copying or aliasing files."""
    project = project.resolve()
    _, source_events = read_csv_rows(project / "AD-creative/orchestrator/source_events.csv")
    _, references = read_csv_rows(project / "AD-creative/references/reference_cards.csv")
    _, assets = read_csv_rows(project / "AD-creative/visual_assets/asset_manifest.csv")
    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    _, version_rows = read_csv_rows(project / "AD-creative/orchestrator/version_map.csv")
    truth_path = project / "AD-creative/orchestrator/current_truth.md"
    truth_text = truth_path.read_text(encoding="utf-8") if truth_path.is_file() else ""
    exact_ids = {
        key: current_truth_value(truth_text, key)
        for key in CURRENT_ARTIFACT_TRUTH_KEYS
        if current_truth_value(truth_text, key)
    }
    exact_rank = {artifact_id: rank for rank, artifact_id in enumerate(exact_ids.values())}
    artifacts_by_id = {
        row.get("artifact_id", "").strip(): row
        for row in artifacts
        if row.get("artifact_id", "").strip()
    }
    superseded_ids = {
        ref
        for row in artifacts
        for ref in split_asset_refs(row.get("supersedes_artifact_id", ""))
    }
    attention: list[str] = []
    current_version_id = current_truth_value(truth_text, "current_version_id")
    if current_version_id:
        current_versions = [row for row in version_rows if row.get("version_id", "").strip() == current_version_id]
        if len(current_versions) != 1:
            attention.append(f"current_version_id `{current_version_id}` 必须且只能匹配一条 version_map 记录；当前匹配 {len(current_versions)} 条。")
        elif (
            current_versions[0].get("status", "").strip().lower()
            not in CURRENT_VIEW_VERSION_STATUSES
        ):
            attention.append(
                f"current_version_id `{current_version_id}` 对应 version_map 状态不属于 current-view 状态。"
            )
    for truth_key, artifact_id in exact_ids.items():
        row = artifacts_by_id.get(artifact_id)
        if not row:
            attention.append(f"{truth_key} 指向未登记 artifact `{artifact_id}`。")
            continue
        lifecycle = normalized_artifact_lifecycle(row, superseded_ids=superseded_ids)
        if lifecycle != "active":
            attention.append(
                f"{truth_key} 指向 non-active artifact `{artifact_id}`（{lifecycle}）。"
            )
        raw_exact_path = row.get("path") or row.get("original_path") or ""
        if project_relative_path_has_symlink_component(project, raw_exact_path):
            attention.append(
                f"{truth_key} 的 exact-current target 包含 symlink，已拒绝链接：`{raw_exact_path}`。"
            )
        rel = canonical_project_relative_path(project, raw_exact_path)
        if not rel or not (project / rel).is_file():
            attention.append(f"{truth_key} 的 exact-current target 缺失：`{row.get('path') or row.get('original_path') or '<blank>'}`。")

    entries_by_folder: dict[str, dict[str, dict[str, object]]] = {
        folder: {} for folder, _, _ in HUMAN_WORKSPACE_SPECS
    }

    def add_entry(
        folder: str,
        *,
        entry_id: str,
        label: str,
        target: str,
        lifecycle: str = "active",
        exact: bool = False,
        registered: bool = True,
    ) -> None:
        unsafe_symlink = project_relative_path_has_symlink_component(project, target)
        rel = "" if unsafe_symlink else canonical_project_relative_path(project, target)
        key = f"path:{rel}" if rel else f"id:{entry_id}:{target}"
        existing = entries_by_folder[folder].get(key)
        candidate = {
            "ids": {entry_id} if entry_id else set(),
            "label": label or "-",
            "target": target,
            "rel": rel,
            "lifecycle": lifecycle,
            "exact": exact,
            "rank": exact_rank.get(entry_id, 9999),
            "registered": registered,
            "unsafe_symlink": unsafe_symlink,
        }
        if existing is None:
            entries_by_folder[folder][key] = candidate
            return
        existing_registered = bool(existing["registered"])
        candidate_registered = bool(candidate["registered"])
        if existing_registered and not candidate_registered:
            merged_ids = set(existing["ids"])
        elif candidate_registered and not existing_registered:
            merged_ids = {
                value for value in set(existing["ids"]) if value != "LOCAL/UNREGISTERED"
            } | set(candidate["ids"])
        else:
            merged_ids = set(existing["ids"]) | set(candidate["ids"])
        existing["ids"] = merged_ids
        existing["exact"] = bool(existing["exact"] or exact)
        existing["registered"] = bool(existing["registered"] or registered)
        existing["unsafe_symlink"] = bool(
            existing.get("unsafe_symlink") or unsafe_symlink
        )
        existing["rank"] = min(int(existing["rank"]), int(candidate["rank"]))
        if lifecycle not in ARTIFACT_INACTIVE_LIFECYCLE_VALUES:
            existing["lifecycle"] = lifecycle
        if (label.casefold(), label) < (str(existing["label"]).casefold(), str(existing["label"])):
            existing["label"] = label

    for row in source_events:
        paths = split_registered_paths(row.get("file_paths")) or [row.get("status", "")]
        for target in paths:
            add_entry(
                "00_项目资料_ProjectMaterials",
                entry_id=row.get("source_event_id", ""),
                label=first_nonempty(row.get("raw_summary"), row.get("source_type"), default="资料"),
                target=target,
            )
    for row in references:
        add_entry(
            "01_参考资料_References",
            entry_id=row.get("reference_id", ""),
            label=first_nonempty(row.get("title"), row.get("platform"), default="参考"),
            target=first_nonempty(row.get("url"), row.get("role"), default=row.get("reference_type", "")),
            lifecycle=normalized_artifact_lifecycle(row),
        )
    for row in assets:
        add_entry(
            "02_重要素材_KeyAssets",
            entry_id=row.get("asset_id", ""),
            label=first_nonempty(row.get("slot_id"), row.get("asset_type"), default="素材"),
            target=first_nonempty(row.get("path"), row.get("status"), default=row.get("visibility", "")),
            lifecycle=normalized_artifact_lifecycle(row),
        )
    folder_for_kind = {
        "wip": "03_阶段成果_WorkInProgress",
        "client": "04_客户审阅_ClientReview",
        "final": "05_最终交付_FinalDelivery",
    }
    for row in artifacts:
        artifact_id = row.get("artifact_id", "").strip()
        lifecycle = normalized_artifact_lifecycle(row, superseded_ids=superseded_ids)
        target = row.get("path") or row.get("original_path") or row.get("cleanup_ref") or ""
        for kind in artifact_workspace_kinds(row):
            add_entry(
                folder_for_kind[kind],
                entry_id=artifact_id,
                label=first_nonempty(row.get("artifact_type"), row.get("stage"), default="产物"),
                target=target,
                lifecycle=lifecycle,
                exact=artifact_id in exact_rank,
            )

    for folder, _, _ in HUMAN_WORKSPACE_SPECS:
        for rel in physical_human_files(project, folder):
            if folder == "05_最终交付_FinalDelivery" and is_final_delivery_metadata("", rel):
                continue
            add_entry(
                folder,
                entry_id="LOCAL/UNREGISTERED",
                label=Path(rel).name,
                target=rel,
                lifecycle="active",
                registered=False,
            )
        for rel in physical_human_unsafe_files(project, folder):
            add_entry(
                folder,
                entry_id="LOCAL/UNSAFE_SYMLINK",
                label=Path(rel).name,
                target=rel,
                lifecycle="pending",
                registered=False,
            )

    def table(entries: list[dict[str, object]], index_path: Path, empty: str) -> str:
        if not entries:
            return f"| - | - | {md_cell(empty)} | - |"
        lines: list[str] = []
        for entry in entries:
            ids = sorted(set(entry["ids"]), key=lambda value: (value.casefold(), value))
            registered = bool(entry["registered"])
            exact = bool(entry["exact"])
            rel = str(entry["rel"])
            missing = bool(rel and not (project / rel).is_file())
            status = "EXACT CURRENT" if exact else "REGISTERED" if registered else "LOCAL/UNREGISTERED"
            if bool(entry.get("unsafe_symlink")):
                status = "LOCAL/UNSAFE_SYMLINK"
            if missing:
                status += "/MISSING"
            lines.append(
                "| "
                + " | ".join(
                    [
                        md_cell(status),
                        md_cell(";".join(ids) or "LOCAL/UNREGISTERED"),
                        md_cell(str(entry["label"])),
                        human_index_link(project, index_path, str(entry["target"])),
                    ]
                )
                + " |"
            )
        return "\n".join(lines)

    written: list[Path] = []
    for folder, title, purpose in HUMAN_WORKSPACE_SPECS:
        index_path = project / folder / "目录索引.md"
        entries = list(entries_by_folder[folder].values())
        current_entries = [
            entry
            for entry in entries
            if bool(entry["exact"])
            or str(entry["lifecycle"]) not in ARTIFACT_INACTIVE_LIFECYCLE_VALUES
        ]
        history_entries = [
            entry
            for entry in entries
            if not bool(entry["exact"])
            and str(entry["lifecycle"]) in ARTIFACT_INACTIVE_LIFECYCLE_VALUES
        ]
        sort_key = lambda entry: (
            0 if bool(entry["exact"]) else 1 if bool(entry["registered"]) else 2,
            int(entry["rank"]),
            str(entry["rel"] or entry["target"]).casefold(),
            str(entry["rel"] or entry["target"]),
        )
        current_entries.sort(key=sort_key)
        history_entries.sort(key=sort_key)
        attention_text = ""
        if folder in {"03_阶段成果_WorkInProgress", "04_客户审阅_ClientReview", "05_最终交付_FinalDelivery"} and attention:
            attention_text = "\n## Exact-current 注意事项（P0）\n\n" + "\n".join(f"- {item}" for item in attention) + "\n"
        history_text = ""
        if history_entries:
            history_text = (
                "\n## 历史（inactive，保留追溯）\n\n"
                "| 状态 | ID | 内容 | 原路径或清理证据 |\n"
                "|---|---|---|---|\n"
                + table(history_entries, index_path, "暂无历史记录。")
                + "\n"
            )
        write_text(
            index_path,
            f"""# {title}

用途：{purpose}

当前视图以 `Current Version Truth` 与 `version_map.csv` 为先；同一路径只出现一次。索引只写链接和状态，不复制、移动、创建软链或别名。
{attention_text}
## 当前

| 状态 | ID | 内容 | 路径或状态 |
|---|---|---|---|
{table(current_entries, index_path, "暂无当前记录。")}
{history_text}""",
        )
        written.append(index_path)
    return written


TEXT_SUFFIXES = {".md", ".txt", ".csv", ".json", ".yml", ".yaml"}
SKIP_MATERIAL_LINES = {"text", "```"}
REQUIREMENT_TRIGGER_PATTERN = re.compile(
    r"项目|客户希望|希望|要求|交付|产品|品牌|方向|关键画面|关键帧|moodboard|参考|PPT|不要|不能|必须|需要|新增要求|客户明确"
)


def material_files(path: Path) -> list[Path]:
    if path.is_file():
        return [path] if path.suffix.lower() in TEXT_SUFFIXES else []
    if not path.is_dir():
        return []
    files = [
        item
        for item in path.rglob("*")
        if item.is_file() and item.suffix.lower() in TEXT_SUFFIXES and "AD-creative" not in item.parts
    ]
    return sorted(files)


def read_material_text(path: Path, max_chars: int = 12000) -> str:
    try:
        text = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        text = path.read_text(encoding="utf-8", errors="ignore")
    return text[:max_chars]


def clean_material_line(line: str) -> str:
    line = line.strip().strip("|").strip()
    line = re.sub(r"^[-*#>`\s]+", "", line).strip()
    line = re.sub(r"\s+", " ", line)
    return line


def classify_requirement(statement: str) -> tuple[str, str, str]:
    text = statement.lower()
    if any(token in statement for token in ["不要", "禁区", "不能", "未经授权", "假 logo", "假logo"]):
        return "constraint", "high", "client_review"
    if any(token in statement for token in ["交付", "PPT", "ppt", "可编辑", "SlideSpec"]):
        return "delivery", "high", "ppt_gate"
    if any(token in statement for token in ["参考", "moodboard", "视频链接", "摄影参考"]):
        return "research", "high", "reference_research"
    if any(token in statement for token in ["画面", "视觉", "关键帧", "产品", "logo", "人物", "场景", "颜色"]):
        return "visual", "high", "visual_plan"
    if any(token in statement for token in ["方向", "主张", "创意", "情绪", "功能"]):
        return "creative", "high", "creative"
    if "ai" in text or "image" in text:
        return "visual_policy", "medium", "visual_review"
    return "brief", "medium", "intake"


def extract_requirement_statements(materials: list[tuple[Path, str]]) -> list[tuple[Path, str]]:
    candidates: list[tuple[Path, str]] = []
    seen: set[str] = set()
    for path, text in materials:
        for raw_line in text.splitlines():
            line = clean_material_line(raw_line)
            if not line or line in SKIP_MATERIAL_LINES:
                continue
            if len(line) < 6 or len(line) > 180:
                continue
            if not REQUIREMENT_TRIGGER_PATTERN.search(line):
                continue
            normalized = line.lower()
            if normalized in seen:
                continue
            seen.add(normalized)
            candidates.append((path, line))
    return candidates[:16]


def gap_templates(requirements: list[dict[str, str]], all_text: str) -> list[dict[str, str]]:
    """Deprecated legacy heuristic; production gaps come only from Fact Inventory."""
    del requirements, all_text
    return []


def normalize_profile_key(value: str) -> str:
    normalized = re.sub(r"\s+", "", value.strip().lower())
    return re.sub(r"[^\w\u4e00-\u9fff]+", "", normalized)


def stable_profile_id(prefix: str, *parts: str) -> str:
    raw = "|".join(normalize_profile_key(part) for part in parts if part)
    digest = hashlib.sha1(raw.encode("utf-8")).hexdigest()[:10].upper()
    return f"{prefix}-{digest}"


def join_unique_values(*values: str) -> str:
    seen: set[str] = set()
    merged: list[str] = []
    for value in values:
        for item in value.split(";"):
            cleaned = item.strip()
            if cleaned and cleaned not in seen:
                seen.add(cleaned)
                merged.append(cleaned)
    return ";".join(merged)


def quote_excerpt(value: str, limit: int = 96) -> str:
    cleaned = re.sub(r"\s+", " ", value.strip())
    if len(cleaned) <= limit:
        return cleaned
    return cleaned[: limit - 1].rstrip() + "…"


def split_speaker_line(line: str) -> tuple[str, str] | None:
    cleaned = clean_material_line(line)
    match = SPEAKER_LINE_PATTERN.match(cleaned)
    if not match:
        return None
    speaker = (match.group("bracket") or match.group("label") or "").strip()
    utterance = (match.group("body") or "").strip()
    if not speaker or not utterance or speaker in PROFILE_SOURCE_LABELS:
        return None
    if len(utterance) < 4:
        return None
    return speaker, utterance


def collect_source_materials(project: Path, source_ids: list[str]) -> tuple[list[tuple[dict[str, str], Path, str]], list[str]]:
    _, source_rows = read_csv_rows(project / "AD-creative/orchestrator/source_events.csv")
    source_id_set = set(source_ids)
    target_sources = [
        row for row in source_rows if not source_id_set or row.get("source_event_id") in source_id_set
    ]
    materials: list[tuple[dict[str, str], Path, str]] = []
    resolved_source_ids: list[str] = []
    for source in target_sources:
        source_id = source.get("source_event_id", "")
        if source_id:
            resolved_source_ids.append(source_id)
        for file_path, _ in source_row_files(project, source):
            materials.append((source, file_path, read_material_text(file_path, max_chars=24000)))
    return materials, resolved_source_ids


def ensure_profile_knowledge_base(project: Path) -> None:
    profile_dir = project / "AD-creative/orchestrator/profile_knowledge"
    profile_dir.mkdir(parents=True, exist_ok=True)
    csv_specs = {
        "profile_subjects.csv": PROFILE_SUBJECT_FIELDS,
        "meeting_voice_map.csv": PROFILE_VOICE_FIELDS,
        "profile_insights.csv": PROFILE_INSIGHT_FIELDS,
        "profile_conflicts.csv": PROFILE_CONFLICT_FIELDS,
    }
    for filename, fields in csv_specs.items():
        path = profile_dir / filename
        if path.exists():
            ensure_csv_fields(path, fields)
        else:
            write_csv_rows(path, fields, [])
    current_truth = profile_dir / "profile_current_truth.md"
    if not current_truth.exists():
        write_text(
            current_truth,
            """# Profile Current Truth

## Participant Profiles
- 暂无会议画像。

## Brand / Company Profiles
- 暂无品牌或公司画像。

## Demand And Decision Map
- 暂无。

## Conflicts
- 暂无。

## Next Confirmation
- 等待会议资料或客户资料输入。
""",
        )


def profile_signal(value: str, pattern: re.Pattern[str]) -> str:
    return "yes" if pattern.search(value) else ""


def profile_insight_types(text: str, *, subject_type: str) -> list[str]:
    insight_types: list[str] = []
    if NEED_SIGNAL_PATTERN.search(text):
        insight_types.append("need")
    if PREFERENCE_SIGNAL_PATTERN.search(text):
        insight_types.append("preference")
    if CONCERN_SIGNAL_PATTERN.search(text):
        insight_types.append("concern")
    if DECISION_SIGNAL_PATTERN.search(text):
        insight_types.append("decision_signal")
    if subject_type == "brand" or BRAND_SIGNAL_PATTERN.search(text):
        insight_types.append("brand_trait")
    if subject_type == "company" or COMPANY_SIGNAL_PATTERN.search(text):
        insight_types.append("company_context")
    return insight_types or ["observation"]


def infer_decision_power(name: str, utterances: list[str]) -> str:
    combined = f"{name}\n" + "\n".join(utterances)
    if DECISION_SIGNAL_PATTERN.search(combined) or re.search(r"老板|CEO|CMO|负责人|总经理|总监|创始人", combined):
        return "high"
    if re.search(r"经理|主管|客户|甲方|品牌", combined):
        return "medium"
    return "unknown"


def infer_influence_level(count: int, total: int, decision_power: str) -> str:
    if decision_power == "high":
        return "high"
    if total <= 0:
        return "unknown"
    share = count / total
    if count >= 4 or share >= 0.35:
        return "high"
    if count >= 2 or share >= 0.15:
        return "medium"
    return "low"


def infer_traits(utterances: list[str]) -> str:
    joined = "\n".join(utterances)
    traits: list[str] = []
    if CONCERN_SIGNAL_PATTERN.search(joined):
        traits.append("谨慎，关注风险边界")
    if DECISION_SIGNAL_PATTERN.search(joined):
        traits.append("会推动拍板或最终确认")
    if BRAND_SIGNAL_PATTERN.search(joined):
        traits.append("重视品牌调性和产品露出")
    if re.search(r"时间|预算|交付|效率|周期", joined):
        traits.append("关注执行效率和交付约束")
    if PREFERENCE_SIGNAL_PATTERN.search(joined):
        traits.append("对风格有明确偏好")
    return "；".join(traits) if traits else "信息不足，需继续观察"


def summarize_signals(utterances: list[str], pattern: re.Pattern[str], fallback: str) -> str:
    matches = [quote_excerpt(item, 60) for item in utterances if pattern.search(item)]
    return "；".join(matches[:3]) if matches else fallback


def profile_subject_row(
    *,
    subject_type: str,
    name: str,
    source_ids: str,
    now: str,
    utterances: list[str],
    total_utterances: int = 0,
    organization: str = "",
    role_or_title: str = "",
) -> dict[str, str]:
    decision_power = infer_decision_power(name, utterances) if subject_type == "participant" else "unknown"
    influence_level = (
        infer_influence_level(len(utterances), total_utterances or len(utterances), decision_power)
        if subject_type == "participant"
        else "unknown"
    )
    return {
        "subject_id": stable_profile_id("PROF", subject_type, name),
        "subject_type": subject_type,
        "name": name,
        "role_or_title": role_or_title,
        "organization": organization,
        "source_event_ids": source_ids,
        "first_seen_at": now,
        "last_seen_at": now,
        "profile_status": "candidate",
        "influence_level": influence_level,
        "decision_power": decision_power,
        "traits": infer_traits(utterances) if utterances else "待补充",
        "needs": summarize_signals(utterances, NEED_SIGNAL_PATTERN, "待补充"),
        "preferences": summarize_signals(utterances, PREFERENCE_SIGNAL_PATTERN, "待补充"),
        "concerns": summarize_signals(utterances, CONCERN_SIGNAL_PATTERN, "待补充"),
        "notes": "由会议/客户资料自动整理；需要人工确认后升级为 confirmed。",
    }


def merge_profile_subject(existing: dict[str, str], incoming: dict[str, str]) -> dict[str, str]:
    merged = dict(existing)
    for key in ["name", "subject_type", "role_or_title", "organization"]:
        if incoming.get(key):
            merged[key] = incoming[key]
    merged["source_event_ids"] = join_unique_values(existing.get("source_event_ids", ""), incoming.get("source_event_ids", ""))
    merged["last_seen_at"] = incoming.get("last_seen_at", existing.get("last_seen_at", ""))
    for key in ["traits", "needs", "preferences", "concerns", "notes"]:
        old = existing.get(key, "")
        new = incoming.get(key, "")
        if old in {"", "待补充", "信息不足，需继续观察"}:
            merged[key] = new
        elif new and new not in old:
            merged[key] = join_unique_values(old.replace("；", ";"), new.replace("；", ";")).replace(";", "；")
    level_order = {"unknown": 0, "low": 1, "medium": 2, "high": 3}
    for key in ["influence_level", "decision_power"]:
        if level_order.get(incoming.get(key, "unknown"), 0) > level_order.get(existing.get(key, "unknown"), 0):
            merged[key] = incoming[key]
    return merged


def profile_priority(insight_type: str) -> str:
    if insight_type in {"decision_signal", "concern"}:
        return "high"
    if insight_type in {"need", "brand_trait", "company_context"}:
        return "medium"
    return "low"


def add_profile_insight(
    rows_by_id: dict[str, dict[str, str]],
    *,
    subject_id: str,
    subject_type: str,
    source_id: str,
    file_path: str,
    insight_type: str,
    statement: str,
    evidence: str,
    now: str,
) -> tuple[int, int]:
    insight_id = stable_profile_id("INS", subject_id, insight_type, statement)
    row = {
        "insight_id": insight_id,
        "subject_id": subject_id,
        "subject_type": subject_type,
        "source_event_id": source_id,
        "file_path": file_path,
        "insight_type": insight_type,
        "statement": statement,
        "evidence_quote": quote_excerpt(evidence),
        "confidence": "0.68" if insight_type in {"need", "preference", "concern", "decision_signal"} else "0.56",
        "status": "candidate",
        "priority": profile_priority(insight_type),
        "linked_requirement_ids": "",
        "supersedes_insight_id": "",
        "created_at": now,
        "updated_at": now,
    }
    if insight_id in rows_by_id:
        existing = rows_by_id[insight_id]
        existing["source_event_id"] = source_id or existing.get("source_event_id", "")
        existing["file_path"] = file_path or existing.get("file_path", "")
        existing["evidence_quote"] = row["evidence_quote"] or existing.get("evidence_quote", "")
        existing["updated_at"] = now
        return 0, 1
    rows_by_id[insight_id] = row
    return 1, 0


def profile_current_truth_content(
    project: Path,
    subjects: list[dict[str, str]],
    insights: list[dict[str, str]],
    conflicts: list[dict[str, str]],
) -> str:
    participant_rows = [row for row in subjects if row.get("subject_type") == "participant"]
    org_rows = [row for row in subjects if row.get("subject_type") in {"brand", "company", "client_group"}]
    top_participants = sorted(
        participant_rows,
        key=lambda row: (
            {"high": 3, "medium": 2, "low": 1}.get(row.get("decision_power", ""), 0),
            {"high": 3, "medium": 2, "low": 1}.get(row.get("influence_level", ""), 0),
        ),
        reverse=True,
    )[:8]
    demand_insights = [row for row in insights if row.get("insight_type") in {"need", "preference", "concern", "decision_signal"}]
    participant_table = "\n".join(
        f"| {md_cell(row.get('name', ''))} | {row.get('influence_level', '')} | {row.get('decision_power', '')} | {md_cell(row.get('needs', ''))} | {md_cell(row.get('concerns', ''))} |"
        for row in top_participants
    ) or "| - | - | - | 暂无 | 暂无 |"
    org_table = "\n".join(
        f"| {md_cell(row.get('subject_type', ''))} | {md_cell(row.get('name', ''))} | {md_cell(row.get('traits', ''))} | {md_cell(row.get('preferences', ''))} |"
        for row in org_rows[:8]
    ) or "| - | - | 暂无 | 暂无 |"
    demand_rows = "\n".join(
        f"- {row.get('insight_type')}: {row.get('statement')}（证据：{row.get('evidence_quote')}）"
        for row in demand_insights[:10]
    ) or "- 暂无"
    conflict_rows = "\n".join(
        f"- {row.get('topic')}: {row.get('conflict_summary')}；建议：{row.get('recommended_resolution')}"
        for row in conflicts[:8]
    ) or "- 暂无显式分歧"
    return f"""# Profile Current Truth

## Project
{project.name}

## Participant Profiles
| 人 | 影响力 | 决策权 | 想要什么 | 担心什么 |
|---|---|---|---|---|
{participant_table}

## Brand / Company Profiles
| 类型 | 名称 | 特点 | 偏好 |
|---|---|---|---|
{org_table}

## Demand And Decision Map
{demand_rows}

## Conflicts
{conflict_rows}

## How To Use
- 先满足高决策权/高影响力人物的明确诉求。
- 有分歧时不要硬合并成一句口号，先列出双方关心点，再给折中方案。
- 证据不足的画像只能作为 candidate，客户确认后再升级为 confirmed。

## Next Confirmation
- 请确认关键发言人姓名/职务、最终拍板人、品牌禁区，以及哪些分歧已经内部统一。
"""


def profile_handoff_content(
    subjects: list[dict[str, str]],
    insights: list[dict[str, str]],
    conflicts: list[dict[str, str]],
) -> str:
    decision_people = [
        row for row in subjects if row.get("subject_type") == "participant" and row.get("decision_power") in {"high", "medium"}
    ][:6]
    decision_rows = "\n".join(
        f"| {md_cell(row.get('name', ''))} | {row.get('decision_power', '')} | {row.get('influence_level', '')} | {md_cell(row.get('traits', ''))} |"
        for row in decision_people
    ) or "| - | - | - | 暂无 |"
    need_rows = "\n".join(
        f"- {row.get('statement')}" for row in insights if row.get("insight_type") in {"need", "preference"}
    ) or "- 暂无明确诉求"
    concern_rows = "\n".join(
        f"- {row.get('statement')}" for row in insights if row.get("insight_type") == "concern"
    ) or "- 暂无明确担心"
    conflict_rows = "\n".join(
        f"- {row.get('conflict_summary')} 建议：{row.get('recommended_resolution')}" for row in conflicts
    ) or "- 暂无显式分歧"
    return f"""# 画像分析简报

## 这次会议看出了什么
{need_rows}

## 谁更需要重点照顾
| 人 | 决策权 | 影响力 | 特点 |
|---|---|---|---|
{decision_rows}

## 他们担心什么
{concern_rows}

## 分歧怎么合
{conflict_rows}

## 下一步怎么用
- 写方案前先看高决策权人物的诉求。
- 研究阶段优先补齐有争议、证据不足、会影响客户拍板的信息。
- 方案里同时照顾品牌特点、公司内部共识和具体发言人的担心点。
"""


def analyze_profiles(
    project: Path,
    *,
    source_ids: list[str] | None = None,
    work_id: str = "",
    goal: str = "",
    brand: str = "",
    company: str = "",
    client: str = "",
) -> dict[str, object]:
    ensure_profile_knowledge_base(project)
    source_ids = source_ids or []
    materials, resolved_source_ids = collect_source_materials(project, source_ids)
    now = now_iso()
    linked_source_ids = ";".join(source_ids or resolved_source_ids)
    profile_dir = project / "AD-creative/orchestrator/profile_knowledge"

    subject_fields, subject_rows = read_csv_rows(profile_dir / "profile_subjects.csv")
    voice_fields, voice_rows = read_csv_rows(profile_dir / "meeting_voice_map.csv")
    insight_fields, insight_rows = read_csv_rows(profile_dir / "profile_insights.csv")
    conflict_fields, conflict_rows = read_csv_rows(profile_dir / "profile_conflicts.csv")
    subjects_by_id = {row.get("subject_id", ""): row for row in subject_rows if row.get("subject_id")}
    voices_by_id = {row.get("voice_id", ""): row for row in voice_rows if row.get("voice_id")}
    insights_by_id = {row.get("insight_id", ""): row for row in insight_rows if row.get("insight_id")}
    conflicts_by_id = {row.get("conflict_id", ""): row for row in conflict_rows if row.get("conflict_id")}

    speaker_utterances: dict[str, list[str]] = {}
    speaker_source_ids: dict[str, str] = {}
    speaker_files: dict[str, str] = {}
    new_voices = 0
    deduped = 0
    new_insights = 0
    updated_insights = 0
    conflict_candidates: list[tuple[str, str, str, str]] = []

    for source, file_path, text in materials:
        source_id = source.get("source_event_id", "")
        rel_file = source_path_label(project, source, file_path)
        for raw_line in text.splitlines():
            line = clean_material_line(raw_line)
            if not line or len(line) < 4:
                continue
            speaker_line = split_speaker_line(line)
            if speaker_line:
                speaker, utterance = speaker_line
                speaker_utterances.setdefault(speaker, []).append(utterance)
                speaker_source_ids[speaker] = join_unique_values(speaker_source_ids.get(speaker, ""), source_id)
                speaker_files[speaker] = join_unique_values(speaker_files.get(speaker, ""), rel_file)
                voice_id = stable_profile_id("VOICE", source_id, speaker, utterance)
                voice_row = {
                    "voice_id": voice_id,
                    "source_event_id": source_id,
                    "file_path": rel_file,
                    "speaker": speaker,
                    "utterance": utterance,
                    "need_signal": profile_signal(utterance, NEED_SIGNAL_PATTERN),
                    "preference_signal": profile_signal(utterance, PREFERENCE_SIGNAL_PATTERN),
                    "concern_signal": profile_signal(utterance, CONCERN_SIGNAL_PATTERN),
                    "decision_signal": profile_signal(utterance, DECISION_SIGNAL_PATTERN),
                    "influence_level": "",
                    "decision_power": "",
                    "evidence_quote": quote_excerpt(utterance),
                    "confidence": "0.72",
                    "status": "candidate",
                }
                if voice_id in voices_by_id:
                    voices_by_id[voice_id].update(voice_row)
                    deduped += 1
                else:
                    voices_by_id[voice_id] = voice_row
                    new_voices += 1
                if CONFLICT_SIGNAL_PATTERN.search(utterance):
                    conflict_candidates.append((source_id, rel_file, speaker, utterance))
            elif BRAND_SIGNAL_PATTERN.search(line) or COMPANY_SIGNAL_PATTERN.search(line) or CONFLICT_SIGNAL_PATTERN.search(line):
                subject_type = "brand" if BRAND_SIGNAL_PATTERN.search(line) else "company"
                subject_name = brand if subject_type == "brand" and brand else company if company else client or "未命名客户"
                subject_id = stable_profile_id("PROF", subject_type, subject_name)
                if subject_id not in subjects_by_id:
                    subjects_by_id[subject_id] = profile_subject_row(
                        subject_type=subject_type,
                        name=subject_name,
                        source_ids=source_id,
                        now=now,
                        utterances=[line],
                        organization=company or client,
                    )
                else:
                    incoming = profile_subject_row(
                        subject_type=subject_type,
                        name=subject_name,
                        source_ids=source_id,
                        now=now,
                        utterances=[line],
                        organization=company or client,
                    )
                    subjects_by_id[subject_id] = merge_profile_subject(subjects_by_id[subject_id], incoming)
                for insight_type in profile_insight_types(line, subject_type=subject_type):
                    added, updated = add_profile_insight(
                        insights_by_id,
                        subject_id=subject_id,
                        subject_type=subject_type,
                        source_id=source_id,
                        file_path=rel_file,
                        insight_type=insight_type,
                        statement=line,
                        evidence=line,
                        now=now,
                    )
                    new_insights += added
                    updated_insights += updated
                if CONFLICT_SIGNAL_PATTERN.search(line):
                    conflict_candidates.append((source_id, rel_file, subject_name, line))

    total_voice_count = sum(len(items) for items in speaker_utterances.values())
    for speaker, utterances in speaker_utterances.items():
        subject_id = stable_profile_id("PROF", "participant", speaker)
        incoming = profile_subject_row(
            subject_type="participant",
            name=speaker,
            source_ids=speaker_source_ids.get(speaker, linked_source_ids),
            now=now,
            utterances=utterances,
            total_utterances=total_voice_count,
            organization=company or client,
        )
        if subject_id in subjects_by_id:
            subjects_by_id[subject_id] = merge_profile_subject(subjects_by_id[subject_id], incoming)
        else:
            subjects_by_id[subject_id] = incoming
        for voice in voices_by_id.values():
            if voice.get("speaker") == speaker:
                voice["influence_level"] = subjects_by_id[subject_id].get("influence_level", "")
                voice["decision_power"] = subjects_by_id[subject_id].get("decision_power", "")
        for utterance in utterances:
            for insight_type in profile_insight_types(utterance, subject_type="participant"):
                added, updated = add_profile_insight(
                    insights_by_id,
                    subject_id=subject_id,
                    subject_type="participant",
                    source_id=speaker_source_ids.get(speaker, "").split(";")[0],
                    file_path=speaker_files.get(speaker, "").split(";")[0],
                    insight_type=insight_type,
                    statement=utterance,
                    evidence=utterance,
                    now=now,
                )
                new_insights += added
                updated_insights += updated

    for subject_type, name in [("brand", brand), ("company", company), ("client_group", client)]:
        if not name:
            continue
        subject_id = stable_profile_id("PROF", subject_type, name)
        incoming = profile_subject_row(
            subject_type=subject_type,
            name=name,
            source_ids=linked_source_ids,
            now=now,
            utterances=[],
            organization=company if subject_type != "company" else name,
        )
        if subject_id in subjects_by_id:
            subjects_by_id[subject_id] = merge_profile_subject(subjects_by_id[subject_id], incoming)
        else:
            subjects_by_id[subject_id] = incoming

    for source_id, rel_file, speaker, utterance in conflict_candidates:
        conflict_id = stable_profile_id("CONF", source_id, utterance)
        subject_id = stable_profile_id("PROF", "participant", speaker)
        row = {
            "conflict_id": conflict_id,
            "topic": "会议分歧 / 需要统一",
            "source_event_ids": source_id,
            "subject_ids": subject_id if subject_id in subjects_by_id else "",
            "conflict_summary": utterance,
            "recommended_resolution": "先确认最终拍板人；把双方担心点拆成必须满足、可折中、暂缓三类。",
            "status": "candidate",
            "confidence": "0.62",
            "evidence_quotes": quote_excerpt(f"{speaker}: {utterance}"),
            "created_at": now,
            "updated_at": now,
        }
        if conflict_id in conflicts_by_id:
            conflicts_by_id[conflict_id]["updated_at"] = now
            deduped += 1
        else:
            conflicts_by_id[conflict_id] = row

    subjects = sorted(subjects_by_id.values(), key=lambda row: (row.get("subject_type", ""), row.get("name", "")))
    voices = sorted(voices_by_id.values(), key=lambda row: row.get("voice_id", ""))
    insights = sorted(insights_by_id.values(), key=lambda row: (row.get("priority", ""), row.get("insight_id", "")), reverse=True)
    conflicts = sorted(conflicts_by_id.values(), key=lambda row: row.get("conflict_id", ""))
    write_csv_rows(profile_dir / "profile_subjects.csv", subject_fields, subjects)
    write_csv_rows(profile_dir / "meeting_voice_map.csv", voice_fields, voices)
    write_csv_rows(profile_dir / "profile_insights.csv", insight_fields, insights)
    write_csv_rows(profile_dir / "profile_conflicts.csv", conflict_fields, conflicts)

    profile_truth_path = profile_dir / "profile_current_truth.md"
    handoff_path = project / "AD-creative/handoff/画像分析简报.md"
    write_text(profile_truth_path, profile_current_truth_content(project, subjects, insights, conflicts))
    write_text(handoff_path, profile_handoff_content(subjects, insights, conflicts))
    delivery_surface = project_surface(project) == DELIVERY_SURFACE
    if delivery_surface:
        for artifact_id, artifact_type, rel_path in [
            ("ART-AUTO-PROFILE-TRUTH", "profile_current_truth", safe_rel(project, profile_truth_path)),
            ("ART-AUTO-PROFILE-BRIEF", "profile_handoff_brief", safe_rel(project, handoff_path)),
        ]:
            update_artifact(
                project,
                artifact_id,
                artifact_type,
                rel_path,
                "intake",
                source_event_ids=linked_source_ids,
                linked_work_items=work_id,
                gate_status="PARTIAL_PASS",
            )
        append_gate(
            project,
            "GATE-AUTO-PROFILE-001",
            "intake",
            "PARTIAL_PASS",
            "70",
            "ART-AUTO-PROFILE-TRUTH;ART-AUTO-PROFILE-BRIEF",
            "画像均为 candidate，需人工确认关键角色和最终拍板人。",
            "确认关键发言人、决策权、分歧是否已统一。",
            "请确认谁是最终拍板人、哪些画像结论可升级为 confirmed。",
            "research_plan",
            "ad_creative_operator",
        )
        append_event(
            project,
            {
                "event_id": f"EVT-PROFILE-{now}",
                "event_type": "profile_analysis_completed",
                "created_at": now,
                "actor": "ad_creative_operator",
                "source_event_ids": linked_source_ids,
                "goal": goal,
                "subjects": len(subjects),
                "insights": len(insights),
                "conflicts": len(conflicts),
            },
        )
    if delivery_surface and work_id:
        work_path = project / "AD-creative/orchestrator/work_items.csv"
        work_fields, work_rows = read_csv_rows(work_path)
        for row in work_rows:
            if row.get("work_id") == work_id:
                row["status"] = "done"
                row["output_artifacts"] = "ART-AUTO-PROFILE-TRUTH;ART-AUTO-PROFILE-BRIEF"
                row["linked_source_events"] = join_unique_values(row.get("linked_source_events", ""), linked_source_ids)
                row["updated_at"] = now
                break
        write_csv_rows(work_path, work_fields, work_rows)
    return {
        "materials": len(materials),
        "subjects": len(subjects),
        "voices": len(voices),
        "insights": len(insights),
        "conflicts": len(conflicts),
        "new_voices": new_voices,
        "new_insights": new_insights,
        "updated_insights": updated_insights,
        "deduped": deduped,
        "profile_current_truth": profile_truth_path,
        "handoff": handoff_path,
        "governance_records_written": int(delivery_surface),
        "source_ids": linked_source_ids,
    }


def unique_rows(rows: list[dict[str, str]], key: str) -> set[str]:
    return {row.get(key, "").strip() for row in rows if row.get(key, "").strip()}


def owner_for_statement(statement: str) -> str:
    return "client" if CLIENT_OWNER_PATTERN.search(statement) else "operator"


def first_existing_id(project: Path, rel_path: str, key: str) -> str:
    _, rows = read_csv_rows(project / rel_path)
    for row in rows:
        value = row.get(key, "").strip()
        if value:
            return value
    return ""


def update_artifact(
    project: Path,
    artifact_id: str,
    artifact_type: str,
    rel_path: str,
    stage: str,
    *,
    status: str = "done",
    visibility: str = "internal_only",
    source_event_ids: str = "",
    linked_requirements: str = "",
    linked_work_items: str = "",
    linked_references: str = "",
    linked_assets: str = "",
    gate_status: str = "PASS",
    version: str = "v001",
    derived_from_artifact_id: str = "",
    derived_from_sha256: str = "",
    sha256: str = "",
    size_bytes: str = "",
) -> None:
    update_or_append_csv_row(
        project / "AD-creative/orchestrator/artifact_index.csv",
        "artifact_id",
        {
            "artifact_id": artifact_id,
            "artifact_type": artifact_type,
            "path": rel_path,
            "stage": stage,
            "version": version,
            "status": status,
            "visibility": visibility,
            "source_event_ids": source_event_ids,
            "linked_requirements": linked_requirements,
            "linked_work_items": linked_work_items,
            "linked_references": linked_references,
            "linked_assets": linked_assets,
            "gate_status": gate_status,
            "supersedes_artifact_id": "",
            "created_at": now_iso(),
            "updated_at": now_iso(),
            "sha256": sha256,
            "size_bytes": size_bytes,
            "derived_from_artifact_id": derived_from_artifact_id,
            "derived_from_sha256": derived_from_sha256,
            "lifecycle_state": normalized_artifact_lifecycle({"status": status}),
            "original_path": rel_path,
            "cleanup_ref": "",
            "removed_at": "",
            "removal_reason": "",
            "superseded_by": "",
            "status_reason": status,
        },
    )


def append_gate(
    project: Path,
    gate_id: str,
    stage: str,
    status: str,
    score: str,
    checked_artifacts: str,
    blocking_issues: str,
    revision_items: str,
    questions: str,
    next_state: str,
    owner: str,
    target_ref: str = "",
    target_sha256: str = "",
) -> None:
    path = project / "AD-creative/orchestrator/gate_log.csv"
    fields, rows = normalize_gate_log_schema(path)
    if not target_ref:
        _, artifacts = read_csv_rows(
            project / "AD-creative/orchestrator/artifact_index.csv"
        )
        artifacts_by_id = {
            row.get("artifact_id", "").strip(): row
            for row in artifacts
            if row.get("artifact_id", "").strip()
        }
        for artifact_id in split_asset_refs(checked_artifacts):
            artifact = artifacts_by_id.get(artifact_id)
            if not artifact:
                continue
            try:
                candidate = contained_project_path(
                    project,
                    artifact.get("path", ""),
                    f"gate target artifact {artifact_id}",
                )
            except ValueError:
                continue
            if candidate.is_file():
                target_ref = safe_rel(project, candidate)
                target_sha256 = file_sha256(candidate)
                break
    previous = next(
        (row for row in reversed(rows) if row.get("gate_id") == gate_id), None
    )
    run_id = next_id(rows, "gate_run_id", "GATE-RUN")
    evidence_snapshot_ref = ""
    evidence_snapshot_sha256 = ""
    if target_ref:
        target = contained_project_path(project, target_ref, "gate target_ref")
        if not target.is_file():
            raise ValueError(f"gate target does not exist: {target_ref}")
        actual_target_sha = file_sha256(target)
        if target_sha256 and target_sha256 != actual_target_sha:
            raise ValueError("gate target_sha256 does not match target_ref")
        target_sha256 = actual_target_sha
        snapshot_suffix = target.suffix if target.suffix else ".evidence"
        snapshot = (
            project
            / "AD-creative/gates/history"
            / safe_artifact_suffix(gate_id)
            / f"{safe_artifact_suffix(run_id)}_{target_sha256[:12]}{snapshot_suffix}"
        )
        snapshot.parent.mkdir(parents=True, exist_ok=True)
        if snapshot.exists():
            if file_sha256(snapshot) != target_sha256:
                raise FileExistsError(
                    f"immutable gate evidence snapshot collision: {snapshot}"
                )
        else:
            shutil.copy2(target, snapshot)
        evidence_snapshot_ref = safe_rel(project, snapshot)
        evidence_snapshot_sha256 = file_sha256(snapshot)
    rows.append(
        {
            "gate_id": gate_id,
            "gate_run_id": run_id,
            "stage": stage,
            "status": status,
            "score": score,
            "checked_artifacts": checked_artifacts,
            "target_ref": target_ref,
            "target_sha256": target_sha256,
            "evidence_snapshot_ref": evidence_snapshot_ref,
            "evidence_snapshot_sha256": evidence_snapshot_sha256,
            "blocking_issues": blocking_issues,
            "revision_items": revision_items,
            "questions": questions,
            "next_state": next_state,
            "created_at": now_iso(),
            "owner": owner,
            "supersedes_gate_run_id": (
                previous.get("gate_run_id", "") if previous else ""
            ),
        },
    )
    write_csv_rows(path, fields, rows)


POLLUTION_DIR_NAMES = {"__pycache__", ".pytest_cache", ".mypy_cache", ".ruff_cache"}
POLLUTION_FILE_SUFFIXES = {".pyc", ".pyo"}
POLLUTION_FILE_NAMES = {".DS_Store"}
THREAD_TERMINAL_STATES = {"archived", "closed"}


def find_pollution_paths(root: Path, limit: int = 80) -> list[str]:
    if not root.exists():
        return []
    findings: list[str] = []
    for path in root.rglob("*"):
        if ".git" in path.parts:
            continue
        if path.name in POLLUTION_DIR_NAMES or path.name in POLLUTION_FILE_NAMES or path.suffix in POLLUTION_FILE_SUFFIXES:
            findings.append(safe_rel(root, path))
            if len(findings) >= limit:
                break
    return findings


def git_status_for(root: Path) -> tuple[str, list[str], list[str]]:
    try:
        git_root = subprocess.run(
            ["git", "-C", str(root), "rev-parse", "--show-toplevel"],
            check=True,
            text=True,
            capture_output=True,
        ).stdout.strip()
    except (subprocess.CalledProcessError, FileNotFoundError):
        return "", [], []
    completed = subprocess.run(
        ["git", "-C", git_root, "status", "--short", "--untracked-files=all"],
        check=True,
        text=True,
        capture_output=True,
    )
    lines = [line for line in completed.stdout.splitlines() if line.strip()]
    tracked = [line for line in lines if not line.startswith("??")]
    untracked = [line for line in lines if line.startswith("??")]
    return git_root, tracked, untracked


def active_thread_registry_rows(project: Path) -> list[dict[str, str]]:
    _, rows = read_csv_rows(project / "AD-creative/orchestrator/thread_registry.csv")
    active: list[dict[str, str]] = []
    for row in rows:
        thread_id = row.get("thread_id", "").strip()
        archived = row.get("archived", "").strip().lower()
        lifecycle = row.get("lifecycle_state", "").strip().lower()
        if not thread_id or thread_id.startswith("planned:"):
            continue
        if archived in {"true", "yes", "1"} or lifecycle in THREAD_TERMINAL_STATES:
            continue
        active.append(row)
    return active


def workspace_hygiene_report(project: Path) -> dict[str, object]:
    project = project.resolve()
    git_root, tracked, untracked = git_status_for(project)
    pollution = find_pollution_paths(project)
    active_threads = active_thread_registry_rows(project)
    issues: list[str] = []
    if tracked:
        issues.append(f"tracked git changes: {len(tracked)}")
    if untracked:
        issues.append(f"untracked git files: {len(untracked)}")
    if pollution:
        issues.append(f"cache/temp pollution paths: {len(pollution)}")
    if active_threads:
        issues.append(f"active thread registry rows: {len(active_threads)}")
    plan = [
        "把验证和草稿生成放在 /tmp 或 AD-creative/workspaces/<work_id>/，不要写到仓库根目录。",
        "每次大任务结束跑 git status --short --untracked-files=all，并清理 __pycache__ / .pytest_cache / .mypy_cache / .ruff_cache / *.pyc / *.pyo / .DS_Store。",
        "Codex Thread 结果合并后立即归档，并把 cleanup_action / archived 写回 thread_registry.csv。",
        "源码模板和 packaged mirror 必须一起更新，并跑 check_packaged_assets。",
        "用户未要求 commit 前，只报告修改文件；不要 reset、checkout 或删除用户资料。",
    ]
    return {
        "status": "PASS" if not issues else "CHECK",
        "project": str(project),
        "git_root": git_root,
        "tracked_changes": tracked,
        "untracked_files": untracked,
        "pollution_paths": pollution,
        "active_threads": [
            {
                "thread_id": row.get("thread_id", ""),
                "title": row.get("title", ""),
                "role": row.get("role", ""),
                "lifecycle_state": row.get("lifecycle_state", ""),
                "archived": row.get("archived", ""),
            }
            for row in active_threads
        ],
        "issues": issues,
        "plan": plan,
    }


def host_platform(url: str) -> str:
    host = urllib.parse.urlparse(url).netloc.lower()
    if not host:
        return "unknown"
    return host.removeprefix("www.")


def check_url(url: str, timeout: float = 8.0) -> tuple[bool, str]:
    parsed = urllib.parse.urlparse(url)
    if parsed.scheme != "https" or not parsed.netloc:
        return False, "URL 必须是 https。"
    request = urllib.request.Request(url, method="HEAD", headers={"User-Agent": "ad-creative-orchestrator/1.0"})
    try:
        with urllib.request.urlopen(request, timeout=timeout) as response:
            status = getattr(response, "status", 200)
            return 200 <= int(status) < 400, f"HTTP {status}"
    except urllib.error.HTTPError as exc:
        if exc.code == 405:
            try:
                fallback = urllib.request.Request(url, headers={"User-Agent": "ad-creative-orchestrator/1.0"})
                with urllib.request.urlopen(fallback, timeout=timeout) as response:
                    status = getattr(response, "status", 200)
                    return 200 <= int(status) < 400, f"HTTP {status}"
            except Exception as fallback_exc:  # noqa: BLE001 - report exact network failure
                return False, str(fallback_exc)
        return False, f"HTTP {exc.code}"
    except Exception as exc:  # noqa: BLE001 - report exact network failure
        return False, str(exc)


def render_reference_shortlist(project: Path) -> Path:
    _, refs = read_csv_rows(project / "AD-creative/references/reference_cards.csv")
    rows = "\n".join(
        "| {reference_id} | {title} | {platform} | {role} | {url} | {client_visible} |".format(
            **{key: row.get(key, "") for key in [
                "reference_id",
                "title",
                "platform",
                "role",
                "url",
                "client_visible",
            ]}
        )
        for row in refs
    )
    if not rows:
        rows = "| - | 暂无真实参考 | - | - | - | - |"
    path = project / "AD-creative/references/reference_shortlist.md"
    write_text(
        path,
        f"""# Reference Shortlist

status: real_links_registered
visibility: internal_only

| ID | 标题 | 平台 | 角色 | URL | 客户可见 |
| --- | --- | --- | --- | --- | --- |
{rows}

## Rules

- 只登记 https 链接。
- 客户稿引用前仍需确认版权、使用范围和截图来源。
- 未过 Research Gate 的参考只用于内部方向判断。
""",
    )
    return path


def add_reference(
    project: Path,
    url: str,
    title: str,
    role: str,
    reference_type: str,
    *,
    source_owner: str = "official_or_public",
    why_relevant: str = "",
    borrow: str = "",
    do_not_copy: str = "",
    client_visible: str = "false",
    live_check: bool = True,
) -> tuple[str, str]:
    ok, note = check_url(url) if live_check else (True, "live_check_skipped")
    if not ok:
        raise ValueError(f"reference URL check failed: {url}: {note}")
    ref_path = project / "AD-creative/references/reference_cards.csv"
    fields, refs = read_csv_rows(ref_path)
    if not fields:
        raise FileNotFoundError(f"CSV header not found: {ref_path}")
    for row in refs:
        if row.get("url") == url:
            ref_id = row.get("reference_id", "")
            return ref_id, "existing"
    ref_id = next_id(refs, "reference_id", "REF")
    source_id = first_existing_id(project, "AD-creative/orchestrator/source_events.csv", "source_event_id")
    refs.append(
        {
            "reference_id": ref_id,
            "source_event_id": source_id,
            "platform": host_platform(url),
            "url": url,
            "title": title or url,
            "source_owner": source_owner,
            "reference_type": reference_type,
            "role": role,
            "why_relevant": why_relevant,
            "borrow": borrow,
            "do_not_copy": do_not_copy or "不要复制构图、标志、人物或可识别画面；只抽取方向逻辑。",
            "client_visible": client_visible,
            "notes": note,
        }
    )
    write_csv_rows(ref_path, fields, refs)
    shortlist = render_reference_shortlist(project)
    update_artifact(
        project,
        "ART-AUTO-REFERENCE-SHORTLIST",
        "reference_pack",
        safe_rel(project, shortlist),
        "reference_research",
        linked_references=ref_id,
        gate_status="PASS",
    )
    append_gate(
        project,
        "GATE-AUTO-REFERENCE-001",
        "reference_research",
        "PASS",
        "90",
        "ART-AUTO-REFERENCE-SHORTLIST",
        "",
        "客户稿引用前确认版权和截图范围。",
        "",
        "ready_for_internal_creative",
        "ad_creative_operator",
    )
    append_event(
        project,
        {
            "event_id": f"EVT-{ref_id}",
            "event_type": "reference_registered",
            "created_at": now_iso(),
            "reference_id": ref_id,
            "url": url,
            "live_check": note,
        },
    )
    return ref_id, "created"


def compact_evidence(value: str, limit: int = 140) -> str:
    cleaned = re.sub(r"\s+", " ", clean_material_line(value)).strip()
    if len(cleaned) <= limit:
        return cleaned
    return cleaned[: limit - 1].rstrip() + "…"


def first_evidence_line(lines: list[tuple[str, str]], patterns: list[str], fallback: str) -> tuple[str, str]:
    for pattern in patterns:
        regex = re.compile(re.escape(pattern), re.IGNORECASE)
        for source_id, line in lines:
            if regex.search(line):
                return compact_evidence(line), source_id
    return fallback, ""


def open_question(label: str) -> str:
    return f"待确认：{label}"


def material_evidence_lines(source_materials: list[tuple[dict[str, str], Path, str]]) -> list[tuple[str, str]]:
    evidence: list[tuple[str, str]] = []
    for source, _, text in source_materials:
        source_id = source.get("source_event_id", "")
        for raw_line in text.splitlines():
            if re.match(r"^\s{0,3}#{1,6}\s+", raw_line):
                continue
            line = clean_material_line(raw_line)
            if not line or line in SKIP_MATERIAL_LINES:
                continue
            if 6 <= len(line) <= 220:
                evidence.append((source_id, line))
    return evidence


PROPOSAL_EVIDENCE_LABELS = {
    "business_problem": ("business problem", "business challenge", "商业问题", "业务问题", "创意问题", "挑战"),
    "client_objective": ("objective", "client objective", "目标", "客户目标", "商业任务", "项目目标", "任务目标"),
    "audience": ("audience", "target audience", "受众", "目标人群", "人群"),
    "barrier": ("barrier", "behavior barrier", "痛点", "障碍", "顾虑"),
    "insight": ("insight", "consumer insight", "洞察", "消费者洞察"),
    "product_feature": ("feature", "product feature", "卖点", "产品卖点", "产品功能", "产品事实", "已锁定产品事实", "功能"),
    "visual_reference": ("visual direction", "visual", "direction", "方向", "场景", "画面方向", "视觉方向", "关键视觉"),
}

PROPOSAL_EVIDENCE_SECTIONS = {
    "business_problem": ("business problem", "business challenge", "商业问题", "业务问题", "创意问题"),
    "client_objective": ("objective", "client objective", "商业任务", "项目目标", "客户目标", "任务目标"),
    "audience": ("audience", "target audience", "目标人群", "受众"),
    "barrier": ("barrier", "behavior barrier", "行为阻力", "痛点", "障碍"),
    "insight": ("insight", "consumer insight", "消费者洞察", "用户洞察"),
    "visual_reference": ("visual direction", "画面方向", "视觉方向", "关键视觉"),
}


def labeled_proposal_evidence(
    source_materials: list[tuple[dict[str, str], Path, str]],
) -> dict[str, tuple[str, str]]:
    """Extract only explicitly labelled brief facts; ignore deliverable/list pollution."""
    found: dict[str, tuple[str, str]] = {}
    label_map = {
        label.lower(): key
        for key, labels in PROPOSAL_EVIDENCE_LABELS.items()
        for label in labels
    }
    for source, _, text in source_materials:
        source_id = source.get("source_event_id", "")
        section = ""
        section_key = ""
        for raw in text.splitlines():
            heading = re.match(r"^\s{0,3}#{1,6}\s+(.+?)\s*$", raw)
            if heading:
                section = heading.group(1).strip().lower()
                section_key = next(
                    (
                        key
                        for key, names in PROPOSAL_EVIDENCE_SECTIONS.items()
                        if any(name in section for name in names)
                    ),
                    "",
                )
                continue
            match = re.match(r"^\s*(?:[-*]\s*)?([^:：]{1,40})[:：]\s*(.+?)\s*$", raw)
            if any(token in section for token in ("deliverable", "交付", "output")):
                continue
            if match:
                label = match.group(1).strip().lower()
                value = clean_material_line(match.group(2))
                key = label_map.get(label)
                if key and key not in found and value:
                    found[key] = (compact_evidence(value), source_id)
                continue
            value = clean_material_line(raw)
            if not value or re.match(r"^\d+[.)、]\s+", value):
                continue
            if "希望观众记住" in value and "insight" not in found:
                found["insight"] = (compact_evidence(value), source_id)
            elif section_key and section_key not in found:
                found[section_key] = (compact_evidence(value), source_id)
    return found


def collect_proposal_evidence(project: Path) -> dict[str, object]:
    _, requirement_rows = read_csv_rows(project / "AD-creative/orchestrator/requirements.csv")
    _, gap_rows = read_csv_rows(project / "AD-creative/orchestrator/gaps.csv")
    _, reference_rows = read_csv_rows(project / "AD-creative/references/reference_cards.csv")
    _, profile_insights = read_csv_rows(project / "AD-creative/orchestrator/profile_knowledge/profile_insights.csv")
    source_materials, resolved_source_ids = collect_source_materials(project, [])
    source_lines = material_evidence_lines(source_materials)
    labeled = labeled_proposal_evidence(source_materials)
    requirement_lines: list[tuple[str, str]] = []
    gap_lines: list[tuple[str, str]] = []
    for requirement in requirement_rows:
        statement = requirement.get("statement", "").strip()
        if statement:
            requirement_lines.append((requirement.get("source_event_id", ""), statement))
    for gap in gap_rows:
        description = gap.get("description", "").strip()
        question = gap.get("question_for_client", "").strip()
        if description:
            gap_lines.append((gap.get("linked_requirement_id", ""), description))
        if question:
            gap_lines.append((gap.get("linked_requirement_id", ""), question))
    creative_lines = source_lines + requirement_lines + gap_lines
    source_or_requirement_lines = source_lines + requirement_lines

    business_problem, business_source = labeled.get("business_problem", (open_question("business problem (internal-only)"), ""))
    client_objective, objective_source = labeled.get("client_objective", (open_question("client real objective (internal-only)"), ""))
    audience, audience_source = labeled.get("audience", (open_question("target audience (internal-only)"), ""))
    barrier, barrier_source = labeled.get("barrier", (open_question("target behavior barrier (internal-only)"), ""))
    product_feature, feature_source = labeled.get("product_feature", (open_question("product feature (internal-only)"), ""))
    visual_reference, visual_source = labeled.get("visual_reference", (open_question("visual/action evidence (internal-only)"), ""))
    insight, insight_source = labeled.get("insight", (open_question("consumer insight (internal-only)"), ""))
    for row in profile_insights:
        if not insight.startswith("TBD"):
            break
        if row.get("insight_type") in {"need", "preference", "concern", "brand_trait"}:
            insight = compact_evidence(row.get("statement", ""))
            insight_source = row.get("source_event_id", "")
            break
    reference_notes = [
        f"{ref.get('reference_id')}: {ref.get('title') or ref.get('url')} ({ref.get('role')})"
        for ref in reference_rows
        if ref.get("reference_id")
    ]
    competitor_notes = "; ".join(reference_notes[:5]) or open_question("brand/category/competitor notes")
    source_ids = ";".join(sorted({item for item in resolved_source_ids if item}))
    return {
        "business_problem": business_problem,
        "business_source": business_source,
        "client_objective": client_objective,
        "objective_source": objective_source,
        "audience": audience,
        "audience_source": audience_source,
        "barrier": barrier,
        "barrier_source": barrier_source,
        "insight": insight,
        "insight_source": insight_source,
        "product_feature": product_feature,
        "feature_source": feature_source,
        "visual_reference": visual_reference,
        "visual_source": visual_source,
        "competitor_notes": competitor_notes,
        "source_ids": source_ids,
        "requirement_ids": ";".join(row.get("requirement_id", "") for row in requirement_rows if row.get("requirement_id")),
        "reference_ids": ";".join(row.get("reference_id", "") for row in reference_rows if row.get("reference_id")),
    }


def proposal_evidence_ref(source: str) -> str:
    return source or "TBD"


def proposal_display_phrase(label: str, value: str) -> str:
    cleaned = value.strip().strip("。.")
    return cleaned


def build_creative_direction_rows(context: dict[str, object]) -> list[dict[str, str]]:
    """Deprecated compatibility hook; deterministic templates are not creative output."""
    del context
    return []


def legacy_build_creative_direction_rows(context: dict[str, object]) -> list[dict[str, str]]:
    """Legacy test fixture retained for migration comparisons; never called in production."""
    feature = str(context["product_feature"])
    feature_phrase = proposal_display_phrase("product_feature", feature)
    objective_phrase = proposal_display_phrase("client_objective", str(context["client_objective"]))
    audience_phrase = proposal_display_phrase("audience", str(context["audience"]))
    barrier_phrase = proposal_display_phrase("barrier", str(context["barrier"]))
    insight_phrase = proposal_display_phrase("insight", str(context["insight"]))
    visual_phrase = proposal_display_phrase("visual_reference", str(context["visual_reference"]))
    benefit = (
        "把已证实的产品卖点翻译成具体使用价值。"
        if not feature.startswith("TBD")
        else open_question("communication benefit")
    )
    return [
        {
            "direction_id": "DIR-01",
            "name": "场景接续",
            "role": "把产品功能落到真实使用场景",
            "strategy_path": "product_feature_to_behavior_moment",
            "creative_proposition": f"用{visual_phrase}说明{feature_phrase}如何进入{audience_phrase}的真实行动。",
            "core_message": f"{feature_phrase} -> {benefit}",
            "target_feeling": "真实、清爽、可信",
            "product_feature": feature,
            "communication_benefit": benefit,
            "behavior_barrier": barrier_phrase,
            "key_visual_or_action": visual_phrase,
            "title_or_use_case": "待资料确认的使用场景与连续动作",
            "reference_ids": str(context["reference_ids"]),
            "risk": "缺少产品高清图时只能保留 internal_only placeholder。",
            "why_choose": "适合先说明产品如何进入真实行为，不依赖竞品或案例背书。",
            "evidence_refs": ";".join(
                filter(
                    None,
                    [
                        proposal_evidence_ref(str(context["audience_source"])),
                        proposal_evidence_ref(str(context["feature_source"])),
                        proposal_evidence_ref(str(context["visual_source"])),
                    ],
                )
            ),
            "status": "draft",
            "notes": "internal traceable draft",
        },
        {
            "direction_id": "DIR-02",
            "name": "选择理由显性化",
            "role": "把客户目标翻译成可比较的提案路径",
            "strategy_path": "client_objective_to_choice_rationale",
            "creative_proposition": f"围绕客户真实目标：{objective_phrase}，把每个方向的取舍说清楚。",
            "core_message": "客户选择靠证据、动作和风险边界，不靠口号气势。",
            "target_feeling": "清晰、有判断依据、可推进",
            "product_feature": feature,
            "communication_benefit": "让产品利益、执行方式、风险边界能同时被审阅。",
            "behavior_barrier": barrier_phrase,
            "key_visual_or_action": "一页对比矩阵 + 每条方向一张关键动作图或 placeholder slot。",
            "title_or_use_case": "内部评审会方向选择页",
            "reference_ids": str(context["reference_ids"]),
            "risk": "如果客户目标证据不足，本方向必须保持内部待确认。",
            "why_choose": "适合客户还在内部统一意见时使用。",
            "evidence_refs": proposal_evidence_ref(str(context["objective_source"])),
            "status": "draft",
            "notes": "internal traceable draft",
        },
        {
            "direction_id": "DIR-03",
            "name": "阻力转译",
            "role": "把受众阻力转成创意动作",
            "strategy_path": "audience_barrier_to_execution",
            "creative_proposition": f"承认阻力：{barrier_phrase}，用更具体的行动画面降低理解成本。",
            "core_message": f"{insight_phrase} -> 把顾虑变成一个可拍的动作。",
            "target_feeling": "直接、具体、少解释",
            "product_feature": feature,
            "communication_benefit": "让受众先理解为什么需要它，再记住产品。",
            "behavior_barrier": barrier_phrase,
            "key_visual_or_action": (
                "把已确认的行为阻力转成前后可验证动作；具体动作待资料确认。"
                if not barrier_phrase.startswith("TBD")
                else open_question("barrier-to-action execution")
            ),
            "title_or_use_case": "受众痛点页或短片第一幕",
            "reference_ids": str(context["reference_ids"]),
            "risk": "若 insight 未被资料支持，只能作为假设方向，不可客户可见。",
            "why_choose": "适合资料里已有明确担心、禁区或使用障碍时推进。",
            "evidence_refs": ";".join(
                filter(
                    None,
                    [
                        proposal_evidence_ref(str(context["barrier_source"])),
                        proposal_evidence_ref(str(context["insight_source"])),
                    ],
                )
            ),
            "status": "draft",
            "notes": "internal traceable draft",
        },
    ]


def render_creative_directions_content(context: dict[str, object], rows: list[dict[str, str]]) -> str:
    overview = "\n".join(
        "| {direction_id} | {name} | {role} | {strategy_path} | {core_message} | {why_choose} |".format(
            **{key: md_cell(row.get(key, "")) for key in [
                "direction_id",
                "name",
                "role",
                "strategy_path",
                "core_message",
                "why_choose",
            ]}
        )
        for row in rows
    )
    detail_sections = "\n\n".join(
        f"""## {row['direction_id']} {row['name']}

- creative proposition: {row['creative_proposition']}
- core message: {row['core_message']}
- key visual/action: {row['key_visual_or_action']}
- title/use case: {row['title_or_use_case']}
- risk: {row['risk']}
- why choose: {row['why_choose']}
- evidence refs: {row['evidence_refs'] or 'TBD'}
"""
        for row in rows
    )
    return f"""# Creative Directions

status: draft
visibility: internal_only
artifact_role: traceable_internal_creative_proposal_draft

## Evidence Boundaries
- Do not fabricate insight, competitors, audience barriers, or case-study facts.
- Missing facts stay internal as unconfirmed gaps and must not become client-visible claims.
- Video/storyboard execution goes to dircreative; image/KV/backgrounds go to imagegen or Creative Production; fixed templates go to Template Creator.

## Proposal Inputs
- business problem: {context['business_problem']} [source: {proposal_evidence_ref(str(context['business_source']))}]
- client real objective: {context['client_objective']} [source: {proposal_evidence_ref(str(context['objective_source']))}]
- target audience: {context['audience']} [source: {proposal_evidence_ref(str(context['audience_source']))}]
- behavior barrier: {context['barrier']} [source: {proposal_evidence_ref(str(context['barrier_source']))}]
- consumer insight: {context['insight']} [source: {proposal_evidence_ref(str(context['insight_source']))}]
- feature to benefit: {context['product_feature']} -> {rows[0]['communication_benefit']} [source: {proposal_evidence_ref(str(context['feature_source']))}]
- brand/category/competitor notes: {context['competitor_notes']}
- strategy path: product evidence -> audience barrier -> differentiated direction -> client choice rationale

## Direction Overview

| Direction | Name | Role | Strategy Path | Core Message | Why Choose |
|---|---|---|---|---|---|
{overview}

{detail_sections}
## Open Questions
- Confirm unsupported TBD fields before client-facing use.
- Confirm competitor/category evidence before naming any real competitor.
- Confirm which direction should become a PPT/client review lane.
"""


def render_proposal_structure_content(context: dict[str, object], rows: list[dict[str, str]]) -> str:
    direction_pages = "\n".join(
        f"- {row['direction_id']} {row['name']}: proposition / core message / key visual-action / use case / risk / why choose."
        for row in rows
    )
    return f"""# Proposal Structure

status: draft
visibility: internal_only
artifact_role: traceable_internal_proposal_architecture

## Client Review Goal
client real objective: {context['client_objective']}

## Business Problem
{context['business_problem']}

## Audience And Insight
- target audience: {context['audience']}
- behavior barrier: {context['barrier']}
- consumer insight: {context['insight']}

## Feature To Benefit
- product feature: {context['product_feature']}
- communication benefit: {rows[0]['communication_benefit']}

## Brand/Category/Competitor Notes
{context['competitor_notes']}

## Recommended Page Flow
1. 目标和问题界定
2. 受众与行为阻力
3. 产品功能到传播利益
4. 策略路径
5. 2-3 条创意方向对比
6. 推荐方向与选择理由
7. 风险、缺口、待确认问题

## Direction Pages
{direction_pages}

## Reference Pages
- Only cite registered REF rows or explicit TBD search targets.
- Do not use case-study facts until source evidence exists.

## Visual Asset Slots
- DIR-01 key action frame: internal placeholder until visual-quality-gate.
- DIR-02 matrix/choice slide: editable text first.
- DIR-03 barrier/action contrast: internal placeholder until asset slot is bound.

## Proposal Outline
The PPT/proposal outline must preserve problem, objective, audience, insight, feature-to-benefit, direction choices, visual/action execution, risks, and unconfirmed questions.

## Open Questions
- Which TBD fields must be confirmed before client review?
- Which visual route should be delegated to imagegen/Creative Production?
- Which video/storyboard route should be delegated to dircreative?
"""


def render_slide_spec_content(rows: list[dict[str, str]]) -> str:
    slide_rows = [
        ("1", "Problem", "Business problem + client real objective", "none", "internal_only"),
        ("2", "Audience Insight", "Target audience + behavior barrier + consumer insight", "none", "internal_only"),
        ("3", "Feature To Benefit", "Product feature translated to communication benefit", "none", "internal_only"),
        ("4", "Direction Matrix", "2-3 differentiated directions with why choose", "none", "internal_only"),
    ]
    for index, row in enumerate(rows, start=5):
        slide_rows.append(
            (
                str(index),
                row["name"],
                f"{row['creative_proposition']} / {row['core_message']} / risk: {row['risk']}",
                f"{row['direction_id']}-KEY-ACTION",
                "internal_only",
            )
        )
    table = "\n".join(
        f"| {num} | {md_cell(purpose)} | {md_cell(content)} | {md_cell(slot)} | {visibility} |"
        for num, purpose, content, slot, visibility in slide_rows
    )
    return f"""# Slide Spec

status: draft
visibility: internal_only
artifact_role: internal_editable_proposal_outline

## Rules

```text
Text must remain editable.
Images must have asset IDs or placeholder IDs.
No internal notes in client-visible slides.
No fake logo or fake case evidence.
Do not treat VALIDATION=PASS as creative quality approval.
```

## Slides

| Slide | Purpose | Content | Asset Slot | Visibility |
|---|---|---|---|---|
{table}
"""


def render_client_outline_rows(context: dict[str, object], rows: list[dict[str, str]]) -> list[dict[str, str]]:
    def readable(value: object, fallback: str) -> str:
        text = re.sub(r"\s+", " ", str(value or "")).strip()
        if (
            not non_placeholder(text, min_chars=4)
            or find_client_language_hits(text)
            or re.search(r"internal[_ -]?only|placeholder|source[_ -]?brief", text, flags=re.IGNORECASE)
        ):
            return fallback
        return text

    objective = readable(
        context.get("client_objective"),
        "本轮先对齐方案要解决的业务问题、客户需要做出的决定，以及下一阶段深化范围。",
    )
    problem = readable(
        context.get("business_problem"),
        "现有资料需要被整理成更清楚的传播问题和可比较的创意选择。",
    )
    audience = readable(
        context.get("audience"),
        "目标受众的具体生活场景和使用动机，将在首次方案讨论中共同锁定。",
    )
    barrier = readable(
        context.get("barrier"),
        "受众为何犹豫、忽略或不相信现有表达，是本轮需要共同判断的关键。",
    )
    insight = readable(
        context.get("insight"),
        "洞察页先呈现可讨论的行为假设，并明确哪些判断仍需要客户事实支持。",
    )
    product_feature = readable(
        context.get("product_feature"),
        "产品事实、卖点优先级和不可夸张边界，将由客户资料与首次讨论共同锁定。",
    )
    outline_rows: list[dict[str, str]] = [
        {
            "slide_id": "1",
            "page_title": "客户目标与本轮问题",
            "body_copy": f"本轮先把客户真实目标和需要解决的传播问题讲清楚：{objective} 同时聚焦这一核心问题：{problem}",
            "client_confirmation_point": "确认本轮客户审阅是先定方向，还是需要接近可发送版本。",
            "material_role": "source_brief_to_decision_context",
            "visual_slot": "文字页，保留一处客户目标或 brief 摘要占位。",
            "visual_asset_status": "text_only",
            "asset_ids": "",
            "visibility": "client_visible_pending",
            "status": "draft",
            "notes": "must be reviewed before PPT export",
        },
        {
            "slide_id": "2",
            "page_title": "受众洞察与行为阻力",
            "body_copy": f"把受众、阻力和洞察放在同一页，避免方案只剩口号。受众：{audience} 行为阻力：{barrier} 洞察假设：{insight}",
            "client_confirmation_point": "确认目标受众和行为阻力是否符合客户团队判断。",
            "material_role": "audience_insight",
            "visual_slot": "受众场景图或情绪参考图；素材锁定前使用低保真示意画面。",
            "visual_asset_status": "placeholder",
            "asset_ids": "",
            "visibility": "client_visible_pending",
            "status": "draft",
            "notes": "customer-readable story logic",
        },
        {
            "slide_id": "3",
            "page_title": "产品事实到传播利益",
            "body_copy": f"先讲清楚产品事实怎样转成客户能判断的传播利益：{product_feature}。这一页只保留有资料依据的事实和可讨论的利益表达。",
            "client_confirmation_point": "确认产品事实、卖点优先级和不可夸张边界。",
            "material_role": "product_fact_to_benefit",
            "visual_slot": "产品事实图、功能示意或已有产品图。",
            "visual_asset_status": "to_generate",
            "asset_ids": "",
            "visibility": "client_visible_pending",
            "status": "draft",
            "notes": "must remain editable",
        },
        {
            "slide_id": "4",
            "page_title": "方向比较与选择理由",
            "body_copy": "用同一套标准比较方向：策略路径、核心信息、关键画面、风险和为什么值得选，帮助客户做决策。",
            "client_confirmation_point": "确认客户希望保留几个方向进入下一轮深化。",
            "material_role": "decision_matrix",
            "visual_slot": "方向矩阵，不放大图，保留小比例视觉槽位。",
            "visual_asset_status": "text_only",
            "asset_ids": "",
            "visibility": "client_visible_pending",
            "status": "draft",
            "notes": "no production-side language",
        },
    ]
    for index, row in enumerate(rows, start=5):
        proposition = readable(
            row.get("creative_proposition"),
            "本方向从真实使用场景出发，把产品利益转成受众能理解和记住的行动画面。",
        )
        core_message = readable(
            row.get("core_message"),
            "核心信息将围绕产品事实、受众动机和可感知利益展开。",
        )
        key_visual = readable(
            row.get("key_visual_or_action"),
            "以一张关键动作画面说明人物、产品与使用场景之间的关系。",
        )
        why_choose = readable(
            row.get("why_choose"),
            "这个方向便于客户判断信息重点、画面潜力和下一阶段深化成本。",
        )
        outline_rows.append(
            {
                "slide_id": str(index),
                "page_title": row["name"],
                "body_copy": f"{proposition} 核心信息是：{core_message} 关键画面或行动：{key_visual} 选择理由：{why_choose}",
                "client_confirmation_point": f"确认是否选择 {row['name']} 继续深化。",
                "material_role": "creative_direction_story_page",
                "visual_slot": f"{key_visual}，素材锁定前先使用可替换的示意画面。",
                "visual_asset_status": "to_generate",
                "asset_ids": "",
                "visibility": "client_visible_pending",
                "status": "draft",
                "notes": f"direction_id={row['direction_id']}",
            }
        )
    return outline_rows


def render_client_review_outline_content(context: dict[str, object], rows: list[dict[str, str]]) -> str:
    outline_rows = render_client_outline_rows(context, rows)
    page_lines = "\n".join(
        f"{row['slide_id']}. {row['page_title']} - {row['body_copy']}"
        for row in outline_rows
    )
    confirmation_lines = "\n".join(
        f"- P{row['slide_id']}: {row['client_confirmation_point']}"
        for row in outline_rows
    )
    return f"""# Client Review Outline

status: draft
visibility: internal_only
artifact_role: client_readable_text_framework_before_ppt

## Review Goal

{context['client_objective']}

## Content Flow

{page_lines}

## Client Confirmation Points

{confirmation_lines}

## Material Logic

- 每页必须有客户可读标题、正文、客户确认点和素材角色。
- 图片或视频素材只作为支撑客户判断的证据，不替代故事逻辑。
- 进入 PPT builder 前先跑 `adco client-outline-gate <project_dir>`。
"""


def render_creative_brief(project: Path, *, work_id: str = "") -> dict[str, object]:
    ensure_project(project)
    delivery_surface = project_surface(project) == DELIVERY_SURFACE
    source_ids = ""
    requirement_ids = ""
    if delivery_surface:
        _, requirements = read_csv_rows(
            project / "AD-creative/orchestrator/requirements.csv"
        )
        _, sources = read_csv_rows(
            project / "AD-creative/orchestrator/source_events.csv"
        )
        source_ids = ";".join(
            row.get("source_event_id", "")
            for row in sources
            if row.get("source_event_id")
        )
        requirement_ids = ";".join(
            row.get("requirement_id", "")
            for row in requirements
            if row.get("requirement_id")
        )
        objective = next(
            (row.get("statement", "") for row in requirements if row.get("statement")),
            "Evidence-bound creative brief",
        )
        work_id = ensure_creative_proposal_work(
            project,
            work_id,
            source_ids,
            objective,
        )
    result = create_creative_brief(project)
    brief_artifacts = [
        ("ART-AUTO-CREATIVE-BRIEF-SNAPSHOT", "creative_brief_snapshot", BRIEF_SNAPSHOT_REL),
        ("ART-AUTO-CREATIVE-BRIEF-CONTRACT", "creative_brief_contract", BRIEF_CONTRACT_REL),
        ("ART-AUTO-CREATIVE-CANDIDATE-SCHEMA", "creative_candidate_schema", CANDIDATE_SCHEMA_REL),
        ("ART-AUTO-CREATIVE-GENERATION-REQUEST", "creative_generation_request", GENERATION_REQUEST_REL),
        ("ART-AUTO-CREATIVE-OPEN-GAPS", "creative_open_evidence_gaps", OPEN_GAPS_REL),
    ]
    artifact_ids: list[str] = []
    for artifact_id, artifact_type, rel_path in brief_artifacts:
        artifact_ids.append(artifact_id)
        if delivery_surface:
            update_artifact(
                project,
                artifact_id,
                artifact_type,
                str(rel_path),
                "creative",
                visibility="internal_only",
                source_event_ids=source_ids,
                linked_requirements=requirement_ids,
                linked_work_items=work_id,
                linked_references="",
                gate_status="NOT_RUN",
            )
    if delivery_surface:
        append_event(
            project,
            {
                "event_id": f"EVT-CREATIVE-BRIEF-{datetime.now().strftime('%Y%m%d%H%M%S')}",
                "event_type": "creative_brief_created",
                "created_at": now_iso(),
                "work_id": work_id,
                "artifacts": artifact_ids,
                "brief_snapshot_sha256": result.snapshot_sha256,
            },
        )
    return {
        "project": str(project),
        "work_id": work_id,
        "artifact_ids": artifact_ids,
        "paths": [str(path) for path in result.paths],
        "brief_snapshot_sha256": result.snapshot_sha256,
        "evidence_refs": result.evidence_refs,
        "open_evidence_gaps": result.open_gaps,
        "directions_generated": 0,
    }


def render_creative_proposal(project: Path, *, work_id: str = "") -> dict[str, object]:
    """Deprecated alias for render_creative_brief; never synthesizes directions."""
    payload = render_creative_brief(project, work_id=work_id)
    payload["deprecated_alias"] = "creative-brief"
    return payload


def creative_proposal_scan_files(project: Path) -> list[Path]:
    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    candidate_mode = (project / CURRENT_CANDIDATE_REL).is_file()
    paths = (
        [
            project / CURRENT_CANDIDATE_REL,
            project / CREATIVE_DIRECTIONS_REL,
            project / OPTION_MATRIX_REL,
            project / CRITIC_RECEIPT_REL,
        ]
        if candidate_mode
        else [project / rel_path for _, _, rel_path in CREATIVE_PROPOSAL_ARTIFACTS]
    )
    creative_types = {
        "creative_directions",
        "creative_option_matrix",
        "proposal_structure",
        "slide_spec",
        "creative_proposal",
        "proposal_outline",
        "creative_candidate",
        "creative_critic_receipt",
    }
    for artifact in artifacts:
        if artifact.get("artifact_type", "").strip() in creative_types:
            rel_path = artifact.get("path", "").strip()
            if rel_path:
                paths.append(project / rel_path)
    deduped: list[Path] = []
    seen: set[Path] = set()
    for path in paths:
        resolved = path.resolve()
        if resolved not in seen:
            seen.add(resolved)
            deduped.append(path)
    return deduped


def read_proposal_texts(paths: list[Path]) -> tuple[str, list[str]]:
    chunks: list[str] = []
    missing: list[str] = []
    for path in paths:
        if not path.exists():
            missing.append(str(path))
            continue
        if path.suffix.lower() not in TEXT_SUFFIXES:
            continue
        chunks.append(path.read_text(encoding="utf-8", errors="ignore"))
    return "\n\n".join(chunks), missing


def load_direction_rows(project: Path) -> list[dict[str, str]]:
    _, rows = read_csv_rows(project / "AD-creative/creative/option_matrix.csv")
    return [row for row in rows if row.get("direction_id", "").strip()]


def field_is_tbd(value: str) -> bool:
    cleaned = value.strip().lower()
    return cleaned in {"", "tbd", "n/a", "-"} or "tbd" in cleaned or "open question" in cleaned or "待补充" in value or "暂无" in value


def has_supported_case_claim(line: str) -> bool:
    if not CASE_CLAIM_PATTERN.search(line):
        return True
    if re.search(r"不依赖|不是|非|不能|不要|without|not a|not rely|excluded", line, re.IGNORECASE):
        return True
    return bool(re.search(r"\b(REF|SRC|ART)-[A-Z0-9-]+|https://|TBD|open question|待补充", line))


def has_trace_marker(line: str) -> bool:
    return bool(re.search(r"\b(REF|SRC|REQ|ART)-[A-Z0-9-]+|https://|\[source:", line, re.IGNORECASE))


def collect_humanizer_writing_risks(combined_text: str, client_texts: list[tuple[Path, str]], project: Path) -> list[tuple[str, str]]:
    risks: list[tuple[str, str]] = []
    for code, pattern in CLIENT_WRITING_RISK_CODES.items():
        hits: list[str] = []
        for line in combined_text.splitlines():
            if not pattern.search(line):
                continue
            if code == "VAGUE_AUTHORITY_CLAIM" and has_trace_marker(line):
                continue
            hits.append(compact_evidence(line, 120))
        if hits:
            risks.append((code, "; ".join(hits[:3])))

    vocab_hits = [match.group(0) for match in GENERIC_AI_VOCABULARY_PATTERN.finditer(combined_text)]
    unique_vocab_hits = sorted(set(vocab_hits), key=str.lower)
    if len(vocab_hits) >= 4 or len(unique_vocab_hits) >= 3:
        risks.append(("GENERIC_AI_VOCABULARY", ", ".join(unique_vocab_hits[:8])))

    dash_hits = re.findall(r"—|–| -- ", combined_text)
    client_dash_locations = [
        safe_rel(project, path)
        for path, text in client_texts
        if re.search(r"—|–| -- ", text)
    ]
    if len(dash_hits) >= 2:
        risks.append(("DASH_OVERUSE", f"{len(dash_hits)} dash-like separators in proposal text"))
    elif client_dash_locations:
        risks.append(("DASH_OVERUSE", "client-facing dash-like separators: " + ", ".join(client_dash_locations[:4])))
    return risks


def client_facing_scan_paths(project: Path, files: list[Path]) -> list[Path]:
    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    client_paths: set[Path] = set()
    for artifact in artifacts:
        if artifact.get("visibility", "").lower() in CLIENT_VISIBLE_VALUES:
            rel_path = artifact.get("path", "").strip()
            if rel_path:
                client_paths.add((project / rel_path).resolve())
    for path in files:
        if not path.exists() or path.suffix.lower() not in TEXT_SUFFIXES:
            continue
        head = "\n".join(path.read_text(encoding="utf-8", errors="ignore").splitlines()[:12]).lower()
        if "visibility: client_visible" in head or "visibility: client_visible_ready" in head:
            client_paths.add(path.resolve())
    return sorted(client_paths)


def review_creative_quality(project: Path) -> tuple[str, list[str], Path]:
    files = creative_proposal_scan_files(project)
    candidate_mode = (project / CURRENT_CANDIDATE_REL).is_file()
    active_artifacts = (
        [
            ("ART-AUTO-CREATIVE-CANDIDATE", "creative_candidate", CURRENT_CANDIDATE_REL),
            ("ART-AUTO-CREATIVE-DIRECTIONS", "creative_directions", CREATIVE_DIRECTIONS_REL),
            ("ART-AUTO-CREATIVE-OPTION-MATRIX", "creative_option_matrix", OPTION_MATRIX_REL),
        ]
        if candidate_mode
        else CREATIVE_PROPOSAL_ARTIFACTS
    )
    for artifact_id, artifact_type, rel_path in active_artifacts:
        if (project / rel_path).exists():
            update_artifact(
                project,
                artifact_id,
                artifact_type,
                str(rel_path),
                "creative",
                visibility="internal_only",
                gate_status="NOT_RUN",
            )
    combined_text, missing_files = read_proposal_texts(files)
    lower_text = combined_text.lower()
    direction_rows = load_direction_rows(project)
    issues: list[str] = []
    warnings: list[str] = []
    reason_codes: list[str] = []
    evidence: list[str] = [
        f"scanned_files={len(files)}",
        f"direction_rows={len(direction_rows)}",
    ]

    if missing_files:
        reason_codes.append("MISSING_PROPOSAL_ARTIFACT")
        issues.extend(f"creative proposal artifact missing: {safe_rel(project, Path(path))}" for path in missing_files[:6])
    if len(combined_text.strip()) < 900 or not direction_rows:
        reason_codes.append("EMPTY_SKELETON")
        issues.append("创意/提案文件仍是空骨架或缺少 option_matrix 方向行。")
    required_labels = (
        [
            "human tension",
            "brand truth",
            "audience truth",
            "single-minded proposition",
            "creative mechanism",
            "key visual",
            "story or behavior",
            "product role",
            "why brand can own it",
            "production risk",
            "evidence refs",
        ]
        if candidate_mode
        else CREATIVE_PROPOSAL_REQUIRED_LABELS
    )
    for label in required_labels:
        if label not in lower_text:
            reason_codes.append("MISSING_REQUIRED_FIELD")
            issues.append(f"缺少必要提案字段: {label}")
            break

    insight_label = "human tension" if candidate_mode else "consumer insight"
    insight_lines = [line for line in combined_text.splitlines() if insight_label in line.lower()]
    if not insight_lines or all(field_is_tbd(line) for line in insight_lines):
        reason_codes.append("WEAK_OR_MISSING_INSIGHT")
        issues.append("consumer insight 缺失、过薄或仍是 TBD。")
    feature_labels = ("product role", "single-minded proposition") if candidate_mode else ("feature to benefit",)
    feature_lines = [
        line
        for line in combined_text.splitlines()
        if any(label in line.lower() for label in feature_labels)
    ]
    if not feature_lines or all(field_is_tbd(line) for line in feature_lines):
        reason_codes.append("NO_PRODUCT_TO_BENEFIT")
        issues.append("缺少产品功能到传播利益的明确翻译。")

    if len(direction_rows) < 2:
        reason_codes.append("TOO_FEW_DIRECTIONS")
        issues.append("创意方向少于 2 条。")
    else:
        signatures = {
            re.sub(
                r"\s+",
                " ",
                " ".join(
                    row.get(field, "").strip().lower()
                    for field in ["strategy_path", "creative_proposition", "core_message", "key_visual_or_action"]
                ),
            )
            for row in direction_rows
        }
        if len(signatures) < len(direction_rows):
            reason_codes.append("UNDIFFERENTIATED_DIRECTIONS")
            issues.append("创意方向之间不可区分，strategy/proposition/message/action 过度重复。")
    if direction_rows and all(field_is_tbd(row.get("key_visual_or_action", "")) for row in direction_rows):
        reason_codes.append("NO_KEY_VISUAL_OR_ACTION")
        issues.append("所有方向都缺少 key visual/actionable execution。")
    if direction_rows and all(field_is_tbd(row.get("why_choose", "")) for row in direction_rows):
        reason_codes.append("NO_CLIENT_CHOICE_RATIONALE")
        issues.append("所有方向都缺少客户选择理由。")
    if direction_rows and any(field_is_tbd(row.get("creative_proposition", "")) for row in direction_rows):
        reason_codes.append("THIN_CREATIVE_PROPOSITION")
        issues.append("至少一条方向的 creative proposition 仍是 TBD。")

    generic_hits = sorted({match.group(0) for match in GENERIC_CREATIVE_PATTERN.finditer(combined_text)})
    if generic_hits:
        reason_codes.append("GENERIC_AI_CLICHE")
        issues.append("提案含泛化口号/AI cliche: " + ", ".join(generic_hits[:8]))

    unsupported_claims = [
        compact_evidence(line, 120)
        for line in combined_text.splitlines()
        if CASE_CLAIM_PATTERN.search(line) and not has_supported_case_claim(line)
    ]
    if unsupported_claims:
        reason_codes.append("UNSUPPORTED_REFERENCE_OR_CASE_CLAIM")
        issues.extend(f"未追溯案例/参考/数据声明: {line}" for line in unsupported_claims[:6])

    client_paths = client_facing_scan_paths(project, files)
    client_texts: list[tuple[Path, str]] = []
    for path in client_paths:
        text = path.read_text(encoding="utf-8", errors="ignore")
        client_texts.append((path, text))
        risky = sorted({match.group(0) for match in INTERNAL_LANGUAGE_PATTERN.finditer(text)})
        if risky:
            reason_codes.append("INTERNAL_LANGUAGE_LEAK")
            issues.append(f"client-facing material contains internal language: {safe_rel(project, path)}: {', '.join(risky[:6])}")

    for code, evidence_item in collect_humanizer_writing_risks(combined_text, client_texts, project):
        reason_codes.append(code)
        issues.append(f"humanizer writing risk ({code}): {evidence_item}")

    tbd_count = len(re.findall(r"\bTBD\b|open question|待补充|暂无", combined_text, flags=re.IGNORECASE))
    if tbd_count:
        reason_codes.append("OPEN_EVIDENCE_GAPS")
        warnings.append(f"仍有 {tbd_count} 个 TBD/open question，只能作为内部草案或 PARTIAL。")

    if candidate_mode:
        try:
            critic_result = review_creative_candidate(project)
        except ValueError as exc:
            reason_codes.append("CRITIC_RECEIPT_INVALID")
            issues.append(f"Critic Receipt could not be produced: {exc}")
        else:
            assert critic_result.receipt_path is not None
            evidence.append(f"critic_receipt={safe_rel(project, critic_result.receipt_path)}")
            if critic_result.blocking_issues:
                reason_codes.append("CRITIC_STRUCTURE_BLOCKED")
                issues.extend(critic_result.blocking_issues)
            warnings.extend(critic_result.warnings)
        reason_codes.append("INDEPENDENT_CREATIVE_CRITIC_REQUIRED")
        warnings.append(
            "确定性结构/语言 Lint 不能批准创意质量；仍需绑定 exact candidate 的独立创意 Critic。"
        )
    status = "PASS" if not issues and not warnings else "PARTIAL_PASS" if not issues else "BLOCKED"
    if status == "PASS":
        status = "PARTIAL_PASS"
        reason_codes.append("INDEPENDENT_CREATIVE_CRITIC_REQUIRED")
        warnings.append(
            "确定性结构/语言 Lint 已通过，但独立创意 Critic 尚未提供完整创意质量批准。"
        )
    status = enforce_adversarial_gate_policy(project, "creative", status, warnings, evidence)
    if status == "PARTIAL_PASS" and "ADVERSARIAL_COUNCIL_MISSING" not in reason_codes and any("反驳性议会" in item for item in warnings):
        reason_codes.append("ADVERSARIAL_COUNCIL_MISSING")

    report_path = project / "AD-creative/gates/GATE-AUTO-CREATIVE-QUALITY-001_report.md"
    issue_text = "\n".join(f"- {issue}" for issue in issues) or "- 无"
    warning_text = "\n".join(f"- {warning}" for warning in warnings) or "- 无"
    reason_text = "\n".join(f"- {code}" for code in sorted(set(reason_codes))) or "- NONE"
    evidence_text = "\n".join(f"- {item}" for item in evidence)
    write_text(
        report_path,
        f"""# Creative Quality Gate

status: {status}
visibility: internal_only
checked_at: {now_iso()}

## Reason Codes

{reason_text}

## Evidence

{evidence_text}

## Blocking Issues

{issue_text}

## Warnings

{warning_text}

## Rules

- Gate checks proposal traceability and deterministic structure/language lint, not subjective creative quality.
- A structural PASS is capped at PARTIAL_PASS until an independent creative Critic is bound to the exact candidate.
- PASS/PARTIAL_PASS/BLOCKED are reason-code based; score alone is never approval.
- Missing facts stay internal as unconfirmed gaps and prevent client-ready claims.
- Blocks empty skeletons, generic slogans, weak insight, undifferentiated directions, missing feature-to-benefit, missing key visual/action, missing why-choose, unsupported case/reference claims, internal language leaks, and humanizer writing risks.
- `VALIDATION=PASS` is structural only and never replaces this creative-quality-gate.
""",
    )
    update_artifact(
        project,
        "ART-AUTO-CREATIVE-QUALITY-GATE",
        "creative_quality_gate_report",
        safe_rel(project, report_path),
        "creative",
        status="done" if status != "BLOCKED" else "blocked",
        visibility="internal_only",
        gate_status=status,
    )
    checked_artifacts = ";".join([artifact_id for artifact_id, _, _ in active_artifacts] + ["ART-AUTO-CREATIVE-QUALITY-GATE"])
    append_gate(
        project,
        "GATE-AUTO-CREATIVE-QUALITY-001",
        "creative",
        status,
        "90" if status == "PASS" else "65" if status == "PARTIAL_PASS" else "30",
        checked_artifacts,
        ";".join(sorted(set(reason_codes))) if issues else "",
        ";".join(warnings[:8]) or ("修正创意提案后重跑 creative-quality-gate。" if issues else ""),
        "",
        "ready_for_internal_ppt" if status == "PASS" else "resolve_creative_gaps" if status == "PARTIAL_PASS" else "revise_creative_proposal",
        "ad_creative_operator",
    )
    append_event(
        project,
        {
            "event_id": "EVT-AUTO-CREATIVE-QUALITY-GATE",
            "event_type": "creative_quality_gate_run",
            "created_at": now_iso(),
            "status": status,
            "reason_codes": sorted(set(reason_codes)),
            "issues": issues[:12],
            "warnings": warnings[:12],
        },
    )
    return status, issues + warnings, report_path


def review_reference_pack(project: Path, *, live_check: bool = False) -> tuple[str, list[str], Path]:
    _, refs = read_csv_rows(project / "AD-creative/references/reference_cards.csv")
    issues: list[str] = []
    warnings: list[str] = []
    evidence: list[str] = [f"references={len(refs)}", f"live_check={live_check}"]

    if not refs:
        issues.append("参考包为空。")

    for ref in refs:
        ref_id = ref.get("reference_id", "")
        url = ref.get("url", "").strip()
        title = ref.get("title", "").strip()
        source_owner = ref.get("source_owner", "").strip().lower()
        role = ref.get("role", "").strip()
        why_relevant = ref.get("why_relevant", "").strip()
        borrow = ref.get("borrow", "").strip()
        do_not_copy = ref.get("do_not_copy", "").strip()
        client_visible = ref.get("client_visible", "").strip().lower() == "true"

        if not title:
            warnings.append(f"{ref_id} 缺少标题。")
        if not role:
            warnings.append(f"{ref_id} 缺少 role。")
        if not why_relevant:
            warnings.append(f"{ref_id} 缺少 why_relevant。")
        if not borrow:
            warnings.append(f"{ref_id} 缺少 borrow。")
        if not do_not_copy:
            message = f"{ref_id} 缺少 do_not_copy。"
            if client_visible:
                issues.append(message)
            else:
                warnings.append(message)
        if not source_owner or source_owner == "unknown":
            message = f"{ref_id} 来源归属未明确。"
            if client_visible:
                issues.append(message)
            else:
                warnings.append(message)

        if url == "TBD" or not url:
            message = f"{ref_id} 尚未登记真实 https URL。"
            if client_visible:
                issues.append(message)
            else:
                warnings.append(message)
            continue

        if not url.startswith("https://"):
            message = f"{ref_id} 不是 https URL。"
            if client_visible:
                issues.append(message)
            else:
                warnings.append(message)
            continue

        if live_check:
            ok, note = check_url(url)
            evidence.append(f"{ref_id}_live={note}")
            if not ok:
                message = f"{ref_id} live check failed: {note}"
                if client_visible:
                    issues.append(message)
                else:
                    warnings.append(message)

        if client_visible and source_owner in {"research", "ugc", "unknown"}:
            issues.append(f"{ref_id} 客户可见但 source_owner={source_owner}。")

    status = "PASS" if not issues and not warnings else "PARTIAL_PASS" if not issues else "BLOCKED"
    status = enforce_adversarial_gate_policy(
        project, "reference_research", status, warnings, evidence
    )
    report_path = project / "AD-creative/gates/GATE-AUTO-REFERENCE-PACK-001_report.md"
    issue_text = "\n".join(f"- {issue}" for issue in issues) or "- 无"
    warning_text = "\n".join(f"- {warning}" for warning in warnings) or "- 无"
    evidence_text = "\n".join(f"- {item}" for item in evidence)
    write_text(
        report_path,
        f"""# Reference Pack Gate

status: {status}
visibility: internal_only
checked_at: {now_iso()}

## Evidence

{evidence_text}

## Blocking Issues

{issue_text}

## Warnings

{warning_text}

## Rules

- 客户可见参考必须是真实 https URL。
- 客户可见参考必须说明 source_owner、role、why_relevant、borrow、do_not_copy。
- `TBD` 只能作为内部搜索计划，不能进入客户稿引用。
- UGC / research / unknown 不能冒充官方证据。
""",
    )
    update_artifact(
        project,
        "ART-AUTO-REFERENCE-PACK-GATE",
        "reference_pack_gate_report",
        safe_rel(project, report_path),
        "reference_research",
        status="done" if status != "BLOCKED" else "blocked",
        visibility="internal_only",
        linked_references=";".join(ref.get("reference_id", "") for ref in refs if ref.get("reference_id")),
        gate_status=status,
    )
    append_gate(
        project,
        "GATE-AUTO-REFERENCE-PACK-001",
        "reference_research",
        status,
        "90" if status == "PASS" else "65" if status == "PARTIAL_PASS" else "35",
        "ART-AUTO-REFERENCE-PACK-GATE",
        ";".join(issues[:8]),
        ";".join(warnings[:8]) or "补充真实公开参考后重跑 reference-pack-gate。",
        "",
        "ready_for_internal_creative" if status != "BLOCKED" else "fix_reference_pack",
        "ad_creative_operator",
    )
    append_event(
        project,
        {
            "event_id": "EVT-AUTO-REFERENCE-PACK-GATE",
            "event_type": "reference_pack_gate_run",
            "created_at": now_iso(),
            "status": status,
            "issues": issues[:12],
            "warnings": warnings[:12],
        },
    )
    return status, issues + warnings, report_path


def search_plan_files(project: Path, artifacts: list[dict[str, str]]) -> list[Path]:
    paths: list[Path] = []
    for artifact in artifacts:
        artifact_type = artifact.get("artifact_type", "").lower()
        rel_path = artifact.get("path", "").strip()
        if artifact_type == "search_plan" and rel_path:
            paths.append(project / rel_path)
    ref_dir = project / "AD-creative/references"
    if ref_dir.exists():
        paths.extend(sorted(ref_dir.glob("*search_plan*.md")))
    deduped: list[Path] = []
    seen: set[Path] = set()
    for path in paths:
        resolved = path.resolve()
        if resolved not in seen:
            seen.add(resolved)
            deduped.append(path)
    return deduped


def review_search_quality(project: Path) -> tuple[str, list[str], Path]:
    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    _, refs = read_csv_rows(project / "AD-creative/references/reference_cards.csv")
    plans = search_plan_files(project, artifacts)
    issues: list[str] = []
    warnings: list[str] = []
    evidence: list[str] = [f"search_plans={len(plans)}", f"references={len(refs)}"]

    if not plans:
        warnings.append("未找到 search_plan / official_search_plan。")

    for plan in plans:
        rel_plan = safe_rel(project, plan)
        if not plan.exists():
            issues.append(f"搜索计划文件不存在: {rel_plan}")
            continue
        text = plan.read_text(encoding="utf-8", errors="ignore")
        lowered = text.lower()
        evidence.append(f"plan={rel_plan}")
        if "why search" not in lowered and "为什么" not in text and "why" not in lowered:
            warnings.append(f"{rel_plan} 缺少搜索原因。")
        if (
            "search scope" not in lowered
            and "suggested platforms" not in lowered
            and "搜索范围" not in text
            and "平台" not in text
        ):
            warnings.append(f"{rel_plan} 缺少搜索范围/平台。")
        if "expected output" not in lowered and "output" not in lowered and "产出" not in text:
            warnings.append(f"{rel_plan} 缺少预期产物。")
        if "user decision needed" in lowered or "needs_user_input" in lowered or "批准" in text:
            warnings.append(f"{rel_plan} 仍需用户批准或缩小范围。")
        if "do_not_copy" not in lowered and "do not copy" not in lowered and "禁止复制" not in text:
            warnings.append(f"{rel_plan} 未明确 do_not_copy 输出。")

    search_targets = [ref for ref in refs if ref.get("role", "").strip() == "search_target"]
    if search_targets:
        evidence.append(f"search_targets={len(search_targets)}")
    for ref in search_targets:
        ref_id = ref.get("reference_id", "")
        url = ref.get("url", "").strip()
        if url == "TBD" or not url:
            warnings.append(f"{ref_id} 仍是搜索目标，尚无真实 URL。")
        if not ref.get("why_relevant", "").strip():
            warnings.append(f"{ref_id} 缺少 why_relevant。")
        if not ref.get("borrow", "").strip():
            warnings.append(f"{ref_id} 缺少 borrow。")
        if not ref.get("do_not_copy", "").strip():
            warnings.append(f"{ref_id} 缺少 do_not_copy。")
        if ref.get("client_visible", "").strip().lower() == "true":
            issues.append(f"{ref_id} 搜索目标不能标记客户可见。")

    client_visible_refs = [
        ref for ref in refs if ref.get("client_visible", "").strip().lower() == "true"
    ]
    for ref in client_visible_refs:
        url = ref.get("url", "").strip()
        ref_id = ref.get("reference_id", "")
        if not url.startswith("https://"):
            issues.append(f"{ref_id} 客户可见参考不是真实 https URL。")
        if ref.get("source_owner", "").strip().lower() in {"research", "unknown", "ugc", ""}:
            issues.append(f"{ref_id} 客户可见参考来源归属不可信。")

    status = "PASS" if not issues and not warnings else "PARTIAL_PASS" if not issues else "BLOCKED"
    status = enforce_adversarial_gate_policy(
        project, "reference_research", status, warnings, evidence
    )
    report_path = project / "AD-creative/gates/GATE-AUTO-SEARCH-QUALITY-001_report.md"
    issue_text = "\n".join(f"- {issue}" for issue in issues) or "- 无"
    warning_text = "\n".join(f"- {warning}" for warning in warnings) or "- 无"
    evidence_text = "\n".join(f"- {item}" for item in evidence)
    write_text(
        report_path,
        f"""# Search Quality Gate

status: {status}
visibility: internal_only
checked_at: {now_iso()}

## Evidence

{evidence_text}

## Blocking Issues

{issue_text}

## Warnings

{warning_text}

## Rules

- 搜索计划必须说明 why / scope / platform / expected output / do_not_copy。
- `NEEDS_USER_INPUT` 或 `TBD` 只能内部推进，不能当成客户可见参考。
- role=search_target 的参考不能标记客户可见。
- 客户可见参考必须是真实 https URL 且来源归属可信。
""",
    )
    update_artifact(
        project,
        "ART-AUTO-SEARCH-QUALITY-GATE",
        "search_quality_gate_report",
        safe_rel(project, report_path),
        "reference_research",
        status="done" if status != "BLOCKED" else "blocked",
        visibility="internal_only",
        linked_references=";".join(ref.get("reference_id", "") for ref in refs if ref.get("reference_id")),
        gate_status=status,
    )
    append_gate(
        project,
        "GATE-AUTO-SEARCH-QUALITY-001",
        "reference_research",
        status,
        "90" if status == "PASS" else "65" if status == "PARTIAL_PASS" else "35",
        "ART-AUTO-SEARCH-QUALITY-GATE",
        ";".join(issues[:8]),
        ";".join(warnings[:8]) or "补齐搜索计划和真实参考后重跑 search-quality-gate。",
        "",
        "ready_for_reference_pack" if status != "BLOCKED" else "fix_search_quality",
        "ad_creative_operator",
    )
    append_event(
        project,
        {
            "event_id": "EVT-AUTO-SEARCH-QUALITY-GATE",
            "event_type": "search_quality_gate_run",
            "created_at": now_iso(),
            "status": status,
            "issues": issues[:12],
            "warnings": warnings[:12],
        },
    )
    return status, issues + warnings, report_path


def add_visual_asset(
    project: Path,
    file_path: Path,
    slot_id: str,
    requirement_id: str,
    reference_id: str,
    asset_type: str,
    visibility: str,
    qa_status: str,
    risk_level: str,
    prompt_or_edit_ref: str,
    notes: str,
    selected: bool,
) -> tuple[str, Path]:
    if not file_path.exists():
        raise FileNotFoundError(f"asset not found: {file_path}")
    target_dir = project / ("AD-creative/visual_assets/selected" if selected else "AD-creative/visual_assets/raw")
    target_dir.mkdir(parents=True, exist_ok=True)
    fields, assets = read_csv_rows(project / "AD-creative/visual_assets/asset_manifest.csv")
    if not fields:
        raise FileNotFoundError("asset_manifest.csv header not found")
    asset_id = next_id(assets, "asset_id", "IMG")
    suffix = file_path.suffix.lower() or ".bin"
    target = target_dir / f"{asset_id}{suffix}"
    shutil.copy2(file_path, target)
    image_note = ""
    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        try:
            from PIL import Image

            with Image.open(target) as image:
                image_note = f"{image.width}x{image.height}"
        except Exception as exc:  # noqa: BLE001 - non-fatal QA evidence
            image_note = f"image_probe_failed={exc}"
    status = "selected" if selected else "registered"
    assets.append(
        {
            "asset_id": asset_id,
            "slot_id": slot_id,
            "requirement_id": requirement_id,
            "reference_id": reference_id,
            "path": safe_rel(project, target),
            "asset_type": asset_type,
            "stage": "visual_review",
            "version": "v001",
            "status": status,
            "visibility": visibility,
            "qa_status": qa_status,
            "risk_level": risk_level,
            "prompt_or_edit_ref": prompt_or_edit_ref,
            "notes": "; ".join(item for item in [notes, image_note] if item),
        }
    )
    write_csv_rows(project / "AD-creative/visual_assets/asset_manifest.csv", fields, assets)
    refresh_asset_current_manifest(project)
    review_path = project / "AD-creative/visual_review/review_matrix.csv"
    review_fields, review_rows = read_csv_rows(review_path)
    if review_fields:
        review_rows.append(
            {
                "review_id": next_id(review_rows, "review_id", "VR"),
                "asset_id": asset_id,
                "slot_id": slot_id,
                "direction": "registered_asset",
                "check_type": "file_traceability",
                "result": qa_status,
                "score": "88" if qa_status.upper() == "PASS" else "60",
                "issue": "" if qa_status.upper() == "PASS" else "需要人工视觉筛选。",
                "action": "可用于内部排版；客户可见前再跑最终 Gate。",
                "reviewer": "ad_creative_operator",
                "created_at": now_iso(),
            }
        )
        write_csv_rows(review_path, review_fields, review_rows)
    update_artifact(
        project,
        f"ART-{asset_id}",
        "visual_asset",
        safe_rel(project, target),
        "visual_review",
        status=status,
        visibility=visibility,
        linked_requirements=requirement_id,
        linked_references="" if reference_id == "pending" else reference_id,
        linked_assets=asset_id,
        gate_status=qa_status,
    )
    append_gate(
        project,
        "GATE-AUTO-VISUAL-ASSET-001",
        "visual_review",
        "PASS" if qa_status.upper() == "PASS" else "PARTIAL_PASS",
        "88" if qa_status.upper() == "PASS" else "60",
        f"ART-{asset_id}",
        "" if qa_status.upper() == "PASS" else "图片尚需人工筛选。",
        "客户可见前确认 AI / 版权 / 品牌资产边界。",
        "",
        "ready_for_internal_ppt" if qa_status.upper() == "PASS" else "visual_review_needed",
        "ad_creative_operator",
    )
    append_event(
        project,
        {
            "event_id": f"EVT-{asset_id}",
            "event_type": "visual_asset_registered",
            "created_at": now_iso(),
            "asset_id": asset_id,
            "path": safe_rel(project, target),
            "qa_status": qa_status,
        },
    )
    return asset_id, target


def probe_image(path: Path) -> tuple[int, int, str]:
    try:
        from PIL import Image

        with Image.open(path) as image:
            return image.width, image.height, image.format or path.suffix.lower().strip(".")
    except Exception as exc:  # noqa: BLE001 - gate report should carry the concrete probe failure
        return 0, 0, f"probe_failed={exc}"


def review_visual_quality(
    project: Path,
    *,
    min_long_edge: int = 720,
    min_short_edge: int = 480,
) -> tuple[str, list[str], Path]:
    _, assets = read_csv_rows(project / "AD-creative/visual_assets/asset_manifest.csv")
    current_manifest_path = sync_asset_current_manifest(project)
    _, current_assets = read_csv_rows(current_manifest_path)
    current_by_id = {
        row.get("asset_id", "").strip(): row
        for row in current_assets
        if row.get("asset_id", "").strip()
    }
    _, authorization_rows = read_csv_rows(
        project / "AD-creative/visual_assets/asset_authorizations.csv"
    )
    _, references = read_csv_rows(project / "AD-creative/references/reference_cards.csv")
    reference_ids = unique_rows(references, "reference_id")
    issues: list[str] = []
    warnings: list[str] = []
    evidence: list[str] = [
        f"assets={len(assets)}",
        f"asset_authorizations={len(authorization_rows)}",
        f"min_long_edge={min_long_edge}",
        f"min_short_edge={min_short_edge}",
    ]

    if not assets:
        warnings.append("尚未登记视觉资产。")

    for asset in assets:
        asset_id = asset.get("asset_id", "").strip()
        status = asset.get("status", "").strip().lower()
        visibility = asset.get("visibility", "").strip().lower()
        qa_status = asset.get("qa_status", "").strip().upper()
        risk_level = asset.get("risk_level", "").strip().lower()
        asset_type = asset.get("asset_type", "").strip().lower()
        rel_path = asset.get("path", "").strip()
        prompt_ref = asset.get("prompt_or_edit_ref", "").strip()
        reference_id = asset.get("reference_id", "").strip()
        notes = asset.get("notes", "").strip()
        notes_lower = notes.lower()
        active = status in ACTIVE_ASSET_STATUSES and qa_status != "NOT_RUN"
        selected = status in SELECTED_ASSET_STATUSES
        client_visible = visibility in CLIENT_VISIBLE_VALUES

        if not asset_id:
            issues.append("asset_manifest 存在空 asset_id。")
            continue
        if not asset.get("slot_id", "").strip():
            issues.append(f"{asset_id} 缺少 slot_id。")
        if not asset.get("requirement_id", "").strip():
            warnings.append(f"{asset_id} 缺少 requirement_id。")
        if asset_type == "generated_image" and not prompt_ref:
            issues.append(f"{asset_id} 是生成图但缺少 prompt_or_edit_ref。")
        if prompt_ref and not (project / prompt_ref).exists() and not Path(prompt_ref).is_absolute():
            warnings.append(f"{asset_id} prompt_or_edit_ref 未在项目中找到: {prompt_ref}")
        if reference_id and reference_id != "pending" and reference_id not in reference_ids:
            issues.append(f"{asset_id} reference_id 未登记: {reference_id}")
        if selected and qa_status != "PASS":
            issues.append(f"{asset_id} 已选中但 QA 不是 PASS。")
        if selected and risk_level == "high":
            warnings.append(f"{asset_id} 已选中但 risk_level=high。")
        if client_visible:
            if qa_status != "PASS":
                issues.append(f"{asset_id} 客户可见但 QA 未 PASS。")
            if not selected:
                issues.append(f"{asset_id} 客户可见但未进入 selected/approved/done。")
            current = current_by_id.get(asset_id, {})
            current_sha = current.get("sha256", "").strip()
            if not current_sha:
                issues.append(f"{asset_id} 客户可见但 current manifest 缺少 sha256。")
            elif not matching_asset_authorization(
                project, asset_id, current_sha, authorization_rows
            ):
                issues.append(
                    f"{asset_id} 客户可见但缺少匹配 asset hash/scope 的独立授权 receipt。"
                )
            if reference_id in {"", "pending"}:
                warnings.append(f"{asset_id} 客户可见但未绑定真实 reference_id。")
        if client_visible:
            risk_match = VISUAL_RISK_PATTERN.search(notes_lower)
            if risk_match:
                issues.append(f"{asset_id} 客户可见图片 notes 含风险词: {risk_match.group(0)}")

        if active:
            if not rel_path:
                issues.append(f"{asset_id} active asset 缺少 path。")
                continue
            path = project / rel_path
            if not path.exists():
                issues.append(f"{asset_id} 文件不存在: {rel_path}")
                continue
            if path.suffix.lower() not in GENERATED_IMAGE_SUFFIXES:
                warnings.append(f"{asset_id} 不是常见图片格式: {path.suffix}")
                continue
            width, height, image_format = probe_image(path)
            evidence.append(f"{asset_id}={width}x{height} {image_format}")
            if not width or not height:
                issues.append(f"{asset_id} 图片无法读取: {image_format}")
                continue
            if max(width, height) < min_long_edge or min(width, height) < min_short_edge:
                issues.append(f"{asset_id} 尺寸过低: {width}x{height}")

    status = "PASS" if not issues and not warnings else "PARTIAL_PASS" if not issues else "BLOCKED"
    status = enforce_adversarial_gate_policy(
        project, "visual_review", status, warnings, evidence
    )
    report_path = project / "AD-creative/gates/GATE-AUTO-VISUAL-QUALITY-001_report.md"
    issue_text = "\n".join(f"- {issue}" for issue in issues) or "- 无"
    warning_text = "\n".join(f"- {warning}" for warning in warnings) or "- 无"
    evidence_text = "\n".join(f"- {item}" for item in evidence)
    write_text(
        report_path,
        f"""# Visual Quality Gate

status: {status}
visibility: internal_only
checked_at: {now_iso()}

## Evidence

{evidence_text}

## Blocking Issues

{issue_text}

## Warnings

{warning_text}

## Rules

- active 图片必须存在且可读取。
- selected / approved / done 图片必须 QA PASS。
- 生成图必须有 prompt_or_edit_ref。
- 客户可见图片必须有绑定 exact asset SHA-256、scope、确认者、时间与证据的独立授权 receipt；notes token 不算授权。
- 客户可见图片不能是 contact sheet、低质拼贴、placeholder-only、假 logo。
- 默认最低尺寸：长边 {min_long_edge}px，短边 {min_short_edge}px。
""",
    )
    update_artifact(
        project,
        "ART-AUTO-VISUAL-QUALITY-GATE",
        "visual_quality_gate_report",
        safe_rel(project, report_path),
        "visual_review",
        status="done" if status != "BLOCKED" else "blocked",
        visibility="internal_only",
        linked_assets=";".join(asset.get("asset_id", "") for asset in assets if asset.get("asset_id")),
        gate_status=status,
    )
    append_gate(
        project,
        "GATE-AUTO-VISUAL-QUALITY-001",
        "visual_review",
        status,
        "90" if status == "PASS" else "65" if status == "PARTIAL_PASS" else "35",
        "ART-AUTO-VISUAL-QUALITY-GATE",
        ";".join(issues[:8]),
        ";".join(warnings[:8]) or "修正视觉资产后重跑 visual-quality-gate。",
        "",
        "ready_for_internal_ppt" if status != "BLOCKED" else "fix_visual_assets",
        "ad_creative_operator",
    )
    append_event(
        project,
        {
            "event_id": "EVT-AUTO-VISUAL-QUALITY-GATE",
            "event_type": "visual_quality_gate_run",
            "created_at": now_iso(),
            "status": status,
            "issues": issues[:12],
            "warnings": warnings[:12],
        },
    )
    return status, issues + warnings, report_path


def review_film_quality(project: Path) -> tuple[str, list[str], Path]:
    counts = read_counts(project)
    _, requirements = read_csv_rows(project / "AD-creative/orchestrator/requirements.csv")
    _, references = read_csv_rows(project / "AD-creative/references/reference_cards.csv")
    _, assets = read_csv_rows(project / "AD-creative/visual_assets/asset_manifest.csv")
    current_assets, _ = refresh_asset_current_manifest(project)
    current_by_id = {
        row.get("asset_id", "").strip(): row
        for row in current_assets
        if row.get("asset_id", "").strip()
    }
    _, authorization_rows = read_csv_rows(
        project / "AD-creative/visual_assets/asset_authorizations.csv"
    )
    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    _, exchanges = read_csv_rows(
        project / "AD-creative/orchestrator/specialist_exchange/exchange_index.csv"
    )
    issues: list[str] = []
    warnings: list[str] = []
    evidence: list[str] = [
        f"requirements={counts['requirements']}",
        f"references={counts['references']}",
        f"assets={counts['assets']}",
        f"artifacts={counts['artifacts']}",
    ]
    physical_film_kinds = {
        "film.story_package",
        "film.treatment",
        "film.script",
        "film.shot_plan",
        "film.visual_bible",
        "film.reference_prompt_plan",
    }
    scanned_paths: set[str] = set()
    for artifact in artifacts:
        kind = artifact.get("artifact_type", "").strip()
        if kind not in physical_film_kinds or normalized_artifact_lifecycle(artifact) != "active":
            continue
        artifact_id = artifact.get("artifact_id", "").strip()
        rel_path = artifact.get("path", "").strip()
        expected_sha = artifact.get("sha256", "").strip()
        try:
            physical = contained_project_path(project, rel_path, f"film artifact {artifact_id}")
        except ValueError as exc:
            issues.append(str(exc))
            continue
        if not physical.is_file() or physical.is_symlink():
            issues.append(f"film artifact missing/non-regular: {artifact_id} {rel_path}")
            continue
        actual_sha = file_sha256(physical)
        if not expected_sha or actual_sha != expected_sha:
            issues.append(f"film artifact hash mismatch: {artifact_id} {rel_path}")
            continue
        scanned_paths.add(rel_path)
        evidence.append(
            f"film_artifact id={artifact_id} kind={kind} path={rel_path} sha256={actual_sha}"
        )
    artifacts_by_id = {
        artifact.get("artifact_id", "").strip(): artifact
        for artifact in artifacts
        if artifact.get("artifact_id", "").strip()
    }
    for exchange in exchanges:
        if exchange.get("adoption_decision", "").strip() not in {"adopt", "partial_adopt"}:
            continue
        adoption_rel = exchange.get("adoption_path", "").strip()
        try:
            adoption_path = contained_project_path(
                project, adoption_rel, "specialist adoption"
            )
            if not adoption_path.is_file() or adoption_path.is_symlink():
                raise ValueError("specialist adoption is missing or non-regular")
            adoption_sha = file_sha256(adoption_path)
            expected_adoption_sha = exchange.get("adoption_sha256", "").strip()
            if not expected_adoption_sha or adoption_sha != expected_adoption_sha:
                raise ValueError("specialist adoption hash does not match exchange index")
            adoption = read_json_object(adoption_path, "specialist adoption")
        except (ValueError, FileNotFoundError, json.JSONDecodeError) as exc:
            issues.append(f"invalid adopted specialist exchange: {exc}")
            continue
        if adoption.get("decision") != exchange.get("adoption_decision", "").strip():
            issues.append("invalid adopted specialist exchange: adoption decision mismatch")
            continue
        for output in adoption.get("adopted_outputs", []):
            if isinstance(output, dict):
                target_id = str(output.get("target_artifact_id") or "").strip()
                target_path = str(output.get("target_path") or "").strip()
                target_sha = str(output.get("sha256") or "").strip()
                target_row = artifacts_by_id.get(target_id)
                if not target_row:
                    issues.append(f"adopted specialist output missing artifact_index row: {target_id}")
                    continue
                target_kind = target_row.get("artifact_type", "").strip()
                if target_path != target_row.get("path", "").strip() or target_sha != target_row.get(
                    "sha256", ""
                ).strip():
                    issues.append(
                        f"adopted specialist output path/hash does not match artifact_index: {target_id}"
                    )
                    continue
                if target_kind == "domain.film_qa":
                    try:
                        qa_path = contained_project_path(
                            project, target_path, f"film domain QA {target_id}"
                        )
                    except ValueError as exc:
                        issues.append(str(exc))
                        continue
                    if (
                        not qa_path.is_file()
                        or qa_path.is_symlink()
                        or file_sha256(qa_path) != target_sha
                    ):
                        issues.append(f"film domain QA path/hash mismatch: {target_id}")
                    else:
                        evidence.append(
                            f"film_domain_qa id={target_id} path={target_path} sha256={target_sha}"
                        )
                elif target_kind in physical_film_kinds and target_path not in scanned_paths:
                    issues.append(
                        f"adopted specialist output not scanned by Film Gate: {target_path}"
                    )

    creative_gate = latest_gate(project, gate_id="GATE-AUTO-CREATIVE-QUALITY-001")
    if creative_gate and creative_gate.get("status", "").strip() == "BLOCKED":
        issues.append("Creative Quality Gate 为 BLOCKED，Film Gate 不得推进下游。")

    creative_path = project / "AD-creative/creative/creative_directions.md"
    proposal_path = project / "AD-creative/proposal_architecture/proposal_structure.md"
    treatment_path = project / "AD-creative/film/treatment_packet.md"
    shot_plan_path = project / "AD-creative/film/shot_list_storyboard_plan.md"
    constraints_path = project / "AD-creative/film/production_constraints.md"

    if not requirements:
        issues.append("缺少 requirements，影视/商业判断没有客户事实基线。")
    if not references:
        warnings.append("缺少 reference pack，影视调性和拍法证据不足。")
    if not assets:
        warnings.append("缺少 visual asset，关键画面和资产锁未验证。")
    for label, path in [
        ("creative_directions", creative_path),
        ("proposal_structure", proposal_path),
        ("treatment_packet", treatment_path),
        ("shot_list_storyboard_plan", shot_plan_path),
        ("production_constraints", constraints_path),
    ]:
        if path.exists():
            evidence.append(f"{label}={safe_rel(project, path)}")
        else:
            warnings.append(f"缺少 {safe_rel(project, path)}。")

    visual_gate = latest_gate(project, gate_id="GATE-AUTO-VISUAL-QUALITY-001")
    if not visual_gate:
        warnings.append("尚未运行 Visual Quality Gate。")
    elif visual_gate.get("status") == "BLOCKED":
        issues.append("Visual Quality Gate 为 BLOCKED，不能进入影视级客户审阅。")
    else:
        evidence.append(f"visual_quality_gate={visual_gate.get('status')}")

    client_visible_generated = []
    for asset in assets:
        if asset.get("visibility", "").strip().lower() not in CLIENT_VISIBLE_VALUES:
            continue
        if asset.get("asset_type", "").strip().lower() != "generated_image":
            continue
        asset_id = asset.get("asset_id", "").strip()
        current_sha = current_by_id.get(asset_id, {}).get("sha256", "").strip()
        if not current_sha or not matching_asset_authorization(
            project, asset_id, current_sha, authorization_rows
        ):
            client_visible_generated.append(asset_id)
    if client_visible_generated:
        issues.append(
            "客户可见生成图缺少匹配 asset hash/scope 的独立授权 receipt: "
            + ";".join(client_visible_generated)
        )

    client_visible_artifacts = [
        artifact.get("artifact_id", "")
        for artifact in artifacts
        if artifact.get("visibility", "").strip().lower() in CLIENT_VISIBLE_VALUES
        and artifact.get("gate_status", "").strip().lower() not in PASS_GATE_VALUES
    ]
    if client_visible_artifacts:
        issues.append("客户可见产物 Gate 未 PASS: " + ";".join(client_visible_artifacts))

    commercial_requirements = [
        row for row in requirements
        if row.get("category", "").strip().lower() in {"delivery", "creative", "visual", "research"}
    ]
    if commercial_requirements:
        evidence.append(f"commercial_requirement_rows={len(commercial_requirements)}")
    else:
        warnings.append("未识别到 delivery/creative/visual/research 类商业需求。")

    status = "PASS" if not issues and not warnings else "PARTIAL_PASS" if not issues else "BLOCKED"
    status = enforce_adversarial_gate_policy(
        project, "film_quality", status, warnings, evidence
    )
    report_path = project / "AD-creative/gates/GATE-AUTO-FILM-QUALITY-001_report.md"
    issue_text = "\n".join(f"- {issue}" for issue in issues) or "- 无"
    warning_text = "\n".join(f"- {warning}" for warning in warnings) or "- 无"
    evidence_text = "\n".join(f"- {item}" for item in evidence)
    write_text(
        report_path,
        f"""# Film / Commercial Quality Gate

status: {status}
visibility: internal_only
checked_at: {now_iso()}

## Evidence

{evidence_text}

## Blocking Issues

{issue_text}

## Warnings

{warning_text}

## Review Dimensions

- cinematic_clarity: treatment / shot list / key-frame logic can explain the film idea.
- commercial_message: client value, product role, and campaign output are traceable to requirements.
- brand_fit: visual direction and references do not invent unsupported brand facts.
- product_truth: product, packaging, logo, and claims are not faked.
- production_feasibility: production constraints and known blockers are visible.
- client_risk: client-visible generated assets and artifacts remain gated.
""",
    )
    update_artifact(
        project,
        "ART-AUTO-FILM-QUALITY-GATE",
        "film_quality_gate_report",
        safe_rel(project, report_path),
        "film_quality",
        status="done" if status != "BLOCKED" else "blocked",
        visibility="internal_only",
        linked_assets=";".join(asset.get("asset_id", "") for asset in assets if asset.get("asset_id")),
        gate_status=status,
    )
    append_gate(
        project,
        "GATE-AUTO-FILM-QUALITY-001",
        "film_quality",
        status,
        "90" if status == "PASS" else "68" if status == "PARTIAL_PASS" else "35",
        "ART-AUTO-FILM-QUALITY-GATE",
        ";".join(issues[:8]),
        ";".join(warnings[:8]) or "补齐 treatment / shot list / production constraints 后复核。",
        "",
        "ready_for_internal_ppt_review" if status != "BLOCKED" else "fix_film_quality",
        "ad_creative_operator",
    )
    append_event(
        project,
        {
            "event_id": "EVT-AUTO-FILM-QUALITY-GATE",
            "event_type": "film_quality_gate_run",
            "created_at": now_iso(),
            "status": status,
            "issues": issues[:12],
            "warnings": warnings[:12],
        },
    )
    return status, issues + warnings, report_path


def latest_generated_image(root: Path) -> Path:
    if not root.exists():
        raise FileNotFoundError(f"generated_images root not found: {root}")
    candidates = [
        path
        for path in root.rglob("*")
        if path.is_file() and path.suffix.lower() in GENERATED_IMAGE_SUFFIXES
    ]
    if not candidates:
        raise FileNotFoundError(f"no generated image found under: {root}")
    return max(candidates, key=lambda path: path.stat().st_mtime)


def resolve_imagegen_source(raw_file: str) -> tuple[Path, Path]:
    root = codex_generated_images_root().resolve()
    source = Path(raw_file).expanduser().resolve() if raw_file else latest_generated_image(root)
    if source.suffix.lower() not in GENERATED_IMAGE_SUFFIXES:
        raise ValueError(f"unsupported generated image suffix: {source.suffix}")
    if not source.is_file():
        raise FileNotFoundError(f"generated image not found: {source}")
    try:
        source.relative_to(root)
    except ValueError as exc:
        raise ValueError(f"source must be under {root}: {source}") from exc
    return source, root


def default_prompt_ref(project: Path) -> str:
    for rel_path in [
        "AD-creative/image_jobs/image_prompt_pack.json",
        "AD-creative/image_jobs/image_generation_policy.md",
        "AD-creative/image_jobs/image_job_spec_template.md",
    ]:
        if (project / rel_path).exists():
            return rel_path
    return ""


def md_cell(value: str) -> str:
    return value.replace("|", "\\|").replace("\n", " ").strip()


def append_imagegen_import_log(
    project: Path,
    source_root: Path,
    source: Path,
    asset_id: str,
    target: Path,
    prompt_ref: str,
    qa_status: str,
    notes: str,
) -> Path:
    path = project / "AD-creative/image_jobs/imagegen_import_log.md"
    path.parent.mkdir(parents=True, exist_ok=True)
    should_write_header = not path.exists()
    with path.open("a", encoding="utf-8") as handle:
        if should_write_header:
            handle.write(
                "# ImageGen Import Log\n\n"
                "visibility: internal_only\n\n"
                "| imported_at | asset_id | source | target | prompt_ref | qa_status | notes |\n"
                "| --- | --- | --- | --- | --- | --- | --- |\n"
            )
        handle.write(
            "| {imported_at} | {asset_id} | {source} | {target} | {prompt_ref} | {qa_status} | {notes} |\n".format(
                imported_at=now_iso(),
                asset_id=md_cell(asset_id),
                source=md_cell(safe_rel_to(source_root, source)),
                target=md_cell(safe_rel(project, target)),
                prompt_ref=md_cell(prompt_ref),
                qa_status=md_cell(qa_status),
                notes=md_cell(notes),
            )
        )

    status = "PASS" if qa_status.upper() == "PASS" else "PARTIAL_PASS"
    update_artifact(
        project,
        "ART-AUTO-IMAGEGEN-IMPORT",
        "imagegen_import_log",
        safe_rel(project, path),
        "visual_review",
        status="done",
        visibility="internal_only",
        linked_assets=asset_id,
        gate_status=status,
    )
    append_event(
        project,
        {
            "event_id": f"EVT-IMAGEGEN-IMPORT-{asset_id}",
            "event_type": "imagegen_asset_imported",
            "created_at": now_iso(),
            "asset_id": asset_id,
            "source": safe_rel_to(source_root, source),
            "target": safe_rel(project, target),
            "qa_status": qa_status,
        },
    )
    return path


def perform_intake(
    project: Path,
    source_ids: list[str],
    goal: str,
    *,
    max_total_chars: int = 2_000_000,
) -> dict[str, int]:
    source_path = project / "AD-creative/orchestrator/source_events.csv"
    _, source_rows = read_csv_rows(source_path)
    source_id_set = set(source_ids)
    target_sources = [
        row for row in source_rows if not source_id_set or row.get("source_event_id") in source_id_set
    ]
    intake_result = run_evidence_intake(
        project,
        target_sources,
        max_total_chars=max_total_chars,
    )
    requirement_path = project / "AD-creative/orchestrator/requirements.csv"
    _, requirement_rows = read_csv_rows(requirement_path)
    gap_path = project / "AD-creative/orchestrator/gaps.csv"
    _, gap_rows = read_csv_rows(gap_path)
    delivery_surface = project_surface(project) == DELIVERY_SURFACE
    blocking_gap_rows = [
        row
        for row in gap_rows
        if row.get("status", "").strip().lower() not in {"closed", "resolved", "done"}
        and row.get("impact", "").strip().lower() in {"blocking", "high_impact"}
    ]
    current_truth_path = project / "AD-creative/orchestrator/current_truth.md"
    confirmed = "\n".join(f"- {row['statement']}" for row in requirement_rows[:12]) or "- 暂无已抽取需求"
    open_questions = "\n".join(
        f"- {row.get('question_for_client') or row.get('description')}" for row in gap_rows[:8]
    ) or "- 暂无"
    if delivery_surface:
        inferred = (
            "- 当前处于 intake；已从本地资料抽取第一轮需求和缺口。\n"
            "- 客户可见稿前需要独立质量与发送准备检查；生成图默认 internal_only。"
        )
        next_action = "按缺口向客户/内部负责人追问；确认文本框架后再进入视觉与 PPT。"
    else:
        inferred = "- 当前处于内容理解阶段；先形成可用判断，再按风险决定是否进入交付治理。"
        next_action = (
            "处理真正阻塞的缺口，同时继续不受影响的内部内容工作。"
            if blocking_gap_rows
            else "基于当前证据完成本轮内部广告内容产出。"
        )
    update_markdown_sections(
        current_truth_path,
        {
            "Project": project.name,
            "Confirmed": confirmed,
            "Inferred": inferred,
            "Conflicted": "- 暂无自动识别冲突。",
            "Deprecated": "- 暂无。",
            "Open Questions": open_questions,
            "Current Stage": "intake",
            "Next Action": next_action,
        },
    )
    if delivery_surface:
        new_requirements = intake_result.new_requirements
        new_gaps = intake_result.new_gaps
        linked_req_ids = ";".join(
            row["requirement_id"] for row in new_requirements
        ) or ";".join(row.get("requirement_id", "") for row in requirement_rows[:8])
        linked_source_ids = ";".join(source_ids) if source_ids else ";".join(
            row.get("source_event_id", "") for row in target_sources
        )
        question_rows = "\n".join(
            f"| {row['gap_id']} | {row.get('question_for_client') or row['description']} | "
            f"{row['recommended_action']} | {row['impact']} | 客户补充 / 先内部推进 |"
            for row in gap_rows[:8]
        ) or "| - | 暂无 | - | - | - |"
        write_text(
            project / "AD-creative/handoff/客户追问话术.md",
            f"""# 客户追问话术

## 建议询问
| ID | 问题 | 推荐动作 | 影响 | 可选动作 |
|---|---|---|---|---|
{question_rows}

## 可复制话术
我们已先把资料整理成需求和缺口。为避免误判，请补充品牌资产、产品高清图/包装规范、参考片边界、AI 图是否可用于客户审阅，以及本轮希望看到的方案精度。
""",
        )

        work_path = project / "AD-creative/orchestrator/work_items.csv"
        work_fields, work_rows = read_csv_rows(work_path)
        intake_work_id = "WORK-001"
        for row in work_rows:
            if row.get("stage") == "intake" and row.get("title") == "需求整理与缺口判断":
                intake_work_id = row.get("work_id", intake_work_id) or intake_work_id
                row["status"] = "done"
                row["output_artifacts"] = (
                    "ART-AUTO-CURRENT-TRUTH;ART-AUTO-CLIENT-QUESTIONS"
                )
                row["linked_requirements"] = linked_req_ids
                row["linked_source_events"] = linked_source_ids
                row["updated_at"] = now_iso()
                break
        write_csv_rows(work_path, work_fields, work_rows)

        artifact_path = project / "AD-creative/orchestrator/artifact_index.csv"
        recorded_at = now_iso()
        for artifact_id, artifact_type, rel_path in [
            (
                "ART-AUTO-CURRENT-TRUTH",
                "intake_report",
                "AD-creative/orchestrator/current_truth.md",
            ),
            (
                "ART-AUTO-CLIENT-QUESTIONS",
                "client_questions",
                "AD-creative/handoff/客户追问话术.md",
            ),
        ]:
            update_or_append_csv_row(
                artifact_path,
                "artifact_id",
                {
                    "artifact_id": artifact_id,
                    "artifact_type": artifact_type,
                    "path": rel_path,
                    "stage": "intake",
                    "version": "v001",
                    "status": "done",
                    "visibility": "internal_only",
                    "source_event_ids": linked_source_ids,
                    "linked_requirements": linked_req_ids,
                    "linked_work_items": intake_work_id,
                    "linked_references": "",
                    "linked_assets": "",
                    "gate_status": "PARTIAL_PASS",
                    "supersedes_artifact_id": "",
                    "created_at": recorded_at,
                    "updated_at": recorded_at,
                },
            )
        append_gate(
            project,
            "GATE-AUTO-BRIEF-001",
            "intake",
            "PARTIAL_PASS" if gap_rows else "PASS",
            "72",
            "ART-AUTO-CURRENT-TRUTH;ART-AUTO-CLIENT-QUESTIONS",
            ";".join(row["description"] for row in gap_rows[:5]),
            "补齐缺口后进入 research_plan。",
            ";".join(row.get("question_for_client", "") for row in gap_rows[:5]),
            "research_plan",
            "ad_creative_operator",
        )
        append_event(
            project,
            {
                "event_id": f"EVT-INTAKE-{recorded_at}",
                "event_type": "intake_completed",
                "created_at": recorded_at,
                "actor": "ad_creative_operator",
                "summary": (
                    f"Extracted {len(new_requirements)} requirements and "
                    f"{len(new_gaps)} gaps from local materials."
                ),
            },
        )
    return intake_result.stats()


def render_handoff(project: Path, goal: str, source_ids: list[str]) -> dict[str, object]:
    _, requirement_rows = read_csv_rows(
        project / "AD-creative/orchestrator/requirements.csv"
    )
    _, gap_rows = read_csv_rows(project / "AD-creative/orchestrator/gaps.csv")
    open_gaps = [
        row
        for row in gap_rows
        if row.get("status", "").strip().lower() not in {"closed", "resolved", "done"}
    ]
    blocking_gaps = [
        row
        for row in open_gaps
        if row.get("impact", "").strip().lower() in {"blocking", "high_impact"}
    ]
    non_blocking_gaps = [row for row in open_gaps if row not in blocking_gaps]
    facts = load_fact_inventory(project)
    fact_lines = [
        f"{fact.fact_key}: {fact.value or fact.state}"
        for fact in facts
        if fact.state == "present"
    ][:8]
    current_source_ids = set(source_ids)
    ordered_requirement_rows = (
        [
            row
            for row in requirement_rows
            if row.get("source_event_id", "") in current_source_ids
        ]
        + [
            row
            for row in requirement_rows
            if row.get("source_event_id", "") not in current_source_ids
        ]
        if current_source_ids
        else requirement_rows
    )
    requirement_lines = [
        row.get("statement", "").strip()
        for row in ordered_requirement_rows
        if row.get("statement", "").strip()
    ][:8]
    blocking_gap_lines = [
        row.get("description", "").strip()
        for row in blocking_gaps
        if row.get("description", "").strip()
    ][:6]
    non_blocking_gap_lines = [
        row.get("description", "").strip()
        for row in non_blocking_gaps
        if row.get("description", "").strip()
    ][:6]
    requirement_types = {
        row.get("requirement_type", "").strip().lower()
        for row in ordered_requirement_rows
    }
    proceed: list[str] = []
    if requirement_lines:
        proceed.append("可以基于已锁定要求继续内部分析和方案构思，不需要先建立交付账本。")
    if requirement_types & {"creative", "visual", "research"}:
        proceed.append("下一步应把证据转成内容判断或创意 brief，并优先检查真实素材语义。")
    if "delivery" in requirement_types:
        proceed.append("交付格式已被识别，但只有真正进入客户可见版本时才升级到 Delivery Surface。")
    if not proceed:
        proceed.append("材料已完成证据化读取；需要补充一个明确产出目标，才能继续内容工作。")
    if blocking_gap_lines:
        next_action = "先确认最小阻塞项；不受其影响的内部内容工作可以继续。"
    elif requirement_lines:
        next_action = "基于当前证据生成或更新 creative brief，再进入专业内容推理。"
    else:
        next_action = "补充本轮希望得到的具体广告产出。"

    facts_markdown = "\n".join(f"- {item}" for item in fact_lines) or "- 暂无可确认的结构化事实"
    requirements_markdown = (
        "\n".join(f"- {item}" for item in requirement_lines)
        or "- 暂无可确认的明确要求"
    )
    gaps_markdown = (
        "\n".join(f"- {item}" for item in blocking_gap_lines)
        or "- 无真实阻塞缺口"
    )
    non_blocking_markdown = (
        "\n".join(f"- {item}" for item in non_blocking_gap_lines)
        or "- 无"
    )
    proceed_markdown = "\n".join(f"- {item}" for item in proceed)
    objective = goal.strip() or "理解材料并推进当前广告任务"
    answer_markdown = f"""## 当前目标
{objective}

## 材料事实
{facts_markdown}

## 明确要求
{requirements_markdown}

## 真正阻塞
{gaps_markdown}

## 非阻塞未知
{non_blocking_markdown}

## 现在可以推进
{proceed_markdown}

## 下一步
{next_action}"""
    board_path = project / "AD-creative/handoff/项目看板.md"
    write_text(board_path, "# 工作摘要\n\n" + answer_markdown + "\n")
    confirmation_gaps = (
        open_gaps
        if project_surface(project) == DELIVERY_SURFACE
        else blocking_gaps
    )
    question_rows = "\n".join(
        "| {id} | {question} | {action} |".format(
            id=row.get("gap_id", ""),
            question=(
                row.get("question_for_user")
                or row.get("question_for_client")
                or row.get("description", "")
            ),
            action=row.get("recommended_action", ""),
        )
        for row in confirmation_gaps[:8]
    ) or "| - | 无阻塞确认 | 继续内部内容工作 |"
    write_text(
        project / "AD-creative/handoff/待你确认.md",
        "# 待你确认\n\n| ID | 问题 | 推荐动作 |\n|---|---|---|\n"
        + question_rows
        + "\n",
    )
    return {
        "objective": objective,
        "facts": fact_lines,
        "requirements": requirement_lines,
        "blocking_gaps": blocking_gap_lines,
        "non_blocking_unknowns": non_blocking_gap_lines,
        "can_proceed": proceed,
        "next_action": next_action,
        "markdown": answer_markdown,
        "path": str(board_path),
        "source_ids": source_ids,
    }


def cell(value: str | None) -> str:
    return html.escape(value or "")


def first_nonempty(*values: str | None, default: str = "-") -> str:
    for value in values:
        if value:
            return value
    return default


def project_stage(project: Path) -> str:
    path = project / "AD-creative/orchestrator/project.yml"
    if not path.exists():
        return "unknown"
    for line in path.read_text(encoding="utf-8").splitlines():
        stripped = line.strip()
        if stripped.startswith("stage:"):
            return stripped.split(":", 1)[1].strip() or "intake"
    return "intake"


def rows_to_table(rows: list[dict[str, str]], columns: list[tuple[str, str]], empty: str) -> str:
    if not rows:
        return f"<tr><td colspan=\"{len(columns)}\" class=\"empty\">{html.escape(empty)}</td></tr>"
    body = []
    for row in rows:
        cells = "".join(
            f'<td data-label="{cell(label)}">{cell(row.get(key))}</td>' for key, label in columns
        )
        body.append(f"<tr>{cells}</tr>")
    return "\n".join(body)


def html_json(value: object) -> str:
    return json.dumps(value, ensure_ascii=False).replace("<", "\\u003c")


def parse_confirmation_rows(path: Path) -> list[dict[str, str]]:
    if not path.exists():
        return []
    rows: list[dict[str, str]] = []
    for line in path.read_text(encoding="utf-8").splitlines():
        if not line.startswith("|"):
            continue
        parts = [part.strip() for part in line.strip("|").split("|")]
        if len(parts) < 3 or parts[0] in CONFIRMATION_HEADER_CELLS or set(parts[0]) == {"-"}:
            continue
        rows.append(
            {
                "id": parts[0],
                "question": parts[1],
                "recommendation": parts[2],
                "impact": parts[3] if len(parts) > 3 else "",
                "actions": parts[4] if len(parts) > 4 else parts[2],
            }
        )
    return rows


def extract_markdown_section(text: str, heading: str) -> str:
    marker = f"## {heading}"
    start = text.find(marker)
    if start == -1:
        return ""
    body = text[start + len(marker):]
    next_heading = body.find("\n## ")
    if next_heading != -1:
        body = body[:next_heading]
    return body.strip()


def first_section_line(text: str, heading: str) -> str:
    section = extract_markdown_section(text, heading)
    for line in section.splitlines():
        stripped = line.strip().strip("-").strip()
        if stripped and not stripped.startswith("|"):
            return stripped
    return ""


def parse_key_value_line(text: str, key: str) -> str:
    prefix = f"{key}:"
    for line in text.splitlines():
        if line.startswith(prefix):
            return line.split(":", 1)[1].strip()
    return ""


def goal_iteration_rows(project: Path) -> list[dict[str, str]]:
    goal_dir = project / GOAL_ITERATIONS_REL
    if not goal_dir.exists():
        return []
    rows: list[dict[str, str]] = []
    for path in sorted(goal_dir.glob("*.md"), key=lambda item: item.stat().st_mtime, reverse=True):
        text = path.read_text(encoding="utf-8", errors="ignore")
        rows.append(
            {
                "goal_id": parse_key_value_line(text, "goal_id") or path.stem,
                "status": parse_key_value_line(text, "status") or "unknown",
                "title": parse_key_value_line(text, "goal_title") or path.stem,
                "owner": parse_key_value_line(text, "owner") or "Main Controller",
                "updated_at": parse_key_value_line(text, "updated_at"),
                "next": first_section_line(text, "Next Iteration Queue") or "按下一阶段 Gate 推进",
                "path": safe_rel(project, path),
            }
        )
    return rows


def latest_goal_row(project: Path) -> dict[str, str] | None:
    rows = goal_iteration_rows(project)
    return rows[0] if rows else None


def resolve_goal_row(project: Path, goal_id: str) -> dict[str, str] | None:
    rows = goal_iteration_rows(project)
    if not rows:
        return None
    if goal_id in {"", "latest"}:
        return rows[0]
    for row in rows:
        if row.get("goal_id") == goal_id:
            return row
    return None


def slug_text(value: str, default: str = "run") -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "-", value.strip().lower()).strip("-")
    return slug or default


def creative_production_root_candidates() -> list[Path]:
    candidates: list[Path] = []
    env_root = os.environ.get("ADCO_CREATIVE_PRODUCTION_ROOT", "").strip()
    if env_root:
        candidates.append(Path(env_root).expanduser())

    cache_root = Path.home() / ".codex/plugins/cache"
    for pattern in [
        "openai-curated-remote/creative-production/*",
        "openai-curated/creative-production/*",
        "creative-production/*",
    ]:
        if cache_root.exists():
            candidates.extend(sorted(cache_root.glob(pattern), reverse=True))

    source = source_root()
    if source:
        candidates.extend(
            [
                source / "plugins/creative-production",
                source / "creative-production",
            ]
        )
    return candidates


def is_creative_production_root(path: Path) -> bool:
    return all(
        (path / rel).exists()
        for rel in [
            "skills/ads-explorer/scripts/build_ads_explorer.py",
            "skills/shot-explorer/scripts/create_shot_explorer.py",
            "skills/moodboard-explorer/scripts/create_mood_board.py",
            "scripts/review_renderer.py",
        ]
    )


def resolve_creative_production_root() -> Path | None:
    for candidate in creative_production_root_candidates():
        root = candidate.expanduser().resolve()
        if is_creative_production_root(root):
            return root
    return None


def creative_production_script(root: Path, kind: str) -> Path:
    scripts = {
        "ads": root / "skills/ads-explorer/scripts/build_ads_explorer.py",
        "shots": root / "skills/shot-explorer/scripts/create_shot_explorer.py",
        "moodboard": root / "skills/moodboard-explorer/scripts/create_mood_board.py",
    }
    return scripts[kind]


def creative_doctor_report() -> tuple[str, list[str], list[str], list[str]]:
    issues: list[str] = []
    warnings: list[str] = []
    evidence: list[str] = []
    root = resolve_creative_production_root()
    if not root:
        issues.append("Creative Production plugin root not found.")
        evidence.append("ADCO_CREATIVE_PRODUCTION_ROOT=" + os.environ.get("ADCO_CREATIVE_PRODUCTION_ROOT", ""))
        evidence.append("searched_candidates=" + str(len(creative_production_root_candidates())))
        return "BLOCKED", issues, warnings, evidence

    evidence.append(f"creative_production_root={root}")
    for kind in sorted(CREATIVE_PRODUCTION_KINDS):
        script = creative_production_script(root, kind)
        evidence.append(f"{kind}_script={script}")
        if not script.exists():
            issues.append(f"missing {kind} script: {script}")
    try:
        import PIL  # noqa: F401
    except Exception as exc:  # noqa: BLE001 - optional bridge should report exact reason
        warnings.append(f"Pillow import check failed: {exc}")
    return ("PASS" if not issues else "BLOCKED"), issues, warnings, evidence


def creative_output_dir(project: Path, kind: str, work_id: str) -> Path:
    stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    return project / "AD-creative/creative_production/runs" / f"{kind}-{slug_text(work_id)}-{stamp}"


def brief_title(brief_file: Path, work_id: str) -> str:
    try:
        for line in brief_file.read_text(encoding="utf-8", errors="ignore").splitlines():
            stripped = line.strip("# ").strip()
            if stripped:
                return stripped[:80]
    except FileNotFoundError:
        pass
    return work_id


def moodboard_spec_from_brief(brief_file: Path, work_id: str, output: Path) -> Path:
    text = brief_file.read_text(encoding="utf-8", errors="ignore")[:2400]
    title = brief_title(brief_file, work_id)
    cues = [
        ("audience-feel", "Audience feel", "single editorial photograph showing the target audience mood and occasion"),
        ("brand-world", "Brand world", "single cinematic brand-world reference image with setting, light, and material cues"),
        ("product-moment", "Product moment", "single commercial product-in-use visual reference image"),
        ("texture-light", "Texture and light", "single material and lighting study with premium production value"),
        ("story-beat", "Story beat", "single film-still style moment that can become a key frame"),
        ("campaign-kv", "Campaign key visual", "single clean advertising key visual reference without readable text"),
    ]
    items = []
    for index, (item_id, item_title, prompt_seed) in enumerate(cues, start=1):
        items.append(
            {
                "id": f"{item_id}-{index}",
                "title": item_title,
                "caption": f"{item_title} for {title}",
                "source": "ADCO brief",
                "tone": "commercial cinematic",
                "motif": item_title,
                "palette": "derived from brief",
                "prompt": (
                    f"{prompt_seed}. Preserve this business brief: {text}. "
                    "No logos, no readable copy, no collage, no contact sheet, no internal notes."
                ),
            }
        )
    spec = {
        "meta": {
            "title": f"{title} mood board",
            "summary": "ADCO review-only mood board generated from project brief.",
            "source": str(brief_file),
        },
        "signals": {
            "work_id": work_id,
            "brief_excerpt": text[:800],
            "visibility": "internal_only",
        },
        "items": items,
    }
    spec_path = output.parent / f"{output.name}-spec.json"
    write_text(spec_path, json.dumps(spec, ensure_ascii=False, indent=2))
    return spec_path


def register_agent_run(
    project: Path,
    *,
    work_id: str,
    role: str,
    status: str,
    input_files: str,
    output_files: str,
    gate_id: str = "",
    summary: str = "",
    next_action: str = "",
) -> str:
    path = project / "AD-creative/orchestrator/agent_runs.csv"
    fields, rows = read_csv_rows(path)
    if not fields:
        return ""
    run_id = next_id(rows, "run_id", "RUN")
    rows.append(
        {
            "run_id": run_id,
            "work_id": work_id,
            "agent_role": role,
            "status": status,
            "started_at": now_iso(),
            "completed_at": now_iso(),
            "input_files": input_files,
            "output_files": output_files,
            "gate_id": gate_id,
            "summary": summary,
            "next_action": next_action,
        }
    )
    write_csv_rows(path, fields, rows)
    return run_id


def work_id_exists(project: Path, work_id: str) -> bool:
    _, rows = read_csv_rows(project / "AD-creative/orchestrator/work_items.csv")
    return any(row.get("work_id") == work_id for row in rows)


def run_creative_production(
    project: Path,
    *,
    kind: str,
    work_id: str,
    brief_file: Path,
    base_asset: Path | None = None,
    generate: bool = False,
    force: bool = False,
) -> tuple[Path, list[str]]:
    root = resolve_creative_production_root()
    if not root:
        raise FileNotFoundError("Creative Production plugin root not found.")
    if kind not in CREATIVE_PRODUCTION_KINDS:
        raise ValueError(f"unknown creative production kind: {kind}")
    if not brief_file.exists():
        raise FileNotFoundError(f"brief file not found: {brief_file}")
    if not work_id_exists(project, work_id):
        raise ValueError(f"unknown work_id: {work_id}")

    out_dir = creative_output_dir(project, kind, work_id)
    out_dir.parent.mkdir(parents=True, exist_ok=True)
    script = creative_production_script(root, kind)
    cmd: list[str]
    if kind == "ads":
        cmd = [
            sys.executable,
            str(script),
            "--ad-name",
            brief_title(brief_file, work_id),
            "--ad-brief-file",
            str(brief_file),
            "--out-dir",
            str(out_dir),
            "--force",
        ]
        if base_asset:
            cmd.extend(["--reference-image", str(base_asset)])
        cmd.append("--generate" if generate else "--review-only")
    elif kind == "shots":
        cmd = [
            sys.executable,
            str(script),
            "--output",
            str(out_dir),
            "--force",
            "--generate" if generate else "--review-only",
        ]
        if base_asset:
            cmd.extend(["--base-asset", str(base_asset)])
    else:
        spec_path = moodboard_spec_from_brief(brief_file, work_id, out_dir)
        cmd = [
            sys.executable,
            str(script),
            "--spec",
            str(spec_path),
            "--output",
            str(out_dir),
            "--force",
        ]

    if out_dir.exists() and force:
        shutil.rmtree(out_dir)
    completed = subprocess.run(cmd, check=False, text=True, capture_output=True)
    log_path = out_dir.parent / f"{out_dir.name}-command.json"
    write_text(
        log_path,
        json.dumps(
            {
                "command": cmd,
                "returncode": completed.returncode,
                "stdout": completed.stdout[-4000:],
                "stderr": completed.stderr[-4000:],
                "generate": generate,
            },
            ensure_ascii=False,
            indent=2,
        ),
    )
    if completed.returncode != 0:
        raise RuntimeError(f"Creative Production {kind} failed; see {log_path}")

    review = out_dir / ("mood-board.html" if kind == "moodboard" else "review-board.html")
    if not review.exists() and kind != "moodboard":
        review = out_dir
    update_artifact(
        project,
        f"ART-CP-{safe_artifact_suffix(out_dir.name)}",
        "creative_production_run",
        safe_rel(project, review),
        "visual_plan",
        status="done",
        visibility="internal_only",
        linked_work_items=work_id,
        gate_status="PARTIAL_PASS",
    )
    register_agent_run(
        project,
        work_id=work_id,
        role=f"creative_production_{kind}",
        status="done",
        input_files=safe_rel(project, brief_file),
        output_files=safe_rel(project, out_dir),
        summary=f"Creative Production {kind} {'generated' if generate else 'review-only'} run.",
        next_action="Import with adco import-creative-production before visual gate.",
    )
    append_event(
        project,
        {
            "event_id": f"EVT-CP-{safe_artifact_suffix(out_dir.name)}",
            "event_type": "creative_production_run",
            "created_at": now_iso(),
            "kind": kind,
            "work_id": work_id,
            "run_dir": safe_rel(project, out_dir),
            "generate": generate,
        },
    )
    render_dashboard(project)
    return out_dir, [safe_rel(project, log_path)]


def load_json_file(path: Path) -> object:
    return json.loads(path.read_text(encoding="utf-8"))


def creative_manifest_paths(run_dir: Path, kind: str) -> dict[str, Path]:
    if kind == "shots":
        return {
            "manifest": run_dir / "data/prompts-manifest.json",
            "jobs": run_dir / "data/jobs.jsonl",
            "review": run_dir / "review-board.html",
            "widget": run_dir / "moodboard-widget-payload.json",
        }
    if kind == "moodboard":
        return {
            "manifest": run_dir / "data/stream.json",
            "jobs": run_dir / "data/stream-static.json",
            "review": run_dir / "mood-board.html",
            "widget": run_dir / "data/stream.json",
        }
    return {
        "manifest": run_dir / "prompts-manifest.json",
        "jobs": run_dir / "jobs.jsonl",
        "review": run_dir / "review-board.html",
        "widget": run_dir / "moodboard-widget-payload.json",
    }


def prompt_items_from_manifest(path: Path, kind: str) -> list[dict[str, object]]:
    if not path.exists():
        return []
    raw = load_json_file(path)
    if isinstance(raw, list):
        return [item for item in raw if isinstance(item, dict)]
    if isinstance(raw, dict):
        if isinstance(raw.get("items"), list):
            return [item for item in raw["items"] if isinstance(item, dict)]
        if isinstance(raw.get("routes"), list):
            return [item for item in raw["routes"] if isinstance(item, dict)]
    return []


def item_image_path(run_dir: Path, item: dict[str, object]) -> Path | None:
    for key in ("output", "src", "image", "imageUrl", "previewImageUrl", "sourceImageUrl", "path", "url"):
        value = str(item.get(key) or "").strip()
        if not value or value.startswith("data:") or value.startswith("http"):
            continue
        value = value.lstrip("/")
        path = run_dir / value
        if path.exists() and path.suffix.lower() in GENERATED_IMAGE_SUFFIXES:
            return path
    return None


def copy_creative_metadata(project: Path, run_dir: Path, kind: str, manifest_paths: dict[str, Path]) -> tuple[Path, list[str]]:
    target_dir = project / "AD-creative/image_jobs/creative_production" / slug_text(run_dir.name)
    target_dir.mkdir(parents=True, exist_ok=True)
    copied: list[str] = []
    for label, source in manifest_paths.items():
        if source.exists() and source.is_file():
            target = target_dir / f"{label}{source.suffix or '.txt'}"
            shutil.copy2(source, target)
            copied.append(safe_rel(project, target))
    summary = {
        "kind": kind,
        "source_run_dir": str(run_dir.resolve()),
        "copied_files": copied,
        "imported_at": now_iso(),
        "visibility": "internal_only",
    }
    write_text(target_dir / "adco_import_summary.json", json.dumps(summary, ensure_ascii=False, indent=2))
    return target_dir, copied


def import_creative_production_run(
    project: Path,
    *,
    run_dir: Path,
    kind: str,
    slot_prefix: str,
    requirement_id: str = "",
    reference_id: str = "pending",
) -> tuple[list[str], Path]:
    if kind not in CREATIVE_PRODUCTION_KINDS:
        raise ValueError(f"unknown creative production kind: {kind}")
    if not run_dir.exists():
        raise FileNotFoundError(f"run dir not found: {run_dir}")

    manifest_paths = creative_manifest_paths(run_dir, kind)
    metadata_dir, copied = copy_creative_metadata(project, run_dir, kind, manifest_paths)
    manifest_rel = safe_rel(project, metadata_dir / "manifest.json") if (metadata_dir / "manifest.json").exists() else ""
    manifest = prompt_items_from_manifest(manifest_paths["manifest"], kind)
    imported_asset_ids: list[str] = []

    for index, item in enumerate(manifest, start=1):
        image_path = item_image_path(run_dir, item)
        if not image_path:
            continue
        slot_id = f"{slot_prefix}-{index:02d}"
        prompt_ref = manifest_rel or safe_rel(project, metadata_dir)
        label = str(item.get("title") or item.get("label") or item.get("id") or slot_id)
        asset_id, _ = add_visual_asset(
            project,
            image_path,
            slot_id,
            requirement_id,
            reference_id,
            "generated_image",
            "internal_only",
            "PARTIAL_PASS",
            "medium",
            prompt_ref,
            f"Imported from Creative Production {kind} run {run_dir.name}; route={label}; internal_only.",
            False,
        )
        imported_asset_ids.append(asset_id)

    checked_artifacts = []
    run_artifact_id = f"ART-CP-IMPORT-{safe_artifact_suffix(run_dir.name)}"
    review_path = manifest_paths.get("review")
    if review_path and review_path.exists():
        rel_review = safe_rel(project, review_path)
    else:
        rel_review = safe_rel(project, metadata_dir)
    update_artifact(
        project,
        run_artifact_id,
        "creative_production_import",
        rel_review,
        "visual_plan",
        status="done",
        visibility="internal_only",
        linked_assets=";".join(imported_asset_ids),
        gate_status="PARTIAL_PASS",
    )
    checked_artifacts.append(run_artifact_id)
    append_gate(
        project,
        f"GATE-CP-IMPORT-{safe_artifact_suffix(run_dir.name)}",
        "visual_plan",
        "PARTIAL_PASS",
        "70",
        ";".join(checked_artifacts),
        "",
        "Creative Production outputs are internal_only until visual-quality-gate and human review pass.",
        "",
        "run_visual_quality_gate",
        "ad_creative_operator",
    )
    append_event(
        project,
        {
            "event_id": f"EVT-CP-IMPORT-{safe_artifact_suffix(run_dir.name)}",
            "event_type": "creative_production_imported",
            "created_at": now_iso(),
            "kind": kind,
            "run_dir": str(run_dir.resolve()),
            "metadata_dir": safe_rel(project, metadata_dir),
            "imported_assets": imported_asset_ids,
        },
    )
    render_dashboard(project)
    return imported_asset_ids, metadata_dir


def latest_gate(project: Path, stage: str | None = None, gate_id: str | None = None) -> dict[str, str] | None:
    _, gates = read_csv_rows(project / "AD-creative/orchestrator/gate_log.csv")
    candidates = []
    for gate in gates:
        if gate_id and gate.get("gate_id") != gate_id:
            continue
        if stage and gate.get("stage") != stage:
            continue
        candidates.append(gate)
    if not candidates:
        return None
    return candidates[-1]


def has_gate(
    project: Path,
    gate_id: str,
    allowed_statuses: set[str] | None = None,
) -> bool:
    gate = latest_gate(project, gate_id=gate_id)
    if not gate:
        return False
    allowed = allowed_statuses or {"PASS", "PARTIAL_PASS"}
    if gate.get("status", "").strip().upper() not in allowed:
        return False
    if gate_id.startswith("GATE-AUTO-") and not gate_target_is_fresh(project, gate):
        return False
    return True


def gate_target_is_fresh(project: Path, gate: dict[str, str]) -> bool:
    target_ref = gate.get("target_ref", "").strip()
    target_sha256 = gate.get("target_sha256", "").strip()
    if not target_ref or not re.fullmatch(r"[0-9a-f]{64}", target_sha256):
        return False
    try:
        target = contained_project_path(project, target_ref, "gate target_ref")
    except ValueError:
        return False
    return target.is_file() and file_sha256(target) == target_sha256


def latest_gate_rows(project: Path) -> list[dict[str, str]]:
    _, gates = read_csv_rows(project / "AD-creative/orchestrator/gate_log.csv")
    latest_by_id: dict[str, dict[str, str]] = {}
    order: list[str] = []
    for gate in gates:
        gate_id = gate.get("gate_id", "").strip()
        if not gate_id:
            continue
        if gate_id not in latest_by_id:
            order.append(gate_id)
        latest_by_id[gate_id] = gate
    return [latest_by_id[gate_id] for gate_id in order]


def derive_goal_phase(project: Path) -> str:
    counts = read_counts(project)
    if counts["source_events"] == 0 or counts["requirements"] == 0:
        return "P0"
    outline = client_outline_rows(project)
    if client_outline_content_issues(outline, require_confirmed_state=False):
        return "P1"
    if client_outline_confirmation_errors(project):
        return "P2"
    _, artifacts = read_csv_rows(
        project / "AD-creative/orchestrator/artifact_index.csv"
    )
    truth_path = project / "AD-creative/orchestrator/current_truth.md"
    truth = truth_path.read_text(encoding="utf-8") if truth_path.is_file() else ""
    if not current_truth_value(truth, "current_pptx_artifact_id"):
        if not has_gate(project, "GATE-AUTO-CREATIVE-QUALITY-001"):
            return "P3"
        return "P4"
    required_package_gates = [
        "GATE-AUTO-CLIENT-OUTLINE-001",
        "GATE-AUTO-CLIENT-LANGUAGE-001",
        "GATE-AUTO-VISUAL-LAYOUT-001",
        "GATE-AUTO-VISUAL-QUALITY-001",
        "GATE-AUTO-PPT-001",
    ]
    if any(not has_gate(project, gate_id, {"PASS"}) for gate_id in required_package_gates):
        return "P5"
    binding_errors, _ = current_client_pack_binding_errors(project, artifacts)
    if binding_errors or not has_gate(project, "GATE-AUTO-CLIENT-PACK-001", {"PASS"}):
        return "P6"
    if not has_gate(
        project, "GATE-AUTO-CLIENT-SEND-READINESS-001", {"PASS"}
    ):
        return "P7"
    return "P8"


def gate_status_for_stages(project: Path, stages: tuple[str, ...]) -> str:
    statuses = [
        gate.get("status", "")
        for gate in latest_gate_rows(project)
        if gate.get("stage", "").strip() in stages and gate.get("status", "").strip()
    ]
    if not statuses:
        return "missing"
    if any(status == "BLOCKED" for status in statuses):
        return "blocked"
    if any(status == "PARTIAL_PASS" for status in statuses):
        return "partial"
    if any(status == "PASS" for status in statuses):
        return "pass"
    return statuses[-1].lower()


def goal_lane_states(project: Path) -> dict[str, str]:
    counts = read_counts(project)
    return {
        "text_framework": "confirmed" if not client_outline_confirmation_errors(project) else "active" if client_outline_rows(project) else "needs_material",
        "specialist_assets": "active" if counts["assets"] else "optional",
        "gates": gate_status_for_stages(project, ("visual_review", "film_quality", "client_review", "final_delivery")),
        "client_pack": "fresh" if has_gate(project, "GATE-AUTO-CLIENT-PACK-001", {"PASS"}) else "internal_only",
        "send_readiness": "ready_not_sent" if has_gate(project, "GATE-AUTO-CLIENT-SEND-READINESS-001", {"PASS"}) else "not_ready",
    }


def goal_completion_readiness(project: Path, errors: list[str], confirmations: list[dict[str, str]]) -> dict[str, object]:
    latest_gates = latest_gate_rows(project)
    blocking_gates = [
        gate.get("gate_id", "")
        for gate in latest_gates
        if gate.get("status", "").strip().upper() == "BLOCKED"
    ]
    required = [
        "GATE-AUTO-CLIENT-OUTLINE-001",
        "GATE-AUTO-CLIENT-LANGUAGE-001",
        "GATE-AUTO-VISUAL-QUALITY-001",
        "GATE-AUTO-VISUAL-LAYOUT-001",
        "GATE-AUTO-PPT-001",
        "GATE-AUTO-CLIENT-PACK-001",
    ]
    missing = [gate_id for gate_id in required if not has_gate(project, gate_id, {"PASS"})]
    ready = not errors and not confirmations and not blocking_gates and not missing
    return {
        "status": "READY_FOR_INDEPENDENT_REVIEW" if ready else "NOT_READY",
        "missing_gates": missing,
        "blocking_gates": blocking_gates,
        "send_readiness": latest_gate_status(
            project, "GATE-AUTO-CLIENT-SEND-READINESS-001"
        )
        or "NOT_RUN",
        "validation_errors": len(errors),
        "pending_confirmations": len(confirmations),
    }


def goal_stop_reason(payload: dict[str, object]) -> str:
    if payload["validation"] != "PASS":
        return "VALIDATION_CHECK"
    if payload["pending_confirmation_count"]:
        return "WAITING_FOR_CONFIRMATION"
    if payload["blocking_gap_count"]:
        return "BLOCKING_GAP"
    readiness = payload.get("completion_readiness", {})
    if isinstance(readiness, dict) and readiness.get("blocking_gates"):
        return "BLOCKED_GATE"
    if payload["counts"]["source_events"] == 0:
        return "NEEDS_MATERIAL"
    return ""


def item_title(row: dict[str, str], *keys: str, default: str = "-") -> str:
    for key in keys:
        value = row.get(key, "").strip()
        if value:
            return value
    return default


def render_dashboard(
    project: Path,
    *,
    validation_errors: list[str] | None = None,
    validation_status: str | None = None,
) -> Path:
    render_human_workspace_indexes(project)
    counts = read_counts(project)
    _, work_items = read_csv_rows(project / "AD-creative/orchestrator/work_items.csv")
    _, gaps = read_csv_rows(project / "AD-creative/orchestrator/gaps.csv")
    _, gates = read_csv_rows(project / "AD-creative/orchestrator/gate_log.csv")
    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    _, source_events = read_csv_rows(project / "AD-creative/orchestrator/source_events.csv")
    _, references = read_csv_rows(project / "AD-creative/references/reference_cards.csv")
    _, visual_assets = read_csv_rows(project / "AD-creative/visual_assets/asset_manifest.csv")
    confirmations = parse_confirmation_rows(project / "AD-creative/handoff/待你确认.md")
    goals = goal_iteration_rows(project)
    creative_runs = [
        row
        for row in artifacts
        if row.get("artifact_type", "").strip() in {"creative_production_run", "creative_production_import"}
    ]

    stage = project_stage(project)
    phase = derive_goal_phase(project)
    lane_states = goal_lane_states(project)
    if validation_errors is None:
        validation_errors, _ = validate(project)
    if validation_status is None:
        validation_status = "PASS" if not validation_errors else "CHECK"
    active_work = [row for row in work_items if row.get("status", "").lower() not in {"done", "closed"}]
    next_action = first_nonempty(active_work[0].get("title") if active_work else "", "读取资料并生成需求/缺口")
    payload = {
        "project": project.name,
        "stage": stage,
        "validation": validation_status,
        "updated": now_iso(),
        "counts": counts,
        "work": work_items,
        "gaps": gaps,
        "gates": gates,
        "artifacts": artifacts,
        "sources": source_events,
        "references": references,
        "visualAssets": visual_assets,
        "goals": goals,
        "creativeRuns": creative_runs,
        "goalPhase": {"id": phase, "name": GOAL_PHASE_NAMES.get(phase, phase)},
        "laneStates": lane_states,
        "completionReadiness": goal_completion_readiness(project, validation_errors, confirmations),
        "confirmations": confirmations,
        "validationErrors": validation_errors,
        "nextAction": next_action,
    }
    work_columns = [
        ("work_id", "工作"),
        ("status", "状态"),
        ("title", "标题"),
        ("owner_agent", "负责人"),
        ("gate_required", "关卡"),
        ("client_visibility", "可见性"),
    ]
    gap_columns = [
        ("gap_id", "缺口"),
        ("impact", "影响"),
        ("status", "状态"),
        ("recommended_action", "动作"),
    ]
    initial_row = active_work[0] if active_work else (work_items[0] if work_items else {})
    statuses = sorted({row.get("status", "") for row in work_items if row.get("status", "")})
    status_options = '<option value="all">全部状态</option>' + "".join(
        f'<option value="{cell(status)}">{cell(status)}</option>' for status in statuses
    )
    table_head = "<tr>" + "".join(f"<th>{cell(label)}</th>" for _, label in work_columns) + "</tr>"
    table_body = rows_to_table(work_items[:12], work_columns, "暂无工作项")
    gap_body = rows_to_table(gaps[:8], gap_columns, "暂无缺口")
    decision_html = (
        "\n".join(
            f"""<div class="decision"><strong>{cell(row['id'])} · {cell(row['question'])}</strong><div class="muted">{cell(row['recommendation'])}</div></div>"""
            for row in confirmations[:4]
        )
        or '<div class="decision"><strong>暂无待确认</strong><div class="muted">内部整理可继续。</div></div>'
    )
    selected_props = "\n".join(
        f'<div class="prop"><span>{cell(label)}</span><strong>{cell(value)}</strong></div>'
        for label, value in [
            ("ID", item_title(initial_row, "work_id", "artifact_id", "source_event_id", "gate_id")),
            ("状态", item_title(initial_row, "status", "gate_status", "trust_level")),
            ("阶段", item_title(initial_row, "stage", "affected_stage")),
            ("关卡", item_title(initial_row, "gate_required", "gate_status", "gate_id")),
            ("可见性", item_title(initial_row, "client_visibility", "visibility")),
            ("更新时间", item_title(initial_row, "updated_at", "created_at", default=payload["updated"])),
        ]
    )
    selected_objective = item_title(initial_row, "objective", "notes", "blocking_issues")

    html_doc = """<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>广告创意操作台</title>
<style>
:root {
  --bg: #f6f7f8;
  --ink: #202124;
  --muted: #6b7078;
  --soft: #8b9097;
  --line: #d8dde2;
  --panel: #ffffff;
  --panel-2: #f8fafb;
  --rail: #23262c;
  --rail-2: #30343b;
  --rail-muted: #aeb4bd;
  --blue: #285f89;
  --amber: #a96b18;
  --green: #237255;
  --red: #ad3b35;
}
* { box-sizing: border-box; }
body {
  margin: 0;
  background: var(--bg);
  color: var(--ink);
  font: 14px/1.45 -apple-system, BlinkMacSystemFont, "Segoe UI", "PingFang SC", sans-serif;
}
button, input, select { font: inherit; }
button { cursor: pointer; }
.app { min-height: 100vh; display: grid; grid-template-columns: 224px minmax(520px, 1fr) 382px; }
.sidebar { background: var(--rail); color: #f3f4f6; padding: 18px 14px; }
.brand { display: flex; align-items: center; gap: 10px; font-size: 13px; letter-spacing: .08em; text-transform: uppercase; color: #fff; margin-bottom: 22px; }
.mark { width: 18px; height: 18px; border: 1px solid rgba(255,255,255,.55); border-radius: 4px; display: grid; place-items: center; font-size: 10px; }
.nav { display: grid; gap: 4px; }
.nav button { width: 100%; display: flex; align-items: center; gap: 10px; border: 0; background: transparent; color: var(--rail-muted); padding: 8px 10px; border-radius: 6px; text-align: left; }
.nav button.active, .nav button:hover { color: #fff; background: rgba(255,255,255,.08); }
.nav .dot { width: 7px; height: 7px; border-radius: 999px; background: currentColor; opacity: .7; }
.main { min-width: 0; padding: 17px 22px 74px; }
.topbar { height: 42px; display: flex; align-items: center; justify-content: space-between; gap: 14px; border-bottom: 1px solid var(--line); margin-bottom: 14px; min-width: 0; }
.crumb { font-weight: 700; min-width: 0; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }
.tabs { display: flex; flex-shrink: 0; gap: 7px; color: var(--muted); font-size: 12px; }
.tabs button, .btn { border: 1px solid var(--line); border-radius: 6px; padding: 5px 8px; background: rgba(255,255,255,.65); color: var(--ink); }
.tabs button.active { color: #fff; background: var(--rail-2); border-color: var(--rail-2); }
	.summary { display: grid; grid-template-columns: repeat(6, minmax(98px, 1fr)); gap: 10px; margin-bottom: 14px; }
.metric { border: 1px solid var(--line); background: var(--panel); border-radius: 6px; padding: 9px 11px; min-width: 0; }
.metric small { display: block; color: var(--muted); font-size: 11px; }
.metric strong { display: block; font-size: 20px; line-height: 1.18; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }
.toolbar { display: flex; align-items: center; gap: 8px; margin-bottom: 10px; }
.search { flex: 1; min-width: 160px; border: 1px solid var(--line); border-radius: 6px; background: #fff; padding: 7px 9px; }
.select { border: 1px solid var(--line); border-radius: 6px; background: #fff; padding: 7px 8px; color: var(--ink); max-width: 150px; }
.section { border: 1px solid var(--line); background: var(--panel); border-radius: 6px; overflow: hidden; margin-bottom: 13px; }
.section-head { display: flex; align-items: center; justify-content: space-between; gap: 10px; border-bottom: 1px solid var(--line); padding: 9px 11px; }
.section h2 { margin: 0; font-size: 13px; }
.count { color: var(--muted); font-size: 12px; }
table { width: 100%; border-collapse: collapse; table-layout: fixed; }
th, td { border-bottom: 1px solid #ece9e2; padding: 8px 11px; text-align: left; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }
th { color: var(--muted); font-weight: 650; font-size: 11px; }
tr:last-child td { border-bottom: 0; }
tbody tr { outline: 0; }
tbody tr:hover, tbody tr.selected { background: #f3f7f8; }
.empty { color: var(--muted); padding: 18px 11px; }
.chip { display: inline-flex; align-items: center; min-height: 22px; padding: 2px 7px; border-radius: 999px; font-size: 11px; border: 1px solid var(--line); background: #fff; max-width: 100%; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }
.pass { color: var(--green); border-color: rgba(35,114,85,.35); }
.warn { color: var(--amber); border-color: rgba(169,107,24,.35); }
.check { color: var(--red); border-color: rgba(173,59,53,.35); }
.muted { color: var(--muted); }
.inspector { border-left: 1px solid var(--line); background: var(--panel-2); padding: 17px 16px 74px; min-width: 0; }
.inspector h1 { font-size: 19px; line-height: 1.25; margin: 0 0 5px; overflow-wrap: anywhere; }
.sub { color: var(--muted); margin-bottom: 12px; }
.props { display: grid; gap: 8px; margin: 15px 0; }
.prop { display: flex; justify-content: space-between; gap: 12px; border-bottom: 1px solid var(--line); padding-bottom: 8px; }
.prop span:first-child { color: var(--muted); }
.prop strong { text-align: right; min-width: 0; overflow: hidden; text-overflow: ellipsis; }
.detail-box { border: 1px solid var(--line); border-radius: 6px; background: #fff; padding: 10px 11px; margin-top: 12px; }
.detail-box h3 { margin: 0 0 8px; font-size: 12px; color: var(--muted); text-transform: uppercase; letter-spacing: .03em; }
.detail-box p { margin: 0; color: var(--ink); overflow-wrap: anywhere; }
.decision-list { display: grid; gap: 8px; }
.decision { border: 1px solid var(--line); border-radius: 6px; background: #fff; padding: 8px; }
.decision strong { display: block; margin-bottom: 3px; }
.rail { position: fixed; left: 224px; right: 0; bottom: 0; min-height: 50px; background: #fff; border-top: 1px solid var(--line); display: flex; align-items: center; justify-content: space-between; gap: 12px; padding: 8px 18px; }
.rail strong { color: var(--blue); }
.rail-actions { display: flex; gap: 8px; flex-shrink: 0; }
	@media (max-width: 1060px) {
	  .app { grid-template-columns: 196px minmax(430px, 1fr); }
	  .inspector { display: none; }
	  .rail { left: 196px; }
	  .summary { grid-template-columns: repeat(3, minmax(0, 1fr)); }
	}
@media (max-width: 760px) {
  .app { grid-template-columns: 1fr; }
  .sidebar { display: none; }
  .main { padding: 12px 12px 86px; }
  .topbar { align-items: flex-start; height: auto; flex-direction: column; padding-bottom: 10px; }
  .tabs { flex-wrap: wrap; }
	  .summary { grid-template-columns: repeat(2, minmax(0, 1fr)); }
	  .toolbar { flex-wrap: wrap; }
	  .search, .select { flex: 1 1 100%; max-width: none; }
	  .section { overflow: visible; }
	  table { min-width: 0; }
	  thead { display: none; }
	  tbody, tr, td { display: block; width: 100%; }
	  tbody tr { padding: 8px 0; border-bottom: 1px solid #ece9e2; }
	  tbody tr:last-child { border-bottom: 0; }
	  td { border-bottom: 0; padding: 5px 11px; white-space: normal; display: grid; grid-template-columns: 86px minmax(0, 1fr); gap: 10px; align-items: start; }
	  td::before { content: attr(data-label); color: var(--muted); font-size: 11px; font-weight: 650; }
	  td.empty { display: block; }
	  td.empty::before { content: ""; display: none; }
	  .rail { position: static; left: auto; right: auto; bottom: auto; align-items: flex-start; flex-direction: column; }
	}
</style>
</head>
<body>
<script>const DATA = __PAYLOAD__;</script>
<div class="app">
  <aside class="sidebar">
    <div class="brand"><span class="mark">AD</span><span>广告创意操作台</span></div>
    <nav class="nav">
	      <button class="active" data-nav="work"><span class="dot"></span>工作台</button>
	      <button data-nav="sources"><span class="dot"></span>资料</button>
	      <button data-nav="references"><span class="dot"></span>参考</button>
	      <button data-nav="visualAssets"><span class="dot"></span>图片</button>
	      <button data-nav="creativeRuns"><span class="dot"></span>创意运行</button>
	      <button data-nav="goals"><span class="dot"></span>Goal</button>
	      <button data-nav="artifacts"><span class="dot"></span>产物</button>
	      <button data-nav="gates"><span class="dot"></span>关卡</button>
	      <button data-nav="confirmations"><span class="dot"></span>待确认</button>
    </nav>
  </aside>
  <main class="main">
    <div class="topbar">
      <div class="crumb">项目 / __PROJECT__</div>
      <div class="tabs">
	        <button class="active" data-view="work">工作</button>
	        <button data-view="sources">资料</button>
	        <button data-view="references">参考</button>
	        <button data-view="visualAssets">图片</button>
	        <button data-view="creativeRuns">创意运行</button>
	        <button data-view="goals">Goal</button>
	        <button data-view="artifacts">产物</button>
	        <button data-view="gates">关卡</button>
	        <button data-view="confirmations">待确认</button>
      </div>
    </div>
    <div class="summary">
	      <div class="metric"><small>阶段</small><strong>__STAGE__</strong></div>
	      <div class="metric"><small>校验</small><strong>__VALIDATION__</strong></div>
	      <div class="metric"><small>工作项</small><strong>__WORK_COUNT__</strong></div>
	      <div class="metric"><small>参考</small><strong>__REFERENCE_COUNT__</strong></div>
	      <div class="metric"><small>图片</small><strong>__ASSET_COUNT__</strong></div>
	      <div class="metric"><small>Goal阶段</small><strong>__GOAL_PHASE__</strong></div>
	      <div class="metric"><small>Goal</small><strong>__GOAL_COUNT__</strong></div>
	      <div class="metric"><small>产物</small><strong>__ARTIFACT_COUNT__</strong></div>
    </div>
    <section class="section">
      <div class="section-head">
        <h2 id="contentTitle">工作</h2>
        <span class="count" id="rowCount">__ROW_COUNT__</span>
      </div>
      <div class="toolbar">
	        <input id="searchInput" class="search" placeholder="搜索工作、资料、参考、图片、产物">
        <select id="statusFilter" class="select">__STATUS_OPTIONS__</select>
        <select id="riskFilter" class="select">
          <option value="all">全部</option>
          <option value="blocked">阻塞</option>
          <option value="client">客户可见</option>
          <option value="open">进行中</option>
        </select>
      </div>
      <table>
        <thead id="tableHead">__TABLE_HEAD__</thead>
        <tbody id="tableBody">__TABLE_BODY__</tbody>
      </table>
    </section>
    <section class="section">
      <div class="section-head"><h2>缺口</h2><span class="count" id="gapCount">__GAP_COUNT__</span></div>
      <table>
        <thead><tr><th>缺口</th><th>影响</th><th>状态</th><th>动作</th></tr></thead>
        <tbody id="gapBody">__GAP_BODY__</tbody>
      </table>
    </section>
  </main>
  <aside class="inspector">
    <h1 id="selectedTitle">__NEXT_ACTION__</h1>
    <div class="sub">当前选中项</div>
    <span class="chip __VALIDATION_CLASS__">校验 __VALIDATION__</span>
    <span class="chip warn">内部预览</span>
    <div class="props" id="selectedProps">__SELECTED_PROPS__</div>
    <div class="detail-box">
      <h3>目标 / 说明</h3>
      <p id="selectedObjective">__SELECTED_OBJECTIVE__</p>
    </div>
    <div class="detail-box">
      <h3>待确认</h3>
      <div class="decision-list" id="decisionList">__DECISIONS__</div>
    </div>
  </aside>
</div>
<div class="rail">
  <div><strong>需要你确认</strong> · 客户稿发送、付费/登录平台、上传资料、全局安装前需要明确确认</div>
  <div class="rail-actions"><button class="btn" data-view="confirmations">待确认</button><button class="btn" onclick="location.reload()">刷新</button></div>
</div>
<script>
const VIEWS = {
  work: {
    title: "工作",
    rows: DATA.work,
    columns: [
      ["work_id", "工作"], ["status", "状态"], ["title", "标题"], ["owner_agent", "负责人"], ["gate_required", "关卡"], ["client_visibility", "可见性"]
    ],
    id: "work_id",
    titleKey: "title",
    objectiveKey: "objective"
  },
	  sources: {
	    title: "资料",
    rows: DATA.sources,
    columns: [
      ["source_event_id", "ID"], ["declared_semantics", "语义"], ["source_type", "类型"], ["file_paths", "路径"], ["trust_level", "可信度"]
    ],
    id: "source_event_id",
	    titleKey: "raw_summary",
	    objectiveKey: "notes"
	  },
	  references: {
	    title: "参考",
	    rows: DATA.references,
	    columns: [
	      ["reference_id", "ID"], ["platform", "平台"], ["title", "标题"], ["role", "角色"], ["client_visible", "客户可见"]
	    ],
	    id: "reference_id",
	    titleKey: "title",
	    objectiveKey: "why_relevant"
	  },
	  visualAssets: {
	    title: "图片",
	    rows: DATA.visualAssets,
	    columns: [
	      ["asset_id", "图片"], ["slot_id", "槽位"], ["status", "状态"], ["visibility", "可见性"], ["qa_status", "QA"], ["risk_level", "风险"], ["prompt_or_edit_ref", "Trace"]
	    ],
	    id: "asset_id",
	    titleKey: "path",
	    objectiveKey: "notes"
	  },
	  creativeRuns: {
	    title: "创意运行",
	    rows: DATA.creativeRuns,
	    columns: [
	      ["artifact_id", "运行"], ["artifact_type", "类型"], ["status", "状态"], ["path", "路径"], ["gate_status", "关卡"]
	    ],
	    id: "artifact_id",
	    titleKey: "path",
	    objectiveKey: "stage"
	  },
	  goals: {
	    title: "Goal",
	    rows: DATA.goals,
	    columns: [
	      ["goal_id", "Goal"], ["status", "状态"], ["title", "标题"], ["owner", "负责人"], ["next", "下一步"]
	    ],
	    id: "goal_id",
	    titleKey: "title",
	    objectiveKey: "path"
	  },
	  artifacts: {
	    title: "产物",
	    rows: DATA.artifacts,
	    columns: [
	      ["artifact_id", "产物"], ["artifact_type", "类型"], ["status", "状态"], ["visibility", "可见性"], ["gate_status", "关卡"]
	    ],
    id: "artifact_id",
    titleKey: "path",
    objectiveKey: "stage"
  },
  gates: {
    title: "关卡",
    rows: DATA.gates,
    columns: [
      ["gate_id", "关卡"], ["stage", "阶段"], ["status", "状态"], ["score", "分数"], ["next_state", "下一步"]
    ],
    id: "gate_id",
    titleKey: "stage",
    objectiveKey: "blocking_issues"
  },
  confirmations: {
    title: "待确认",
    rows: DATA.confirmations,
    columns: [
      ["id", "ID"], ["question", "问题"], ["recommendation", "建议"], ["impact", "影响"], ["actions", "动作"]
    ],
    id: "id",
    titleKey: "question",
    objectiveKey: "recommendation"
  }
};
let currentView = "work";
let selectedKey = "";
const $ = (id) => document.getElementById(id);
function value(row, key) { return (row && row[key]) ? String(row[key]) : ""; }
function escapeHtml(text) {
  return String(text ?? "").replace(/[&<>"']/g, (char) => ({ "&": "&amp;", "<": "&lt;", ">": "&gt;", '"': "&quot;", "'": "&#39;" }[char]));
}
function statusClass(text) {
  const lower = String(text || "").toLowerCase();
  if (lower.includes("pass") || lower.includes("done") || lower.includes("resolved")) return "pass";
  if (lower.includes("block") || lower.includes("error") || lower.includes("reject")) return "check";
  if (lower.includes("partial") || lower.includes("pending") || lower.includes("wait") || lower.includes("draft")) return "warn";
  return "";
}
function searchable(row) { return Object.values(row).join(" ").toLowerCase(); }
function rowMatchesRisk(row, risk) {
  const text = searchable(row);
  if (risk === "all") return true;
  if (risk === "blocked") return text.includes("block") || text.includes("wait") || text.includes("needs_user");
  if (risk === "client") return text.includes("client_visible") || text.includes("客户");
  if (risk === "open") return !text.includes("done") && !text.includes("closed") && !text.includes("pass,");
  return true;
}
function rowStatus(row) {
  return value(row, "status") || value(row, "gate_status") || value(row, "trust_level") || value(row, "impact") || "";
}
function syncStatusOptions() {
  const view = VIEWS[currentView];
  const statuses = Array.from(new Set(view.rows.map(rowStatus).filter(Boolean))).sort();
  $("statusFilter").innerHTML = `<option value="all">全部状态</option>` + statuses.map((status) => `<option value="${escapeHtml(status)}">${escapeHtml(status)}</option>`).join("");
}
function filteredRows() {
  const view = VIEWS[currentView];
  const query = $("searchInput").value.trim().toLowerCase();
  const status = $("statusFilter").value;
  const risk = $("riskFilter").value;
  return view.rows.filter((row) => {
    const matchesQuery = !query || searchable(row).includes(query);
    const matchesStatus = status === "all" || rowStatus(row) === status;
    return matchesQuery && matchesStatus && rowMatchesRisk(row, risk);
  });
}
function renderTable() {
  const view = VIEWS[currentView];
  const rows = filteredRows();
  $("contentTitle").textContent = view.title;
  $("rowCount").textContent = `${rows.length} 条`;
  $("tableHead").innerHTML = `<tr>${view.columns.map(([, label]) => `<th>${escapeHtml(label)}</th>`).join("")}</tr>`;
  if (!rows.length) {
    $("tableBody").innerHTML = `<tr><td class="empty" colspan="${view.columns.length}">暂无记录</td></tr>`;
    renderInspector(null);
    return;
  }
  if (!selectedKey || !rows.some((row) => value(row, view.id) === selectedKey)) selectedKey = value(rows[0], view.id);
  $("tableBody").innerHTML = rows.map((row) => {
    const key = value(row, view.id);
    const selected = key === selectedKey ? ' class="selected"' : "";
    const cells = view.columns.map(([field]) => {
      const raw = value(row, field);
      const chip = field.includes("status") || field.includes("visibility") || field === "trust_level" || field === "impact";
      const content = chip && raw ? `<span class="chip ${statusClass(raw)}">${escapeHtml(raw)}</span>` : escapeHtml(raw || "-");
      const label = view.columns.find(([name]) => name === field)?.[1] || field;
      return `<td data-label="${escapeHtml(label)}" title="${escapeHtml(raw)}">${content}</td>`;
    }).join("");
    return `<tr data-key="${escapeHtml(key)}"${selected}>${cells}</tr>`;
  }).join("");
  document.querySelectorAll("#tableBody tr[data-key]").forEach((rowEl) => {
    rowEl.addEventListener("click", () => {
      selectedKey = rowEl.dataset.key;
      renderTable();
    });
  });
  renderInspector(rows.find((row) => value(row, view.id) === selectedKey) || rows[0]);
}
function renderGaps() {
  const rows = DATA.gaps.slice(0, 8);
  $("gapCount").textContent = `${DATA.gaps.length} 条`;
  $("gapBody").innerHTML = rows.length ? rows.map((row) => `
    <tr>
      <td data-label="缺口" title="${escapeHtml(value(row, "description"))}">${escapeHtml(value(row, "gap_id") || "-")}</td>
      <td data-label="影响"><span class="chip ${statusClass(value(row, "impact"))}">${escapeHtml(value(row, "impact") || "-")}</span></td>
      <td data-label="状态">${escapeHtml(value(row, "status") || "-")}</td>
      <td data-label="动作" title="${escapeHtml(value(row, "recommended_action"))}">${escapeHtml(value(row, "recommended_action") || "-")}</td>
    </tr>`).join("") : `<tr><td class="empty" colspan="4">暂无缺口</td></tr>`;
}
function renderInspector(row) {
  const view = VIEWS[currentView];
  if (!row) {
    $("selectedTitle").textContent = DATA.nextAction;
    $("selectedObjective").textContent = "-";
    $("selectedProps").innerHTML = "";
    return;
  }
  const title = value(row, view.titleKey) || value(row, view.id) || DATA.nextAction;
  $("selectedTitle").textContent = title;
  $("selectedObjective").textContent = value(row, view.objectiveKey) || value(row, "notes") || value(row, "blocking_issues") || "-";
  const props = [
    ["ID", value(row, view.id)],
    ["状态", rowStatus(row) || "-"],
    ["阶段", value(row, "stage") || value(row, "affected_stage") || "-"],
    ["关卡", value(row, "gate_required") || value(row, "gate_status") || value(row, "gate_id") || "-"],
    ["可见性", value(row, "client_visibility") || value(row, "visibility") || "-"],
    ["更新时间", value(row, "updated_at") || value(row, "created_at") || DATA.updated]
  ];
  $("selectedProps").innerHTML = props.map(([label, val]) => `<div class="prop"><span>${escapeHtml(label)}</span><strong title="${escapeHtml(val)}">${escapeHtml(val)}</strong></div>`).join("");
}
function renderDecisions() {
  const rows = DATA.confirmations;
  $("decisionList").innerHTML = rows.length ? rows.slice(0, 4).map((row) => `
    <div class="decision">
      <strong>${escapeHtml(value(row, "id"))} · ${escapeHtml(value(row, "question"))}</strong>
      <div class="muted">${escapeHtml(value(row, "recommendation"))}</div>
    </div>`).join("") : `<div class="decision"><strong>暂无待确认</strong><div class="muted">内部整理可继续。</div></div>`;
}
function setView(nextView) {
  currentView = nextView;
  selectedKey = "";
  document.querySelectorAll("[data-view]").forEach((button) => button.classList.toggle("active", button.dataset.view === nextView));
  document.querySelectorAll("[data-nav]").forEach((button) => button.classList.toggle("active", button.dataset.nav === nextView));
  syncStatusOptions();
  renderTable();
}
document.querySelectorAll("[data-view]").forEach((button) => button.addEventListener("click", () => setView(button.dataset.view)));
document.querySelectorAll("[data-nav]").forEach((button) => button.addEventListener("click", () => setView(button.dataset.nav)));
$("searchInput").addEventListener("input", renderTable);
$("statusFilter").addEventListener("change", renderTable);
$("riskFilter").addEventListener("change", renderTable);
renderDecisions();
renderGaps();
setView("work");
</script>
</body>
</html>
"""
    replacements = {
        "__PAYLOAD__": html_json(payload),
        "__PROJECT__": cell(project.name),
        "__STAGE__": cell(stage),
	        "__VALIDATION__": validation_status,
	        "__WORK_COUNT__": str(counts["work_items"]),
	        "__REFERENCE_COUNT__": str(counts["references"]),
	        "__ASSET_COUNT__": str(counts["assets"]),
	        "__GOAL_PHASE__": cell(phase),
	        "__GOAL_COUNT__": str(len(goals)),
	        "__ARTIFACT_COUNT__": str(counts["artifacts"]),
        "__VALIDATION_CLASS__": "pass" if validation_status == "PASS" else "check",
        "__NEXT_ACTION__": cell(next_action),
        "__STATUS_OPTIONS__": status_options,
        "__ROW_COUNT__": f"{len(work_items)} 条",
        "__TABLE_HEAD__": table_head,
        "__TABLE_BODY__": table_body,
        "__GAP_COUNT__": f"{len(gaps)} 条",
        "__GAP_BODY__": gap_body,
        "__SELECTED_PROPS__": selected_props,
        "__SELECTED_OBJECTIVE__": cell(selected_objective),
        "__DECISIONS__": decision_html,
    }
    for placeholder, value in replacements.items():
        html_doc = html_doc.replace(placeholder, value)
    dashboard_path = project / DASHBOARD_REL
    write_text(dashboard_path, html_doc)
    return dashboard_path


def audit_dashboard(project: Path) -> list[str]:
    path = project / DASHBOARD_REL
    if not path.exists():
        return ["操作台 HTML 尚未生成。"]
    text = path.read_text(encoding="utf-8")
    issues: list[str] = []
    required_markers = {
        "data-view=\"work\"": "缺少工作 Tab。",
        "data-view=\"sources\"": "缺少资料 Tab。",
        "data-view=\"references\"": "缺少参考 Tab。",
        "data-view=\"visualAssets\"": "缺少图片 Tab。",
        "data-view=\"creativeRuns\"": "缺少创意运行 Tab。",
        "data-view=\"goals\"": "缺少 Goal Tab。",
        "data-view=\"artifacts\"": "缺少产物 Tab。",
        "data-view=\"gates\"": "缺少关卡 Tab。",
        "data-view=\"confirmations\"": "缺少待确认 Tab。",
        "id=\"searchInput\"": "缺少搜索框。",
        "id=\"statusFilter\"": "缺少状态筛选。",
        "id=\"riskFilter\"": "缺少风险筛选。",
        "当前选中项": "缺少右侧检查器。",
        "<h3>待确认</h3>": "缺少待确认区。",
        "@media (max-width: 760px)": "缺少移动端降级。",
        "data-label=": "缺少移动端卡片式表格标签。",
        "const DATA =": "缺少嵌入数据。",
    }
    for marker, message in required_markers.items():
        if marker not in text:
            issues.append(message)
    unresolved = [
        marker
        for marker in [
            "__PAYLOAD__",
            "__PROJECT__",
            "__TABLE_BODY__",
            "__SELECTED_PROPS__",
            "__DECISIONS__",
        ]
        if marker in text
    ]
    if unresolved:
        issues.append(f"操作台存在未替换占位符：{', '.join(unresolved)}")
    if "<tbody id=\"tableBody\"></tbody>" in text:
        issues.append("首屏 Work 表格没有静态行。")
    if "广告创意操作台" not in text or "需要你确认" not in text:
        issues.append("非开发者关键入口文案缺失。")
    return issues


@dataclass
class CouncilResult:
    name: str
    status: str
    evidence: list[str]
    issues: list[str]


def council_status(results: list[CouncilResult]) -> str:
    if any(result.status == "BLOCKED" for result in results):
        return "BLOCKED"
    if any(result.status == "PARTIAL" for result in results):
        return "PARTIAL_PASS"
    return "PASS"


def run_council(project: Path) -> tuple[str, list[CouncilResult], Path]:
    errors, stats = validate(project)
    counts = read_counts(project)
    dashboard_path = project / DASHBOARD_REL
    board_path = project / "AD-creative/handoff/项目看板.md"
    confirm_path = project / "AD-creative/handoff/待你确认.md"
    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    _, visual_assets = read_csv_rows(project / "AD-creative/visual_assets/asset_manifest.csv")

    strategy_issues = []
    if counts["source_events"] == 0:
        strategy_issues.append("未登记客户资料或项目输入。")
    if counts["work_items"] == 0:
        strategy_issues.append("没有可执行工作项。")
    strategy = CouncilResult(
        "Strategy Council",
        "PASS" if not strategy_issues else "PARTIAL",
        [
            f"source_events={counts['source_events']}",
            f"work_items={counts['work_items']}",
            "下一步聚焦需求/缺口/客户追问。",
        ],
        strategy_issues,
    )

    ops_issues = list(errors)
    for artifact in artifacts:
        visibility = artifact.get("visibility", "").lower()
        gate_status = artifact.get("gate_status", "").lower()
        is_client_visible = visibility in CLIENT_VISIBLE_VALUES
        has_pass_gate = gate_status in PASS_GATE_VALUES
        if is_client_visible and not has_pass_gate:
            ops_issues.append(
                f"client-visible artifact without pass gate: {artifact.get('artifact_id')}"
            )
    for asset in visual_assets:
        visibility = asset.get("visibility", "").lower()
        qa_status = asset.get("qa_status", "").upper()
        if visibility in CLIENT_VISIBLE_VALUES and qa_status != "PASS":
            ops_issues.append(f"client-visible asset without QA PASS: {asset.get('asset_id')}")
    operations = CouncilResult(
        "Operations Council",
        "PASS" if not ops_issues else "BLOCKED",
        [
            f"validation_errors={len(errors)}",
            f"artifacts={counts['artifacts']}",
            f"references={counts['references']}",
            f"assets={counts['assets']}",
            "客户可见产物必须有 Gate。",
            "客户可见图片必须 QA PASS。",
        ],
        ops_issues,
    )

    craft_issues = audit_dashboard(project)
    if not board_path.exists() or len(board_path.read_text(encoding="utf-8")) < 180:
        craft_issues.append("项目看板内容不足。")
    confirm_text = confirm_path.read_text(encoding="utf-8") if confirm_path.exists() else ""
    has_clear_empty_state = "当前无必须立即确认" in confirm_text
    has_decision_rows = bool(re.search(r"\|\s*(Q|DEC)-", confirm_text))
    if not confirm_path.exists() or not (has_clear_empty_state or has_decision_rows):
        craft_issues.append("待确认信息不够清晰。")
    craft = CouncilResult(
        "Craft Council",
        "PASS" if not craft_issues else "PARTIAL",
        [
            "操作台采用可搜索/筛选列表、右侧检查器、底部决策栏。",
            "操作台显示工作、资料、参考、图片、产物、关卡、待确认。",
            "操作台有静态首屏和浏览器交互增强。",
            "非开发者优先看 handoff，不直接编辑 CSV。",
        ],
        craft_issues,
    )

    results = [strategy, operations, craft]
    overall = council_status(results)
    report_path = project / COUNCIL_REPORT_REL
    report_sections = []
    for result in results:
        evidence = "\n".join(f"- {item}" for item in result.evidence)
        issues = "\n".join(f"- {item}" for item in result.issues) or "- 无"
        report_sections.append(
            f"""## {result.name}
Status: {result.status}

Evidence:
{evidence}

Issues:
{issues}
"""
        )

    write_text(
        report_path,
        f"""# Three Council Readiness Report

Generated: {now_iso()}
Overall: {overall}

{''.join(report_sections)}
## Self-Approval Boundary
Allowed after PASS:
- 初始化模板
- 登记资料
- 内部需求/缺口整理
- 公开官方来源搜索计划
- 内部视觉方向草图规划
- 只读操作台刷新
- 内部 Gate 初检

Still requires explicit confirmation:
- 发送客户稿
- 付费、登录、私密账号、KYC、钱包或凭据
- 上传客户资料到外部平台
- 全局安装 Skill
- 覆盖或删除旧版本
- 将 AI 图标记为客户可见
""",
    )

    append_gate(
        project,
        "GATE-THREE-COUNCIL-READINESS",
        "project_readiness",
        overall,
        "3/3" if overall == "PASS" else "",
        "",
        ";".join(issue for result in results for issue in result.issues),
        "",
        "",
        "ready_for_non_developer_operation" if overall == "PASS" else "revise_readiness",
        "three_council",
        target_ref=safe_rel(project, report_path),
        target_sha256=file_sha256(report_path),
    )
    return overall, results, report_path


def status_payload(
    project: Path,
    *,
    validation_errors: list[str] | None = None,
) -> dict[str, object]:
    project = project.resolve()
    counts = read_counts(project)
    if validation_errors is None:
        validation_issues, validation_stats = validate_issues(project)
        errors = [
            issue.message for issue in validation_issues if issue.scope != "legacy"
        ]
    else:
        validation_issues = []
        validation_stats = {"legacy_debt": 0}
        errors = list(validation_errors)
    dashboard = project / DASHBOARD_REL
    report = project / COUNCIL_REPORT_REL
    surface = project_surface(project)
    content_surface = surface == CONTENT_SURFACE
    _, work_items = read_csv_rows(project / "AD-creative/orchestrator/work_items.csv")
    _, gaps = read_csv_rows(project / "AD-creative/orchestrator/gaps.csv")
    confirmations = parse_confirmation_rows(project / "AD-creative/handoff/待你确认.md")
    active_work = [
        row
        for row in work_items
        if row.get("status", "").strip().lower() not in {"done", "closed", "resolved"}
    ]
    open_gaps = [
        row
        for row in gaps
        if row.get("status", "").strip().lower() not in {"resolved", "closed", "done"}
    ]
    blocking_gaps = [
        row
        for row in open_gaps
        if row.get("impact", "").strip().lower() in {"blocking", "high_impact"}
    ]
    if content_surface:
        known_gap_ids = {row.get("gap_id", "") for row in open_gaps}
        blocking_gap_ids = {row.get("gap_id", "") for row in blocking_gaps}
        confirmations = [
            row
            for row in confirmations
            if row.get("id", "") not in known_gap_ids
            or row.get("id", "") in blocking_gap_ids
        ]
    current_truth_path = project / "AD-creative/orchestrator/current_truth.md"
    recorded_next_action = (
        first_section_line(current_truth_path.read_text(encoding="utf-8"), "Next Action")
        if current_truth_path.is_file()
        else ""
    )
    if errors:
        next_status = "VALIDATION_CHECK"
        next_action = "Fix validation errors before continuing."
    elif content_surface and blocking_gaps:
        next_status = "BLOCKING_CONTENT_GAP"
        next_action = first_nonempty(
            blocking_gaps[0].get("recommended_action"),
            blocking_gaps[0].get("description"),
            default="Resolve the minimum content blocker.",
        )
    elif confirmations:
        next_status = "WAITING_FOR_CONFIRMATION"
        next_action = first_nonempty(confirmations[0].get("question"), default="Resolve pending confirmation.")
    elif content_surface and counts["source_events"] == 0:
        next_status = "NEEDS_MATERIAL"
        next_action = "Run adco run <project> --material <brief_file_or_folder>."
    elif content_surface:
        next_status = "READY_FOR_CONTENT_WORK"
        next_action = recorded_next_action or "Use the current evidence to produce the requested internal advertising outcome."
    elif active_work:
        next_status = "ACTIVE_WORK"
        next_action = first_nonempty(active_work[0].get("title"), active_work[0].get("objective"), default="Continue active work item.")
    elif blocking_gaps:
        next_status = "BLOCKING_GAP"
        next_action = first_nonempty(
            blocking_gaps[0].get("recommended_action"),
            blocking_gaps[0].get("description"),
            default="Resolve blocking gap.",
        )
    elif open_gaps:
        next_status = "OPEN_GAP"
        next_action = first_nonempty(
            open_gaps[0].get("recommended_action"),
            open_gaps[0].get("description"),
            default="Review open gap.",
        )
    elif counts["source_events"] == 0:
        next_status = "NEEDS_MATERIAL"
        next_action = "Run adco run <project> --material <brief_file_or_folder>."
    else:
        next_status = "READY_FOR_NEXT_GATE"
        next_action = f"Run adco next {project} or an explicit stage Gate."
    goal = latest_goal_row(project) if not content_surface else None
    phase = "CONTENT" if content_surface else derive_goal_phase(project)
    next_command = ""
    if next_status == "NEEDS_MATERIAL":
        next_command = f"adco run {project} --material <brief_file_or_folder>"
    elif not content_surface and next_status == "READY_FOR_NEXT_GATE":
        next_command = f"adco goal-run {project} --goal-id latest --max-steps 1"
    elif not content_surface and next_status == "ACTIVE_WORK":
        next_command = f"adco goal-run {project} --goal-id latest --max-steps 1"

    payload: dict[str, object] = {
        "project": str(project),
        "surface": surface,
        "stage": project_stage(project),
        "phase": phase,
        "phase_name": GOAL_PHASE_NAMES.get(phase, phase),
        "goal_id": goal.get("goal_id") if goal else "",
        "goal": goal or {},
        "lane_states": {} if content_surface else goal_lane_states(project),
        "validation": "PASS" if not errors else "CHECK",
        "validation_issues": [issue.as_dict() for issue in validation_issues],
        "legacy_debt": validation_stats.get("legacy_debt", 0),
        "counts": counts,
        "active_work_count": len(active_work),
        "open_gap_count": len(open_gaps),
        "blocking_gap_count": len(blocking_gaps),
        "pending_confirmation_count": len(confirmations),
        "next_status": next_status,
        "next_action": next_action,
        "active_work": active_work[:5],
        "open_gaps": open_gaps[:5],
        "blocking_gaps": blocking_gaps[:5],
        "pending_confirmations": confirmations[:5],
        "dashboard": str(dashboard) if dashboard.exists() else None,
        "council_report": str(report) if report.exists() else None,
        "errors": errors,
        "next_command": next_command,
    }
    if content_surface:
        payload["completion_readiness"] = {
            "status": (
                "READY_FOR_CONTENT_WORK"
                if not errors and not blocking_gaps and not confirmations and counts["source_events"]
                else "NEEDS_INPUT"
            ),
            "delivery_gates_required": False,
            "validation_errors": len(errors),
            "pending_confirmations": len(confirmations),
        }
    else:
        payload["completion_readiness"] = goal_completion_readiness(project, errors, confirmations)
    payload["stop_reason"] = goal_stop_reason(payload)
    return payload


def print_status(project: Path) -> None:
    payload = status_payload(project)
    counts = payload["counts"]
    print(f"PROJECT={payload['project']}")
    print(f"PROJECT_SURFACE={payload['surface']}")
    print(f"STAGE={payload['stage']}")
    if payload["surface"] == DELIVERY_SURFACE:
        print(f"PHASE={payload['phase']}")
        print(f"GOAL_ID={payload['goal_id']}")
    print(f"VALIDATION={payload['validation']}")
    print(f"SOURCE_EVENTS={counts['source_events']}")
    print(f"REQUIREMENTS={counts['requirements']}")
    print(f"GAPS={counts['gaps']}")
    if payload["surface"] == DELIVERY_SURFACE:
        print(f"WORK_ITEMS={counts['work_items']}")
        print(f"ACTIVE_WORK={payload['active_work_count']}")
    print(f"OPEN_GAPS={payload['open_gap_count']}")
    print(f"BLOCKING_GAPS={payload['blocking_gap_count']}")
    print(f"PENDING_CONFIRMATIONS={payload['pending_confirmation_count']}")
    if payload["surface"] == DELIVERY_SURFACE:
        print(f"REFERENCES={counts['references']}")
        print(f"ASSETS={counts['assets']}")
        print(f"ARTIFACTS={counts['artifacts']}")
        print(f"GATES={counts['gates']}")
    print(f"NEXT_STATUS={payload['next_status']}")
    print(f"NEXT_ACTION={payload['next_action']}")
    print(f"NEXT_COMMAND={payload['next_command']}")
    print(f"STOP_REASON={payload['stop_reason'] or 'NONE'}")
    print(f"DASHBOARD={payload['dashboard'] or 'NOT_RUN'}")
    if payload["surface"] == DELIVERY_SURFACE:
        print(f"COUNCIL_REPORT={payload['council_report'] or 'NOT_RUN'}")
    errors = payload["errors"]
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")


def all_source_event_ids(project: Path) -> list[str]:
    _, rows = read_csv_rows(project / "AD-creative/orchestrator/source_events.csv")
    return [row.get("source_event_id", "") for row in rows if row.get("source_event_id", "")]


def goal_run_step(project: Path, *, allow_generate: bool) -> tuple[str, str, str]:
    payload = status_payload(project, validation_errors=[])
    counts = payload["counts"]
    if isinstance(counts, dict) and counts.get("source_events", 0) == 0:
        return "stop", "NEEDS_MATERIAL", str(payload.get("next_action", ""))
    if payload["pending_confirmation_count"]:
        return "stop", "WAITING_FOR_CONFIRMATION", str(payload.get("next_action", ""))
    if payload["blocking_gap_count"]:
        return "stop", "BLOCKING_GAP", str(payload.get("next_action", ""))

    evidence_path = project / "AD-creative/orchestrator/evidence_chunks.jsonl"
    if (
        isinstance(counts, dict)
        and counts.get("source_events", 0)
        and (counts.get("requirements", 0) == 0 or not evidence_path.is_file())
    ):
        source_ids = all_source_event_ids(project)
        perform_intake(project, source_ids, "goal-run 自动执行本地 intake。")
        render_handoff(project, "goal-run 自动执行本地 intake。", source_ids)
        return "intake", "PASS", "Extracted requirements and gaps from registered materials."

    del allow_generate
    return (
        "stop",
        "SAFE_DECISION_POINT",
        "Local intake is current. Use adco next to choose the next explicit Gate or creative brief action.",
    )


def run_goal(
    project: Path,
    *,
    goal_id: str,
    max_steps: int,
    allow_generate: bool,
) -> dict[str, object]:
    ensure_project(project)
    goal = resolve_goal_row(project, goal_id)
    if not goal:
        created_goal_id = default_goal_id()
        plan = render_goal_iteration_plan(
            project,
            goal_id=created_goal_id,
            title="Auto-created goal-run",
            objective="Run local deterministic ADCO goal steps until a safe stop condition.",
            owner="Main Controller",
        )
        goal = resolve_goal_row(project, created_goal_id) or {"goal_id": created_goal_id, "path": safe_rel(project, plan)}

    append_event(
        project,
        {
            "event_id": f"EVT-GOAL-RUN-{safe_artifact_suffix(goal.get('goal_id', 'GOAL'))}-{datetime.now().strftime('%H%M%S')}",
            "event_type": "goal_run_started",
            "created_at": now_iso(),
            "goal_id": goal.get("goal_id", ""),
            "max_steps": max_steps,
            "allow_generate": allow_generate,
        },
    )
    steps: list[dict[str, str]] = []
    stop_reason = ""
    for index in range(1, max(1, max_steps) + 1):
        action, status, detail = goal_run_step(project, allow_generate=allow_generate)
        steps.append(
            {
                "step": str(index),
                "action": action,
                "status": status,
                "detail": detail,
            }
        )
        append_event(
            project,
            {
                "event_id": f"EVT-GOAL-RUN-STEP-{index}-{datetime.now().strftime('%H%M%S')}",
                "event_type": "goal_run_step",
                "created_at": now_iso(),
                "goal_id": goal.get("goal_id", ""),
                "action": action,
                "status": status,
                "detail": detail,
            },
        )
        if action == "stop" or status == "BLOCKED":
            stop_reason = status
            break
    else:
        stop_reason = "MAX_STEPS_REACHED"

    intake_changed = any(step["action"] == "intake" for step in steps)
    incremental = run_incremental_validation(
        project,
        changed_artifact_ids=(
            [
                "ART-AUTO-EVIDENCE-CHUNKS",
                "ART-AUTO-FACT-INVENTORY",
                "ART-AUTO-REQUIREMENTS",
                "ART-AUTO-GAPS",
            ]
            if intake_changed
            else []
        ),
        changed_file_paths=(
            [
                "AD-creative/orchestrator/evidence_chunks.jsonl",
                "AD-creative/orchestrator/fact_inventory.jsonl",
                "AD-creative/orchestrator/requirements.csv",
                "AD-creative/orchestrator/gaps.csv",
            ]
            if intake_changed
            else []
        ),
    )
    dashboard = render_dashboard(
        project,
        validation_errors=incremental.errors,
        validation_status=(
            "SCOPED_PASS" if not incremental.errors else "SCOPED_CHECK"
        ),
    )
    payload = status_payload(project, validation_errors=incremental.errors)
    if not stop_reason:
        stop_reason = str(payload.get("stop_reason") or "READY_FOR_NEXT_STEP")
    append_event(
        project,
        {
            "event_id": f"EVT-GOAL-RUN-STOP-{datetime.now().strftime('%H%M%S')}",
            "event_type": "goal_run_stopped",
            "created_at": now_iso(),
            "goal_id": goal.get("goal_id", ""),
            "stop_reason": stop_reason,
            "steps": steps,
            "dashboard": safe_rel(project, dashboard),
        },
    )
    return {
        "project": str(project),
        "goal_id": goal.get("goal_id", ""),
        "stop_reason": stop_reason,
        "steps": steps,
        "dashboard": str(dashboard),
        "status": payload,
        "incremental_validation": incremental.as_dict(),
        "full_validation": "NOT_RUN",
        "council_run_count": 0,
    }


def inspect_pptx(path: Path) -> dict[str, int | bool | str]:
    if not path.exists():
        raise FileNotFoundError(f"pptx not found: {path}")
    slide_files: list[str] = []
    text_runs = 0
    slide_text_runs: list[int] = []
    image_refs = 0
    has_presentation = False
    with zipfile.ZipFile(path) as archive:
        names = archive.namelist()
        name_set = set(names)
        has_presentation = "ppt/presentation.xml" in name_set
        slide_files = sorted(
            name
            for name in names
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
        media_files = [name for name in names if name.startswith("ppt/media/")]
        image_refs = len(media_files)
        for slide_name in slide_files:
            root = ET.fromstring(archive.read(slide_name))
            slide_runs = sum(
                1
                for text_node in root.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}t")
                if text_node.text and text_node.text.strip()
            )
            slide_text_runs.append(slide_runs)
            text_runs += slide_runs
    flattened_slides = [
        index for index, count in enumerate(slide_text_runs, start=1) if count == 0
    ]
    editable = bool(
        has_presentation
        and slide_files
        and text_runs
        and not flattened_slides
    )
    return {
        "slides": len(slide_files),
        "editable_text_runs": text_runs,
        "editable_slides": len(slide_files) - len(flattened_slides),
        "flattened_slides": ",".join(str(item) for item in flattened_slides),
        "embedded_media": image_refs,
        "has_presentation": has_presentation,
        "editable": editable,
    }


def pptx_text_content(path: Path) -> str:
    if not path.exists():
        return ""
    chunks: list[str] = []
    try:
        with zipfile.ZipFile(path) as archive:
            for slide_name in sorted(
                name
                for name in archive.namelist()
                if name.startswith("ppt/slides/slide") and name.endswith(".xml")
            ):
                root = ET.fromstring(archive.read(slide_name))
                for text_node in root.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}t"):
                    if text_node.text:
                        chunks.append(text_node.text)
    except Exception:
        return ""
    return "\n".join(chunks)


def pptx_layout_findings(path: Path) -> list[str]:
    findings: list[str] = []
    if not path.exists():
        return [f"PPTX 文件不存在: {path}"]
    try:
        from PIL import Image
    except Exception:
        return findings
    try:
        with zipfile.ZipFile(path) as archive:
            names = set(archive.namelist())
            media_sizes: dict[str, tuple[int, int]] = {}
            for name in names:
                if not name.startswith("ppt/media/"):
                    continue
                try:
                    with archive.open(name) as handle:
                        with Image.open(handle) as image:
                            media_sizes[name] = (image.width, image.height)
                except Exception:
                    continue
            for slide_name in sorted(
                name for name in names if name.startswith("ppt/slides/slide") and name.endswith(".xml")
            ):
                rels_name = slide_name.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
                rels: dict[str, str] = {}
                if rels_name in names:
                    rel_root = ET.fromstring(archive.read(rels_name))
                    for rel in rel_root:
                        rel_id = rel.attrib.get("Id", "")
                        target = rel.attrib.get("Target", "")
                        if rel_id and target.startswith("../media/"):
                            rels[rel_id] = "ppt/media/" + target.removeprefix("../media/")
                slide_root = ET.fromstring(archive.read(slide_name))
                for pic in slide_root.iter("{http://schemas.openxmlformats.org/presentationml/2006/main}pic"):
                    embed = ""
                    for blip in pic.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}blip"):
                        embed = blip.attrib.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed", "")
                    ext = next(pic.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}ext"), None)
                    media_name = rels.get(embed, "")
                    if not media_name or media_name not in media_sizes or ext is None:
                        continue
                    try:
                        box_w = int(ext.attrib.get("cx", "0"))
                        box_h = int(ext.attrib.get("cy", "0"))
                    except ValueError:
                        continue
                    img_w, img_h = media_sizes[media_name]
                    if not box_w or not box_h or not img_w or not img_h:
                        continue
                    box_ratio = box_w / box_h
                    img_ratio = img_w / img_h
                    if abs(box_ratio - img_ratio) / img_ratio > 0.18:
                        findings.append(
                            f"{Path(slide_name).stem} 图片可能被拉伸: {Path(media_name).name} image_ratio={img_ratio:.2f} box_ratio={box_ratio:.2f}"
                        )
    except Exception as exc:  # noqa: BLE001 - layout gate should report parse failures as findings
        findings.append(f"PPTX layout parse failed: {exc}")
    return findings




def find_client_language_hits(text: str) -> list[str]:
    hits = {match.group(0) for match in CLIENT_LANGUAGE_ASCII_PATTERN.finditer(text)}
    hits.update(match.group(0) for match in CLIENT_LANGUAGE_CJK_PATTERN.finditer(text))
    return sorted(hits, key=lambda item: item.lower())




def write_pptx_check(project: Path, pptx_path: Path, stats: dict[str, int | bool | str]) -> Path:
    sha256 = file_sha256(pptx_path)
    check_path = project / f"AD-creative/ppt/checks/{pptx_path.stem}_{sha256[:12]}_editability.md"
    if not check_path.exists():
        write_text(
            check_path,
            pptx_editability_report_content(
                project,
                pptx_path,
                stats,
                sha256=sha256,
            ),
        )
    return check_path


RISKY_CLIENT_COPY_PATTERNS = [
    "internal_only",
    "internal note",
    "内部",
    "TODO",
    "TBD",
    "fake logo",
    "假 logo",
    "假logo",
    "placeholder-only",
    "simulated",
    "模拟",
    "do not send",
]
RISKY_CLIENT_COPY_PATTERNS_LOWER = tuple(pattern.lower() for pattern in RISKY_CLIENT_COPY_PATTERNS)
RISKY_CLIENT_COPY_PATTERN = re.compile(
    "|".join(
        re.escape(pattern)
        for pattern in sorted(RISKY_CLIENT_COPY_PATTERNS_LOWER, key=len, reverse=True)
    )
)


def candidate_client_files(project: Path, artifacts: list[dict[str, str]]) -> list[Path]:
    files: list[Path] = []
    for artifact in artifacts:
        visibility = artifact.get("visibility", "").lower()
        if visibility not in CLIENT_VISIBLE_VALUES:
            continue
        rel_path = artifact.get("path", "").strip()
        try:
            path = contained_project_path(project, rel_path, "client-visible artifact path")
        except ValueError:
            continue
        if is_private_local_state_path(project, path):
            continue
        if path.is_file() and path.suffix.lower() in {".md", ".txt", ".csv", ".pptx", ".pdf"}:
            files.append(path)
    return sorted(files)


def candidate_client_language_files(project: Path, artifacts: list[dict[str, str]]) -> list[Path]:
    files: set[Path] = set(candidate_client_files(project, artifacts))
    for root in [
        project / "AD-creative/client_review",
        project / "04_客户审阅_ClientReview",
        project / "05_最终交付_FinalDelivery",
    ]:
        if not root.exists():
            continue
        for path in root.rglob("*"):
            if path.is_file() and is_reliably_client_language_file(project, path):
                files.add(path)
    return sorted(files)


def artifact_path_by_id(project: Path, artifacts: list[dict[str, str]], artifact_id: str) -> Path | None:
    for artifact in artifacts:
        if artifact.get("artifact_id", "").strip() == artifact_id:
            rel_path = artifact.get("path", "").strip()
            return project / rel_path if rel_path else None
    return None


def current_pptx_path(project: Path, artifacts: list[dict[str, str]]) -> tuple[Path | None, bool]:
    current_truth = project / "AD-creative/orchestrator/current_truth.md"
    if not current_truth.exists():
        return None, False
    artifact_id = current_truth_value(
        current_truth.read_text(encoding="utf-8"), "current_pptx_artifact_id"
    )
    if not artifact_id:
        return None, False
    return artifact_path_by_id(project, artifacts, artifact_id), True


def current_delivery_paths(
    project: Path,
    artifacts: list[dict[str, str]],
) -> dict[str, Path | None]:
    truth_path = project / "AD-creative/orchestrator/current_truth.md"
    truth = truth_path.read_text(encoding="utf-8") if truth_path.exists() else ""
    keys = {
        "pptx": "current_pptx_artifact_id",
        "pdf": "current_pdf_artifact_id",
        "preview": "current_preview_artifact_id",
        "text_extract": "current_text_extract_artifact_id",
        "editability": "current_ppt_editability_artifact_id",
    }
    return {
        name: artifact_path_by_id(project, artifacts, current_truth_value(truth, key))
        if current_truth_value(truth, key)
        else None
        for name, key in keys.items()
    }


def _build_client_pack_input_manifest_bound(
    project: Path,
    artifacts: list[dict[str, str]],
    project_root_fd: int,
) -> tuple[dict[str, object], str, list[str]]:
    truth_path = project / "AD-creative/orchestrator/current_truth.md"
    try:
        truth = _read_project_text(
            project,
            "AD-creative/orchestrator/current_truth.md",
            project_root_fd=project_root_fd,
        )
    except FileNotFoundError:
        truth = ""
    version_id = current_truth_value(truth, "current_version_id")
    pointer_keys = [
        "current_pptx_artifact_id",
        "current_pdf_artifact_id",
        "current_preview_artifact_id",
        "current_text_extract_artifact_id",
        "current_ppt_editability_artifact_id",
    ]
    artifact_by_id = {
        row.get("artifact_id", ""): row
        for row in artifacts
        if row.get("artifact_id", "")
    }
    pointers = {key: current_truth_value(truth, key) for key in pointer_keys}
    errors: list[str] = []
    paths: set[Path] = {
        truth_path,
        project / "AD-creative/orchestrator/version_map.csv",
        project / "AD-creative/client_review/client_outline.csv",
        project / CLIENT_OUTLINE_CONFIRMATION_REL,
        project / "AD-creative/visual_assets/asset_manifest.csv",
        project / "AD-creative/visual_assets/asset_current_manifest.csv",
        project / "AD-creative/visual_assets/asset_authorizations.csv",
        project / "AD-creative/visual_review/review_matrix.csv",
        project / "AD-creative/references/reference_cards.csv",
        project / "AD-creative/feedback/feedback_map.csv",
    }
    for key, artifact_id in pointers.items():
        row = artifact_by_id.get(artifact_id)
        if not artifact_id or not row:
            errors.append(f"{key} 未解析到 exact-current artifact")
            continue
        try:
            path = contained_project_path(
                project, row.get("path", ""), f"{key} artifact path"
            )
        except ValueError as exc:
            errors.append(str(exc))
            continue
        paths.add(path)
    paths.update(candidate_client_language_files(project, artifacts))
    _, current_assets = _read_project_csv_rows(
        project,
        "AD-creative/visual_assets/asset_current_manifest.csv",
        project_root_fd=project_root_fd,
    )
    for row in current_assets:
        if not (
            normalized_bool(row.get("direct_client_use"))
            or row.get("used_in_slide", "").strip()
        ):
            continue
        try:
            paths.add(
                contained_project_path(
                    project,
                    row.get("path", ""),
                    f"client-used asset {row.get('asset_id', '')}",
                )
            )
        except ValueError as exc:
            errors.append(str(exc))
    entries: list[dict[str, object]] = []
    try:
        private_markers: list[bytes] | None = private_source_path_markers(
            project,
            project_root_fd=project_root_fd,
        )
    except ValueError as exc:
        private_markers = None
        errors.append(str(exc))
    for path in sorted(paths):
        try:
            rel_path = canonical_project_relative(project, path)
        except ValueError:
            errors.append(f"client-pack input escapes project: {path}")
            continue
        if Path(rel_path).parts[0] == PRIVATE_LOCAL_STATE_REL.name:
            errors.append(f"client-pack input uses private local state: {rel_path}")
            continue
        if private_markers is None:
            continue
        privacy_issue, digest, size_bytes = private_source_candidate_evidence(
            project,
            rel_path,
            private_markers,
            project_root_fd=project_root_fd,
        )
        if privacy_issue:
            errors.append(
                f"client-pack input privacy blocked: {rel_path}: {privacy_issue}"
            )
            continue
        entries.append(
            {
                "path": rel_path,
                "sha256": digest,
                "size_bytes": size_bytes,
            }
        )
    payload: dict[str, object] = {
        "protocol_id": "adco.client-pack-input-manifest",
        "version": "1.0",
        "current_version_id": version_id,
        "artifact_pointers": pointers,
        "files": entries,
    }
    digest = hashlib.sha256(
        json.dumps(payload, ensure_ascii=False, sort_keys=True).encode("utf-8")
    ).hexdigest()
    return payload, digest, errors


def build_client_pack_input_manifest(
    project: Path,
    artifacts: list[dict[str, str]],
) -> tuple[dict[str, object], str, list[str]]:
    pointer_keys = [
        "current_pptx_artifact_id",
        "current_pdf_artifact_id",
        "current_preview_artifact_id",
        "current_text_extract_artifact_id",
        "current_ppt_editability_artifact_id",
    ]

    def failed_result(message: str) -> tuple[dict[str, object], str, list[str]]:
        payload: dict[str, object] = {
            "protocol_id": "adco.client-pack-input-manifest",
            "version": "1.0",
            "current_version_id": "",
            "artifact_pointers": {key: "" for key in pointer_keys},
            "files": [],
        }
        digest = hashlib.sha256(
            json.dumps(payload, ensure_ascii=False, sort_keys=True).encode("utf-8")
        ).hexdigest()
        return payload, digest, [message]

    project_root_fd: int | None = None
    try:
        project_root_fd = _open_project_root_fd(project)
        payload, digest, errors = _build_client_pack_input_manifest_bound(
            project,
            artifacts,
            project_root_fd,
        )
        if not _project_root_fd_is_current(project, project_root_fd):
            return failed_result(
                "client-pack project root changed during manifest construction"
            )
        return payload, digest, errors
    except (FileNotFoundError, OSError, UnicodeError, ValueError) as exc:
        return failed_result(f"client-pack project binding failed: {exc}")
    finally:
        if project_root_fd is not None:
            os.close(project_root_fd)


def write_client_pack_binding(
    project: Path,
    *,
    payload: dict[str, object],
    package_digest: str,
    report_path: Path,
    status: str,
) -> tuple[Path, Path]:
    version_id = validate_specialist_token(
        str(payload.get("current_version_id", "")), "current_version_id"
    )
    manifest_path = (
        project
        / "AD-creative/delivery/client_pack_manifests"
        / f"{version_id}_{package_digest}.json"
    )
    manifest_core = {**payload, "package_digest": package_digest}
    manifest_payload = {**manifest_core, "created_at": now_iso()}
    if manifest_path.exists():
        existing = read_json_object(manifest_path, "client pack manifest")
        existing_core = {
            key: value for key, value in existing.items() if key != "created_at"
        }
        if existing_core != manifest_core or not str(
            existing.get("created_at", "")
        ).strip():
            raise ValueError("immutable client pack manifest collision")
    else:
        write_json_object(manifest_path, manifest_payload)
    binding_path = project / "AD-creative/delivery/client_pack_binding.json"
    write_json_object(
        binding_path,
        {
            "protocol_id": "adco.client-pack-binding",
            "version": "1.0",
            "status": status,
            "current_version_id": version_id,
            "package_digest": package_digest,
            "manifest_path": safe_rel(project, manifest_path),
            "manifest_sha256": file_sha256(manifest_path),
            "report_path": safe_rel(project, report_path),
            "report_sha256": file_sha256(report_path),
            "checked_at": now_iso(),
        },
    )
    return manifest_path, binding_path


def current_client_pack_binding_errors(
    project: Path,
    artifacts: list[dict[str, str]],
) -> tuple[list[str], str]:
    binding_path = project / "AD-creative/delivery/client_pack_binding.json"
    if not binding_path.is_file():
        return ["缺少 current client-pack binding；先重跑 client-pack-gate"], ""
    try:
        binding = read_json_object(binding_path, "client pack binding")
    except ValueError as exc:
        return [str(exc)], ""
    errors: list[str] = []
    if binding.get("status") != "PASS":
        errors.append("current client-pack binding status 不是 PASS")
    current_payload, current_digest, manifest_errors = build_client_pack_input_manifest(
        project, artifacts
    )
    errors.extend(manifest_errors)
    bound_digest = str(binding.get("package_digest", ""))
    if not re.fullmatch(r"[0-9a-f]{64}", bound_digest):
        errors.append("client-pack binding package_digest 无效")
    elif current_digest != bound_digest:
        errors.append("client-pack binding 已过期：exact-current package inputs changed")
    checked_paths: dict[str, Path] = {}
    for key in ["manifest", "report"]:
        try:
            path = contained_project_path(
                project, str(binding.get(f"{key}_path", "")), f"client-pack {key}"
            )
        except ValueError as exc:
            errors.append(str(exc))
            continue
        checked_paths[key] = path
        if not path.is_file() or file_sha256(path) != binding.get(f"{key}_sha256"):
            errors.append(f"client-pack {key} missing or hash mismatch")
    manifest_path = checked_paths.get("manifest")
    if manifest_path and manifest_path.is_file():
        try:
            manifest = read_json_object(manifest_path, "client pack manifest")
        except ValueError as exc:
            errors.append(str(exc))
        else:
            expected_core = {**current_payload, "package_digest": current_digest}
            actual_core = {
                key: value for key, value in manifest.items() if key != "created_at"
            }
            if actual_core != expected_core:
                errors.append(
                    "client-pack immutable manifest content does not match exact-current inputs"
                )
            expected_name = (
                f"{current_payload.get('current_version_id')}_{current_digest}.json"
            )
            if manifest_path.name != expected_name:
                errors.append("client-pack manifest path is not digest-addressed")
    return errors, bound_digest






def review_client_pack(project: Path, pptx_path: Path | None = None) -> tuple[str, list[str], Path]:
    migrate_control_plane(project)
    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    package_paths = current_delivery_paths(project, artifacts)
    _, version_map = read_csv_rows(project / "AD-creative/orchestrator/version_map.csv")
    _, feedback_rows = read_csv_rows(project / "AD-creative/feedback/feedback_map.csv")
    _, assets = read_csv_rows(project / "AD-creative/visual_assets/asset_manifest.csv")
    current_manifest_path = sync_asset_current_manifest(project)
    _, current_assets = read_csv_rows(current_manifest_path)
    current_by_id = {
        row.get("asset_id", "").strip(): row
        for row in current_assets
        if row.get("asset_id", "").strip()
    }
    _, authorization_rows = read_csv_rows(
        project / "AD-creative/visual_assets/asset_authorizations.csv"
    )
    _, references = read_csv_rows(project / "AD-creative/references/reference_cards.csv")
    issues: list[str] = []
    warnings: list[str] = []
    evidence: list[str] = []

    outline_status, outline_findings, outline_report = review_client_outline(project)
    evidence.append(f"client_outline_gate={outline_status} {safe_rel(project, outline_report)}")
    if outline_status == "BLOCKED":
        issues.extend(f"Client Outline Gate: {item}" for item in outline_findings[:8])
    elif outline_findings:
        warnings.extend(f"Client Outline Gate: {item}" for item in outline_findings[:8])

    for row in client_outline_rows(project):
        slide_id = row.get("slide_id", "<missing>")
        if row.get("visibility", "").strip().lower() not in CLIENT_VISIBLE_VALUES:
            issues.append(f"Client Outline: {slide_id} 仍是 internal/pending，不能满足客户包准备。")
        if row.get("status", "").strip().lower() not in {"ready", "approved", "done"}:
            issues.append(f"Client Outline: {slide_id} status 未 ready/approved/done。")

    language_status, language_findings, language_report = review_client_language(
        project,
        extra_paths=[
            path
            for name, path in package_paths.items()
            if name in {"pptx", "pdf", "text_extract"} and path is not None
        ],
    )
    evidence.append(f"client_language_gate={language_status} {safe_rel(project, language_report)}")
    if language_status == "BLOCKED":
        issues.extend(f"Client Language Gate: {item}" for item in language_findings[:8])
    elif language_findings:
        warnings.extend(f"Client Language Gate: {item}" for item in language_findings[:8])

    layout_status, layout_findings, layout_report = review_visual_layout(project)
    evidence.append(f"visual_layout_gate={layout_status} {safe_rel(project, layout_report)}")
    if layout_status == "BLOCKED":
        issues.extend(f"Visual Layout Gate: {item}" for item in layout_findings[:8])
    elif layout_findings:
        warnings.extend(f"Visual Layout Gate: {item}" for item in layout_findings[:8])

    for artifact in artifacts:
        visibility = artifact.get("visibility", "").lower()
        gate_status = artifact.get("gate_status", "").lower()
        if visibility in CLIENT_VISIBLE_VALUES and gate_status not in PASS_GATE_VALUES:
            issues.append(f"客户可见产物未通过 Gate: {artifact.get('artifact_id')}")

    issues.extend(
        validate_client_delivery_readiness(
            project, artifacts, version_map, feedback_rows
        )
    )

    for asset in assets:
        visibility = asset.get("visibility", "").lower()
        qa_status = asset.get("qa_status", "").upper()
        status = asset.get("status", "").lower()
        reference_id = asset.get("reference_id", "").strip()
        if visibility in CLIENT_VISIBLE_VALUES:
            if qa_status != "PASS":
                issues.append(f"客户可见图片 QA 未 PASS: {asset.get('asset_id')}")
            if status not in CLIENT_REVIEW_ASSET_STATUSES:
                issues.append(f"客户可见图片未进入 selected/approved: {asset.get('asset_id')}")
            asset_id = asset.get("asset_id", "").strip()
            current = current_by_id.get(asset_id, {})
            current_sha = current.get("sha256", "").strip()
            if not current_sha:
                issues.append(f"客户可见图片 current manifest 缺少 sha256: {asset_id}")
            elif not matching_asset_authorization(
                project, asset_id, current_sha, authorization_rows
            ):
                issues.append(
                    f"客户可见图片缺少匹配 asset hash/scope 的独立授权 receipt: {asset_id}"
                )
            if reference_id in {"", "pending"}:
                issues.append(f"客户可见图片未绑定真实参考: {asset.get('asset_id')}")

    for reference in references:
        if reference.get("client_visible", "").lower() == "true":
            url = reference.get("url", "")
            if not url.startswith("https://"):
                issues.append(f"客户可见参考不是 https: {reference.get('reference_id')}")
            if reference.get("do_not_copy", "").strip() == "":
                issues.append(f"客户可见参考缺少 do_not_copy: {reference.get('reference_id')}")

    exact_current_pptx, current_pptx_declared = current_pptx_path(project, artifacts)
    if current_pptx_declared and not exact_current_pptx:
        issues.append("current_pptx_artifact_id 未解析到带 path 的 PPTX artifact。")
    if pptx_path and exact_current_pptx and pptx_path.resolve() != exact_current_pptx.resolve():
        issues.append(
            f"PPTX 检查目标不是 current_pptx_artifact_id 指向的当前文件: {safe_rel(project, pptx_path)}"
        )
    if not current_pptx_declared:
        issues.append("缺少 current_pptx_artifact_id；client-pack-gate 只检查已登记的 exact current PPTX。")
    check_target = exact_current_pptx if current_pptx_declared else None
    pptx_stats: dict[str, int | bool | str] | None = None
    if check_target:
        if not check_target.exists():
            issues.append(f"PPTX 文件不存在: {safe_rel(project, check_target)}")
            evidence.append(
                f"pptx={safe_rel(project, check_target)} missing=true exact_current={current_pptx_declared}"
            )
        else:
            pptx_stats = inspect_pptx(check_target)
            if not pptx_stats["editable"]:
                issues.append("PPTX 缺少可编辑文本层。")
            evidence.append(
                f"pptx={safe_rel(project, check_target)} slides={pptx_stats['slides']} editable_text_runs={pptx_stats['editable_text_runs']} exact_current={current_pptx_declared}"
            )
    else:
        issues.append("未找到可审稿 PPTX。")

    scanned_files = candidate_client_files(project, artifacts)
    risky_hits: list[str] = []
    for path in scanned_files:
        try:
            text = client_language_text_for_path(path)
        except RuntimeError as exc:
            issues.append(f"客户稿候选无法可靠解析: {safe_rel(project, path)}: {exc}")
            continue
        lowered = text.lower()
        risky_patterns = {match.group(0) for match in RISKY_CLIENT_COPY_PATTERN.finditer(lowered)}
        risky_hits.extend(f"{safe_rel(project, path)}: {pattern}" for pattern in risky_patterns)
    if risky_hits:
        issues.extend(f"客户稿候选含风险词: {hit}" for hit in sorted(risky_hits)[:12])

    manifest_payload, package_digest, manifest_issues = build_client_pack_input_manifest(
        project, artifacts
    )
    issues.extend(manifest_issues)
    evidence.append(f"package_digest={package_digest}")
    adversarial_target = write_adversarial_target_snapshot(
        project,
        stage="final_delivery",
        payload=manifest_payload,
        target_digest=package_digest,
    )
    evidence.append(f"adversarial_target={safe_rel(project, adversarial_target)}")
    status = "PASS" if not issues else "BLOCKED"
    status = enforce_adversarial_gate_policy(
        project,
        "final_delivery",
        status,
        warnings,
        evidence,
        expected_targets=[adversarial_target],
    )
    version_id = str(manifest_payload.get("current_version_id") or "UNVERSIONED")
    report_path = (
        project
        / "AD-creative/gates"
        / f"GATE-AUTO-CLIENT-PACK-001_{safe_artifact_suffix(version_id)}_{package_digest[:12]}.md"
    )
    evidence.extend(
        [
            f"client_candidate_files={len(scanned_files)}",
            f"artifacts={len(artifacts)}",
            f"assets={len(assets)}",
            f"references={len(references)}",
        ]
    )
    issue_text = "\n".join(f"- {issue}" for issue in issues) or "- 无"
    warning_text = "\n".join(f"- {warning}" for warning in warnings) or "- 无"
    evidence_text = "\n".join(f"- {item}" for item in evidence)
    write_text(
        report_path,
        f"""# Client Pack Gate

status: {status}
checked_at: {now_iso()}
visibility: internal_only
current_version_id: {version_id}
package_digest: {package_digest}

## Evidence

{evidence_text}

## Issues

{issue_text}

## Warnings

{warning_text}

## Rules

- 客户可见产物必须 Gate PASS。
- 客户可见图片必须 QA PASS 且进入 selected/approved。
- 客户可见参考必须是 https，并保留 do_not_copy。
- PPTX 必须存在可编辑文本层。
- 客户稿候选不能含内部注释、模拟标记、TODO/TBD、假 logo、placeholder-only。
- 当前交付包必须具备 current version、PPTX、PDF、preview、text extract、PPT editability 和 feedback closure 证据。
""",
    )
    manifest_payload_for_binding = {
        **manifest_payload,
        "current_version_id": version_id,
    }
    manifest_path, binding_path = write_client_pack_binding(
        project,
        payload=manifest_payload_for_binding,
        package_digest=package_digest,
        report_path=report_path,
        status=status,
    )
    manifest_artifact_id = f"ART-CLIENT-PACK-MANIFEST-{package_digest[:12].upper()}"
    update_artifact(
        project,
        manifest_artifact_id,
        "client_pack_input_manifest",
        safe_rel(project, manifest_path),
        "final_delivery",
        status="done",
        visibility="internal_only",
        gate_status=status,
        version=version_id,
        derived_from_sha256=package_digest,
        sha256=file_sha256(manifest_path),
        size_bytes=str(manifest_path.stat().st_size),
    )
    update_artifact(
        project,
        "ART-AUTO-CLIENT-PACK-BINDING",
        "client_pack_binding",
        safe_rel(project, binding_path),
        "final_delivery",
        status="done" if status == "PASS" else "blocked",
        visibility="internal_only",
        gate_status=status,
        version=version_id,
        derived_from_artifact_id=manifest_artifact_id,
        derived_from_sha256=file_sha256(manifest_path),
        sha256=file_sha256(binding_path),
        size_bytes=str(binding_path.stat().st_size),
    )
    update_artifact(
        project,
        "ART-AUTO-CLIENT-PACK-GATE",
        "client_pack_gate_report",
        safe_rel(project, report_path),
        "final_delivery",
        status="done" if status != "BLOCKED" else "blocked",
        visibility="internal_only",
        gate_status=status,
        version=version_id,
        derived_from_artifact_id=manifest_artifact_id,
        derived_from_sha256=package_digest,
        sha256=file_sha256(report_path),
        size_bytes=str(report_path.stat().st_size),
    )
    append_gate(
        project,
        "GATE-AUTO-CLIENT-PACK-001",
        "final_delivery",
        status,
        "92" if status == "PASS" else "65" if status == "PARTIAL_PASS" else "40",
        f"ART-AUTO-CLIENT-PACK-GATE;{manifest_artifact_id};ART-AUTO-CLIENT-PACK-BINDING",
        ";".join(issues[:8]),
        ";".join(warnings[:8]) or "修正客户稿风险后重跑 client-pack-gate。",
        "",
        "ready_for_independent_human_review" if status == "PASS" else "revise_client_pack",
        "ad_creative_operator",
        target_ref=safe_rel(project, binding_path),
        target_sha256=file_sha256(binding_path),
    )
    append_event(
        project,
        {
            "event_id": "EVT-AUTO-CLIENT-PACK-GATE",
            "event_type": "client_pack_gate_run",
            "created_at": now_iso(),
            "status": status,
            "issues": issues[:12],
            "warnings": warnings[:12],
            "current_version_id": version_id,
            "package_digest": package_digest,
            "binding": safe_rel(project, binding_path),
        },
    )
    return status, issues + warnings, report_path


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def ensure_csv_file(path: Path, fields: list[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    if not path.exists() or not path.read_text(encoding="utf-8").strip():
        write_csv_rows(path, fields, [])
        return
    ensure_csv_fields(path, fields)


def canonical_row_sha256(row: dict[str, str]) -> str:
    payload = json.dumps(
        {str(key): row.get(key, "") or "" for key in sorted(key for key in row if key is not None)},
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
    ).encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def legacy_cleanup_reference_path(row: dict[str, str]) -> bool:
    lifecycle = normalized_artifact_lifecycle(row)
    if lifecycle not in ARTIFACT_INACTIVE_LIFECYCLE_VALUES:
        return False
    path = (row.get("path") or "").strip().lower().replace("-", "_")
    artifact_type = (row.get("artifact_type") or "").strip().lower()
    cleanup_like = any(
        token in path
        for token in ("cleanup", "dedupe", "/gates/", "gate_report", "removal_report")
    )
    artifact_is_report = any(
        token in artifact_type
        for token in ("cleanup", "dedupe", "gate_report", "removal_report")
    )
    return cleanup_like and not artifact_is_report


def migrate_artifact_lifecycle_rows(
    rows: list[dict[str, str]],
) -> tuple[list[dict[str, str]], list[dict[str, object]]]:
    reverse_supersession: dict[str, list[str]] = {}
    for row in rows:
        new_id = (row.get("artifact_id") or "").strip()
        for old_id in split_asset_refs(row.get("supersedes_artifact_id", "")):
            reverse_supersession.setdefault(old_id, []).append(new_id)
    migrated: list[dict[str, str]] = []
    legacy_evidence: list[dict[str, object]] = []
    superseded_ids = set(reverse_supersession)
    for index, raw_row in enumerate(rows, start=2):
        row = dict(raw_row)
        original_lifecycle = (row.get("lifecycle_state") or "").strip()
        lifecycle = normalized_artifact_lifecycle(row, superseded_ids=superseded_ids)
        cleanup_ref_is_path = legacy_cleanup_reference_path(row)
        changed = not original_lifecycle or original_lifecycle != lifecycle
        if cleanup_ref_is_path:
            lifecycle = "legacy_unresolved_tombstone"
            if not (row.get("cleanup_ref") or "").strip():
                row["cleanup_ref"] = (row.get("path") or "").strip()
                changed = True
            # A cleanup report is evidence, never a fabricated original artifact path.
            if (row.get("original_path") or "").strip() == (row.get("path") or "").strip():
                row["original_path"] = ""
                changed = True
        elif not (row.get("original_path") or "").strip() and (row.get("path") or "").strip():
            row["original_path"] = (row.get("path") or "").strip()
            changed = True
        row["lifecycle_state"] = lifecycle
        if not (row.get("status_reason") or "").strip() and (row.get("status") or "").strip():
            row["status_reason"] = (row.get("status") or "").strip()
            changed = True
        if lifecycle in {"removed", "withdrawn", "rejected", "legacy_unresolved_tombstone"}:
            if not (row.get("removal_reason") or "").strip():
                row["removal_reason"] = first_nonempty(
                    row.get("status_reason"), row.get("status"), default="legacy lifecycle classification"
                )
                changed = True
        artifact_id = (row.get("artifact_id") or "").strip()
        superseded_by = ";".join(
            sorted(set(filter(None, reverse_supersession.get(artifact_id, []))))
        )
        if superseded_by and (row.get("superseded_by") or "").strip() != superseded_by:
            row["superseded_by"] = superseded_by
            changed = True
        if changed:
            legacy_evidence.append(
                {
                    "row_number": index,
                    "evidence_sha256": canonical_row_sha256(raw_row),
                    "raw": raw_row,
                    "classification": lifecycle,
                }
            )
        migrated.append(row)
    return migrated, legacy_evidence


def markdown_named_sections(text: str, heading: str) -> list[str]:
    return re.findall(
        rf"(?ims)^##[ \t]+{re.escape(heading)}[ \t]*\n(.*?)(?=^##[ \t]+|\Z)",
        text,
    )


def section_key_values(section: str) -> dict[str, str]:
    values: dict[str, str] = {}
    for line in section.splitlines():
        match = re.match(r"^[ \t]*[-*]?[ \t]*([A-Za-z][A-Za-z0-9_ -]{1,80})[ \t]*:[ \t]*(.*?)[ \t]*$", line)
        if match:
            values[match.group(1).strip().lower().replace(" ", "_")] = match.group(2).strip().strip("`")
    return values


def artifact_id_for_legacy_value(
    project: Path, value: str, artifacts: list[dict[str, str]]
) -> tuple[str, list[str]]:
    if not value:
        return "", []
    direct = [
        row.get("artifact_id", "").strip()
        for row in artifacts
        if row.get("artifact_id", "").strip() == value
    ]
    if len(direct) == 1:
        return direct[0], direct
    rel = canonical_project_relative_path(project, value)
    matches = [
        row.get("artifact_id", "").strip()
        for row in artifacts
        if row.get("artifact_id", "").strip()
        and rel
        and canonical_project_relative_path(
            project, row.get("path") or row.get("original_path")
        )
        == rel
    ]
    if len(matches) == 1:
        return matches[0], matches
    return "", sorted(set(direct + matches))


def legacy_current_package_backfill(
    project: Path,
    truth_text: str,
    artifacts: list[dict[str, str]],
    versions: list[dict[str, str]],
) -> tuple[dict[str, str], list[dict[str, object]], list[str]]:
    sections = markdown_named_sections(truth_text, "Current Package")
    if not sections:
        return {}, [], []
    raw_sections = sections[:]
    if len(sections) != 1:
        return {}, [
            {
                "severity": "P0",
                "scope": "current",
                "code": "ambiguous_legacy_current_package_sections",
                "message": f"legacy Current Package has {len(sections)} sections; migration refused to choose",
            }
        ], raw_sections
    values = section_key_values(sections[0])
    aliases = {
        "current_pptx_artifact_id": ("current_pptx_artifact_id", "pptx_artifact_id", "pptx", "pptx_path"),
        "current_pdf_artifact_id": ("current_pdf_artifact_id", "pdf_artifact_id", "pdf", "pdf_path"),
        "current_preview_artifact_id": ("current_preview_artifact_id", "preview_artifact_id", "preview", "preview_path"),
        "current_text_extract_artifact_id": ("current_text_extract_artifact_id", "text_extract_artifact_id", "text_extract", "text_extract_path"),
        "current_ppt_editability_artifact_id": ("current_ppt_editability_artifact_id", "ppt_editability_artifact_id", "editability", "editability_path"),
    }
    result: dict[str, str] = {}
    blockers: list[dict[str, object]] = []
    for target_key, source_keys in aliases.items():
        raw_value = next((values[key] for key in source_keys if values.get(key)), "")
        if not raw_value:
            continue
        artifact_id, candidates = artifact_id_for_legacy_value(project, raw_value, artifacts)
        if not artifact_id:
            blockers.append(
                {
                    "severity": "P0",
                    "scope": "current",
                    "code": "ambiguous_legacy_current_package_artifact",
                    "message": f"{target_key} legacy value {raw_value!r} did not resolve to exactly one artifact",
                    "candidates": candidates,
                }
            )
        else:
            result[target_key] = artifact_id
    raw_version_id = first_nonempty(
        values.get("current_version_id"), values.get("version_id"), default=""
    )
    version_candidates = [
        row
        for row in versions
        if raw_version_id and row.get("version_id", "").strip() == raw_version_id
    ]
    if not version_candidates and result.get("current_pptx_artifact_id"):
        version_candidates = [
            row
            for row in versions
            if row.get("artifact_id", "").strip() == result["current_pptx_artifact_id"]
            and row.get("status", "").strip().lower()
            in CURRENT_VIEW_VERSION_STATUSES
        ]
    if len(version_candidates) != 1:
        blockers.append(
            {
                "severity": "P0",
                "scope": "current",
                "code": "ambiguous_legacy_current_package_version",
                "message": f"legacy Current Package did not resolve to exactly one current version; candidates={len(version_candidates)}",
                "candidates": [row.get("version_id", "") for row in version_candidates],
            }
        )
    else:
        version = version_candidates[0]
        result["current_version_id"] = version.get("version_id", "").strip()
        result["version_map_status"] = version.get("status", "").strip()
        result.setdefault("current_pptx_artifact_id", version.get("artifact_id", "").strip())
    result["last_archive_before_edit"] = values.get("last_archive_before_edit", "")
    return result, blockers, raw_sections


def replace_current_version_truth(
    text: str, values: dict[str, str]
) -> str:
    body = "```text\n" + "".join(
        f"{key}: {values.get(key, '').strip()}\n" for key in CURRENT_VERSION_TRUTH_KEYS
    ) + "```\n"
    pattern = re.compile(
        r"(?ims)^##[ \t]+Current Version Truth[ \t]*\n.*?(?=^##[ \t]+|\Z)"
    )
    section = "## Current Version Truth\n\n" + body
    matches = list(pattern.finditer(text))
    if len(matches) == 1:
        return pattern.sub(section, text, count=1)
    if not matches:
        return text.rstrip() + "\n\n" + section
    return text


def project_yml_with_schema_version(text: str) -> str:
    if re.search(r"(?m)^  schema_version:[ \t]*", text):
        return re.sub(
            r"(?m)^  schema_version:[ \t]*.*$",
            f'  schema_version: "{CONTROL_PLANE_SCHEMA_VERSION}"',
            text,
            count=1,
        )
    project_match = re.search(r"(?m)^project:[ \t]*$", text)
    if project_match:
        insert_at = project_match.end()
        return text[:insert_at] + f'\n  schema_version: "{CONTROL_PLANE_SCHEMA_VERSION}"' + text[insert_at:]
    return f'project:\n  schema_version: "{CONTROL_PLANE_SCHEMA_VERSION}"\n' + text


def current_version_truth_binding_valid(
    truth_text: str,
    artifacts: list[dict[str, str]],
    versions: list[dict[str, str]],
) -> bool:
    sections = markdown_named_sections(truth_text, "Current Version Truth")
    if len(sections) != 1:
        return False
    values = section_key_values(sections[0])
    version_id = values.get("current_version_id", "").strip()
    artifact_id = values.get("current_pptx_artifact_id", "").strip()
    if not version_id or not artifact_id:
        return False
    version_matches = [
        row for row in versions if row.get("version_id", "").strip() == version_id
    ]
    artifact_matches = [
        row for row in artifacts if row.get("artifact_id", "").strip() == artifact_id
    ]
    if len(version_matches) != 1 or len(artifact_matches) != 1:
        return False
    if version_matches[0].get("artifact_id", "").strip() != artifact_id:
        return False
    version_status = version_matches[0].get("status", "").strip().lower()
    if version_status not in CURRENT_VIEW_VERSION_STATUSES:
        return False
    truth_status = current_truth_value(truth_text, "version_map_status").strip().lower()
    if truth_status and truth_status != version_status:
        return False
    return normalized_artifact_lifecycle(artifact_matches[0]) == "active"


def canonical_payload_sha256(payload: object) -> str:
    canonical = json.dumps(
        payload,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
    ).encode("utf-8")
    return hashlib.sha256(canonical).hexdigest()


def normalized_migration_blockers(
    blockers: list[dict[str, object]],
) -> list[dict[str, object]]:
    normalized = [
        json.loads(json.dumps(blocker, ensure_ascii=False, sort_keys=True))
        for blocker in blockers
    ]
    return sorted(
        normalized,
        key=lambda blocker: (
            str(blocker.get("severity") or ""),
            str(blocker.get("scope") or ""),
            str(blocker.get("code") or ""),
            str(blocker.get("message") or ""),
            canonical_payload_sha256(blocker),
        ),
    )


def updated_migration_manifest(
    existing: dict[str, object] | None,
    *,
    source_hashes: dict[str, str],
    changes: list[str],
    active_blockers: list[dict[str, object]],
    raw_legacy_evidence: dict[str, object],
    observed_hashes: dict[str, str],
) -> dict[str, object]:
    """Preserve raw evidence while recording only state-changing migration attempts."""
    manifest = dict(existing or {})
    legacy_blockers = manifest.pop("blockers", [])
    if not isinstance(legacy_blockers, list):
        legacy_blockers = []
    normalized_active = normalized_migration_blockers(active_blockers)
    attempt_id = canonical_payload_sha256(
        {
            "schema_version": CONTROL_PLANE_SCHEMA_VERSION,
            "observed_hashes": observed_hashes,
            "active_blockers": normalized_active,
        }
    )

    raw_history = manifest.get("blocker_history", [])
    history_by_fingerprint: dict[str, dict[str, object]] = {}
    if isinstance(raw_history, list):
        for entry in raw_history:
            if not isinstance(entry, dict):
                continue
            blocker = entry.get("blocker")
            if not isinstance(blocker, dict):
                continue
            fingerprint = str(entry.get("fingerprint") or canonical_payload_sha256(blocker))
            history_by_fingerprint[fingerprint] = dict(entry)

    historical_candidates = [
        blocker for blocker in legacy_blockers if isinstance(blocker, dict)
    ] + normalized_active
    for blocker in historical_candidates:
        normalized = normalized_migration_blockers([blocker])[0]
        fingerprint = canonical_payload_sha256(normalized)
        history_by_fingerprint.setdefault(
            fingerprint,
            {
                "fingerprint": fingerprint,
                "blocker": normalized,
                "first_seen_attempt_id": attempt_id,
                "last_seen_attempt_id": attempt_id,
                "resolved_attempt_id": "",
            },
        )

    active_fingerprints = {
        canonical_payload_sha256(blocker) for blocker in normalized_active
    }
    for fingerprint, entry in history_by_fingerprint.items():
        if fingerprint in active_fingerprints:
            entry["last_seen_attempt_id"] = attempt_id
            entry["resolved_attempt_id"] = ""
        elif not entry.get("resolved_attempt_id"):
            entry["resolved_attempt_id"] = attempt_id

    attempts = manifest.get("attempts", [])
    if not isinstance(attempts, list):
        attempts = []
    if not any(
        isinstance(attempt, dict) and attempt.get("attempt_id") == attempt_id
        for attempt in attempts
    ):
        attempts = [
            *attempts,
            {
                "attempt_id": attempt_id,
                "recorded_at": now_iso(),
                "observed_hashes": observed_hashes,
                "changes": list(changes),
                "active_blockers": normalized_active,
            },
        ]

    existing_raw = manifest.get("raw_legacy_evidence", {})
    merged_raw: dict[str, object] = (
        dict(existing_raw) if isinstance(existing_raw, dict) else {}
    )
    for key, value in raw_legacy_evidence.items():
        existing_value = merged_raw.get(key)
        if isinstance(value, dict):
            merged_mapping = (
                dict(existing_value) if isinstance(existing_value, dict) else {}
            )
            for item_key, item_value in value.items():
                merged_mapping.setdefault(item_key, item_value)
            merged_raw[key] = merged_mapping
        elif isinstance(value, list):
            merged_list = (
                list(existing_value) if isinstance(existing_value, list) else []
            )
            seen = {canonical_payload_sha256(item) for item in merged_list}
            for item in value:
                fingerprint = canonical_payload_sha256(item)
                if fingerprint not in seen:
                    merged_list.append(item)
                    seen.add(fingerprint)
            merged_raw[key] = merged_list
        elif key not in merged_raw:
            merged_raw[key] = value

    manifest["manifest_id"] = "adco-control-plane-v2"
    manifest["schema_version"] = CONTROL_PLANE_SCHEMA_VERSION
    manifest.setdefault("created_at", now_iso())
    manifest.setdefault("source_hashes", source_hashes)
    manifest.setdefault("changes", list(changes))
    manifest["raw_legacy_evidence"] = merged_raw
    manifest["active_blockers"] = normalized_active
    manifest["blocker_history"] = [
        history_by_fingerprint[key] for key in sorted(history_by_fingerprint)
    ]
    manifest["attempts"] = attempts
    manifest["latest_attempt_id"] = attempt_id
    return manifest


def _read_project_file_bytes(
    project: Path,
    relative_path: str | Path,
    *,
    project_root_fd: int | None = None,
) -> bytes:
    file_fd: int | None = None
    directory_fds: list[int] = []
    try:
        file_fd, opened, directory_fds, parts = _open_project_relative_regular_file(
            project,
            relative_path,
            project_root_fd=project_root_fd,
        )
        chunks: list[bytes] = []
        while True:
            chunk = os.read(file_fd, PRIVATE_MARKER_SCAN_CHUNK_BYTES)
            if not chunk:
                break
            chunks.append(chunk)
        if not _project_file_binding_is_current(
            project,
            file_fd,
            opened,
            directory_fds,
            parts,
        ):
            raise ValueError(f"project file changed while reading: {relative_path}")
        return b"".join(chunks)
    finally:
        if file_fd is not None:
            os.close(file_fd)
        _close_project_fd_chain(directory_fds)


def _read_project_text(
    project: Path,
    relative_path: str | Path,
    *,
    project_root_fd: int | None = None,
) -> str:
    return _read_project_file_bytes(
        project,
        relative_path,
        project_root_fd=project_root_fd,
    ).decode("utf-8")


def _project_regular_file_exists(project: Path, relative_path: str | Path) -> bool:
    directory_fds: list[int] = []
    try:
        directory_fds, parts = _open_project_parent_chain(project, relative_path)
        try:
            visible = os.stat(
                parts[-1],
                dir_fd=directory_fds[-1],
                follow_symlinks=False,
            )
        except FileNotFoundError:
            return False
        if not stat.S_ISREG(visible.st_mode):
            raise ValueError(f"project file is not regular: {relative_path}")
        return True
    except FileNotFoundError:
        return False
    finally:
        _close_project_fd_chain(directory_fds)


def _project_file_sha256(project: Path, relative_path: str | Path) -> str:
    file_fd: int | None = None
    directory_fds: list[int] = []
    try:
        file_fd, opened, directory_fds, parts = _open_project_relative_regular_file(
            project,
            relative_path,
        )
        digest, _, _ = _hash_open_fd(
            file_fd,
            limit=max(opened.st_size, 0),
        )
        if not _project_file_binding_is_current(
            project,
            file_fd,
            opened,
            directory_fds,
            parts,
        ):
            raise ValueError(f"project file changed while hashing: {relative_path}")
        return digest
    finally:
        if file_fd is not None:
            os.close(file_fd)
        _close_project_fd_chain(directory_fds)


def _write_project_file_bytes(
    project: Path,
    relative_path: str | Path,
    data: bytes,
) -> None:
    directory_fds: list[int] = []
    temp_name = ""
    try:
        directory_fds, parts = _open_project_parent_chain(
            project,
            relative_path,
            create=True,
        )
        parent_fd = directory_fds[-1]
        try:
            current = os.stat(
                parts[-1],
                dir_fd=parent_fd,
                follow_symlinks=False,
            )
        except FileNotFoundError:
            current = None
        if current is not None and not stat.S_ISREG(current.st_mode):
            raise ValueError(f"unsafe migration target: {relative_path}")

        for _ in range(32):
            candidate = ".adco-migration-" + os.urandom(12).hex()
            try:
                temp_fd = os.open(
                    candidate,
                    os.O_WRONLY
                    | os.O_CREAT
                    | os.O_EXCL
                    | getattr(os, "O_NOFOLLOW", 0),
                    mode=0o600,
                    dir_fd=parent_fd,
                )
                temp_name = candidate
                break
            except FileExistsError:
                continue
        else:
            raise ValueError("cannot allocate private migration temp file")
        try:
            view = memoryview(data)
            while view:
                written = os.write(temp_fd, view)
                if written <= 0:
                    raise OSError("migration temp-file write made no progress")
                view = view[written:]
            os.fchmod(temp_fd, 0o644)
            os.fsync(temp_fd)
        finally:
            os.close(temp_fd)

        if not _project_parent_chain_is_current(project, directory_fds, parts):
            raise ValueError(f"migration managed path changed: {relative_path}")
        try:
            final_current = os.stat(
                parts[-1],
                dir_fd=parent_fd,
                follow_symlinks=False,
            )
        except FileNotFoundError:
            final_current = None
        if final_current is not None and not stat.S_ISREG(final_current.st_mode):
            raise ValueError(f"unsafe migration target: {relative_path}")
        os.replace(
            temp_name,
            parts[-1],
            src_dir_fd=parent_fd,
            dst_dir_fd=parent_fd,
        )
        temp_name = ""
        os.fsync(parent_fd)
    finally:
        if temp_name and directory_fds:
            try:
                os.unlink(temp_name, dir_fd=directory_fds[-1])
            except FileNotFoundError:
                pass
        _close_project_fd_chain(directory_fds)


def _write_project_text(
    project: Path,
    relative_path: str | Path,
    content: str,
) -> None:
    _write_project_file_bytes(
        project,
        relative_path,
        (content.rstrip() + "\n").encode("utf-8"),
    )


def _read_project_csv_rows(
    project: Path,
    relative_path: str | Path,
    *,
    project_root_fd: int | None = None,
) -> tuple[list[str], list[dict[str, str]]]:
    try:
        text = _read_project_text(
            project,
            relative_path,
            project_root_fd=project_root_fd,
        )
    except FileNotFoundError:
        return [], []
    reader = csv.DictReader(io.StringIO(text, newline=""))
    return list(reader.fieldnames or []), list(reader)


def _write_project_csv_rows(
    project: Path,
    relative_path: str | Path,
    fieldnames: list[str],
    rows: Iterable[dict[str, str]],
) -> None:
    output = io.StringIO(newline="")
    writer = csv.DictWriter(output, fieldnames=fieldnames, lineterminator="\n")
    writer.writeheader()
    for row in rows:
        writer.writerow({key: row.get(key, "") for key in fieldnames})
    _write_project_file_bytes(project, relative_path, output.getvalue().encode("utf-8"))


def _ensure_project_csv(
    project: Path,
    relative_path: str | Path,
    required_fields: list[str],
) -> None:
    fieldnames, rows = _read_project_csv_rows(project, relative_path)
    if not fieldnames:
        _write_project_csv_rows(project, relative_path, required_fields, [])
        return
    missing = [field for field in required_fields if field not in fieldnames]
    if missing or csv_rows_need_normalization(fieldnames, rows):
        _write_project_csv_rows(
            project,
            relative_path,
            [*fieldnames, *missing],
            rows,
        )


def _validate_project_relative_parent(
    project: Path,
    relative_path: str | Path,
) -> None:
    directory_fds, parts = _open_project_parent_chain(project, relative_path)
    try:
        if not _project_parent_chain_is_current(project, directory_fds, parts):
            raise ValueError(f"migration managed path changed: {relative_path}")
    finally:
        _close_project_fd_chain(directory_fds)


def _snapshot_open_project_directory(
    directory_fd: int,
    *,
    prefix: str,
    snapshot: dict[str, str],
) -> None:
    with os.scandir(directory_fd) as iterator:
        entries = sorted(iterator, key=lambda entry: entry.name)
    for entry in entries:
        name = entry.name
        rel = f"{prefix}/{name}" if prefix else name
        if rel == ".git" or rel.startswith(".git/"):
            continue
        visible = os.stat(name, dir_fd=directory_fd, follow_symlinks=False)
        if stat.S_ISLNK(visible.st_mode):
            snapshot[rel] = "symlink:" + os.readlink(name, dir_fd=directory_fd)
            continue
        if stat.S_ISDIR(visible.st_mode):
            child_fd = os.open(
                name,
                os.O_RDONLY
                | getattr(os, "O_DIRECTORY", 0)
                | getattr(os, "O_NOFOLLOW", 0),
                dir_fd=directory_fd,
            )
            try:
                opened = os.fstat(child_fd)
                if (visible.st_dev, visible.st_ino) != (opened.st_dev, opened.st_ino):
                    raise ValueError(f"migration directory changed during snapshot: {rel}")
                _snapshot_open_project_directory(
                    child_fd,
                    prefix=rel,
                    snapshot=snapshot,
                )
            finally:
                os.close(child_fd)
            continue
        if stat.S_ISREG(visible.st_mode):
            file_fd = os.open(
                name,
                os.O_RDONLY | getattr(os, "O_NOFOLLOW", 0),
                dir_fd=directory_fd,
            )
            try:
                opened = os.fstat(file_fd)
                if (visible.st_dev, visible.st_ino) != (opened.st_dev, opened.st_ino):
                    raise ValueError(f"migration file changed during snapshot: {rel}")
                digest = hashlib.sha256()
                while True:
                    chunk = os.read(file_fd, PRIVATE_MARKER_SCAN_CHUNK_BYTES)
                    if not chunk:
                        break
                    digest.update(chunk)
                final_opened = os.fstat(file_fd)
                final_visible = os.stat(
                    name,
                    dir_fd=directory_fd,
                    follow_symlinks=False,
                )
                if (
                    (final_visible.st_dev, final_visible.st_ino)
                    != (opened.st_dev, opened.st_ino)
                    or (
                        final_opened.st_size,
                        final_opened.st_mtime_ns,
                        final_opened.st_ctime_ns,
                    )
                    != (opened.st_size, opened.st_mtime_ns, opened.st_ctime_ns)
                ):
                    raise ValueError(f"migration file changed during snapshot: {rel}")
                snapshot[rel] = digest.hexdigest()
            finally:
                os.close(file_fd)


def _migration_input_snapshot(project: Path) -> dict[str, str]:
    try:
        root_fd = _open_project_root_fd(project)
    except (FileNotFoundError, ValueError):
        if not project.exists() and not project.is_symlink():
            return {}
        raise
    try:
        snapshot: dict[str, str] = {}
        _snapshot_open_project_directory(root_fd, prefix="", snapshot=snapshot)
        return snapshot
    finally:
        os.close(root_fd)


def _migration_plan_sha256(
    *,
    input_snapshot: dict[str, str],
    changes: list[str],
    blockers: list[dict[str, object]],
    desired_truth: str,
    desired_project_yml: str,
    desired_schema: dict[str, object],
) -> str:
    return canonical_payload_sha256(
        {
            "protocol_id": "adco.control-plane-migration-plan",
            "version": "1.0",
            "input_snapshot": input_snapshot,
            "changes": changes,
            "blockers": blockers,
            "desired_truth_sha256": hashlib.sha256(
                desired_truth.encode("utf-8")
            ).hexdigest(),
            "desired_project_yml_sha256": hashlib.sha256(
                desired_project_yml.encode("utf-8")
            ).hexdigest(),
            "desired_schema": desired_schema,
        }
    )


def _migration_project_identity(
    project: Path,
) -> tuple[Path, str, tuple[int, int] | None]:
    lexical = project.expanduser().absolute()
    try:
        canonical = lexical.resolve(
            strict=lexical.exists() or lexical.is_symlink()
        )
    except OSError as exc:
        raise ValueError("migration project path cannot be resolved safely") from exc
    expected_inode: tuple[int, int] | None = None
    if canonical.exists():
        try:
            current = os.stat(canonical, follow_symlinks=False)
        except OSError as exc:
            raise ValueError("migration project changed while resolving") from exc
        if not stat.S_ISDIR(current.st_mode):
            raise ValueError("migration project must be a directory")
        expected_inode = (current.st_dev, current.st_ino)
    # Keep one lock identity across the nonexistent -> created transition while
    # still collapsing lexical and symlink aliases onto the resolved path.
    identity = {"canonical_path": str(canonical)}
    lock_key = canonical_payload_sha256(identity)
    return canonical, lock_key, expected_inode


def _migration_lock_directory_path() -> Path:
    return Path(tempfile.gettempdir()).resolve() / (
        f"adco-migration-locks-{os.getuid()}"
    )


def _open_migration_lock(lock_key: str) -> int:
    no_follow = getattr(os, "O_NOFOLLOW", None)
    directory_flag = getattr(os, "O_DIRECTORY", None)
    if no_follow is None or directory_flag is None:
        raise ValueError("safe migration locks require O_NOFOLLOW and O_DIRECTORY")
    base = Path(tempfile.gettempdir()).resolve()
    lock_dir_name = _migration_lock_directory_path().name
    base_fd: int | None = None
    lock_dir_fd: int | None = None
    lock_fd: int | None = None
    try:
        base_fd = os.open(base, os.O_RDONLY | directory_flag | no_follow)
        try:
            os.mkdir(lock_dir_name, mode=0o700, dir_fd=base_fd)
        except FileExistsError:
            pass
        lock_dir_fd = os.open(
            lock_dir_name,
            os.O_RDONLY | directory_flag | no_follow,
            dir_fd=base_fd,
        )
        visible_dir = os.stat(
            lock_dir_name,
            dir_fd=base_fd,
            follow_symlinks=False,
        )
        opened_dir = os.fstat(lock_dir_fd)
        if (
            not stat.S_ISDIR(visible_dir.st_mode)
            or (visible_dir.st_dev, visible_dir.st_ino)
            != (opened_dir.st_dev, opened_dir.st_ino)
            or opened_dir.st_uid != os.getuid()
        ):
            raise ValueError("migration lock directory is unsafe")
        os.fchmod(lock_dir_fd, 0o700)

        lock_name = f"{lock_key}.lock"
        lock_fd = os.open(
            lock_name,
            os.O_RDWR | os.O_CREAT | no_follow,
            mode=0o600,
            dir_fd=lock_dir_fd,
        )
        visible_lock = os.stat(
            lock_name,
            dir_fd=lock_dir_fd,
            follow_symlinks=False,
        )
        opened_lock = os.fstat(lock_fd)
        if (
            not stat.S_ISREG(visible_lock.st_mode)
            or (visible_lock.st_dev, visible_lock.st_ino)
            != (opened_lock.st_dev, opened_lock.st_ino)
            or opened_lock.st_uid != os.getuid()
            or opened_lock.st_nlink != 1
        ):
            raise ValueError("migration lock file is unsafe")
        os.fchmod(lock_fd, 0o600)
        result = lock_fd
        lock_fd = None
        return result
    except OSError as exc:
        raise ValueError("cannot safely open migration lock") from exc
    finally:
        if lock_fd is not None:
            os.close(lock_fd)
        if lock_dir_fd is not None:
            os.close(lock_dir_fd)
        if base_fd is not None:
            os.close(base_fd)


def _verify_migration_project_identity(
    project: Path,
    expected_inode: tuple[int, int] | None,
) -> None:
    if expected_inode is None:
        if project.is_symlink():
            raise ValueError("migration project changed to a symlink before execution")
        return
    try:
        current = os.stat(project, follow_symlinks=False)
    except OSError as exc:
        raise ValueError("migration project changed before execution") from exc
    if (
        not stat.S_ISDIR(current.st_mode)
        or (current.st_dev, current.st_ino) != expected_inode
    ):
        raise ValueError("migration project changed before execution")


def _open_migration_project_dir(
    project: Path,
    expected_inode: tuple[int, int],
) -> int:
    no_follow = getattr(os, "O_NOFOLLOW", None)
    directory_flag = getattr(os, "O_DIRECTORY", None)
    if no_follow is None or directory_flag is None:
        raise ValueError("safe migration requires O_NOFOLLOW and O_DIRECTORY")
    try:
        fd = os.open(project, os.O_RDONLY | directory_flag | no_follow)
    except OSError as exc:
        raise ValueError("cannot safely open migration project") from exc
    try:
        opened = os.fstat(fd)
        visible = os.stat(project, follow_symlinks=False)
        if (
            not stat.S_ISDIR(visible.st_mode)
            or (opened.st_dev, opened.st_ino) != expected_inode
            or (visible.st_dev, visible.st_ino) != expected_inode
        ):
            raise ValueError("migration project changed while opening")
        return fd
    except BaseException:
        os.close(fd)
        raise


def _migration_project_binding_is_current(
    project: Path,
    project_fd: int,
    expected_inode: tuple[int, int],
) -> bool:
    try:
        opened = os.fstat(project_fd)
        visible = os.stat(project, follow_symlinks=False)
    except OSError:
        return False
    return bool(
        stat.S_ISDIR(visible.st_mode)
        and (opened.st_dev, opened.st_ino) == expected_inode
        and (visible.st_dev, visible.st_ino) == expected_inode
    )


def _normalize_migration_result_paths(
    result: dict[str, object],
    project: Path,
) -> dict[str, object]:
    result["project"] = str(project)
    if result.get("migration_manifest"):
        result["migration_manifest"] = str(
            project / CONTROL_PLANE_MIGRATION_MANIFEST_REL
        )
    return result


def migrate_control_plane(project: Path, *, dry_run: bool = False) -> dict[str, object]:
    canonical_project, lock_key, expected_inode = _migration_project_identity(project)
    lock_fd = _open_migration_lock(lock_key)
    with os.fdopen(lock_fd, "a+", encoding="utf-8") as lock_handle:
        fcntl.flock(lock_handle.fileno(), fcntl.LOCK_EX)
        try:
            _verify_migration_project_identity(canonical_project, expected_inode)
            if expected_inode is None:
                if not dry_run:
                    raise ValueError("migration project must already exist")
                return _migrate_control_plane_locked(
                    canonical_project,
                    dry_run=True,
                )

            project_fd = _open_migration_project_dir(
                canonical_project,
                expected_inode,
            )
            original_cwd_fd: int | None = None
            try:
                with MIGRATION_CWD_LOCK:
                    original_cwd_fd = os.open(
                        ".",
                        os.O_RDONLY | getattr(os, "O_DIRECTORY", 0),
                    )
                    try:
                        os.fchdir(project_fd)
                        result = _migrate_control_plane_locked(
                            Path("."),
                            dry_run=dry_run,
                            pre_write_guard=lambda: _migration_project_binding_is_current(
                                canonical_project,
                                project_fd,
                                expected_inode,
                            ),
                        )
                    finally:
                        os.fchdir(original_cwd_fd)
                        os.close(original_cwd_fd)
                        original_cwd_fd = None
                if not _migration_project_binding_is_current(
                    canonical_project,
                    project_fd,
                    expected_inode,
                ):
                    raise ValueError("migration project changed during execution")
                return _normalize_migration_result_paths(result, canonical_project)
            finally:
                if original_cwd_fd is not None:
                    os.close(original_cwd_fd)
                os.close(project_fd)
        finally:
            fcntl.flock(lock_handle.fileno(), fcntl.LOCK_UN)


def _migrate_control_plane_locked(
    project: Path,
    *,
    dry_run: bool,
    pre_write_guard: Callable[[], bool] | None = None,
) -> dict[str, object]:
    input_snapshot = _migration_input_snapshot(project)
    input_snapshot_sha256 = canonical_payload_sha256(input_snapshot)
    targets = [
        ("AD-creative/client_review/client_outline.csv", CLIENT_OUTLINE_FIELDS),
        ("AD-creative/visual_assets/asset_current_manifest.csv", ASSET_CURRENT_FIELDS),
        ("AD-creative/visual_assets/asset_authorizations.csv", ASSET_AUTHORIZATION_FIELDS),
        ("AD-creative/orchestrator/final_delivery_lock.csv", FINAL_DELIVERY_LOCK_FIELDS),
        ("AD-creative/orchestrator/artifact_index.csv", ARTIFACT_INDEX_FIELDS),
        ("AD-creative/orchestrator/gate_log.csv", GATE_LOG_FIELDS),
        (
            "AD-creative/orchestrator/specialist_exchange/exchange_index.csv",
            SPECIALIST_EXCHANGE_INDEX_FIELDS,
        ),
        ("AD-creative/orchestrator/agency/specialist_preflight.csv", ["preflight_id", "work_id", "requested_skill", "skill_path", "rules_read", "derived_gates", "status", "blocked_reason", "created_at"]),
        ("AD-creative/orchestrator/agency/asset_preflight.csv", ["preflight_id", "work_id", "source_scope", "local_manifest_checked", "browser_checked", "browser_tool", "download_method", "imported_asset_ids", "replacement_generation_allowed", "status", "blocked_reason", "created_at"]),
        ("AD-creative/orchestrator/agency/skill_scout.csv", ["scout_id", "work_id", "skill_name", "skill_path", "match_reason", "selected", "status", "created_at", "notes"]),
        ("AD-creative/orchestrator/agency/agent_scout.csv", ["scout_id", "work_id", "agent_name", "agent_path", "match_reason", "selected", "status", "created_at", "notes"]),
        ("AD-creative/orchestrator/thread_registry.csv", THREADOPS_REGISTRY_FIELDS),
        ("AD-creative/orchestrator/agent_runs.csv", THREADOPS_AGENT_RUN_FIELDS),
    ]
    static_files = {
        "AD-creative/orchestrator/agency/maintenance_heartbeat.md": "# Maintenance Heartbeat\n\nstatus: active\nvisibility: internal_only\n\nRecord repeated failures, cleanup freezes, and skill-hardening follow-ups here.\n",
        "AD-creative/orchestrator/agency/self_improvement_log.md": "# Self Improvement Log\n\nstatus: active\nvisibility: internal_only\n\nUse this only for verified project failures that resulted in reusable ADCO rule or tool changes.\n",
    }
    changes: list[str] = []
    warnings: list[str] = []
    blockers: list[dict[str, object]] = []
    pre_csv: dict[str, tuple[list[str], list[dict[str, str]]]] = {
        rel_path: _read_project_csv_rows(project, rel_path)
        for rel_path, _ in targets
    }
    source_hashes = {
        rel_path: _project_file_sha256(project, rel_path)
        for rel_path, _ in targets
        if _project_regular_file_exists(project, rel_path)
    }
    baseline_source_hashes = {
        rel_path: _project_file_sha256(project, rel_path)
        for rel_path in LEGACY_BASELINE_SOURCE_RELS
        if _project_regular_file_exists(project, rel_path)
    }
    source_hashes.update(baseline_source_hashes)
    schema_rel = CONTROL_PLANE_SCHEMA_REL.as_posix()
    schema_path = project / schema_rel
    schema_payload: dict[str, object] = {}
    if _project_regular_file_exists(project, schema_rel):
        try:
            schema_payload = json.loads(_read_project_text(project, schema_rel))
        except json.JSONDecodeError:
            schema_payload = {}
    project_yml_rel = "AD-creative/orchestrator/project.yml"
    project_yml = project / project_yml_rel
    project_yml_text = (
        _read_project_text(project, project_yml_rel)
        if _project_regular_file_exists(project, project_yml_rel)
        else ""
    )
    pre_migration_schema_v2 = bool(
        schema_payload.get("schema_version") == CONTROL_PLANE_SCHEMA_VERSION
        and re.search(
            rf'(?m)^  schema_version:[ \t]*["\']?{re.escape(CONTROL_PLANE_SCHEMA_VERSION)}["\']?[ \t]*$',
            project_yml_text,
        )
    )
    legacy_messages = (
        sorted(
            {
                message
                for message in string_validation_errors(project)
                if legacy_baseline_message_allowed(message)
            }
        )
        if not pre_migration_schema_v2
        else []
    )
    legacy_string_error_snapshot: dict[str, object] = {}
    legacy_string_error_baseline: list[dict[str, str]] = []
    if legacy_messages:
        legacy_string_error_snapshot = {
            "captured_from_schema_state": "pre_v2",
            "source_hashes": baseline_source_hashes,
            "messages": legacy_messages,
        }
        evidence_sha256 = canonical_payload_sha256(
            legacy_string_error_snapshot
        )
        source_hashes_sha256 = canonical_payload_sha256(
            baseline_source_hashes
        )
        legacy_string_error_baseline = [
            {
                "fingerprint": canonical_payload_sha256(
                    {
                        "message": message,
                        "evidence_sha256": evidence_sha256,
                        "source_hashes_sha256": source_hashes_sha256,
                    }
                ),
                "message": message,
                "evidence_sha256": evidence_sha256,
            }
            for message in legacy_messages
        ]
    for rel_path, fields in targets:
        if not _project_regular_file_exists(
            project,
            rel_path,
        ) or not _read_project_text(project, rel_path).strip():
            changes.append(f"create_csv:{rel_path}")
        else:
            existing, rows = _read_project_csv_rows(project, rel_path)
            missing = [field for field in fields if field not in existing]
            if missing:
                changes.append(f"add_csv_fields:{rel_path}:{','.join(missing)}")
            elif csv_rows_need_normalization(existing, rows):
                changes.append(f"normalize_csv_rows:{rel_path}")
    for rel_path in static_files:
        if not _project_regular_file_exists(project, rel_path):
            changes.append(f"create_file:{rel_path}")

    artifact_rel = "AD-creative/orchestrator/artifact_index.csv"
    raw_artifact_fields, raw_artifact_rows = pre_csv[artifact_rel]
    preview_artifacts, artifact_evidence = migrate_artifact_lifecycle_rows(raw_artifact_rows)
    if preview_artifacts != raw_artifact_rows:
        changes.append(f"migrate_artifact_lifecycle:{artifact_rel}:{len(artifact_evidence)}")

    registry_rel = "AD-creative/orchestrator/thread_registry.csv"
    raw_registry_fields, raw_registry_rows = pre_csv[registry_rel]
    legacy_thread_schema = bool(raw_registry_rows) and (
        not pre_migration_schema_v2
        or any(field not in raw_registry_fields for field in THREADOPS_REGISTRY_FIELDS[:-4])
    )
    unclassified_thread_indices = {
        index
        for index, row in enumerate(raw_registry_rows)
        if (row.get("schema_state") or "").strip()
        not in {"current", "legacy_quarantined"}
    }
    thread_rows_to_quarantine = (
        set(range(len(raw_registry_rows)))
        if legacy_thread_schema
        else unclassified_thread_indices
    )
    thread_rows_need_state = bool(thread_rows_to_quarantine)
    thread_evidence_by_sha = {
        canonical_row_sha256(raw_registry_rows[index]): {
            "row_number": index + 2,
            "evidence_sha256": canonical_row_sha256(raw_registry_rows[index]),
            "raw": raw_registry_rows[index],
        }
        for index in sorted(thread_rows_to_quarantine)
    }
    if thread_rows_need_state:
        changes.append(
            f"quarantine_threadops_rows:{registry_rel}:{len(thread_rows_to_quarantine)}"
        )

    truth_rel = "AD-creative/orchestrator/current_truth.md"
    truth_path = project / truth_rel
    truth_text = (
        _read_project_text(project, truth_rel)
        if _project_regular_file_exists(project, truth_rel)
        else ""
    )
    truth_sections = markdown_named_sections(truth_text, "Current Version Truth")
    _, version_rows = pre_csv.get(
        "AD-creative/orchestrator/version_map.csv",
        _read_project_csv_rows(
            project,
            "AD-creative/orchestrator/version_map.csv",
        ),
    )
    raw_package_sections = markdown_named_sections(truth_text, "Current Package")
    if current_version_truth_binding_valid(
        truth_text, raw_artifact_rows, version_rows
    ):
        backfill: dict[str, str] = {}
        package_blockers: list[dict[str, object]] = []
    else:
        backfill, package_blockers, raw_package_sections = legacy_current_package_backfill(
            project, truth_text, raw_artifact_rows, version_rows
        )
    blockers.extend(package_blockers)
    desired_truth = truth_text
    if len(truth_sections) > 1:
        blockers.append(
            {
                "severity": "P0",
                "scope": "current",
                "code": "duplicate_current_version_truth",
                "message": "current_truth has duplicate Current Version Truth sections; migration refused to choose",
            }
        )
    elif backfill and not package_blockers:
        existing_values = (
            section_key_values(truth_sections[0]) if len(truth_sections) == 1 else {}
        )
        merged_values = {
            key: first_nonempty(existing_values.get(key), backfill.get(key), default="")
            for key in CURRENT_VERSION_TRUTH_KEYS
        }
        desired_truth = replace_current_version_truth(truth_text, merged_values)
        if desired_truth != truth_text:
            changes.append("backfill_current_version_truth_from_legacy_current_package")
    elif len(truth_sections) == 0 and not package_blockers:
        desired_truth = replace_current_version_truth(truth_text, {})
        changes.append("add_section:AD-creative/orchestrator/current_truth.md:Current Version Truth")
    elif len(truth_sections) == 1:
        desired_truth, missing_truth_keys = normalize_current_version_truth_section(truth_text)
        if missing_truth_keys:
            changes.append(
                "add_current_truth_keys:AD-creative/orchestrator/current_truth.md:"
                + ",".join(missing_truth_keys)
            )

    if _project_regular_file_exists(project, schema_rel) and not schema_payload:
        blockers.append(
            {
                "severity": "P0",
                "scope": "current",
                "code": "malformed_control_plane_schema",
                "message": f"cannot parse {CONTROL_PLANE_SCHEMA_REL}",
            }
        )
    desired_schema = {
        "schema_id": "adco.control-plane",
        "schema_version": CONTROL_PLANE_SCHEMA_VERSION,
        "migration_manifest": CONTROL_PLANE_MIGRATION_MANIFEST_REL.as_posix(),
        "lifecycle_values": sorted(ARTIFACT_LIFECYCLE_VALUES),
    }
    schema_current = schema_payload.get("schema_version") == CONTROL_PLANE_SCHEMA_VERSION
    if not schema_current:
        changes.append(f"set_control_plane_schema:{CONTROL_PLANE_SCHEMA_VERSION}")
    elif schema_payload != desired_schema:
        changes.append("normalize_control_plane_schema")
    desired_project_yml = project_yml_with_schema_version(project_yml_text)
    if desired_project_yml != project_yml_text:
        changes.append(f"set_project_schema_version:{CONTROL_PLANE_SCHEMA_VERSION}")

    manifest_rel = CONTROL_PLANE_MIGRATION_MANIFEST_REL.as_posix()
    manifest_path = project / manifest_rel
    existing_manifest: dict[str, object] | None = None
    manifest_malformed = False
    manifest_exists = _project_regular_file_exists(project, manifest_rel)
    if manifest_exists:
        try:
            loaded_manifest = json.loads(_read_project_text(project, manifest_rel))
            if not isinstance(loaded_manifest, dict):
                raise ValueError("migration manifest must be a JSON object")
            existing_manifest = loaded_manifest
        except (json.JSONDecodeError, ValueError):
            manifest_malformed = True
            blockers.append(
                {
                    "severity": "P0",
                    "scope": "current",
                    "code": "malformed_migration_manifest",
                    "message": f"cannot safely update {CONTROL_PLANE_MIGRATION_MANIFEST_REL}",
                }
            )
    blockers = normalized_migration_blockers(blockers)
    needs_manifest = bool(
        manifest_exists
        or not schema_current
        or artifact_evidence
        or legacy_thread_schema
        or thread_rows_to_quarantine
        or raw_package_sections
        or legacy_string_error_baseline
        or blockers
    )
    if needs_manifest and not manifest_exists:
        changes.append(f"create_migration_manifest:{CONTROL_PLANE_MIGRATION_MANIFEST_REL}")

    plan_sha256 = _migration_plan_sha256(
        input_snapshot=input_snapshot,
        changes=changes,
        blockers=blockers,
        desired_truth=desired_truth,
        desired_project_yml=desired_project_yml,
        desired_schema=desired_schema,
    )

    if not dry_run and blockers:
        return {
            "project": str(project),
            "dry_run": False,
            "applied": False,
            "blocked_before_write": True,
            "changes": changes,
            "warnings": warnings,
            "blockers": blockers,
            "schema_version": CONTROL_PLANE_SCHEMA_VERSION,
            "migration_manifest": str(manifest_path) if manifest_exists else "",
            "input_snapshot_sha256": input_snapshot_sha256,
            "plan_sha256": plan_sha256,
        }

    if not dry_run:
        if pre_write_guard is not None and not pre_write_guard():
            root_blocker = {
                "severity": "P0",
                "scope": "current",
                "code": "migration_project_changed",
                "message": "migration project root changed before first write",
            }
            return {
                "project": str(project),
                "dry_run": False,
                "applied": False,
                "blocked_before_write": True,
                "changes": changes,
                "warnings": warnings,
                "blockers": [root_blocker],
                "schema_version": CONTROL_PLANE_SCHEMA_VERSION,
                "migration_manifest": "",
                "input_snapshot_sha256": input_snapshot_sha256,
                "plan_sha256": plan_sha256,
            }
        current_snapshot = _migration_input_snapshot(project)
        if current_snapshot != input_snapshot:
            changed_paths = sorted(
                {
                    *input_snapshot.keys(),
                    *current_snapshot.keys(),
                }
                - {
                    key
                    for key in {*input_snapshot.keys(), *current_snapshot.keys()}
                    if input_snapshot.get(key) == current_snapshot.get(key)
                }
            )
            concurrent_blocker = {
                "severity": "P0",
                "scope": "current",
                "code": "migration_inputs_changed",
                "message": (
                    "migration inputs changed before first write: "
                    + ";".join(changed_paths[:20])
                ),
            }
            return {
                "project": str(project),
                "dry_run": False,
                "applied": False,
                "blocked_before_write": True,
                "changes": changes,
                "warnings": warnings,
                "blockers": [concurrent_blocker],
                "schema_version": CONTROL_PLANE_SCHEMA_VERSION,
                "migration_manifest": str(manifest_path) if manifest_exists else "",
                "input_snapshot_sha256": input_snapshot_sha256,
                "plan_sha256": plan_sha256,
            }
        template_created_files, _ = ensure_delivery_project(project)
        surface_match = re.search(
            r'(?m)^  surface:\s*["\']?([^"\'\n]+)',
            project_yml_text,
        )
        template_materialization_changed = bool(
            template_created_files
            or surface_match is None
            or surface_match.group(1).strip() != DELIVERY_SURFACE
        )
        managed_parent_paths = [
            *(rel_path for rel_path, _ in targets),
            *static_files,
            truth_rel,
            project_yml_rel,
            schema_rel,
        ]
        try:
            for rel_path in managed_parent_paths:
                _validate_project_relative_parent(project, rel_path)
        except (OSError, ValueError) as exc:
            managed_path_blocker = {
                "severity": "P0",
                "scope": "current",
                "code": "migration_managed_path_changed",
                "message": f"migration managed path changed before apply: {exc}",
            }
            return {
                "project": str(project),
                "dry_run": False,
                "applied": False,
                "blocked_before_write": not template_materialization_changed,
                "partial_apply": template_materialization_changed,
                "changes": changes,
                "warnings": warnings,
                "blockers": [managed_path_blocker],
                "schema_version": CONTROL_PLANE_SCHEMA_VERSION,
                "migration_manifest": "",
                "input_snapshot_sha256": input_snapshot_sha256,
                "plan_sha256": plan_sha256,
            }
        for rel_path, fields in targets:
            _ensure_project_csv(project, rel_path, fields)
        gate_rel = "AD-creative/orchestrator/gate_log.csv"
        gate_fields, gate_rows = _read_project_csv_rows(project, gate_rel)
        normalized_gate_fields, normalized_gate_rows = normalized_gate_log_data(
            gate_fields,
            gate_rows,
        )
        if gate_fields != normalized_gate_fields or gate_rows != normalized_gate_rows:
            _write_project_csv_rows(
                project,
                gate_rel,
                normalized_gate_fields,
                normalized_gate_rows,
            )
        for rel_path, content in static_files.items():
            if not _project_regular_file_exists(project, rel_path):
                _write_project_text(project, rel_path, content)

        artifact_fields, artifact_rows = _read_project_csv_rows(
            project,
            artifact_rel,
        )
        migrated_artifacts, _ = migrate_artifact_lifecycle_rows(artifact_rows)
        if migrated_artifacts != artifact_rows:
            _write_project_csv_rows(
                project,
                artifact_rel,
                artifact_fields,
                migrated_artifacts,
            )

        registry_fields, registry_rows = _read_project_csv_rows(
            project,
            registry_rel,
        )
        if registry_rows:
            normalized_registry: list[dict[str, str]] = []
            for index, row in enumerate(registry_rows):
                normalized = dict(row)
                if index in thread_rows_to_quarantine:
                    raw_row = raw_registry_rows[index] if index < len(raw_registry_rows) else row
                    evidence_sha = canonical_row_sha256(raw_row)
                    normalized["schema_state"] = "legacy_quarantined"
                    normalized["legacy_evidence_sha256"] = evidence_sha
                    normalized["legacy_quarantine_reason"] = (
                        "pre_v2_threadops_row_missing_proof_columns"
                        if legacy_thread_schema
                        else "v2_unclassified_threadops_row_without_writer_proof"
                    )
                    normalized["legacy_raw_ref"] = (
                        f"{CONTROL_PLANE_MIGRATION_MANIFEST_REL.as_posix()}"
                        f"#/raw_legacy_evidence/thread_rows_by_sha/{evidence_sha}"
                    )
                normalized_registry.append(normalized)
            if normalized_registry != registry_rows:
                _write_project_csv_rows(
                    project,
                    registry_rel,
                    registry_fields,
                    normalized_registry,
                )

        truth_exists = _project_regular_file_exists(project, truth_rel)
        if truth_exists and desired_truth != truth_text and not blockers:
            _write_project_text(project, truth_rel, desired_truth)
        elif not truth_exists and not blockers:
            _write_project_text(
                project,
                truth_rel,
                desired_truth
                or replace_current_version_truth("# Current Truth\n", {}),
            )

        if desired_project_yml != project_yml_text:
            _write_project_text(project, project_yml_rel, desired_project_yml)
        if schema_payload != desired_schema:
            _write_project_text(
                project,
                schema_rel,
                json.dumps(
                    desired_schema,
                    ensure_ascii=False,
                    indent=2,
                    sort_keys=True,
                ),
            )

        if needs_manifest and not manifest_malformed:
            evidence = {
                "artifact_rows": artifact_evidence,
                "thread_rows": [
                    entry
                    for _, entry in sorted(thread_evidence_by_sha.items())
                ]
                if thread_rows_to_quarantine
                else [],
                "thread_rows_by_sha": thread_evidence_by_sha,
                "agent_run_rows": pre_csv[
                    "AD-creative/orchestrator/agent_runs.csv"
                ][1]
                if thread_rows_to_quarantine
                else [],
                "current_package_sections": raw_package_sections,
                "string_error_baseline": legacy_string_error_baseline,
                "string_error_snapshot": legacy_string_error_snapshot,
            }
            observed_rel_paths = sorted(
                {
                    artifact_rel,
                    "AD-creative/orchestrator/current_truth.md",
                    "AD-creative/orchestrator/version_map.csv",
                    "AD-creative/orchestrator/project.yml",
                    CONTROL_PLANE_SCHEMA_REL.as_posix(),
                }
            )
            observed_hashes = {
                rel_path: (
                    _project_file_sha256(project, rel_path)
                    if _project_regular_file_exists(project, rel_path)
                    else "<missing>"
                )
                for rel_path in observed_rel_paths
            }
            updated_manifest = updated_migration_manifest(
                existing_manifest,
                source_hashes=source_hashes,
                changes=changes,
                active_blockers=blockers,
                raw_legacy_evidence=evidence,
                observed_hashes=observed_hashes,
            )
            if updated_manifest != existing_manifest:
                _write_project_text(
                    project,
                    manifest_rel,
                    json.dumps(
                        updated_manifest,
                        ensure_ascii=False,
                        indent=2,
                        sort_keys=True,
                    ),
                )
                manifest_exists = True
                if existing_manifest is not None:
                    changes.append("update_migration_manifest_active_state")
    return {
        "project": str(project),
        "dry_run": dry_run,
        "applied": not dry_run,
        "blocked_before_write": False,
        "changes": changes,
        "warnings": warnings,
        "blockers": blockers,
        "schema_version": CONTROL_PLANE_SCHEMA_VERSION,
        "migration_manifest": str(manifest_path) if manifest_exists or needs_manifest else "",
        "input_snapshot_sha256": input_snapshot_sha256,
        "plan_sha256": plan_sha256,
    }


def row_text(row: dict[str, str], keys: Iterable[str] | None = None) -> str:
    selected = keys or row.keys()
    return " ".join(row.get(key, "") for key in selected)


def agency_audit_report(project: Path) -> dict[str, object]:
    ensure_project(project)
    issues: list[dict[str, str]] = []
    stats: dict[str, int] = {}
    _, work_items = read_csv_rows(project / "AD-creative/orchestrator/work_items.csv")
    _, registry = read_csv_rows(project / "AD-creative/orchestrator/thread_registry.csv")
    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    _, gates = read_csv_rows(project / "AD-creative/orchestrator/gate_log.csv")
    stats.update(
        {
            "work_items": len(work_items),
            "thread_registry": len(registry),
            "artifacts": len(artifacts),
            "gates": len(gates),
        }
    )
    for work in work_items:
        work_id = work.get("work_id", "").strip()
        text = row_text(work, ("title", "objective", "gate_required", "output_artifacts"))
        named_ids = sorted(set(re.findall(r"\bWORK-\d+\b", text)))
        if work_id.startswith("WORK-GOAL") and named_ids:
            issues.append(
                {
                    "severity": "P1",
                    "code": "goal_work_id_hides_named_work_id",
                    "message": f"{work_id} text names {', '.join(named_ids)}",
                    "evidence": text[:180],
                    "fix": "Use explicit --work-id or split the goal shell from delivery work item.",
                }
            )
    for row in registry:
        thread_id = row.get("thread_id", "").strip()
        lifecycle = row.get("lifecycle_state", "").strip().lower()
        dispatch_status = row.get("dispatch_status", "").strip().lower()
        real_thread_id = row.get("real_thread_id", "").strip()
        if lifecycle in {"dispatched", "running", "returned", "reconciled"} or dispatch_status in {"dispatched", "running"}:
            missing = [
                name
                for name in ["real_thread_id", "title_verified_at", "dispatch_receipt_path", "dispatch_evidence"]
                if not row.get(name, "").strip()
            ]
            if thread_id.startswith("planned:") or not real_thread_id or missing:
                issues.append(
                    {
                        "severity": "P1",
                        "code": "missing_real_thread_dispatch_proof",
                        "message": f"{row.get('lane_id') or thread_id} lacks real Thread dispatch proof",
                        "evidence": f"thread_id={thread_id}; missing={','.join(missing)}",
                        "fix": "Create/reuse a real Codex Thread, read back title, then run dispatch-record.",
                    }
                )
    return {
        "status": "PASS" if not issues else "CHECK",
        "project": str(project),
        "p1": sum(1 for issue in issues if issue["severity"] == "P1"),
        "issues": issues,
        "stats": stats,
    }


def format_agency_issue(issue: dict[str, str]) -> str:
    return (
        f"{issue.get('severity', '')} {issue.get('code', '')}: "
        f"{issue.get('message', '')} | evidence={issue.get('evidence', '')} | fix={issue.get('fix', '')}"
    )


def write_specialist_preflight(
    project: Path,
    *,
    work_id: str,
    requested_skill: str,
    skill_path: str,
    rules_read: str,
    derived_gates: str,
    status: str,
    blocked_reason: str,
) -> str:
    migrate_control_plane(project)
    path = project / "AD-creative/orchestrator/agency/specialist_preflight.csv"
    fields, rows = read_csv_rows(path)
    preflight_id = next_id(rows, "preflight_id", "SPF")
    rows.append(
        {
            "preflight_id": preflight_id,
            "work_id": work_id,
            "requested_skill": requested_skill,
            "skill_path": skill_path,
            "rules_read": rules_read,
            "derived_gates": derived_gates,
            "status": status,
            "blocked_reason": blocked_reason,
            "created_at": now_iso(),
        }
    )
    write_csv_rows(path, fields, rows)
    return preflight_id


def read_json_object(path: Path, label: str) -> dict[str, object]:
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise ValueError(f"{label} is not valid JSON: {path}: {exc}") from exc
    if not isinstance(payload, dict):
        raise ValueError(f"{label} must be a JSON object: {path}")
    return payload


def write_json_object(path: Path, payload: dict[str, object]) -> None:
    write_text(path, json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True))


def contained_project_path(project: Path, raw_path: str, label: str) -> Path:
    candidate = Path(raw_path)
    if candidate.is_absolute():
        raise ValueError(f"{label} must be a project-relative path: {raw_path}")
    resolved = (project / candidate).resolve()
    try:
        resolved.relative_to(project.resolve())
    except ValueError as exc:
        raise ValueError(f"{label} escapes project scope: {raw_path}") from exc
    return resolved


def project_relative_path_has_symlink_component(project: Path, raw_path: str) -> bool:
    candidate = Path(raw_path.replace("\\", "/"))
    if candidate.is_absolute() or ".." in candidate.parts:
        return False
    current = project.resolve()
    for part in candidate.parts:
        if part in {"", "."}:
            continue
        current = current / part
        if current.is_symlink():
            return True
    return False


def contained_thread_scope_baseline_path(
    project: Path, raw_path: str, label: str
) -> Path:
    raw = (raw_path or "").strip()
    candidate = Path(raw.replace("\\", "/"))
    lexical = candidate.as_posix()
    if project_relative_path_has_symlink_component(project, lexical):
        raise ValueError(f"{label} contains a symlink component: {raw_path}")
    if not raw or candidate.is_absolute() or ".." in candidate.parts:
        raise ValueError(f"{label} must be a project-relative path: {raw_path}")
    parent = candidate.parent.as_posix()
    if parent not in {"", "."} and project_relative_path_has_symlink_component(
        project, parent
    ):
        raise ValueError(f"{label} parent contains a symlink component: {raw_path}")
    resolved = (project / candidate).resolve()
    try:
        resolved.relative_to(project.resolve())
    except ValueError as exc:
        raise ValueError(f"{label} escapes project scope: {raw_path}") from exc
    return resolved


def validate_specialist_token(value: str, label: str) -> str:
    token = value.strip()
    if (
        not SPECIALIST_ID_PATTERN.fullmatch(token)
        or token in {".", ".."}
        or ".." in token
    ):
        raise ValueError(f"{label} must be a safe protocol token: {value}")
    return token


def relative_path_is_within(relative_path: str, root: str) -> bool:
    path_parts = Path(relative_path).parts
    root_parts = Path(root.rstrip("/")).parts
    return len(path_parts) >= len(root_parts) and path_parts[: len(root_parts)] == root_parts


def require_regular_control_file(path: Path, label: str) -> None:
    if not path.is_file() or path.is_symlink() or path.stat().st_nlink != 1:
        raise ValueError(f"{label} must be a regular non-symlink file: {path}")


def validate_scope_manifest_path(raw_path: object, label: str) -> str:
    if not isinstance(raw_path, str):
        raise ValueError(f"{label} must be a string path")
    path = raw_path.strip()
    if not path or path != raw_path:
        raise ValueError(f"{label} must be a non-empty canonical POSIX path")
    if "\\" in path or re.match(r"^[A-Za-z]:", path):
        raise ValueError(f"{label} must be a project-relative POSIX path: {raw_path}")
    parts = path.split("/")
    if path.startswith("/") or any(part in {"", ".", ".."} for part in parts):
        raise ValueError(f"{label} must be a canonical project-relative path: {raw_path}")
    canonical = "/".join(parts)
    if canonical != path:
        raise ValueError(f"{label} is not canonical: {raw_path}")
    return canonical


def validate_scope_manifest_exclusion_roots(excluded_roots: object) -> list[str]:
    if not isinstance(excluded_roots, list):
        raise ValueError("thread scope baseline exclusions are invalid")
    roots: list[str] = []
    for raw_root in excluded_roots:
        root = validate_scope_manifest_path(raw_root, "thread scope exclusion root")
        if root == "AD-creative":
            raise ValueError("thread scope exclusion root cannot cover AD-creative")
        roots.append(root)
    return list(dict.fromkeys(roots))


def validate_scope_manifest_files(files: object) -> dict[str, str]:
    if not isinstance(files, dict):
        raise ValueError("thread scope baseline files are invalid")
    validated: dict[str, str] = {}
    for raw_path, raw_digest in files.items():
        path = validate_scope_manifest_path(raw_path, "thread scope manifest key")
        if not isinstance(raw_digest, str) or not re.fullmatch(
            r"[0-9a-fA-F]{64}", raw_digest
        ):
            raise ValueError(
                f"thread scope manifest digest must be a 64-character SHA-256 hex value: {raw_path}"
            )
        validated[path] = raw_digest.lower()
    return validated


def filtered_scope_manifest_files(
    files: object, *, excluded_roots: list[str]
) -> dict[str, str]:
    normalized_roots = validate_scope_manifest_exclusion_roots(excluded_roots)
    validated_files = validate_scope_manifest_files(files)
    return {
        path: digest
        for path, digest in validated_files.items()
        if not any(
            relative_path_is_within(path, root) for root in normalized_roots
        )
    }


def derive_thread_scope_baseline(
    project: Path,
    source_path: Path,
    source_payload: dict[str, object],
    *,
    suffix: str,
    added_exclusions: list[str],
    binding_kind: str,
    binding_ref: str,
) -> tuple[Path, dict[str, object]]:
    try:
        source_rel = source_path.relative_to(project).as_posix()
    except ValueError:
        source_rel = source_path.relative_to(project.resolve()).as_posix()
    source_path = contained_thread_scope_baseline_path(
        project, source_rel, "source thread scope baseline"
    )
    require_regular_control_file(source_path, "source thread scope baseline")
    target_path = source_path.with_name(
        f"{source_path.stem}{suffix}{source_path.suffix}"
    )
    target_rel = target_path.relative_to(project.resolve()).as_posix()
    target_path = contained_thread_scope_baseline_path(
        project, target_rel, "derived thread scope baseline"
    )
    if target_path.exists():
        raise ValueError(f"derived thread scope baseline already exists: {target_path}")
    source_sha256 = file_sha256(source_path)
    existing_exclusions = source_payload.get("excluded_roots")
    if not isinstance(existing_exclusions, list):
        raise ValueError("thread scope baseline exclusions are invalid")
    exclusions = list(
        dict.fromkeys(
            [str(item) for item in existing_exclusions]
            + [str(item) for item in added_exclusions if str(item).strip()]
        )
    )
    derived_payload = dict(source_payload)
    derived_files = filtered_scope_manifest_files(
        source_payload.get("files"), excluded_roots=exclusions
    )
    derived_payload["excluded_roots"] = exclusions
    derived_payload["files"] = derived_files
    derived_payload["manifest_sha256"] = specialist_manifest_digest(derived_files)
    derived_payload["derived_from_baseline_path"] = safe_rel(project, source_path)
    derived_payload["derived_from_baseline_sha256"] = source_sha256
    derived_payload["binding_kind"] = binding_kind
    derived_payload["binding_ref"] = binding_ref
    derived_payload["derived_at"] = now_iso()
    write_json_object(target_path, derived_payload)
    require_regular_control_file(target_path, "derived thread scope baseline")
    return target_path, derived_payload


def canonical_project_relative(project: Path, path: Path) -> str:
    return path.resolve().relative_to(project.resolve()).as_posix()


def specialist_scope_manifest(
    project: Path,
    *,
    excluded_roots: list[str],
) -> dict[str, str]:
    roots = [root.strip().rstrip("/") for root in excluded_roots if root.strip()]
    files: dict[str, str] = {}
    for path in sorted(project.rglob("*")):
        rel = path.relative_to(project).as_posix()
        if (
            rel == ".git"
            or rel.startswith(".git/")
            or rel == PRIVATE_LOCAL_STATE_REL.as_posix()
            or rel.startswith(PRIVATE_LOCAL_STATE_REL.as_posix() + "/")
        ):
            continue
        if any(relative_path_is_within(rel, root) for root in roots):
            continue
        if path.is_symlink():
            files[rel] = "symlink:" + os.readlink(path)
        elif path.is_file():
            files[rel] = file_sha256(path)
    return files


def specialist_manifest_digest(files: dict[str, str]) -> str:
    canonical = json.dumps(files, ensure_ascii=False, sort_keys=True).encode("utf-8")
    return hashlib.sha256(canonical).hexdigest()


def with_project_advisory_lock(
    project: Path,
    lock_name: str,
    operation: Callable[[], tuple[dict[str, object], Path | None]],
) -> tuple[dict[str, object], Path | None]:
    lock_path = project / "AD-creative/orchestrator/specialist_exchange" / lock_name
    lock_path.parent.mkdir(parents=True, exist_ok=True)
    with lock_path.open("a+", encoding="utf-8") as lock_handle:
        fcntl.flock(lock_handle.fileno(), fcntl.LOCK_EX)
        try:
            return operation()
        finally:
            fcntl.flock(lock_handle.fileno(), fcntl.LOCK_UN)


def validate_specialist_descriptor(
    descriptor: dict[str, object],
    *,
    profile_id: str,
    required_capabilities: list[str],
    contract_version: str | None = None,
) -> tuple[str, dict[str, object]]:
    if descriptor.get("protocol_id") != SPECIALIST_EXCHANGE_PROTOCOL:
        raise ValueError("descriptor protocol_id mismatch")
    if descriptor.get("message_type") != "descriptor":
        raise ValueError("descriptor message_type mismatch")
    descriptor_version = str(descriptor.get("descriptor_version", ""))
    if not re.fullmatch(r"[12]\.\d+", descriptor_version):
        raise ValueError("unsupported descriptor_version")
    selected_version = contract_version or negotiate_contract_version(descriptor)
    supported = descriptor.get("supported_contract_versions")
    if not isinstance(supported, list) or selected_version not in supported:
        raise ValueError(
            f"descriptor does not support contract_version {selected_version}"
        )
    provider = descriptor.get("provider")
    if not isinstance(provider, dict) or not str(provider.get("id", "")).strip():
        raise ValueError("descriptor provider.id missing")
    provider_id = validate_specialist_token(str(provider.get("id", "")), "provider.id")
    profile_id = validate_specialist_token(profile_id, "profile_id")
    profiles = descriptor.get("profiles")
    if not isinstance(profiles, list):
        raise ValueError("descriptor profiles missing")
    profile = next(
        (
            item
            for item in profiles
            if isinstance(item, dict) and item.get("profile_id") == profile_id
        ),
        None,
    )
    if not isinstance(profile, dict):
        raise ValueError(f"descriptor profile not found: {profile_id}")
    capabilities = profile.get("capabilities")
    if not isinstance(capabilities, list):
        raise ValueError("descriptor profile capabilities missing")
    capability_set = {
        validate_specialist_token(str(item), "descriptor capability")
        for item in capabilities
    }
    missing = sorted(set(required_capabilities) - capability_set)
    if missing:
        raise ValueError("descriptor missing capabilities: " + ",".join(missing))
    authority = profile.get("authority")
    if not isinstance(authority, dict):
        raise ValueError("descriptor profile authority missing")
    for key in [
        "client_interaction",
        "artifact_adoption",
        "client_readiness",
        "final_export",
        "nested_dispatch",
    ]:
        if authority.get(key) is not False:
            raise ValueError(f"descriptor authority escalation: {key}")
    for key, value in authority.items():
        if value is not False:
            raise ValueError(f"descriptor authority escalation: {key}")
    receipt_extension = profile.get("receipt_extension")
    if receipt_extension is not None:
        if not isinstance(receipt_extension, dict):
            raise ValueError("descriptor receipt_extension must be an object")
        validate_specialist_token(
            str(receipt_extension.get("id", "")), "receipt_extension.id"
        )
        if not re.fullmatch(r"\d+\.\d+", str(receipt_extension.get("version", ""))):
            raise ValueError("receipt_extension.version must be major.minor")
        if not isinstance(receipt_extension.get("required"), bool):
            raise ValueError("receipt_extension.required must be boolean")
    validate_specialist_payload(
        "descriptor", descriptor, schema_version=selected_version
    )
    return provider_id, profile


def create_specialist_handoff(
    project: Path,
    *,
    work_id: str,
    profile_id: str,
    objective: str,
    input_artifact_ids: list[str],
    expected_output_kinds: list[str],
    required_capabilities: list[str],
    descriptor_path: Path | None,
    execution_mode: str,
    workspace_mode: str,
    lane_id: str = "",
    generation_mode: str = "prompt_only",
    generation_authorized: bool = False,
    authorization_ref: str = "",
) -> tuple[dict[str, object], Path]:
    result, path = with_project_advisory_lock(
        project,
        ".exchange.lock",
        lambda: _create_specialist_handoff_locked(
            project,
            work_id=work_id,
            profile_id=profile_id,
            objective=objective,
            input_artifact_ids=input_artifact_ids,
            expected_output_kinds=expected_output_kinds,
            required_capabilities=required_capabilities,
            descriptor_path=descriptor_path,
            execution_mode=execution_mode,
            workspace_mode=workspace_mode,
            lane_id=lane_id,
            generation_mode=generation_mode,
            generation_authorized=generation_authorized,
            authorization_ref=authorization_ref,
        ),
    )
    if path is None:
        raise RuntimeError("specialist handoff did not produce a handoff path")
    return result, path


def _create_specialist_handoff_v2_locked(
    project: Path,
    *,
    descriptor: dict[str, object],
    work_id: str,
    profile_id: str,
    objective: str,
    input_artifact_ids: list[str],
    expected_output_kinds: list[str],
    required_capabilities: list[str],
    execution_mode: str,
    lane_id: str,
    generation_mode: str,
    generation_authorized: bool,
    authorization_ref: str,
) -> tuple[dict[str, object], Path]:
    """Persist a minimal provider-facing v2 handoff and ADCO-local control data."""
    if execution_mode != "inline" or lane_id:
        raise ValueError("specialist exchange v2 supports inline execution only")
    if generation_mode != "prompt_only" or generation_authorized or authorization_ref:
        raise ValueError(
            "specialist exchange v2 does not carry real-media authorization; use v1"
        )
    if not work_id.strip() or not objective.strip():
        raise ValueError("work_id and objective are required")
    if not input_artifact_ids or not expected_output_kinds:
        raise ValueError("at least one input artifact and expected output kind are required")

    work_id = validate_specialist_token(work_id, "work_id")
    profile_id = validate_specialist_token(profile_id, "profile_id")
    input_artifact_ids = list(
        dict.fromkeys(
            validate_specialist_token(item, "input_artifact_id")
            for item in input_artifact_ids
        )
    )
    expected_output_kinds = list(
        dict.fromkeys(
            validate_specialist_token(item, "expected_output_kind")
            for item in expected_output_kinds
        )
    )
    required_capabilities = list(
        dict.fromkeys(
            validate_specialist_token(item, "required_capability")
            for item in [*required_capabilities, *expected_output_kinds]
        )
    )
    provider_id, profile = validate_specialist_descriptor(
        descriptor,
        profile_id=profile_id,
        required_capabilities=required_capabilities,
        contract_version=V2_CONTRACT_VERSION,
    )
    modes = {str(item) for item in profile.get("execution_modes", [])}
    if "inline" not in modes:
        raise ValueError("descriptor does not support inline execution")

    _, work_items = read_csv_rows(
        project / "AD-creative/orchestrator/work_items.csv"
    )
    if work_id not in {row.get("work_id", "") for row in work_items}:
        raise ValueError(f"specialist handoff work_id is not registered: {work_id}")
    _, artifacts = read_csv_rows(
        project / "AD-creative/orchestrator/artifact_index.csv"
    )
    artifact_by_id = {
        row.get("artifact_id", ""): row
        for row in artifacts
        if row.get("artifact_id", "")
    }
    locked_decisions: list[dict[str, object]] = []
    for artifact_id in input_artifact_ids:
        row = artifact_by_id.get(artifact_id)
        if row is None:
            raise ValueError(f"input artifact is not indexed: {artifact_id}")
        rel_path = row.get("path", "").strip()
        path = contained_project_path(project, rel_path, f"input artifact {artifact_id}")
        if not path.is_file():
            raise ValueError(f"input artifact file missing: {artifact_id}: {rel_path}")
        actual_sha = file_sha256(path)
        registered_sha = row.get("sha256", "").strip()
        if registered_sha and registered_sha != actual_sha:
            raise ValueError(f"stale_input_artifact: {artifact_id}")
        artifact_type = validate_specialist_token(
            row.get("artifact_type", "").strip() or "input",
            "input artifact type",
        )
        locked_decisions.append(
            {
                "artifact_id": artifact_id,
                "type": artifact_type,
                "path": canonical_project_relative(project, path),
                "sha256": actual_sha,
            }
        )

    index_path = (
        project
        / "AD-creative/orchestrator/specialist_exchange/exchange_index.csv"
    )
    index_fields, index_rows = read_csv_rows(index_path)
    exchange_id = next_id(index_rows, "exchange_id", "SPX")
    handoff_id = "SPH-" + exchange_id.rsplit("-", 1)[-1]
    output_root_rel = (
        f"AD-creative/workspaces/{work_id}/specialists/{handoff_id}/outputs"
    )
    receipt_rel = (
        f"AD-creative/workspaces/{work_id}/specialists/{handoff_id}/receipt.json"
    )
    baseline_rel = (
        f"AD-creative/orchestrator/specialist_exchange/baselines/{handoff_id}.json"
    )
    handoff_rel = (
        f"AD-creative/orchestrator/specialist_exchange/handoffs/{handoff_id}.json"
    )
    requested_outputs = [
        {
            "output_id": f"OUT-{number:02d}",
            "type": kind,
            "path_root": output_root_rel,
        }
        for number, kind in enumerate(expected_output_kinds, start=1)
    ]
    handoff = build_v2_handoff(
        task=objective,
        brief_snapshot=str(locked_decisions[0]["path"]),
        locked_decisions=locked_decisions,
        requested_outputs=requested_outputs,
        quality_targets=required_capabilities,
        execution_mode=execution_mode,
    )
    validate_specialist_payload(
        "handoff", handoff, schema_version=V2_CONTRACT_VERSION
    )

    descriptor_canonical = json.dumps(
        descriptor, ensure_ascii=False, sort_keys=True
    ).encode("utf-8")
    descriptor_sha = hashlib.sha256(descriptor_canonical).hexdigest()
    descriptor_snapshot = (
        project
        / "AD-creative/orchestrator/specialist_exchange/descriptors"
        / f"descriptor_{descriptor_sha}.json"
    )
    baseline_exclusions = [
        "AD-creative/orchestrator/specialist_exchange",
        output_root_rel,
        receipt_rel,
    ]
    baseline_files = v2_scope_manifest(
        project, excluded_roots=baseline_exclusions
    )
    created_at = now_iso()
    baseline_payload: dict[str, object] = {
        "protocol_id": SPECIALIST_EXCHANGE_PROTOCOL,
        "contract_version": V2_CONTRACT_VERSION,
        "message_type": "host_scope_baseline",
        "handoff_id": handoff_id,
        "excluded_roots": baseline_exclusions,
        "files": baseline_files,
        "manifest_sha256": v2_manifest_digest(baseline_files),
        "created_at": created_at,
    }
    baseline_path = project / baseline_rel
    write_json_object(baseline_path, baseline_payload)
    baseline_sha = file_sha256(baseline_path)
    if not descriptor_snapshot.exists():
        write_json_object(descriptor_snapshot, descriptor)
    handoff_path = project / handoff_rel
    write_json_object(handoff_path, handoff)
    handoff_sha = file_sha256(handoff_path)
    index_rows.append(
        {
            "exchange_id": exchange_id,
            "handoff_id": handoff_id,
            "attempt": "1",
            "work_id": work_id,
            "provider_id": provider_id,
            "profile_id": profile_id,
            "contract_version": V2_CONTRACT_VERSION,
            "descriptor_sha256": descriptor_sha,
            "handoff_sha256": handoff_sha,
            "baseline_path": baseline_rel,
            "baseline_sha256": baseline_sha,
            "compatibility_status": "compatible",
            "execution_mode": "inline",
            "lane_id": "",
            "thread_id": "",
            "handoff_path": handoff_rel,
            "receipt_path": receipt_rel,
            "receipt_sha256": "",
            "outcome": "pending",
            "adoption_path": "",
            "adoption_sha256": "",
            "adoption_decision": "",
            "thread_reconciliation_ref": "",
            "created_at": created_at,
            "updated_at": created_at,
        }
    )
    write_csv_rows(index_path, index_fields, index_rows)
    return handoff, handoff_path


def _create_specialist_handoff_locked(
    project: Path,
    *,
    work_id: str,
    profile_id: str,
    objective: str,
    input_artifact_ids: list[str],
    expected_output_kinds: list[str],
    required_capabilities: list[str],
    descriptor_path: Path | None,
    execution_mode: str,
    workspace_mode: str,
    lane_id: str = "",
    generation_mode: str = "prompt_only",
    generation_authorized: bool = False,
    authorization_ref: str = "",
) -> tuple[dict[str, object], Path]:
    migrate_control_plane(project)
    if descriptor_path:
        descriptor = read_json_object(descriptor_path, "descriptor")
        if negotiate_contract_version(descriptor) == V2_CONTRACT_VERSION:
            return _create_specialist_handoff_v2_locked(
                project,
                descriptor=descriptor,
                work_id=work_id,
                profile_id=profile_id,
                objective=objective,
                input_artifact_ids=input_artifact_ids,
                expected_output_kinds=expected_output_kinds,
                required_capabilities=required_capabilities,
                execution_mode=execution_mode,
                lane_id=lane_id,
                generation_mode=generation_mode,
                generation_authorized=generation_authorized,
                authorization_ref=authorization_ref,
            )
    if not work_id.strip() or not objective.strip():
        raise ValueError("work_id and objective are required")
    if not input_artifact_ids or not expected_output_kinds:
        raise ValueError("at least one input artifact and expected output kind are required")
    work_id = validate_specialist_token(work_id, "work_id")
    profile_id = validate_specialist_token(profile_id, "profile_id")
    input_artifact_ids = [
        validate_specialist_token(item, "input_artifact_id")
        for item in input_artifact_ids
    ]
    expected_output_kinds = list(
        dict.fromkeys(
            validate_specialist_token(item, "expected_output_kind")
            for item in expected_output_kinds
        )
    )
    required_capabilities = list(
        dict.fromkeys(
            validate_specialist_token(item, "required_capability")
            for item in [*required_capabilities, *expected_output_kinds]
        )
    )
    if execution_mode not in {"inline", "codex_thread", "external_handoff"}:
        raise ValueError(f"unsupported execution_mode: {execution_mode}")
    if workspace_mode not in {"read_only", "isolated_workspace", "worktree"}:
        raise ValueError(f"unsupported workspace_mode: {workspace_mode}")
    authorization_block: dict[str, object] = {
        "generation_mode": generation_mode,
        "authorized": generation_authorized,
        "authorization_ref": authorization_ref.strip() or None,
        "external_upload": False,
    }
    authorization_errors = specialist_generation_authorization_errors(
        project,
        authorization=authorization_block,
        work_id=work_id,
        profile_id=profile_id,
        input_artifact_ids=input_artifact_ids,
        expected_output_kinds=expected_output_kinds,
    )
    if authorization_errors:
        raise ValueError("; ".join(authorization_errors))
    if generation_mode == "real_media":
        authorization_path = contained_project_path(
            project,
            str(authorization_block["authorization_ref"]),
            "generation authorization_ref",
        )
        authorization_block["authorization_ref"] = canonical_project_relative(
            project, authorization_path
        )

    descriptor_sha = ""
    compatibility_status = "unverified"
    provider_id = validate_specialist_token(
        profile_id.split(".", 1)[0], "provider_id"
    )
    profile: dict[str, object] | None = None
    descriptor_ref: dict[str, object] | None = None
    descriptor_snapshot: Path | None = None
    descriptor_snapshot_payload: dict[str, object] | None = None
    required_receipt_extensions: list[dict[str, str]] = []
    if descriptor_path:
        descriptor = read_json_object(descriptor_path, "descriptor")
        provider_id, profile = validate_specialist_descriptor(
            descriptor,
            profile_id=profile_id,
            required_capabilities=required_capabilities,
        )
        modes = {str(item) for item in profile.get("execution_modes", [])}
        workspaces = {str(item) for item in profile.get("workspace_modes", [])}
        if execution_mode not in modes or workspace_mode not in workspaces:
            raise ValueError("descriptor does not support requested execution/workspace mode")
        canonical = json.dumps(descriptor, ensure_ascii=False, sort_keys=True).encode("utf-8")
        descriptor_sha = hashlib.sha256(canonical).hexdigest()
        descriptor_snapshot = (
            project
            / "AD-creative/orchestrator/specialist_exchange/descriptors"
            / f"descriptor_{descriptor_sha}.json"
        )
        descriptor_snapshot_payload = descriptor
        descriptor_ref = {"provider_id": provider_id, "sha256": descriptor_sha}
        receipt_extension = profile.get("receipt_extension")
        if isinstance(receipt_extension, dict) and receipt_extension.get("required") is True:
            required_receipt_extensions.append(
                {
                    "id": str(receipt_extension["id"]),
                    "version": str(receipt_extension["version"]),
                }
            )
        compatibility_status = "compatible"

    generation_modes = {"prompt_only"}
    if profile is not None:
        raw_generation_modes = profile.get("generation_modes", ["prompt_only"])
        if isinstance(raw_generation_modes, list):
            generation_modes = {str(item) for item in raw_generation_modes}
    if generation_mode not in generation_modes:
        raise ValueError(
            "specialist profile does not support generation_mode: " + generation_mode
        )

    artifact_path = project / "AD-creative/orchestrator/artifact_index.csv"
    _, artifacts = read_csv_rows(artifact_path)
    _, work_items = read_csv_rows(project / "AD-creative/orchestrator/work_items.csv")
    if work_id not in {row.get("work_id", "") for row in work_items}:
        raise ValueError(f"specialist handoff work_id is not registered: {work_id}")
    artifact_by_id = {
        row.get("artifact_id", ""): row for row in artifacts if row.get("artifact_id", "")
    }
    source_artifacts: list[dict[str, object]] = []
    for artifact_id in input_artifact_ids:
        row = artifact_by_id.get(artifact_id)
        if not row:
            raise ValueError(f"input artifact is not indexed: {artifact_id}")
        rel_path = row.get("path", "").strip()
        path = contained_project_path(project, rel_path, f"input artifact {artifact_id}")
        if not path.is_file():
            raise ValueError(f"input artifact file missing: {artifact_id}: {rel_path}")
        actual_sha = file_sha256(path)
        registered_sha = row.get("sha256", "").strip()
        if registered_sha and registered_sha != actual_sha:
            raise ValueError(f"stale_input_artifact: {artifact_id}")
        source_artifacts.append(
            {
                "artifact_id": artifact_id,
                "version": row.get("version", "") or "unversioned",
                "path": rel_path,
                "sha256": actual_sha,
                "visibility": row.get("visibility", "") or "internal_only",
            }
        )

    index_path = project / "AD-creative/orchestrator/specialist_exchange/exchange_index.csv"
    index_fields, index_rows = read_csv_rows(index_path)
    exchange_id = next_id(index_rows, "exchange_id", "SPX")
    handoff_id = "SPH-" + exchange_id.rsplit("-", 1)[-1]
    output_root_rel = f"AD-creative/workspaces/{work_id}/specialists/{handoff_id}/outputs"
    receipt_rel = f"AD-creative/workspaces/{work_id}/specialists/{handoff_id}/receipt.json"
    baseline_rel = f"AD-creative/orchestrator/specialist_exchange/baselines/{handoff_id}.json"
    handoff_rel = f"AD-creative/orchestrator/specialist_exchange/handoffs/{handoff_id}.json"
    thread_id: str | None = None
    lane_run_id: str | None = None
    dispatch_attempt = 1
    thread_control_exclusions: list[str] = []
    if execution_mode == "codex_thread":
        if not lane_id:
            raise ValueError("codex_thread execution requires lane_id")
        validate_specialist_token(lane_id, "lane_id")
        if workspace_mode != "isolated_workspace":
            raise ValueError("specialist exchange v1 codex_thread requires isolated_workspace")
        registry_path = project / "AD-creative/orchestrator/thread_registry.csv"
        registry_fields, registry = read_csv_rows(registry_path)
        row = thread_registry_target(registry, lane_id=lane_id, work_id=work_id)
        candidate = row.get("real_thread_id", "").strip()
        validate_real_thread_id(candidate)
        if row.get("dispatch_status", "").strip().lower() not in {"dispatched", "running"}:
            raise ValueError("codex_thread lane is not an active verified dispatch")
        if row.get("receipt_status", "").strip().lower() in {"received", "rejected"}:
            raise ValueError("codex_thread lane already has a terminal receipt")
        if row.get("reconciliation_status", "").strip().lower() not in {"", "pending"}:
            raise ValueError("codex_thread lane is stale or already reconciled")
        if row.get("mode", "").strip() != "execution_worker":
            raise ValueError("codex_thread specialist lane must be execution_worker")
        if row.get("environment", "").strip() != "isolated_workspace":
            raise ValueError("codex_thread specialist lane must use isolated_workspace")
        lane_run_id = row.get("lane_run_id", "").strip()
        if lane_run_id != f"{work_id}:{lane_id}":
            raise ValueError("codex_thread lane_run_id mismatch")
        dispatch_rel = row.get("dispatch_receipt_path", "").strip()
        dispatch_path = contained_project_path(
            project, dispatch_rel, "codex_thread dispatch receipt"
        )
        if not dispatch_path.is_file() or receipt_thread_ids(
            dispatch_path.read_text(encoding="utf-8", errors="ignore")
        ) != [candidate]:
            raise ValueError("codex_thread lane lacks verified dispatch receipt")
        dispatch_attempt = dispatch_attempt_from_receipt(dispatch_path) or 1
        write_scope_values = [
            item.strip()
            for item in row.get("write_scope", "").split(";")
            if item.strip()
        ]
        if len(write_scope_values) != 1:
            raise ValueError(
                "codex_thread specialist lane requires one bounded workspace-directory write_scope"
            )
        lane_scope = contained_project_path(
            project, write_scope_values[0], "codex_thread write_scope"
        )
        lane_scope_rel = canonical_project_relative(project, lane_scope)
        expected_workspace_rel = f"AD-creative/workspaces/{work_id}/{lane_id}"
        if lane_scope_rel != expected_workspace_rel:
            raise ValueError(
                "codex_thread specialist lane write_scope must equal its isolated workspace directory"
            )
        output_root_rel = f"{lane_scope_rel}/specialist_exchange/{handoff_id}/outputs"
        receipt_rel = f"{lane_scope_rel}/specialist_exchange/{handoff_id}/receipt.json"
        thread_baseline_path = contained_thread_scope_baseline_path(
            project,
            row.get("scope_baseline_path", "").strip(),
            "codex_thread scope baseline",
        )
        if (
            not thread_baseline_path.is_file()
            or file_sha256(thread_baseline_path)
            != row.get("scope_baseline_sha256", "").strip()
        ):
            raise ValueError("codex_thread scope baseline is missing or stale")
        require_regular_control_file(thread_baseline_path, "codex_thread scope baseline")
        thread_baseline = read_json_object(
            thread_baseline_path, "codex_thread scope baseline"
        )
        host_control_paths = [
            "AD-creative/orchestrator/specialist_exchange",
            receipt_rel,
        ]
        if descriptor_path:
            try:
                host_control_paths.append(
                    canonical_project_relative(project, descriptor_path)
                )
            except ValueError:
                pass
        thread_baseline_path, thread_baseline = derive_thread_scope_baseline(
            project,
            thread_baseline_path,
            thread_baseline,
            suffix=f"_handoff-{handoff_id}",
            added_exclusions=host_control_paths,
            binding_kind="specialist_handoff",
            binding_ref=handoff_id,
        )
        row["receipt_path"] = receipt_rel
        row["scope_baseline_path"] = safe_rel(project, thread_baseline_path)
        row["scope_baseline_sha256"] = file_sha256(thread_baseline_path)
        write_csv_rows(registry_path, registry_fields, registry)
        update_thread_agent_run(
            project,
            lane_id=lane_id,
            work_id=work_id,
            updates={
                "receipt_path": receipt_rel,
                "scope_baseline_path": safe_rel(project, thread_baseline_path),
                "scope_baseline_sha256": file_sha256(thread_baseline_path),
                "proof_status": "dispatch_verified_specialist_handoff_bound",
            },
        )
        thread_control_exclusions = [
            "AD-creative/orchestrator/thread_registry.csv",
            "AD-creative/orchestrator/agent_runs.csv",
            "AD-creative/orchestrator/thread_lane_plan.md",
            f"AD-creative/orchestrator/thread_cleanup_{work_id}.md",
            "AD-creative/orchestrator/thread_scope_baselines",
            "AD-creative/orchestrator/thread_scope_proofs",
            f"AD-creative/orchestrator/thread_convergence_{safe_artifact_suffix(work_id)}.md",
            dispatch_rel,
            f"AD-creative/orchestrator/thread_rescue_dispatch_{safe_artifact_suffix(work_id)}_{safe_artifact_suffix(lane_id)}.md",
        ]
        thread_id = candidate
    elif lane_id:
        raise ValueError("lane_id is allowed only for codex_thread execution")

    writable_output_roots = (
        [] if workspace_mode == "read_only" else [output_root_rel]
    )
    baseline_exclusions = [
        "AD-creative/orchestrator/specialist_exchange",
        *writable_output_roots,
        receipt_rel,
        *thread_control_exclusions,
    ]
    baseline_files = specialist_scope_manifest(
        project, excluded_roots=baseline_exclusions
    )
    if generation_mode == "real_media":
        authorization_rel = str(authorization_block["authorization_ref"])
        authorization_path = contained_project_path(
            project, authorization_rel, "generation authorization_ref"
        )
        if baseline_files.get(authorization_rel) != file_sha256(authorization_path):
            raise ValueError(
                "generation authorization evidence must stay inside the monitored host scope"
            )
    exchange_created_at = now_iso()
    baseline_payload: dict[str, object] = {
        "protocol_id": SPECIALIST_EXCHANGE_PROTOCOL,
        "contract_version": SPECIALIST_EXCHANGE_VERSION,
        "message_type": "host_scope_baseline",
        "handoff_id": handoff_id,
        "excluded_roots": baseline_exclusions,
        "files": baseline_files,
        "manifest_sha256": specialist_manifest_digest(baseline_files),
        "created_at": exchange_created_at,
    }
    baseline_path = project / baseline_rel
    write_json_object(baseline_path, baseline_payload)
    baseline_sha = file_sha256(baseline_path)

    handoff: dict[str, object] = {
        "protocol_id": SPECIALIST_EXCHANGE_PROTOCOL,
        "contract_version": SPECIALIST_EXCHANGE_VERSION,
        "message_type": "handoff",
        "exchange_id": exchange_id,
        "handoff_id": handoff_id,
        "attempt": dispatch_attempt,
        "supersedes_handoff_id": None,
        "work_id": work_id,
        "provider_id": provider_id,
        "profile_id": profile_id,
        "descriptor_ref": descriptor_ref,
        "task": {
            "objective": objective,
            "required_capabilities": required_capabilities,
            "expected_output_kinds": expected_output_kinds,
        },
        "source_truth": {"artifacts": source_artifacts, "locks": []},
        "execution": {
            "mode": execution_mode,
            "workspace_mode": workspace_mode,
            "lane_id": lane_id or None,
            "lane_run_id": lane_run_id,
            "thread_id": thread_id,
            "nested_dispatch_allowed": False,
        },
        "scope": {
            "read": [item["path"] for item in source_artifacts],
            "write": [*writable_output_roots, receipt_rel],
            "receipt_path": receipt_rel,
            "host_baseline": {
                "path": safe_rel(project, baseline_path),
                "sha256": baseline_sha,
                "manifest_sha256": baseline_payload["manifest_sha256"],
            },
            "forbidden": [
                "AD-creative/orchestrator/current_truth.md",
                "AD-creative/orchestrator/version_map.csv",
                "AD-creative/orchestrator/artifact_index.csv",
                "AD-creative/orchestrator/gate_log.csv",
                "AD-creative/ppt/exports/",
                "05_最终交付_FinalDelivery/",
            ],
        },
        "authorization": authorization_block,
        "acceptance": {
            "visibility": "internal_only",
            "provider_recommendation_only": True,
            "required_receipt_extensions": required_receipt_extensions,
            "stop_on": ["needs_user", "blocked", "failed"],
        },
    }
    validate_specialist_payload("handoff", handoff)
    if (
        descriptor_snapshot is not None
        and descriptor_snapshot_payload is not None
        and not descriptor_snapshot.exists()
    ):
        write_json_object(descriptor_snapshot, descriptor_snapshot_payload)
    handoff_path = project / handoff_rel
    write_json_object(handoff_path, handoff)
    handoff_sha = file_sha256(handoff_path)
    now = exchange_created_at
    index_rows.append(
        {
            "exchange_id": exchange_id,
            "handoff_id": handoff_id,
            "attempt": str(dispatch_attempt),
            "work_id": work_id,
            "provider_id": provider_id,
            "profile_id": profile_id,
            "contract_version": SPECIALIST_EXCHANGE_VERSION,
            "descriptor_sha256": descriptor_sha,
            "handoff_sha256": handoff_sha,
            "baseline_path": safe_rel(project, baseline_path),
            "baseline_sha256": baseline_sha,
            "compatibility_status": compatibility_status,
            "execution_mode": execution_mode,
            "lane_id": lane_id,
            "thread_id": thread_id or "",
            "handoff_path": safe_rel(project, handoff_path),
            "receipt_path": receipt_rel,
            "receipt_sha256": "",
            "outcome": "pending",
            "adoption_path": "",
            "adoption_sha256": "",
            "adoption_decision": "",
            "thread_reconciliation_ref": "",
            "created_at": now,
            "updated_at": now,
        }
    )
    write_csv_rows(index_path, index_fields, index_rows)
    return handoff, handoff_path


def adopt_specialist_receipt(
    project: Path,
    *,
    handoff_path: Path,
    receipt_path: Path,
    decision: str,
    reason: str,
    output_mappings: dict[str, str],
    dry_run: bool = False,
) -> tuple[dict[str, object], Path | None]:
    operation = lambda: _adopt_specialist_receipt_locked(
        project,
        handoff_path=handoff_path,
        receipt_path=receipt_path,
        decision=decision,
        reason=reason,
        output_mappings=output_mappings,
        dry_run=dry_run,
    )
    if dry_run:
        return operation()
    return with_project_advisory_lock(
        project,
        ".exchange.lock",
        operation,
    )


def _adopt_specialist_receipt_v2_locked(
    project: Path,
    *,
    handoff: dict[str, object],
    handoff_path: Path,
    receipt_path: Path,
    decision: str,
    reason: str,
    output_mappings: dict[str, str],
    dry_run: bool,
) -> tuple[dict[str, object], Path | None]:
    """Validate a minimal provider receipt and persist ADCO adoption separately."""
    validate_specialist_payload(
        "handoff", handoff, schema_version=V2_CONTRACT_VERSION
    )
    boundary_errors = v2_boundary_errors(handoff, message_type="handoff")
    if boundary_errors:
        raise ValueError("; ".join(boundary_errors))

    index_path = (
        project
        / "AD-creative/orchestrator/specialist_exchange/exchange_index.csv"
    )
    index_fields, index_rows = read_csv_rows(index_path)
    handoff_rel = canonical_project_relative(project, handoff_path)
    matches = [row for row in index_rows if row.get("handoff_path") == handoff_rel]
    if len(matches) != 1:
        raise ValueError("specialist exchange v2 index identity is missing or ambiguous")
    exchange_row = matches[0]
    handoff_id = exchange_row.get("handoff_id", "")
    if exchange_row.get("contract_version") != V2_CONTRACT_VERSION:
        raise ValueError("specialist exchange v2 index contract mismatch")
    if exchange_row.get("execution_mode") != "inline":
        raise ValueError("specialist exchange v2 index execution must be inline")
    if exchange_row.get("lane_id") or exchange_row.get("thread_id"):
        raise ValueError("specialist exchange v2 must not bind nested dispatch")
    if exchange_row.get("compatibility_status") != "compatible":
        raise ValueError("specialist exchange v2 descriptor is unverified")
    if file_sha256(handoff_path) != exchange_row.get("handoff_sha256"):
        raise ValueError("handoff_hash_mismatch")
    control_errors = specialist_control_plane_errors(project, index_rows)
    if control_errors:
        raise ValueError(
            "specialist_control_plane_write: " + "; ".join(control_errors)
        )

    receipt_rel = exchange_row.get("receipt_path", "")
    expected_receipt = contained_v2_project_path(
        project, receipt_rel, "v2 receipt"
    )
    if canonical_project_relative(project, receipt_path) != receipt_rel:
        raise ValueError("receipt path does not match v2 exchange index")
    if receipt_path.resolve() != expected_receipt or not expected_receipt.is_file():
        raise ValueError("specialist receipt file is missing")
    receipt_stat = expected_receipt.stat()
    if receipt_stat.st_size == 0 or receipt_stat.st_nlink != 1:
        raise ValueError("specialist receipt must be non-empty and not hardlinked")
    receipt = read_json_object(expected_receipt, "receipt")
    validate_specialist_payload(
        "receipt", receipt, schema_version=V2_CONTRACT_VERSION
    )
    output_by_id, output_errors = validate_v2_receipt_outputs(
        project, handoff, receipt
    )
    if output_errors:
        raise ValueError("; ".join(output_errors))

    status = str(receipt.get("status", ""))
    if decision == "adopt" and status != "completed":
        raise ValueError("only a completed v2 receipt can be fully adopted")
    if decision in {"adopt", "partial_adopt"}:
        if status in {"blocked", "failed"}:
            raise ValueError("blocked/failed v2 receipt cannot be adopted")
        domain_qa = receipt.get("domain_qa")
        if not isinstance(domain_qa, dict) or domain_qa.get("status") != "pass":
            raise ValueError("v2 adoption requires domain_qa.status=pass")

    if decision == "adopt" and set(output_mappings) != set(output_by_id):
        raise ValueError("full adoption requires a target mapping for every v2 output")
    if decision == "partial_adopt" and not output_mappings:
        raise ValueError("partial adoption requires at least one v2 output mapping")
    if decision in {"reject", "defer"} and output_mappings:
        raise ValueError("reject/defer adoption must not map v2 outputs")
    unknown_mappings = set(output_mappings) - set(output_by_id)
    if unknown_mappings:
        raise ValueError(
            "mapping references unknown v2 output: "
            + ",".join(sorted(unknown_mappings))
        )

    baseline_path = contained_v2_project_path(
        project, exchange_row.get("baseline_path", ""), "v2 host scope baseline"
    )
    if not baseline_path.is_file() or file_sha256(baseline_path) != exchange_row.get(
        "baseline_sha256", ""
    ):
        raise ValueError("v2 host scope baseline is missing or stale")
    baseline = read_json_object(baseline_path, "v2 host scope baseline")
    baseline_files = baseline.get("files")
    excluded_roots = baseline.get("excluded_roots")
    if not isinstance(baseline_files, dict) or not isinstance(excluded_roots, list):
        raise ValueError("v2 host scope baseline is malformed")
    current_files = v2_scope_manifest(
        project, excluded_roots=[str(item) for item in excluded_roots]
    )
    if current_files != {str(key): str(value) for key, value in baseline_files.items()}:
        raise ValueError("specialist_scope_violation: host files changed outside v2 scope")
    observed_manifest_sha = v2_manifest_digest(current_files)

    forbidden_roots = [
        contained_project_path(project, value, "v2 adoption forbidden scope")
        for value in [
            "AD-creative/orchestrator",
            "AD-creative/ppt/exports",
            "05_最终交付_FinalDelivery",
        ]
    ]
    prepared: list[tuple[str, dict[str, object], Path, Path]] = []
    prepared_targets: set[Path] = set()
    for output_id, target_rel in output_mappings.items():
        target = contained_project_path(project, target_rel, "v2 adoption target")
        if any(target == root or root in target.parents for root in forbidden_roots):
            raise ValueError(f"v2 adoption target is forbidden: {target_rel}")
        if target in prepared_targets:
            raise ValueError(f"multiple v2 outputs map to one target: {target_rel}")
        if target.exists():
            raise ValueError(f"v2 adoption target already exists: {target_rel}")
        prepared_targets.add(target)
        item, source = output_by_id[output_id]
        prepared.append((output_id, item, source, target))

    receipt_sha = file_sha256(expected_receipt)
    adoption_id = "SPA-" + safe_artifact_suffix(handoff_id)
    adoption_rel = (
        Path("AD-creative/orchestrator/specialist_exchange/adoptions")
        / f"{adoption_id}.json"
    )
    adoption_path = project / adoption_rel
    adopted_outputs: list[dict[str, object]] = []
    for number, (output_id, item, _, target) in enumerate(prepared, start=1):
        adopted_outputs.append(
            {
                "output_id": output_id,
                "target_artifact_id": f"ART-{safe_artifact_suffix(adoption_id)}-{number:02d}",
                "type": item.get("type"),
                "target_path": safe_rel(project, target),
                "sha256": item.get("sha256"),
                "visibility": "internal_only",
            }
        )
    adoption: dict[str, object] = {
        "protocol_id": "adco.specialist-adoption",
        "version": "1.0",
        "contract_version": V2_CONTRACT_VERSION,
        "adoption_id": adoption_id,
        "handoff_id": handoff_id,
        "receipt_sha256": receipt_sha,
        "decision_owner": "adco",
        "decision": decision,
        "reason": reason,
        "adopted_outputs": adopted_outputs,
        "rejected_outputs": sorted(set(output_by_id) - set(output_mappings)),
        "limitations_carried_forward": (
            receipt.get("domain_qa", {}).get("limitations", [])
            if isinstance(receipt.get("domain_qa"), dict)
            else []
        ),
        "adco_validation": ["path", "hash", "type", "domain_qa", "scope"],
        "host_scope_proof": {
            "baseline_path": exchange_row.get("baseline_path", ""),
            "baseline_sha256": exchange_row.get("baseline_sha256", ""),
            "baseline_manifest_sha256": baseline.get("manifest_sha256", ""),
            "observed_manifest_sha256": observed_manifest_sha,
            "changed_paths": [],
        },
        "gate_effect": {
            "advance_allowed": status == "completed"
            and decision in {"adopt", "partial_adopt"},
            "next_gate": "creative-quality-gate",
        },
        "created_at": now_iso(),
    }
    if dry_run:
        return adoption, None
    if adoption_path.exists():
        raise FileExistsError(
            f"specialist adoption already exists for handoff: {adoption_rel}"
        )

    artifact_path = project / "AD-creative/orchestrator/artifact_index.csv"
    artifact_snapshot = artifact_path.read_bytes()
    index_snapshot = index_path.read_bytes()
    created_targets: list[Path] = []
    try:
        artifact_fields, artifact_rows = read_csv_rows(artifact_path)
        for adopted, (_, item, source, target) in zip(adopted_outputs, prepared):
            target.parent.mkdir(parents=True, exist_ok=True)
            with target.open("xb") as handle:
                handle.write(source.read_bytes())
            created_targets.append(target)
            artifact_rows.append(
                {
                    "artifact_id": str(adopted["target_artifact_id"]),
                    "artifact_type": str(item.get("type", "specialist_output")),
                    "path": safe_rel(project, target),
                    "stage": "specialist_adoption",
                    "version": "1",
                    "status": "internal_review",
                    "visibility": "internal_only",
                    "source_event_ids": "",
                    "linked_requirements": "",
                    "linked_work_items": exchange_row.get("work_id", ""),
                    "linked_references": "",
                    "linked_assets": "",
                    "gate_status": "NOT_RUN",
                    "supersedes_artifact_id": "",
                    "created_at": now_iso(),
                    "updated_at": now_iso(),
                    "sha256": str(item.get("sha256", "")),
                    "size_bytes": str(target.stat().st_size),
                    "derived_from_artifact_id": "",
                    "derived_from_sha256": "",
                }
            )
        write_csv_rows(artifact_path, artifact_fields, artifact_rows)
        write_json_object(adoption_path, adoption)
        exchange_row.update(
            {
                "receipt_sha256": receipt_sha,
                "outcome": status,
                "adoption_path": str(adoption_rel),
                "adoption_sha256": file_sha256(adoption_path),
                "adoption_decision": decision,
                "updated_at": now_iso(),
            }
        )
        write_csv_rows(index_path, index_fields, index_rows)
        validation_errors = validate_v2_exchange_row(project, exchange_row)
        if validation_errors:
            raise ValueError(
                "specialist v2 adoption would be invalid: "
                + "; ".join(validation_errors[:12])
            )
    except Exception:
        artifact_path.write_bytes(artifact_snapshot)
        index_path.write_bytes(index_snapshot)
        adoption_path.unlink(missing_ok=True)
        for target in created_targets:
            target.unlink(missing_ok=True)
        raise
    return adoption, adoption_path


def _adopt_specialist_receipt_locked(
    project: Path,
    *,
    handoff_path: Path,
    receipt_path: Path,
    decision: str,
    reason: str,
    output_mappings: dict[str, str],
    dry_run: bool = False,
) -> tuple[dict[str, object], Path | None]:
    migrate_control_plane(project, dry_run=dry_run)
    for evidence_path, label in [
        (handoff_path.resolve(), "handoff"),
        (receipt_path.resolve(), "receipt"),
    ]:
        try:
            evidence_path.relative_to(project.resolve())
        except ValueError as exc:
            raise ValueError(f"{label} must stay inside the project") from exc
    decision = decision.strip().lower()
    if decision not in {"adopt", "partial_adopt", "reject", "defer"}:
        raise ValueError(f"unsupported adoption decision: {decision}")
    if not reason.strip():
        raise ValueError("adoption reason is required")
    handoff = read_json_object(handoff_path, "handoff")
    if handoff.get("contract_version") == V2_CONTRACT_VERSION:
        return _adopt_specialist_receipt_v2_locked(
            project,
            handoff=handoff,
            handoff_path=handoff_path,
            receipt_path=receipt_path,
            decision=decision,
            reason=reason,
            output_mappings=output_mappings,
            dry_run=dry_run,
        )
    validate_specialist_payload("handoff", handoff)
    receipt_scope = handoff.get("scope")
    if not isinstance(receipt_scope, dict):
        raise ValueError("handoff scope missing")
    receipt_rel = str(receipt_scope.get("receipt_path", ""))
    expected_receipt_lexical = project / receipt_rel
    supplied_receipt_lexical = Path(os.path.abspath(receipt_path))
    allowed_receipt_lexical_paths = {
        Path(os.path.abspath(expected_receipt_lexical)),
        project.resolve() / receipt_rel,
    }
    if supplied_receipt_lexical not in allowed_receipt_lexical_paths:
        raise ValueError("receipt path does not match exact handoff receipt_path")
    if "\\" in receipt_rel or project_relative_path_has_symlink_component(
        project, receipt_rel
    ):
        raise ValueError("receipt_path must use a non-symlink POSIX project path")
    if not expected_receipt_lexical.is_file():
        raise ValueError("specialist receipt file is missing")
    receipt_stat = expected_receipt_lexical.stat()
    if receipt_stat.st_size == 0 or receipt_stat.st_nlink != 1:
        raise ValueError("specialist receipt must be non-empty and not hardlinked")
    receipt = read_json_object(expected_receipt_lexical, "receipt")
    for payload, message_type in [(handoff, "handoff"), (receipt, "receipt")]:
        if payload.get("protocol_id") != SPECIALIST_EXCHANGE_PROTOCOL:
            raise ValueError(f"{message_type} protocol_id mismatch")
        if payload.get("contract_version") != SPECIALIST_EXCHANGE_VERSION:
            raise ValueError(f"unsupported_contract_version: {message_type}")
        if payload.get("message_type") != message_type:
            raise ValueError(f"{message_type} message_type mismatch")
    index_path = project / "AD-creative/orchestrator/specialist_exchange/exchange_index.csv"
    _, exchange_rows = read_csv_rows(index_path)
    exchange_matches = [
        row
        for row in exchange_rows
        if row.get("handoff_id", "") == str(handoff.get("handoff_id", ""))
    ]
    if len(exchange_matches) != 1:
        raise ValueError("specialist exchange index identity is missing or ambiguous")
    exchange_row = exchange_matches[0]
    control_errors = specialist_control_plane_errors(project, exchange_rows)
    if control_errors:
        raise ValueError("specialist_control_plane_write: " + "; ".join(control_errors))
    handoff_execution = handoff.get("execution")
    handoff_scope = handoff.get("scope")
    descriptor_ref_for_index = handoff.get("descriptor_ref")
    baseline_ref_for_index = (
        handoff_scope.get("host_baseline")
        if isinstance(handoff_scope, dict)
        else None
    )
    expected_index_values = {
        "exchange_id": str(handoff.get("exchange_id", "")),
        "handoff_id": str(handoff.get("handoff_id", "")),
        "attempt": str(handoff.get("attempt", "")),
        "work_id": str(handoff.get("work_id", "")),
        "provider_id": str(handoff.get("provider_id", "")),
        "profile_id": str(handoff.get("profile_id", "")),
        "contract_version": str(handoff.get("contract_version", "")),
        "descriptor_sha256": (
            str(descriptor_ref_for_index.get("sha256", ""))
            if isinstance(descriptor_ref_for_index, dict)
            else ""
        ),
        "baseline_path": (
            str(baseline_ref_for_index.get("path", ""))
            if isinstance(baseline_ref_for_index, dict)
            else ""
        ),
        "baseline_sha256": (
            str(baseline_ref_for_index.get("sha256", ""))
            if isinstance(baseline_ref_for_index, dict)
            else ""
        ),
        "compatibility_status": (
            "compatible" if isinstance(descriptor_ref_for_index, dict) else "unverified"
        ),
        "execution_mode": (
            str(handoff_execution.get("mode", ""))
            if isinstance(handoff_execution, dict)
            else ""
        ),
        "lane_id": (
            str(handoff_execution.get("lane_id") or "")
            if isinstance(handoff_execution, dict)
            else ""
        ),
        "thread_id": (
            str(handoff_execution.get("thread_id") or "")
            if isinstance(handoff_execution, dict)
            else ""
        ),
        "handoff_path": safe_rel(project, handoff_path),
        "receipt_path": (
            str(handoff_scope.get("receipt_path", ""))
            if isinstance(handoff_scope, dict)
            else ""
        ),
        "receipt_sha256": "",
        "outcome": "pending",
        "adoption_path": "",
        "adoption_sha256": "",
        "adoption_decision": "",
        "thread_reconciliation_ref": "",
    }
    for field, expected_value in expected_index_values.items():
        if exchange_row.get(field, "") != expected_value:
            raise ValueError(f"specialist exchange index {field} binding mismatch")
    if exchange_row.get("handoff_path", "") != safe_rel(project, handoff_path):
        raise ValueError("handoff path does not match exchange index")
    actual_handoff_sha = file_sha256(handoff_path)
    if exchange_row.get("handoff_sha256", "") != actual_handoff_sha:
        raise ValueError("handoff_hash_mismatch")
    if decision in {"adopt", "partial_adopt"} and exchange_row.get(
        "compatibility_status", ""
    ) != "compatible":
        raise ValueError("unverified specialist descriptor cannot be adopted")

    for key in ["exchange_id", "handoff_id", "work_id", "provider_id", "profile_id"]:
        if receipt.get(key) != handoff.get(key):
            raise ValueError(f"receipt {key} mismatch")
    descriptor_ref = handoff.get("descriptor_ref")
    descriptor_sha = (
        str(descriptor_ref.get("sha256", ""))
        if isinstance(descriptor_ref, dict)
        else ""
    )
    if receipt.get("descriptor_sha256") != descriptor_sha:
        raise ValueError("receipt descriptor_sha256 mismatch")
    if receipt.get("handoff_sha256") != actual_handoff_sha:
        raise ValueError("receipt handoff_sha256 mismatch")
    if exchange_row.get("provider_id", "") != str(handoff.get("provider_id", "")):
        raise ValueError("provider identity does not match exchange index")
    if exchange_row.get("descriptor_sha256", "") != descriptor_sha:
        raise ValueError("descriptor identity does not match exchange index")
    claims = receipt.get("claims")
    if (
        not isinstance(claims, dict)
        or set(claims) != set(SPECIALIST_RESERVED_CLAIMS)
        or any(claims.get(key) is not False for key in SPECIALIST_RESERVED_CLAIMS)
    ):
        raise ValueError("authority escalation in specialist receipt claims")

    acceptance = handoff.get("acceptance")
    extensions = receipt.get("extensions")
    if not isinstance(acceptance, dict) or not isinstance(extensions, list):
        raise ValueError("specialist receipt extensions missing")
    extension_pairs = {
        (str(item.get("id", "")), str(item.get("version", "")))
        for item in extensions
        if isinstance(item, dict)
    }
    for required in acceptance.get("required_receipt_extensions", []):
        if not isinstance(required, dict):
            raise ValueError("invalid required receipt extension negotiation")
        pair = (str(required.get("id", "")), str(required.get("version", "")))
        if pair not in extension_pairs:
            raise ValueError(
                f"required receipt extension missing: {pair[0]}@{pair[1]}"
            )
    scope = handoff.get("scope")
    if not isinstance(scope, dict):
        raise ValueError("handoff scope missing")
    source_truth_for_auth = handoff.get("source_truth")
    task_for_auth = handoff.get("task")
    if not isinstance(source_truth_for_auth, dict) or not isinstance(task_for_auth, dict):
        raise ValueError("handoff authorization context is missing")
    authorization_errors = specialist_generation_authorization_errors(
        project,
        authorization=handoff.get("authorization"),
        work_id=str(handoff.get("work_id", "")),
        profile_id=str(handoff.get("profile_id", "")),
        input_artifact_ids=[
            str(item.get("artifact_id", ""))
            for item in source_truth_for_auth.get("artifacts", [])
            if isinstance(item, dict)
        ],
        expected_output_kinds=[
            str(item) for item in task_for_auth.get("expected_output_kinds", [])
        ],
    )
    if authorization_errors:
        raise ValueError("; ".join(authorization_errors))
    expected_receipt = contained_project_path(
        project, str(scope.get("receipt_path", "")), "handoff receipt_path"
    )
    if receipt_path.resolve() != expected_receipt:
        raise ValueError("receipt path does not match handoff receipt_path")
    baseline_ref = scope.get("host_baseline")
    if not isinstance(baseline_ref, dict):
        raise ValueError("handoff host scope baseline missing")
    baseline_path = contained_project_path(
        project, str(baseline_ref.get("path", "")), "host scope baseline"
    )
    if not baseline_path.is_file() or file_sha256(baseline_path) != baseline_ref.get(
        "sha256"
    ):
        raise ValueError("host scope baseline hash mismatch")
    baseline = read_json_object(baseline_path, "host scope baseline")
    baseline_created_at = str(baseline.get("created_at", ""))
    if (
        not baseline_created_at
        or exchange_row.get("created_at", "") != baseline_created_at
        or exchange_row.get("updated_at", "") != baseline_created_at
    ):
        raise ValueError("specialist exchange index timestamp binding mismatch")
    baseline_files = baseline.get("files")
    excluded_roots = baseline.get("excluded_roots")
    if not isinstance(baseline_files, dict) or not isinstance(excluded_roots, list):
        raise ValueError("host scope baseline is malformed")
    baseline_files = {str(key): str(value) for key, value in baseline_files.items()}
    authorization = handoff.get("authorization")
    if isinstance(authorization, dict) and authorization.get("generation_mode") == "real_media":
        authorization_rel = str(authorization.get("authorization_ref", ""))
        authorization_path = contained_project_path(
            project, authorization_rel, "generation authorization_ref"
        )
        if baseline_files.get(authorization_rel) != file_sha256(authorization_path):
            raise ValueError(
                "generation authorization evidence is not bound by the host baseline"
            )
    current_files = specialist_scope_manifest(
        project, excluded_roots=[str(item) for item in excluded_roots]
    )
    if current_files != baseline_files:
        changed_paths = sorted(
            path
            for path in set(current_files) | set(baseline_files)
            if current_files.get(path) != baseline_files.get(path)
        )
        raise ValueError("out_of_scope_changes: " + ",".join(changed_paths[:12]))
    host_scope_proof = {
        "baseline_path": safe_rel(project, baseline_path),
        "baseline_sha256": str(baseline_ref.get("sha256", "")),
        "baseline_manifest_sha256": specialist_manifest_digest(baseline_files),
        "observed_manifest_sha256": specialist_manifest_digest(current_files),
        "changed_paths": [],
    }
    execution = handoff.get("execution")
    evidence = receipt.get("execution_evidence")
    if not isinstance(execution, dict) or not isinstance(evidence, dict):
        raise ValueError("execution evidence missing")
    if evidence.get("mode") != execution.get("mode"):
        raise ValueError("receipt execution mode mismatch")
    if evidence.get("nested_dispatch_used") is not False:
        raise ValueError("nested dispatch is not allowed")
    if evidence.get("out_of_scope_writes") not in ([], None):
        raise ValueError("receipt reports out-of-scope writes")
    thread_reconciliation_ref: dict[str, object] | None = None
    if execution.get("mode") == "codex_thread":
        registry_path = project / "AD-creative/orchestrator/thread_registry.csv"
        _, registry = read_csv_rows(registry_path)
        row = thread_registry_target(
            registry,
            lane_id=str(execution.get("lane_id") or ""),
            work_id=str(handoff.get("work_id") or ""),
        )
        valid_ids = {
            value
            for value in [row.get("real_thread_id", ""), row.get("rescue_thread_id", "")]
            if value
        }
        if evidence.get("thread_id") not in valid_ids:
            raise ValueError("invalid_worker_thread_id")
        if execution.get("lane_run_id") != row.get("lane_run_id", ""):
            raise ValueError("codex_thread lane_run_id changed after handoff")
        if decision in {"adopt", "partial_adopt"}:
            if row.get("receipt_status", "").strip().lower() != "received":
                raise ValueError("codex_thread receipt is not host-received")
            if row.get("reconciliation_status", "").strip().lower() != "reconciled":
                raise ValueError("codex_thread receipt is not host-reconciled")
            if row.get("receipt_thread_id", "").strip() != evidence.get("thread_id"):
                raise ValueError("codex_thread reconciliation identity mismatch")
            registry_receipt = contained_project_path(
                project,
                row.get("receipt_path", "").strip(),
                "codex_thread reconciled receipt",
            )
            if registry_receipt != receipt_path.resolve():
                raise ValueError("codex_thread reconciled a different receipt")
            if row.get("adoption_decision", "").strip().upper() not in {
                "ADOPT",
                "PARTIAL_ADOPT",
            }:
                raise ValueError("codex_thread main adoption decision is missing")
            if not normalized_bool(row.get("archived")) or not row.get(
                "archived_at", ""
            ).strip():
                raise ValueError("codex_thread cleanup/archive is incomplete")
            if not row.get("cleanup_action", "").strip():
                raise ValueError("codex_thread cleanup action is missing")
            scope_proof_rel = row.get("scope_proof_path", "").strip()
            scope_proof_path = contained_project_path(
                project, scope_proof_rel, "codex_thread host scope proof"
            )
            if not scope_proof_path.is_file() or file_sha256(
                scope_proof_path
            ) != row.get("scope_proof_sha256", "").strip():
                raise ValueError("codex_thread host scope proof is missing or stale")
            scope_proof = read_json_object(
                scope_proof_path, "codex_thread host scope proof"
            )
            if scope_proof.get("decision") not in {"ADOPT", "PARTIAL_ADOPT"}:
                raise ValueError("codex_thread host scope proof lacks adoption decision")
            if scope_proof.get("validation_success") is not True:
                raise ValueError("codex_thread host scope proof validation failed")
            dispatch_ref = row.get("dispatch_receipt_path", "").strip()
            if evidence.get("thread_id") == row.get("rescue_thread_id", "").strip():
                dispatch_ref = row.get("rescue_dispatch_receipt_path", "").strip()
            dispatch_path = contained_project_path(
                project, dispatch_ref, "codex_thread dispatch proof"
            )
            if not dispatch_path.is_file() or evidence.get("thread_id") not in dispatch_path.read_text(
                encoding="utf-8", errors="ignore"
            ):
                raise ValueError("codex_thread dispatch proof is missing")
            thread_reconciliation_ref = {
                "lane_id": row.get("lane_id", ""),
                "lane_run_id": row.get("lane_run_id", ""),
                "thread_id": evidence.get("thread_id"),
                "receipt_path": safe_rel(project, receipt_path),
                "receipt_sha256": file_sha256(receipt_path),
                "dispatch_receipt_path": safe_rel(project, dispatch_path),
                "dispatch_receipt_sha256": file_sha256(dispatch_path),
                "registry_sha256": file_sha256(registry_path),
                "scope_proof_path": safe_rel(project, scope_proof_path),
                "scope_proof_sha256": file_sha256(scope_proof_path),
                "archived_at": row.get("archived_at", ""),
                "cleanup_action": row.get("cleanup_action", ""),
            }
    elif evidence.get("thread_id") not in {None, ""}:
        raise ValueError("non-thread exchange must not claim thread_id")

    outcome = str(receipt.get("outcome", ""))
    if outcome not in {"completed", "needs_user", "blocked", "failed"}:
        raise ValueError("invalid specialist outcome")
    open_questions = receipt.get("open_questions")
    if outcome == "needs_user":
        if not isinstance(open_questions, list) or not open_questions:
            raise ValueError("needs_user receipt lacks open_questions")
        if any(
            not isinstance(item, dict)
            or not str(item.get("id", "")).strip()
            or not str(item.get("question", "")).strip()
            for item in open_questions
        ):
            raise ValueError("needs_user receipt open_questions must contain id and question")
        question_ids = [str(item["id"]).strip() for item in open_questions]
        if len(question_ids) != len(set(question_ids)):
            raise ValueError("needs_user receipt contains duplicate question id")
    if decision == "adopt" and outcome != "completed":
        raise ValueError("only completed receipt can be fully adopted")
    if decision in {"adopt", "partial_adopt"}:
        if outcome in {"blocked", "failed"}:
            raise ValueError("blocked/failed receipt cannot be adopted")
        if bool(receipt.get("simulated")):
            raise ValueError("simulated receipt cannot be adopted")
        qa = receipt.get("qa")
        if not isinstance(qa, dict) or qa.get("status") != "pass":
            raise ValueError("adoption requires specialist qa.status=pass")

    source_truth = handoff.get("source_truth")
    consumed = receipt.get("consumed_inputs")
    if not isinstance(source_truth, dict) or not isinstance(consumed, list):
        raise ValueError("receipt consumed_inputs missing")
    consumed_by_id = {
        str(item.get("artifact_id")): item
        for item in consumed
        if isinstance(item, dict)
    }
    for item in source_truth.get("artifacts", []):
        if not isinstance(item, dict):
            continue
        artifact_id = str(item.get("artifact_id"))
        consumed_item = consumed_by_id.get(artifact_id)
        if not consumed_item or consumed_item.get("sha256") != item.get("sha256"):
            raise ValueError(f"stale_input_artifact: {artifact_id}")
        source_path = contained_project_path(project, str(item.get("path")), artifact_id)
        if not source_path.is_file() or file_sha256(source_path) != item.get("sha256"):
            raise ValueError(f"stale_input_artifact: {artifact_id}")

    scope = handoff.get("scope")
    if not isinstance(scope, dict):
        raise ValueError("handoff scope missing")
    allowed_roots = [
        contained_project_path(project, str(item), "handoff write scope")
        for item in scope.get("write", [])
        if str(item) != str(scope.get("receipt_path"))
    ]
    outputs = receipt.get("output_artifacts")
    if not isinstance(outputs, list):
        raise ValueError("receipt output_artifacts missing")
    write_scope = [str(item) for item in scope.get("write", [])]
    if execution.get("workspace_mode") == "read_only":
        if write_scope != [str(scope.get("receipt_path", ""))]:
            raise ValueError(
                "read_only handoff may write only its exact receipt_path"
            )
        if outputs:
            raise ValueError("read_only receipt must not return output artifacts")
    task = handoff.get("task")
    if not isinstance(task, dict):
        raise ValueError("handoff task missing")
    expected_kinds = {
        str(item) for item in task.get("expected_output_kinds", []) if str(item)
    }
    if not expected_kinds:
        raise ValueError("handoff expected_output_kinds missing")
    if outcome == "completed" and not outputs:
        raise ValueError("completed specialist receipt has no outputs")
    output_by_id: dict[str, tuple[dict[str, object], Path]] = {}
    returned_kinds: set[str] = set()
    returned_paths: set[str] = set()
    returned_inodes: set[tuple[int, int]] = set()
    source_input_ids = {
        str(item.get("artifact_id"))
        for item in source_truth.get("artifacts", [])
        if isinstance(item, dict)
    }
    for item in outputs:
        if not isinstance(item, dict):
            raise ValueError("invalid output artifact entry")
        provider_artifact_id = validate_specialist_token(
            str(item.get("provider_artifact_id", "")), "provider_artifact_id"
        )
        if provider_artifact_id in output_by_id:
            raise ValueError(f"duplicate provider_artifact_id: {provider_artifact_id}")
        kind = validate_specialist_token(str(item.get("kind", "")), "output kind")
        if kind not in expected_kinds:
            raise ValueError(f"unexpected specialist output kind: {kind}")
        if kind in returned_kinds:
            raise ValueError(f"duplicate specialist output kind: {kind}")
        returned_kinds.add(kind)
        if item.get("visibility") != "internal_only":
            raise ValueError("specialist output visibility must remain internal_only")
        output_sources = item.get("source_input_ids")
        if (
            not isinstance(output_sources, list)
            or not output_sources
            or not {str(source) for source in output_sources}.issubset(source_input_ids)
        ):
            raise ValueError(f"specialist output source_input_ids invalid: {provider_artifact_id}")
        raw_output_path = str(item.get("path", ""))
        if "\\" in raw_output_path:
            raise ValueError(
                f"specialist output must use POSIX path separators: {provider_artifact_id}"
            )
        if project_relative_path_has_symlink_component(project, raw_output_path):
            raise ValueError(
                f"specialist output must not use symlink path: {provider_artifact_id}"
            )
        output_path = contained_project_path(
            project, raw_output_path, provider_artifact_id
        )
        if not output_path.is_file():
            raise ValueError(f"specialist output missing: {provider_artifact_id}")
        if not any(output_path == root or root in output_path.parents for root in allowed_roots):
            raise ValueError(f"specialist output outside write scope: {provider_artifact_id}")
        canonical_output_path = canonical_project_relative(project, output_path)
        if canonical_output_path in returned_paths:
            raise ValueError(
                f"duplicate specialist output path: {canonical_output_path}"
            )
        returned_paths.add(canonical_output_path)
        output_stat = output_path.stat()
        if output_stat.st_size == 0 or output_stat.st_nlink != 1:
            raise ValueError(
                f"specialist output must be non-empty and not hardlinked: {provider_artifact_id}"
            )
        physical_id = (output_stat.st_dev, output_stat.st_ino)
        if physical_id in returned_inodes:
            raise ValueError(
                f"specialist output physical file reused: {provider_artifact_id}"
            )
        returned_inodes.add(physical_id)
        if file_sha256(output_path) != item.get("sha256"):
            raise ValueError(f"output_hash_mismatch: {provider_artifact_id}")
        output_by_id[provider_artifact_id] = (item, output_path)
    if outcome == "completed" and not expected_kinds.issubset(returned_kinds):
        missing_kinds = sorted(expected_kinds - returned_kinds)
        raise ValueError("completed receipt missing expected output kinds: " + ",".join(missing_kinds))
    validate_specialist_payload("receipt", receipt)

    if decision == "adopt" and set(output_mappings) != set(output_by_id):
        raise ValueError("full adoption requires a target mapping for every output")
    if decision == "partial_adopt" and not output_mappings:
        raise ValueError("partial adoption requires at least one output mapping")
    if decision in {"reject", "defer"} and output_mappings:
        raise ValueError("reject/defer adoption must not map specialist outputs")
    unknown_mappings = set(output_mappings) - set(output_by_id)
    if unknown_mappings:
        raise ValueError("mapping references unknown output: " + ",".join(sorted(unknown_mappings)))

    forbidden_roots = [
        contained_project_path(project, value, "adoption forbidden scope")
        for value in [
            "AD-creative/orchestrator",
            "AD-creative/ppt/exports",
            "05_最终交付_FinalDelivery",
            *[str(item) for item in scope.get("forbidden", [])],
        ]
    ]
    prepared: list[tuple[str, dict[str, object], Path, Path]] = []
    prepared_targets: set[Path] = set()
    for provider_artifact_id, target_rel in output_mappings.items():
        target = contained_project_path(project, target_rel, "adoption target")
        if any(target == root or root in target.parents for root in forbidden_roots):
            raise ValueError(f"adoption target is forbidden control/final scope: {target_rel}")
        if target in prepared_targets:
            raise ValueError(f"multiple outputs map to the same adoption target: {target_rel}")
        prepared_targets.add(target)
        if target.exists():
            raise ValueError(f"adoption target already exists; overwrite forbidden: {target_rel}")
        item, source = output_by_id[provider_artifact_id]
        prepared.append((provider_artifact_id, item, source, target))

    receipt_sha = file_sha256(receipt_path)
    adoption_id = "SPA-" + safe_artifact_suffix(str(handoff.get("handoff_id")))
    adoption_rel = (
        Path("AD-creative/orchestrator/specialist_exchange/adoptions")
        / f"{adoption_id}.json"
    )
    adoption_path = project / adoption_rel
    adopted_outputs: list[dict[str, object]] = []
    for number, (provider_artifact_id, item, _, target) in enumerate(prepared, start=1):
        adopted_outputs.append(
            {
                "provider_artifact_id": provider_artifact_id,
                "target_artifact_id": f"ART-{safe_artifact_suffix(adoption_id)}-{number:02d}",
                "target_path": safe_rel(project, target),
                "sha256": item.get("sha256"),
                "visibility": "internal_only",
            }
        )
    adoption: dict[str, object] = {
        "protocol_id": SPECIALIST_EXCHANGE_PROTOCOL,
        "contract_version": SPECIALIST_EXCHANGE_VERSION,
        "message_type": "adoption",
        "adoption_id": adoption_id,
        "handoff_id": handoff.get("handoff_id"),
        "receipt_id": receipt.get("receipt_id"),
        "receipt_sha256": receipt_sha,
        "decision_owner": "adco",
        "decision": decision,
        "reason": reason,
        "adopted_outputs": adopted_outputs,
        "rejected_outputs": sorted(set(output_by_id) - set(output_mappings)),
        "limitations_carried_forward": (receipt.get("qa") or {}).get("limitations", []) if isinstance(receipt.get("qa"), dict) else [],
        "adco_validation": [
            "protocol",
            "identity",
            "scope",
            "host_scope_manifest",
            "hash",
            "authority",
            "output_contract",
        ],
        "host_scope_proof": host_scope_proof,
        "gate_effect": {
            "advance_allowed": outcome == "completed" and decision in {"adopt", "partial_adopt"},
            "next_gate": "creative-quality-gate",
        },
        "thread_reconciliation_ref": thread_reconciliation_ref,
        "created_at": now_iso(),
    }
    validate_specialist_payload("adoption", adoption)
    if dry_run:
        return adoption, None
    if adoption_path.exists():
        raise FileExistsError(
            f"specialist adoption already exists for handoff: {adoption_rel}"
        )

    artifact_path = project / "AD-creative/orchestrator/artifact_index.csv"
    index_path = project / "AD-creative/orchestrator/specialist_exchange/exchange_index.csv"
    artifact_snapshot = artifact_path.read_bytes()
    index_snapshot = index_path.read_bytes()
    created_targets: list[Path] = []
    try:
        artifact_fields, artifact_rows = read_csv_rows(artifact_path)
        for adopted, (_, item, source, target) in zip(adopted_outputs, prepared):
            target.parent.mkdir(parents=True, exist_ok=True)
            with target.open("xb") as handle:
                handle.write(source.read_bytes())
            created_targets.append(target)
            artifact_rows.append(
                {
                    "artifact_id": str(adopted["target_artifact_id"]),
                    "artifact_type": str(item.get("kind", "specialist_output")),
                    "path": safe_rel(project, target),
                    "stage": "specialist_adoption",
                    "version": str(item.get("version", "1")),
                    "status": "internal_review",
                    "visibility": "internal_only",
                    "source_event_ids": "",
                    "linked_requirements": "",
                    "linked_work_items": str(handoff.get("work_id", "")),
                    "linked_references": "",
                    "linked_assets": "",
                    "gate_status": "NOT_RUN",
                    "supersedes_artifact_id": "",
                    "created_at": now_iso(),
                    "updated_at": now_iso(),
                    "sha256": str(item.get("sha256", "")),
                    "size_bytes": str(target.stat().st_size),
                    "derived_from_artifact_id": "",
                    "derived_from_sha256": "",
                }
            )
        write_csv_rows(artifact_path, artifact_fields, artifact_rows)
        write_json_object(adoption_path, adoption)
        index_fields, index_rows = read_csv_rows(index_path)
        row = next(
            item
            for item in index_rows
            if item.get("handoff_id") == handoff.get("handoff_id")
        )
        if row.get("compatibility_status") != "compatible" and decision in {"adopt", "partial_adopt"}:
            raise ValueError("unverified specialist descriptor cannot be adopted")
        row.update(
            {
                "receipt_path": safe_rel(project, receipt_path),
                "receipt_sha256": receipt_sha,
                "outcome": outcome,
                "adoption_path": str(adoption_rel),
                "adoption_sha256": file_sha256(adoption_path),
                "adoption_decision": decision,
                "thread_reconciliation_ref": (
                    json.dumps(
                        thread_reconciliation_ref,
                        ensure_ascii=False,
                        sort_keys=True,
                    )
                    if thread_reconciliation_ref
                    else ""
                ),
                "updated_at": now_iso(),
            }
        )
        write_csv_rows(index_path, index_fields, index_rows)
        validation_errors, _ = validate(project)
        if validation_errors:
            raise ValueError(
                "specialist adoption would leave project invalid: "
                + "; ".join(validation_errors[:12])
            )
    except Exception:
        artifact_path.write_bytes(artifact_snapshot)
        index_path.write_bytes(index_snapshot)
        adoption_path.unlink(missing_ok=True)
        for target in created_targets:
            target.unlink(missing_ok=True)
        raise
    return adoption, adoption_path


def write_asset_preflight(
    project: Path,
    *,
    work_id: str,
    source_scope: str,
    local_manifest_checked: str,
    browser_checked: str,
    browser_tool: str,
    download_method: str,
    imported_asset_ids: str,
    replacement_generation_allowed: str,
    status: str,
    blocked_reason: str,
) -> str:
    migrate_control_plane(project)
    path = project / "AD-creative/orchestrator/agency/asset_preflight.csv"
    fields, rows = read_csv_rows(path)
    preflight_id = next_id(rows, "preflight_id", "APF")
    rows.append(
        {
            "preflight_id": preflight_id,
            "work_id": work_id,
            "source_scope": source_scope,
            "local_manifest_checked": local_manifest_checked,
            "browser_checked": browser_checked,
            "browser_tool": browser_tool,
            "download_method": download_method,
            "imported_asset_ids": imported_asset_ids,
            "replacement_generation_allowed": replacement_generation_allowed,
            "status": status,
            "blocked_reason": blocked_reason,
            "created_at": now_iso(),
        }
    )
    write_csv_rows(path, fields, rows)
    return preflight_id


def parse_thread_timestamp(value: str, field: str) -> datetime:
    normalized = value.strip()
    if normalized.endswith("Z"):
        normalized = normalized[:-1] + "+00:00"
    try:
        parsed = datetime.fromisoformat(normalized)
    except ValueError as exc:
        raise ValueError(f"{field} must be an ISO-8601 timestamp with timezone") from exc
    if parsed.tzinfo is None:
        raise ValueError(f"{field} must include a timezone")
    return parsed


def validate_real_thread_id(value: str, field: str = "real_thread_id") -> str:
    thread_id = value.strip()
    if not THREADOPS_REAL_THREAD_ID_PATTERN.fullmatch(thread_id):
        raise ValueError(f"{field} must be a real Codex Thread UUID, not a planned or narrative id")
    return thread_id


def thread_registry_target(
    rows: list[dict[str, str]], *, lane_id: str, work_id: str
) -> dict[str, str]:
    matches = [
        row
        for row in rows
        if row.get("lane_id", "").strip() == lane_id
        and row.get("work_id", "").strip() == work_id
    ]
    if len(matches) != 1:
        raise ValueError(
            f"expected exactly one registry row for work_id={work_id} lane_id={lane_id}; found {len(matches)}"
        )
    return matches[0]


def update_thread_agent_run(
    project: Path,
    *,
    lane_id: str,
    work_id: str,
    updates: dict[str, str],
) -> None:
    path = project / "AD-creative/orchestrator/agent_runs.csv"
    fields, rows = read_csv_rows(path)
    matches = [
        row
        for row in rows
        if row.get("lane_id", "").strip() == lane_id
        and row.get("work_id", "").strip() == work_id
    ]
    if len(matches) != 1:
        raise ValueError(
            f"expected exactly one agent_runs row for work_id={work_id} lane_id={lane_id}; found {len(matches)}"
        )
    updates = dict(updates)
    baseline_path = updates.pop("scope_baseline_path", "").strip()
    baseline_sha256 = updates.pop("scope_baseline_sha256", "").strip()
    if baseline_path or baseline_sha256:
        if not baseline_path or not baseline_sha256:
            raise ValueError("agent run scope baseline pointer requires path and hash")
        summary = matches[0].get("summary", "")
        summary = re.sub(
            r"(?:^|\s*\|\s*)scope_baseline_path=[^;|]*;scope_baseline_sha256=[^;|]*",
            "",
            summary,
        ).strip(" |;")
        pointer = (
            f"scope_baseline_path={baseline_path};"
            f"scope_baseline_sha256={baseline_sha256}"
        )
        updates["summary"] = f"{summary} | {pointer}" if summary else pointer
    matches[0].update(updates)
    write_csv_rows(path, fields, rows)


def append_thread_convergence_event(
    project: Path,
    row: dict[str, str],
    *,
    state: str,
    observed_at: str,
    evidence: str,
) -> Path:
    work_id = row.get("work_id", "") or "unknown-work"
    path = project / f"AD-creative/orchestrator/thread_convergence_{safe_artifact_suffix(work_id)}.md"
    if path.exists():
        text = path.read_text(encoding="utf-8").rstrip() + "\n"
    else:
        text = (
            "# Thread Convergence Log\n\n"
            "Fixed poll counts are inspection budgets, not automatic failure. "
            "Progress permits at most one reasoned bounded extension; rescue_count must stay <= 1.\n\n"
            "| observed_at | work_id | lane_id | real_thread_id | state | absolute_deadline_at | "
            "extension_used | rescue_count | receipt_thread_id | evidence |\n"
            "|---|---|---|---|---|---|---|---|---|---|\n"
        )
    clean_evidence = evidence.replace("|", "/").replace("\n", " ").strip()
    text += (
        f"| {observed_at} | {work_id} | {row.get('lane_id', '')} | "
        f"{row.get('real_thread_id', '')} | {state} | {row.get('absolute_deadline_at', '')} | "
        f"{row.get('bounded_extension_used', 'false')} | {row.get('rescue_count', '0')} | "
        f"{row.get('receipt_thread_id', '')} | {clean_evidence} |\n"
    )
    write_text(path, text)
    return path


def refresh_threadops_projections(project: Path, work_id: str) -> None:
    """Refresh human views from the registry; the registry remains authoritative."""
    _, registry = read_csv_rows(project / "AD-creative/orchestrator/thread_registry.csv")
    rows = [row for row in registry if row.get("work_id", "").strip() == work_id]
    plan_path = project / "AD-creative/orchestrator/thread_lane_plan.md"
    if plan_path.is_file() and rows:
        lines = plan_path.read_text(encoding="utf-8").splitlines()
        by_lane = {row.get("lane_id", ""): row for row in rows}
        for heading in ("## Lane Map", "## Thread Registry"):
            try:
                heading_index = lines.index(heading)
            except ValueError:
                continue
            header_index = next(
                (
                    index
                    for index in range(heading_index + 1, len(lines))
                    if lines[index].startswith("| lane_id |")
                    or (heading == "## Thread Registry" and lines[index].startswith("| thread_id |"))
                ),
                None,
            )
            if header_index is None:
                continue
            headers = [cell.strip() for cell in lines[header_index].strip("|").split("|")]
            for index in range(header_index + 2, len(lines)):
                if not lines[index].startswith("|"):
                    break
                cells = [cell.strip() for cell in lines[index].strip("|").split("|")]
                projected = dict(zip(headers, cells))
                source = by_lane.get(projected.get("lane_id", ""))
                if not source:
                    continue
                updates = {header: source.get(header, "") for header in headers if header in source}
                if "lifecycle_status" in headers:
                    updates["lifecycle_status"] = source.get("lifecycle_state", "")
                projected.update({key: value for key, value in updates.items() if key in headers})
                lines[index] = "| " + " | ".join(
                    markdown_table_cell(projected.get(header, "")) for header in headers
                ) + " |"
        write_text(plan_path, "\n".join(lines) + "\n")
    cleanup_path = project / f"AD-creative/orchestrator/thread_cleanup_{work_id}.md"
    if cleanup_path.is_file() and rows:
        state = "archived" if all(
            row.get("archived", "").strip().lower() in {"1", "true", "yes", "y", "on"}
            for row in rows
        ) else "active"
        text = re.sub(r"(?m)^status:\s*\S+", f"status: {state}", cleanup_path.read_text(encoding="utf-8"), count=1)
        write_text(cleanup_path, text)


def receipt_thread_ids(text: str) -> list[str]:
    json_ids: list[str] = []
    try:
        payload = json.loads(text)
    except json.JSONDecodeError:
        payload = None
    if isinstance(payload, dict):
        direct = payload.get("thread_id")
        evidence = payload.get("execution_evidence")
        nested = evidence.get("thread_id") if isinstance(evidence, dict) else None
        for value in [direct, nested]:
            if isinstance(value, str) and value.strip():
                json_ids.append(value.strip())
        if json_ids:
            return list(dict.fromkeys(json_ids))
    pattern = re.compile(
        r"(?im)^\s*(?:[-*]\s*)?(?:receipt\.)?(?:thread_id|real_thread_id)\s*[:=]\s*([^\s,;]+)\s*$"
    )
    return list(
        dict.fromkeys(match.group(1).strip() for match in pattern.finditer(text))
    )


def bind_dispatch_identity_file(
    path: Path,
    *,
    real_thread_id: str,
    absolute_deadline_at: str,
    receipt_envelope: bool,
) -> None:
    if not path.is_file():
        raise ValueError(f"dispatch identity target missing: {path}")
    text = path.read_text(encoding="utf-8")
    text = text.replace("TBD_DISPATCH_RECORD_REQUIRED", real_thread_id)
    text = text.replace("TBD_DISPATCH_DEADLINE_REQUIRED", absolute_deadline_at)
    if receipt_envelope and not receipt_thread_ids(text):
        text = re.sub(
            r"(?im)^\s*thread_id\s*:\s*TBD\s*$",
            f"thread_id: {real_thread_id}",
            text,
            count=1,
        )
    if receipt_envelope and receipt_thread_ids(text) != [real_thread_id]:
        raise ValueError(f"receipt envelope could not be bound to {real_thread_id}: {path}")
    if not receipt_envelope and real_thread_id not in text:
        lane_match = re.search(r"(?m)^Lane id:.*$", text)
        if not lane_match:
            raise ValueError(f"worker prompt lacks Lane id anchor: {path}")
        insert_at = lane_match.end()
        text = (
            text[:insert_at]
            + f"\nExpected real thread id: {real_thread_id}\nAbsolute deadline: {absolute_deadline_at}"
            + text[insert_at:]
        )
    write_text(path, text)


def dispatch_attempt_from_receipt(path: Path) -> int:
    if not path.is_file():
        return 0
    match = re.search(r"(?im)^dispatch_attempt:\s*(\d+)\s*$", path.read_text(encoding="utf-8", errors="ignore"))
    return int(match.group(1)) if match else 1


def dispatch_identity_paths(
    project: Path, target: dict[str, str], *, work_id: str, lane_id: str
) -> tuple[Path, Path]:
    dispatch_rel = target.get("dispatch_receipt_path", "").strip()
    if dispatch_rel:
        dispatch_path = contained_project_path(
            project, dispatch_rel, "current dispatch receipt"
        )
        if dispatch_path.is_file():
            dispatch_text = dispatch_path.read_text(encoding="utf-8", errors="ignore")
            prompt_match = re.search(r"(?im)^prompt_path:\s*(\S+)\s*$", dispatch_text)
            receipt_match = re.search(
                r"(?im)^worker_receipt_path:\s*(\S+)\s*$", dispatch_text
            )
            if prompt_match and receipt_match:
                return (
                    contained_project_path(project, prompt_match.group(1), "worker prompt"),
                    contained_project_path(
                        project, receipt_match.group(1), "worker receipt envelope"
                    ),
                )
    prompt_rel = target.get("notes", "").partition("prompt=")[2].strip()
    if not prompt_rel:
        raise ValueError("registry row lacks worker prompt path")
    worker_receipt_rel = target.get("receipt_path", "").strip()
    if (
        not worker_receipt_rel
        or "/specialist_exchange/" in worker_receipt_rel
        or worker_receipt_rel.startswith("AD-creative/workspaces/")
    ):
        worker_receipt_rel = (
            f"AD-creative/agents/receipts/{work_id}/{lane_id}_receipt.md"
        )
    return (
        contained_project_path(project, prompt_rel, "worker prompt"),
        contained_project_path(project, worker_receipt_rel, "worker receipt envelope"),
    )


def attempt_identity_path(path: Path, attempt: int) -> Path:
    if attempt == 1:
        return path
    return path.with_name(f"{path.stem}_attempt-{attempt:02d}{path.suffix}")


def update_agent_run_path_list(value: str, old_path: str, new_path: str) -> str:
    paths = [item.strip() for item in value.split(";") if item.strip()]
    replaced = False
    updated: list[str] = []
    for item in paths:
        if item == old_path:
            updated.append(new_path)
            replaced = True
        else:
            updated.append(item)
    if new_path not in updated:
        updated.append(new_path)
    if not replaced and old_path:
        updated = [item for item in updated if item != old_path]
    return ";".join(dict.fromkeys(updated))


def record_thread_dispatch(
    project: Path,
    *,
    lane_id: str,
    work_id: str,
    real_thread_id: str,
    title_action: str,
    title_verified_at: str,
    dispatch_evidence: str,
    dispatch_status: str,
    absolute_deadline_at: str,
) -> dict[str, str]:
    migrate_control_plane(project)
    real_thread_id = validate_real_thread_id(real_thread_id)
    title_verified_time = parse_thread_timestamp(title_verified_at, "title_verified_at")
    absolute_deadline_time = parse_thread_timestamp(absolute_deadline_at, "absolute_deadline_at")
    if absolute_deadline_time <= title_verified_time:
        raise ValueError("absolute_deadline_at must be later than title_verified_at")
    registry_path = project / "AD-creative/orchestrator/thread_registry.csv"
    fields, rows = read_csv_rows(registry_path)
    target = thread_registry_target(rows, lane_id=lane_id, work_id=work_id)
    terminal_reasons: list[str] = []
    if normalized_bool(target.get("archived")):
        terminal_reasons.append("archived")
    if target.get("lifecycle_state", "").strip().lower() in {
        "archived", "reconciled", "adopted", "closed"
    }:
        terminal_reasons.append(target.get("lifecycle_state", "").strip().lower())
    if target.get("receipt_status", "").strip().lower() in {"received", "rejected"}:
        terminal_reasons.append(target.get("receipt_status", "").strip().lower())
    if target.get("reconciliation_status", "").strip().lower() in {
        "reconciled", "rejected_evidence", "adopted", "archived"
    }:
        terminal_reasons.append(target.get("reconciliation_status", "").strip().lower())
    if target.get("adoption_decision", "").strip().upper() in {
        "ADOPT", "PARTIAL_ADOPT", "REJECT"
    }:
        terminal_reasons.append(target.get("adoption_decision", "").strip().upper())
    if terminal_reasons:
        raise ValueError(
            "cannot redispatch terminal lane; create a new work_id/lane_id "
            f"(state={','.join(dict.fromkeys(terminal_reasons))})"
        )
    planned_thread_id = target.get("planned_thread_id", "").strip() or target.get("thread_id", "").strip()
    source_prompt_path, source_worker_receipt_path = dispatch_identity_paths(
        project, target, work_id=work_id, lane_id=lane_id
    )
    dispatch_dir = project / "AD-creative/orchestrator/dispatch_receipts"
    stem = f"{safe_artifact_suffix(work_id)}_{safe_artifact_suffix(lane_id)}"
    existing_dispatches = sorted(
        dispatch_dir.glob(f"thread_dispatch_{stem}_attempt-*.md")
    )
    previous_attempt = max(
        [dispatch_attempt_from_receipt(path) for path in existing_dispatches]
        + ([1] if target.get("real_thread_id", "").strip() else [0])
    )
    if previous_attempt >= 2:
        raise ValueError(
            "redispatch attempt limit exceeded; create a new work_id/lane_id"
        )
    attempt = previous_attempt + 1 if previous_attempt else 1
    while (dispatch_dir / f"thread_dispatch_{stem}_attempt-{attempt:02d}.md").exists():
        attempt += 1
    prompt_path = attempt_identity_path(source_prompt_path, attempt)
    worker_receipt_path = attempt_identity_path(source_worker_receipt_path, attempt)
    if attempt > 1:
        existing_baseline_rel = target.get("scope_baseline_path", "").strip()
        existing_baseline_sha = target.get("scope_baseline_sha256", "").strip()
        if existing_baseline_rel:
            existing_baseline = contained_thread_scope_baseline_path(
                project, existing_baseline_rel, "existing thread scope baseline"
            )
            if not existing_baseline.is_file() or file_sha256(existing_baseline) != existing_baseline_sha:
                raise ValueError("thread scope baseline is missing or stale before redispatch")
        if not source_prompt_path.is_file() or not source_worker_receipt_path.is_file():
            raise ValueError("redispatch requires the prior worker prompt and receipt envelope")
        if prompt_path.exists() or worker_receipt_path.exists():
            raise ValueError("attempt-specific dispatch identity already exists")
        previous_thread_id = target.get("real_thread_id", "").strip()
        previous_deadline = target.get("absolute_deadline_at", "").strip()
        for source_path, copied_path in [
            (source_prompt_path, prompt_path),
            (source_worker_receipt_path, worker_receipt_path),
        ]:
            copied_text = source_path.read_text(encoding="utf-8")
            for old_value, new_value in [
                (previous_thread_id, real_thread_id),
                (previous_deadline, absolute_deadline_at),
            ]:
                if old_value:
                    copied_text = copied_text.replace(old_value, new_value)
            write_text(copied_path, copied_text)
    bind_dispatch_identity_file(
        prompt_path,
        real_thread_id=real_thread_id,
        absolute_deadline_at=absolute_deadline_at,
        receipt_envelope=False,
    )
    bind_dispatch_identity_file(
        worker_receipt_path,
        real_thread_id=real_thread_id,
        absolute_deadline_at=absolute_deadline_at,
        receipt_envelope=True,
    )
    prompt_rel = safe_rel(project, prompt_path)
    worker_receipt_rel = safe_rel(project, worker_receipt_path)
    old_prompt_rel = safe_rel(project, source_prompt_path)
    old_worker_receipt_rel = safe_rel(project, source_worker_receipt_path)
    supersedes_thread_id = target.get("real_thread_id", "").strip() if attempt > 1 else ""
    lane_plan_path = project / "AD-creative/orchestrator/thread_lane_plan.md"
    if lane_plan_path.is_file() and planned_thread_id:
        write_text(
            lane_plan_path,
            lane_plan_path.read_text(encoding="utf-8").replace(planned_thread_id, real_thread_id),
        )
    receipt_path = dispatch_dir / f"thread_dispatch_{stem}_attempt-{attempt:02d}.md"
    baseline_rel = (
        "AD-creative/orchestrator/thread_scope_baselines/"
        f"{safe_artifact_suffix(work_id)}_{safe_artifact_suffix(lane_id)}"
        f"{'_attempt-' + format(attempt, '02d') if attempt > 1 else ''}.json"
    )
    baseline_path = contained_thread_scope_baseline_path(
        project, baseline_rel, "dispatch scope baseline"
    )
    baseline_exclusions = [
        "AD-creative/orchestrator/thread_scope_baselines",
        "AD-creative/orchestrator/thread_scope_proofs",
        "AD-creative/orchestrator/thread_registry.csv",
        "AD-creative/orchestrator/agent_runs.csv",
        "AD-creative/orchestrator/thread_lane_plan.md",
        f"AD-creative/orchestrator/thread_cleanup_{work_id}.md",
        safe_rel(project, receipt_path),
        safe_rel(project, worker_receipt_path),
        f"AD-creative/orchestrator/thread_convergence_{safe_artifact_suffix(work_id)}.md",
        f"AD-creative/orchestrator/thread_rescue_dispatch_{safe_artifact_suffix(work_id)}_{safe_artifact_suffix(lane_id)}.md",
    ]
    baseline_files = specialist_scope_manifest(
        project, excluded_roots=baseline_exclusions
    )
    if baseline_path.exists():
        raise ValueError(f"dispatch scope baseline already exists: {baseline_path}")
    write_json_object(
        baseline_path,
        {
            "protocol_id": "adco.thread-scope-baseline",
            "version": "1.0",
            "work_id": work_id,
            "lane_id": lane_id,
            "real_thread_id": real_thread_id,
            "write_scope": target.get("write_scope", ""),
            "excluded_roots": baseline_exclusions,
            "files": baseline_files,
            "manifest_sha256": specialist_manifest_digest(baseline_files),
            "created_at": now_iso(),
        },
    )
    require_regular_control_file(baseline_path, "dispatch scope baseline")
    baseline_sha = file_sha256(baseline_path)
    target.update(
        {
            "thread_id": real_thread_id,
            "lane_run_id": f"{work_id}:{lane_id}",
            "planned_thread_id": target.get("planned_thread_id") or f"planned:{lane_id}",
            "real_thread_id": real_thread_id,
            "dispatch_status": dispatch_status,
            "title_action": title_action,
            "title_verified_at": title_verified_at,
            "dispatch_receipt_path": safe_rel(project, receipt_path),
            "dispatch_evidence": dispatch_evidence,
            "notes": f"prompt={prompt_rel}",
            "receipt_path": worker_receipt_rel,
            "scope_baseline_path": safe_rel(project, baseline_path),
            "scope_baseline_sha256": baseline_sha,
            "scope_proof_path": "",
            "scope_proof_sha256": "",
            "rescue_dispatch_receipt_path": target.get("rescue_dispatch_receipt_path") or "",
            "rescue_dispatch_evidence": target.get("rescue_dispatch_evidence") or "",
            "lifecycle_state": dispatch_status,
            "updated_at": now_iso(),
            "last_seen_at": now_iso(),
            "convergence_state": "awaiting_first_readback",
            "reconciliation_status": "pending",
            "last_progress_at": "",
            "absolute_deadline_at": absolute_deadline_at,
            "bounded_extension_used": target.get("bounded_extension_used") or "false",
            "extension_reason": target.get("extension_reason") or "",
            "convergence_reminder_at": target.get("convergence_reminder_at") or "",
            "convergence_reason": "",
            "rescue_count": target.get("rescue_count") or "0",
            "rescue_thread_id": target.get("rescue_thread_id") or "",
            "receipt_thread_id": target.get("receipt_thread_id") or "",
            "adoption_decision": target.get("adoption_decision") or "",
            "rejection_reason": target.get("rejection_reason") or "",
            "schema_state": "current",
            "legacy_evidence_sha256": "",
            "legacy_quarantine_reason": "",
            "legacy_raw_ref": "",
        }
    )
    write_csv_rows(registry_path, fields, rows)
    agent_fields, agent_rows = read_csv_rows(
        project / "AD-creative/orchestrator/agent_runs.csv"
    )
    agent_match = next(
        row
        for row in agent_rows
        if row.get("lane_id", "").strip() == lane_id
        and row.get("work_id", "").strip() == work_id
    )
    update_thread_agent_run(
        project,
        lane_id=lane_id,
        work_id=work_id,
        updates={
            "thread_id": real_thread_id,
            "status": dispatch_status,
            "started_at": title_verified_at,
            "input_files": update_agent_run_path_list(
                agent_match.get("input_files", ""), old_prompt_rel, prompt_rel
            ),
            "output_files": update_agent_run_path_list(
                agent_match.get("output_files", ""), old_worker_receipt_rel, worker_receipt_rel
            ),
            "receipt_path": worker_receipt_rel,
            "scope_baseline_path": safe_rel(project, baseline_path),
            "scope_baseline_sha256": baseline_sha,
            "proof_status": "dispatch_verified",
            "reconciliation_status": "pending",
        },
    )
    write_text(
        receipt_path,
        f"""# Thread Dispatch Receipt

lane_id: {lane_id}
work_id: {work_id}
real_thread_id: {real_thread_id}
dispatch_attempt: {attempt}
dispatch_status: {dispatch_status}
title_action: {title_action}
title_verified_at: {title_verified_at}
absolute_deadline_at: {absolute_deadline_at}
prompt_path: {prompt_rel}
worker_receipt_path: {worker_receipt_rel}
supersedes_thread_id: {supersedes_thread_id}
scope_baseline_path: {safe_rel(project, baseline_path)}
scope_baseline_sha256: {baseline_sha}

## Evidence

{dispatch_evidence}

## Rule

planned:* ids are placeholders only. This receipt records the real Codex Thread id and title readback evidence.
""",
    )
    refresh_threadops_projections(project, work_id)
    return {
        "lane_id": lane_id,
        "work_id": work_id,
        "real_thread_id": real_thread_id,
        "dispatch_status": dispatch_status,
        "dispatch_receipt_path": safe_rel(project, receipt_path),
        "absolute_deadline_at": absolute_deadline_at,
        "dispatch_attempt": str(attempt),
        "scope_baseline_path": safe_rel(project, baseline_path),
        "scope_baseline_sha256": baseline_sha,
    }


def record_thread_observation(
    project: Path,
    *,
    lane_id: str,
    work_id: str,
    state: str,
    observed_at: str,
    evidence: str,
    absolute_deadline_at: str = "",
    extension_reason: str = "",
    convergence_reminder_sent: bool = False,
    rescue_thread_id: str = "",
) -> dict[str, str]:
    migrate_control_plane(project)
    if state not in THREADOPS_OBSERVATION_STATES:
        raise ValueError(f"unknown convergence state: {state}")
    observed_time = parse_thread_timestamp(observed_at, "observed_at")
    registry_path = project / "AD-creative/orchestrator/thread_registry.csv"
    fields, rows = read_csv_rows(registry_path)
    target = thread_registry_target(rows, lane_id=lane_id, work_id=work_id)
    validate_real_thread_id(target.get("real_thread_id", ""))
    current_deadline = target.get("absolute_deadline_at", "").strip()
    deadline_time = (
        parse_thread_timestamp(current_deadline, "absolute_deadline_at")
        if current_deadline
        else None
    )
    previous_state = target.get("convergence_state", "").strip()
    new_deadline = absolute_deadline_at.strip()

    if state in THREADOPS_PROGRESS_STATES:
        if new_deadline and new_deadline != current_deadline:
            if normalized_bool(target.get("bounded_extension_used")):
                raise ValueError("bounded_extension_already_used")
            if not extension_reason.strip():
                raise ValueError("bounded extension requires --extension-reason")
            new_deadline_time = parse_thread_timestamp(new_deadline, "absolute_deadline_at")
            if deadline_time and new_deadline_time <= deadline_time:
                raise ValueError("bounded extension deadline must be later than the current absolute deadline")
            if new_deadline_time <= observed_time:
                raise ValueError("bounded extension deadline must be later than observed_at")
            target["absolute_deadline_at"] = new_deadline
            target["bounded_extension_used"] = "true"
            target["extension_reason"] = extension_reason.strip()
            deadline_time = new_deadline_time
        if deadline_time and observed_time > deadline_time:
            raise ValueError("absolute_deadline_exceeded_without_available_extension")
        target["last_progress_at"] = observed_at
        # New visible activity invalidates any older silence/reminder evidence.
        target["convergence_reminder_at"] = ""
        target["convergence_reason"] = ""
    elif state == "silent":
        if new_deadline and new_deadline != current_deadline:
            raise ValueError("silent observation cannot extend the absolute deadline")
        last_progress = target.get("last_progress_at", "").strip()
        if last_progress and observed_time <= parse_thread_timestamp(
            last_progress, "last_progress_at"
        ):
            raise ValueError("silent observation must be later than last_progress_at")
        if convergence_reminder_sent:
            target["convergence_reminder_at"] = observed_at
    elif state == "thread_not_converged":
        if not deadline_time or observed_time <= deadline_time:
            raise ValueError("thread_not_converged requires observed_at after absolute_deadline_at")
        if previous_state != "silent":
            raise ValueError("thread_not_converged requires a fresh prior silent observation")
        reminder_at = target.get("convergence_reminder_at", "").strip()
        last_progress = target.get("last_progress_at", "").strip()
        if reminder_at and last_progress and parse_thread_timestamp(
            reminder_at, "convergence_reminder_at"
        ) <= parse_thread_timestamp(last_progress, "last_progress_at"):
            raise ValueError("convergence reminder must be later than last_progress_at")
        target["convergence_reason"] = (
            "reminder_no_receipt"
            if reminder_at
            else "silent_past_absolute_deadline"
        )
        target["lifecycle_state"] = "thread_not_converged"
        target["reconciliation_status"] = "thread_not_converged"
    elif state == "rescue_dispatched":
        try:
            rescue_count = int(target.get("rescue_count", "0") or "0")
        except ValueError as exc:
            raise ValueError("rescue_count must be an integer") from exc
        if rescue_count >= 1:
            raise ValueError("rescue_limit_exceeded")
        if previous_state != "thread_not_converged":
            raise ValueError("rescue requires thread_not_converged state")
        rescue_thread_id = validate_real_thread_id(rescue_thread_id, "rescue_thread_id")
        if rescue_thread_id == target.get("real_thread_id", "").strip():
            raise ValueError("rescue_thread_id must differ from real_thread_id")
        if not new_deadline:
            raise ValueError("rescue requires a new --absolute-deadline-at")
        rescue_deadline = parse_thread_timestamp(new_deadline, "absolute_deadline_at")
        if rescue_deadline <= observed_time:
            raise ValueError("rescue absolute deadline must be later than observed_at")
        if not evidence.strip():
            raise ValueError("rescue dispatch requires title/readback dispatch evidence")
        original_prompt, original_receipt = dispatch_identity_paths(
            project, target, work_id=work_id, lane_id=lane_id
        )
        if not original_prompt.is_file() or not original_receipt.is_file():
            raise ValueError("rescue dispatch requires existing bound prompt and receipt envelope")
        rescue_prompt = original_prompt.with_name(
            original_prompt.stem + "_rescue" + original_prompt.suffix
        )
        rescue_receipt = original_receipt.with_name(
            original_receipt.stem + "_rescue" + original_receipt.suffix
        )
        primary_id = target.get("real_thread_id", "").strip()
        prompt_text = original_prompt.read_text(encoding="utf-8")
        receipt_text = original_receipt.read_text(encoding="utf-8")
        for old, new in [
            (primary_id, rescue_thread_id),
            (current_deadline, new_deadline),
        ]:
            if old:
                prompt_text = prompt_text.replace(old, new)
                receipt_text = receipt_text.replace(old, new)
        write_text(rescue_prompt, prompt_text)
        write_text(rescue_receipt, receipt_text)
        if rescue_thread_id not in rescue_prompt.read_text(encoding="utf-8"):
            raise ValueError("rescue prompt could not be bound to rescue_thread_id")
        if receipt_thread_ids(rescue_receipt.read_text(encoding="utf-8")) != [
            rescue_thread_id
        ]:
            raise ValueError("rescue receipt envelope could not be bound to rescue_thread_id")
        rescue_dispatch_path = (
            project
            / "AD-creative/orchestrator"
            / f"thread_rescue_dispatch_{safe_artifact_suffix(work_id)}_{safe_artifact_suffix(lane_id)}.md"
        )
        source_baseline_path = contained_thread_scope_baseline_path(
            project,
            target.get("scope_baseline_path", "").strip(),
            "thread scope baseline",
        )
        if not source_baseline_path.is_file() or file_sha256(source_baseline_path) != target.get(
            "scope_baseline_sha256", ""
        ).strip():
            raise ValueError("thread scope baseline is missing or stale before rescue")
        require_regular_control_file(source_baseline_path, "thread scope baseline")
        baseline = read_json_object(source_baseline_path, "thread scope baseline")
        baseline_path, baseline = derive_thread_scope_baseline(
            project,
            source_baseline_path,
            baseline,
            suffix=f"_rescue-{safe_artifact_suffix(rescue_thread_id)}",
            added_exclusions=[
                safe_rel(project, rescue_prompt),
                safe_rel(project, rescue_receipt),
                safe_rel(project, rescue_dispatch_path),
            ],
            binding_kind="rescue",
            binding_ref=rescue_thread_id,
        )
        write_text(
            rescue_dispatch_path,
            f"""# Thread Rescue Dispatch Receipt

work_id: {work_id}
lane_id: {lane_id}
real_thread_id: {rescue_thread_id}
supersedes_thread_id: {primary_id}
dispatched_at: {observed_at}
absolute_deadline_at: {new_deadline}
rescue_prompt_path: {safe_rel(project, rescue_prompt)}
rescue_receipt_path: {safe_rel(project, rescue_receipt)}
scope_baseline_path: {safe_rel(project, baseline_path)}
scope_baseline_sha256: {file_sha256(baseline_path)}

## Readback Evidence

{evidence}
""",
        )
        target["rescue_count"] = "1"
        target["rescue_thread_id"] = rescue_thread_id
        target["receipt_path"] = safe_rel(project, rescue_receipt)
        target["scope_baseline_path"] = safe_rel(project, baseline_path)
        target["scope_baseline_sha256"] = file_sha256(baseline_path)
        target["rescue_dispatch_receipt_path"] = safe_rel(
            project, rescue_dispatch_path
        )
        target["rescue_dispatch_evidence"] = evidence.strip()
        target["absolute_deadline_at"] = new_deadline
        target["lifecycle_state"] = "rescue_dispatched"
        target["reconciliation_status"] = "rescue_dispatched"

    target["convergence_state"] = state
    target["last_seen_at"] = observed_at
    target["updated_at"] = now_iso()
    write_csv_rows(registry_path, fields, rows)
    if state == "rescue_dispatched":
        update_thread_agent_run(
            project,
            lane_id=lane_id,
            work_id=work_id,
            updates={
                "receipt_path": target.get("receipt_path", ""),
                "scope_baseline_path": target.get("scope_baseline_path", ""),
                "scope_baseline_sha256": target.get("scope_baseline_sha256", ""),
                "status": "rescue_dispatched",
                "proof_status": "rescue_dispatch_verified",
                "reconciliation_status": "rescue_dispatched",
            },
        )
    refresh_threadops_projections(project, work_id)
    log_path = append_thread_convergence_event(
        project, target, state=state, observed_at=observed_at, evidence=evidence
    )
    return {
        "lane_id": lane_id,
        "work_id": work_id,
        "state": state,
        "absolute_deadline_at": target.get("absolute_deadline_at", ""),
        "bounded_extension_used": target.get("bounded_extension_used", "false"),
        "convergence_reminder_at": target.get("convergence_reminder_at", ""),
        "rescue_count": target.get("rescue_count", "0"),
        "log_path": safe_rel(project, log_path),
    }


def thread_receipt_json(text: str) -> dict[str, object] | None:
    try:
        payload = json.loads(text)
    except json.JSONDecodeError:
        return None
    return payload if isinstance(payload, dict) else None


def receipt_field_text(text: str, key: str) -> str:
    values = receipt_proof_values(
        text, key, THREADOPS_RECEIPT_REQUIRED_PROOF.get(key, ())
    )
    return "\n".join(value.strip() for value in values if value.strip())


def declared_thread_changed_paths(text: str) -> set[str]:
    payload = thread_receipt_json(text)
    if payload and payload.get("protocol_id") == SPECIALIST_EXCHANGE_PROTOCOL:
        outputs = payload.get("output_artifacts")
        if not isinstance(outputs, list):
            return set()
        return {
            str(item.get("path", "")).strip()
            for item in outputs
            if isinstance(item, dict) and str(item.get("path", "")).strip()
        }
    raw = receipt_field_text(text, "files_changed")
    paths: set[str] = set()
    for line in raw.splitlines():
        cleaned = re.sub(r"^[\s>*-]+", "", line).strip().strip("`")
        if not cleaned:
            continue
        for item in cleaned.split(";"):
            candidate = item.strip().strip("`")
            if candidate.lower() in {
                "none",
                "no files changed",
                "no file changes",
                "not_applicable",
                "n/a",
            }:
                continue
            if " | " in candidate:
                candidate = candidate.split(" | ", 1)[0].strip()
            paths.add(candidate)
    return paths


def validation_result_is_success(text: str) -> bool:
    payload = thread_receipt_json(text)
    if payload and payload.get("protocol_id") == SPECIALIST_EXCHANGE_PROTOCOL:
        qa = payload.get("qa")
        return isinstance(qa, dict) and qa.get("status") == "pass"
    value = receipt_field_text(text, "validation_result")
    lowered = value.lower()
    if re.search(r"\b(fail(?:ed)?|error|blocked|not[_ ]?run|exit\s*[=:]?\s*[1-9])\b", lowered):
        return False
    return bool(re.search(r"\b(pass(?:ed)?|success|ok|exit\s*[=:]?\s*0)\b", lowered))


def receipt_loop_state_is_complete(text: str) -> bool:
    payload = thread_receipt_json(text)
    if payload and payload.get("protocol_id") == SPECIALIST_EXCHANGE_PROTOCOL:
        return payload.get("outcome") == "completed"
    value = receipt_field_text(text, "loop_state").lower()
    return bool(re.search(r"\b(returned|reconciled|archived|completed|success)\b", value)) and not bool(
        re.search(r"\b(blocked|failed|error|frozen|replay_requested)\b", value)
    )


def validate_thread_receipt_scope_and_semantics(
    project: Path,
    row: dict[str, str],
    receipt_text: str,
    *,
    decision: str,
    cleanup_action: str,
    archived_at: str,
) -> dict[str, object]:
    if not archived_at:
        raise ValueError("thread reconciliation requires archived_at cleanup evidence")
    if re.search(r"\b(none|not[_ ]?done|did not|failed|pending|todo|tbd)\b", cleanup_action, re.IGNORECASE):
        raise ValueError("thread cleanup_action is not complete")
    baseline_rel = row.get("scope_baseline_path", "").strip()
    baseline_path = contained_thread_scope_baseline_path(
        project, baseline_rel, "thread scope baseline"
    )
    if not baseline_path.is_file() or file_sha256(baseline_path) != row.get(
        "scope_baseline_sha256", ""
    ).strip():
        raise ValueError("thread scope baseline missing or hash mismatch")
    baseline = read_json_object(baseline_path, "thread scope baseline")
    baseline_files = baseline.get("files")
    excluded_roots = baseline.get("excluded_roots")
    if not isinstance(baseline_files, dict) or not isinstance(excluded_roots, list):
        raise ValueError("thread scope baseline is malformed")
    baseline_files = {str(key): str(value) for key, value in baseline_files.items()}
    current_files = specialist_scope_manifest(
        project, excluded_roots=[str(item) for item in excluded_roots]
    )
    actual_changed = {
        path
        for path in set(baseline_files) | set(current_files)
        if baseline_files.get(path) != current_files.get(path)
    }
    declared_raw = declared_thread_changed_paths(receipt_text)
    declared: set[str] = set()
    for raw_path in declared_raw:
        path = contained_project_path(project, raw_path, "receipt files_changed")
        declared.add(canonical_project_relative(project, path))
    host_attestation_prefix = FINAL_DELIVERY_HOST_ATTESTATION_ROOT.as_posix() + "/"
    forbidden_host_changes = sorted(
        rel_path
        for rel_path in actual_changed | declared
        if rel_path == FINAL_DELIVERY_HOST_ATTESTATION_ROOT.as_posix()
        or rel_path.startswith(host_attestation_prefix)
    )
    if forbidden_host_changes:
        raise ValueError(
            "worker changed host-only attestation path: "
            + ",".join(forbidden_host_changes)
        )
    mode = row.get("mode", "").strip()
    if mode in THREADOPS_EXECUTION_MODES:
        write_scope_values = [
            item.strip()
            for item in row.get("write_scope", "").split(";")
            if item.strip()
        ]
        if not write_scope_values:
            raise ValueError("execution worker write_scope is missing")
        scope_roots = [
            contained_project_path(project, item, "thread write_scope")
            for item in write_scope_values
        ]
        for rel_path in actual_changed | declared:
            path = contained_project_path(project, rel_path, "thread changed path")
            if not any(path == root or root in path.parents for root in scope_roots):
                raise ValueError(f"thread changed path outside write_scope: {rel_path}")
        if actual_changed != declared:
            missing = sorted(actual_changed - declared)
            invented = sorted(declared - actual_changed)
            raise ValueError(
                "thread changed-path proof mismatch: "
                f"undeclared={','.join(missing) or 'none'}; "
                f"not_observed={','.join(invented) or 'none'}"
            )
        if decision in {"ADOPT", "PARTIAL_ADOPT"} and not actual_changed:
            raise ValueError("execution worker adoption requires observed file output")
        for rel_path in declared:
            if decision in {"ADOPT", "PARTIAL_ADOPT"} and not (
                project / rel_path
            ).is_file():
                raise ValueError(f"adopted worker output missing: {rel_path}")
    elif actual_changed or declared:
        raise ValueError("read-only/research thread changed project files")

    if decision in {"ADOPT", "PARTIAL_ADOPT"}:
        if not validation_result_is_success(receipt_text):
            raise ValueError("worker validation_result is not successful")
        if not receipt_loop_state_is_complete(receipt_text):
            raise ValueError("worker loop_state is not complete")
        payload = thread_receipt_json(receipt_text)
        if not payload:
            recommendation = receipt_field_text(
                receipt_text, "worker_recommendation"
            ).upper()
            if not any(
                value in recommendation for value in ["ADOPT", "PARTIAL_ADOPT"]
            ):
                raise ValueError("worker recommendation does not support adoption")
            dirty_state = receipt_field_text(receipt_text, "dirty_state_impact")
            if not dirty_state or re.search(
                r"\b(unknown|failed|outside|unbounded|pending|tbd)\b",
                dirty_state,
                re.IGNORECASE,
            ):
                raise ValueError("worker dirty_state_impact is unsafe or missing")
            receipt_cleanup = receipt_field_text(receipt_text, "cleanup_actions")
            if not receipt_cleanup or re.search(
                r"\b(none|not[_ ]?done|did not|failed|pending|tbd)\b",
                receipt_cleanup,
                re.IGNORECASE,
            ):
                raise ValueError("worker receipt cleanup_actions are incomplete")
    proof: dict[str, object] = {
        "protocol_id": "adco.thread-scope-proof",
        "version": "1.0",
        "work_id": row.get("work_id", ""),
        "lane_id": row.get("lane_id", ""),
        "lane_run_id": row.get("lane_run_id", ""),
        "thread_id": row.get("receipt_thread_id", "") or row.get("real_thread_id", ""),
        "baseline_path": baseline_rel,
        "baseline_sha256": row.get("scope_baseline_sha256", ""),
        "baseline_manifest_sha256": specialist_manifest_digest(baseline_files),
        "observed_manifest_sha256": specialist_manifest_digest(current_files),
        "observed_changed_paths": sorted(actual_changed),
        "receipt_declared_paths": sorted(declared),
        "decision": decision,
        "validation_success": validation_result_is_success(receipt_text),
        "cleanup_action": cleanup_action,
        "archived_at": archived_at,
        "created_at": now_iso(),
    }
    return proof


def reconcile_thread_receipt(
    project: Path,
    *,
    lane_id: str,
    work_id: str,
    receipt_path_value: str,
    adoption_decision: str,
    rejection_reason: str,
    reconciled_at: str,
    cleanup_action: str,
    archived_at: str = "",
) -> dict[str, str]:
    migrate_control_plane(project)
    parse_thread_timestamp(reconciled_at, "reconciled_at")
    if archived_at:
        parse_thread_timestamp(archived_at, "archived_at")
    decision = adoption_decision.strip().upper()
    if decision not in THREADOPS_ADOPTION_DECISIONS:
        raise ValueError(f"unknown adoption decision: {adoption_decision}")
    if decision != "ADOPT" and not rejection_reason.strip():
        raise ValueError("rejection_reason is required unless adoption_decision is ADOPT")
    if not cleanup_action.strip():
        raise ValueError("cleanup_action is required")

    registry_path = project / "AD-creative/orchestrator/thread_registry.csv"
    fields, rows = read_csv_rows(registry_path)
    target = thread_registry_target(rows, lane_id=lane_id, work_id=work_id)
    registered_receipt_rel = target.get("receipt_path", "").strip()
    registered_receipt_path = contained_project_path(
        project, registered_receipt_rel, "registered worker receipt"
    )
    receipt_path = Path(receipt_path_value).expanduser()
    if not receipt_path.is_absolute():
        receipt_path = project / receipt_path
    if not receipt_path.is_file():
        raise ValueError(f"receipt file does not exist: {receipt_path_value}")
    receipt_text = receipt_path.read_text(encoding="utf-8", errors="ignore")
    identities = receipt_thread_ids(receipt_text)
    expected_ids = {
        value
        for value in [
            target.get("real_thread_id", "").strip(),
            target.get("rescue_thread_id", "").strip(),
        ]
        if value
    }
    identity = identities[0] if len(identities) == 1 else ";".join(identities)
    target["receipt_thread_id"] = identity
    target["returned_at"] = reconciled_at
    target["updated_at"] = now_iso()
    if len(identities) != 1 or identities[0] not in expected_ids:
        target["receipt_status"] = "rejected"
        target["reconciliation_status"] = "rejected_evidence"
        target["convergence_state"] = "receipt_rejected"
        target["adoption_decision"] = "REJECT"
        target["rejection_reason"] = "invalid_worker_thread_id"
        write_csv_rows(registry_path, fields, rows)
        update_thread_agent_run(
            project,
            lane_id=lane_id,
            work_id=work_id,
            updates={
                "status": "rejected_evidence",
                "completed_at": reconciled_at,
                "proof_status": "invalid_worker_thread_id",
                "reconciliation_status": "rejected_evidence",
            },
        )
        log_path = append_thread_convergence_event(
            project,
            target,
            state="receipt_rejected",
            observed_at=reconciled_at,
            evidence="invalid_worker_thread_id",
        )
        return {
            "status": "rejected_evidence",
            "error": "invalid_worker_thread_id",
            "receipt_thread_id": identity,
            "log_path": safe_rel(project, log_path),
        }

    if receipt_path.resolve() != registered_receipt_path:
        semantic_error = "receipt path does not match the bound registry receipt_path"
    else:
        target["receipt_thread_id"] = identity
        try:
            scope_proof = validate_thread_receipt_scope_and_semantics(
                project,
                target,
                receipt_text,
                decision=decision,
                cleanup_action=cleanup_action,
                archived_at=archived_at,
            )
        except ValueError as exc:
            semantic_error = str(exc)
        else:
            semantic_error = ""
    if semantic_error:
        target["receipt_status"] = "rejected"
        target["reconciliation_status"] = "rejected_evidence"
        target["convergence_state"] = "receipt_rejected"
        target["adoption_decision"] = "REJECT"
        target["rejection_reason"] = semantic_error
        write_csv_rows(registry_path, fields, rows)
        update_thread_agent_run(
            project,
            lane_id=lane_id,
            work_id=work_id,
            updates={
                "status": "rejected_evidence",
                "completed_at": reconciled_at,
                "proof_status": "receipt_semantics_rejected",
                "reconciliation_status": "rejected_evidence",
            },
        )
        log_path = append_thread_convergence_event(
            project,
            target,
            state="receipt_rejected",
            observed_at=reconciled_at,
            evidence=semantic_error,
        )
        return {
            "status": "rejected_evidence",
            "error": semantic_error,
            "receipt_thread_id": identity,
            "log_path": safe_rel(project, log_path),
        }

    proof_path = (
        project
        / "AD-creative/orchestrator/thread_scope_proofs"
        / f"{safe_artifact_suffix(work_id)}_{safe_artifact_suffix(lane_id)}_{file_sha256(receipt_path)[:12]}.json"
    )
    agent_runs_path = project / "AD-creative/orchestrator/agent_runs.csv"
    convergence_log_path = (
        project
        / f"AD-creative/orchestrator/thread_convergence_{safe_artifact_suffix(work_id)}.md"
    )
    lane_plan_path = project / "AD-creative/orchestrator/thread_lane_plan.md"
    cleanup_plan_path = project / f"AD-creative/orchestrator/thread_cleanup_{work_id}.md"
    rollback_files = {
        path: path.read_bytes() if path.is_file() else None
        for path in [
            registry_path,
            agent_runs_path,
            convergence_log_path,
            proof_path,
            lane_plan_path,
            cleanup_plan_path,
        ]
    }
    write_json_object(proof_path, scope_proof)

    target["receipt_status"] = "received"
    target["reconciliation_status"] = (
        "reconciled" if decision in {"ADOPT", "PARTIAL_ADOPT"} else "rejected_evidence" if decision == "REJECT" else "blocked"
    )
    target["convergence_state"] = "receipt_received"
    target["adoption_decision"] = decision
    target["rejection_reason"] = rejection_reason.strip()
    target["reconciled_at"] = reconciled_at
    target["cleanup_action"] = cleanup_action.strip()
    target["archived"] = "true"
    target["archived_at"] = archived_at
    target["lifecycle_state"] = "archived"
    target["scope_proof_path"] = safe_rel(project, proof_path)
    target["scope_proof_sha256"] = file_sha256(proof_path)
    write_csv_rows(registry_path, fields, rows)
    update_thread_agent_run(
        project,
        lane_id=lane_id,
        work_id=work_id,
        updates={
            "status": "reconciled" if decision in {"ADOPT", "PARTIAL_ADOPT"} else decision.lower(),
            "completed_at": reconciled_at,
            "receipt_path": registered_receipt_rel,
            "proof_status": "receipt_identity_and_host_scope_verified",
            "reconciliation_status": target["reconciliation_status"],
        },
    )
    log_path = append_thread_convergence_event(
        project,
        target,
        state="receipt_received",
        observed_at=reconciled_at,
        evidence=f"main_adoption_decision={decision}; cleanup_action={cleanup_action}",
    )
    refresh_threadops_projections(project, work_id)
    post_validation_errors, _ = validate(project)
    if post_validation_errors:
        for path, content in rollback_files.items():
            if content is None:
                path.unlink(missing_ok=True)
            else:
                path.parent.mkdir(parents=True, exist_ok=True)
                path.write_bytes(content)
        raise ValueError(
            "project validation failed for candidate thread reconciliation; control-plane changes rolled back: "
            + "; ".join(post_validation_errors[:12])
        )
    return {
        "status": target["reconciliation_status"],
        "receipt_thread_id": identity,
        "adoption_decision": decision,
        "scope_proof_path": safe_rel(project, proof_path),
        "log_path": safe_rel(project, log_path),
    }


def client_outline_rows(project: Path) -> list[dict[str, str]]:
    _, rows = read_csv_rows(project / "AD-creative/client_review/client_outline.csv")
    return rows


def client_outline_content_issues(
    rows: list[dict[str, str]], *, require_confirmed_state: bool
) -> list[str]:
    issues: list[str] = []
    if not rows:
        return ["缺少客户可读文本框架：client_outline.csv 没有任何页。"]
    for row in rows:
        slide = row.get("slide_id", "<missing>")
        for field, min_chars in [
            ("page_title", 3),
            ("body_copy", 22),
            ("client_confirmation_point", 8),
            ("material_role", 4),
            ("visual_slot", 4),
            ("visual_asset_status", 4),
        ]:
            if not non_placeholder(row.get(field), min_chars=min_chars):
                issues.append(f"{slide} 缺少可客户阅读的 {field}。")
        body_copy = row.get("body_copy", "").strip()
        if len(body_copy) > CLIENT_OUTLINE_BODY_MAX_CHARS:
            issues.append(
                f"{slide} 正文过密：body_copy 超过 {CLIENT_OUTLINE_BODY_MAX_CHARS} 字，需拆页或降密度。"
            )
        visual_status = row.get("visual_asset_status", "").strip().lower()
        if visual_status and visual_status not in CLIENT_OUTLINE_VISUAL_STATUSES:
            issues.append(
                f"{slide} visual_asset_status 无效: {row.get('visual_asset_status')}。"
            )
        if visual_status in EXISTING_IMAGE_STATUSES and not row.get(
            "asset_ids", ""
        ).strip():
            issues.append(f"{slide} 标记已有图但 asset_ids 为空。")
        hits = find_client_language_hits(row_text(row))
        if hits:
            issues.append(
                f"{slide} 客户文本框架含内部/执行侧词: {', '.join(hits[:8])}"
            )
        if require_confirmed_state:
            if row.get("visibility", "").strip().lower() not in CLIENT_VISIBLE_VALUES:
                issues.append(f"{slide} 文本仍是 pending/internal，尚未确认进入客户版。")
            if row.get("status", "").strip().lower() not in {
                "ready",
                "approved",
                "done",
            }:
                issues.append(f"{slide} status 不是 ready/approved/done。")
    return issues


def client_outline_confirmed_content_sha256(
    fields: list[str], rows: list[dict[str, str]]
) -> str:
    """Hash client-reviewed content while excluding host-owned workflow state."""
    content_fields = [
        field for field in fields if field not in {"visibility", "status"}
    ]
    payload = {
        "protocol_id": "adco.client-outline-confirmed-content",
        "version": "1.0",
        "fields": content_fields,
        "rows": [
            {field: row.get(field, "") for field in content_fields} for row in rows
        ],
    }
    return hashlib.sha256(
        json.dumps(
            payload,
            ensure_ascii=False,
            sort_keys=True,
            separators=(",", ":"),
        ).encode("utf-8")
    ).hexdigest()


def client_outline_confirmation_errors(project: Path) -> list[str]:
    outline_path = project / "AD-creative/client_review/client_outline.csv"
    receipt_path = project / CLIENT_OUTLINE_CONFIRMATION_REL
    if not receipt_path.is_file():
        return ["缺少 hash-bound client outline 人工确认 receipt。"]
    try:
        receipt = read_json_object(receipt_path, "client outline confirmation")
    except ValueError as exc:
        return [str(exc)]
    errors: list[str] = []
    if receipt.get("decision") != "approved_for_ppt":
        errors.append("client outline confirmation decision 不是 approved_for_ppt。")
    confirmed_by = str(receipt.get("confirmed_by", "")).strip()
    if confirmed_by.lower() in {
        "",
        "ad_creative_operator",
        "automation",
        "worker",
        "main controller",
    }:
        errors.append("client outline confirmation 缺少独立的人类确认者。")
    try:
        parse_thread_timestamp(
            str(receipt.get("confirmed_at", "")), "client_outline.confirmed_at"
        )
    except ValueError as exc:
        errors.append(str(exc))
    evidence_ref = str(receipt.get("evidence_ref", "")).strip()
    if not (
        evidence_ref.startswith("user_confirmation:")
        or evidence_ref.startswith("client_confirmation:")
    ):
        errors.append("client outline confirmation evidence_ref 必须绑定用户/客户确认。")
    if receipt.get("outline_path") != safe_rel(project, outline_path):
        errors.append("client outline confirmation outline_path 不匹配。")
    if not outline_path.is_file() or receipt.get("outline_sha256") != file_sha256(
        outline_path
    ):
        errors.append("client outline confirmation 已过期：outline hash 不匹配。")
    if not re.fullmatch(
        r"[0-9a-f]{64}", str(receipt.get("presented_outline_sha256", ""))
    ):
        errors.append("client outline confirmation 缺少确认前原始文件 hash。")
    if receipt.get("confirmation_basis") != (
        "exact_presented_file_and_canonical_content_excluding_host_state"
    ):
        errors.append("client outline confirmation confirmation_basis 不匹配。")
    if outline_path.is_file():
        fields, rows = read_csv_rows(outline_path)
        if receipt.get(
            "confirmed_content_sha256"
        ) != client_outline_confirmed_content_sha256(fields, rows):
            errors.append(
                "client outline confirmation 已过期：客户确认内容 digest 不匹配。"
            )
    return errors


def confirm_client_outline(
    project: Path,
    *,
    confirmed_by: str,
    confirmed_at: str,
    evidence_ref: str,
) -> Path:
    migrate_control_plane(project)
    parse_thread_timestamp(confirmed_at, "confirmed_at")
    if confirmed_by.strip().lower() in {
        "",
        "ad_creative_operator",
        "automation",
        "worker",
        "main controller",
    }:
        raise ValueError("confirmed_by must identify the human/client confirmer")
    if not (
        evidence_ref.startswith("user_confirmation:")
        or evidence_ref.startswith("client_confirmation:")
    ):
        raise ValueError("evidence_ref must start with user_confirmation: or client_confirmation:")
    outline_path = project / "AD-creative/client_review/client_outline.csv"
    fields, rows = read_csv_rows(outline_path)
    issues = client_outline_content_issues(rows, require_confirmed_state=False)
    if issues:
        raise ValueError("client outline content is not confirmable: " + "; ".join(issues[:8]))
    presented_outline_sha256 = file_sha256(outline_path)
    confirmed_content_sha256 = client_outline_confirmed_content_sha256(fields, rows)
    for row in rows:
        row["visibility"] = "client_visible_ready"
        row["status"] = "approved"
    write_csv_rows(outline_path, fields, rows)
    receipt_path = project / CLIENT_OUTLINE_CONFIRMATION_REL
    write_json_object(
        receipt_path,
        {
            "confirmation_id": "CLIENT-OUTLINE-CONFIRMATION-001",
            "outline_path": safe_rel(project, outline_path),
            "presented_outline_sha256": presented_outline_sha256,
            "confirmed_content_sha256": confirmed_content_sha256,
            "outline_sha256": file_sha256(outline_path),
            "confirmation_basis": (
                "exact_presented_file_and_canonical_content_excluding_host_state"
            ),
            "decision": "approved_for_ppt",
            "confirmed_by": confirmed_by.strip(),
            "confirmed_at": confirmed_at,
            "evidence_ref": evidence_ref,
            "scope": "current_client_outline_all_rows",
        },
    )
    return receipt_path


def review_client_outline(project: Path) -> tuple[str, list[str], Path]:
    migrate_control_plane(project)
    rows = client_outline_rows(project)
    issues = client_outline_content_issues(rows, require_confirmed_state=True)
    evidence = [f"outline_rows={len(rows)}"]
    confirmation_issues = client_outline_confirmation_errors(project)
    issues.extend(confirmation_issues)
    evidence.append(
        "outline_confirmation=" + ("valid" if not confirmation_issues else "missing_or_stale")
    )
    status = "PASS" if not issues else "BLOCKED"
    report_path = project / "AD-creative/gates/GATE-AUTO-CLIENT-OUTLINE-001_report.md"
    write_text(
        report_path,
        f"""# Client Outline Gate

status: {status}
checked_at: {now_iso()}
visibility: internal_only

## Evidence

{chr(10).join(f"- {item}" for item in evidence)}

## Blocking Issues

{chr(10).join(f"- {issue}" for issue in issues) or "- 无"}

## Rules

- PPT builder 前必须有客户可读文本框架。
- 每页必须有 page_title、body_copy、client_confirmation_point、material_role、visual_slot、visual_asset_status。
- 详细客户方案允许 22-45+ 页；Gate 约束每页低密度，不把详细方案压成短 pitch。
- visual_asset_status 必须说明画面是已有图、占位、待生成、无图或纯文字页。
- 客户可见页必须显式 ready/approved/done。
- PPT builder 前必须有绑定当前 client_outline.csv SHA-256 的人工确认 receipt。
""",
    )
    update_artifact(project, "ART-AUTO-CLIENT-OUTLINE-GATE", "client_outline_gate_report", safe_rel(project, report_path), "client_review", status="done" if status == "PASS" else "blocked", visibility="internal_only", gate_status=status)
    append_gate(project, "GATE-AUTO-CLIENT-OUTLINE-001", "client_review", status, "90" if status == "PASS" else "35", "ART-AUTO-CLIENT-OUTLINE-GATE", ";".join(issues[:8]), "补齐客户可读大纲后重跑 client-outline-gate。", "", "ready_for_ppt_builder" if status == "PASS" else "fix_client_outline", "ad_creative_operator")
    return status, issues, report_path


def text_hits_for_blocklist(label: str, text: str) -> list[str]:
    hits = find_client_language_hits(text)
    return [f"{label}: {hit}" for hit in hits]


def review_client_language(project: Path, extra_paths: list[Path] | None = None) -> tuple[str, list[str], Path]:
    migrate_control_plane(project)
    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    candidate_files = set(candidate_client_language_files(project, artifacts))
    forced_files = set(extra_paths or [])
    issues: list[str] = []
    for row in client_outline_rows(project):
        if row_is_client_visible(row):
            issues.extend(text_hits_for_blocklist(f"client_outline {row.get('slide_id')}", row_text(row)))
    for path in sorted(candidate_files):
        if not path.exists():
            continue
        try:
            text = client_language_text_for_path(path)
        except RuntimeError as exc:
            issues.append(f"{safe_rel(project, path)} 无法可靠提取客户文本: {exc}")
            continue
        if not text.strip():
            issues.append(f"{safe_rel(project, path)} 客户文本为空或不可解析")
            continue
        issues.extend(text_hits_for_blocklist(safe_rel(project, path), text))
    for path in sorted(forced_files - candidate_files):
        if not path.exists():
            continue
        if path.suffix.lower() in {".pptx", ".pdf"}:
            try:
                text = client_language_text_for_path(path)
            except RuntimeError as exc:
                issues.append(f"{safe_rel(project, path)} 无法可靠提取客户文本: {exc}")
                continue
        elif path.suffix.lower() in TEXT_CLIENT_SCAN_SUFFIXES:
            try:
                text = path.read_text(encoding="utf-8")
            except UnicodeDecodeError:
                text = path.read_text(encoding="utf-8", errors="ignore")
        else:
            continue
        if not text.strip():
            issues.append(f"{safe_rel(project, path)} 客户文本为空或不可解析")
            continue
        issues.extend(text_hits_for_blocklist(safe_rel(project, path), text))
    status = "PASS" if not issues else "BLOCKED"
    report_path = project / "AD-creative/gates/GATE-AUTO-CLIENT-LANGUAGE-001_report.md"
    write_text(
        report_path,
        f"""# Client Language Gate

status: {status}
checked_at: {now_iso()}
visibility: internal_only

## Scanned

- client_outline_rows: {sum(1 for row in client_outline_rows(project) if row_is_client_visible(row))}
- text_files: {len(candidate_files) + len(forced_files - candidate_files)}

## Blocking Issues

{chr(10).join(f"- {issue}" for issue in sorted(set(issues))[:80]) or "- 无"}

## Rules

- 客户稿不能出现 prompt/thread/worker/AI/gate/内部/执行过程/需确认等执行侧语言。
- 命中即阻断客户版导出。
""",
    )
    update_artifact(project, "ART-AUTO-CLIENT-LANGUAGE-GATE", "client_language_gate_report", safe_rel(project, report_path), "client_review", status="done" if status == "PASS" else "blocked", visibility="internal_only", gate_status=status)
    append_gate(project, "GATE-AUTO-CLIENT-LANGUAGE-001", "client_review", status, "90" if status == "PASS" else "30", "ART-AUTO-CLIENT-LANGUAGE-GATE", ";".join(sorted(set(issues))[:8]), "清理客户稿内部/执行侧语言后重跑 client-language-gate。", "", "ready_for_client_visual_review" if status == "PASS" else "rewrite_client_language", "ad_creative_operator")
    return status, issues, report_path


def refresh_asset_current_manifest(project: Path) -> tuple[list[dict[str, str]], Path]:
    migrate_control_plane(project)
    path = sync_asset_current_manifest(project)
    _, refreshed = read_csv_rows(path)
    return refreshed, path


def update_current_asset_metadata(
    project: Path,
    asset_ids: list[str],
    *,
    source: str = "",
    platform: str = "",
    conversation: str = "",
    qa_flags: str = "",
) -> None:
    if not asset_ids:
        return
    rows, path = refresh_asset_current_manifest(project)
    wanted = set(asset_ids)
    changed = False
    for row in rows:
        if row.get("asset_id", "") not in wanted:
            continue
        if source:
            row["source"] = source
        if platform:
            row["platform"] = platform
        if conversation:
            row["conversation"] = conversation
        row["local_file"] = row.get("local_file") or row.get("path", "")
        if qa_flags:
            row["qa_flags"] = qa_flags
        elif not row.get("qa_flags", "").strip():
            row["qa_flags"] = "registered_without_visual_approval"
        changed = True
    if changed:
        write_csv_rows(path, ASSET_CURRENT_FIELDS, rows)


def matching_asset_authorization(
    project: Path,
    asset_id: str,
    asset_sha256: str,
    rows: list[dict[str, str]],
) -> dict[str, str] | None:
    for row in rows:
        if row.get("asset_id", "").strip() != asset_id:
            continue
        if row.get("asset_sha256", "").strip() != asset_sha256:
            continue
        if row.get("approval_scope", "").strip().lower() not in {
            "client_review",
            "client_delivery",
            "client_visible",
        }:
            continue
        if row.get("status", "").strip().lower() != "approved" or row.get("revoked_at", "").strip():
            continue
        approved_by = row.get("approved_by", "").strip()
        if approved_by.lower() in {
            "",
            "ad_creative_operator",
            "automation",
            "worker",
            "main controller",
        }:
            continue
        evidence_ref = row.get("evidence_ref", "").strip()
        try:
            parse_thread_timestamp(row.get("approved_at", ""), "asset approved_at")
        except ValueError:
            continue
        if evidence_ref.startswith(
            ("user_confirmation:", "client_confirmation:")
        ):
            return row
        try:
            evidence_path = contained_project_path(
                project, evidence_ref, "asset authorization evidence_ref"
            )
        except ValueError:
            continue
        if (
            evidence_path.is_file()
            and row.get("evidence_sha256", "").strip() == file_sha256(evidence_path)
        ):
            return row
    return None


def review_visual_layout(
    project: Path,
    *,
    min_long_edge: int = 900,
    min_short_edge: int = 600,
    pptx_path: Path | None = None,
    preview_path: Path | None = None,
) -> tuple[str, list[str], Path]:
    current_manifest_path = project / "AD-creative/visual_assets/asset_current_manifest.csv"
    _, previous_current_rows = read_csv_rows(current_manifest_path)
    previous_current_by_id = {
        row.get("asset_id", "").strip(): row
        for row in previous_current_rows
        if row.get("asset_id", "").strip()
    }
    manifest, _ = refresh_asset_current_manifest(project)
    _, authorization_rows = read_csv_rows(
        project / "AD-creative/visual_assets/asset_authorizations.csv"
    )
    outline = client_outline_rows(project)
    assets_by_id = {row.get("asset_id", "").strip(): row for row in manifest if row.get("asset_id", "").strip()}
    slide_usage: dict[str, list[str]] = {}
    issues: list[str] = []
    warnings: list[str] = []
    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    package_paths = current_delivery_paths(project, artifacts)
    pptx_path = pptx_path or package_paths.get("pptx")
    preview_path = preview_path or package_paths.get("preview")
    if not pptx_path or not pptx_path.exists():
        issues.append("缺少 exact current PPTX，visual-layout-gate 不能验证真实页面。")
    else:
        try:
            pptx_stats = inspect_pptx(pptx_path)
            if not pptx_stats["editable"]:
                issues.append("exact current PPTX 缺少可编辑文本层。")
            issues.extend(pptx_layout_findings(pptx_path))
        except (OSError, zipfile.BadZipFile, ET.ParseError) as exc:
            issues.append(f"exact current PPTX 无法解析: {exc}")
    if not preview_path or not preview_path.exists():
        issues.append("缺少 exact current preview，visual-layout-gate 不能验证真实渲染。")
    else:
        width, height, image_format = probe_image(preview_path)
        if not width or not height or not image_format:
            issues.append("exact current preview 不是可解析图片。")
    if not outline:
        issues.append("缺少 client_outline.csv，无法判断页面阅读顺序和素材角色。")
    for row in outline:
        slide = row.get("slide_id", "<missing>")
        body = row.get("body_copy", "").strip()
        if not non_placeholder(body, min_chars=22):
            issues.append(f"{slide} 客户页正文过短或未成故事。")
        if len(body) > CLIENT_OUTLINE_BODY_MAX_CHARS:
            issues.append(f"{slide} 客户页正文过密，需拆页或降密度。")
        if not non_placeholder(row.get("material_role"), min_chars=4):
            issues.append(f"{slide} 缺少素材角色，无法判断素材服务哪段故事。")
        if not non_placeholder(row.get("visual_slot"), min_chars=4):
            issues.append(f"{slide} 缺少画面占位/已有图/待生成图说明。")
        visual_status = row.get("visual_asset_status", "").strip().lower()
        if visual_status not in CLIENT_OUTLINE_VISUAL_STATUSES:
            issues.append(f"{slide} visual_asset_status 无效或缺失: {row.get('visual_asset_status') or 'missing'}。")
        asset_ids = split_asset_refs(row.get("asset_ids"))
        if visual_status in EXISTING_IMAGE_STATUSES and not asset_ids:
            issues.append(f"{slide} 标记已有图但没有登记 asset_ids。")
        for asset_id in asset_ids:
            slide_usage.setdefault(asset_id, []).append(slide)
            if asset_id not in assets_by_id:
                issues.append(f"{slide} 引用了未登记资产 {asset_id}。")
                continue
            asset = assets_by_id[asset_id]
            rel_path = asset.get("path", "")
            path = project / rel_path
            if rel_path and path.suffix.lower() in GENERATED_IMAGE_SUFFIXES and path.exists():
                width, height, _ = probe_image(path)
                slot_text = row.get("visual_slot", "").lower()
                if width and height:
                    if re.search(r"竖屏|portrait|vertical|9:16", slot_text, flags=re.IGNORECASE) and width > height:
                        issues.append(f"{slide} 需要竖屏图但 {asset_id} 是横图: {width}x{height}。")
                    if re.search(r"横屏|landscape|horizontal|16:9", slot_text, flags=re.IGNORECASE) and height > width:
                        issues.append(f"{slide} 需要横屏图但 {asset_id} 是竖图: {width}x{height}。")
            qa_blob = " ".join([asset.get("qa_flags", ""), asset.get("notes", "")]).lower()
            if VISUAL_LAYOUT_RISK_PATTERN.search(qa_blob):
                issues.append(f"{slide} 使用的 {asset_id} QA flags/notes 命中视觉风险。")
        if VISUAL_LAYOUT_RISK_PATTERN.search(row_text(row).lower()):
            issues.append(f"{slide} outline 命中视觉/版式风险词。")
    for asset_id, slides in slide_usage.items():
        if len(set(slides)) <= 1:
            continue
        asset = assets_by_id.get(asset_id, {})
        reuse_context = " ".join(
            [
                asset.get("notes", ""),
                asset.get("qa_flags", ""),
                " ".join(row_text(row) for row in outline if row.get("slide_id") in slides),
            ]
        ).lower()
        if not re.search(r"intentional_reuse|repeat_ok|系列主视觉|贯穿主视觉|重复使用已确认", reuse_context):
            issues.append(f"{asset_id} 被多个页面使用但未标记 intentional_reuse/repeat_ok: {', '.join(sorted(set(slides)))}。")
    for asset in manifest:
        asset_id = asset.get("asset_id", "").strip() or "<missing asset_id>"
        rel_path = asset.get("path", "")
        path = project / rel_path if rel_path else Path()
        direct_client_use = normalized_bool(asset.get("direct_client_use"))
        usage_slides = sorted(
            set(split_asset_refs(asset.get("used_in_slide")) + slide_usage.get(asset_id, []))
        )
        client_used = direct_client_use or bool(usage_slides)
        if client_used:
            status_value = asset.get("status", "").strip().lower()
            if status_value in UNAVAILABLE_ASSET_STATUSES:
                issues.append(f"{asset_id} 用于客户页但 status={asset.get('status') or 'missing'} 不可用。")
            if not rel_path:
                issues.append(f"{asset_id} 用于客户页但缺少 asset path。")
            elif not path.exists():
                issues.append(f"{asset_id} 用于客户页但文件不存在: {rel_path}。")
            elif not path.is_file():
                issues.append(f"{asset_id} 用于客户页但 path 不是文件: {rel_path}。")
            else:
                actual_sha = file_sha256(path)
                recorded_sha = asset.get("sha256", "").strip()
                previous_sha = previous_current_by_id.get(asset_id, {}).get("sha256", "").strip()
                if not recorded_sha:
                    issues.append(f"{asset_id} 用于客户页但缺少 sha256。")
                elif recorded_sha != actual_sha:
                    issues.append(f"{asset_id} 用于客户页但 sha256 过期: manifest={recorded_sha} actual={actual_sha}。")
                elif previous_sha and previous_sha != actual_sha:
                    issues.append(f"{asset_id} 用于客户页但运行前 sha256 过期: manifest={previous_sha} actual={actual_sha}。")
            recorded_sha = asset.get("sha256", "").strip()
            if not matching_asset_authorization(
                project, asset_id, recorded_sha, authorization_rows
            ):
                slides = ",".join(usage_slides) or "direct_client_use"
                issues.append(
                    f"{asset_id} 用于客户页但缺少匹配 asset hash/scope 的独立授权 receipt: {slides}。"
                )
        if path.suffix.lower() in GENERATED_IMAGE_SUFFIXES and path.exists():
            width, height, image_format = probe_image(path)
            if width and height and (max(width, height) < min_long_edge or min(width, height) < min_short_edge):
                issues.append(f"{asset_id} 图像尺寸不足用于 PPT 主视觉: {width}x{height} {image_format}")
        if VISUAL_LAYOUT_RISK_PATTERN.search(row_text(asset).lower()):
            issues.append(f"{asset_id} manifest 命中视觉/版式风险词。")
        if client_used and not asset.get("qa_flags", "").strip():
            issues.append(f"{asset_id} 用于客户页但缺少 qa_flags。")
    status = "PASS" if not issues and not warnings else "PARTIAL_PASS" if not issues else "BLOCKED"
    report_path = project / "AD-creative/gates/GATE-AUTO-VISUAL-LAYOUT-001_report.md"
    write_text(
        report_path,
        f"""# Visual Layout Gate

status: {status}
checked_at: {now_iso()}
visibility: internal_only

## Evidence

- outline_rows: {len(outline)}
- current_manifest_assets: {len(manifest)}
- current_pptx: {safe_rel(project, pptx_path) if pptx_path else 'missing'}
- current_preview: {safe_rel(project, preview_path) if preview_path else 'missing'}
- min_long_edge: {min_long_edge}
- min_short_edge: {min_short_edge}

## Blocking Issues

{chr(10).join(f"- {issue}" for issue in issues) or "- 无"}

## Warnings

{chr(10).join(f"- {warning}" for warning in warnings) or "- 无"}

## Rules

- 检查图片拉伸、裁切、图像大小、页面拥挤、卡片套卡片、报告感、文字过短、客户阅读顺序。
- 检查图片与文案/画面槽位匹配、同图重复误用、竖屏/横屏比例不当。
- 没有 exact current PPTX 与真实 preview 时不得 PASS。
- `direct_client_use=yes` 的图片还必须有与文件 hash 绑定的授权证据。
""",
    )
    update_artifact(project, "ART-AUTO-VISUAL-LAYOUT-GATE", "visual_layout_gate_report", safe_rel(project, report_path), "ppt_gate", status="done" if status != "BLOCKED" else "blocked", visibility="internal_only", gate_status=status)
    append_gate(project, "GATE-AUTO-VISUAL-LAYOUT-001", "ppt_gate", status, "90" if status == "PASS" else "65" if status == "PARTIAL_PASS" else "35", "ART-AUTO-VISUAL-LAYOUT-GATE", ";".join(issues[:8]), ";".join(warnings[:8]) or "修正版式风险后重跑 visual-layout-gate。", "", "ready_for_client_pack_gate" if status != "BLOCKED" else "fix_visual_layout", "ad_creative_operator")
    return status, issues + warnings, report_path


def final_delivery_lock(project: Path) -> tuple[list[dict[str, str]], Path]:
    migrate_control_plane(project)
    return final_delivery_lock_snapshot(project, protected_value="yes")


def cleanup_category(path: Path) -> str:
    lower = str(path).lower()
    if "05_最终交付_finaldelivery".lower() in lower:
        return "protected_final_delivery"
    if "contact" in lower and "sheet" in lower:
        return "contact_sheet"
    if "cache" in lower or ".pptx-cache" in lower:
        return "cache"
    if "preview" in lower:
        return "preview"
    if "download" in lower or "tmp" in lower or "temp" in lower:
        return "temporary_download"
    if "exports" in lower or path.suffix.lower() in {".pptx", ".pdf"}:
        return "old_export"
    if "selected" in lower:
        return "important_crop_or_selected"
    if "generated" in lower or "imagegen" in lower or "grok" in lower or "chatgpt" in lower:
        return "generated_original"
    if "raw" in lower or "original" in lower:
        return "original"
    return "derived_or_unclassified"


def dedupe_audit(project: Path) -> tuple[list[dict[str, str]], Path, Path]:
    migrate_control_plane(project)
    rows: list[dict[str, str]] = []
    seen: dict[str, str] = {}
    roots = [
        project / "AD-creative/visual_assets",
        project / "AD-creative/ppt",
        project / "03_阶段成果_WorkInProgress",
        project / "04_客户审阅_ClientReview",
        project / "05_最终交付_FinalDelivery",
    ]
    for root in roots:
        if not root.exists():
            continue
        for path in sorted(root.rglob("*")):
            if not path.is_file():
                continue
            sha = file_sha256(path)
            rel_path = safe_rel(project, path)
            rows.append(
                {
                    "path": rel_path,
                    "sha256": sha,
                    "category": cleanup_category(path),
                    "duplicate_of": seen.get(sha, ""),
                    "recommended_action": "protect" if cleanup_category(path) == "protected_final_delivery" else "review",
                }
            )
            seen.setdefault(sha, rel_path)
    csv_path = project / "AD-creative/gates/dedupe_audit.csv"
    write_csv_rows(csv_path, ["path", "sha256", "category", "duplicate_of", "recommended_action"], rows)
    report_path = project / "AD-creative/gates/GATE-AUTO-DEDUPE-AUDIT-001_report.md"
    duplicates = [row for row in rows if row["duplicate_of"]]
    write_text(
        report_path,
        f"""# Dedupe Audit

status: PASS
checked_at: {now_iso()}
visibility: internal_only

## Evidence

- files_scanned: {len(rows)}
- duplicates_by_hash: {len(duplicates)}
- csv: {safe_rel(project, csv_path)}

## Rule

This command does not delete files. It classifies originals, important crops, derived images, old exports, cache, previews, contact sheets, and protected final delivery files for human cleanup planning.
""",
    )
    update_artifact(project, "ART-AUTO-DEDUPE-AUDIT", "dedupe_audit_report", safe_rel(project, report_path), "workspace_hygiene", visibility="internal_only", gate_status="PASS")
    append_gate(project, "GATE-AUTO-DEDUPE-AUDIT-001", "workspace_hygiene", "PASS", "90", "ART-AUTO-DEDUPE-AUDIT", "", "Use cleanup-plan for non-destructive cleanup decisions.", "", "ready_for_cleanup_plan", "ad_creative_operator")
    return rows, csv_path, report_path


def cleanup_plan(project: Path) -> tuple[Path, list[str]]:
    locked, lock_path = final_delivery_lock(project)
    audit_rows, audit_csv, _ = dedupe_audit(project)
    actions: list[str] = []
    for row in audit_rows:
        category = row["category"]
        if category == "protected_final_delivery":
            action = "LOCKED_DO_NOT_MOVE_OR_DELETE"
        elif row["duplicate_of"]:
            action = "REVIEW_DUPLICATE_KEEP_BEST_SOURCE"
        elif category in {"cache", "preview", "contact_sheet"}:
            action = "REVIEW_CAN_ARCHIVE_AFTER_CONFIRMATION"
        else:
            action = "KEEP_OR_REVIEW"
        actions.append(f"{row['path']} => {action}")
    plan_path = project / "AD-creative/gates/CLEANUP-PLAN.md"
    write_text(
        plan_path,
        f"""# Cleanup Plan

status: REVIEW_ONLY
checked_at: {now_iso()}
visibility: internal_only

## Inputs

- dedupe_audit: {safe_rel(project, audit_csv)}
- final_delivery_lock: {safe_rel(project, lock_path)}
- protected_final_delivery_files: {len(locked)}

## Proposed Actions

{chr(10).join(f"- {action}" for action in actions[:200]) or "- 无文件需要分类。"}

## Rule

This plan never deletes, moves, or overwrites files. `05_最终交付_FinalDelivery` files are protected by default and only hash-registered.
""",
    )
    update_artifact(project, "ART-AUTO-CLEANUP-PLAN", "cleanup_plan", safe_rel(project, plan_path), "workspace_hygiene", visibility="internal_only", gate_status="PASS")
    return plan_path, actions


def latest_gate_status(project: Path, gate_id: str) -> str:
    _, gates = read_csv_rows(project / "AD-creative/orchestrator/gate_log.csv")
    for row in reversed(gates):
        if row.get("gate_id") == gate_id:
            return row.get("status", "").strip().upper()
    return ""


def classify_cleanup_path(project: Path, path: Path) -> str:
    rel = safe_rel(project, path)
    lowered = rel.lower()
    suffix = path.suffix.lower()
    if "/05_" in f"/{rel}" or rel.startswith("05_最终交付_FinalDelivery/"):
        return "protected_final_delivery"
    if any(part in path.parts for part in ("__pycache__", ".pytest_cache", ".mypy_cache", ".ruff_cache")) or suffix in {".pyc", ".pyo"}:
        return "cache"
    if "contact" in lowered and suffix in GENERATED_IMAGE_SUFFIXES:
        return "contact_sheet"
    if "preview" in lowered and suffix in GENERATED_IMAGE_SUFFIXES:
        return "preview"
    if "version_archive" in lowered or re.search(r"(old|backup|archive|v\d{1,3})", lowered) and suffix in {".pptx", ".pdf"}:
        return "old_export"
    if "/visual_assets/raw/" in lowered:
        return "original_image"
    if "/visual_assets/selected/" in lowered or "crop" in lowered or "裁切" in rel:
        return "important_crop_or_selected"
    if suffix in GENERATED_IMAGE_SUFFIXES:
        return "derived_image"
    return "other"


def cleanup_file_inventory(project: Path) -> list[dict[str, str]]:
    rows: list[dict[str, str]] = []
    for path in sorted(project.rglob("*")):
        if ".git" in path.parts or not path.is_file():
            continue
        try:
            size = path.stat().st_size
        except OSError:
            continue
        category = classify_cleanup_path(project, path)
        sha = ""
        if size <= 80 * 1024 * 1024 and category != "cache":
            try:
                sha = file_sha256(path)
            except OSError:
                sha = ""
        rows.append(
            {
                "path": safe_rel(project, path),
                "category": category,
                "sha256": sha,
                "size_bytes": str(size),
                "protected": "true" if category == "protected_final_delivery" else "false",
            }
        )
    return rows


def write_dedupe_audit(project: Path) -> tuple[str, list[str], Path]:
    register_final_delivery_locks(project)
    inventory = cleanup_file_inventory(project)
    by_hash: dict[str, list[dict[str, str]]] = {}
    for row in inventory:
        sha = row.get("sha256", "")
        if sha:
            by_hash.setdefault(sha, []).append(row)
    duplicate_groups = [rows for rows in by_hash.values() if len(rows) > 1]
    issues = [
        "duplicate group includes protected final delivery file; review manually only"
        for rows in duplicate_groups
        if any(row.get("protected") == "true" for row in rows)
    ]
    category_counts: dict[str, int] = {}
    for row in inventory:
        category_counts[row["category"]] = category_counts.get(row["category"], 0) + 1
    duplicate_text = []
    for index, rows in enumerate(duplicate_groups[:40], start=1):
        duplicate_text.append(f"### DUP-{index:03d}")
        duplicate_text.extend(
            f"- {row['category']} | protected={row['protected']} | {row['path']}"
            for row in rows
        )
        duplicate_text.append("")
    report_path = project / "AD-creative/orchestrator/dedupe_audit.md"
    write_text(
        report_path,
        f"""# Dedupe Audit

status: {"CHECK" if issues else "PASS"}
visibility: internal_only
checked_at: {now_iso()}
delete_action: none

## Category Counts

{chr(10).join(f"- {key}: {value}" for key, value in sorted(category_counts.items())) or "- none"}

## Duplicate Hash Groups

{chr(10).join(duplicate_text).rstrip() or "- none"}

## Issues

{chr(10).join(f"- {issue}" for issue in issues) or "- 无"}

## Rule

本报告只分类和估算重复，不删除文件。清理必须通过 cleanup-plan，且 `05_最终交付_FinalDelivery` 下用户手动放入的 PPT/PDF 默认 protected。
""",
    )
    update_artifact(
        project,
        "ART-AUTO-DEDUPE-AUDIT",
        "dedupe_audit_report",
        safe_rel(project, report_path),
        "workspace_hygiene",
        visibility="internal_only",
        gate_status="CHECK" if issues else "PASS",
    )
    return ("CHECK" if issues else "PASS"), issues, report_path


def write_cleanup_plan(project: Path) -> tuple[str, list[str], Path]:
    lock_path = register_final_delivery_locks(project)
    inventory = cleanup_file_inventory(project)
    recommendations: list[str] = []
    for category, action in [
        ("protected_final_delivery", "register only; do not move, overwrite, or delete"),
        ("original_image", "keep as source of truth unless manually superseded"),
        ("important_crop_or_selected", "keep if referenced by slide/client outline; otherwise mark for manual review"),
        ("derived_image", "keep if selected or referenced; otherwise candidate for archive"),
        ("old_export", "keep in version_archive or move only after hash registration"),
        ("cache", "safe candidate for automated cleanup after confirmation"),
        ("preview", "regenerate only after current PPT/PDF hash is registered"),
        ("contact_sheet", "archive separately; never use as client-visible image"),
    ]:
        count = sum(1 for row in inventory if row["category"] == category)
        recommendations.append(f"- {category}: {count} files -> {action}")
    report_path = project / "AD-creative/orchestrator/cleanup_plan.md"
    write_text(
        report_path,
        f"""# Cleanup Plan

status: planned
visibility: internal_only
created_at: {now_iso()}
delete_action: none
final_delivery_lock: {safe_rel(project, lock_path)}

## Recommendations

{chr(10).join(recommendations)}

## Protected Files

{chr(10).join(f"- {row['path']}" for row in inventory if row['protected'] == 'true') or "- none"}

## Rule

不要按 hash duplicate 直接删除。先按原图、重要裁切、派生图、旧导出、缓存、预览、contact sheet 分层；最终交付目录只登记 hash，不移动、不覆盖。
""",
    )
    update_artifact(
        project,
        "ART-AUTO-CLEANUP-PLAN",
        "cleanup_plan",
        safe_rel(project, report_path),
        "workspace_hygiene",
        visibility="internal_only",
        gate_status="PASS",
    )
    return "PASS", [], report_path


def normalize_thread_status(value: str | None) -> str:
    return (value or "").strip().lower().replace("-", "_").replace(" ", "_")


def review_thread_discipline(project: Path) -> tuple[str, list[str], Path]:
    _, rows = read_csv_rows(project / "AD-creative/orchestrator/thread_registry.csv")
    issues: list[str] = []
    warnings: list[str] = []
    evidence: list[str] = [f"thread_rows={len(rows)}"]
    for row in rows:
        thread_id = row.get("thread_id", "").strip()
        owner = thread_id or row.get("lane_id", "") or "<missing thread>"
        receipt_received = normalize_thread_status(row.get("receipt_status")) not in {"", "missing", "pending", "planned", "todo", "tbd"}
        lifecycle = normalize_thread_status(row.get("lifecycle_state"))
        adopted_or_consumed = receipt_received or lifecycle in {"returned", "reconciled", "archived", "complete", "completed"}
        if not adopted_or_consumed:
            continue
        if not thread_id or thread_id.startswith("planned:"):
            issues.append(f"{owner} 被当作已用 worker 但没有真实 thread_id。")
        receipt_path = row.get("receipt_path", "").strip()
        receipt_text = ""
        if not receipt_path:
            issues.append(f"{owner} 缺少 receipt_path。")
        else:
            try:
                path = contained_project_path(project, receipt_path, "thread receipt")
            except ValueError as exc:
                issues.append(f"{owner} {exc}")
                path = None
            if path is None or not path.is_file():
                issues.append(f"{owner} receipt 文件不存在: {receipt_path}")
            else:
                receipt_text = path.read_text(encoding="utf-8", errors="ignore")
        if receipt_text:
            payload = thread_receipt_json(receipt_text)
            has_worker_recommendation = bool(
                payload.get("specialist_recommendation")
                if payload and payload.get("protocol_id") == SPECIALIST_EXCHANGE_PROTOCOL
                else receipt_field_text(receipt_text, "worker_recommendation")
            )
            if not has_worker_recommendation:
                issues.append(f"{owner} receipt 缺少 worker_recommendation。")
        if not row.get("adoption_decision", "").strip():
            issues.append(f"{owner} registry 缺少 main adoption_decision。")
        if receipt_text:
            identities = receipt_thread_ids(receipt_text)
            expected_ids = {
                value
                for value in [
                    row.get("real_thread_id", "").strip(),
                    row.get("rescue_thread_id", "").strip(),
                ]
                if value
            }
            if len(identities) != 1 or identities[0] not in expected_ids:
                issues.append(
                    f"{owner} invalid_worker_thread_id: receipt={';'.join(identities) or 'missing'}。"
                )
        cleanup_values = " ".join(
            row.get(key, "")
            for key in ["cleanup_action", "cleanup_reason", "archived_at", "notes"]
        ).strip()
        if not cleanup_values:
            issues.append(f"{owner} 缺少 cleanup 记录。")
        archived = normalized_bool(row.get("archived"))
        if not archived or not row.get("archived_at", "").strip():
            issues.append(f"{owner} cleanup 尚未由 archived=true + archived_at 证实。")
    if not rows:
        warnings.append("thread_registry 为空；本项目未声称使用 Codex Threads 时可接受。")
    status = "PASS" if not issues and not warnings else "PARTIAL_PASS" if not issues else "BLOCKED"
    report_path = project / "AD-creative/gates/GATE-AUTO-THREAD-DISCIPLINE-001_report.md"
    write_text(
        report_path,
        f"""# Thread Discipline Gate

status: {status}
visibility: internal_only
checked_at: {now_iso()}

## Evidence

{chr(10).join(f"- {item}" for item in evidence)}

## Blocking Issues

{chr(10).join(f"- {issue}" for issue in issues) or "- 无"}

## Warnings

{chr(10).join(f"- {warning}" for warning in warnings) or "- 无"}

## Rule

没有真实 thread_id、dispatch record、receipt、adoption/rejection 和 cleanup 的 worker，不得被主线程声称采用。`planned:*`、标题创建和空 receipt 都不是执行证明。
""",
    )
    update_artifact(
        project,
        "ART-AUTO-THREAD-DISCIPLINE-GATE",
        "thread_discipline_gate_report",
        safe_rel(project, report_path),
        "threadops",
        visibility="internal_only",
        gate_status=status,
    )
    append_gate(
        project,
        "GATE-AUTO-THREAD-DISCIPLINE-001",
        "threadops",
        status,
        "90" if status == "PASS" else "65" if status == "PARTIAL_PASS" else "30",
        "ART-AUTO-THREAD-DISCIPLINE-GATE",
        ";".join(issues[:8]),
        ";".join(warnings[:8]) or "补齐真实 worker receipt 与 cleanup 后重跑。",
        "",
        "thread_evidence_ready" if status != "BLOCKED" else "fix_thread_evidence",
        "ad_creative_operator",
    )
    return status, issues + warnings, report_path


def normalize_stage(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", "_", value.strip().lower()).strip("_")


def markdown_cells(line: str) -> list[str]:
    return [cell.strip() for cell in line.strip().strip("|").split("|")]


def is_markdown_separator(cells: list[str]) -> bool:
    return bool(cells) and all(set(cell.replace(" ", "")) <= {"-", ":"} for cell in cells)


def text_value_after_colon(line: str) -> str:
    if ":" not in line:
        return ""
    return line.split(":", 1)[1].strip()


def adversarial_report_paths(project: Path) -> list[Path]:
    gate_dir = project / "AD-creative/gates"
    if not gate_dir.exists():
        return []
    return sorted(gate_dir.glob("ADVERSARIAL_REVIEW_*.md"))


def has_adversarial_row_for_stage(text: str, stage: str) -> bool:
    target_stage = normalize_stage(stage)
    metadata = {
        key.lower(): value.strip()
        for line in text.splitlines()
        if ":" in line
        for key, value in [line.split(":", 1)]
    }
    report_stage = normalize_stage(metadata.get("stage", ""))
    reviewer_id = metadata.get("reviewer_id", "").strip().lower()
    if report_stage != target_stage:
        return False
    if reviewer_id in {"", "main controller", "ad_creative_operator", "automation", "worker"}:
        return False
    if not metadata.get("reviewer_role", "").strip():
        return False
    if metadata.get("independent", "").strip().lower() not in {"true", "yes"}:
        return False
    try:
        parse_thread_timestamp(metadata.get("reviewed_at", ""), "reviewed_at")
    except ValueError:
        return False
    if not metadata.get("target_ref", "").strip():
        return False
    if not re.fullmatch(r"[0-9a-f]{64}", metadata.get("target_sha256", "").strip(), re.IGNORECASE):
        return False

    for line in text.splitlines():
        if not line.lstrip().startswith("|"):
            continue
        cells = markdown_cells(line)
        if is_markdown_separator(cells):
            continue
        joined = " ".join(cells).lower()
        if "objection" in joined or "反对意见" in joined:
            continue

        if len(cells) >= 9:
            objection = cells[2]
            rebuttal = cells[6]
            revision = cells[7]
            gate_status = cells[8]
            if (
                objection
                and rebuttal
                and revision
                and gate_status.strip().upper() in {"PASS", "APPROVED"}
            ):
                return True
        elif len(cells) >= 5:
            row_stage = normalize_stage(cells[0])
            if row_stage and row_stage != target_stage:
                continue
            objection = cells[1]
            rebuttal = cells[2]
            revision = cells[3]
            gate_status = cells[4]
            if (
                objection
                and rebuttal
                and revision
                and gate_status.strip().upper() in {"PASS", "APPROVED"}
            ):
                return True
    return False


def default_adversarial_targets(project: Path, stage: str) -> list[Path]:
    targets = {
        "creative": [
            "AD-creative/creative/creative_directions.md",
            "AD-creative/proposal_architecture/proposal_structure.md",
        ],
        "reference_research": ["AD-creative/references/reference_cards.csv"],
        "visual_review": [
            "AD-creative/visual_assets/asset_current_manifest.csv",
            "AD-creative/visual_assets/asset_manifest.csv",
        ],
        "film_quality": [
            "AD-creative/film/treatment_packet.md",
            "AD-creative/film/shot_list_storyboard_plan.md",
        ],
        "final_delivery": [
            "AD-creative/delivery/client_pack_binding.json",
        ],
    }
    return [
        project / rel_path
        for rel_path in targets.get(normalize_stage(stage), [])
        if (project / rel_path).is_file()
    ]


def write_adversarial_target_snapshot(
    project: Path,
    *,
    stage: str,
    payload: dict[str, object],
    target_digest: str,
) -> Path:
    path = (
        project
        / "AD-creative/gates/adversarial_targets"
        / f"{normalize_stage(stage)}_{target_digest}.json"
    )
    if not path.exists():
        write_json_object(
            path,
            {
                "protocol_id": "adco.adversarial-review-target",
                "version": "1.0",
                "stage": normalize_stage(stage),
                "target_digest": target_digest,
                "payload": payload,
                "created_at": now_iso(),
            },
        )
    return path


def adversarial_council_evidence(
    project: Path,
    stage: str,
    *,
    expected_targets: list[Path] | None = None,
) -> tuple[bool, list[str]]:
    evidence: list[str] = []
    allowed_targets = {
        path.resolve()
        for path in (
            expected_targets
            if expected_targets is not None
            else default_adversarial_targets(project, stage)
        )
        if path.is_file()
    }
    if not allowed_targets:
        return False, []
    for path in adversarial_report_paths(project):
        text = path.read_text(encoding="utf-8", errors="ignore")
        if not has_adversarial_row_for_stage(text, stage):
            continue
        target_match = re.search(r"(?im)^target_ref:\s*(.+?)\s*$", text)
        sha_match = re.search(r"(?im)^target_sha256:\s*([0-9a-f]{64})\s*$", text)
        if not target_match or not sha_match:
            continue
        try:
            target = contained_project_path(
                project, target_match.group(1).strip(), "adversarial target_ref"
            )
        except ValueError:
            continue
        if target not in allowed_targets:
            continue
        if not target.is_file() or file_sha256(target) != sha_match.group(1).lower():
            continue
        reviewed_match = re.search(r"(?im)^reviewed_at:\s*(.+?)\s*$", text)
        if not reviewed_match:
            continue
        try:
            reviewed_at = parse_thread_timestamp(
                reviewed_match.group(1), "reviewed_at"
            )
        except ValueError:
            continue
        if reviewed_at.timestamp() + 5 < target.stat().st_mtime:
            continue
        evidence.append(safe_rel(project, path))
    return bool(evidence), evidence


def enforce_adversarial_gate_policy(
    project: Path,
    stage: str,
    status: str,
    warnings: list[str],
    evidence: list[str],
    *,
    expected_targets: list[Path] | None = None,
) -> str:
    has_record, records = adversarial_council_evidence(
        project, stage, expected_targets=expected_targets
    )
    evidence.append(
        "adversarial_council="
        + (";".join(records) if has_record else "missing")
    )
    if status == "PASS" and not has_record:
        warnings.append("缺少有效反驳性议会记录，Gate 最高只能 PARTIAL_PASS。")
        return "PARTIAL_PASS"
    return status


def default_goal_id() -> str:
    return datetime.now().strftime("GOAL-%Y%m%d-%H%M%S")


def safe_artifact_suffix(value: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9]+", "-", value).strip("-").upper()
    return cleaned[:48] or "GOAL"


@dataclass(frozen=True)
class ThreadRoleSpec:
    role_id: str
    label: str
    professional_identity: str
    purpose: str
    default_environment: str
    default_write_scope: str
    read_first: tuple[str, ...]
    allowed_actions: tuple[str, ...]
    forbidden_actions: tuple[str, ...]
    required_output: tuple[str, ...]
    stop_condition: str
    validation_proof: str


THREADOPS_ROLE_ORDER = (
    "brand_client",
    "copy_creative",
    "film_director",
    "art_design",
    "producer_risk",
    "qa_review",
)
THREADOPS_DEFAULT_ROLES = ("brand_client", "copy_creative", "qa_review")
THREADOPS_REGISTRY_FIELDS = [
    "thread_id",
    "title",
    "role",
    "lane_id",
    "lane_run_id",
    "work_id",
    "lifecycle_state",
    "pinned",
    "archived",
    "created_at",
    "updated_at",
    "cleanup_action",
    "notes",
    "goal_id",
    "mode",
    "environment",
    "workspace_path",
    "write_scope",
    "professional_identity",
    "receipt_path",
    "receipt_status",
    "reconciliation_status",
    "assigned_at",
    "returned_at",
    "reconciled_at",
    "archived_at",
    "cleanup_reason",
    "last_seen_at",
    "duplicate_of",
    "planned_thread_id",
    "dispatch_status",
    "real_thread_id",
    "title_action",
    "title_verified_at",
    "dispatch_receipt_path",
    "dispatch_evidence",
    "scope_baseline_path",
    "scope_baseline_sha256",
    "scope_proof_path",
    "scope_proof_sha256",
    "rescue_dispatch_receipt_path",
    "rescue_dispatch_evidence",
    "convergence_state",
    "last_progress_at",
    "absolute_deadline_at",
    "bounded_extension_used",
    "extension_reason",
    "convergence_reminder_at",
    "convergence_reason",
    "rescue_count",
    "rescue_thread_id",
    "receipt_thread_id",
    "adoption_decision",
    "rejection_reason",
    "schema_state",
    "legacy_evidence_sha256",
    "legacy_quarantine_reason",
    "legacy_raw_ref",
]
THREADOPS_PROGRESS_STATES = {"active_with_progress", "finalizing_receipt"}
THREADOPS_OBSERVATION_STATES = {
    *THREADOPS_PROGRESS_STATES,
    "silent",
    "thread_not_converged",
    "rescue_dispatched",
}
THREADOPS_ADOPTION_DECISIONS = {"ADOPT", "PARTIAL_ADOPT", "REJECT", "BLOCKED"}
THREADOPS_REAL_THREAD_ID_PATTERN = re.compile(
    r"^[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}$",
    re.IGNORECASE,
)
THREADOPS_AGENT_RUN_FIELDS = [
    "run_id",
    "work_id",
    "agent_role",
    "status",
    "started_at",
    "completed_at",
    "input_files",
    "output_files",
    "gate_id",
    "summary",
    "next_action",
    "thread_id",
    "lane_id",
    "receipt_path",
    "proof_status",
    "reconciliation_status",
]
THREADOPS_ROLE_SPECS = {
    "brand_client": ThreadRoleSpec(
        "BRAND_CLIENT",
        "BrandClient",
        "Brand strategist and client-demand interpreter.",
        "Check client intent, misunderstanding risk, brand fit, and missing decisions.",
        "read_only",
        "receipt only",
        (
            "AD-creative/orchestrator/current_truth.md",
            "AD-creative/orchestrator/requirements.csv",
            "AD-creative/orchestrator/gaps.csv",
            "AD-creative/handoff/待你确认.md",
        ),
        (
            "Inspect project truth, requirements, gaps, and handoff files.",
            "Propose questions, blockers, and client-risk guardrails.",
        ),
        (
            "Do not edit current_truth, version_map, artifact_index, gate_log, final exports, or client-visible files.",
            "Do not approve client send or mark AI assets client-visible.",
        ),
        (
            "client intent summary",
            "misunderstanding risks",
            "blocking questions",
            "recommended next action",
        ),
        "All client-intent risks and missing decisions are listed with evidence.",
        "Receipt references exact project files and rows inspected.",
    ),
    "copy_creative": ThreadRoleSpec(
        "COPY_CREATIVE",
        "CopyCreative",
        "Advertising copy lead for concepts, campaign lines, and tone.",
        "Review concept clarity, message hierarchy, headlines, and proposal language.",
        "isolated_workspace",
        "AD-creative/workspaces/{work_id}/{lane_id}",
        (
            "AD-creative/creative/creative_directions.md",
            "AD-creative/copywriting/message_line_candidates.md",
            "AD-creative/client_review/client_review_outline.md",
            "AD-creative/client_review/slide_spec.md",
        ),
        (
            "Inspect copy and proposal structure.",
            "Draft rewrites only in the lane workspace and receipt path declared by write_scope.",
        ),
        (
            "Do not rewrite client-visible files directly unless the lane plan grants an exact write path.",
            "Do not remove story, brand mapping, timing, or key dialogue to simplify layout.",
        ),
        (
            "copy diagnosis",
            "candidate wording",
            "message hierarchy risks",
            "revision checklist",
        ),
        "Copy risks and proposed replacements are concrete enough for the master thread to merge.",
        "Receipt includes old/new wording or target section references.",
    ),
    "film_director": ThreadRoleSpec(
        "FILM_DIRECTOR",
        "FilmDirector",
        "Commercial film director focused on story logic and scene rhythm.",
        "Review treatment, story beats, timing, scene transitions, and dialogue labels.",
        "isolated_workspace",
        "AD-creative/workspaces/{work_id}/{lane_id}",
        (
            "AD-creative/film/treatment_packet.md",
            "AD-creative/film/shot_list_storyboard_plan.md",
            "AD-creative/film/production_constraints.md",
            "AD-creative/client_review/slide_spec.md",
        ),
        (
            "Inspect film/story files and slide spec.",
            "Draft story, timing, and shot-level fixes only in the lane workspace and receipt path declared by write_scope.",
        ),
        (
            "Do not export PPT/PDF.",
            "Do not invent unavailable production facts or talent approvals.",
        ),
        (
            "story logic review",
            "timing and rhythm notes",
            "shot/scene risks",
            "fix recommendations",
        ),
        "Story and scene issues are listed with file/section evidence.",
        "Receipt identifies affected scenes, slides, or treatment sections.",
    ),
    "art_design": ThreadRoleSpec(
        "ART_DESIGN",
        "ArtDesign",
        "Art director for visual system, layout, typography, and image quality.",
        "Review visual direction, asset slots, layout risks, and client-visible image safety.",
        "isolated_workspace",
        "AD-creative/workspaces/{work_id}/{lane_id}",
        (
            "AD-creative/ppt/ppt_visual_system.md",
            "AD-creative/visual_assets/asset_manifest.csv",
            "AD-creative/visual_assets/visual_asset_slots.csv",
            "AD-creative/visual_review/visual_qa_checklist.md",
        ),
        (
            "Inspect visual system, asset manifest, asset slots, and QA checklist.",
            "Draft visual risks and exact gate/checklist additions only in the lane workspace and receipt path declared by write_scope.",
        ),
        (
            "Do not mark generated assets client-visible.",
            "Do not replace source images or final exports.",
        ),
        (
            "visual system diagnosis",
            "asset safety risks",
            "layout QA issues",
            "gate recommendations",
        ),
        "Visual risks are tied to asset ids, slots, or checklist items.",
        "Receipt names asset/slot ids and required gate checks.",
    ),
    "producer_risk": ThreadRoleSpec(
        "PRODUCER_RISK",
        "ProducerRisk",
        "Producer and risk controller for feasibility, scope, legal, and platform constraints.",
        "Review feasibility, budget/scope exposure, rights risk, platform risk, and stop points.",
        "read_only",
        "receipt only",
        (
            "AD-creative/orchestrator/decisions.csv",
            "AD-creative/orchestrator/resolutions.csv",
            "AD-creative/film/production_constraints.md",
            "AD-creative/handoff/待你确认.md",
        ),
        (
            "Inspect decisions, resolutions, constraints, and pending confirmations.",
            "Return production risks, stop points, and owner assignments.",
        ),
        (
            "Do not approve paid, login, upload, send, identity, KYC, wallet, or private-account actions.",
            "Do not downgrade blockers without evidence.",
        ),
        (
            "feasibility risk list",
            "stop conditions",
            "owner/action matrix",
            "scope guardrails",
        ),
        "Every high-risk action has an owner, stop condition, and evidence path.",
        "Receipt lists blocker severity and affected delivery stage.",
    ),
    "qa_review": ThreadRoleSpec(
        "QA_REVIEW",
        "QAReview",
        "Independent adversarial reviewer and evidence checker.",
        "Cold-review the proposed handoff, validation evidence, thread cleanup, and missing tests.",
        "read_only",
        "receipt only",
        (
            "AD-creative/orchestrator/thread_lane_plan.md",
            "AD-creative/orchestrator/thread_registry.csv",
            "AD-creative/orchestrator/artifact_index.csv",
            "AD-creative/orchestrator/gate_log.csv",
        ),
        (
            "Inspect changed control-plane files and validation output.",
            "Return severity-ranked findings and final readiness gate.",
        ),
        (
            "Do not edit files.",
            "Do not mark final readiness without validation and cleanup proof.",
        ),
        (
            "severity-ranked findings",
            "missing validation",
            "thread cleanup issues",
            "approve/block recommendation",
        ),
        "Review is complete when blocker, risk, and validation gaps are enumerated.",
        "Receipt includes validation command/output references and cleanup status.",
    ),
}
THREADOPS_ROLE_ALIASES = {
    "brand": "brand_client",
    "brand_client": "brand_client",
    "client": "brand_client",
    "copy": "copy_creative",
    "copy_creative": "copy_creative",
    "creative": "copy_creative",
    "film": "film_director",
    "film_director": "film_director",
    "director": "film_director",
    "art": "art_design",
    "art_design": "art_design",
    "design": "art_design",
    "visual": "art_design",
    "producer": "producer_risk",
    "producer_risk": "producer_risk",
    "risk": "producer_risk",
    "qa": "qa_review",
    "qa_review": "qa_review",
    "review": "qa_review",
}

READ_ONLY_THREADOPS_ROLES = {"brand_client", "producer_risk"}
THREADOPS_DEFAULT_LOOP_MODE = "sequential"
THREADOPS_LOOP_MODES = ("sequential", "rfc_dag", "continuous_pr", "infinite")
THREADOPS_EXECUTION_MODES = {
    "execution_worker",
    "isolated_worktree_execution_worker",
}
THREADOPS_OBSERVATION_CONTRACT = (
    "status success|warning|error; summary one-line result; artifacts file paths or ids; "
    "next_actions actionable follow-ups; evidence_refs exact files, rows, or commands"
)
THREADOPS_ERROR_RECOVERY_CONTRACT = (
    "include root_cause_hint, safe_retry_instruction, and explicit_stop_condition; "
    "freeze after repeated same root cause or thread confusion"
)
THREADOPS_CONTEXT_BUDGET = (
    "role brief, lane plan, and declared read_first files only; load extra context only with a reason"
)
THREADOPS_ITERATION_BUDGET = (
    "max 3 internal passes per worker; no unbounded retries; main/control may lower the budget"
)
THREADOPS_REPLAY_TRIGGER = (
    "validation_result FAIL, missing required receipt field, stale evidence, out-of-scope edit, "
    "or reopened user/client feedback"
)
THREADOPS_FREEZE_TRIGGER = (
    "thread confusion, wrong-thread behavior, heat/cost spike, worker budget exceeded, "
    "same root cause twice, or cleanup request"
)
THREADOPS_PROMPT_ONLY_FORBIDDEN = (
    "prompt-only output is invalid for production workers; execution_worker must produce declared files "
    "or return BLOCKED with evidence"
)
THREADOPS_HELPER_DEFAULT_MODE = "none"
THREADOPS_HELPER_MODE = "stateless_secondary_helper"
THREADOPS_ALLOWED_HELPER_KINDS = (
    "image_generation",
    "ocr",
    "layout_lint",
    "asset_resize",
    "reference_extraction",
)
THREADOPS_HELPER_POLICY = (
    "optional bounded stateless helper/subagent-style local invocation inside a real Codex Thread worker; "
    "helper is stateless and cannot replace the worker/reviewer layer"
)
THREADOPS_HELPER_WRITE_BOUNDARY = (
    "helper has no thread_id, no thread_registry row, and no write_scope; "
    "any kept artifact must be written or imported by the L1 worker inside declared write_scope"
)
THREADOPS_HELPER_EVIDENCE_REQUIRED = (
    "when helper_mode is stateless_secondary_helper, receipt must include helper_invocations, "
    "helper_input_refs, helper_output_refs, helper_artifacts, helper_validation_result, "
    "helper_adopted_by_worker, helper_failure_reason if any, and worker_synthesis"
)
THREADOPS_HELPER_FAILURE_POLICY = (
    "helper failure cannot bypass worker validation; worker records failure reason, "
    "continues without helper evidence only when safe, or returns BLOCKED/PARTIAL"
)


def threadops_action_space(mode: str, write_scope: str) -> str:
    if mode == "execution_worker":
        return (
            "read declared inputs; write only declared workspace and receipt paths; "
            f"respect write_scope [{write_scope}]; no master truth, final export, send, upload, login, or paid action"
        )
    return (
        "read declared inputs; write receipt only; propose changes and blockers; "
        "do not edit artifacts, master truth, final exports, or client-visible files"
    )


def threadops_harness_contract(
    *,
    mode: str,
    write_scope: str,
    validation_proof: str,
    stop_condition: str,
) -> dict[str, str]:
    return {
        "action_space": threadops_action_space(mode, write_scope),
        "observation_contract": THREADOPS_OBSERVATION_CONTRACT,
        "error_recovery_contract": THREADOPS_ERROR_RECOVERY_CONTRACT,
        "context_budget": THREADOPS_CONTEXT_BUDGET,
        "iteration_budget": THREADOPS_ITERATION_BUDGET,
        "eval_gate": validation_proof,
        "adoption_decision": "pending_main_control_adoption",
        "rejection_reason": "required_if_adoption_decision_is_reject_or_partial",
        "loop_state": "planned",
        "replay_trigger": THREADOPS_REPLAY_TRIGGER,
        "freeze_trigger": THREADOPS_FREEZE_TRIGGER,
        "stop_condition": stop_condition,
    }


def format_threadops_contract(contract: dict[str, str]) -> str:
    return "\n".join(f"{key}: {value}" for key, value in contract.items())


def threadops_helper_contract() -> dict[str, str]:
    return {
        "helper_mode": THREADOPS_HELPER_DEFAULT_MODE,
        "helper_policy": THREADOPS_HELPER_POLICY,
        "allowed_helper_kinds": ",".join(THREADOPS_ALLOWED_HELPER_KINDS),
        "helper_write_boundary": THREADOPS_HELPER_WRITE_BOUNDARY,
        "helper_evidence_required": THREADOPS_HELPER_EVIDENCE_REQUIRED,
        "helper_failure_policy": THREADOPS_HELPER_FAILURE_POLICY,
        "helper_invocations": "none",
        "helper_input_refs": "none",
        "helper_output_refs": "none",
        "helper_artifacts": "none",
        "helper_validation_result": "not_applicable",
        "helper_adopted_by_worker": "no",
        "helper_failure_reason": "none",
        "worker_synthesis": "worker must synthesize and adopt/reject helper output before returning receipt",
    }


def format_threadops_helper_contract(contract: dict[str, str] | None = None) -> str:
    return format_threadops_contract(contract or threadops_helper_contract())


def threadops_loop_mode_contract() -> str:
    return """loop_mode: sequential
allowed_loop_modes: sequential,rfc_dag,continuous_pr,infinite
sequential: default; finish one bounded lane step before the next handoff
rfc_dag: use only when an RFC-style dependency graph is written in the lane plan
continuous_pr: use only for controlled PR/check cycles with explicit validation commands
infinite: bounded internal exploration only; must declare iteration_budget, freeze_trigger, replay_trigger, and stop_condition
safe_stop: stop before client-visible send, paid/private/upload actions, destructive edits, global install, validation failure, or missing receipt proof
replay_trigger: validation_result FAIL, missing receipt schema, stale evidence, out-of-scope edit, or reopened feedback
freeze_trigger: thread confusion, wrong thread, heat/cost spike, budget exceeded, repeated same root cause, or cleanup request
stop_condition: receipt reconciled, eval_gate passed or blocker recorded, adoption_decision recorded, and cleanup action planned"""


def threadops_lane_mode(role: str, spec: ThreadRoleSpec) -> str:
    if spec.default_environment in {"isolated_workspace", "worktree"}:
        return "execution_worker"
    if role == "brand_client":
        return "research"
    if role == "qa_review":
        return "cold_review"
    if role in READ_ONLY_THREADOPS_ROLES:
        return "read_only_review"
    return "read_only_review"


def threadops_workspace_path(work_id: str, lane_id: str, spec: ThreadRoleSpec) -> str:
    if spec.default_environment == "isolated_workspace":
        return f"AD-creative/workspaces/{work_id}/{lane_id}"
    if spec.default_environment == "worktree":
        return "declared_git_worktree_required"
    return "not_applicable_for_read_only"


def resolve_threadops_write_scope(work_id: str, lane_id: str, spec: ThreadRoleSpec) -> str:
    return spec.default_write_scope.format(work_id=work_id, lane_id=lane_id)


def markdown_table_cell(value: str) -> str:
    return re.sub(r"\s+", " ", value.replace("|", "/")).strip()


def parse_threadops_roles(raw_roles: str) -> list[str]:
    raw_values = [item.strip().lower() for item in raw_roles.split(",") if item.strip()]
    if not raw_values:
        raw_values = list(THREADOPS_DEFAULT_ROLES)
    roles: list[str] = []
    for raw in raw_values:
        role = THREADOPS_ROLE_ALIASES.get(raw)
        if not role:
            choices = ", ".join(THREADOPS_ROLE_ORDER)
            raise ValueError(f"unknown ThreadOps role: {raw}; choices: {choices}")
        if role not in roles:
            roles.append(role)
    return roles


def current_version_id(project: Path) -> str:
    path = project / "AD-creative/orchestrator/current_truth.md"
    if not path.exists():
        return ""
    return current_truth_value(path.read_text(encoding="utf-8"), "current_version_id")


def default_threadops_work_id(goal_id: str) -> str:
    return f"WORK-{safe_artifact_suffix(goal_id)}-THREADS"


def default_task_signature_id(goal_id: str) -> str:
    return f"TS-{safe_artifact_suffix(goal_id)}"


def threadops_role_brief_content(
    *,
    spec: ThreadRoleSpec,
    goal_id: str,
    work_id: str,
    title: str,
    objective: str,
    task_signature: dict[str, str],
    mode: str,
    write_scope: str,
) -> str:
    contract = threadops_harness_contract(
        mode=mode,
        write_scope=write_scope,
        validation_proof=spec.validation_proof,
        stop_condition=spec.stop_condition,
    )
    helper_contract = threadops_helper_contract()
    project_facts = "\n".join(
        f"- {key}: {value or 'TBD'}" for key, value in task_signature.items()
    )
    read_first = "\n".join(f"- `{item}`" for item in spec.read_first)
    allowed = "\n".join(f"- {item}" for item in spec.allowed_actions)
    forbidden = "\n".join(f"- {item}" for item in spec.forbidden_actions)
    output = "\n".join(f"- {item}" for item in spec.required_output)
    return f"""# {spec.role_id} Role Brief

goal_id: {goal_id}
work_id: {work_id}
goal_title: {title}
role_id: {spec.role_id}
role_label: {spec.label}
mode: {mode}
environment: {spec.default_environment}
write_scope: {write_scope}
loop_mode: {THREADOPS_DEFAULT_LOOP_MODE}

## Professional Identity

{spec.professional_identity}

## Role Objective

{spec.purpose}

Goal objective: {objective}

## Project Facts To Honor

{project_facts}

## Read First

{read_first}

## Allowed Actions

{allowed}

## Forbidden Actions

{forbidden}

## Output Contract

{output}

## Harness Contract

```text
{format_threadops_contract(contract)}
```

## Loop Mode Contract

```text
{threadops_loop_mode_contract()}
```

## Stateless Secondary Helper Contract

```text
{format_threadops_helper_contract(helper_contract)}
```

## Acceptance Evidence

- Stop condition: {spec.stop_condition}
- Validation proof: {spec.validation_proof}
- Main/control thread is the only merge owner.
- Final PPT/PDF export is not allowed from this lane.
- Helpers are not Codex Threads, substitute workers, registry rows, or write scopes.
"""


def threadops_worker_prompt_content(
    *,
    project: Path,
    spec: ThreadRoleSpec,
    goal_id: str,
    work_id: str,
    lane_id: str,
    title: str,
    objective: str,
    task_signature_id: str,
    task_signature: dict[str, str],
    role_brief_path: Path,
    receipt_path: Path,
    extra_read_first: list[str],
    mode: str,
    workspace_path: str,
    write_scope: str,
) -> str:
    contract = threadops_harness_contract(
        mode=mode,
        write_scope=write_scope,
        validation_proof=spec.validation_proof,
        stop_condition=spec.stop_condition,
    )
    helper_contract = threadops_helper_contract()
    read_first = list(dict.fromkeys([*spec.read_first, *extra_read_first]))
    read_first_text = "\n".join(f"- {item}" for item in read_first)
    allowed = "\n".join(f"- {item}" for item in spec.allowed_actions)
    forbidden = "\n".join(f"- {item}" for item in spec.forbidden_actions)
    output = "\n".join(f"- {item}" for item in spec.required_output)
    task_signature_text = "\n".join(
        f"- {key}: {value or 'TBD'}" for key, value in task_signature.items()
    )
    execution_receipt = (
        "- write_scope\n"
        "- files_changed (required; prompt-only output is invalid)\n"
        "- validation_result\n"
        "- dirty_state_impact\n"
        "- manifest_index_updates\n"
        "- QA/gate status\n"
        "- adoption/rejection recommendation\n"
        "- recurrence_guard\n"
        "- cleanup_actions"
        if mode == "execution_worker"
        else "- files_inspected\n- proposed_changes\n- no_files_changed_confirmation\n- adoption/rejection recommendation"
    )
    return f"""Repo: {project}
Mode: {mode}
Loop mode: {THREADOPS_DEFAULT_LOOP_MODE}
Task signature: {task_signature_id}
Agent role md: {safe_rel(project, role_brief_path)}
Agency selection id: none
Agency role brief: {safe_rel(project, role_brief_path)}
Agency source agents: none
Source staff count: 0
Goal: {title}
Goal id: {goal_id}
Goal objective: {objective}
Work item: {work_id}
Lane id: {lane_id}
Expected real thread id: TBD_DISPATCH_RECORD_REQUIRED
Absolute deadline: TBD_DISPATCH_DEADLINE_REQUIRED

Task signature details:
{task_signature_text}

Harness contract:
{format_threadops_contract(contract)}

Loop mode contract:
{threadops_loop_mode_contract()}

Codex Thread contract:
- Main/control must create or reuse a real Codex Thread, then run `adco dispatch-record` to bind that id into this prompt and receipt envelope before work starts.
- The worker must verify that `Expected real thread id` equals its own real thread id and return that exact id as `thread_id` in the receipt. `source_thread_id` is lineage only and cannot satisfy identity.
- Codex Threads are not subagents and must not be simulated by role-play inside the control thread.
- Writable work requires an execution_worker lane with isolated_workspace or worktree write_scope.
- If real Codex Thread tooling or isolated writable scope is unavailable for writable work, stop with TOOL_BLOCKED instead of falling back.

Stateless secondary helper invocation contract:
{format_threadops_helper_contract(helper_contract)}
- A helper invocation may be backed by a stateless helper/subagent-style call, but it is not a Codex Thread, not a substitute worker/reviewer, has no thread_id, has no thread_registry.csv row, and has no write_scope.
- Helper output must be recorded in this L1 worker receipt, adopted or rejected by the worker, then exposed to main/control through the worker receipt.
- Do not call real image generation, OCR, or other helpers unless the lane task explicitly needs a bounded local subtask; this repo defines the contract only.

Read first:
{read_first_text}

Environment: {spec.default_environment}
Workspace path: {workspace_path}
Allowed actions:
{allowed}

Forbidden actions:
{forbidden}
- Do not edit files unless the lane plan grants an exact write path.
- Do not update current_truth, version_map, artifact_index, gate_log, or final exports.
- Do not send, upload, purchase, log in, or use private accounts.

Write scope: {write_scope}
Receipt path: {safe_rel(project, receipt_path)}
Stop condition: {spec.stop_condition}
Merge owner: main/control thread
Final export allowed: no
Completion proof: {spec.validation_proof}

Return format:
- summary
- thread_id: exact Expected real thread id from this prompt
- write_scope: {write_scope}
{execution_receipt}
- evidence refs
- QA/gate status
- open questions
- worker_recommendation: ADOPT, PARTIAL_ADOPT, REJECT, or BLOCKED
- worker_rejection_reason: required when recommendation is not ADOPT
- loop_state
- replay_trigger
- freeze_trigger
- helper_mode
- helper_invocations
- helper_input_refs
- helper_output_refs
- helper_artifacts
- helper_validation_result
- helper_adopted_by_worker
- helper_failure_reason
- worker_synthesis
- workflow issue and recurrence guard, if found

Required output:
{output}
"""


def threadops_receipt_template_content(
    *,
    spec: ThreadRoleSpec,
    goal_id: str,
    work_id: str,
    lane_id: str,
    mode: str,
    workspace_path: str,
    write_scope: str,
) -> str:
    contract = threadops_harness_contract(
        mode=mode,
        write_scope=write_scope,
        validation_proof=spec.validation_proof,
        stop_condition=spec.stop_condition,
    )
    helper_contract = threadops_helper_contract()
    files_changed_rule = (
        "files_changed: required_non_empty_for_adopt; if no file changed, set status BLOCKED and explain the blocker"
        if mode == "execution_worker"
        else "files_changed: forbidden_for_read_only; confirm no files changed"
    )
    return f"""# {lane_id} Receipt

status: pending
goal_id: {goal_id}
work_id: {work_id}
role_id: {spec.role_id}
mode: {mode}
environment: {spec.default_environment}
workspace_path: {workspace_path}
write_scope: {write_scope}
thread_id: TBD_DISPATCH_RECORD_REQUIRED
absolute_deadline_at: TBD_DISPATCH_DEADLINE_REQUIRED
harness_id: HARN-{safe_artifact_suffix(work_id)}-{lane_id}
loop_mode: {THREADOPS_DEFAULT_LOOP_MODE}
prompt_only_output: invalid
production_receipt_rule: {THREADOPS_PROMPT_ONLY_FORBIDDEN}
helper_mode: {helper_contract["helper_mode"]}
helper_policy: {helper_contract["helper_policy"]}
allowed_helper_kinds: {helper_contract["allowed_helper_kinds"]}
helper_write_boundary: {helper_contract["helper_write_boundary"]}
helper_evidence_required: {helper_contract["helper_evidence_required"]}
helper_failure_policy: {helper_contract["helper_failure_policy"]}

## Harness Contract

```text
{format_threadops_contract(contract)}
```

## Loop Mode Contract

```text
{threadops_loop_mode_contract()}
```

## Stateless Secondary Helper Contract

```text
{format_threadops_helper_contract(helper_contract)}
```

## Helper Invocation Evidence

helper_mode: {THREADOPS_HELPER_DEFAULT_MODE}
helper_invocations: none
helper_input_refs: none
helper_output_refs: none
helper_artifacts: none
helper_validation_result: not_applicable
helper_adopted_by_worker: no
helper_failure_reason: none
worker_synthesis: none

## Observation

status: pending
summary: pending
artifacts: pending
next_actions: pending
evidence_refs: pending

## Files Changed

{files_changed_rule}

pending

## Validation Result

pending

## Dirty-State Impact

pending

## Manifest / Index Updates

pending

## QA / Gate Status

pending

## Open Questions

pending

## Recurrence Guard

pending

## Adoption / Rejection Recommendation

worker_recommendation: pending
worker_rejection_reason: pending_if_not_adopted
files_merged: pending_main_control_decision

## Cleanup Actions

pending

## Evidence

pending
"""


def render_thread_cleanup_note(
    project: Path,
    *,
    goal_id: str,
    work_id: str,
    roles: list[str],
    max_active: int,
) -> Path:
    output = project / f"AD-creative/orchestrator/thread_cleanup_{work_id}.md"
    role_lines = "\n".join(f"- {THREADOPS_ROLE_SPECS[role].role_id}: planned" for role in roles)
    write_text(
        output,
        f"""# Thread Cleanup Plan

goal_id: {goal_id}
work_id: {work_id}
status: planned
created_at: {now_iso()}

## Active Budget

- max_active_worker_reviewer: {max_active}
- main/control thread remains integration owner.
- worker/reviewer threads are archived after their receipts are reconciled.

## Planned Lanes

{role_lines}

## Audit Steps

1. Run `list_threads` with query `ADCO`.
2. Keep the main/control thread active.
3. Keep only lanes whose lifecycle is `created`, `assigned`, `running`, or `returned`.
4. Archive duplicate, stale, superseded, off-scope, or reconciled employee threads.
5. Record real thread ids and archive decisions in `AD-creative/orchestrator/thread_registry.csv`.
6. Before final status, confirm no consumed worker remains active.
""",
    )
    return output


def render_thread_execution_plan(
    project: Path,
    *,
    goal_id: str,
    title: str,
    objective: str,
    roles: list[str],
    work_id: str = "",
    task_signature_id: str = "",
    brand: str = "",
    product: str = "",
    talent_or_ip: str = "",
    platform_or_channel: str = "",
    deliverable: str = "",
    stage: str = "threadops",
    primary_risks: str = "",
    evidence_needed: str = "",
    master_thread_id: str = "",
    current_version: str = "",
    max_active: int = 3,
    extra_read_first: list[str] | None = None,
    force: bool = False,
) -> dict[str, object]:
    if not roles:
        raise ValueError("at least one ThreadOps role is required")
    if max_active < 1 or max_active > 3:
        raise ValueError("max_active_worker_reviewer must be between 1 and 3")
    if len(roles) > 5:
        raise ValueError("broad council over 5 roles requires explicit user approval")

    ensure_project(project)
    work_id = work_id or default_threadops_work_id(goal_id)
    task_signature_id = task_signature_id or default_task_signature_id(goal_id)
    current_version = current_version or current_version_id(project)
    extra_read_first = extra_read_first or []
    now = now_iso()
    task_signature = {
        "brand": brand,
        "product": product,
        "talent_or_ip": talent_or_ip,
        "platform_or_channel": platform_or_channel,
        "deliverable": deliverable,
        "stage": stage,
        "primary_risks": primary_risks,
        "evidence_needed": evidence_needed,
    }
    plan_path = project / "AD-creative/orchestrator/thread_lane_plan.md"
    if plan_path.exists() and not force:
        raise FileExistsError(f"thread lane plan already exists: {plan_path}; pass --force to replace it")
    ensure_csv_fields(project / "AD-creative/orchestrator/thread_registry.csv", THREADOPS_REGISTRY_FIELDS)
    ensure_csv_fields(project / "AD-creative/orchestrator/agent_runs.csv", THREADOPS_AGENT_RUN_FIELDS)

    role_dir = project / "AD-creative/agents/role_briefs"
    prompt_dir = project / f"AD-creative/agents/thread_prompts/{work_id}"
    receipt_dir = project / f"AD-creative/agents/receipts/{work_id}"
    role_dir.mkdir(parents=True, exist_ok=True)
    prompt_dir.mkdir(parents=True, exist_ok=True)
    receipt_dir.mkdir(parents=True, exist_ok=True)

    lane_rows: list[dict[str, str]] = []
    registry_rows: list[dict[str, str]] = []
    prompt_paths: list[Path] = []
    role_brief_paths: list[Path] = []
    receipt_paths: list[Path] = []

    for index, role in enumerate(roles, 1):
        spec = THREADOPS_ROLE_SPECS[role]
        lane_id = f"LANE-{index:02d}-{spec.role_id}"
        lane_run_id = f"{work_id}:{lane_id}"
        planned_thread_id = f"planned:{lane_run_id}"
        mode = threadops_lane_mode(role, spec)
        workspace_path = threadops_workspace_path(work_id, lane_id, spec)
        write_scope = resolve_threadops_write_scope(work_id, lane_id, spec)
        harness_contract = threadops_harness_contract(
            mode=mode,
            write_scope=write_scope,
            validation_proof=spec.validation_proof,
            stop_condition=spec.stop_condition,
        )
        helper_contract = threadops_helper_contract()
        run_id = f"RUN-{safe_artifact_suffix(work_id)}-{index:02d}"
        thread_title = f"ADCO 员工｜{spec.label}｜{title[:24] or work_id}"
        role_brief_path = role_dir / f"{spec.role_id}_{work_id}.md"
        prompt_path = prompt_dir / f"{lane_id}_prompt.md"
        receipt_path = receipt_dir / f"{lane_id}_receipt.md"
        role_brief_paths.append(role_brief_path)
        prompt_paths.append(prompt_path)
        receipt_paths.append(receipt_path)

        write_text(
            role_brief_path,
            threadops_role_brief_content(
                spec=spec,
                goal_id=goal_id,
                work_id=work_id,
                title=title,
                objective=objective,
                task_signature=task_signature,
                mode=mode,
                write_scope=write_scope,
            ),
        )
        write_text(
            prompt_path,
            threadops_worker_prompt_content(
                project=project,
                spec=spec,
                goal_id=goal_id,
                work_id=work_id,
                lane_id=lane_id,
                title=title,
                objective=objective,
                task_signature_id=task_signature_id,
                task_signature=task_signature,
                role_brief_path=role_brief_path,
                receipt_path=receipt_path,
                extra_read_first=extra_read_first,
                mode=mode,
                workspace_path=workspace_path,
                write_scope=write_scope,
            ),
        )
        write_text(
            receipt_path,
            threadops_receipt_template_content(
                spec=spec,
                goal_id=goal_id,
                work_id=work_id,
                lane_id=lane_id,
                mode=mode,
                workspace_path=workspace_path,
                write_scope=write_scope,
            ),
        )

        lane_rows.append(
            {
                "lane_id": lane_id,
                "lane_run_id": lane_run_id,
                "thread_id": planned_thread_id,
                "thread_title": thread_title,
                "thread_role": spec.role_id,
                "professional_identity": spec.professional_identity,
                "agent_role_md": safe_rel(project, role_brief_path),
                "agency_selection_id": "none",
                "agency_role_brief": safe_rel(project, role_brief_path),
                "agency_source_agents": "none",
                "source_staff_count": "0",
                "work_id": work_id,
                "purpose": spec.purpose,
                "spawn_mode": "create_or_reuse_codex_thread",
                "mode": mode,
                "environment": spec.default_environment,
                "workspace_path": workspace_path,
                "read_first": ";".join([*spec.read_first, *extra_read_first]),
                "write_scope": write_scope,
                "receipt_path": safe_rel(project, receipt_path),
                "receipt_status": "missing",
                "reconciliation_status": "pending",
                "action_space": harness_contract["action_space"],
                "observation_contract": harness_contract["observation_contract"],
                "error_recovery_contract": harness_contract["error_recovery_contract"],
                "context_budget": harness_contract["context_budget"],
                "iteration_budget": harness_contract["iteration_budget"],
                "eval_gate": harness_contract["eval_gate"],
                "adoption_decision": harness_contract["adoption_decision"],
                "rejection_reason": harness_contract["rejection_reason"],
                "loop_state": harness_contract["loop_state"],
                "replay_trigger": harness_contract["replay_trigger"],
                "freeze_trigger": harness_contract["freeze_trigger"],
                "helper_mode": helper_contract["helper_mode"],
                "helper_policy": helper_contract["helper_policy"],
                "allowed_helper_kinds": helper_contract["allowed_helper_kinds"],
                "helper_write_boundary": helper_contract["helper_write_boundary"],
                "helper_evidence_required": helper_contract["helper_evidence_required"],
                "helper_failure_policy": helper_contract["helper_failure_policy"],
                "helper_invocations": helper_contract["helper_invocations"],
                "helper_input_refs": helper_contract["helper_input_refs"],
                "helper_output_refs": helper_contract["helper_output_refs"],
                "helper_artifacts": helper_contract["helper_artifacts"],
                "helper_validation_result": helper_contract["helper_validation_result"],
                "helper_adopted_by_worker": helper_contract["helper_adopted_by_worker"],
                "helper_failure_reason": helper_contract["helper_failure_reason"],
                "worker_synthesis": helper_contract["worker_synthesis"],
                "stop_condition": spec.stop_condition,
                "validation_proof": spec.validation_proof,
                "required_output": ";".join(spec.required_output),
                "merge_owner": "main/control thread",
                "final_export_allowed": "no",
                "lifecycle_status": "planned",
                "cleanup_note": "archive after receipt is reconciled",
                "run_id": run_id,
            }
        )
        registry_rows.append(
            {
                "thread_id": planned_thread_id,
                "title": thread_title,
                "role": spec.role_id,
                "lane_id": lane_id,
                "lane_run_id": lane_run_id,
                "work_id": work_id,
                "lifecycle_state": "planned",
                "pinned": "false",
                "archived": "false",
                "created_at": now,
                "updated_at": now,
                "cleanup_action": "create_or_reuse_thread_then_archive_after_reconcile",
                "notes": f"prompt={safe_rel(project, prompt_path)}",
                "goal_id": goal_id,
                "mode": mode,
                "environment": spec.default_environment,
                "workspace_path": workspace_path,
                "write_scope": write_scope,
                "professional_identity": spec.professional_identity,
                "receipt_path": safe_rel(project, receipt_path),
                "receipt_status": "missing",
                "reconciliation_status": "pending",
                "assigned_at": "",
                "returned_at": "",
                "reconciled_at": "",
                "archived_at": "",
                "cleanup_reason": "",
                "last_seen_at": now,
                "duplicate_of": "",
                "planned_thread_id": planned_thread_id,
                "dispatch_status": "planned",
                "real_thread_id": "",
                "title_action": "",
                "title_verified_at": "",
                "dispatch_receipt_path": "",
                "dispatch_evidence": "",
                "convergence_state": "",
                "last_progress_at": "",
                "absolute_deadline_at": "",
                "bounded_extension_used": "false",
                "extension_reason": "",
                "convergence_reminder_at": "",
                "convergence_reason": "",
                "rescue_count": "0",
                "rescue_thread_id": "",
                "receipt_thread_id": "",
                "adoption_decision": "",
                "rejection_reason": "",
                "schema_state": "current",
                "legacy_evidence_sha256": "",
                "legacy_quarantine_reason": "",
                "legacy_raw_ref": "",
            }
        )
        update_or_append_csv_row(
            project / "AD-creative/orchestrator/agent_runs.csv",
            "run_id",
            {
                "run_id": run_id,
                "work_id": work_id,
                "agent_role": spec.role_id,
                "status": "planned",
                "started_at": "",
                "completed_at": "",
                "input_files": ";".join([safe_rel(project, role_brief_path), safe_rel(project, prompt_path)]),
                "output_files": safe_rel(project, receipt_path),
                "gate_id": "",
                "summary": f"Planned Codex Thread lane for {spec.role_id}.",
                "next_action": "create_or_reuse_codex_thread",
                "thread_id": planned_thread_id,
                "lane_id": lane_id,
                "receipt_path": safe_rel(project, receipt_path),
                "proof_status": "pending",
                "reconciliation_status": "pending",
            },
        )

    for row in registry_rows:
        update_or_append_csv_row(
            project / "AD-creative/orchestrator/thread_registry.csv",
            "thread_id",
            row,
        )

    cleanup_path = render_thread_cleanup_note(
        project,
        goal_id=goal_id,
        work_id=work_id,
        roles=roles,
        max_active=max_active,
    )

    lane_header = [
        "lane_id",
        "lane_run_id",
        "thread_id",
        "thread_title",
        "thread_role",
        "professional_identity",
        "agent_role_md",
        "agency_selection_id",
        "agency_role_brief",
        "agency_source_agents",
        "source_staff_count",
        "work_id",
        "purpose",
        "spawn_mode",
        "mode",
        "environment",
        "workspace_path",
        "read_first",
        "write_scope",
        "receipt_path",
        "receipt_status",
        "reconciliation_status",
        "stop_condition",
        "validation_proof",
        "required_output",
        "merge_owner",
        "final_export_allowed",
        "lifecycle_status",
        "cleanup_note",
    ]
    lane_table = [
        "| " + " | ".join(lane_header) + " |",
        "| " + " | ".join("---" for _ in lane_header) + " |",
    ]
    for row in lane_rows:
        lane_table.append(
            "| " + " | ".join(markdown_table_cell(row.get(column, "")) for column in lane_header) + " |"
        )

    harness_header = [
        "lane_id",
        "action_space",
        "observation_contract",
        "error_recovery_contract",
        "context_budget",
        "iteration_budget",
        "eval_gate",
        "adoption_decision",
        "rejection_reason",
        "loop_state",
        "replay_trigger",
        "freeze_trigger",
        "stop_condition",
    ]
    harness_table = [
        "| " + " | ".join(harness_header) + " |",
        "| " + " | ".join("---" for _ in harness_header) + " |",
    ]
    for row in lane_rows:
        harness_table.append(
            "| " + " | ".join(markdown_table_cell(row.get(column, "")) for column in harness_header) + " |"
        )

    helper_header = [
        "lane_id",
        "helper_mode",
        "helper_policy",
        "allowed_helper_kinds",
        "helper_write_boundary",
        "helper_evidence_required",
        "helper_failure_policy",
        "helper_invocations",
        "helper_input_refs",
        "helper_output_refs",
        "helper_artifacts",
        "helper_validation_result",
        "helper_adopted_by_worker",
        "helper_failure_reason",
        "worker_synthesis",
    ]
    helper_table = [
        "| " + " | ".join(helper_header) + " |",
        "| " + " | ".join("---" for _ in helper_header) + " |",
    ]
    for row in lane_rows:
        helper_table.append(
            "| " + " | ".join(markdown_table_cell(row.get(column, "")) for column in helper_header) + " |"
        )

    registry_header = [
        "thread_id",
        "title",
        "role",
        "mode",
        "lane_id",
        "lane_run_id",
        "lifecycle_state",
        "convergence_state",
        "absolute_deadline_at",
        "bounded_extension_used",
        "rescue_count",
        "receipt_thread_id",
        "adoption_decision",
        "receipt_status",
        "reconciliation_status",
        "pinned",
        "archived",
        "archived_at",
        "cleanup_action",
        "notes",
    ]
    registry_table = [
        "| " + " | ".join(registry_header) + " |",
        "| " + " | ".join("---" for _ in registry_header) + " |",
    ]
    for row in registry_rows:
        registry_table.append(
            "| " + " | ".join(markdown_table_cell(row.get(column, "")) for column in registry_header) + " |"
        )

    task_signature_block = "\n".join(f"{key}: {value or 'TBD'}" for key, value in task_signature.items())
    prompt_list = "\n".join(f"- `{safe_rel(project, path)}`" for path in prompt_paths)
    role_brief_list = "\n".join(f"- `{safe_rel(project, path)}`" for path in role_brief_paths)
    receipt_list = "\n".join(f"- `{safe_rel(project, path)}`" for path in receipt_paths)
    write_text(
        plan_path,
        f"""# Thread Lane Plan

goal_id: {goal_id}
run_id: {work_id}
created_at: {now}
master_thread_id: {master_thread_id or 'TBD'}
project_kind: ppt_material_project
task_signature_id: {task_signature_id}
current_version_id: {current_version or 'TBD'}
loop_mode: {THREADOPS_DEFAULT_LOOP_MODE}

## Goal

```text
user_outcome: {title}
success_standard: {objective}
blocked_if: thread budget exceeds {max_active}, worker edits undeclared files, validation fails, or cleanup proof is missing
```

## Task Signature

```text
{task_signature_block}
```

## Thread Budget

```text
max_active_worker_reviewer: {max_active}
broad_council_requires_user_approval_over: 5
main_thread_only_for: integration,current_truth,version_map,artifact_index,gate_log,final_export,final_status
freeze_trigger: user reports thread confusion / high heat / wrong thread / cleanup request
lane_modes: execution_worker requires exact write_scope; research/read_only_review/cold_review are read-only receipt lanes
poll_rule: fixed poll counts are inspection budgets, not automatic failure
progress_states: active_with_progress,finalizing_receipt
deadline_rule: one reasoned bounded extension with an absolute deadline
failure_rule: only silent/reminder evidence past the deadline may become thread_not_converged
rescue_limit: 1
```

## Loop Mode Contract

```text
{threadops_loop_mode_contract()}
```

## Harness Field Contract

```text
action_space: what the lane may read, write, call, and never do
observation_contract: required receipt observation shape with status, summary, artifacts, next_actions, and evidence_refs
error_recovery_contract: root cause hint, safe retry instruction, and explicit stop condition for failures
context_budget: maximum context allowed before the lane must ask main/control to add more files
iteration_budget: bounded internal pass count before stop, replay, or freeze
eval_gate: validation or gate evidence required before adoption
adoption_decision: ADOPT, PARTIAL_ADOPT, REJECT, or BLOCKED recorded by main/control
rejection_reason: required when adoption_decision is not ADOPT
loop_state: planned, running, blocked, replay_requested, frozen, returned, reconciled, or archived
replay_trigger: condition that forces a replay with tighter acceptance criteria
freeze_trigger: condition that stops new workers until cleanup/audit is complete
stop_condition: lane-specific completion boundary
```

## Stateless Secondary Helper Contract

```text
helper_mode: default none; set to stateless_secondary_helper only when a real worker invokes a bounded local helper
helper_policy: optional stateless helper inside an L1 Codex Thread worker; never a Codex Thread or substitute worker
allowed_helper_kinds: image_generation,ocr,layout_lint,asset_resize,reference_extraction
helper_write_boundary: helper has no thread_id, no thread_registry.csv row, and no write_scope; worker owns any kept artifacts
helper_evidence_required: helper_invocations, helper_input_refs, helper_output_refs, helper_artifacts, helper_validation_result, helper_adopted_by_worker, helper_failure_reason, worker_synthesis
helper_failure_policy: helper failure is recorded by the worker and cannot bypass worker validation or main/control adoption
worker_synthesis: L1 worker adopts/rejects helper output before returning receipt; main/control adopts through the worker receipt
```

## Invocation

1. Run `list_threads` with query `ADCO` and reuse/archive stale project worker threads before spawning.
2. Create at most {max_active} active worker/reviewer Codex Threads at one time.
3. Send each worker exactly one prompt from `AD-creative/agents/thread_prompts/{work_id}/`.
4. Run `adco dispatch-record` with work id, lane id, verified real thread id, and an absolute deadline before work starts.
5. Record readbacks with `adco thread-observe`; progress is not failure and permits at most one reasoned extension.
6. Require each worker to return a receipt matching `AD-creative/agents/receipts/{work_id}/` and its bound thread id.
7. Run `adco thread-reconcile`; main/control records adoption/rejection and confirms archive cleanup.

## Role Briefs

{role_brief_list}

## Worker Prompts

{prompt_list}

## Receipts

{receipt_list}

## Lane Harness Matrix

{chr(10).join(harness_table)}

## Lane Helper Matrix

{chr(10).join(helper_table)}

## Lane Map

{chr(10).join(lane_table)}

## Thread Registry

{chr(10).join(registry_table)}

## Master Thread Rules

```text
Owns current_truth, work_items, agent_runs, gate_log, artifact_index, and final answer.
Reads every worker result before accepting it.
Does not copy worker output blindly.
Runs or records the relevant validation/gate before advancing state.
Lists existing ADCO threads before creating a new employee thread.
Archives duplicate, stale, superseded, or reconciled employee threads.
Uses execution_worker for scoped production/editing work; uses read_only only for explorer, reviewer, research, or cold-review lanes.
Execution worker lanes must declare exact write_scope, files_changed, validation_result, dirty_state_impact, and cleanup_actions in the receipt.
Production worker receipts cannot be prompt-only; execution_worker lanes must produce declared files or return BLOCKED with evidence.
Treats worker_recommendation as advisory; records main adoption_decision and rejection_reason before merging or discarding worker output.
Rejects receipt identity mismatches as invalid_worker_thread_id / rejected_evidence.
Distinguishes active_with_progress, silent, and finalizing_receipt; allows one bounded extension and at most one rescue.
Allows stateless secondary helper invocations only inside real worker threads; helpers are not Codex Threads and have no thread_id, registry row, write_scope, or adoption authority.
Requires helper output to be synthesized and adopted/rejected by the worker, then adopted/rejected by main/control through the worker receipt.
Uses replay_trigger for failed eval gates and freeze_trigger for thread confusion, repeated root cause, or budget breach.
Does not allow more than {max_active} active workers in this plan.
Exports final PPT/PDF only from the main control thread.
```

## Reconciliation Log

| lane_id | main_adoption_decision | main_rejection_reason | files_merged | gate_id | archived_at | notes |
|---|---|---|---|---|---|---|

## Cleanup Log

| checked_at | active_threads_kept | threads_archived | duplicate_titles_fixed | notes |
|---|---|---|---|---|
| {now} | main/control only until worker creation | TBD | TBD | cleanup plan: `{safe_rel(project, cleanup_path)}` |
""",
    )

    artifact_plan_id = f"ART-{safe_artifact_suffix(work_id)}-LANE-PLAN"
    artifact_prompt_id = f"ART-{safe_artifact_suffix(work_id)}-PROMPTS"
    artifact_cleanup_id = f"ART-{safe_artifact_suffix(work_id)}-CLEANUP"
    update_artifact(
        project,
        artifact_plan_id,
        "thread_lane_plan",
        safe_rel(project, plan_path),
        "threadops",
        status="done",
        visibility="internal_only",
        gate_status="PASS",
    )
    update_artifact(
        project,
        artifact_prompt_id,
        "thread_prompt_pack",
        safe_rel(project, prompt_dir),
        "threadops",
        status="done",
        visibility="internal_only",
        gate_status="PASS",
    )
    update_artifact(
        project,
        artifact_cleanup_id,
        "thread_cleanup_plan",
        safe_rel(project, cleanup_path),
        "threadops",
        status="planned",
        visibility="internal_only",
        gate_status="PARTIAL_PASS",
    )
    update_or_append_csv_row(
        project / "AD-creative/orchestrator/work_items.csv",
        "work_id",
        {
            "work_id": work_id,
            "stage": "threadops",
            "title": title,
            "objective": objective,
            "owner_agent": "Main Controller",
            "status": "planned",
            "priority": "high",
            "input_refs": "",
            "output_artifacts": ";".join([artifact_plan_id, artifact_prompt_id, artifact_cleanup_id]),
            "linked_requirements": "",
            "linked_source_events": "",
            "linked_references": "",
            "linked_assets": "",
            "linked_slides": "",
            "blocked_by": "",
            "gate_required": "ThreadOps cleanup audit",
            "client_visibility": "internal_only",
            "created_at": now,
            "updated_at": now,
            "supersedes_work_id": "",
        },
    )
    append_event(
        project,
        {
            "event_id": f"EVT-{work_id}",
            "event_type": "thread_execution_plan_created",
            "created_at": now,
            "goal_id": goal_id,
            "work_id": work_id,
            "roles": [THREADOPS_ROLE_SPECS[role].role_id for role in roles],
            "thread_lane_plan": safe_rel(project, plan_path),
            "prompt_dir": safe_rel(project, prompt_dir),
            "cleanup_plan": safe_rel(project, cleanup_path),
        },
    )
    return {
        "goal_id": goal_id,
        "work_id": work_id,
        "task_signature_id": task_signature_id,
        "thread_lane_plan": plan_path,
        "prompt_dir": prompt_dir,
        "role_briefs": role_brief_paths,
        "prompts": prompt_paths,
        "receipts": receipt_paths,
        "cleanup_plan": cleanup_path,
        "roles": [THREADOPS_ROLE_SPECS[role].role_id for role in roles],
    }


def render_goal_iteration_plan(
    project: Path,
    *,
    goal_id: str,
    title: str,
    objective: str,
    owner: str,
    force: bool = False,
) -> Path:
    template = project / GOAL_PLAN_TEMPLATE_REL
    if not template.exists():
        template = TEMPLATE_ROOT / GOAL_PLAN_TEMPLATE_REL
    if not template.exists():
        raise FileNotFoundError(f"goal iteration template not found: {GOAL_PLAN_TEMPLATE_REL}")

    output_dir = project / GOAL_ITERATIONS_REL
    output_dir.mkdir(parents=True, exist_ok=True)
    output = output_dir / f"{goal_id}.md"
    if output.exists() and not force:
        raise FileExistsError(f"goal plan already exists: {output}")

    now = now_iso()
    write_text(
        output,
        f"""# Goal Iteration Plan

goal_id: {goal_id}
goal_title: {title}
status: active
owner: {owner}
created_at: {now}
updated_at: {now}
loop_mode: {THREADOPS_DEFAULT_LOOP_MODE}

## Objective

{objective}

## Scope

- 先完成客户可读文本框架与人工确认，再进入 PPT、客户包和发送准备。
- 品牌研究、素材准备与 DIRcreative 等 specialist exchange 只作为按需支线，不拥有 ADCO 控制面。
- 阶段完成后直接推进下一步低风险内部任务。
- 每个 Gate 前必须有反驳性议会记录。
- 默认不创建 Codex Thread；只有用户明确要求真实 Threads 或任务确需隔离执行证据时才启用。

## Loop Mode Contract

```text
{threadops_loop_mode_contract()}
```

## Non Scope

- 不自动发送客户稿。
- 不自动上传客户资料到外部平台。
- 不自动将 AI 图标记为客户可见。
- 不自动安装全局 Skill。

## Source Of Truth

- current_truth: `AD-creative/orchestrator/current_truth.md`
- requirements: `AD-creative/orchestrator/requirements.csv`
- gaps: `AD-creative/orchestrator/gaps.csv`
- work_items: `AD-creative/orchestrator/work_items.csv`
- gate_log: `AD-creative/orchestrator/gate_log.csv`
- handoff_board: `AD-creative/handoff/项目看板.md`
- pending_decisions: `AD-creative/handoff/待你确认.md`

## Execution Batches

| batch_id | objective | owner | inputs | outputs | gate | status | exit_condition |
|---|---|---|---|---|---|---|---|
| B1 | 建立本轮 goal 执行记录 | Main Controller | goal objective | goal iteration plan | Adversarial Council | done | plan written |
| B2 | 按双泳道推进下一阶段 | Main Controller | current_truth / work_items | updated artifacts | stage gate | queued | stage gate not BLOCKED |
| B3 | 验证并写入下一轮队列 | Operations Council | gate_log / artifacts | verification evidence | validation | queued | VALIDATION=PASS |

## Text-First Delivery Mapping

| phase | primary work | optional specialist / asset lane | dependency | exit_condition | next_phase |
|---|---|---|---|---|---|
| P0 | 需求、事实、缺口、FinalDelivery hash lock | 素材盘点 | source_events | 结构验证完成，FinalDelivery 未被覆盖/移动/删除 | P1 |
| P1 | 客户可读页级文本框架 | 品牌研究、参考与素材缺口 | P0 truth | client_outline 内容完整且不含内部执行语言 | P2 |
| P2 | 人工/客户确认当前 outline hash | 无 | P1 outline | hash-bound confirmation receipt 有效 | P3 |
| P3 | 创意方向、reference pack、proposal architecture | DIRcreative 可通过 neutral specialist exchange 返回内部 film package；素材仅 internal_only | P2 confirmation | 各自质量 Gate 非 BLOCKED，specialist 仅给 recommendation | P4 |
| P4 | 导出新的 immutable `client_review_vNNN.pptx` | 绑定 exact-current 素材 | P2-P3 | 新版本、editability proof、current_truth/version_map 一致 | P5 |
| P5 | 客户语言、视觉版式、素材授权、PPT 可编辑性 | 独立 adversarial review | P4 exact current package | 所有对应 Gate 有新鲜证据 | P6 |
| P6 | 生成 fresh client-pack manifest/binding | 无 | P5 | package digest 与 exact-current 输入一致且 Client Pack Gate PASS | P7 |
| P7 | 独立人工审阅 + 本轮发送授权 + send-readiness | 无 | P6 binding | receipt/authorization 均绑定同一 package digest；Gate 只判准备，不执行发送 | P8 |
| P8 | 反馈合并、supersede、下一版本计划 | specialist/asset 输出按需重新验证 | feedback | next_version_plan 写入且旧版本保持不变 | next goal |

## Adversarial Council

| stage | objection | rebuttal_path | revision_decision | gate_status |
|---|---|---|---|---|
| global | 自动连续执行可能跳过客户可见风险 | 检查授权策略、Gate 日志、待确认文件 | 只自动推进低风险内部动作；客户稿发送、AI 图客户可见、外部上传仍停 | PASS |

## Pause / Continue / Rollback Rules

continue_when: Gate 非 BLOCKED，且无客户可见/付费/上传/覆盖/安装风险。
pause_when: 客户/导演冲突、AI 图客户可见、客户稿发送、外部上传、全局安装、覆盖旧版本。
rollback_path: 回到产生断链的最近阶段，更新 affected artifacts 和 gate report。
resume_when: 阻塞项关闭，Gate report 和 decisions/resolutions 已更新。

## Verification

| check | method | threshold | result | evidence |
|---|---|---|---|---|
| goal plan exists | file check | exists | pending | `{safe_rel(project, output)}` |
| project validation | validate_project.py | VALIDATION=PASS | pending | run after execution |

## Execution Log

| time | action | artifact | result | next |
|---|---|---|---|---|
| {now} | created goal iteration plan | `{safe_rel(project, output)}` | done | run next batch |

## Next Iteration Queue

| priority | task | owner | trigger | exit_condition |
|---|---|---|---|---|
| P1 | 执行下一阶段 work item | Main Controller | current gate non-blocked | next gate report written |
| P1 | 补充阶段反驳性议会 | QA / Review Council | before each Gate | objection chain recorded |
""",
    )
    artifact_id = f"ART-{safe_artifact_suffix(goal_id)}"
    update_artifact(
        project,
        artifact_id,
        "goal_iteration_plan",
        safe_rel(project, output),
        "goal_planning",
        status="done",
        visibility="internal_only",
        gate_status="PASS",
    )
    append_event(
        project,
        {
            "event_id": f"EVT-{goal_id}",
            "event_type": "goal_iteration_plan_created",
            "created_at": now,
            "goal_id": goal_id,
            "path": safe_rel(project, output),
        },
    )
    return output


def write_manual_review_checklist(project: Path) -> Path:
    path = project / "AD-creative/delivery/manual_review_checklist.md"
    if not path.exists():
        write_text(
            path,
            f"""# Manual Review Checklist

status: pending_human_review
visibility: internal_only
created_at: {now_iso()}

## Search Sampling

- [ ] 随机打开 3 条客户可见候选参考，确认链接可访问。
- [ ] 确认每条客户可见参考不是 UGC 冒充官方来源。
- [ ] 确认 `do_not_copy` 限制已进入客户稿备注。

## Visual Taste

- [ ] 打开 `AD-creative/handoff/操作台.html` 的图片区，确认没有低质拼贴、contact sheet、假 logo。
- [ ] 对 selected 图片做审美判断：构图、光线、产品真实感、品牌气质、文字/标志风险。
- [ ] 客户可见图片必须有绑定 exact asset hash/scope/确认者/时间/证据的独立授权 receipt；notes token 不算授权。

## Client Pack

- [ ] 打开 `current_truth.md` 指向的 exact current PPTX，确认每页文本可编辑。
- [ ] 逐页读客户稿，确认没有内部注释、模拟标记、TODO/TBD、假案例。
- [ ] 最终发送前由负责人确认发送对象、附件、版本号。

## Rule

本清单是人工审阅入口，不替代 `client-pack-gate`、`visual-quality-gate`、`search-quality-gate`。
""",
        )
    update_artifact(
        project,
        "ART-AUTO-MANUAL-REVIEW-CHECKLIST",
        "manual_review_checklist",
        safe_rel(project, path),
        "final_delivery",
        status="pending_human_review",
        visibility="internal_only",
        gate_status="NOT_RUN",
    )
    return path


def manual_review_receipt_errors(
    project: Path,
    artifacts: list[dict[str, str]],
    *,
    package_digest: str = "",
) -> list[str]:
    receipt_path = project / "AD-creative/delivery/manual_review_receipt.json"
    if not receipt_path.exists():
        return ["缺少独立人工审阅 receipt: AD-creative/delivery/manual_review_receipt.json"]
    try:
        receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError) as exc:
        return [f"独立人工审阅 receipt 无法解析: {exc}"]
    errors: list[str] = []
    truth_path = project / "AD-creative/orchestrator/current_truth.md"
    truth = truth_path.read_text(encoding="utf-8") if truth_path.exists() else ""
    version_id = current_truth_value(truth, "current_version_id")
    pptx_artifact_id = current_truth_value(truth, "current_pptx_artifact_id")
    pptx_row = next(
        (row for row in artifacts if row.get("artifact_id") == pptx_artifact_id),
        None,
    )
    reviewer_id = str(receipt.get("reviewer_id", "")).strip().lower()
    if reviewer_id in {"", "main controller", "ad_creative_operator", "automation", "worker"}:
        errors.append("人工审阅 reviewer_id 缺失或不是独立审阅者")
    if receipt.get("independent") is not True:
        errors.append("人工审阅 receipt 未声明 independent=true")
    if not str(receipt.get("review_id", "")).strip():
        errors.append("人工审阅 receipt 缺少 review_id")
    if not str(receipt.get("reviewer_role", "")).strip():
        errors.append("人工审阅 receipt 缺少 reviewer_role")
    try:
        parse_thread_timestamp(str(receipt.get("reviewed_at", "")), "reviewed_at")
    except ValueError as exc:
        errors.append(str(exc))
    evidence_ref = str(receipt.get("evidence_ref", "")).strip()
    if not evidence_ref.startswith(
        ("user_confirmation:", "client_confirmation:", "review_record:")
    ):
        errors.append("人工审阅 receipt 缺少可追溯 evidence_ref")
    if receipt.get("decision") != "approved":
        errors.append("人工审阅 decision 不是 approved")
    if receipt.get("version_id") != version_id:
        errors.append("人工审阅 receipt version_id 不是 exact current version")
    if receipt.get("pptx_artifact_id") != pptx_artifact_id:
        errors.append("人工审阅 receipt PPTX artifact 不是 exact current")
    if not pptx_row or receipt.get("pptx_sha256") != pptx_row.get("sha256"):
        errors.append("人工审阅 receipt 未绑定 exact current PPTX hash")
    if not package_digest or receipt.get("package_digest") != package_digest:
        errors.append("人工审阅 receipt 未绑定 fresh exact-current package digest")
    checks = receipt.get("checks")
    required_checks = {
        "client_language",
        "visual_layout",
        "asset_authorization",
        "ppt_editability",
    }
    if not isinstance(checks, dict) or any(
        checks.get(key) is not True for key in required_checks
    ):
        errors.append("人工审阅 receipt 检查项不完整或存在未通过项")
    return errors


def send_authorization_errors(
    project: Path,
    artifacts: list[dict[str, str]],
    *,
    package_digest: str = "",
) -> list[str]:
    path = project / "AD-creative/delivery/send_authorization.json"
    if not path.exists():
        return ["缺少本轮发送授权: AD-creative/delivery/send_authorization.json"]
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        return [f"发送授权无法解析: {exc}"]
    errors: list[str] = []
    truth_path = project / "AD-creative/orchestrator/current_truth.md"
    truth = truth_path.read_text(encoding="utf-8") if truth_path.exists() else ""
    version_id = current_truth_value(truth, "current_version_id")
    pptx_artifact_id = current_truth_value(truth, "current_pptx_artifact_id")
    pptx_row = next(
        (row for row in artifacts if row.get("artifact_id") == pptx_artifact_id),
        None,
    )
    authorized_by = str(payload.get("authorized_by", "")).strip().lower()
    if authorized_by in {"", "ad_creative_operator", "automation", "worker", "main controller"}:
        errors.append("发送授权 authorized_by 缺失或来自执行面")
    if payload.get("decision") != "authorized":
        errors.append("发送授权 decision 不是 authorized")
    if not str(payload.get("authorization_id", "")).strip():
        errors.append("发送授权缺少 authorization_id")
    try:
        parse_thread_timestamp(
            str(payload.get("authorized_at", "")), "authorized_at"
        )
    except ValueError as exc:
        errors.append(str(exc))
    evidence_ref = str(payload.get("evidence_ref", "")).strip()
    if not evidence_ref.startswith(
        ("user_confirmation:", "client_confirmation:", "send_record:")
    ):
        errors.append("发送授权缺少可追溯 evidence_ref")
    if not str(payload.get("recipient_scope", "")).strip():
        errors.append("发送授权缺少 recipient_scope")
    if payload.get("version_id") != version_id:
        errors.append("发送授权 version_id 不是 exact current version")
    if payload.get("pptx_artifact_id") != pptx_artifact_id:
        errors.append("发送授权 PPTX artifact 不是 exact current")
    if not pptx_row or payload.get("pptx_sha256") != pptx_row.get("sha256"):
        errors.append("发送授权未绑定 exact current PPTX hash")
    if not package_digest or payload.get("package_digest") != package_digest:
        errors.append("发送授权未绑定 fresh exact-current package digest")
    return errors


def review_client_send_readiness(
    project: Path,
) -> tuple[str, list[str], Path]:
    issues: list[str] = []
    evidence: list[str] = []
    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    _, versions = read_csv_rows(project / "AD-creative/orchestrator/version_map.csv")
    _, feedback = read_csv_rows(project / "AD-creative/feedback/feedback_map.csv")
    validation_errors, _ = validate(project)
    issues.extend(validation_errors)
    issues.extend(validate_client_delivery_readiness(project, artifacts, versions, feedback))
    client_pack_status = latest_gate_status(project, "GATE-AUTO-CLIENT-PACK-001")
    evidence.append(f"client_pack_gate={client_pack_status or 'MISSING'}")
    if client_pack_status != "PASS":
        issues.append("client-pack-gate 不是 fresh PASS；先重跑当前包检查")
    binding_errors, package_digest = current_client_pack_binding_errors(
        project, artifacts
    )
    issues.extend(binding_errors)
    evidence.append(f"package_digest={package_digest or 'MISSING'}")
    evidence.append(f"client_pack_binding_errors={len(binding_errors)}")
    manual_errors = manual_review_receipt_errors(
        project, artifacts, package_digest=package_digest
    )
    authorization_errors = send_authorization_errors(
        project, artifacts, package_digest=package_digest
    )
    issues.extend(manual_errors)
    issues.extend(authorization_errors)
    evidence.append(f"manual_review_errors={len(manual_errors)}")
    evidence.append(f"send_authorization_errors={len(authorization_errors)}")
    status = "PASS" if not issues else "BLOCKED"
    report_path = project / "AD-creative/gates/GATE-AUTO-CLIENT-SEND-READINESS-001_report.md"
    write_text(
        report_path,
        f"""# Client Send Readiness Gate

status: {status}
visibility: internal_only
checked_at: {now_iso()}

## Evidence

{chr(10).join(f'- {item}' for item in evidence)}

## Blocking Issues

{chr(10).join(f'- {item}' for item in issues[:80]) or '- 无'}

## Scope

本 Gate 只判断 exact current package 是否具备结构、真实文件、哈希、客户语言、视觉、素材授权、独立人工审阅和本轮发送授权。它不会执行发送。
""",
    )
    update_artifact(
        project,
        "ART-AUTO-CLIENT-SEND-READINESS-GATE",
        "client_send_readiness_gate_report",
        safe_rel(project, report_path),
        "final_delivery",
        status="done" if status == "PASS" else "blocked",
        visibility="internal_only",
        gate_status=status,
    )
    binding_path = project / "AD-creative/delivery/client_pack_binding.json"
    gate_target = binding_path if binding_path.is_file() else report_path
    append_gate(
        project,
        "GATE-AUTO-CLIENT-SEND-READINESS-001",
        "final_delivery",
        status,
        "100" if status == "PASS" else "0",
        "ART-AUTO-CLIENT-SEND-READINESS-GATE",
        ";".join(issues[:8]),
        "补齐 exact current package 的独立审阅与本轮发送授权后重跑。",
        "",
        "send_authorized_but_not_sent" if status == "PASS" else "not_send_ready",
        "ad_creative_operator",
        target_ref=safe_rel(project, gate_target),
        target_sha256=file_sha256(gate_target),
    )
    return status, issues, report_path


def review_handoff_readiness(project: Path) -> tuple[str, list[str], list[str], Path]:
    blockers: list[str] = []
    warnings: list[str] = []
    evidence: list[str] = []

    errors, stats = validate(project)
    evidence.append(f"validation_errors={len(errors)}")
    if errors:
        blockers.extend(errors[:12])

    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    dashboard = render_dashboard(project)
    dashboard_issues = audit_dashboard(project)
    evidence.append(f"dashboard={safe_rel(project, dashboard)}")
    if dashboard_issues:
        blockers.extend(dashboard_issues)

    checklist = write_manual_review_checklist(project)
    evidence.append(f"manual_review_checklist={safe_rel(project, checklist)}")

    content_gate_visibility = {
        "GATE-THREE-COUNCIL-READINESS": {"PASS"},
        "GATE-AUTO-CLIENT-PACK-001": {"PASS"},
        "GATE-AUTO-VISUAL-QUALITY-001": {"PASS"},
        "GATE-AUTO-SEARCH-QUALITY-001": {"PASS", "PARTIAL_PASS"},
        "GATE-AUTO-REFERENCE-PACK-001": {"PASS", "PARTIAL_PASS"},
    }
    for gate_id, allowed in content_gate_visibility.items():
        status = latest_gate_status(project, gate_id)
        evidence.append(f"{gate_id}={status or 'MISSING'}")
        if not status:
            warnings.append(f"{gate_id} 尚未运行。")
        elif status not in allowed:
            warnings.append(f"{gate_id} status={status}；内部交接可继续，但不得据此宣称交付/发送就绪。")

    pptx, pptx_declared = current_pptx_path(project, artifacts)
    if pptx_declared and (not pptx or not pptx.exists()):
        blockers.append("current_pptx_artifact_id 已声明但未解析到 exact-current PPTX。")
    elif pptx and pptx.exists():
        pptx_stats = inspect_pptx(pptx)
        evidence.append(
            f"pptx={safe_rel(project, pptx)} slides={pptx_stats['slides']} editable_text_runs={pptx_stats['editable_text_runs']}"
        )
        if not pptx_stats["editable"]:
            blockers.append("PPTX 缺少可编辑文本层。")
    else:
        warnings.append("尚未生成 PPTX；text-first 内部交接仍可进行，PPT/客户包/发送 Gate 均未就绪。")

    launcher = REPO_ROOT / "启动广告创意项目.command"
    if source_root() is None:
        evidence.append("launcher=not_applicable_installed_package; use adco CLI entry point")
    elif launcher.exists() and os.access(launcher, os.X_OK):
        evidence.append(f"launcher={launcher} executable=true")
    else:
        blockers.append("双击启动脚本不存在或不可执行。")

    skill = check_global_skill()
    evidence.append(f"skill_install_match={skill['match']}")
    evidence.append(f"skill_target={skill['target']}")
    if not skill["match"]:
        warnings.append("全局 Skill 未安装或与 canonical 草稿不一致；这不影响当前项目证据质量。")

    status = "PASS" if not blockers else "BLOCKED"
    report_path = project / "AD-creative/gates/GATE-AUTO-HANDOFF-READINESS-001_report.md"
    evidence_text = "\n".join(f"- {item}" for item in evidence)
    blocker_text = "\n".join(f"- {item}" for item in blockers) or "- 无"
    warning_text = "\n".join(f"- {item}" for item in warnings) or "- 无"
    write_text(
        report_path,
        f"""# Non-Developer Handoff Readiness Gate

status: {status}
visibility: internal_only
checked_at: {now_iso()}

## Evidence

{evidence_text}

## Blockers

{blocker_text}

## Warnings

{warning_text}

## Scope

此 Gate 证明项目可交给非开发广告创意者继续内部操作。
它不要求已有 PPT、client-pack PASS、manual review receipt 或 send authorization。
它不代表客户包、最终交付或发送准备，后者只能由各自独立 Gate 证明。
""",
    )
    update_artifact(
        project,
        "ART-AUTO-HANDOFF-READINESS-GATE",
        "handoff_readiness_gate_report",
        safe_rel(project, report_path),
        "operations_handoff",
        status="done" if status != "BLOCKED" else "blocked",
        visibility="internal_only",
        gate_status=status,
    )
    append_gate(
        project,
        "GATE-AUTO-HANDOFF-READINESS-001",
        "operations_handoff",
        status,
        "95" if status == "PASS" else "65" if status == "PARTIAL_PASS" else "40",
        "ART-AUTO-HANDOFF-READINESS-GATE;ART-AUTO-MANUAL-REVIEW-CHECKLIST",
        ";".join(blockers[:8]),
        ";".join(warnings[:8]) or "保持 gate 定期重跑。",
        "",
        "internal_operator_handoff_ready" if status == "PASS" else "fix_handoff_blockers",
        "ad_creative_operator",
    )
    append_event(
        project,
        {
            "event_id": "EVT-AUTO-HANDOFF-READINESS-GATE",
            "event_type": "handoff_readiness_gate_run",
            "created_at": now_iso(),
            "status": status,
            "blockers": blockers[:12],
            "warnings": warnings[:12],
        },
    )
    for key, value in stats.items():
        evidence.append(f"{key}={value}")
    return status, blockers, warnings, report_path


def next_ppt_export_identity(project: Path) -> tuple[str, str, str, str, int]:
    _, artifacts = read_csv_rows(project / "AD-creative/orchestrator/artifact_index.csv")
    _, versions = read_csv_rows(project / "AD-creative/orchestrator/version_map.csv")
    highest = 0
    for row in artifacts:
        if row.get("artifact_type", "").strip().lower() != "pptx":
            continue
        match = re.fullmatch(r"v(\d+)", row.get("version", "").strip(), flags=re.IGNORECASE)
        if match:
            highest = max(highest, int(match.group(1)))
    export_root = project / "AD-creative/ppt/exports"
    for path in export_root.glob("client_review_v*.pptx") if export_root.exists() else []:
        match = re.search(r"_v(\d+)\.pptx$", path.name, flags=re.IGNORECASE)
        if match:
            highest = max(highest, int(match.group(1)))
    number = highest + 1
    while True:
        label = f"v{number:03d}"
        version_id = f"VER-PPT-{number:03d}"
        pptx_artifact_id = f"ART-PPTX-{number:03d}"
        check_artifact_id = f"ART-PPT-EDITABILITY-{number:03d}"
        used_ids = {row.get("artifact_id", "") for row in artifacts}
        used_versions = {row.get("version_id", "") for row in versions}
        if pptx_artifact_id not in used_ids and check_artifact_id not in used_ids and version_id not in used_versions:
            return label, version_id, pptx_artifact_id, check_artifact_id, number
        number += 1


def pptx_editability_report_content(
    project: Path,
    pptx_path: Path,
    stats: dict[str, int | bool | str],
    *,
    sha256: str = "",
) -> str:
    status = "PASS" if stats["editable"] else "BLOCKED"
    return f"""# PPT Editability Check

status: {status}
visibility: internal_only
checked_at: {now_iso()}
pptx: {safe_rel(project, pptx_path)}
sha256: {sha256 or file_sha256(pptx_path)}

## Result

| Check | Value |
| --- | --- |
| has_presentation_xml | {stats['has_presentation']} |
| slides | {stats['slides']} |
| editable_text_runs | {stats['editable_text_runs']} |
| editable_slides | {stats.get('editable_slides', '')} |
| flattened_slides | {stats.get('flattened_slides', '') or 'none'} |
| embedded_media | {stats['embedded_media']} |
| editable | {stats['editable']} |

## Rules

- `editable=true` 只代表 PPTX 内存在可编辑文本层。
- 图片页若存在，必须在客户稿前说明用途、来源和可替换性。
- 客户可见前仍需独立视觉、语言、授权与发送准备检查。
"""


def export_editable_pptx(project: Path, output: Path | None = None) -> Path:
    lock_path = project / "AD-creative/orchestrator/.version.lock"
    lock_path.parent.mkdir(parents=True, exist_ok=True)
    with lock_path.open("a+", encoding="utf-8") as lock_handle:
        fcntl.flock(lock_handle.fileno(), fcntl.LOCK_EX)
        try:
            return _export_editable_pptx_locked(project, output)
        finally:
            fcntl.flock(lock_handle.fileno(), fcntl.LOCK_UN)


def _export_editable_pptx_locked(project: Path, output: Path | None = None) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as exc:  # noqa: BLE001 - actionable dependency error
        raise RuntimeError(f"python-pptx unavailable: {exc}") from exc

    outline_status, outline_findings, outline_report = review_client_outline(project)
    if outline_status == "BLOCKED":
        raise RuntimeError(
            "client-outline-gate BLOCKED; fix the customer-readable text framework before PPT builder. "
            f"report={safe_rel(project, outline_report)} findings={'; '.join(outline_findings[:5])}"
        )

    version_label, version_id, pptx_artifact_id, check_artifact_id, version_number = next_ppt_export_identity(project)
    outline_rows = client_outline_rows(project)
    canonical_output = (
        project / f"AD-creative/ppt/exports/client_review_{version_label}.pptx"
    ).resolve()
    if output is not None and output.resolve() != canonical_output:
        raise RuntimeError(
            "PPTX output is immutable and must use the canonical version path: "
            f"{canonical_output}"
        )
    output = canonical_output
    if output.exists():
        raise FileExistsError(f"refusing to overwrite immutable PPTX export: {output}")
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title(slide, title: str, subtitle: str = "") -> None:
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(11.8), Inches(0.7))
        frame = title_box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.size = Pt(26)
        paragraph.font.bold = True
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(1.18), Inches(11.5), Inches(0.35))
            sub_frame = sub_box.text_frame
            sub_frame.text = subtitle
            sub_frame.paragraphs[0].font.size = Pt(11)

    def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float) -> None:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        frame.clear()
        for index, item in enumerate(items):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = item
            paragraph.level = 0
            paragraph.font.size = Pt(15)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    client_project_name = project.name if not find_client_language_hits(project.name) else ""
    cover_title = f"{client_project_name} 广告创意方向提案" if client_project_name else "广告创意方向提案"
    add_title(slide, cover_title, "用于方向讨论、关键判断与下一阶段深化。")
    add_bullets(
        slide,
        [
            "从客户目标和传播问题出发",
            "对比不同创意方向的核心信息与关键画面",
            "锁定下一轮需要深化的内容与素材",
        ],
        0.8,
        2.0,
        10.8,
        2.8,
    )

    for row in outline_rows:
        slide = prs.slides.add_slide(blank)
        add_title(
            slide,
            row.get("page_title", "方向讨论"),
            f"本页希望对齐：{row.get('client_confirmation_point', '')}",
        )
        add_bullets(
            slide,
            [
                row.get("body_copy", ""),
                f"画面建议：{row.get('visual_slot', '')}",
            ],
            0.8,
            1.7,
            11.4,
            4.8,
        )

    temp_output = output.with_name(f".{output.name}.{os.getpid()}.tmp.pptx")
    try:
        prs.save(temp_output)
        stats = inspect_pptx(temp_output)
        if not stats["editable"]:
            raise RuntimeError("generated PPTX lacks editable text layers")
        os.link(temp_output, output)
    finally:
        temp_output.unlink(missing_ok=True)

    pptx_sha = file_sha256(output)
    check_path = project / f"AD-creative/ppt/exports/client_review_{version_label}_editability.md"
    if check_path.exists():
        output.unlink(missing_ok=True)
        raise FileExistsError(f"refusing to overwrite immutable PPTX check: {check_path}")

    artifact_path = project / "AD-creative/orchestrator/artifact_index.csv"
    version_path = project / "AD-creative/orchestrator/version_map.csv"
    truth_path = project / "AD-creative/orchestrator/current_truth.md"
    gate_path = project / "AD-creative/orchestrator/gate_log.csv"
    event_path = project / "AD-creative/orchestrator/events.jsonl"
    transaction_paths = [artifact_path, version_path, truth_path, gate_path, event_path]
    snapshots = {path: path.read_bytes() if path.exists() else None for path in transaction_paths}
    try:
        write_text(
            check_path,
            pptx_editability_report_content(project, output, stats, sha256=pptx_sha),
        )
        artifact_fields, artifact_rows = read_csv_rows(artifact_path)
        version_fields, version_rows = read_csv_rows(version_path)
        truth_text = truth_path.read_text(encoding="utf-8") if truth_path.exists() else ""
        previous_version_id = current_truth_value(truth_text, "current_version_id")
        previous_pptx_artifact_id = current_truth_value(truth_text, "current_pptx_artifact_id")
        previous_pptx_path = ""
        if previous_pptx_artifact_id:
            for row in artifact_rows:
                if row.get("artifact_id") == previous_pptx_artifact_id and row.get("artifact_type", "").lower() == "pptx":
                    previous_pptx_path = row.get("path", "")
                    row["status"] = "superseded"
                    row["updated_at"] = now_iso()
                    break
        if previous_version_id:
            for row in version_rows:
                if row.get("version_id") == previous_version_id:
                    row["status"] = "superseded"
                    break
        created_at = now_iso()
        artifact_rows.extend(
            [
                {
                    "artifact_id": pptx_artifact_id,
                    "artifact_type": "pptx",
                    "path": safe_rel(project, output),
                    "stage": "ppt_gate",
                    "version": version_label,
                    "status": "done",
                    "visibility": "internal_only",
                    "source_event_ids": "",
                    "linked_requirements": "",
                    "linked_work_items": "",
                    "linked_references": "",
                    "linked_assets": "",
                    "gate_status": "PASS",
                    "supersedes_artifact_id": previous_pptx_artifact_id,
                    "created_at": created_at,
                    "updated_at": created_at,
                    "sha256": pptx_sha,
                    "size_bytes": str(output.stat().st_size),
                },
                {
                    "artifact_id": check_artifact_id,
                    "artifact_type": "ppt_editability_check",
                    "path": safe_rel(project, check_path),
                    "stage": "ppt_gate",
                    "version": version_label,
                    "status": "done",
                    "visibility": "internal_only",
                    "source_event_ids": "",
                    "linked_requirements": "",
                    "linked_work_items": "",
                    "linked_references": "",
                    "linked_assets": "",
                    "gate_status": "PASS",
                    "supersedes_artifact_id": "",
                    "created_at": created_at,
                    "updated_at": created_at,
                    "sha256": file_sha256(check_path),
                    "size_bytes": str(check_path.stat().st_size),
                    "derived_from_artifact_id": pptx_artifact_id,
                    "derived_from_sha256": pptx_sha,
                },
            ]
        )
        version_rows.append(
            {
                "version_id": version_id,
                "artifact_id": pptx_artifact_id,
                "version": version_label,
                "status": "draft",
                "created_at": created_at,
                "source_event_ids": "",
                "supersedes_version_id": previous_version_id,
                "notes": f"immutable_pptx_export;sha256={pptx_sha}",
            }
        )
        write_csv_rows(artifact_path, artifact_fields, artifact_rows)
        write_csv_rows(version_path, version_fields, version_rows)
        update_markdown_sections(
            truth_path,
            {
                "Current Version Truth": f"""```text
current_version_id: {version_id}
current_pptx_artifact_id: {pptx_artifact_id}
current_pdf_artifact_id:
current_preview_artifact_id:
current_text_extract_artifact_id:
current_ppt_editability_artifact_id: {check_artifact_id}
version_map_status: draft
last_archive_before_edit: {('preserved:' + previous_pptx_path) if previous_pptx_path else 'not_required_first_immutable_export'}
```""",
                "Current Stage": "ppt_internal_review",
                "Next Action": "完成视觉、客户语言、素材授权、PDF/preview/text extract 与人工复核后，方可进入 client-send readiness。",
            },
        )
        append_gate(
            project,
            "GATE-AUTO-PPT-001",
            "ppt_gate",
            "PASS",
            "90",
            f"{pptx_artifact_id};{check_artifact_id}",
            "",
            "可编辑性只代表结构检查；客户可见前仍需独立视觉、语言、授权和发送准备检查。",
            "",
            "ready_for_internal_review",
            "ad_creative_operator",
        )
        append_event(
            project,
            {
                "event_id": f"EVT-PPTX-EXPORT-{version_number:03d}",
                "event_type": "immutable_pptx_exported",
                "created_at": created_at,
                "version_id": version_id,
                "artifact_id": pptx_artifact_id,
                "pptx": safe_rel(project, output),
                "sha256": pptx_sha,
                "stats": stats,
            },
        )
        validation_errors, _ = validate(project)
        if validation_errors:
            raise RuntimeError(
                "PPTX export transaction would leave project invalid: "
                + "; ".join(validation_errors[:12])
            )
    except Exception:
        for path, snapshot in snapshots.items():
            if snapshot is None:
                path.unlink(missing_ok=True)
            else:
                path.parent.mkdir(parents=True, exist_ok=True)
                path.write_bytes(snapshot)
        check_path.unlink(missing_ok=True)
        output.unlink(missing_ok=True)
        raise
    return output


def command_goal_plan(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    goal_id = args.goal_id or default_goal_id()
    title = args.title or goal_id
    objective = args.objective or "按双泳道 goal 模式推进下一轮低风险内部工作。"
    plan = render_goal_iteration_plan(
        project,
        goal_id=goal_id,
        title=title,
        objective=objective,
        owner=args.owner,
        force=args.force,
    )
    dashboard = render_dashboard(project)
    errors, stats = validate(project)
    print(f"GOAL_ID={goal_id}")
    print(f"GOAL_PLAN={plan}")
    print(f"DASHBOARD={dashboard}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def command_thread_plan(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    goal_id = args.goal_id or default_goal_id()
    try:
        roles = parse_threadops_roles(args.roles)
        payload = render_thread_execution_plan(
            project,
            goal_id=goal_id,
            title=args.title or goal_id,
            objective=args.objective or "Use Codex Threads as controlled execution lanes for this delivery goal.",
            roles=roles,
            work_id=args.work_id,
            task_signature_id=args.task_signature_id,
            brand=args.brand,
            product=args.product,
            talent_or_ip=args.talent_or_ip,
            platform_or_channel=args.platform_or_channel,
            deliverable=args.deliverable,
            stage=args.stage,
            primary_risks=args.primary_risks,
            evidence_needed=args.evidence_needed,
            master_thread_id=args.master_thread_id,
            current_version=args.current_version_id,
            max_active=args.max_active,
            extra_read_first=args.read_first,
            force=args.force,
        )
    except Exception as exc:  # noqa: BLE001 - CLI should surface actionable plan failure
        print("THREAD_PLAN=CHECK")
        print(f"ERROR={exc}")
        return 1
    dashboard = render_dashboard(project)
    errors, stats = validate(project)
    if args.json:
        output = {
            "thread_plan": "PASS" if not errors else "CHECK",
            "project": str(project),
            "goal_id": payload["goal_id"],
            "work_id": payload["work_id"],
            "task_signature_id": payload["task_signature_id"],
            "thread_lane_plan": str(payload["thread_lane_plan"]),
            "prompt_dir": str(payload["prompt_dir"]),
            "cleanup_plan": str(payload["cleanup_plan"]),
            "roles": payload["roles"],
            "dashboard": str(dashboard),
            "validation": "PASS" if not errors else "CHECK",
            "stats": stats,
            "errors": errors,
        }
        print(json.dumps(output, ensure_ascii=False, indent=2, sort_keys=True))
        return 0 if not errors else 1
    print(f"THREAD_PLAN={'PASS' if not errors else 'CHECK'}")
    print(f"PROJECT={project}")
    print(f"GOAL_ID={payload['goal_id']}")
    print(f"WORK_ID={payload['work_id']}")
    print(f"TASK_SIGNATURE_ID={payload['task_signature_id']}")
    print(f"THREAD_LANE_PLAN={payload['thread_lane_plan']}")
    print(f"PROMPT_DIR={payload['prompt_dir']}")
    print(f"CLEANUP_PLAN={payload['cleanup_plan']}")
    print("ROLES=" + ";".join(payload["roles"]))
    print(f"DASHBOARD={dashboard}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def command_goal_run(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    result = run_goal(
        project,
        goal_id=args.goal_id,
        max_steps=args.max_steps,
        allow_generate=args.allow_generate,
    )
    if args.json:
        print(json.dumps(result, ensure_ascii=False, indent=2, sort_keys=True))
    else:
        print(f"PROJECT={result['project']}")
        print(f"GOAL_ID={result['goal_id']}")
        print(f"STOP_REASON={result['stop_reason']}")
        print(f"DASHBOARD={result['dashboard']}")
        for step in result["steps"]:
            print(f"STEP={step['step']} ACTION={step['action']} STATUS={step['status']} DETAIL={step['detail']}")
    return 1 if result["stop_reason"] in {"VALIDATION_CHECK", "BLOCKED_GATE"} else 0


def command_creative_doctor(args: argparse.Namespace) -> int:
    status, issues, warnings, evidence = creative_doctor_report()
    if args.json:
        print(
            json.dumps(
                {
                    "status": status,
                    "issues": issues,
                    "warnings": warnings,
                    "evidence": evidence,
                },
                ensure_ascii=False,
                indent=2,
                sort_keys=True,
            )
        )
    else:
        print(f"CREATIVE_PRODUCTION={status}")
        for item in evidence:
            print(f"EVIDENCE={item}")
        if warnings:
            print("WARNINGS:")
            for warning in warnings:
                print(f"- {warning}")
        if issues:
            print("ISSUES:")
            for issue in issues:
                print(f"- {issue}")
    return 0 if status == "PASS" else 1


def command_creative_run(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    if args.generate and args.review_only:
        print("CREATIVE_RUN=CHECK")
        print("ERROR=Choose either --review-only or --generate, not both.")
        return 1
    brief_file = Path(args.brief_file).expanduser().resolve()
    base_asset = Path(args.base_asset).expanduser().resolve() if args.base_asset else None
    if not brief_file.is_file():
        print("CREATIVE_RUN=BLOCKED")
        print(f"ERROR=brief file not found: {brief_file}")
        return 1
    if base_asset is not None and not base_asset.is_file():
        print("CREATIVE_RUN=BLOCKED")
        print(f"ERROR=base asset not found: {base_asset}")
        return 1
    ensure_project(project)
    try:
        run_dir, logs = run_creative_production(
            project,
            kind=args.kind,
            work_id=args.work_id,
            brief_file=brief_file,
            base_asset=base_asset,
            generate=args.generate,
            force=args.force,
        )
    except Exception as exc:  # noqa: BLE001 - CLI should surface bridge failure
        print("CREATIVE_RUN=BLOCKED")
        print(f"ERROR={exc}")
        return 1
    errors, stats = validate(project)
    print("CREATIVE_RUN=PASS")
    print(f"KIND={args.kind}")
    print(f"RUN_DIR={run_dir}")
    for log in logs:
        print(f"LOG={log}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def command_import_creative_production(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    run_dir = Path(args.run_dir).expanduser().resolve()
    if not run_dir.is_dir():
        print("CREATIVE_PRODUCTION_IMPORT=BLOCKED")
        print(f"ERROR=run directory not found: {run_dir}")
        return 1
    ensure_project(project)
    try:
        asset_ids, metadata_dir = import_creative_production_run(
            project,
            run_dir=run_dir,
            kind=args.kind,
            slot_prefix=args.slot_prefix,
            requirement_id=args.requirement_id,
            reference_id=args.reference_id,
        )
    except Exception as exc:  # noqa: BLE001 - CLI should surface import failure
        print("CREATIVE_PRODUCTION_IMPORT=BLOCKED")
        print(f"ERROR={exc}")
        return 1
    errors, stats = validate(project)
    print("CREATIVE_PRODUCTION_IMPORT=PASS")
    print(f"KIND={args.kind}")
    print(f"METADATA_DIR={metadata_dir}")
    print(f"IMPORTED_ASSETS={len(asset_ids)}")
    if asset_ids:
        print("ASSET_IDS=" + ";".join(asset_ids))
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def _command_creative_brief(args: argparse.Namespace, *, deprecated: bool) -> int:
    project = Path(args.project).resolve()
    payload = (
        render_creative_proposal(project, work_id=args.work_id)
        if deprecated
        else render_creative_brief(project, work_id=args.work_id)
    )
    incremental = run_incremental_validation(
        project,
        changed_artifact_ids=payload["artifact_ids"],
        changed_file_paths=payload["paths"],
    )
    errors = list(incremental.errors)
    payload.update(
        {
            "creative_brief": "PASS" if not errors else "CHECK",
            "deprecated_alias": "creative-brief" if deprecated else None,
            "dashboard": "",
            "incremental_validation": incremental.as_dict(),
            "full_validation": "NOT_RUN",
            "errors": errors,
        }
    )
    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True))
        return 0 if not errors else 1
    if deprecated:
        print("DEPRECATED=creative-proposal is an alias for creative-brief")
    print(f"CREATIVE_BRIEF={'PASS' if not errors else 'CHECK'}")
    print(f"PROJECT={project}")
    if args.work_id:
        print(f"WORK_ID={args.work_id}")
    print("ARTIFACT_IDS=" + ";".join(payload["artifact_ids"]))
    print(f"BRIEF_SNAPSHOT_SHA256={payload['brief_snapshot_sha256']}")
    print(f"DIRECTIONS_GENERATED={payload['directions_generated']}")
    for path in payload["paths"]:
        print(f"ARTIFACT_PATH={path}")
    print("DASHBOARD=NOT_RUN")
    print("VALIDATORS_RUN=" + ";".join(incremental.validators_run))
    print(f"INCREMENTAL_VALIDATION={'PASS' if not errors else 'CHECK'}")
    print("FULL_VALIDATION=NOT_RUN")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def command_creative_brief(args: argparse.Namespace) -> int:
    return _command_creative_brief(args, deprecated=False)


def command_creative_proposal(args: argparse.Namespace) -> int:
    return _command_creative_brief(args, deprecated=True)


def command_creative_import(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    try:
        result = import_creative_candidate(
            project,
            Path(args.file).expanduser().resolve(),
        )
    except ValueError as exc:
        print("CREATIVE_IMPORT=BLOCKED")
        print(f"ERROR={exc}")
        return 1
    if project_surface(project) == DELIVERY_SURFACE:
        for artifact_id, artifact_type, path in [
            ("ART-AUTO-CREATIVE-CANDIDATE", "creative_candidate", result.current_path),
            (
                "ART-AUTO-CREATIVE-CANDIDATE-RECEIPT",
                "creative_candidate_import_receipt",
                result.receipt_path,
            ),
            ("ART-AUTO-CREATIVE-DIRECTIONS", "creative_directions", result.directions_path),
            ("ART-AUTO-CREATIVE-OPTION-MATRIX", "creative_option_matrix", result.matrix_path),
        ]:
            update_artifact(
                project,
                artifact_id,
                artifact_type,
                safe_rel(project, path),
                "creative",
                visibility="internal_only",
                gate_status="NOT_RUN",
            )
        append_event(
            project,
            {
                "event_id": f"EVT-CREATIVE-IMPORT-{datetime.now().strftime('%Y%m%d%H%M%S')}",
                "event_type": "creative_candidate_imported",
                "created_at": now_iso(),
                "candidate_sha256": result.candidate_sha256,
                "direction_count": result.direction_count,
                "warnings": result.warnings,
            },
        )
    incremental = run_incremental_validation(
        project,
        changed_artifact_ids=[
            "ART-AUTO-CREATIVE-CANDIDATE",
            "ART-AUTO-CREATIVE-DIRECTIONS",
        ],
        changed_file_paths=[
            result.current_path,
            result.directions_path,
            result.matrix_path,
        ],
    )
    payload = {
        "creative_import": "PASS" if not incremental.errors else "CHECK",
        "candidate": str(result.candidate_path),
        "current_candidate": str(result.current_path),
        "candidate_sha256": result.candidate_sha256,
        "directions": result.direction_count,
        "warnings": result.warnings,
        "creative_quality": "NOT_EVALUATED",
        "incremental_validation": incremental.as_dict(),
        "full_validation": "NOT_RUN",
    }
    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True))
    else:
        print(f"CREATIVE_IMPORT={'PASS' if not incremental.errors else 'CHECK'}")
        print(f"CANDIDATE={result.candidate_path}")
        print(f"CANDIDATE_SHA256={result.candidate_sha256}")
        print(f"DIRECTIONS={result.direction_count}")
        print(f"WARNINGS={len(result.warnings)}")
        print("CREATIVE_QUALITY=NOT_EVALUATED")
        print("VALIDATORS_RUN=" + ";".join(incremental.validators_run))
        print("FULL_VALIDATION=NOT_RUN")
        if incremental.errors:
            print("ERRORS:")
            for error in incremental.errors:
                print(f"- {error}")
    return 0 if not incremental.errors else 1


def command_creative_review(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    delivery_surface = project_surface(project) == DELIVERY_SURFACE
    try:
        result = review_creative_candidate(
            project,
            independent_critic_required=delivery_surface,
        )
    except ValueError as exc:
        print("CREATIVE_REVIEW=BLOCKED")
        print(f"ERROR={exc}")
        return 1
    if delivery_surface:
        assert result.receipt_path is not None
        update_artifact(
            project,
            "ART-AUTO-CREATIVE-CRITIC-RECEIPT",
            "creative_critic_receipt",
            safe_rel(project, result.receipt_path),
            "creative",
            visibility="internal_only",
            gate_status=result.status,
        )
        append_event(
            project,
            {
                "event_id": f"EVT-CREATIVE-REVIEW-{datetime.now().strftime('%Y%m%d%H%M%S')}",
                "event_type": "creative_candidate_reviewed",
                "created_at": now_iso(),
                "status": result.status,
                "receipt": safe_rel(project, result.receipt_path),
                "verdict": result.receipt["verdict"],
            },
        )
    payload = {
        "creative_review": result.status,
        "critic_receipt": str(result.receipt_path) if result.receipt_path else "",
        "verdict": result.receipt["verdict"],
        "blocking_issues": result.blocking_issues,
        "warnings": result.warnings,
        "independent_critic_required": bool(
            result.receipt["independent_critic_required"]
        ),
    }
    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True))
    else:
        print(f"CREATIVE_REVIEW={result.status}")
        print(
            "CRITIC_RECEIPT="
            + (str(result.receipt_path) if result.receipt_path else "NOT_RECORDED")
        )
        print(f"VERDICT={result.receipt['verdict']}")
        print(f"BLOCKING_ISSUES={len(result.blocking_issues)}")
        print(f"WARNINGS={len(result.warnings)}")
        print(
            "INDEPENDENT_CRITIC_REQUIRED="
            + str(int(bool(result.receipt["independent_critic_required"])))
        )
    return 1 if result.status == "BLOCKED" else 0


def command_init(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().absolute()
    template = Path(args.template).expanduser().resolve() if args.template else TEMPLATE_ROOT
    if not template.exists():
        print("INIT=CHECK")
        print(f"ERROR=template not found: {template}")
        return 1
    try:
        use_delivery = args.full or project_surface(project) == DELIVERY_SURFACE
        copy = copy_template if use_delivery else copy_content_template
        created, skipped = copy(template, project)
    except (OSError, RuntimeError, ValueError) as exc:
        print("INIT=CHECK")
        print(f"ERROR={exc}")
        return 1
    agents_status = agents_policy_status(project)
    errors, stats = validate(project)
    print(f"PROJECT={project}")
    print(f"TEMPLATE={template}")
    print(f"CREATED_FILES={created}")
    print(f"SKIPPED_EXISTING_FILES={skipped}")
    print(f"PROJECT_SURFACE={project_surface(project)}")
    print(f"AGENTS_MD={agents_status}")
    print(f"INIT={'PASS' if not errors else 'CHECK'}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def command_run(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().absolute()
    materials = [Path(item).expanduser().absolute() for item in args.material]
    delivery_surface = project_surface(project) == DELIVERY_SURFACE
    project_existed = project.exists()
    try:
        result = execute_lightweight_run(
            project,
            materials=materials,
            goal=args.goal,
            max_total_chars=args.max_total_chars,
            ensure_project=ensure_project,
            register_materials=register_materials,
            ensure_intake_work=ensure_intake_work if delivery_surface else None,
            perform_intake=perform_intake,
            render_handoff=render_handoff,
            render_dashboard=render_dashboard,
            render_optional_dashboard=args.dashboard,
        )
    except (OSError, RuntimeError, ValueError) as exc:
        code = exc.code if isinstance(exc, RunPreflightError) else "runtime_error"
        error = {"code": code, "message": str(exc)}
        payload = {
            "run": "CHECK",
            "project": str(project),
            "project_created": not project_existed and project.exists(),
            "error": error,
            "errors": [str(exc)],
        }
        if args.json:
            print(json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True))
        else:
            print("RUN=CHECK")
            print(f"ERROR_CODE={code}")
            print(f"ERROR={exc}")
            print(f"PROJECT_CREATED={int(payload['project_created'])}")
        return 1
    result["agents_policy"] = agents_policy_status(project)
    incremental = result["incremental_validation"]
    intake_stats = result["intake"]
    errors = list(incremental["errors"])
    if intake_stats.get("over_budget_files", 0):
        errors.append("intake total character budget exceeded; inspect intake-evidence report")
    if intake_stats.get("parser_errors", 0):
        errors.append("one or more material parsers failed; inspect intake-evidence report")
    result["run"] = "PASS" if not errors else "CHECK"
    result["errors"] = errors
    if args.json:
        print(json.dumps(result, ensure_ascii=False, indent=2, sort_keys=True))
        return 0 if not errors else 1
    print("CONTENT_ANSWER:")
    print(result["content_answer"]["markdown"])
    print(f"PROJECT={project}")
    print(f"PROJECT_SURFACE={project_surface(project)}")
    print(f"CREATED_FILES={result['created_files']}")
    print(f"SKIPPED_EXISTING_FILES={result['skipped_existing_files']}")
    print(f"AGENTS_MD={result['agents_policy']}")
    print(f"REGISTERED_SOURCES={result['registered_sources']}")
    print(f"INTAKE_MATERIALS={intake_stats['materials']}")
    print(f"INTAKE_REQUIREMENTS={intake_stats['requirements']}")
    print(f"INTAKE_GAPS={intake_stats['gaps']}")
    print(f"INTAKE_CHARACTERS_READ={intake_stats['characters_read']}")
    print(f"INTAKE_EVIDENCE_CHUNKS={intake_stats['evidence_chunks']}")
    print(f"DASHBOARD={result['dashboard'] or 'NOT_RUN'}")
    print(f"DASHBOARD_RENDER_COUNT={result['dashboard_render_count']}")
    print("COUNCIL=NOT_RUN")
    print(f"COUNCIL_RUN_COUNT={result['council_run_count']}")
    print(f"SPECIALIST_HANDOFF_COUNT={result['specialist_handoff_count']}")
    print(f"CLIENT_PACK_RUN_COUNT={result['client_pack_run_count']}")
    print(f"FULL_VALIDATION_RUN_COUNT={result['full_validation_run_count']}")
    print(f"NEXT_COMMAND={result['next_command']}")
    print(f"PPT_AUTO_GENERATED={result['ppt_auto_generated']}")
    print("VALIDATORS_RUN=" + ";".join(incremental["validators_run"]))
    print("VALIDATORS_SKIPPED=" + ";".join(incremental["validators_skipped"]))
    for key, value in result["timings"].items():
        print(f"{key.upper()}={value}")
    print(f"INCREMENTAL_VALIDATION={'PASS' if not errors else 'CHECK'}")
    print("FULL_VALIDATION=NOT_RUN")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def _run_sample_content(project: Path, *, force_material: bool) -> dict[str, object]:
    created, skipped = ensure_project(project)
    agents_status = agents_policy_status(project)
    material, material_action = write_sample_brief(project, force=force_material)
    source_ids = existing_source_ids_for_material(project, material)
    registered_sources = 0
    if not source_ids:
        source_ids = register_materials(project, [material], SAMPLE_GOAL)
        registered_sources = len(source_ids)
    if project_surface(project) == DELIVERY_SURFACE:
        ensure_intake_work(project, source_ids, SAMPLE_GOAL)
    intake_stats = perform_intake(project, source_ids, SAMPLE_GOAL)
    content_answer = render_handoff(project, SAMPLE_GOAL, source_ids)
    errors, stats = validate(project)
    return {
        "project": str(project),
        "surface": project_surface(project),
        "created_files": created,
        "skipped_existing_files": skipped,
        "agents_policy": agents_status,
        "sample_material": str(material),
        "sample_material_action": material_action,
        "registered_sources": registered_sources,
        "source_ids": source_ids,
        "intake": intake_stats,
        "content_answer": content_answer,
        "stats": stats,
        "errors": errors,
    }


def command_sample(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    result = _run_sample_content(project, force_material=args.force_material)
    errors = result["errors"]
    print(f"SAMPLE={'PASS' if not errors else 'CHECK'}")
    print("CONTENT_ANSWER:")
    print(result["content_answer"]["markdown"])
    print(f"PROJECT={result['project']}")
    print(f"PROJECT_SURFACE={result['surface']}")
    print(f"CREATED_FILES={result['created_files']}")
    print(f"SKIPPED_EXISTING_FILES={result['skipped_existing_files']}")
    print(f"AGENTS_MD={result['agents_policy']}")
    print(f"SAMPLE_MATERIAL={result['sample_material']}")
    print(f"SAMPLE_MATERIAL_ACTION={result['sample_material_action']}")
    print(f"REGISTERED_SOURCES={result['registered_sources']}")
    print(f"SOURCE_IDS={';'.join(result['source_ids'])}")
    print(f"INTAKE_MATERIALS={result['intake']['materials']}")
    print(f"INTAKE_REQUIREMENTS={result['intake']['requirements']}")
    print(f"INTAKE_GAPS={result['intake']['gaps']}")
    print("DASHBOARD=NOT_RUN")
    print("COUNCIL=NOT_RUN")
    for key, value in result["stats"].items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def command_demo(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve() if args.project else DEFAULT_DEMO_PROJECT
    result = _run_sample_content(project, force_material=args.force_material)
    dashboard = render_dashboard(project)
    open_status = "SKIPPED"
    if not args.no_open:
        open_status = "PASS" if webbrowser.open(dashboard.as_uri()) else "CHECK"
    errors = result["errors"]
    print(f"DEMO={'PASS' if not errors and open_status != 'CHECK' else 'CHECK'}")
    print("CONTENT_ANSWER:")
    print(result["content_answer"]["markdown"])
    print(f"PROJECT={result['project']}")
    print(f"PROJECT_SURFACE={result['surface']}")
    print(f"CREATED_FILES={result['created_files']}")
    print(f"SKIPPED_EXISTING_FILES={result['skipped_existing_files']}")
    print(f"AGENTS_MD={result['agents_policy']}")
    print(f"SAMPLE_MATERIAL={result['sample_material']}")
    print(f"SAMPLE_MATERIAL_ACTION={result['sample_material_action']}")
    print(f"REGISTERED_SOURCES={result['registered_sources']}")
    print(f"SOURCE_IDS={';'.join(result['source_ids'])}")
    print(f"INTAKE_MATERIALS={result['intake']['materials']}")
    print(f"INTAKE_REQUIREMENTS={result['intake']['requirements']}")
    print(f"INTAKE_GAPS={result['intake']['gaps']}")
    print(f"DASHBOARD={dashboard}")
    print(f"DASHBOARD_OPEN={open_status}")
    print("COUNCIL=NOT_RUN")
    for key, value in result["stats"].items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors or open_status == "CHECK":
        if errors:
            print("ERRORS:")
            for error in errors:
                print(f"- {error}")
        return 1
    return 0


def command_quickstart(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve() if args.project else DEFAULT_DEMO_PROJECT
    result = _run_sample_content(project, force_material=args.force_material)
    dashboard = render_dashboard(project)
    errors = result["errors"]
    open_status = "SKIPPED"
    if not args.no_open:
        open_status = "PASS" if webbrowser.open(dashboard.as_uri()) else "CHECK"
    quickstart_status = "PASS" if not errors and open_status != "CHECK" else "CHECK"
    payload = {
        "quickstart": quickstart_status,
        **result,
        "dashboard": str(dashboard),
        "dashboard_open": open_status,
        "council": "NOT_RUN",
        "next_command": f"adco next {project}",
        "status_command": f"adco status {project}",
        "validate_command": f"adco validate {project}",
        "real_project_command": "adco run <project_dir> --material <material_file_or_folder>",
        "validation": "PASS" if not errors else "CHECK",
    }
    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True))
        return 0 if quickstart_status == "PASS" else 1
    print(f"QUICKSTART={quickstart_status}")
    print("CONTENT_ANSWER:")
    print(result["content_answer"]["markdown"])
    print(f"PROJECT={result['project']}")
    print(f"PROJECT_SURFACE={result['surface']}")
    print(f"CREATED_FILES={result['created_files']}")
    print(f"SKIPPED_EXISTING_FILES={result['skipped_existing_files']}")
    print(f"AGENTS_MD={result['agents_policy']}")
    print(f"SAMPLE_MATERIAL={result['sample_material']}")
    print(f"SAMPLE_MATERIAL_ACTION={result['sample_material_action']}")
    print(f"REGISTERED_SOURCES={result['registered_sources']}")
    print(f"SOURCE_IDS={';'.join(result['source_ids'])}")
    print(f"INTAKE_MATERIALS={result['intake']['materials']}")
    print(f"INTAKE_REQUIREMENTS={result['intake']['requirements']}")
    print(f"INTAKE_GAPS={result['intake']['gaps']}")
    print(f"DASHBOARD={dashboard}")
    print(f"DASHBOARD_OPEN={open_status}")
    print("COUNCIL=NOT_RUN")
    print(f"NEXT_COMMAND=adco next {project}")
    print(f"STATUS_COMMAND=adco status {project}")
    print(f"VALIDATE_COMMAND=adco validate {project}")
    print("REAL_PROJECT_COMMAND=adco run <project_dir> --material <material_file_or_folder>")
    for key, value in result["stats"].items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors or open_status == "CHECK":
        if errors:
            print("ERRORS:")
            for error in errors:
                print(f"- {error}")
        return 1
    return 0


def command_render_dashboard(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    dashboard = render_dashboard(project)
    print(f"DASHBOARD={dashboard}")
    return 0


def command_open_dashboard(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    dashboard = render_dashboard(project)
    errors, stats = validate(project)
    open_status = "SKIPPED"
    if not args.no_open:
        open_status = "PASS" if webbrowser.open(dashboard.as_uri()) else "CHECK"
    print(f"DASHBOARD={dashboard}")
    print(f"DASHBOARD_OPEN={open_status}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors or open_status == "CHECK":
        if errors:
            print("ERRORS:")
            for error in errors:
                print(f"- {error}")
        return 1
    return 0


def command_intake(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    source_ids = args.source_id or []
    stats = perform_intake(
        project,
        source_ids,
        args.goal,
        max_total_chars=args.max_total_chars,
    )
    content_answer = render_handoff(project, args.goal, source_ids)
    incremental = run_incremental_validation(
        project,
        changed_artifact_ids=[
            "ART-AUTO-EVIDENCE-CHUNKS",
            "ART-AUTO-FACT-INVENTORY",
            "ART-AUTO-REQUIREMENTS",
            "ART-AUTO-GAPS",
        ],
        changed_file_paths=[
            "AD-creative/orchestrator/evidence_chunks.jsonl",
            "AD-creative/orchestrator/fact_inventory.jsonl",
            "AD-creative/orchestrator/requirements.csv",
            "AD-creative/orchestrator/gaps.csv",
        ],
    )
    errors = list(incremental.errors)
    if stats["over_budget_files"]:
        errors.append("intake total character budget exceeded")
    if stats["parser_errors"]:
        errors.append("one or more material parsers failed")
    print("CONTENT_ANSWER:")
    print(content_answer["markdown"])
    print(f"PROJECT={project}")
    print(f"PROJECT_SURFACE={project_surface(project)}")
    print(f"INTAKE_MATERIALS={stats['materials']}")
    print(f"INTAKE_REQUIREMENTS={stats['requirements']}")
    print(f"INTAKE_GAPS={stats['gaps']}")
    print(f"INTAKE_CHARACTERS_READ={stats['characters_read']}")
    print(f"INTAKE_EVIDENCE_CHUNKS={stats['evidence_chunks']}")
    print(f"INTAKE_OVER_BUDGET_FILES={stats['over_budget_files']}")
    print(f"INTAKE_PARSER_ERRORS={stats['parser_errors']}")
    print("DASHBOARD=NOT_RUN")
    print("VALIDATORS_RUN=" + ";".join(incremental.validators_run))
    print("VALIDATORS_SKIPPED=" + ";".join(incremental.validators_skipped))
    print(f"VALIDATION_MS={incremental.validation_ms}")
    print(f"INCREMENTAL_VALIDATION={'PASS' if not errors else 'CHECK'}")
    print("FULL_VALIDATION=NOT_RUN")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def selected_source_rows(project: Path, source_ids: list[str]) -> list[dict[str, str]]:
    _, source_rows = read_csv_rows(
        project / "AD-creative/orchestrator/source_events.csv"
    )
    selected = set(source_ids)
    return [
        row
        for row in source_rows
        if not selected or row.get("source_event_id", "") in selected
    ]


def command_intake_evidence(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    report = ingest_source_rows(
        project,
        selected_source_rows(project, args.source_id or []),
        max_total_chars=args.max_total_chars,
    )
    payload = report.as_dict()
    payload["intake_evidence"] = (
        "PASS" if not report.over_budget and not report.parser_errors else "CHECK"
    )
    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True))
    else:
        print(f"INTAKE_EVIDENCE={payload['intake_evidence']}")
        print(f"EVIDENCE_PATH={project / payload['evidence_path']}")
        print(f"FILES_PROCESSED={payload['files_processed']}")
        print(f"CHARACTERS_READ={payload['characters_read']}")
        print(f"EVIDENCE_CHUNKS={payload['evidence_chunks']}")
        print(f"OVER_BUDGET_FILES={len(report.over_budget)}")
        print(f"PARSER_ERRORS={len(report.parser_errors)}")
        for item in report.over_budget:
            print(
                "OVER_BUDGET="
                + json.dumps(item, ensure_ascii=False, sort_keys=True)
            )
        for item in report.parser_errors:
            print(
                "PARSER_ERROR="
                + json.dumps(item, ensure_ascii=False, sort_keys=True)
            )
    return 0 if payload["intake_evidence"] == "PASS" else 1


def command_export_intake_analysis_request(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    try:
        payload, path = export_intake_analysis_request(project)
    except ValueError as exc:
        print("INTAKE_ANALYSIS_REQUEST=BLOCKED")
        print(f"ERROR={exc}")
        return 1
    if args.json:
        print(
            json.dumps(
                {
                    "intake_analysis_request": "PASS",
                    "path": str(path),
                    "evidence_chunks": len(payload["evidence_chunks"]),
                },
                ensure_ascii=False,
                indent=2,
                sort_keys=True,
            )
        )
    else:
        print("INTAKE_ANALYSIS_REQUEST=PASS")
        print(f"PATH={path}")
        print(f"EVIDENCE_CHUNKS={len(payload['evidence_chunks'])}")
    return 0


def command_import_intake_analysis(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    try:
        facts, gaps, path = import_intake_analysis(
            project,
            Path(args.file).expanduser().resolve(),
        )
    except ValueError as exc:
        print("INTAKE_ANALYSIS_IMPORT=BLOCKED")
        print(f"ERROR={exc}")
        return 1
    payload = {
        "intake_analysis_import": "PASS",
        "fact_inventory": str(path),
        "facts": len(facts),
        "new_gaps": len(gaps),
    }
    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True))
    else:
        print("INTAKE_ANALYSIS_IMPORT=PASS")
        print(f"FACT_INVENTORY={path}")
        print(f"FACTS={len(facts)}")
        print(f"NEW_GAPS={len(gaps)}")
    return 0


def command_profile_analyze(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    source_ids = args.source_id or []
    work_id = ensure_profile_work(project, source_ids, args.goal)
    try:
        stats = analyze_profiles(
            project,
            source_ids=source_ids,
            work_id=work_id,
            goal=args.goal,
            brand=args.brand,
            company=args.company,
            client=args.client,
        )
    except Exception as exc:  # noqa: BLE001 - CLI should surface actionable analysis failure
        print("PROFILE_ANALYSIS=CHECK")
        print(f"ERROR={exc}")
        return 1
    dashboard = render_dashboard(project) if args.dashboard else None
    errors, validate_stats = validate(project)
    if args.json:
        output = {
            "profile_analysis": "PASS" if not errors else "CHECK",
            "project": str(project),
            "work_id": work_id,
            "profile_current_truth": str(stats["profile_current_truth"]),
            "handoff": str(stats["handoff"]),
            "stats": {key: value for key, value in stats.items() if isinstance(value, (int, str))},
            "dashboard": str(dashboard) if dashboard else "",
            "validation": "PASS" if not errors else "CHECK",
            "validate_stats": validate_stats,
            "errors": errors,
        }
        print(json.dumps(output, ensure_ascii=False, indent=2))
        return 0 if not errors else 1
    print(f"PROJECT={project}")
    print(f"WORK_ID={work_id or 'NOT_CREATED'}")
    print(f"PROFILE_MATERIALS={stats['materials']}")
    print(f"PROFILE_SUBJECTS={stats['subjects']}")
    print(f"PROFILE_VOICES={stats['voices']}")
    print(f"PROFILE_INSIGHTS={stats['insights']}")
    print(f"PROFILE_CONFLICTS={stats['conflicts']}")
    print(f"PROFILE_DEDUPED={stats['deduped']}")
    print(f"PROFILE_CURRENT_TRUTH={stats['profile_current_truth']}")
    print(f"HANDOFF={stats['handoff']}")
    print(f"DASHBOARD={dashboard or 'NOT_RUN'}")
    for key, value in validate_stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def command_hygiene(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    report = workspace_hygiene_report(project)
    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2))
        return 0 if report["status"] == "PASS" or not args.strict else 1
    print(f"WORKSPACE_HYGIENE={report['status']}")
    print(f"PROJECT={report['project']}")
    print(f"GIT_ROOT={report['git_root'] or 'none'}")
    print(f"TRACKED_CHANGES={len(report['tracked_changes'])}")
    print(f"UNTRACKED_FILES={len(report['untracked_files'])}")
    print(f"POLLUTION_PATHS={len(report['pollution_paths'])}")
    print(f"ACTIVE_THREADS={len(report['active_threads'])}")
    if report["issues"]:
        print("ISSUES:")
        for issue in report["issues"]:
            print(f"- {issue}")
    if report["tracked_changes"]:
        print("TRACKED:")
        for line in report["tracked_changes"][:30]:
            print(f"- {line}")
    if report["untracked_files"]:
        print("UNTRACKED:")
        for line in report["untracked_files"][:30]:
            print(f"- {line}")
    if report["pollution_paths"]:
        print("POLLUTION:")
        for path in report["pollution_paths"][:30]:
            print(f"- {path}")
    print("PLAN:")
    for item in report["plan"]:
        print(f"- {item}")
    return 0 if report["status"] == "PASS" or not args.strict else 1


def command_council(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    if args.render_dashboard:
        render_dashboard(project)
    overall, results, report = run_council(project)
    print(f"COUNCIL={overall}")
    print(f"REPORT={report}")
    for result in results:
        print(f"{result.name.upper().replace(' ', '_')}={result.status}")
    return 0 if overall == "PASS" else 1


def command_status(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    if args.json:
        print(json.dumps(status_payload(project), ensure_ascii=False, indent=2, sort_keys=True))
    else:
        print_status(project)
    return 0


def command_next(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    if args.render_dashboard:
        ensure_project(project)
        render_dashboard(project)
    payload = status_payload(project)
    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True))
    else:
        print(f"PROJECT={payload['project']}")
        print(f"PROJECT_SURFACE={payload['surface']}")
        print(f"NEXT_STATUS={payload['next_status']}")
        print(f"NEXT_ACTION={payload['next_action']}")
        print(f"VALIDATION={payload['validation']}")
        if payload["surface"] == DELIVERY_SURFACE:
            print(f"ACTIVE_WORK={payload['active_work_count']}")
        print(f"OPEN_GAPS={payload['open_gap_count']}")
        print(f"BLOCKING_GAPS={payload['blocking_gap_count']}")
        print(f"PENDING_CONFIRMATIONS={payload['pending_confirmation_count']}")
        print(f"DASHBOARD={payload['dashboard'] or 'NOT_RUN'}")
        if payload["next_status"] == "NEEDS_MATERIAL":
            print(f"SUGGESTED_COMMAND=adco run {payload['project']} --material <brief_file_or_folder>")
        elif payload["next_status"] == "READY_FOR_NEXT_GATE":
            print(f"SUGGESTED_COMMAND=adco goal-plan {payload['project']} --title <title> --objective <objective>")
        elif payload["next_status"] == "READY_FOR_CONTENT_WORK":
            print(f"SUGGESTED_ACTION={payload['next_action']}")
    return 1 if payload["next_status"] == "VALIDATION_CHECK" else 0


def command_validate(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    issues, stats = validate_issues(project, strict_legacy=args.strict_legacy)
    blocking_issues = [
        issue for issue in issues if args.strict_legacy or issue.scope != "legacy"
    ]
    errors = [issue.message for issue in blocking_issues]
    status = "PASS" if not errors else "CHECK"
    if args.json:
        print(
            json.dumps(
                {
                    "project": str(project),
                    "validation": status,
                    "validation_scope": "structure_and_traceability_only",
                    "validation_not_creative_quality": True,
                    "validation_not_client_language": True,
                    "validation_not_visual_approval": True,
                    "strict_legacy": args.strict_legacy,
                    "stats": stats,
                    "issues": [issue.as_dict() for issue in issues],
                    "errors": errors,
                },
                ensure_ascii=False,
                indent=2,
                sort_keys=True,
            )
        )
    else:
        for key, value in stats.items():
            print(f"{key.upper()}={value}")
        print(f"VALIDATION={status}")
        print("VALIDATION_SCOPE=structure_and_traceability_only")
        print("VALIDATION_NOT_CREATIVE_QUALITY=1")
        print("VALIDATION_NOT_CLIENT_LANGUAGE=1")
        print("VALIDATION_NOT_VISUAL_APPROVAL=1")
        if issues:
            print("ISSUES:")
            for issue in issues:
                print(
                    f"- [{issue.severity}][{issue.scope}][{issue.code}] {issue.message}"
                )
    if errors:
        if not args.json:
            print("ERRORS:")
            for error in errors:
                print(f"- {error}")
        return 1
    return 0


def command_check(args: argparse.Namespace) -> int:
    import run_checks

    return run_checks.main()


def command_doctor(args: argparse.Namespace) -> int:
    status, issues, warnings, evidence = doctor_report()
    if args.json:
        print(
            json.dumps(
                {
                    "status": status,
                    "issues": issues,
                    "warnings": warnings,
                    "evidence": evidence,
                },
                ensure_ascii=False,
                indent=2,
                sort_keys=True,
            )
        )
    else:
        print(f"ADCO_DOCTOR={status}")
        print(f"ISSUES={len(issues)}")
        print(f"WARNINGS={len(warnings)}")
        for item in evidence:
            print(f"EVIDENCE={item}")
        if warnings:
            print("DOCTOR_WARNINGS:")
            for warning in warnings:
                print(f"- {warning}")
        if issues:
            print("DOCTOR_ISSUES:")
            for issue in issues:
                print(f"- {issue}")
    return 0 if not issues else 1


def command_release_status(args: argparse.Namespace) -> int:
    payload = release_status_payload()
    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True))
    else:
        print(f"RELEASE_STATUS={payload['release_status']}")
        print(f"DOCTOR={payload['doctor_status']}")
        print(f"REMOTE={payload['remote_status']}")
        print(f"MODE={payload['mode']}")
        print(f"GIT_REMOTE={payload['git_remote']}")
        print(f"VERIFY_COMMAND={payload['verify_command']}")
        print(f"NEXT_ACTION={payload['next_action']}")
        warnings = payload["warnings"]
        issues = payload["issues"]
        print(f"WARNINGS={len(warnings)}")
        print(f"ISSUES={len(issues)}")
        if warnings:
            print("RELEASE_WARNINGS:")
            for warning in warnings:
                print(f"- {warning}")
        if issues:
            print("RELEASE_ISSUES:")
            for issue in issues:
                print(f"- {issue}")
    if args.strict and payload["release_status"] != "READY_FOR_REMOTE_CHECKS":
        return 1
    return 1 if payload["release_status"] == "CHECK" else 0


def command_docs(args: argparse.Namespace) -> int:
    payload = docs_payload()
    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True))
    else:
        print(f"DOCS_MODE={payload['mode']}")
        print(f"SOURCE_ROOT={payload['source_root'] or 'NOT_AVAILABLE'}")
        print(f"TEMPLATE_ROOT={payload['template_root']}")
        print(f"SKILL_DRAFT={payload['skill_draft']}")
        docs = payload["docs"]
        if docs:
            print("DOCS:")
            for item in docs:
                exists = "PASS" if item["exists"] else "MISSING"
                print(f"- {item['label']}={item['path']} [{exists}]")
        else:
            print("DOCS=INSTALLED_PACKAGE_COMMAND_HELP_ONLY")
        print("QUICKSTART:")
        for command in payload["quickstart"]:
            print(f"- {command}")
    missing = [item for item in payload["docs"] if not item["exists"]]
    return 1 if missing else 0


def command_support_bundle(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    if not project.exists():
        if args.json:
            print(json.dumps({
                "support_bundle": "CHECK",
                "project": str(project),
                "report": None,
                "validation": "CHECK",
                "stats": {},
                "errors": [f"project not found: {project}"],
            }, ensure_ascii=False, indent=2, sort_keys=True))
            return 1
        print("SUPPORT_BUNDLE=CHECK")
        print(f"ERROR=project not found: {project}")
        return 1
    try:
        report = render_support_bundle(project)
    except ValueError as exc:
        error = str(exc)
        if args.json:
            print(json.dumps({
                "support_bundle": "BLOCKED",
                "project": str(project),
                "report": None,
                "validation": "NOT_RUN",
                "stats": {},
                "errors": [error],
            }, ensure_ascii=False, indent=2, sort_keys=True))
            return 1
        print("SUPPORT_BUNDLE=BLOCKED")
        print(f"ERROR={error}")
        return 1
    errors, stats = validate(project)
    payload = {
        "support_bundle": "PASS" if not errors else "CHECK",
        "project": str(project),
        "report": str(report),
        "validation": "PASS" if not errors else "CHECK",
        "stats": stats,
        "errors": errors,
    }
    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True))
        return 0 if not errors else 1
    print(f"SUPPORT_BUNDLE={'PASS' if not errors else 'CHECK'}")
    print(f"REPORT={report}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def command_audit_dashboard(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    if args.render:
        ensure_project(project)
        render_dashboard(project)
    issues = audit_dashboard(project)
    payload = {
        "dashboard_audit": "PASS" if not issues else "CHECK",
        "project": str(project),
        "dashboard": str(project / DASHBOARD_REL),
        "rendered": bool(args.render),
        "issues": issues,
    }
    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True))
        return 0 if not issues else 1
    print(f"DASHBOARD_AUDIT={'PASS' if not issues else 'CHECK'}")
    print(f"DASHBOARD={project / DASHBOARD_REL}")
    if issues:
        print("ISSUES:")
        for issue in issues:
            print(f"- {issue}")
        return 1
    return 0


def command_add_reference(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    ref_id, action = add_reference(
        project,
        args.url,
        args.title or args.url,
        args.role,
        args.reference_type,
        source_owner=args.source_owner,
        why_relevant=args.why_relevant,
        borrow=args.borrow,
        do_not_copy=args.do_not_copy,
        client_visible=str(args.client_visible).lower(),
        live_check=not args.no_live_check,
    )
    dashboard = render_dashboard(project)
    errors, stats = validate(project)
    print(f"REFERENCE_ID={ref_id}")
    print(f"REFERENCE_ACTION={action}")
    print(f"DASHBOARD={dashboard}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def command_reference_pack_gate(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    status, items, report = review_reference_pack(project, live_check=args.live_check)
    dashboard = render_dashboard(project)
    errors, stats = validate(project)
    print(f"REFERENCE_PACK_GATE={status}")
    print(f"REPORT={report}")
    print(f"FINDINGS={len(items)}")
    print(f"DASHBOARD={dashboard}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors or status == "BLOCKED":
        if items:
            print("GATE_ISSUES:")
            for item in items:
                print(f"- {item}")
        if errors:
            print("ERRORS:")
            for error in errors:
                print(f"- {error}")
        return 1
    return 0


def command_search_quality_gate(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    status, items, report = review_search_quality(project)
    dashboard = render_dashboard(project)
    errors, stats = validate(project)
    print(f"SEARCH_QUALITY_GATE={status}")
    print(f"REPORT={report}")
    print(f"FINDINGS={len(items)}")
    print(f"DASHBOARD={dashboard}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors or status == "BLOCKED":
        if items:
            print("GATE_ISSUES:")
            for item in items:
                print(f"- {item}")
        if errors:
            print("ERRORS:")
            for error in errors:
                print(f"- {error}")
        return 1
    return 0


def command_creative_quality_gate(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    status, items, report = review_creative_quality(project)
    dashboard = render_dashboard(project)
    errors, stats = validate(project)
    print(f"CREATIVE_QUALITY_GATE={status}")
    print(f"REPORT={report}")
    print(f"FINDINGS={len(items)}")
    print(f"DASHBOARD={dashboard}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors or status == "BLOCKED":
        if items:
            print("GATE_ISSUES:")
            for item in items:
                print(f"- {item}")
        if errors:
            print("ERRORS:")
            for error in errors:
                print(f"- {error}")
        return 1
    return 0


def command_add_asset(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    asset_id, target = add_visual_asset(
        project,
        Path(args.file).expanduser().resolve(),
        args.slot_id,
        args.requirement_id,
        args.reference_id,
        args.asset_type,
        args.visibility,
        args.qa_status,
        args.risk_level,
        args.prompt_or_edit_ref,
        args.notes,
        args.selected,
    )
    dashboard = render_dashboard(project)
    errors, stats = validate(project)
    print(f"ASSET_ID={asset_id}")
    print(f"ASSET_PATH={target}")
    print(f"DASHBOARD={dashboard}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def command_import_imagegen(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    source, source_root = resolve_imagegen_source(args.file)
    requirement_id = args.requirement_id or first_existing_id(
        project, "AD-creative/orchestrator/requirements.csv", "requirement_id"
    )
    prompt_ref = args.prompt_or_edit_ref or default_prompt_ref(project)
    asset_id, target = add_visual_asset(
        project,
        source,
        args.slot_id,
        requirement_id,
        args.reference_id,
        "generated_image",
        args.visibility,
        args.qa_status,
        args.risk_level,
        prompt_ref,
        args.notes,
        args.selected,
    )
    log_path = append_imagegen_import_log(
        project,
        source_root,
        source,
        asset_id,
        target,
        prompt_ref,
        args.qa_status,
        args.notes,
    )
    dashboard = render_dashboard(project)
    errors, stats = validate(project)
    print(f"IMAGEGEN_SOURCE_ROOT={source_root}")
    print(f"IMAGEGEN_SOURCE={source}")
    print(f"ASSET_ID={asset_id}")
    print(f"ASSET_PATH={target}")
    print(f"IMAGEGEN_IMPORT_LOG={log_path}")
    print(f"DASHBOARD={dashboard}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def command_visual_quality_gate(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    status, items, report = review_visual_quality(
        project,
        min_long_edge=args.min_long_edge,
        min_short_edge=args.min_short_edge,
    )
    dashboard = render_dashboard(project)
    errors, stats = validate(project)
    print(f"VISUAL_QUALITY_GATE={status}")
    print(f"REPORT={report}")
    print(f"FINDINGS={len(items)}")
    print(f"DASHBOARD={dashboard}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors or status == "BLOCKED":
        if items:
            print("GATE_ISSUES:")
            for item in items:
                print(f"- {item}")
        if errors:
            print("ERRORS:")
            for error in errors:
                print(f"- {error}")
        return 1
    return 0


def command_film_quality_gate(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    status, items, report = review_film_quality(project)
    dashboard = render_dashboard(project)
    errors, stats = validate(project)
    print(f"FILM_QUALITY_GATE={status}")
    print(f"REPORT={report}")
    print(f"FINDINGS={len(items)}")
    print(f"DASHBOARD={dashboard}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors or status == "BLOCKED":
        if items:
            print("GATE_ISSUES:")
            for item in items:
                print(f"- {item}")
        if errors:
            print("ERRORS:")
            for error in errors:
                print(f"- {error}")
        return 1
    return 0


def command_export_pptx(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    output = Path(args.output).expanduser().resolve() if args.output else None
    pptx_path = export_editable_pptx(project, output)
    stats = inspect_pptx(pptx_path)
    dashboard = render_dashboard(project)
    errors, validate_stats = validate(project)
    print(f"PPTX={pptx_path}")
    print(f"PPTX_SLIDES={stats['slides']}")
    print(f"PPTX_EDITABLE_TEXT_RUNS={stats['editable_text_runs']}")
    print(f"PPTX_EDITABLE={'PASS' if stats['editable'] else 'CHECK'}")
    print(f"DASHBOARD={dashboard}")
    for key, value in validate_stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors or not stats["editable"]:
        if errors:
            print("ERRORS:")
            for error in errors:
                print(f"- {error}")
        return 1
    return 0


def command_check_pptx(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    pptx_path = Path(args.file).expanduser().resolve()
    stats = inspect_pptx(pptx_path)
    check_path = write_pptx_check(project, pptx_path, stats)
    dashboard = render_dashboard(project)
    errors, validate_stats = validate(project)
    print(f"PPTX={pptx_path}")
    print(f"PPTX_CHECK={check_path}")
    print(f"PPTX_SLIDES={stats['slides']}")
    print(f"PPTX_EDITABLE_TEXT_RUNS={stats['editable_text_runs']}")
    print(f"PPTX_EDITABLE={'PASS' if stats['editable'] else 'CHECK'}")
    print(f"DASHBOARD={dashboard}")
    for key, value in validate_stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors or not stats["editable"]:
        if errors:
            print("ERRORS:")
            for error in errors:
                print(f"- {error}")
        return 1
    return 0


def command_agency_audit(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    report = agency_audit_report(project)
    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2, sort_keys=True))
    else:
        print(f"AGENCY_AUDIT={report['status']}")
        print(f"PROJECT={report['project']}")
        print(f"P1={report['p1']}")
        print(f"ISSUES={len(report['issues'])}")
        if report["issues"]:
            print("AGENCY_ISSUES:")
            for issue in report["issues"]:
                print(f"- {format_agency_issue(issue)}")
        for key, value in report["stats"].items():
            print(f"{key.upper()}={value}")
    return 1 if args.strict and report["status"] != "PASS" else 0


def command_migrate_control_plane(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    result = migrate_control_plane(project, dry_run=args.dry_run)
    if args.json:
        print(json.dumps(result, ensure_ascii=False, indent=2, sort_keys=True))
    else:
        migration_status = (
            "BLOCKED"
            if result["blockers"]
            else "DRY_RUN"
            if args.dry_run
            else "PASS"
        )
        print(f"MIGRATE_CONTROL_PLANE={migration_status}")
        print(f"PROJECT={result['project']}")
        print(f"SCHEMA_VERSION={result['schema_version']}")
        print(f"APPLIED={int(bool(result['applied']))}")
        print(f"BLOCKED_BEFORE_WRITE={int(bool(result['blocked_before_write']))}")
        print(f"MIGRATION_MANIFEST={result['migration_manifest'] or 'not_required'}")
        print(f"CHANGES={len(result['changes'])}")
        for change in result["changes"]:
            print(f"- {change}")
        if result["blockers"]:
            print("BLOCKERS:")
            for blocker in result["blockers"]:
                print(
                    f"- [{blocker.get('severity')}][{blocker.get('scope')}][{blocker.get('code')}] {blocker.get('message')}"
                )
    return 1 if result["blockers"] else 0


def command_preflight_skill(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    preflight_id = write_specialist_preflight(
        project,
        work_id=args.work_id,
        requested_skill=args.requested_skill,
        skill_path=args.skill_path,
        rules_read=args.rules_read,
        derived_gates=args.derived_gates,
        status=args.status,
        blocked_reason=args.blocked_reason,
    )
    errors, stats = validate(project)
    print(f"SPECIALIST_PREFLIGHT={preflight_id}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    print("VALIDATION_SCOPE=structure_and_traceability_only")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def command_specialist_handoff(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    descriptor = Path(args.descriptor).expanduser().resolve() if args.descriptor else None
    try:
        handoff, path = create_specialist_handoff(
            project,
            work_id=args.work_id,
            profile_id=args.profile_id,
            objective=args.objective,
            input_artifact_ids=args.input_artifact,
            expected_output_kinds=args.expected_output,
            required_capabilities=args.require_capability,
            descriptor_path=descriptor,
            execution_mode=args.execution_mode,
            workspace_mode=args.workspace_mode,
            lane_id=args.lane_id,
            generation_mode=args.generation_mode,
            generation_authorized=args.generation_authorized,
            authorization_ref=args.authorization_ref,
        )
    except (ValueError, OSError) as exc:
        print("SPECIALIST_HANDOFF=BLOCKED")
        print(f"ERROR={exc}")
        return 1
    if handoff.get("contract_version") == V2_CONTRACT_VERSION:
        _, exchange_rows = read_csv_rows(
            project
            / "AD-creative/orchestrator/specialist_exchange/exchange_index.csv"
        )
        row = next(
            item
            for item in exchange_rows
            if item.get("handoff_path") == safe_rel(project, path)
        )
        handoff_id = row.get("handoff_id", "")
        compatibility = row.get("compatibility_status", "")
        execution_mode = row.get("execution_mode", "")
    else:
        handoff_id = str(handoff["handoff_id"])
        compatibility = "compatible" if handoff.get("descriptor_ref") else "unverified"
        execution_mode = str((handoff.get("execution") or {}).get("mode"))
    print("SPECIALIST_HANDOFF=PASS")
    print(f"HANDOFF_ID={handoff_id}")
    print(f"HANDOFF={path}")
    print(f"CONTRACT_VERSION={handoff.get('contract_version')}")
    print(f"COMPATIBILITY={compatibility}")
    print(f"EXECUTION_MODE={execution_mode}")
    return 0


def command_specialist_adopt(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    mappings: dict[str, str] = {}
    try:
        for item in args.map_output:
            provider_artifact_id, target = item.split("=", 1)
            if not provider_artifact_id.strip() or not target.strip():
                raise ValueError
            mappings[provider_artifact_id.strip()] = target.strip()
    except ValueError:
        print("SPECIALIST_ADOPTION=BLOCKED")
        print("ERROR=--map-output must use PROVIDER_ARTIFACT_ID=project/relative/path")
        return 1
    try:
        adoption, path = adopt_specialist_receipt(
            project,
            handoff_path=Path(args.handoff).expanduser().resolve(),
            receipt_path=Path(args.receipt).expanduser().resolve(),
            decision=args.decision,
            reason=args.reason,
            output_mappings=mappings,
            dry_run=args.dry_run,
        )
    except (ValueError, OSError, StopIteration) as exc:
        print("SPECIALIST_ADOPTION=BLOCKED")
        print(f"ERROR={exc}")
        return 1
    print("SPECIALIST_ADOPTION=DRY_RUN" if args.dry_run else "SPECIALIST_ADOPTION=PASS")
    print(f"DECISION={adoption['decision']}")
    print(f"ADOPTION={path or 'not_written'}")
    print(f"ADVANCE_ALLOWED={(adoption.get('gate_effect') or {}).get('advance_allowed')}")
    return 0


def command_preflight_asset(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    browser_assets_exist = normalized_bool(args.browser_checked) or indicates_browser_held_assets(
        args.source_scope,
        args.browser_tool,
        args.download_method,
    )
    imported_asset_ids = split_asset_refs(args.imported_asset_ids)
    status = args.status.strip().upper()
    if browser_assets_exist and not imported_asset_ids:
        if status != "BLOCKED" or not args.blocked_reason.strip():
            print("ASSET_PREFLIGHT=BLOCKED")
            print(
                "ERROR=browser-held assets without imports must use --status BLOCKED "
                "and provide --blocked-reason"
            )
            return 1
    preflight_id = write_asset_preflight(
        project,
        work_id=args.work_id,
        source_scope=args.source_scope,
        local_manifest_checked=args.local_manifest_checked,
        browser_checked=args.browser_checked,
        browser_tool=args.browser_tool,
        download_method=args.download_method,
        imported_asset_ids=args.imported_asset_ids,
        replacement_generation_allowed=args.replacement_generation_allowed,
        status=args.status,
        blocked_reason=args.blocked_reason,
    )
    errors, stats = validate(project)
    print(f"ASSET_PREFLIGHT={preflight_id}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    print("VALIDATION_SCOPE=structure_and_traceability_only")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def command_dispatch_record(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    try:
        result = record_thread_dispatch(
            project,
            lane_id=args.lane_id,
            work_id=args.work_id,
            real_thread_id=args.real_thread_id,
            title_action=args.title_action,
            title_verified_at=args.title_verified_at,
            dispatch_evidence=args.dispatch_evidence,
            dispatch_status=args.dispatch_status,
            absolute_deadline_at=args.absolute_deadline_at,
        )
    except ValueError as exc:
        print("DISPATCH_RECORD=BLOCKED")
        print(f"ERROR={exc}")
        return 1
    errors, stats = validate(project)
    print(f"DISPATCH_RECORD={result['dispatch_status']}")
    print(f"LANE_ID={result['lane_id']}")
    print(f"REAL_THREAD_ID={result['real_thread_id']}")
    print(f"DISPATCH_RECEIPT={result['dispatch_receipt_path']}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    print("VALIDATION_SCOPE=structure_and_traceability_only")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    if errors:
        print("ERRORS:")
        for error in errors:
            print(f"- {error}")
        return 1
    return 0


def command_thread_observe(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    try:
        result = record_thread_observation(
            project,
            lane_id=args.lane_id,
            work_id=args.work_id,
            state=args.state,
            observed_at=args.observed_at,
            evidence=args.evidence,
            absolute_deadline_at=args.absolute_deadline_at,
            extension_reason=args.extension_reason,
            convergence_reminder_sent=args.convergence_reminder_sent,
            rescue_thread_id=args.rescue_thread_id,
        )
    except ValueError as exc:
        print("THREAD_OBSERVE=BLOCKED")
        print(f"ERROR={exc}")
        return 1
    print(f"THREAD_OBSERVE={result['state']}")
    print(f"WORK_ID={result['work_id']}")
    print(f"LANE_ID={result['lane_id']}")
    print(f"ABSOLUTE_DEADLINE_AT={result['absolute_deadline_at']}")
    print(f"BOUNDED_EXTENSION_USED={result['bounded_extension_used']}")
    print(f"RESCUE_COUNT={result['rescue_count']}")
    print(f"CONVERGENCE_LOG={result['log_path']}")
    return 0


def command_thread_reconcile(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    try:
        result = reconcile_thread_receipt(
            project,
            lane_id=args.lane_id,
            work_id=args.work_id,
            receipt_path_value=args.receipt_path,
            adoption_decision=args.adoption_decision,
            rejection_reason=args.rejection_reason,
            reconciled_at=args.reconciled_at,
            cleanup_action=args.cleanup_action,
            archived_at=args.archived_at,
        )
    except ValueError as exc:
        print("THREAD_RECONCILE=BLOCKED")
        print(f"ERROR={exc}")
        return 1
    print(f"THREAD_RECONCILE={result['status']}")
    print(f"RECEIPT_THREAD_ID={result.get('receipt_thread_id', '')}")
    print(f"ADOPTION_DECISION={result.get('adoption_decision', 'REJECT')}")
    print(f"CONVERGENCE_LOG={result['log_path']}")
    if result.get("error"):
        print(f"ERROR={result['error']}")
        return 1
    return 0


def command_client_outline_gate(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    ensure_project(project)
    status, issues, report = review_client_outline(project)
    print(f"CLIENT_OUTLINE_GATE={status}")
    print(f"REPORT={report}")
    print(f"ISSUES={len(issues)}")
    if status != "PASS":
        for issue in issues:
            print(f"- {issue}")
        return 1
    return 0


def command_confirm_client_outline(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    ensure_project(project)
    try:
        receipt = confirm_client_outline(
            project,
            confirmed_by=args.confirmed_by,
            confirmed_at=args.confirmed_at,
            evidence_ref=args.evidence_ref,
        )
        status, issues, report = review_client_outline(project)
    except ValueError as exc:
        print("CLIENT_OUTLINE_CONFIRMATION=BLOCKED")
        print(f"ERROR={exc}")
        return 1
    print(f"CLIENT_OUTLINE_CONFIRMATION={'PASS' if status == 'PASS' else 'BLOCKED'}")
    print(f"RECEIPT={receipt}")
    print(f"CLIENT_OUTLINE_GATE={status}")
    print(f"REPORT={report}")
    if issues:
        for issue in issues:
            print(f"- {issue}")
        return 1
    return 0


def command_client_language_gate(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    extra_paths = [Path(item).expanduser().resolve() for item in args.file]
    status, issues, report = review_client_language(project, extra_paths)
    print(f"CLIENT_LANGUAGE_GATE={status}")
    print(f"REPORT={report}")
    print(f"ISSUES={len(issues)}")
    if status != "PASS":
        for issue in sorted(set(issues))[:80]:
            print(f"- {issue}")
        return 1
    return 0


def command_asset_current_manifest(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    rows, path = refresh_asset_current_manifest(project)
    print("ASSET_CURRENT_MANIFEST=PASS")
    print(f"PATH={path}")
    print(f"ASSETS={len(rows)}")
    return 0


def command_browser_asset_intake(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    if not args.browser_evidence:
        print("BROWSER_ASSET_INTAKE=BLOCKED")
        print("ERROR=--browser-evidence is required when user says browser-held images exist")
        return 1
    if not args.asset_file:
        print("BROWSER_ASSET_INTAKE=BLOCKED")
        print("ERROR=at least one --asset-file is required to register browser-held images")
        return 1
    asset_ids: list[str] = []
    for file_item in args.asset_file:
        asset_id, _ = add_visual_asset(
            project,
            Path(file_item).expanduser().resolve(),
            args.slot_id,
            args.requirement_id,
            args.reference_id,
            args.asset_type,
            "internal_only",
            "PARTIAL_PASS",
            "medium",
            args.browser_evidence,
            f"browser_asset_intake source={args.source}; direct client use requires visual gates",
            selected=False,
        )
        asset_ids.append(asset_id)
    preflight_id = write_asset_preflight(
        project,
        work_id=args.work_id,
        source_scope=args.source,
        local_manifest_checked="yes",
        browser_checked="yes",
        browser_tool=args.browser_tool,
        download_method=args.download_method,
        imported_asset_ids=";".join(asset_ids),
        replacement_generation_allowed="no",
        status="PASS",
        blocked_reason="",
    )
    update_current_asset_metadata(
        project,
        asset_ids,
        source=f"{args.source}: {args.browser_evidence}",
        platform=args.source,
        conversation=args.conversation or args.browser_evidence,
        qa_flags=args.qa_flags or "browser_asset_registered;replacement_generation_blocked;needs_visual_layout_gate",
    )
    print("BROWSER_ASSET_INTAKE=PASS")
    print(f"ASSET_PREFLIGHT={preflight_id}")
    print("ASSET_IDS=" + ";".join(asset_ids))
    return 0


def command_visual_layout_gate(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    status, issues, report = review_visual_layout(
        project,
        min_long_edge=args.min_long_edge,
        min_short_edge=args.min_short_edge,
    )
    print(f"VISUAL_LAYOUT_GATE={status}")
    print(f"REPORT={report}")
    print(f"ISSUES={len(issues)}")
    if status == "BLOCKED":
        for issue in issues:
            print(f"- {issue}")
        return 1
    return 0


def command_dedupe_audit(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    rows, csv_path, report = dedupe_audit(project)
    print("DEDUPE_AUDIT=PASS")
    print(f"CSV={csv_path}")
    print(f"REPORT={report}")
    print(f"FILES={len(rows)}")
    print(f"DUPLICATES={sum(1 for row in rows if row.get('duplicate_of'))}")
    return 0


def command_cleanup_plan(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    plan_path, actions = cleanup_plan(project)
    print("CLEANUP_PLAN=PASS")
    print(f"PLAN={plan_path}")
    print(f"ACTIONS={len(actions)}")
    print("NO_DELETE=1")
    return 0


def command_final_delivery_lock(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    try:
        locked, lock_path = final_delivery_lock(project)
    except RuntimeError as exc:
        print("FINAL_DELIVERY_LOCK=BLOCKED")
        print(
            f"LOCK={project / 'AD-creative/orchestrator/final_delivery_lock.csv'}"
        )
        print("SAFE_PENDING_INVENTORY_PERSISTED=1")
        print(f"ERROR={exc}")
        return 1
    print("FINAL_DELIVERY_LOCK=PASS")
    print(f"LOCK={lock_path}")
    print(f"PROTECTED_FILES={len(locked)}")
    return 0


def command_final_delivery_reconcile(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    ensure_delivery_project(project)
    try:
        row, lock_path = reconcile_final_delivery(
            project,
            old_path=args.old_path,
            new_path=args.new_path,
            kind=args.kind,
            confirmed_by=args.confirmed_by,
            confirmed_at=args.confirmed_at,
            evidence_ref=args.evidence_ref,
            version_id=args.version_id,
        )
    except (RuntimeError, ValueError) as exc:
        print("FINAL_DELIVERY_RECONCILE=BLOCKED")
        print(f"ERROR={exc}")
        return 1
    print("FINAL_DELIVERY_RECONCILE=PASS")
    print(f"LOCK={lock_path}")
    print(f"KIND={row.get('reconciliation_kind')}")
    print(f"NEW_PATH={row.get('path')}")
    print("FILE_MOVE_OR_OVERWRITE=0")
    return 0


def command_client_pack_gate(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    pptx_path = Path(args.pptx).expanduser().resolve() if args.pptx else None
    status, issues, report = review_client_pack(project, pptx_path)
    final_delivery_lock(project)
    dashboard = render_dashboard(project)
    errors, validate_stats = validate(project)
    print(f"CLIENT_PACK_GATE={status}")
    print(f"REPORT={report}")
    print(f"ISSUES={len(issues)}")
    print(f"DASHBOARD={dashboard}")
    for key, value in validate_stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors or status != "PASS":
        if issues:
            print("GATE_ISSUES:")
            for issue in issues:
                print(f"- {issue}")
        if errors:
            print("ERRORS:")
            for error in errors:
                print(f"- {error}")
        return 1
    return 0


def command_client_send_readiness_gate(args: argparse.Namespace) -> int:
    project = Path(args.project).expanduser().resolve()
    ensure_delivery_project(project)
    status, issues, report = review_client_send_readiness(project)
    print(f"CLIENT_SEND_READINESS_GATE={status}")
    print(f"REPORT={report}")
    print(f"ISSUES={len(issues)}")
    print("SEND_EXECUTED=0")
    if status != "PASS":
        for issue in issues[:80]:
            print(f"- {issue}")
        return 1
    return 0


def command_install_skill(args: argparse.Namespace) -> int:
    # Preserve the lexical root so install_global_skill can reject an
    # unapproved symlink before resolving it.
    target = Path(args.target).expanduser() if args.target else DEFAULT_SKILL_INSTALL_DIR
    result = install_global_skill(target)
    print(f"SKILL_INSTALL={'PASS' if result['match'] else 'CHECK'}")
    print(f"SOURCE={result['source']}")
    print(f"TARGET={result['target']}")
    print(f"SOURCE_SHA256={result['source_hash']}")
    print(f"TARGET_SHA256={result['target_hash']}")
    print(f"SOURCE_TREE_SHA256={result['source_tree_hash']}")
    print(f"TARGET_TREE_SHA256={result['target_tree_hash']}")
    print(f"MANAGED_FILES={len(result['managed_files'])}")
    print(f"INSTALL_MANIFEST={result['manifest']}")
    print(f"REMOVED_STALE_FILES={len(result['removed_stale_files'])}")
    print(f"PRESERVED_STALE_FILES={len(result['preserved_stale_files'])}")
    if result["manifest_warning"]:
        print(f"MANIFEST_WARNING={result['manifest_warning']}")
    return 0 if result["match"] else 1


def command_handoff_readiness_gate(args: argparse.Namespace) -> int:
    project = Path(args.project).resolve()
    ensure_project(project)
    status, blockers, warnings, report = review_handoff_readiness(project)
    dashboard = render_dashboard(project)
    errors, stats = validate(project)
    print(f"HANDOFF_READINESS_GATE={status}")
    print(f"REPORT={report}")
    print(f"BLOCKERS={len(blockers)}")
    print(f"WARNINGS={len(warnings)}")
    print(f"DASHBOARD={dashboard}")
    for key, value in stats.items():
        print(f"{key.upper()}={value}")
    print(f"VALIDATION={'PASS' if not errors else 'CHECK'}")
    if errors or status != "PASS":
        if blockers:
            print("GATE_BLOCKERS:")
            for blocker in blockers:
                print(f"- {blocker}")
        if warnings:
            print("GATE_WARNINGS:")
            for warning in warnings:
                print(f"- {warning}")
        if errors:
            print("ERRORS:")
            for error in errors:
                print(f"- {error}")
        return 1
    return 0


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Operate an Ad Creative Orchestrator project without editing CSVs by hand."
    )
    parser.add_argument("--version", action="version", version=f"adco {package_version()}")
    subparsers = parser.add_subparsers(dest="command", required=True)

    init_parser = subparsers.add_parser("init", help="Initialize a project from packaged templates.")
    init_parser.add_argument("project", help="Project directory.")
    init_parser.add_argument("--template", default="", help="Template directory. Defaults to packaged templates.")
    init_parser.add_argument(
        "--full",
        action="store_true",
        help="Initialize the delivery surface instead of the default content surface.",
    )
    init_parser.set_defaults(func=command_init)

    goal_parser = subparsers.add_parser("goal-plan", help="Create a reusable goal iteration execution plan.")
    goal_parser.add_argument("project", help="Project directory.")
    goal_parser.add_argument("--goal-id", default="", help="Stable goal id. Defaults to timestamp.")
    goal_parser.add_argument("--title", default="", help="Human-readable goal title.")
    goal_parser.add_argument("--objective", default="", help="Goal objective.")
    goal_parser.add_argument("--owner", default="Main Controller", help="Goal owner.")
    goal_parser.add_argument("--force", action="store_true", help="Overwrite an existing goal plan with the same id.")
    goal_parser.set_defaults(func=command_goal_plan)

    thread_plan_parser = subparsers.add_parser(
        "thread-plan",
        help="Create a Codex ThreadOps lane plan, role briefs, worker prompts, and registry rows.",
    )
    thread_plan_parser.add_argument("project", help="Project directory.")
    thread_plan_parser.add_argument("--goal-id", default="", help="Stable goal id. Defaults to timestamp.")
    thread_plan_parser.add_argument("--work-id", default="", help="Stable work id. Defaults from goal id.")
    thread_plan_parser.add_argument("--task-signature-id", default="", help="Stable task signature id. Defaults from goal id.")
    thread_plan_parser.add_argument("--title", default="", help="Human-readable goal title.")
    thread_plan_parser.add_argument("--objective", default="", help="Goal objective.")
    thread_plan_parser.add_argument(
        "--roles",
        default=",".join(THREADOPS_DEFAULT_ROLES),
        help="Comma-separated roles. Choices: brand_client, copy_creative, film_director, art_design, producer_risk, qa_review.",
    )
    thread_plan_parser.add_argument("--brand", default="")
    thread_plan_parser.add_argument("--product", default="")
    thread_plan_parser.add_argument("--talent-or-ip", default="")
    thread_plan_parser.add_argument("--platform-or-channel", default="")
    thread_plan_parser.add_argument("--deliverable", default="")
    thread_plan_parser.add_argument("--stage", default="threadops")
    thread_plan_parser.add_argument("--primary-risks", default="")
    thread_plan_parser.add_argument("--evidence-needed", default="")
    thread_plan_parser.add_argument("--master-thread-id", default="")
    thread_plan_parser.add_argument("--current-version-id", default="")
    thread_plan_parser.add_argument("--max-active", type=int, default=3)
    thread_plan_parser.add_argument("--read-first", action="append", default=[], help="Extra read-first file for every lane. Repeatable.")
    thread_plan_parser.add_argument("--force", action="store_true", help="Replace an existing thread_lane_plan.md for this project.")
    thread_plan_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    thread_plan_parser.set_defaults(func=command_thread_plan)

    goal_run_parser = subparsers.add_parser("goal-run", help="Run deterministic local goal steps until a safe stop condition.")
    goal_run_parser.add_argument("project", help="Project directory.")
    goal_run_parser.add_argument("--goal-id", default="latest", help="Goal id or latest.")
    goal_run_parser.add_argument("--max-steps", type=int, default=3, help="Maximum deterministic steps to execute.")
    goal_run_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    goal_run_parser.add_argument("--allow-generate", action="store_true", help="Allow generation-capable goal steps. Defaults to off.")
    goal_run_parser.set_defaults(func=command_goal_run)

    creative_doctor_parser = subparsers.add_parser("creative-doctor", help="Diagnose optional Creative Production bridge availability.")
    creative_doctor_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    creative_doctor_parser.set_defaults(func=command_creative_doctor)

    creative_run_parser = subparsers.add_parser("creative-run", help="Create a Creative Production review-only or generation run.")
    creative_run_parser.add_argument("project", help="Project directory.")
    creative_run_parser.add_argument("--kind", choices=sorted(CREATIVE_PRODUCTION_KINDS), required=True)
    creative_run_parser.add_argument("--work-id", required=True)
    creative_run_parser.add_argument("--brief-file", required=True)
    creative_run_parser.add_argument("--base-asset", default="", help="Optional source image for shots or ads.")
    creative_run_parser.add_argument("--review-only", action="store_true", help="Write manifests/review surface without generation. This is the default.")
    creative_run_parser.add_argument("--generate", action="store_true", help="Run Creative Production generation. Outputs remain internal_only.")
    creative_run_parser.add_argument("--force", action="store_true", help="Replace output directory if needed.")
    creative_run_parser.set_defaults(func=command_creative_run)

    creative_import_parser = subparsers.add_parser("import-creative-production", help="Import Creative Production run metadata and images into ADCO manifests.")
    creative_import_parser.add_argument("project", help="Project directory.")
    creative_import_parser.add_argument("--run-dir", required=True)
    creative_import_parser.add_argument("--kind", choices=sorted(CREATIVE_PRODUCTION_KINDS), required=True)
    creative_import_parser.add_argument("--slot-prefix", default="CP")
    creative_import_parser.add_argument("--requirement-id", default="")
    creative_import_parser.add_argument("--reference-id", default="pending")
    creative_import_parser.set_defaults(func=command_import_creative_production)

    proposal_parser = subparsers.add_parser(
        "creative-proposal",
        help="Deprecated alias for creative-brief; does not generate directions.",
    )
    proposal_parser.add_argument("project", help="Project directory.")
    proposal_parser.add_argument("--work-id", default="", help="Optional work item id to link artifacts.")
    proposal_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    proposal_parser.set_defaults(func=command_creative_proposal)

    creative_brief_parser = subparsers.add_parser(
        "creative-brief",
        help="Create an evidence snapshot, creative contract, candidate schema, generation request, and open gaps.",
    )
    creative_brief_parser.add_argument("project", help="Project directory.")
    creative_brief_parser.add_argument("--work-id", default="")
    creative_brief_parser.add_argument("--json", action="store_true")
    creative_brief_parser.set_defaults(func=command_creative_brief)

    creative_candidate_parser = subparsers.add_parser(
        "creative-import",
        help="Import 2-3 evidence-bound post-Critic creative candidates.",
    )
    creative_candidate_parser.add_argument("project", help="Project directory.")
    creative_candidate_parser.add_argument("--file", required=True)
    creative_candidate_parser.add_argument("--json", action="store_true")
    creative_candidate_parser.set_defaults(func=command_creative_import)

    creative_review_parser = subparsers.add_parser(
        "creative-review",
        help="Write a Critic Receipt for structure, mechanism, ownership, visual clarity, and shootability lint.",
    )
    creative_review_parser.add_argument("project", help="Project directory.")
    creative_review_parser.add_argument("--json", action="store_true")
    creative_review_parser.set_defaults(func=command_creative_review)

    run_parser = subparsers.add_parser(
        "run",
        help="Initialize the content surface, parse evidence, return a useful summary, and run scoped validation.",
    )
    run_parser.add_argument("project", help="Project directory.")
    run_parser.add_argument("--material", action="append", default=[], help="Client material file or folder. Repeatable.")
    run_parser.add_argument("--goal", default="先完成需求整理、缺口判断、客户追问、下一步建议。")
    run_parser.add_argument("--max-total-chars", type=int, default=2_000_000)
    run_parser.add_argument(
        "--dashboard",
        action="store_true",
        help="Render the optional dashboard after the content answer is ready.",
    )
    run_parser.add_argument("--json", action="store_true", help="Print machine-readable output with phase timings.")
    run_parser.set_defaults(func=command_run)

    sample_parser = subparsers.add_parser("sample", help="Create a runnable bundled sample project.")
    sample_parser.add_argument("project", help="Project directory.")
    sample_parser.add_argument("--title", default="Bundled sample dual-lane run", help="Sample goal title.")
    sample_parser.add_argument("--goal-id", default="", help="Stable sample goal id. Defaults to timestamp.")
    sample_parser.add_argument("--force-material", action="store_true", help="Overwrite the bundled sample brief.")
    sample_parser.add_argument("--force-goal", action="store_true", help="Overwrite an existing sample goal plan with the same id.")
    sample_parser.set_defaults(func=command_sample)

    demo_parser = subparsers.add_parser("demo", help="Create a sample project and open the dashboard.")
    demo_parser.add_argument("project", nargs="?", help="Project directory. Defaults to the system temp adco-demo folder.")
    demo_parser.add_argument("--title", default="Bundled sample dual-lane run", help="Sample goal title.")
    demo_parser.add_argument("--goal-id", default="", help="Stable sample goal id. Defaults to timestamp.")
    demo_parser.add_argument("--force-material", action="store_true", help="Overwrite the bundled sample brief.")
    demo_parser.add_argument("--force-goal", action="store_true", help="Overwrite an existing sample goal plan with the same id.")
    demo_parser.add_argument("--no-open", action="store_true", help="Render and validate without opening a browser.")
    demo_parser.set_defaults(func=command_demo)

    quickstart_parser = subparsers.add_parser(
        "quickstart",
        help="Create a demo project, validate it, open the dashboard, and print first next steps.",
    )
    quickstart_parser.add_argument("project", nargs="?", help="Project directory. Defaults to the system temp adco-demo folder.")
    quickstart_parser.add_argument("--title", default="Bundled sample dual-lane run", help="Sample goal title.")
    quickstart_parser.add_argument("--goal-id", default="", help="Stable sample goal id. Defaults to timestamp.")
    quickstart_parser.add_argument("--force-material", action="store_true", help="Overwrite the bundled sample brief.")
    quickstart_parser.add_argument("--force-goal", action="store_true", help="Overwrite an existing sample goal plan with the same id.")
    quickstart_parser.add_argument("--no-open", action="store_true", help="Render and validate without opening a browser.")
    quickstart_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    quickstart_parser.set_defaults(func=command_quickstart)

    status_parser = subparsers.add_parser("status", help="Print current project status.")
    status_parser.add_argument("project", help="Project directory.")
    status_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    status_parser.set_defaults(func=command_status)

    next_parser = subparsers.add_parser("next", help="Print the next safe action for a project.")
    next_parser.add_argument("project", help="Project directory.")
    next_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    next_parser.add_argument("--render-dashboard", action="store_true", help="Render dashboard before deciding.")
    next_parser.set_defaults(func=command_next)

    validate_parser = subparsers.add_parser("validate", help="Validate a project directory.")
    validate_parser.add_argument("project", help="Project directory.")
    validate_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    validate_parser.add_argument(
        "--strict-legacy",
        action="store_true",
        help="Treat grouped legacy-only debt as blocking.",
    )
    validate_parser.set_defaults(func=command_validate)

    check_parser = subparsers.add_parser("check", help="Run the full verification suite.")
    check_parser.set_defaults(func=command_check)

    doctor_parser = subparsers.add_parser("doctor", help="Diagnose installation, templates, optional dependencies, and release blockers.")
    doctor_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    doctor_parser.set_defaults(func=command_doctor)

    release_parser = subparsers.add_parser("release-status", help="Summarize local release readiness and remote blockers.")
    release_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    release_parser.add_argument("--strict", action="store_true", help="Return non-zero unless ready for remote checks.")
    release_parser.set_defaults(func=command_release_status)

    docs_parser = subparsers.add_parser("docs", help="Print local documentation paths and quickstart commands.")
    docs_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    docs_parser.set_defaults(func=command_docs)

    support_parser = subparsers.add_parser("support-bundle", help="Write a sanitized support bundle for bug reports.")
    support_parser.add_argument("project", help="Project directory.")
    support_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    support_parser.set_defaults(func=command_support_bundle)

    dashboard_parser = subparsers.add_parser("render-dashboard", help="Render static operation dashboard.")
    dashboard_parser.add_argument("project", help="Project directory.")
    dashboard_parser.set_defaults(func=command_render_dashboard)

    open_dashboard_parser = subparsers.add_parser("open-dashboard", help="Render and open the operation dashboard.")
    open_dashboard_parser.add_argument("project", help="Project directory.")
    open_dashboard_parser.add_argument("--no-open", action="store_true", help="Render and validate without opening a browser.")
    open_dashboard_parser.set_defaults(func=command_open_dashboard)

    intake_parser = subparsers.add_parser("intake", help="Extract first-pass requirements and gaps from registered materials.")
    intake_parser.add_argument("project", help="Project directory.")
    intake_parser.add_argument("--source-id", action="append", default=[], help="Registered source_event_id to process. Repeatable.")
    intake_parser.add_argument("--goal", default="先完成需求整理、缺口判断、客户追问、下一步建议。")
    intake_parser.add_argument("--max-total-chars", type=int, default=2_000_000, help="Maximum total extracted characters; overflow is reported by file and media type.")
    intake_parser.set_defaults(func=command_intake)

    intake_evidence_parser = subparsers.add_parser(
        "intake-evidence",
        help="Parse registered materials into provenance-bound evidence chunks.",
    )
    intake_evidence_parser.add_argument("project", help="Project directory.")
    intake_evidence_parser.add_argument("--source-id", action="append", default=[])
    intake_evidence_parser.add_argument("--max-total-chars", type=int, default=2_000_000)
    intake_evidence_parser.add_argument("--json", action="store_true")
    intake_evidence_parser.set_defaults(func=command_intake_evidence)

    analysis_export_parser = subparsers.add_parser(
        "export-intake-analysis-request",
        help="Export the evidence snapshot and schema for GPT-5.6 Sol fact analysis.",
    )
    analysis_export_parser.add_argument("project", help="Project directory.")
    analysis_export_parser.add_argument("--json", action="store_true")
    analysis_export_parser.set_defaults(func=command_export_intake_analysis_request)

    analysis_import_parser = subparsers.add_parser(
        "import-intake-analysis",
        help="Import evidence-bound fact analysis and derive only supported gaps.",
    )
    analysis_import_parser.add_argument("project", help="Project directory.")
    analysis_import_parser.add_argument("--file", required=True)
    analysis_import_parser.add_argument("--json", action="store_true")
    analysis_import_parser.set_defaults(func=command_import_intake_analysis)

    profile_parser = subparsers.add_parser(
        "profile-analyze",
        help="Analyze meeting/client materials into participant, brand, company, decision, and conflict profiles.",
    )
    profile_parser.add_argument("project", help="Project directory.")
    profile_parser.add_argument("--source-id", action="append", default=[], help="Registered source_event_id to process. Repeatable.")
    profile_parser.add_argument("--goal", default="分析会议资料中的人物画像、品牌画像、需求权重、决策权和分歧融合路径。")
    profile_parser.add_argument("--brand", default="", help="Brand name to use for brand profile.")
    profile_parser.add_argument("--company", default="", help="Company or client organization name.")
    profile_parser.add_argument("--client", default="", help="Client group name if the company/brand is unclear.")
    profile_parser.add_argument(
        "--dashboard",
        action="store_true",
        help="Render the optional dashboard after profile analysis.",
    )
    profile_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    profile_parser.set_defaults(func=command_profile_analyze)

    hygiene_parser = subparsers.add_parser(
        "hygiene",
        help="Audit workspace cleanliness without deleting files.",
    )
    hygiene_parser.add_argument("project", help="Project or repo directory.")
    hygiene_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    hygiene_parser.add_argument("--strict", action="store_true", help="Return non-zero when hygiene status is CHECK.")
    hygiene_parser.set_defaults(func=command_hygiene)

    agency_audit_parser = subparsers.add_parser(
        "agency-audit",
        help="Audit Agency control-plane, ThreadOps dispatch proof, and open blockers.",
    )
    agency_audit_parser.add_argument("project", help="Project directory.")
    agency_audit_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    agency_audit_parser.add_argument("--strict", action="store_true", help="Return non-zero when audit status is CHECK.")
    agency_audit_parser.set_defaults(func=command_agency_audit)

    migrate_parser = subparsers.add_parser(
        "migrate-control-plane",
        help="Add missing Agency/ThreadOps/client gate control-plane files and CSV columns.",
    )
    migrate_parser.add_argument("project", help="Project directory.")
    migrate_parser.add_argument("--dry-run", action="store_true", help="Report changes without writing files.")
    migrate_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    migrate_parser.set_defaults(func=command_migrate_control_plane)

    audit_parser = subparsers.add_parser("audit-dashboard", help="Audit dashboard usability markers.")
    audit_parser.add_argument("project", help="Project directory.")
    audit_parser.add_argument("--render", action="store_true", help="Render dashboard before auditing.")
    audit_parser.add_argument("--json", action="store_true", help="Print machine-readable JSON.")
    audit_parser.set_defaults(func=command_audit_dashboard)

    council_parser = subparsers.add_parser("council", help="Run three-council readiness audit.")
    council_parser.add_argument("project", help="Project directory.")
    council_parser.add_argument("--render-dashboard", action="store_true", help="Render dashboard before the audit.")
    council_parser.set_defaults(func=command_council)

    ref_parser = subparsers.add_parser("add-reference", help="Register a live https reference link.")
    ref_parser.add_argument("project", help="Project directory.")
    ref_parser.add_argument("--url", required=True, help="HTTPS reference URL.")
    ref_parser.add_argument("--title", default="", help="Human-readable title.")
    ref_parser.add_argument("--role", default="direction_reference", help="Reference role in the project.")
    ref_parser.add_argument("--reference-type", default="official_or_public_reference")
    ref_parser.add_argument("--source-owner", default="official_or_public")
    ref_parser.add_argument("--why-relevant", default="")
    ref_parser.add_argument("--borrow", default="")
    ref_parser.add_argument("--do-not-copy", default="")
    ref_parser.add_argument("--client-visible", action="store_true")
    ref_parser.add_argument("--no-live-check", action="store_true", help="Skip HTTP check for offline validation.")
    ref_parser.set_defaults(func=command_add_reference)

    ref_gate_parser = subparsers.add_parser("reference-pack-gate", help="Audit reference pack quality before client use.")
    ref_gate_parser.add_argument("project", help="Project directory.")
    ref_gate_parser.add_argument("--live-check", action="store_true", help="Run live HTTPS checks for registered URLs.")
    ref_gate_parser.set_defaults(func=command_reference_pack_gate)

    search_gate_parser = subparsers.add_parser("search-quality-gate", help="Audit search plans and search-target references.")
    search_gate_parser.add_argument("project", help="Project directory.")
    search_gate_parser.set_defaults(func=command_search_quality_gate)

    creative_quality_parser = subparsers.add_parser(
        "creative-quality-gate",
        help="Audit creative/proposal artifacts for traceable quality and client-facing safety.",
    )
    creative_quality_parser.add_argument("project", help="Project directory.")
    creative_quality_parser.set_defaults(func=command_creative_quality_gate)

    preflight_skill_parser = subparsers.add_parser(
        "preflight-skill",
        help="Record specialist skill selection, rules read, and derived gates.",
    )
    preflight_skill_parser.add_argument("project", help="Project directory.")
    preflight_skill_parser.add_argument("--work-id", required=True)
    preflight_skill_parser.add_argument("--requested-skill", required=True)
    preflight_skill_parser.add_argument("--skill-path", required=True)
    preflight_skill_parser.add_argument("--rules-read", required=True)
    preflight_skill_parser.add_argument("--derived-gates", default="")
    preflight_skill_parser.add_argument("--status", default="PASS")
    preflight_skill_parser.add_argument("--blocked-reason", default="")
    preflight_skill_parser.set_defaults(func=command_preflight_skill)

    specialist_handoff_parser = subparsers.add_parser(
        "specialist-handoff",
        help="Create a versioned neutral specialist handoff; DIRcreative is one profile, not a runtime dependency.",
    )
    specialist_handoff_parser.add_argument("project", help="Project directory.")
    specialist_handoff_parser.add_argument("--work-id", required=True)
    specialist_handoff_parser.add_argument("--profile-id", default=DIRCREATIVE_PROFILE_ID)
    specialist_handoff_parser.add_argument("--objective", required=True)
    specialist_handoff_parser.add_argument("--input-artifact", action="append", default=[], required=True)
    specialist_handoff_parser.add_argument("--expected-output", action="append", default=[], required=True)
    specialist_handoff_parser.add_argument("--require-capability", action="append", default=[])
    specialist_handoff_parser.add_argument("--descriptor", default="")
    specialist_handoff_parser.add_argument(
        "--execution-mode",
        choices=["inline", "codex_thread", "external_handoff"],
        default="inline",
    )
    specialist_handoff_parser.add_argument(
        "--workspace-mode",
        choices=["read_only", "isolated_workspace", "worktree"],
        default="isolated_workspace",
    )
    specialist_handoff_parser.add_argument("--lane-id", default="")
    specialist_handoff_parser.add_argument("--generation-mode", default="prompt_only")
    specialist_handoff_parser.add_argument("--generation-authorized", action="store_true")
    specialist_handoff_parser.add_argument("--authorization-ref", default="")
    specialist_handoff_parser.set_defaults(func=command_specialist_handoff)

    specialist_adopt_parser = subparsers.add_parser(
        "specialist-adopt",
        help="Validate a specialist receipt and record ADCO-owned adoption separately.",
    )
    specialist_adopt_parser.add_argument("project", help="Project directory.")
    specialist_adopt_parser.add_argument("--handoff", required=True)
    specialist_adopt_parser.add_argument("--receipt", required=True)
    specialist_adopt_parser.add_argument(
        "--decision",
        choices=["adopt", "partial_adopt", "reject", "defer"],
        required=True,
    )
    specialist_adopt_parser.add_argument("--reason", required=True)
    specialist_adopt_parser.add_argument("--map-output", action="append", default=[])
    specialist_adopt_parser.add_argument("--dry-run", action="store_true")
    specialist_adopt_parser.set_defaults(func=command_specialist_adopt)

    preflight_asset_parser = subparsers.add_parser(
        "preflight-asset",
        help="Record local/browser/download/generated asset inventory before replacement generation.",
    )
    preflight_asset_parser.add_argument("project", help="Project directory.")
    preflight_asset_parser.add_argument("--work-id", required=True)
    preflight_asset_parser.add_argument("--source-scope", required=True)
    preflight_asset_parser.add_argument("--local-manifest-checked", default="yes")
    preflight_asset_parser.add_argument("--browser-checked", default="no")
    preflight_asset_parser.add_argument("--browser-tool", default="")
    preflight_asset_parser.add_argument("--download-method", default="")
    preflight_asset_parser.add_argument("--imported-asset-ids", default="")
    preflight_asset_parser.add_argument("--replacement-generation-allowed", default="no")
    preflight_asset_parser.add_argument("--status", default="PASS")
    preflight_asset_parser.add_argument("--blocked-reason", default="")
    preflight_asset_parser.set_defaults(func=command_preflight_asset)

    confirm_outline_parser = subparsers.add_parser(
        "confirm-client-outline",
        help="Record explicit human/client approval bound to the current client outline hash.",
    )
    confirm_outline_parser.add_argument("project", help="Project directory.")
    confirm_outline_parser.add_argument("--confirmed-by", required=True)
    confirm_outline_parser.add_argument("--confirmed-at", required=True)
    confirm_outline_parser.add_argument(
        "--evidence-ref",
        required=True,
        help="user_confirmation:<id> or client_confirmation:<id>",
    )
    confirm_outline_parser.set_defaults(func=command_confirm_client_outline)

    client_outline_parser = subparsers.add_parser(
        "client-outline-gate",
        help="Block PPT builder until client-readable page outline is complete.",
    )
    client_outline_parser.add_argument("project", help="Project directory.")
    client_outline_parser.set_defaults(func=command_client_outline_gate)

    client_language_parser = subparsers.add_parser(
        "client-language-gate",
        help="Block client export when client-visible copy leaks internal execution language.",
    )
    client_language_parser.add_argument("project", help="Project directory.")
    client_language_parser.add_argument("--file", action="append", default=[], help="Extra text file to scan. Repeatable.")
    client_language_parser.set_defaults(func=command_client_language_gate)

    asset_current_parser = subparsers.add_parser(
        "asset-current-manifest",
        help="Refresh the unique current asset manifest with source/hash/approval/use columns.",
    )
    asset_current_parser.add_argument("project", help="Project directory.")
    asset_current_parser.set_defaults(func=command_asset_current_manifest)

    browser_intake_parser = subparsers.add_parser(
        "browser-asset-intake",
        help="Register browser-held Grok/ChatGPT/ImageGen assets before replacement generation.",
    )
    browser_intake_parser.add_argument("project", help="Project directory.")
    browser_intake_parser.add_argument("--work-id", required=True)
    browser_intake_parser.add_argument("--source", required=True, help="Browser asset source, e.g. Grok or ChatGPT.")
    browser_intake_parser.add_argument("--browser-evidence", required=True)
    browser_intake_parser.add_argument("--conversation", default="", help="Browser conversation/project/canvas identifier or URL.")
    browser_intake_parser.add_argument("--qa-flags", default="", help="Initial QA flags for imported browser assets.")
    browser_intake_parser.add_argument("--browser-tool", default="browser")
    browser_intake_parser.add_argument("--download-method", default="manual_supported_download")
    browser_intake_parser.add_argument("--asset-file", action="append", default=[], help="Downloaded asset file to import. Repeatable.")
    browser_intake_parser.add_argument("--slot-id", default="BROWSER-ASSET")
    browser_intake_parser.add_argument("--requirement-id", default="")
    browser_intake_parser.add_argument("--reference-id", default="pending")
    browser_intake_parser.add_argument("--asset-type", default="browser_download")
    browser_intake_parser.set_defaults(func=command_browser_asset_intake)

    visual_layout_parser = subparsers.add_parser(
        "visual-layout-gate",
        help="Audit actual client deck layout risks before client packaging.",
    )
    visual_layout_parser.add_argument("project", help="Project directory.")
    visual_layout_parser.add_argument("--min-long-edge", type=int, default=900)
    visual_layout_parser.add_argument("--min-short-edge", type=int, default=600)
    visual_layout_parser.set_defaults(func=command_visual_layout_gate)

    dedupe_parser = subparsers.add_parser(
        "dedupe-audit",
        help="Classify duplicate/cache/preview/final files without deleting anything.",
    )
    dedupe_parser.add_argument("project", help="Project directory.")
    dedupe_parser.set_defaults(func=command_dedupe_audit)

    cleanup_parser = subparsers.add_parser(
        "cleanup-plan",
        help="Write a non-destructive cleanup plan with FinalDelivery lock evidence.",
    )
    cleanup_parser.add_argument("project", help="Project directory.")
    cleanup_parser.set_defaults(func=command_cleanup_plan)

    final_lock_parser = subparsers.add_parser(
        "final-delivery-lock",
        help="Hash-register user-placed FinalDelivery files as protected.",
    )
    final_lock_parser.add_argument("project", help="Project directory.")
    final_lock_parser.set_defaults(func=command_final_delivery_lock)

    final_reconcile_parser = subparsers.add_parser(
        "final-delivery-reconcile",
        help="Evidence-bind a missing FinalDelivery baseline to an explicit rename or supersession without moving files.",
    )
    final_reconcile_parser.add_argument("project", help="Project directory.")
    final_reconcile_parser.add_argument("--old-path", required=True)
    final_reconcile_parser.add_argument("--new-path", required=True)
    final_reconcile_parser.add_argument(
        "--kind", choices=["rename", "supersession"], required=True
    )
    final_reconcile_parser.add_argument("--confirmed-by", required=True)
    final_reconcile_parser.add_argument("--confirmed-at", required=True)
    final_reconcile_parser.add_argument("--evidence-ref", required=True)
    final_reconcile_parser.add_argument("--version-id", default="")
    final_reconcile_parser.set_defaults(func=command_final_delivery_reconcile)

    dispatch_record_parser = subparsers.add_parser(
        "dispatch-record",
        help="Backfill real Codex Thread dispatch proof into thread_registry.csv and dispatch receipt.",
    )
    dispatch_record_parser.add_argument("project", help="Project directory.")
    dispatch_record_parser.add_argument("--lane-id", required=True)
    dispatch_record_parser.add_argument("--work-id", required=True)
    dispatch_record_parser.add_argument("--real-thread-id", required=True)
    dispatch_record_parser.add_argument("--title-action", default="dispatcher_set")
    dispatch_record_parser.add_argument("--title-verified-at", required=True)
    dispatch_record_parser.add_argument("--dispatch-evidence", required=True)
    dispatch_record_parser.add_argument("--dispatch-status", default="dispatched")
    dispatch_record_parser.add_argument("--absolute-deadline-at", required=True)
    dispatch_record_parser.set_defaults(func=command_dispatch_record)

    thread_observe_parser = subparsers.add_parser(
        "thread-observe",
        help="Record bounded thread readback state without treating a fixed poll count as failure.",
    )
    thread_observe_parser.add_argument("project", help="Project directory.")
    thread_observe_parser.add_argument("--lane-id", required=True)
    thread_observe_parser.add_argument("--work-id", required=True)
    thread_observe_parser.add_argument(
        "--state", choices=sorted(THREADOPS_OBSERVATION_STATES), required=True
    )
    thread_observe_parser.add_argument("--observed-at", required=True)
    thread_observe_parser.add_argument("--evidence", required=True)
    thread_observe_parser.add_argument("--absolute-deadline-at", default="")
    thread_observe_parser.add_argument("--extension-reason", default="")
    thread_observe_parser.add_argument("--convergence-reminder-sent", action="store_true")
    thread_observe_parser.add_argument("--rescue-thread-id", default="")
    thread_observe_parser.set_defaults(func=command_thread_observe)

    thread_reconcile_parser = subparsers.add_parser(
        "thread-reconcile",
        help="Validate worker receipt identity and record main adoption/rejection plus cleanup.",
    )
    thread_reconcile_parser.add_argument("project", help="Project directory.")
    thread_reconcile_parser.add_argument("--lane-id", required=True)
    thread_reconcile_parser.add_argument("--work-id", required=True)
    thread_reconcile_parser.add_argument("--receipt-path", required=True)
    thread_reconcile_parser.add_argument(
        "--adoption-decision", choices=sorted(THREADOPS_ADOPTION_DECISIONS), required=True
    )
    thread_reconcile_parser.add_argument("--rejection-reason", default="")
    thread_reconcile_parser.add_argument("--reconciled-at", required=True)
    thread_reconcile_parser.add_argument("--cleanup-action", required=True)
    thread_reconcile_parser.add_argument("--archived-at", default="")
    thread_reconcile_parser.set_defaults(func=command_thread_reconcile)

    asset_parser = subparsers.add_parser("add-asset", help="Register a real/generated visual asset file.")
    asset_parser.add_argument("project", help="Project directory.")
    asset_parser.add_argument("--file", required=True, help="Image file to copy into AD-creative/visual_assets.")
    asset_parser.add_argument("--slot-id", default="AUTO-SLOT")
    asset_parser.add_argument("--requirement-id", default="")
    asset_parser.add_argument("--reference-id", default="")
    asset_parser.add_argument("--asset-type", default="generated_image")
    asset_parser.add_argument("--visibility", default="internal_only")
    asset_parser.add_argument("--qa-status", default="PASS")
    asset_parser.add_argument("--risk-level", default="medium")
    asset_parser.add_argument("--prompt-or-edit-ref", default="")
    asset_parser.add_argument("--notes", default="")
    asset_parser.add_argument("--selected", action="store_true")
    asset_parser.set_defaults(func=command_add_asset)

    imagegen_parser = subparsers.add_parser(
        "import-imagegen",
        help="Import a built-in image_gen output from CODEX_HOME/generated_images.",
    )
    imagegen_parser.add_argument("project", help="Project directory.")
    imagegen_parser.add_argument(
        "--file",
        default="",
        help="Generated image under CODEX_HOME/generated_images. Uses the latest generated image when omitted.",
    )
    imagegen_parser.add_argument("--slot-id", default="AUTO-IMAGEGEN")
    imagegen_parser.add_argument("--requirement-id", default="")
    imagegen_parser.add_argument("--reference-id", default="pending")
    imagegen_parser.add_argument("--visibility", default="internal_only")
    imagegen_parser.add_argument("--qa-status", default="PARTIAL_PASS")
    imagegen_parser.add_argument("--risk-level", default="medium")
    imagegen_parser.add_argument("--prompt-or-edit-ref", default="")
    imagegen_parser.add_argument(
        "--notes",
        default="Imported from built-in image_gen output; internal review only.",
    )
    imagegen_parser.add_argument("--selected", action="store_true")
    imagegen_parser.set_defaults(func=command_import_imagegen)

    visual_quality_parser = subparsers.add_parser(
        "visual-quality-gate",
        help="Audit visual assets for file quality, traceability, and client-visible safety.",
    )
    visual_quality_parser.add_argument("project", help="Project directory.")
    visual_quality_parser.add_argument("--min-long-edge", type=int, default=720)
    visual_quality_parser.add_argument("--min-short-edge", type=int, default=480)
    visual_quality_parser.set_defaults(func=command_visual_quality_gate)

    film_quality_parser = subparsers.add_parser(
        "film-quality-gate",
        help="Audit cinematic/commercial quality before client-review packaging.",
    )
    film_quality_parser.add_argument("project", help="Project directory.")
    film_quality_parser.set_defaults(func=command_film_quality_gate)

    export_pptx_parser = subparsers.add_parser("export-pptx", help="Create an editable internal PPTX draft and check it.")
    export_pptx_parser.add_argument("project", help="Project directory.")
    export_pptx_parser.add_argument("--output", default="", help="Optional PPTX output path.")
    export_pptx_parser.set_defaults(func=command_export_pptx)

    check_pptx_parser = subparsers.add_parser("check-pptx", help="Check an actual PPTX for editable text layers.")
    check_pptx_parser.add_argument("project", help="Project directory.")
    check_pptx_parser.add_argument("--file", required=True, help="PPTX file to inspect.")
    check_pptx_parser.set_defaults(func=command_check_pptx)

    client_gate_parser = subparsers.add_parser("client-pack-gate", help="Audit client-review candidates before any send.")
    client_gate_parser.add_argument("project", help="Project directory.")
    client_gate_parser.add_argument("--pptx", default="", help="Optional PPTX file to inspect.")
    client_gate_parser.set_defaults(func=command_client_pack_gate)

    client_send_parser = subparsers.add_parser(
        "client-send-readiness-gate",
        help="Fail closed unless the exact current package has independent review and explicit send authorization; never sends.",
    )
    client_send_parser.add_argument("project", help="Project directory.")
    client_send_parser.set_defaults(func=command_client_send_readiness_gate)

    handoff_gate_parser = subparsers.add_parser(
        "handoff-readiness-gate",
        help="Audit whether the project can be handed to a non-developer operator.",
    )
    handoff_gate_parser.add_argument("project", help="Project directory.")
    handoff_gate_parser.set_defaults(func=command_handoff_readiness_gate)

    install_skill_parser = subparsers.add_parser("install-skill", help="Install the project skill into ~/.codex/skills.")
    install_skill_parser.add_argument("--target", default="", help="Optional skill install directory.")
    install_skill_parser.set_defaults(func=command_install_skill)

    return parser


def main() -> int:
    parser = build_parser()
    args = parser.parse_args()
    content_commands = {
        "init",
        "run",
        "intake",
        "intake-evidence",
        "profile-analyze",
        "export-intake-analysis-request",
        "import-intake-analysis",
        "creative-brief",
        "creative-proposal",
        "creative-import",
        "creative-review",
        "sample",
        "demo",
        "quickstart",
        "status",
        "next",
        "validate",
        "hygiene",
        "agency-audit",
        "support-bundle",
        "render-dashboard",
        "open-dashboard",
        "audit-dashboard",
    }
    project_arg = getattr(args, "project", None)
    delivery_command = bool(
        project_arg
        and args.command not in content_commands
        and not getattr(args, "dry_run", False)
    )
    token = DELIVERY_COMMAND_ACTIVE.set(delivery_command)
    try:
        return args.func(args)
    finally:
        DELIVERY_COMMAND_ACTIVE.reset(token)


if __name__ == "__main__":
    raise SystemExit(main())
