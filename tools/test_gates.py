#!/usr/bin/env python3
"""Structured regression checks for Gate review behavior."""

from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import types
from pathlib import Path
from unittest.mock import patch

from ad_creative_operator import (
    adversarial_council_evidence,
    add_visual_asset,
    confirm_client_outline,
    current_client_pack_binding_errors,
    build_client_pack_input_manifest,
    client_language_text_for_path,
    default_adversarial_targets,
    ensure_project,
    export_editable_pptx,
    file_sha256,
    has_gate,
    inspect_pptx,
    now_iso,
    read_csv_rows,
    render_goal_iteration_plan,
    render_creative_proposal,
    register_materials,
    review_client_pack,
    review_client_outline,
    review_client_send_readiness,
    review_creative_quality,
    review_handoff_readiness,
    review_reference_pack,
    review_search_quality,
    review_visual_quality,
    write_csv_rows,
    write_text,
    update_markdown_sections,
    write_manual_review_checklist,
    write_json_object,
    write_adversarial_target_snapshot,
)
from validate_project import current_truth_value, validate


OPTIONAL_SKIPS: list[str] = []


def assert_valid(project: Path) -> None:
    errors, _ = validate(project)
    if errors:
        raise AssertionError("\n".join(errors))


def add_row(project: Path, rel_path: str, row: dict[str, str]) -> None:
    path = project / rel_path
    fields, rows = read_csv_rows(path)
    if not fields:
        raise AssertionError(f"missing CSV header: {path}")
    rows.append(row)
    write_csv_rows(path, fields, rows)


def add_adversarial_record(project: Path, stage: str = "global") -> None:
    stages = (
        ["creative", "reference_research", "visual_review", "film_quality"]
        if stage == "global"
        else [stage]
    )
    for current_stage in stages:
        if current_stage == "final_delivery":
            _, artifacts = read_csv_rows(
                project / "AD-creative/orchestrator/artifact_index.csv"
            )
            payload, digest, _ = build_client_pack_input_manifest(project, artifacts)
            target = write_adversarial_target_snapshot(
                project,
                stage=current_stage,
                payload=payload,
                target_digest=digest,
            )
        else:
            targets = default_adversarial_targets(project, current_stage)
            if not targets:
                continue
            target = targets[0]
        write_text(
            project
            / f"AD-creative/gates/ADVERSARIAL_REVIEW_{current_stage.upper()}.md",
            f"""# Independent Adversarial Review

stage: {current_stage}
reviewer_id: fixture-independent-reviewer
reviewer_role: cold reviewer
independent: true
reviewed_at: {now_iso()}
target_ref: {target.relative_to(project)}
target_sha256: {file_sha256(target)}

| stage | objection | rebuttal_path | revision_decision | gate_status |
|---|---|---|---|---|
| {current_stage} | Challenge the happy path | review evidence | retain only verified claims | PASS |
""",
        )


def add_client_outline(project: Path, *, visibility: str = "internal_only") -> None:
    add_row(
        project,
        "AD-creative/client_review/client_outline.csv",
        {
            "slide_id": "1",
            "page_title": "客户可读开篇",
            "body_copy": "这一页用完整客户可读段落说明传播问题和方案进入点。",
            "client_confirmation_point": "确认是否作为客户审阅开篇。",
            "material_role": "暂用可编辑文本结构，图片后续登记。",
            "visual_slot": "横屏低密度画面占位。",
            "visual_asset_status": "placeholder",
            "asset_ids": "",
            "visibility": visibility,
            "status": "ready",
            "notes": "",
        },
    )
    if visibility == "client_visible_ready":
        confirm_client_outline(
            project,
            confirmed_by="fixture-project-owner",
            confirmed_at="2026-07-05T00:00:00Z",
            evidence_ref="user_confirmation:test-fixture",
        )


def write_safe_client_review_files(project: Path) -> None:
    write_text(
        project / "AD-creative/client_review/client_review_outline.md",
        """# Client Review Outline

status: ready
visibility: client_visible_ready

Morning trail opening page with a clear product benefit and reviewable story flow.
""",
    )
    write_text(
        project / "AD-creative/client_review/slide_spec.md",
        """# Slide Story Notes

status: ready
visibility: client_visible_ready

| Slide | Purpose | Content | Asset Slot | Visibility |
|---|---|---|---|---|
| 1 | Opening story | Morning trail product benefit | none | client_visible_ready |
""",
    )


def optional_module(name: str) -> bool:
    try:
        __import__(name)
        return True
    except Exception as exc:  # noqa: BLE001 - tests report optional fixture coverage
        OPTIONAL_SKIPS.append(f"{name}: {exc}")
        return False


def add_requirement(project: Path, requirement_id: str = "REQ-001") -> None:
    add_row(
        project,
        "AD-creative/orchestrator/requirements.csv",
        {
            "requirement_id": requirement_id,
            "source_event_id": "",
            "owner": "operator",
            "statement": "Use internal visual QA fixture for selected asset validation.",
            "requirement_type": "visual",
            "priority": "high",
            "status": "extracted",
            "confidence": "0.90",
            "scope": "project",
            "affected_stage": "visual_review",
            "linked_artifacts": "",
            "supersedes_requirement_id": "",
            "open_questions": "",
        },
    )


def add_delivery_artifact(
    project: Path,
    artifact_id: str,
    artifact_type: str,
    rel_path: str,
    visibility: str = "internal_only",
    *,
    version: str = "v001",
    derived_from_artifact_id: str = "",
    derived_from_sha256: str = "",
) -> None:
    target = project / rel_path
    target.parent.mkdir(parents=True, exist_ok=True)
    if not target.exists():
        target.write_text("delivery fixture", encoding="utf-8")
    add_row(
        project,
        "AD-creative/orchestrator/artifact_index.csv",
        {
            "artifact_id": artifact_id,
            "artifact_type": artifact_type,
            "path": rel_path,
            "stage": "final_delivery",
            "version": version,
            "status": "done",
            "visibility": visibility,
            "source_event_ids": "",
            "linked_requirements": "",
            "linked_work_items": "",
            "linked_references": "",
            "linked_assets": "",
            "gate_status": "PASS",
            "supersedes_artifact_id": "",
            "created_at": "",
            "updated_at": "",
            "sha256": file_sha256(target),
            "size_bytes": str(target.stat().st_size),
            "derived_from_artifact_id": derived_from_artifact_id,
            "derived_from_sha256": derived_from_sha256,
        },
    )


def write_complete_creative_fixture(
    project: Path,
    *,
    unsupported_claim: bool = False,
    generic: bool = False,
    humanizer_risk: bool = False,
) -> None:
    claim = "案例证明这个打法有效。" if unsupported_claim else "Registered REF-001 is a placeholder reference boundary, not a case-study claim."
    slogan = "unlock next level innovative breakthrough" if generic else "continue the morning trail"
    humanizer_block = (
        """
## Client Copy Draft Risk Fixture
Of course! Here is a route experts say will mark a pivotal moment in a vibrant brand landscape.
It is not just about hydration, but about showcasing a crucial lifestyle shift -- additionally, it can enhance valuable engagement — and underscore momentum.
"""
        if humanizer_risk
        else ""
    )
    write_text(
        project / "AD-creative/creative/creative_directions.md",
        f"""# Creative Directions

status: draft
visibility: internal_only

## Proposal Inputs
- business problem: New product needs a clear internal proposal route. [source: SRC-001]
- client real objective: Choose one reviewable advertising direction. [source: SRC-001]
- target audience: Morning trail runners and light outdoor users. [source: SRC-001]
- behavior barrier: They distrust abstract energy claims unless the usage moment is concrete. [source: SRC-001]
- consumer insight: A replenishment claim becomes credible when it is shown at the exact moment the user keeps moving. [source: SRC-001]
- feature to benefit: Electrolyte hydration feature -> visible confidence to continue the activity. [source: SRC-001]
- brand/category/competitor notes: Category codes are functional refreshment and outdoor credibility; competitor naming is intentionally excluded from this internal fixture.
- strategy path: source evidence -> audience barrier -> product benefit -> execution route -> client choice rationale.

## Direction Overview
| Direction | Name | Role | Strategy Path | Core Message | Why Choose |
|---|---|---|---|---|---|
| DIR-01 | Trail Continuation | Usage proof | product_feature_to_behavior_moment | {slogan} | Most directly links product and behavior. |
| DIR-02 | Choice Rationale | Internal decision | client_objective_to_choice_rationale | Compare routes by evidence, action, and risk. | Best for client alignment. |

## DIR-01 Trail Continuation
- creative proposition: Show the product in the replenishment moment before the user keeps moving.
- core message: Hydration is not a claim; it is the reason the next step happens.
- key visual/action: Runner pauses on a morning trail, drinks, and continues uphill.
- title/use case: Morning trail continuation.
- risk: Requires real product image before client review.
- why choose: Most direct product-to-benefit translation.

## DIR-02 Choice Rationale
- creative proposition: Turn the client review into a structured choice between two routes.
- core message: Pick the direction whose evidence and risk fit the decision.
- key visual/action: Editable comparison page with action frame slots.
- title/use case: Internal review decision slide.
- risk: Needs confirmation of final decision criteria.
- why choose: Useful when stakeholders need alignment.

## Reference Boundary
{claim}
{humanizer_block}

## Confirmation Items
- Confirm final product asset before client review.
""",
    )
    write_csv_rows(
        project / "AD-creative/creative/option_matrix.csv",
        [
            "direction_id",
            "name",
            "role",
            "strategy_path",
            "creative_proposition",
            "core_message",
            "target_feeling",
            "product_feature",
            "communication_benefit",
            "behavior_barrier",
            "key_visual_or_action",
            "title_or_use_case",
            "reference_ids",
            "risk",
            "why_choose",
            "evidence_refs",
            "status",
            "notes",
        ],
        [
            {
                "direction_id": "DIR-01",
                "name": "Trail Continuation",
                "role": "Usage proof",
                "strategy_path": "product_feature_to_behavior_moment",
                "creative_proposition": "Show the product in the replenishment moment before the user keeps moving.",
                "core_message": "Hydration becomes credible when the next step happens.",
                "target_feeling": "credible",
                "product_feature": "Electrolyte hydration",
                "communication_benefit": "Confidence to continue",
                "behavior_barrier": "Distrust of abstract claims",
                "key_visual_or_action": "Runner drinks and continues uphill.",
                "title_or_use_case": "Morning trail continuation",
                "reference_ids": "REF-001",
                "risk": "Needs real product image.",
                "why_choose": "Direct product-to-benefit translation.",
                "evidence_refs": "SRC-001",
                "status": "draft",
                "notes": "",
            },
            {
                "direction_id": "DIR-02",
                "name": "Choice Rationale",
                "role": "Internal decision",
                "strategy_path": "client_objective_to_choice_rationale",
                "creative_proposition": "Make route selection explicit through evidence, action, and risk.",
                "core_message": "Choose by fit, not slogan preference.",
                "target_feeling": "clear",
                "product_feature": "Electrolyte hydration",
                "communication_benefit": "Reviewable decision confidence",
                "behavior_barrier": "Stakeholder alignment",
                "key_visual_or_action": "Editable comparison matrix.",
                "title_or_use_case": "Review decision slide",
                "reference_ids": "REF-001",
                "risk": "Needs decision criteria.",
                "why_choose": "Best for client alignment.",
                "evidence_refs": "SRC-001",
                "status": "draft",
                "notes": "",
            },
        ],
    )
    write_text(
        project / "AD-creative/proposal_architecture/proposal_structure.md",
        """# Proposal Structure

status: draft
visibility: internal_only

## Client Review Goal
Choose a reviewable advertising direction.

## Recommended Page Flow
Business problem, client real objective, target audience, behavior barrier, consumer insight, feature to benefit, strategy path, directions, proposal outline.

## Proposal Outline
PPT/proposal outline includes problem, audience, product benefit, direction comparison, key visual/action, risk, and why choose.
""",
    )
    write_text(
        project / "AD-creative/client_review/slide_spec.md",
        """# Slide Spec

status: draft
visibility: internal_only

## Slides
| Slide | Purpose | Content | Asset Slot | Visibility |
|---|---|---|---|---|
| 1 | Proposal outline | Business problem and client real objective | none | internal_only |
| 2 | Direction | Creative proposition, core message, key visual/action, title/use case, risk, why choose | DIR-01 | internal_only |
""",
    )


def write_current_delivery_truth(project: Path, version_id: str = "VER-CURRENT") -> None:
    write_text(
        project / "AD-creative/orchestrator/current_truth.md",
        f"""# Current Truth

## Current Version Truth

```text
current_version_id: {version_id}
current_pptx_artifact_id: ART-AUTO-PPTX
current_pdf_artifact_id: ART-CURRENT-PDF
current_preview_artifact_id: ART-CURRENT-PREVIEW
current_text_extract_artifact_id: ART-CURRENT-TEXT
current_ppt_editability_artifact_id: ART-AUTO-PPT-EDITABILITY
version_map_status: active
last_archive_before_edit: n/a
```
""",
    )


def add_current_delivery_package(project: Path, version_id: str = "") -> None:
    truth_path = project / "AD-creative/orchestrator/current_truth.md"
    truth = truth_path.read_text(encoding="utf-8")
    version_id = version_id or current_truth_value(truth, "current_version_id")
    pptx_artifact_id = current_truth_value(truth, "current_pptx_artifact_id")
    editability_artifact_id = current_truth_value(
        truth, "current_ppt_editability_artifact_id"
    )
    artifact_path = project / "AD-creative/orchestrator/artifact_index.csv"
    artifact_fields, artifact_rows = read_csv_rows(artifact_path)
    pptx_row = next(row for row in artifact_rows if row.get("artifact_id") == pptx_artifact_id)
    pptx_sha = pptx_row["sha256"]
    version = pptx_row["version"]
    pptx_row["visibility"] = "client_visible_ready"

    from reportlab.pdfgen import canvas

    pdf_path = project / f"AD-creative/ppt/exports/client_review_{version}.pdf"
    pdf_path.parent.mkdir(parents=True, exist_ok=True)
    pdf = canvas.Canvas(str(pdf_path))
    pdf.drawString(72, 760, "Morning trail product benefit and client review story")
    pdf.save()
    preview_path = project / f"AD-creative/ppt/exports/client_review_{version}.png"
    create_png(preview_path)
    text_path = project / f"AD-creative/ppt/exports/client_review_{version}.txt"
    text_path.write_text(
        "Morning trail product benefit and client review story with a clear decision path.",
        encoding="utf-8",
    )
    write_csv_rows(artifact_path, artifact_fields, artifact_rows)
    add_delivery_artifact(
        project,
        f"ART-PDF-{version[1:]}",
        "pdf",
        str(pdf_path.relative_to(project)),
        "client_visible_ready",
        version=version,
        derived_from_artifact_id=pptx_artifact_id,
        derived_from_sha256=pptx_sha,
    )
    add_delivery_artifact(
        project,
        f"ART-PREVIEW-{version[1:]}",
        "preview",
        str(preview_path.relative_to(project)),
        "client_visible_ready",
        version=version,
        derived_from_artifact_id=pptx_artifact_id,
        derived_from_sha256=pptx_sha,
    )
    add_delivery_artifact(
        project,
        f"ART-TEXT-{version[1:]}",
        "text_extract",
        str(text_path.relative_to(project)),
        "client_visible_ready",
        version=version,
        derived_from_artifact_id=pptx_artifact_id,
        derived_from_sha256=pptx_sha,
    )
    version_path = project / "AD-creative/orchestrator/version_map.csv"
    version_fields, version_rows = read_csv_rows(version_path)
    version_row = next(row for row in version_rows if row.get("version_id") == version_id)
    version_row["status"] = "current"
    write_csv_rows(version_path, version_fields, version_rows)
    update_markdown_sections(
        truth_path,
        {
            "Current Version Truth": f"""```text
current_version_id: {version_id}
current_pptx_artifact_id: {pptx_artifact_id}
current_pdf_artifact_id: ART-PDF-{version[1:]}
current_preview_artifact_id: ART-PREVIEW-{version[1:]}
current_text_extract_artifact_id: ART-TEXT-{version[1:]}
current_ppt_editability_artifact_id: {editability_artifact_id}
version_map_status: current
last_archive_before_edit: fixture
```""",
        },
    )


def create_png(path: Path, size: tuple[int, int] = (960, 640)) -> None:
    from PIL import Image

    path.parent.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", size, (42, 96, 137))
    image.save(path)


def test_reference_pack_blocks_client_visible_bad_reference() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-gate-ref-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project)
        add_row(
            project,
            "AD-creative/references/reference_cards.csv",
            {
                "reference_id": "REF-001",
                "source_event_id": "",
                "platform": "unknown",
                "url": "http://example.com",
                "title": "Bad reference",
                "source_owner": "unknown",
                "reference_type": "reference",
                "role": "direction_reference",
                "why_relevant": "",
                "borrow": "",
                "do_not_copy": "",
                "client_visible": "true",
                "notes": "",
            },
        )
        status, findings, _ = review_reference_pack(project)
        assert status == "BLOCKED", (status, findings)
        assert any("https" in item or "do_not_copy" in item for item in findings), findings


def test_search_quality_passes_clean_plan_with_adversarial_record() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-gate-search-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project)
        write_text(
            project / "AD-creative/references/official_search_plan.md",
            """# Official Search Plan

## Why Search

Verify public, traceable direction references before any client-visible citation.

## Search Scope

Official brand sites, official campaign archives, public video platforms.

## Suggested Platforms

Official website, YouTube, Vimeo.

## Expected Output

Reference cards with source owner, relevance, borrow, and do_not_copy fields.

## do_not_copy

Do not copy logos, layouts, people, product markings, or protected composition.
""",
        )
        status, findings, _ = review_search_quality(project)
        assert status == "PASS", (status, findings)
        assert findings == [], findings
        assert_valid(project)


def test_visual_quality_blocks_missing_selected_asset_file() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-gate-visual-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project)
        add_row(
            project,
            "AD-creative/visual_assets/asset_manifest.csv",
            {
                "asset_id": "IMG-001",
                "slot_id": "SLOT-001",
                "requirement_id": "REQ-001",
                "reference_id": "",
                "path": "AD-creative/visual_assets/selected/missing.png",
                "asset_type": "uploaded_image",
                "stage": "visual_review",
                "version": "v001",
                "status": "selected",
                "visibility": "internal_only",
                "qa_status": "PASS",
                "risk_level": "low",
                "prompt_or_edit_ref": "",
                "notes": "",
            },
        )
        status, findings, _ = review_visual_quality(project)
        assert status == "BLOCKED", (status, findings)
        assert any("文件不存在" in item for item in findings), findings


def test_visual_quality_passes_real_internal_selected_asset() -> None:
    if not optional_module("PIL"):
        return
    with tempfile.TemporaryDirectory(prefix="adco-gate-visual-pass-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project)
        add_requirement(project)
        source = project / "fixture/selected.png"
        create_png(source)
        add_visual_asset(
            project,
            source,
            "SLOT-001",
            "REQ-001",
            "",
            "uploaded_image",
            "internal_only",
            "PASS",
            "low",
            "",
            "fixture image for positive visual gate regression",
            selected=True,
        )
        add_adversarial_record(project, "visual_review")
        status, findings, _ = review_visual_quality(project)
        assert status == "PASS", (status, findings)
        assert findings == [], findings
        assert_valid(project)


def test_visual_quality_rejects_asset_self_stamp_without_hash_bound_authorization() -> None:
    if not optional_module("PIL"):
        return
    with tempfile.TemporaryDirectory(prefix="adco-gate-visual-auth-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project)
        add_requirement(project)
        source = project / "fixture/client-visible.png"
        create_png(source)
        add_visual_asset(
            project,
            source,
            "SLOT-CLIENT-001",
            "REQ-001",
            "",
            "uploaded_image",
            "client_visible_ready",
            "PASS",
            "low",
            "",
            "client_visibility_approved",
            selected=True,
        )
        status, findings, _ = review_visual_quality(project)
        assert status == "BLOCKED", (status, findings)
        assert any("独立授权 receipt" in item for item in findings), findings
        _, current_assets = read_csv_rows(
            project / "AD-creative/visual_assets/asset_current_manifest.csv"
        )
        current = next(row for row in current_assets if row.get("asset_id"))
        add_row(
            project,
            "AD-creative/visual_assets/asset_authorizations.csv",
            {
                "authorization_id": "AUTH-FAKE-001",
                "asset_id": current["asset_id"],
                "asset_sha256": current["sha256"],
                "approval_scope": "client_review",
                "approved_by": "made_up_person",
                "approved_at": "never-validated",
                "evidence_ref": "missing/path.md",
                "evidence_sha256": "0" * 64,
                "status": "approved",
                "revoked_at": "",
                "notes": "must not count as authorization",
            },
        )
        status, findings, _ = review_visual_quality(project)
        assert status == "BLOCKED", (status, findings)
        assert any("独立授权 receipt" in item for item in findings), findings


def test_client_pack_blocks_without_editable_pptx() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-gate-client-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project)
        status, findings, _ = review_client_pack(project)
        assert status == "BLOCKED", (status, findings)
        assert any("PPTX" in item for item in findings), findings
        assert_valid(project)


def test_client_pack_passes_editable_internal_pptx() -> None:
    if not optional_module("pptx"):
        return
    with tempfile.TemporaryDirectory(prefix="adco-gate-client-pass-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project)
        add_client_outline(project, visibility="client_visible_ready")
        write_safe_client_review_files(project)
        pptx_path = export_editable_pptx(project)
        add_current_delivery_package(project)
        add_adversarial_record(project, "final_delivery")
        status, findings, _ = review_client_pack(project, pptx_path)
        assert status == "PASS", (status, findings)
        assert findings == [], findings
        assert_valid(project)


def test_client_pack_blocks_fake_preview_even_with_updated_hash() -> None:
    if not optional_module("pptx") or not optional_module("PIL") or not optional_module("reportlab"):
        return
    with tempfile.TemporaryDirectory(prefix="adco-client-fake-preview-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project)
        add_client_outline(project, visibility="client_visible_ready")
        write_safe_client_review_files(project)
        export_editable_pptx(project)
        add_current_delivery_package(project)
        truth = (project / "AD-creative/orchestrator/current_truth.md").read_text(encoding="utf-8")
        preview_id = current_truth_value(truth, "current_preview_artifact_id")
        artifact_path = project / "AD-creative/orchestrator/artifact_index.csv"
        fields, rows = read_csv_rows(artifact_path)
        preview_row = next(row for row in rows if row.get("artifact_id") == preview_id)
        preview = project / preview_row["path"]
        preview.write_bytes(b"delivery fixture")
        preview_row["sha256"] = file_sha256(preview)
        preview_row["size_bytes"] = str(preview.stat().st_size)
        write_csv_rows(artifact_path, fields, rows)
        status, findings, _ = review_client_pack(project)
        assert status == "BLOCKED", (status, findings)
        assert any("preview is not a valid PNG/JPEG" in item for item in findings), findings


def test_client_pack_scans_exact_current_text_extract() -> None:
    if not optional_module("pptx") or not optional_module("PIL") or not optional_module("reportlab"):
        return
    with tempfile.TemporaryDirectory(prefix="adco-client-language-package-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project)
        add_client_outline(project, visibility="client_visible_ready")
        write_safe_client_review_files(project)
        export_editable_pptx(project)
        add_current_delivery_package(project)
        truth = (project / "AD-creative/orchestrator/current_truth.md").read_text(encoding="utf-8")
        text_id = current_truth_value(truth, "current_text_extract_artifact_id")
        artifact_path = project / "AD-creative/orchestrator/artifact_index.csv"
        fields, rows = read_csv_rows(artifact_path)
        text_row = next(row for row in rows if row.get("artifact_id") == text_id)
        text_path = project / text_row["path"]
        write_text(text_path, "Client page with internal Gate worker notes.")
        text_row["sha256"] = file_sha256(text_path)
        text_row["size_bytes"] = str(text_path.stat().st_size)
        write_csv_rows(artifact_path, fields, rows)
        status, findings, _ = review_client_pack(project)
        assert status == "BLOCKED", (status, findings)
        assert any("Client Language Gate" in item and "gate" in item.lower() for item in findings), findings


def test_internal_outline_cannot_satisfy_client_pack_readiness() -> None:
    if not optional_module("pptx") or not optional_module("PIL") or not optional_module("reportlab"):
        return
    with tempfile.TemporaryDirectory(prefix="adco-client-internal-outline-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project)
        add_client_outline(project)
        write_safe_client_review_files(project)
        try:
            export_editable_pptx(project)
        except RuntimeError as exc:
            assert "client-outline-gate BLOCKED" in str(exc)
        else:
            raise AssertionError("unconfirmed/internal outline must not reach PPT export")


def test_pptx_editability_rejects_flattened_slide() -> None:
    if not optional_module("pptx") or not optional_module("PIL"):
        return
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    with tempfile.TemporaryDirectory(prefix="adco-flat-slide-") as raw_project:
        project = Path(raw_project)
        image_path = project / "flat.png"
        Image.new("RGB", (1600, 900), color=(20, 40, 60)).save(image_path)
        deck = project / "flat.pptx"
        presentation = Presentation()
        first = presentation.slides.add_slide(presentation.slide_layouts[6])
        first.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text = "Editable title"
        second = presentation.slides.add_slide(presentation.slide_layouts[6])
        second.shapes.add_picture(str(image_path), Inches(0), Inches(0), width=Inches(10))
        presentation.save(deck)
        stats = inspect_pptx(deck)
        assert stats["editable"] is False
        assert stats["flattened_slides"] == "2"


def test_client_send_readiness_requires_hash_bound_human_and_send_authorization() -> None:
    if not optional_module("pptx") or not optional_module("PIL") or not optional_module("reportlab"):
        return
    with tempfile.TemporaryDirectory(prefix="adco-send-readiness-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project)
        add_client_outline(project, visibility="client_visible_ready")
        write_safe_client_review_files(project)
        export_editable_pptx(project)
        add_current_delivery_package(project)
        add_adversarial_record(project, "final_delivery")
        pack_status, findings, _ = review_client_pack(project)
        assert pack_status == "PASS", findings
        status, issues, _ = review_client_send_readiness(project)
        assert status == "BLOCKED"
        assert any("人工审阅" in item for item in issues)
        assert any("发送授权" in item for item in issues)

        truth = (project / "AD-creative/orchestrator/current_truth.md").read_text(encoding="utf-8")
        version_id = current_truth_value(truth, "current_version_id")
        pptx_id = current_truth_value(truth, "current_pptx_artifact_id")
        _, artifacts = read_csv_rows(
            project / "AD-creative/orchestrator/artifact_index.csv"
        )
        pptx_row = next(row for row in artifacts if row.get("artifact_id") == pptx_id)
        binding = json.loads(
            (project / "AD-creative/delivery/client_pack_binding.json").read_text(
                encoding="utf-8"
            )
        )
        common = {
            "version_id": version_id,
            "pptx_artifact_id": pptx_id,
            "pptx_sha256": pptx_row["sha256"],
            "package_digest": binding["package_digest"],
        }
        write_json_object(
            project / "AD-creative/delivery/manual_review_receipt.json",
            {
                **common,
                "review_id": "REVIEW-001",
                "reviewer_id": "independent-human-reviewer",
                "reviewer_role": "creative cold reviewer",
                "independent": True,
                "reviewed_at": "2026-07-05T01:00:00Z",
                "evidence_ref": "review_record:test-fixture",
                "decision": "approved",
                "checks": {
                    "client_language": True,
                    "visual_layout": True,
                    "asset_authorization": True,
                    "ppt_editability": True,
                },
            },
        )
        write_json_object(
            project / "AD-creative/delivery/send_authorization.json",
            {
                **common,
                "authorization_id": "SEND-001",
                "authorized_by": "project-owner",
                "authorized_at": "2026-07-05T01:10:00Z",
                "evidence_ref": "user_confirmation:send-fixture",
                "recipient_scope": "approved client review group",
                "decision": "authorized",
            },
        )
        status, issues, _ = review_client_send_readiness(project)
        assert status == "PASS", issues

        authorization_path = project / "AD-creative/delivery/send_authorization.json"
        payload = json.loads(authorization_path.read_text(encoding="utf-8"))
        payload["pptx_sha256"] = "0" * 64
        write_json_object(authorization_path, payload)
        status, issues, _ = review_client_send_readiness(project)
        assert status == "BLOCKED"
        assert any("exact current PPTX hash" in item for item in issues)

        payload["pptx_sha256"] = common["pptx_sha256"]
        write_json_object(authorization_path, payload)
        truth = (project / "AD-creative/orchestrator/current_truth.md").read_text(
            encoding="utf-8"
        )
        text_id = current_truth_value(truth, "current_text_extract_artifact_id")
        artifact_path = project / "AD-creative/orchestrator/artifact_index.csv"
        fields, artifacts = read_csv_rows(artifact_path)
        text_row = next(row for row in artifacts if row.get("artifact_id") == text_id)
        text_path = project / text_row["path"]
        write_text(text_path, "Internal worker gate prompt TODO")
        text_row["sha256"] = file_sha256(text_path)
        text_row["size_bytes"] = str(text_path.stat().st_size)
        write_csv_rows(artifact_path, fields, artifacts)
        status, issues, _ = review_client_send_readiness(project)
        assert status == "BLOCKED"
        assert any("binding 已过期" in item for item in issues), issues


def test_client_pack_manifest_tamper_cannot_be_rebound_by_rerunning_gate() -> None:
    if not optional_module("pptx") or not optional_module("PIL") or not optional_module("reportlab"):
        return
    with tempfile.TemporaryDirectory(prefix="adco-client-pack-immutable-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project)
        add_client_outline(project, visibility="client_visible_ready")
        write_safe_client_review_files(project)
        export_editable_pptx(project)
        add_current_delivery_package(project)
        add_adversarial_record(project, "final_delivery")
        status, findings, _ = review_client_pack(project)
        assert status == "PASS", findings

        binding_path = project / "AD-creative/delivery/client_pack_binding.json"
        binding_before = json.loads(binding_path.read_text(encoding="utf-8"))
        manifest_path = project / str(binding_before["manifest_path"])
        mutate = json.loads(manifest_path.read_text(encoding="utf-8"))
        mutate["files"] = []
        write_json_object(manifest_path, mutate)
        _, artifacts = read_csv_rows(
            project / "AD-creative/orchestrator/artifact_index.csv"
        )
        binding_errors, _ = current_client_pack_binding_errors(project, artifacts)
        assert any(
            "immutable manifest content" in issue for issue in binding_errors
        ), binding_errors
        try:
            review_client_pack(project)
        except ValueError as exc:
            assert "immutable client pack manifest collision" in str(exc), exc
        else:
            raise AssertionError("a tampered immutable manifest must not be rebound")
        binding_after = json.loads(binding_path.read_text(encoding="utf-8"))
        assert binding_after == binding_before


def test_validate_blocks_mismatched_current_delivery_version() -> None:
    if (
        not optional_module("pptx")
        or not optional_module("PIL")
        or not optional_module("reportlab")
    ):
        return
    with tempfile.TemporaryDirectory(prefix="adco-gate-current-version-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_client_outline(project, visibility="client_visible_ready")
        export_editable_pptx(project)
        add_current_delivery_package(project)
        truth_path = project / "AD-creative/orchestrator/current_truth.md"
        truth = truth_path.read_text(encoding="utf-8").replace(
            "current_version_id: VER-PPT-001",
            "current_version_id: VER-UNKNOWN",
        )
        write_text(truth_path, truth)
        errors, _ = validate(project)
        assert any("current_version_id" in error for error in errors), errors


def test_validate_rejects_ambiguous_current_version_truth() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-current-truth-ambiguity-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        truth_path = project / "AD-creative/orchestrator/current_truth.md"
        original = truth_path.read_text(encoding="utf-8")
        truth_path.write_text(
            original.replace(
                "current_version_id:\n",
                "current_version_id:\ncurrent_version_id: VER-FORGED\n",
                1,
            ),
            encoding="utf-8",
        )
        errors, _ = validate(project)
        assert any(
            "current_version_id appears 2 times" in error for error in errors
        ), errors

        truth_path.write_text(
            original
            + "\n## Current Version Truth\n\n"
            + "current_version_id: VER-FORGED\n",
            encoding="utf-8",
        )
        errors, _ = validate(project)
        assert any("exactly one '## Current Version Truth'" in error for error in errors), errors


def test_pdf_text_extraction_falls_back_after_pdftotext_failure() -> None:
    class FakeDocument:
        def __enter__(self) -> list[object]:
            return [types.SimpleNamespace(get_text=lambda: "fallback client copy")]

        def __exit__(self, *_args: object) -> None:
            return None

    fake_fitz = types.SimpleNamespace(open=lambda _path: FakeDocument())
    with tempfile.TemporaryDirectory(prefix="adco-pdf-fallback-") as raw_project:
        path = Path(raw_project) / "fixture.pdf"
        path.write_bytes(b"%PDF-1.4\n")
        failed = subprocess.CompletedProcess(
            ["pdftotext", str(path), "-"], 1, stdout="", stderr="parse failed"
        )
        with (
            patch("ad_creative_operator.shutil.which", return_value="/fake/pdftotext"),
            patch("ad_creative_operator.subprocess.run", return_value=failed),
            patch.dict(sys.modules, {"fitz": fake_fitz}),
        ):
            assert client_language_text_for_path(path) == "fallback client copy"


def test_validate_rejects_artifact_path_outside_project_even_with_matching_hash() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-artifact-path-scope-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        outside = Path("/etc/hosts")
        add_row(
            project,
            "AD-creative/orchestrator/artifact_index.csv",
            {
                "artifact_id": "ART-OUTSIDE-001",
                "artifact_type": "text_extract",
                "path": str(outside),
                "stage": "scope_test",
                "version": "v001",
                "status": "done",
                "visibility": "internal_only",
                "gate_status": "PASS",
                "sha256": file_sha256(outside),
                "size_bytes": str(outside.stat().st_size),
            },
        )
        errors, _ = validate(project)
        assert any("project-relative path" in error for error in errors), errors


def test_client_pack_blocks_missing_exact_current_pptx_without_crashing() -> None:
    if (
        not optional_module("pptx")
        or not optional_module("PIL")
        or not optional_module("reportlab")
    ):
        return
    with tempfile.TemporaryDirectory(prefix="adco-gate-missing-current-pptx-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project)
        add_client_outline(project, visibility="client_visible_ready")
        write_safe_client_review_files(project)
        missing_pptx = export_editable_pptx(project)
        add_current_delivery_package(project)
        missing_pptx.unlink()
        status, findings, report = review_client_pack(project)
        assert status == "BLOCKED", (status, findings)
        assert report.exists(), report
        assert any("PPTX 文件不存在" in item for item in findings), findings


def test_handoff_readiness_is_internal_operation_gate_not_send_gate() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-gate-handoff-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project)
        status, blockers, warnings, _ = review_handoff_readiness(project)
        assert status == "PASS", (status, blockers, warnings)
        assert blockers == []
        assert any("PPTX" in item for item in warnings), warnings
        assert any("GATE-" in item for item in warnings), warnings
        assert_valid(project)


def test_goal_plan_cannot_self_stamp_adversarial_review() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-adversarial-self-stamp-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        render_goal_iteration_plan(
            project,
            goal_id="GOAL-SELF-STAMP",
            title="Self stamp must fail",
            objective="A plan is not an independent review.",
            owner="Main Controller",
        )
        found, evidence = adversarial_council_evidence(project, "final_delivery")
        assert found is False
        assert evidence == []


def test_blocked_gate_does_not_count_as_completed_goal_stage() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-blocked-gate-stage-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_row(
            project,
            "AD-creative/orchestrator/gate_log.csv",
            {
                "gate_id": "GATE-AUTO-VISUAL-QUALITY-001",
                "stage": "visual_review",
                "status": "BLOCKED",
                "score": "0",
                "checked_artifacts": "",
                "blocking_issues": "fixture blocker",
                "revision_items": "fix fixture",
                "questions": "",
                "next_state": "fix_visual_assets",
                "created_at": now_iso(),
                "owner": "fixture",
            },
        )
        assert has_gate(project, "GATE-AUTO-VISUAL-QUALITY-001") is False


def test_gate_history_is_append_only_and_latest_target_must_be_fresh() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-gate-history-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_client_outline(project, visibility="client_visible_ready")
        status, findings, report = review_client_outline(project)
        assert status == "PASS", findings
        _, first_rows = read_csv_rows(
            project / "AD-creative/orchestrator/gate_log.csv"
        )
        first = next(
            row
            for row in first_rows
            if row.get("gate_id") == "GATE-AUTO-CLIENT-OUTLINE-001"
        )
        first_snapshot = project / first["evidence_snapshot_ref"]
        first_snapshot_sha = file_sha256(first_snapshot)
        write_text(report, report.read_text(encoding="utf-8") + "\ntampered\n")
        assert not has_gate(
            project, "GATE-AUTO-CLIENT-OUTLINE-001", {"PASS"}
        )
        errors, _ = validate(project)
        assert errors == [], errors

        status, findings, _ = review_client_outline(project)
        assert status == "PASS", findings
        _, rows = read_csv_rows(project / "AD-creative/orchestrator/gate_log.csv")
        runs = [
            row
            for row in rows
            if row.get("gate_id") == "GATE-AUTO-CLIENT-OUTLINE-001"
        ]
        assert len(runs) == 2
        assert runs[1]["supersedes_gate_run_id"] == runs[0]["gate_run_id"]
        assert first_snapshot.is_file()
        assert file_sha256(first_snapshot) == first_snapshot_sha
        assert runs[1]["evidence_snapshot_ref"] != runs[0]["evidence_snapshot_ref"]
        assert has_gate(project, "GATE-AUTO-CLIENT-OUTLINE-001", {"PASS"})
        assert_valid(project)


def test_adversarial_review_rejects_global_irrelevant_or_nonpass_evidence() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-adversarial-forged-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        unrelated = project / "irrelevant.txt"
        write_text(unrelated, "not the reference gate target")
        write_text(
            project / "AD-creative/gates/ADVERSARIAL_REVIEW_FORGED.md",
            f"""# Forged Review

stage: global
reviewer_id: made_up_person
reviewer_role: reviewer
independent: true
reviewed_at: {now_iso()}
target_ref: {unrelated.relative_to(project)}
target_sha256: {file_sha256(unrelated)}

| stage | objection | rebuttal | revision | gate_status |
|---|---|---|---|---|
| global | fake objection | fake rebuttal | no revision | BLOCKED |
""",
        )
        found, evidence = adversarial_council_evidence(
            project, "reference_research"
        )
        assert found is False
        assert evidence == []


def test_manual_review_checklist_starts_pending_not_passed() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-manual-review-pending-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        checklist = write_manual_review_checklist(project)
        assert "status: pending_human_review" in checklist.read_text(encoding="utf-8")
        _, artifacts = read_csv_rows(
            project / "AD-creative/orchestrator/artifact_index.csv"
        )
        row = next(
            item
            for item in artifacts
            if item.get("artifact_id") == "ART-AUTO-MANUAL-REVIEW-CHECKLIST"
        )
        assert row["status"] == "pending_human_review"
        assert row["gate_status"] == "NOT_RUN"


def test_creative_proposal_writes_required_internal_fields() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-creative-proposal-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        payload = render_creative_proposal(project, work_id="WORK-001")
        text = (project / "AD-creative/creative/creative_directions.md").read_text(encoding="utf-8")
        for marker in [
            "business problem",
            "client real objective",
            "target audience",
            "behavior barrier",
            "consumer insight",
            "feature to benefit",
            "brand/category/competitor notes",
            "creative proposition",
            "key visual/action",
            "why choose",
        ]:
            assert marker in text, marker
        assert "ART-AUTO-CREATIVE-DIRECTIONS" in payload["artifact_ids"]
        assert_valid(project)


def test_creative_proposal_prefers_source_brief_over_asset_gaps() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-creative-source-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        brief = project / "briefs/nova-trail.md"
        write_text(
            brief,
            """# NOVA Trail Shell Launch Brief

## Target Audience
- 24-36 year-old city professionals who commute on rainy weekdays and hike short trails on weekends.
- They distrust inflated outdoor claims and want specific use moments.

## Known Product Facts
- 3-layer waterproof fabric.
- Packable hood.
- Two-way underarm vents.
- Reflective side tabs.

## Creative Problem
The draft client language is too generic and must become specific, grounded, and useful.
""",
        )
        register_materials(project, [brief], "Create a grounded internal launch proposal.")
        add_row(
            project,
            "AD-creative/orchestrator/gaps.csv",
            {
                "gap_id": "GAP-999",
                "linked_requirement_id": "",
                "description": "缺少品牌 logo / 字体 / 包装 / 视觉规范，不能进入客户可见稿。",
                "impact": "blocking",
                "owner": "client",
                "question_for_client": "请提供品牌 logo、字体、包装或产品露出规范。",
                "question_for_director": "",
                "question_for_user": "",
                "recommended_action": "向客户索取品牌资产包；没有资产前只做内部方向稿。",
                "status": "open",
            },
        )

        render_creative_proposal(project, work_id="WORK-001")

        text = (project / "AD-creative/creative/creative_directions.md").read_text(encoding="utf-8")
        assert "3-layer waterproof fabric" in text
        assert "3 层防水面料" in text
        assert "24-36 year-old city professionals" in text
        assert "24-36 岁城市职场人" in text
        assert "target audience: Target Audience" not in text
        assert "在24-36 year-old" not in text
        assert "缺少品牌 logo" not in text
        assert_valid(project)


def test_creative_quality_blocks_generic_proposal() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-creative-generic-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project, "creative")
        write_complete_creative_fixture(project, generic=True)
        status, findings, report = review_creative_quality(project)
        assert status == "BLOCKED", (status, findings)
        assert report.exists()
        assert any("GENERIC_AI_CLICHE" in item or "cliche" in item for item in findings), findings
        assert_valid(project)


def test_creative_quality_passes_complete_structured_fixture() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-creative-pass-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        write_complete_creative_fixture(project)
        add_adversarial_record(project, "creative")
        status, findings, _ = review_creative_quality(project)
        assert status == "PASS", (status, findings)
        assert findings == [], findings
        assert_valid(project)


def test_creative_quality_blocks_unsupported_case_claim() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-creative-case-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project, "creative")
        write_complete_creative_fixture(project, unsupported_claim=True)
        status, findings, _ = review_creative_quality(project)
        assert status == "BLOCKED", (status, findings)
        assert any("UNSUPPORTED_REFERENCE_OR_CASE_CLAIM" in item or "案例" in item for item in findings), findings
        assert_valid(project)


def test_creative_quality_blocks_humanizer_writing_risks() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-creative-humanizer-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        add_adversarial_record(project, "creative")
        write_complete_creative_fixture(project, humanizer_risk=True)
        status, findings, _ = review_creative_quality(project)
        assert status == "BLOCKED", (status, findings)
        expected_codes = {
            "CHATBOT_RESIDUE",
            "VAGUE_AUTHORITY_CLAIM",
            "EXAGGERATED_SIGNIFICANCE",
            "FORMULAIC_NOT_ONLY_BUT",
            "GENERIC_AI_VOCABULARY",
            "DASH_OVERUSE",
        }
        for code in expected_codes:
            assert any(code in item for item in findings), (code, findings)
        assert_valid(project)


def test_validation_pass_is_not_creative_quality_pass() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-creative-validation-") as raw_project:
        project = Path(raw_project)
        ensure_project(project)
        errors, _ = validate(project)
        assert errors == [], errors
        status, findings, _ = review_creative_quality(project)
        assert status == "BLOCKED", (status, findings)
        assert any("EMPTY_SKELETON" in item or "空骨架" in item for item in findings), findings
        assert_valid(project)


def main() -> int:
    test_reference_pack_blocks_client_visible_bad_reference()
    test_search_quality_passes_clean_plan_with_adversarial_record()
    test_visual_quality_blocks_missing_selected_asset_file()
    test_visual_quality_passes_real_internal_selected_asset()
    test_visual_quality_rejects_asset_self_stamp_without_hash_bound_authorization()
    test_client_pack_blocks_without_editable_pptx()
    test_client_pack_passes_editable_internal_pptx()
    test_client_pack_blocks_fake_preview_even_with_updated_hash()
    test_client_pack_scans_exact_current_text_extract()
    test_internal_outline_cannot_satisfy_client_pack_readiness()
    test_pptx_editability_rejects_flattened_slide()
    test_client_send_readiness_requires_hash_bound_human_and_send_authorization()
    test_client_pack_manifest_tamper_cannot_be_rebound_by_rerunning_gate()
    test_validate_blocks_mismatched_current_delivery_version()
    test_validate_rejects_ambiguous_current_version_truth()
    test_pdf_text_extraction_falls_back_after_pdftotext_failure()
    test_validate_rejects_artifact_path_outside_project_even_with_matching_hash()
    test_client_pack_blocks_missing_exact_current_pptx_without_crashing()
    test_handoff_readiness_is_internal_operation_gate_not_send_gate()
    test_goal_plan_cannot_self_stamp_adversarial_review()
    test_blocked_gate_does_not_count_as_completed_goal_stage()
    test_gate_history_is_append_only_and_latest_target_must_be_fresh()
    test_adversarial_review_rejects_global_irrelevant_or_nonpass_evidence()
    test_manual_review_checklist_starts_pending_not_passed()
    test_creative_proposal_writes_required_internal_fields()
    test_creative_quality_blocks_generic_proposal()
    test_creative_quality_passes_complete_structured_fixture()
    test_creative_quality_blocks_unsupported_case_claim()
    test_creative_quality_blocks_humanizer_writing_risks()
    test_validation_pass_is_not_creative_quality_pass()
    if OPTIONAL_SKIPS:
        print("TEST_GATES_OPTIONAL_SKIPS=" + "; ".join(OPTIONAL_SKIPS))
    print("TEST_GATES=PASS")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
