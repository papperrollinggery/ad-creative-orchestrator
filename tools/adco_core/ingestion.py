"""Deterministic, provenance-preserving material ingestion for ADCO projects."""

from __future__ import annotations

import csv
import hashlib
import json
import mimetypes
import os
import re
import secrets
import shutil
import stat
import subprocess
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Iterable

from .models import EvidenceChunk
from .safe_write import atomic_write_text


DEFAULT_CHUNK_CHARS = 3000
DEFAULT_OVERLAP_CHARS = 250
DEFAULT_TOTAL_CHAR_BUDGET = 2_000_000
EVIDENCE_REL = Path("AD-creative/orchestrator/evidence_chunks.jsonl")
LOCAL_STATE_REL = Path(".adco-local")
LOCAL_SOURCE_MAP_REL = LOCAL_STATE_REL / "source_paths.json"
LOCAL_SOURCE_PREFIX = "local-source://"
LOCAL_SOURCE_MAP_VERSION = 1
LOCAL_GITIGNORE_TEXT = "*\n!.gitignore\n"
MAX_LOCAL_STATE_FILE_BYTES = 8 * 1024 * 1024

PLAIN_TEXT_SUFFIXES = {".md", ".txt"}
STRUCTURED_TEXT_SUFFIXES = {".csv", ".json", ".yaml", ".yml"}
DOCUMENT_SUFFIXES = {".docx", ".pptx", ".pdf"}
SUBTITLE_SUFFIXES = {".srt", ".vtt"}
IMAGE_SUFFIXES = {".avif", ".bmp", ".gif", ".heic", ".heif", ".jpeg", ".jpg", ".png", ".tif", ".tiff", ".webp"}
VIDEO_SUFFIXES = {".avi", ".m4v", ".mkv", ".mov", ".mp4", ".mpeg", ".mpg", ".webm"}
SUPPORTED_SUFFIXES = (
    PLAIN_TEXT_SUFFIXES
    | STRUCTURED_TEXT_SUFFIXES
    | DOCUMENT_SUFFIXES
    | SUBTITLE_SUFFIXES
    | IMAGE_SUFFIXES
    | VIDEO_SUFFIXES
)

SAFE_EXIF_TAGS = {
    "ColorSpace",
    "ExifImageHeight",
    "ExifImageWidth",
    "Orientation",
    "ResolutionUnit",
    "XResolution",
    "YResolution",
}

SAFE_FFPROBE_STREAM_FIELDS = {
    "index",
    "codec_type",
    "codec_name",
    "width",
    "height",
    "r_frame_rate",
    "channels",
    "sample_rate",
}


@dataclass
class RawEvidence:
    text: str
    page: int | None = None
    slide: int | None = None
    start_line: int | None = None
    end_line: int | None = None
    field_path: str | None = None
    start_time: str | None = None
    end_time: str | None = None
    inspection_status: str = "text_extracted"
    metadata: dict[str, object] = field(default_factory=dict)


@dataclass
class IngestionReport:
    chunks: list[EvidenceChunk]
    files_processed: int
    characters_read: int
    over_budget: list[dict[str, object]]
    parser_errors: list[dict[str, str]]

    def as_dict(self) -> dict[str, object]:
        return {
            "evidence_path": EVIDENCE_REL.as_posix(),
            "files_processed": self.files_processed,
            "characters_read": self.characters_read,
            "evidence_chunks": len(self.chunks),
            "over_budget": self.over_budget,
            "parser_errors": self.parser_errors,
        }


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def sha256_text(text: str) -> str:
    return sha256_bytes(text.encode("utf-8"))


def sha256_file(path: Path, *, chunk_size: int = 1024 * 1024) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(chunk_size), b""):
            digest.update(chunk)
    return digest.hexdigest()


def media_type_for(path: Path) -> str:
    guessed, _ = mimetypes.guess_type(path.name)
    return guessed or {
        ".md": "text/markdown",
        ".yaml": "application/yaml",
        ".yml": "application/yaml",
        ".srt": "application/x-subrip",
        ".vtt": "text/vtt",
    }.get(path.suffix.lower(), "application/octet-stream")


def display_source_path(project: Path, path: Path) -> str:
    try:
        return path.resolve().relative_to(project.resolve()).as_posix()
    except ValueError:
        return f"external-source://{path.name}"


def _open_project_dir(project: Path) -> int:
    if project.is_symlink():
        raise ValueError(f"project root must not be a symlink: {project}")
    try:
        fd = os.open(
            project,
            os.O_RDONLY
            | getattr(os, "O_DIRECTORY", 0)
            | getattr(os, "O_NOFOLLOW", 0),
        )
    except OSError as exc:
        raise ValueError(f"cannot safely open project root: {project}: {exc}") from exc
    try:
        opened = os.fstat(fd)
        visible = os.stat(project, follow_symlinks=False)
        if (
            not stat.S_ISDIR(visible.st_mode)
            or (opened.st_dev, opened.st_ino) != (visible.st_dev, visible.st_ino)
        ):
            raise ValueError("project root changed while opening local source state")
    except BaseException:
        os.close(fd)
        raise
    return fd


def _verify_local_state_binding(project_fd: int, local_fd: int) -> None:
    try:
        visible = os.stat(
            LOCAL_STATE_REL.name,
            dir_fd=project_fd,
            follow_symlinks=False,
        )
    except FileNotFoundError as exc:
        raise ValueError("local source state changed during operation") from exc
    opened = os.fstat(local_fd)
    if (
        not stat.S_ISDIR(visible.st_mode)
        or (opened.st_dev, opened.st_ino) != (visible.st_dev, visible.st_ino)
    ):
        raise ValueError("local source state changed during operation")


def _open_local_state(
    project: Path,
    *,
    create: bool,
) -> tuple[int, int | None]:
    project_fd = _open_project_dir(project)
    local_fd: int | None = None
    flags = (
        os.O_RDONLY
        | getattr(os, "O_DIRECTORY", 0)
        | getattr(os, "O_NOFOLLOW", 0)
    )
    try:
        try:
            local_fd = os.open(LOCAL_STATE_REL.name, flags, dir_fd=project_fd)
        except FileNotFoundError:
            if not create:
                return project_fd, None
            try:
                os.mkdir(LOCAL_STATE_REL.name, mode=0o700, dir_fd=project_fd)
            except FileExistsError:
                pass
            local_fd = os.open(LOCAL_STATE_REL.name, flags, dir_fd=project_fd)
        assert local_fd is not None
        os.fchmod(local_fd, 0o700)
        _verify_local_state_binding(project_fd, local_fd)
        return project_fd, local_fd
    except OSError as exc:
        if local_fd is not None:
            os.close(local_fd)
        os.close(project_fd)
        raise ValueError(f"cannot safely open local source state: {exc}") from exc
    except BaseException:
        if local_fd is not None:
            os.close(local_fd)
        os.close(project_fd)
        raise


def _read_private_text_at(local_fd: int, name: str) -> str | None:
    try:
        visible = os.stat(name, dir_fd=local_fd, follow_symlinks=False)
    except FileNotFoundError:
        return None
    if not stat.S_ISREG(visible.st_mode) or visible.st_nlink != 1:
        raise ValueError(f"local state file must be a private regular file: {name}")
    fd = os.open(
        name,
        os.O_RDONLY | getattr(os, "O_NOFOLLOW", 0),
        dir_fd=local_fd,
    )
    try:
        opened = os.fstat(fd)
        if (
            not stat.S_ISREG(opened.st_mode)
            or opened.st_nlink != 1
            or (opened.st_dev, opened.st_ino) != (visible.st_dev, visible.st_ino)
        ):
            raise ValueError(f"local state file changed while opening: {name}")
        chunks: list[bytes] = []
        total = 0
        while True:
            chunk = os.read(fd, 1024 * 1024)
            if not chunk:
                break
            total += len(chunk)
            if total > MAX_LOCAL_STATE_FILE_BYTES:
                raise ValueError(f"local state file is too large: {name}")
            chunks.append(chunk)
        final_opened = os.fstat(fd)
        try:
            final_visible = os.stat(
                name,
                dir_fd=local_fd,
                follow_symlinks=False,
            )
        except FileNotFoundError as exc:
            raise ValueError(f"local state file changed while reading: {name}") from exc
        if (
            (final_visible.st_dev, final_visible.st_ino)
            != (opened.st_dev, opened.st_ino)
            or (final_opened.st_size, final_opened.st_mtime_ns, final_opened.st_ctime_ns)
            != (opened.st_size, opened.st_mtime_ns, opened.st_ctime_ns)
        ):
            raise ValueError(f"local state file changed while reading: {name}")
        return b"".join(chunks).decode("utf-8")
    finally:
        os.close(fd)


def _load_local_source_map_at(
    local_fd: int,
    *,
    allow_missing: bool = False,
) -> dict[str, str]:
    raw = _read_private_text_at(local_fd, LOCAL_SOURCE_MAP_REL.name)
    if raw is None:
        if allow_missing:
            return {}
        raise ValueError("local source map is missing from existing local state")
    try:
        payload = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise ValueError(f"invalid local source map: {exc}") from exc
    if not isinstance(payload, dict) or payload.get("version") != LOCAL_SOURCE_MAP_VERSION:
        raise ValueError("invalid local source map version")
    sources = payload.get("sources")
    if not isinstance(sources, dict) or not all(
        isinstance(key, str)
        and bool(key)
        and isinstance(value, str)
        and bool(value)
        and Path(value).is_absolute()
        for key, value in sources.items()
    ):
        raise ValueError("invalid local source map entries")
    return dict(sources)


def _load_local_source_map(project: Path) -> dict[str, str]:
    project_fd, local_fd = _open_local_state(project, create=False)
    try:
        if local_fd is None:
            return {}
        sources = _load_local_source_map_at(local_fd)
        _verify_local_state_binding(project_fd, local_fd)
        return sources
    finally:
        if local_fd is not None:
            os.close(local_fd)
        os.close(project_fd)


def load_local_source_paths(project: Path) -> dict[str, str]:
    """Read the private source map through stable directory descriptors."""
    try:
        return _load_local_source_map(project)
    except (OSError, UnicodeError) as exc:
        raise ValueError("local source map is unreadable") from exc


def load_local_source_paths_from_project_fd(project_fd: int) -> dict[str, str]:
    """Read the private source map below an already-bound project root."""
    local_fd: int | None = None
    flags = (
        os.O_RDONLY
        | getattr(os, "O_DIRECTORY", 0)
        | getattr(os, "O_NOFOLLOW", 0)
    )
    try:
        try:
            local_fd = os.open(LOCAL_STATE_REL.name, flags, dir_fd=project_fd)
        except FileNotFoundError:
            return {}
        sources = _load_local_source_map_at(local_fd)
        _verify_local_state_binding(project_fd, local_fd)
        return sources
    except (OSError, UnicodeError) as exc:
        raise ValueError("local source map is unreadable") from exc
    finally:
        if local_fd is not None:
            os.close(local_fd)


def _atomic_write_private_text_at(local_fd: int, name: str, text: str) -> None:
    try:
        current = os.stat(name, dir_fd=local_fd, follow_symlinks=False)
    except FileNotFoundError:
        current = None
    if current is not None and (
        not stat.S_ISREG(current.st_mode) or current.st_nlink != 1
    ):
        raise ValueError(f"local state target must be a private regular file: {name}")

    temp_name = f".adco-private-{secrets.token_hex(12)}"
    temp_fd = os.open(
        temp_name,
        os.O_WRONLY
        | os.O_CREAT
        | os.O_EXCL
        | getattr(os, "O_NOFOLLOW", 0),
        mode=0o600,
        dir_fd=local_fd,
    )
    try:
        try:
            data = text.encode("utf-8")
            view = memoryview(data)
            while view:
                written = os.write(temp_fd, view)
                if written <= 0:
                    raise OSError("private local-state write made no progress")
                view = view[written:]
            os.fchmod(temp_fd, 0o600)
            os.fsync(temp_fd)
        finally:
            os.close(temp_fd)
        try:
            current = os.stat(name, dir_fd=local_fd, follow_symlinks=False)
        except FileNotFoundError:
            current = None
        if current is not None and (
            not stat.S_ISREG(current.st_mode) or current.st_nlink != 1
        ):
            raise ValueError(
                f"local state target raced to an unsafe file: {name}"
            )
        os.replace(
            temp_name,
            name,
            src_dir_fd=local_fd,
            dst_dir_fd=local_fd,
        )
        os.fsync(local_fd)
    finally:
        try:
            os.unlink(temp_name, dir_fd=local_fd)
        except FileNotFoundError:
            pass


def register_local_source_path(project: Path, source_event_id: str, path: Path) -> str:
    """Store an external absolute path only in ignored, owner-readable local state."""
    if not source_event_id:
        raise ValueError("source_event_id is required for external material")
    project_fd, local_fd = _open_local_state(project, create=True)
    assert local_fd is not None
    try:
        if _read_private_text_at(local_fd, ".gitignore") != LOCAL_GITIGNORE_TEXT:
            _atomic_write_private_text_at(
                local_fd,
                ".gitignore",
                LOCAL_GITIGNORE_TEXT,
            )
            _verify_local_state_binding(project_fd, local_fd)

        sources = _load_local_source_map_at(local_fd, allow_missing=True)
        sources[source_event_id] = str(path.resolve())
        payload = {
            "version": LOCAL_SOURCE_MAP_VERSION,
            "sources": dict(sorted(sources.items())),
        }
        _atomic_write_private_text_at(
            local_fd,
            LOCAL_SOURCE_MAP_REL.name,
            json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
        )
        _verify_local_state_binding(project_fd, local_fd)
    finally:
        os.close(local_fd)
        os.close(project_fd)
    return f"{LOCAL_SOURCE_PREFIX}{source_event_id}"


def source_row_material_roots(project: Path, row: dict[str, str]) -> list[Path]:
    raw_paths = [
        item.strip()
        for item in row.get("file_paths", "").split(";")
        if item.strip()
    ]
    local_sources: dict[str, str] | None = None
    roots: list[Path] = []
    for raw in raw_paths:
        if raw.startswith(LOCAL_SOURCE_PREFIX):
            source_event_id = raw.removeprefix(LOCAL_SOURCE_PREFIX).split("/", 1)[0]
            local_sources = local_sources or _load_local_source_map(project)
            mapped = local_sources.get(source_event_id, "")
            if not mapped:
                raise ValueError(f"local source path is unavailable for {source_event_id}")
            root = Path(mapped).expanduser()
        else:
            root = Path(raw).expanduser()
            if not root.is_absolute():
                root = project / root
        roots.append(root)
    return roots


def source_path_label(project: Path, row: dict[str, str], path: Path) -> str:
    source_event_id = row.get("source_event_id", "") or "unregistered"
    for raw, root in zip(
        [item.strip() for item in row.get("file_paths", "").split(";") if item.strip()],
        source_row_material_roots(project, row),
    ):
        try:
            relative = path.resolve().relative_to(root.resolve())
        except ValueError:
            continue
        try:
            root.resolve().relative_to(project.resolve())
            external_root = False
        except ValueError:
            external_root = True
        if raw.startswith(LOCAL_SOURCE_PREFIX) or external_root:
            base = (
                raw
                if raw.startswith(LOCAL_SOURCE_PREFIX)
                else f"external-source://{source_event_id}"
            )
            return base if relative == Path() else f"{base}/{relative.as_posix()}"
        return display_source_path(project, path)
    return f"external-source://{source_event_id}/{path.name}"


def material_files(path: Path) -> list[Path]:
    if path.is_symlink():
        raise ValueError(f"material path must not be a symlink: {path}")
    if path.is_file():
        return [path] if path.suffix.lower() in SUPPORTED_SUFFIXES else []
    if not path.is_dir():
        return []
    files: list[Path] = []
    for item in path.rglob("*"):
        if item.is_symlink():
            raise ValueError(f"material tree must not contain symlinks: {item}")
        if (
            item.is_file()
            and item.suffix.lower() in SUPPORTED_SUFFIXES
            and "AD-creative" not in item.parts
        ):
            files.append(item)
    return sorted(files)


def _safe_read_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="utf-8", errors="replace")


def _line_number(text: str, char_offset: int) -> int:
    return text.count("\n", 0, max(char_offset, 0)) + 1


def _split_text(
    text: str,
    *,
    chunk_chars: int,
    overlap_chars: int,
    base_line: int = 1,
) -> list[RawEvidence]:
    if not text.strip():
        return []
    chunk_chars = max(500, chunk_chars)
    overlap_chars = min(max(0, overlap_chars), chunk_chars // 3)
    boundaries = [
        match.end()
        for match in re.finditer(r"(?m)\n\s*\n|^#{1,6}\s+", text)
    ]
    boundaries.append(len(text))
    chunks: list[RawEvidence] = []
    start = 0
    while start < len(text):
        target = min(len(text), start + chunk_chars)
        candidates = [item for item in boundaries if start + 500 <= item <= target]
        end = max(candidates) if candidates else target
        if end <= start:
            end = min(len(text), start + chunk_chars)
        value = text[start:end].strip()
        if value:
            chunks.append(
                RawEvidence(
                    text=value,
                    start_line=base_line + _line_number(text, start) - 1,
                    end_line=base_line + _line_number(text, end) - 1,
                )
            )
        if end >= len(text):
            break
        start = max(start + 1, end - overlap_chars)
    return chunks


def _split_raw_record(
    record: RawEvidence,
    *,
    chunk_chars: int,
    overlap_chars: int,
) -> list[RawEvidence]:
    if len(record.text) <= chunk_chars:
        return [record]
    pieces = _split_text(
        record.text,
        chunk_chars=chunk_chars,
        overlap_chars=overlap_chars,
        base_line=record.start_line or 1,
    )
    return [
        RawEvidence(
            text=piece.text,
            page=record.page,
            slide=record.slide,
            start_line=piece.start_line,
            end_line=piece.end_line,
            field_path=record.field_path,
            start_time=record.start_time,
            end_time=record.end_time,
            inspection_status=record.inspection_status,
            metadata=dict(record.metadata),
        )
        for piece in pieces
    ]


def parse_plain_text(path: Path, *, chunk_chars: int, overlap_chars: int) -> list[RawEvidence]:
    return _split_text(
        _safe_read_text(path),
        chunk_chars=chunk_chars,
        overlap_chars=overlap_chars,
    )


def parse_csv(path: Path, *, chunk_chars: int, overlap_chars: int) -> list[RawEvidence]:
    records: list[RawEvidence] = []
    with path.open(newline="", encoding="utf-8-sig", errors="replace") as handle:
        reader = csv.DictReader(handle)
        for row_number, row in enumerate(reader, 2):
            values = [
                f"row[{row_number}].{column}: {value}"
                for column, value in row.items()
                if column is not None and value not in {None, ""}
            ]
            if values:
                records.append(
                    RawEvidence(
                        text="\n".join(values),
                        start_line=row_number,
                        end_line=row_number,
                        field_path=f"row[{row_number}]",
                    )
                )
    return [
        piece
        for record in records
        for piece in _split_raw_record(
            record, chunk_chars=chunk_chars, overlap_chars=overlap_chars
        )
    ]


def _json_leaf_records(value: object, path: str = "$") -> Iterable[tuple[str, object]]:
    if isinstance(value, dict):
        for key, item in value.items():
            escaped = str(key).replace("~", "~0").replace("/", "~1")
            yield from _json_leaf_records(item, f"{path}/{escaped}")
    elif isinstance(value, list):
        for index, item in enumerate(value):
            yield from _json_leaf_records(item, f"{path}/{index}")
    else:
        yield path, value


def parse_json(path: Path, *, chunk_chars: int, overlap_chars: int) -> list[RawEvidence]:
    payload = json.loads(_safe_read_text(path))
    records = [
        RawEvidence(
            text=f"{field_path}: {json.dumps(value, ensure_ascii=False)}",
            field_path=field_path,
        )
        for field_path, value in _json_leaf_records(payload)
    ]
    return [
        piece
        for record in records
        for piece in _split_raw_record(
            record, chunk_chars=chunk_chars, overlap_chars=overlap_chars
        )
    ]


def parse_yaml(path: Path, *, chunk_chars: int, overlap_chars: int) -> list[RawEvidence]:
    text = _safe_read_text(path)
    stack: list[tuple[int, str]] = []
    records: list[RawEvidence] = []
    for line_number, raw in enumerate(text.splitlines(), 1):
        if not raw.strip() or raw.lstrip().startswith("#"):
            continue
        indent = len(raw) - len(raw.lstrip(" "))
        match = re.match(r"\s*(?:-\s*)?([^:#][^:]*):(?:\s*(.*))?$", raw)
        if match:
            key = match.group(1).strip()
            while stack and stack[-1][0] >= indent:
                stack.pop()
            stack.append((indent, key))
            field_path = ".".join(item[1] for item in stack)
            value = (match.group(2) or "").strip()
        else:
            field_path = ".".join(item[1] for item in stack) or f"line[{line_number}]"
            value = raw.strip()
        records.append(
            RawEvidence(
                text=f"{field_path}: {value}" if value else f"{field_path}:",
                start_line=line_number,
                end_line=line_number,
                field_path=field_path,
            )
        )
    return [
        piece
        for record in records
        for piece in _split_raw_record(
            record, chunk_chars=chunk_chars, overlap_chars=overlap_chars
        )
    ]


def parse_docx(path: Path, *, chunk_chars: int, overlap_chars: int) -> list[RawEvidence]:
    from docx import Document

    document = Document(path)
    records: list[RawEvidence] = []
    for index, paragraph in enumerate(document.paragraphs, 1):
        if paragraph.text.strip():
            records.append(
                RawEvidence(
                    text=paragraph.text.strip(),
                    start_line=index,
                    end_line=index,
                    field_path=f"paragraph[{index}]",
                )
            )
    for table_index, table in enumerate(document.tables, 1):
        for row_index, row in enumerate(table.rows, 1):
            for column_index, cell in enumerate(row.cells, 1):
                if cell.text.strip():
                    field_path = f"table[{table_index}].row[{row_index}].cell[{column_index}]"
                    records.append(
                        RawEvidence(text=cell.text.strip(), field_path=field_path)
                    )
    return [
        piece
        for record in records
        for piece in _split_raw_record(
            record, chunk_chars=chunk_chars, overlap_chars=overlap_chars
        )
    ]


def parse_pptx(path: Path, *, chunk_chars: int, overlap_chars: int) -> list[RawEvidence]:
    from pptx import Presentation

    presentation = Presentation(path)
    records: list[RawEvidence] = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        for shape_index, shape in enumerate(slide.shapes, 1):
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                records.append(
                    RawEvidence(
                        text=text.strip(),
                        slide=slide_number,
                        field_path=f"slide[{slide_number}].shape[{shape_index}]",
                    )
                )
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows, 1):
                    for cell_index, cell in enumerate(row.cells, 1):
                        if cell.text.strip():
                            records.append(
                                RawEvidence(
                                    text=cell.text.strip(),
                                    slide=slide_number,
                                    field_path=(
                                        f"slide[{slide_number}].shape[{shape_index}]"
                                        f".table.row[{row_index}].cell[{cell_index}]"
                                    ),
                                )
                            )
        try:
            notes_text = slide.notes_slide.notes_text_frame.text
        except (AttributeError, ValueError):
            notes_text = ""
        if notes_text.strip():
            records.append(
                RawEvidence(
                    text=notes_text.strip(),
                    slide=slide_number,
                    field_path=f"slide[{slide_number}].notes",
                )
            )
    return [
        piece
        for record in records
        for piece in _split_raw_record(
            record, chunk_chars=chunk_chars, overlap_chars=overlap_chars
        )
    ]


def parse_pdf(path: Path, *, chunk_chars: int, overlap_chars: int) -> list[RawEvidence]:
    try:
        from pypdf import PdfReader
    except ImportError:
        pdftotext = shutil.which("pdftotext")
        if not pdftotext:
            raise
        completed = subprocess.run(
            [pdftotext, "-layout", str(path), "-"],
            check=False,
            capture_output=True,
            text=True,
            timeout=30,
        )
        if completed.returncode != 0:
            raise ValueError(
                "pdftotext failed: " + completed.stderr.strip()[:1000]
            )
        records: list[RawEvidence] = []
        for page_number, text in enumerate(completed.stdout.split("\f"), 1):
            for piece in _split_text(
                text,
                chunk_chars=chunk_chars,
                overlap_chars=overlap_chars,
            ):
                piece.page = page_number
                piece.field_path = f"page[{page_number}]"
                records.append(piece)
        return records

    reader = PdfReader(path)
    records: list[RawEvidence] = []
    for page_number, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        for piece in _split_text(
            text,
            chunk_chars=chunk_chars,
            overlap_chars=overlap_chars,
        ):
            piece.page = page_number
            piece.field_path = f"page[{page_number}]"
            records.append(piece)
    return records


def parse_subtitle(path: Path, *, chunk_chars: int, overlap_chars: int) -> list[RawEvidence]:
    text = _safe_read_text(path).replace("\r\n", "\n")
    blocks = re.split(r"\n\s*\n", text)
    records: list[RawEvidence] = []
    cursor_line = 1
    for block in blocks:
        lines = [line.strip() for line in block.splitlines() if line.strip()]
        timing_index = next((i for i, line in enumerate(lines) if "-->" in line), None)
        if timing_index is None:
            cursor_line += block.count("\n") + 2
            continue
        start_time, end_time = [item.strip() for item in lines[timing_index].split("-->", 1)]
        caption = "\n".join(lines[timing_index + 1 :]).strip()
        if caption:
            records.append(
                RawEvidence(
                    text=caption,
                    start_line=cursor_line,
                    end_line=cursor_line + block.count("\n"),
                    start_time=start_time,
                    end_time=end_time,
                    field_path=f"cue[{len(records) + 1}]",
                )
            )
        cursor_line += block.count("\n") + 2
    return [
        piece
        for record in records
        for piece in _split_raw_record(
            record, chunk_chars=chunk_chars, overlap_chars=overlap_chars
        )
    ]


def parse_image(
    path: Path,
    *,
    chunk_chars: int,
    overlap_chars: int,
    file_hash: str,
) -> list[RawEvidence]:
    from PIL import ExifTags, Image

    metadata: dict[str, object] = {"file_sha256": file_hash}
    try:
        with Image.open(path) as image:
            metadata.update(
                {
                    "width": image.width,
                    "height": image.height,
                    "format": image.format or path.suffix.lstrip(".").upper(),
                    "mode": image.mode,
                }
            )
            try:
                exif = image.getexif()
            except (AttributeError, OSError):
                exif = {}
            safe_exif = {}
            for key, value in exif.items():
                tag = str(ExifTags.TAGS.get(key, key))
                if tag in SAFE_EXIF_TAGS:
                    safe_exif[tag] = str(value)[:100]
            if safe_exif:
                metadata["exif"] = safe_exif
    except (OSError, ValueError) as exc:
        error = str(exc)
        for raw_path in {str(path), str(path.resolve())}:
            error = error.replace(raw_path, path.name)
        metadata["metadata_error"] = error
    return [
        RawEvidence(
            text=json.dumps(metadata, ensure_ascii=False, sort_keys=True),
            inspection_status="requires_visual_inspection",
            metadata=metadata,
        )
    ]


def parse_video(
    path: Path,
    *,
    chunk_chars: int,
    overlap_chars: int,
    file_hash: str,
) -> list[RawEvidence]:
    metadata: dict[str, object] = {"file_sha256": file_hash}
    ffprobe = shutil.which("ffprobe")
    inspection_status = "requires_media_inspection"
    if ffprobe:
        completed = subprocess.run(
            [
                ffprobe,
                "-v",
                "error",
                "-show_entries",
                "format=duration:stream=index,codec_type,codec_name,width,height,r_frame_rate,channels,sample_rate",
                "-of",
                "json",
                str(path),
            ],
            check=False,
            capture_output=True,
            text=True,
            timeout=30,
        )
        if completed.returncode == 0:
            try:
                payload = json.loads(completed.stdout)
                if not isinstance(payload, dict):
                    raise ValueError("ffprobe payload must be an object")
                safe_format: dict[str, object] = {}
                raw_format = payload.get("format")
                if isinstance(raw_format, dict):
                    duration = str(raw_format.get("duration", "")).strip()
                    if re.fullmatch(r"\d+(?:\.\d+)?", duration):
                        safe_format["duration"] = duration
                safe_streams: list[dict[str, object]] = []
                raw_streams = payload.get("streams")
                if isinstance(raw_streams, list):
                    for raw_stream in raw_streams:
                        if not isinstance(raw_stream, dict):
                            continue
                        stream: dict[str, object] = {}
                        for key in SAFE_FFPROBE_STREAM_FIELDS:
                            value = raw_stream.get(key)
                            if key in {"index", "width", "height", "channels"}:
                                if isinstance(value, int) and value >= 0:
                                    stream[key] = value
                            elif key == "sample_rate":
                                normalized = str(value or "").strip()
                                if re.fullmatch(r"\d{1,7}", normalized):
                                    stream[key] = normalized
                            elif key == "r_frame_rate":
                                normalized = str(value or "").strip()
                                if re.fullmatch(r"\d{1,9}/\d{1,9}", normalized):
                                    stream[key] = normalized
                            else:
                                normalized = str(value or "").strip()
                                if re.fullmatch(r"[A-Za-z0-9_.-]{1,64}", normalized):
                                    stream[key] = normalized
                        if stream:
                            safe_streams.append(stream)
                metadata["ffprobe"] = {
                    "format": safe_format,
                    "streams": safe_streams,
                }
            except (json.JSONDecodeError, ValueError):
                metadata["ffprobe_error"] = "invalid_json"
        else:
            error = completed.stderr.strip()
            for raw_path in {str(path), str(path.resolve())}:
                error = error.replace(raw_path, path.name)
            metadata["ffprobe_error"] = error[:1000]
    else:
        metadata["ffprobe"] = "unavailable"
    return [
        RawEvidence(
            text=json.dumps(metadata, ensure_ascii=False, sort_keys=True),
            inspection_status=inspection_status,
            metadata=metadata,
        )
    ]


Parser = Callable[..., list[RawEvidence]]
PARSER_REGISTRY: dict[str, Parser] = {
    **{suffix: parse_plain_text for suffix in PLAIN_TEXT_SUFFIXES},
    ".csv": parse_csv,
    ".json": parse_json,
    ".yaml": parse_yaml,
    ".yml": parse_yaml,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".pdf": parse_pdf,
    **{suffix: parse_subtitle for suffix in SUBTITLE_SUFFIXES},
    **{suffix: parse_image for suffix in IMAGE_SUFFIXES},
    **{suffix: parse_video for suffix in VIDEO_SUFFIXES},
}


def parse_file(
    project: Path,
    path: Path,
    source_event_id: str,
    *,
    chunk_chars: int = DEFAULT_CHUNK_CHARS,
    overlap_chars: int = DEFAULT_OVERLAP_CHARS,
    source_path_label: str = "",
) -> list[EvidenceChunk]:
    parser = PARSER_REGISTRY.get(path.suffix.lower())
    if parser is None:
        raise ValueError(f"unsupported material type: {path.suffix.lower()}")
    file_hash = sha256_file(path)
    parser_kwargs: dict[str, object] = {}
    if path.suffix.lower() in IMAGE_SUFFIXES | VIDEO_SUFFIXES:
        parser_kwargs["file_hash"] = file_hash
    raw_chunks = parser(
        path,
        chunk_chars=chunk_chars,
        overlap_chars=overlap_chars,
        **parser_kwargs,
    )
    source_path = source_path_label or display_source_path(project, path)
    media_type = media_type_for(path)
    chunks: list[EvidenceChunk] = []
    for index, raw in enumerate(raw_chunks, 1):
        chunk_hash = sha256_text(raw.text)
        identity = sha256_text(
            f"{source_event_id}\0{source_path}\0{index}\0{chunk_hash}"
        )[:20]
        metadata = {"file_sha256": file_hash, **raw.metadata}
        chunks.append(
            EvidenceChunk(
                chunk_id=f"EVC-{identity}",
                source_event_id=source_event_id,
                source_path=source_path,
                media_type=media_type,
                page=raw.page,
                slide=raw.slide,
                start_line=raw.start_line,
                end_line=raw.end_line,
                text=raw.text,
                sha256=(file_hash if raw.inspection_status.startswith("requires_") else chunk_hash),
                inspection_status=raw.inspection_status,
                field_path=raw.field_path,
                start_time=raw.start_time,
                end_time=raw.end_time,
                metadata=metadata,
            )
        )
    return chunks


def load_evidence_chunks(project: Path) -> list[EvidenceChunk]:
    path = project / EVIDENCE_REL
    if not path.is_file():
        return []
    chunks: list[EvidenceChunk] = []
    with path.open(encoding="utf-8") as handle:
        for line_number, line in enumerate(handle, 1):
            if not line.strip():
                continue
            try:
                payload = json.loads(line)
            except json.JSONDecodeError as exc:
                raise ValueError(f"invalid evidence chunk JSONL at line {line_number}: {exc}") from exc
            if not isinstance(payload, dict):
                raise ValueError(f"invalid evidence chunk record at line {line_number}")
            chunks.append(EvidenceChunk.from_dict(payload))
    return chunks


def write_evidence_chunks(project: Path, chunks: Iterable[EvidenceChunk]) -> Path:
    path = project / EVIDENCE_REL
    ordered = sorted(
        chunks,
        key=lambda item: (
            item.source_event_id,
            item.source_path.casefold(),
            item.page or 0,
            item.slide or 0,
            item.start_line or 0,
            item.chunk_id,
        ),
    )
    text = "".join(
        json.dumps(item.as_dict(), ensure_ascii=False, sort_keys=True) + "\n"
        for item in ordered
    )
    atomic_write_text(project, path, text)
    return path


def source_row_files(project: Path, row: dict[str, str]) -> list[tuple[Path, str]]:
    files: dict[Path, str] = {}
    for root in source_row_material_roots(project, row):
        for path in material_files(root):
            files[path.resolve()] = source_path_label(project, row, path)
    return sorted(files.items(), key=lambda item: item[0])


def ingest_source_rows(
    project: Path,
    source_rows: list[dict[str, str]],
    *,
    max_total_chars: int = DEFAULT_TOTAL_CHAR_BUDGET,
    chunk_chars: int = DEFAULT_CHUNK_CHARS,
    overlap_chars: int = DEFAULT_OVERLAP_CHARS,
) -> IngestionReport:
    max_total_chars = max(0, int(max_total_chars))
    processed_source_ids = {
        row.get("source_event_id", "") for row in source_rows if row.get("source_event_id", "")
    }
    preserved = [
        chunk
        for chunk in load_evidence_chunks(project)
        if chunk.source_event_id not in processed_source_ids
    ]
    accepted: list[EvidenceChunk] = []
    over_budget: list[dict[str, object]] = []
    parser_errors: list[dict[str, str]] = []
    files_processed = 0
    characters_read = 0
    seen_files: set[tuple[str, Path]] = set()
    for row in source_rows:
        source_event_id = row.get("source_event_id", "")
        try:
            row_files = source_row_files(project, row)
        except Exception as exc:
            parser_errors.append(
                {
                    "source_path": f"source-event://{source_event_id or 'unknown'}",
                    "media_type": "application/octet-stream",
                    "error": f"{type(exc).__name__}: material source is unavailable or unsafe",
                }
            )
            continue
        for path, path_label in row_files:
            identity = (source_event_id, path.resolve())
            if identity in seen_files:
                continue
            seen_files.add(identity)
            try:
                file_chunks = parse_file(
                    project,
                    path,
                    source_event_id,
                    chunk_chars=chunk_chars,
                    overlap_chars=overlap_chars,
                    source_path_label=path_label,
                )
            except Exception as exc:  # parser boundary must report, not hide, failures
                error = str(exc)
                for raw_path in {str(path), str(path.resolve())}:
                    error = error.replace(raw_path, path_label)
                parser_errors.append(
                    {
                        "source_path": path_label,
                        "media_type": media_type_for(path),
                        "error": f"{type(exc).__name__}: {error}",
                    }
                )
                continue
            file_characters = sum(len(chunk.text) for chunk in file_chunks)
            if characters_read + file_characters > max_total_chars:
                over_budget.append(
                    {
                        "source_path": path_label,
                        "media_type": media_type_for(path),
                        "characters": file_characters,
                        "reason": "total_character_budget_exceeded",
                    }
                )
                continue
            accepted.extend(file_chunks)
            files_processed += 1
            characters_read += file_characters
    write_evidence_chunks(project, [*preserved, *accepted])
    return IngestionReport(
        chunks=accepted,
        files_processed=files_processed,
        characters_read=characters_read,
        over_budget=over_budget,
        parser_errors=parser_errors,
    )
