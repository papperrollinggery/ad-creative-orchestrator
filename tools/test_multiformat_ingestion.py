#!/usr/bin/env python3
"""Regression checks for provenance-preserving multi-format ingestion."""

from __future__ import annotations

import importlib.util
import hashlib
import json
import shutil
import subprocess
import tempfile
from pathlib import Path
from unittest.mock import patch

from adco_core.ingestion import ingest_source_rows, parse_file


def _write_pdf(path: Path, text: str) -> None:
    stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    data = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for number, payload in enumerate(objects, 1):
        offsets.append(len(data))
        data.extend(f"{number} 0 obj\n".encode("ascii"))
        data.extend(payload)
        data.extend(b"\nendobj\n")
    xref = len(data)
    data.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    data.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        data.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    data.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode(
            "ascii"
        )
    )
    path.write_bytes(bytes(data))


def _make_materials(root: Path) -> tuple[list[Path], list[str]]:
    paths: list[Path] = []
    optional_skips: list[str] = []
    long_tail = "END-OF-LONG-BRIEF"
    fixtures = {
        "brief.md": "# Brief\n\n客户必须保留可编辑文本。\n" + ("长资料证据。" * 2600) + long_tail,
        "notes.txt": "客户已提供产品图，不应反推为缺失。",
        "facts.csv": "name,value\nproduct,Trail Shell\naudience,City commuter\n",
        "facts.json": json.dumps({"brand": {"name": "NOVA", "tone": "specific"}}),
        "facts.yaml": "brand:\n  name: NOVA\n  tone: grounded\n",
        "captions.srt": "1\n00:00:00,000 --> 00:00:02,000\nSRT evidence line\n",
        "captions.vtt": "WEBVTT\n\n00:00:02.000 --> 00:00:04.000\nVTT evidence line\n",
    }
    for name, content in fixtures.items():
        path = root / name
        path.write_text(content, encoding="utf-8")
        paths.append(path)

    if importlib.util.find_spec("docx"):
        from docx import Document

        path = root / "brief.docx"
        document = Document()
        document.add_paragraph("DOCX evidence paragraph")
        document.save(path)
        paths.append(path)
    else:
        optional_skips.append("docx")

    if importlib.util.find_spec("pptx"):
        from pptx import Presentation
        from pptx.util import Inches

        path = root / "deck.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text = (
            "PPTX evidence slide"
        )
        presentation.save(path)
        paths.append(path)
    else:
        optional_skips.append("pptx")

    if importlib.util.find_spec("pypdf") or shutil.which("pdftotext"):
        path = root / "source.pdf"
        _write_pdf(path, "PDF evidence page")
        paths.append(path)
    else:
        optional_skips.append("pdf")

    if importlib.util.find_spec("PIL"):
        from PIL import Image

        path = root / "product.png"
        Image.new("RGB", (32, 24), (12, 34, 56)).save(path)
        paths.append(path)
    else:
        optional_skips.append("image")

    video = root / "reference.mp4"
    video.write_bytes(b"not-a-real-video-but-a-valid-media-ingestion-fixture")
    paths.append(video)
    return paths, optional_skips


def test_all_available_formats_and_long_material() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-ingestion-formats-") as raw:
        project = Path(raw)
        materials = project / "materials"
        materials.mkdir()
        paths, optional_skips = _make_materials(materials)
        source_rows = [
            {
                "source_event_id": "SRC-001",
                "file_paths": ";".join(path.relative_to(project).as_posix() for path in paths),
                "declared_semantics": "test fixtures",
            }
        ]
        report = ingest_source_rows(project, source_rows)
        assert report.files_processed == len(paths), report.as_dict()
        assert not report.parser_errors, report.parser_errors
        assert not report.over_budget, report.over_budget
        assert report.characters_read > 12_000, report.characters_read
        assert any("END-OF-LONG-BRIEF" in chunk.text for chunk in report.chunks)
        assert len([chunk for chunk in report.chunks if chunk.source_path.endswith("brief.md")]) > 4
        assert any(chunk.field_path == "$/brand/name" for chunk in report.chunks)
        assert any(chunk.field_path == "row[2]" for chunk in report.chunks)
        assert any(chunk.start_time for chunk in report.chunks if chunk.source_path.endswith((".srt", ".vtt")))
        if "pptx" not in optional_skips:
            assert any(chunk.slide == 1 for chunk in report.chunks)
        if "pdf" not in optional_skips:
            assert any(chunk.page == 1 and "PDF evidence" in chunk.text for chunk in report.chunks)
        if "image" not in optional_skips:
            assert any(chunk.inspection_status == "requires_visual_inspection" for chunk in report.chunks)
        assert any(chunk.inspection_status == "requires_media_inspection" for chunk in report.chunks)
        if optional_skips:
            print("TEST_MULTIFORMAT_OPTIONAL_SKIPS=" + ",".join(optional_skips))


def test_budget_overflow_is_reported_not_silently_truncated() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-ingestion-budget-") as raw:
        project = Path(raw)
        source = project / "large.txt"
        source.write_text("budget evidence " * 1500, encoding="utf-8")
        report = ingest_source_rows(
            project,
            [
                {
                    "source_event_id": "SRC-002",
                    "file_paths": source.name,
                    "declared_semantics": "budget fixture",
                }
            ],
            max_total_chars=1_000,
        )
        assert report.files_processed == 0
        assert report.characters_read == 0
        assert len(report.over_budget) == 1
        assert report.over_budget[0]["reason"] == "total_character_budget_exceeded"


def test_media_hashing_is_streaming_and_exif_is_allowlisted() -> None:
    with tempfile.TemporaryDirectory(prefix="adco-ingestion-media-") as raw:
        project = Path(raw)
        video = project / "reference.mp4"
        video.write_bytes(b"streaming-media-fixture")
        expected_hash = hashlib.sha256(video.read_bytes()).hexdigest()

        with patch.object(Path, "read_bytes", side_effect=AssertionError("whole-file read")):
            video_chunks = parse_file(project, video, "SRC-VIDEO")
        assert video_chunks[0].metadata["file_sha256"] == expected_hash

        private_marker = "PRIVATE-LOCATION-31.2304-121.4737"
        ffprobe_payload = {
            "format": {
                "duration": "12.5",
                "filename": str(video),
                "tags": {
                    "creation_time": "2026-07-21T12:34:56Z",
                    "location": private_marker,
                    "comment": private_marker,
                },
            },
            "streams": [
                {
                    "index": 0,
                    "codec_type": "video",
                    "codec_name": "h264",
                    "width": 1920,
                    "height": 1080,
                    "r_frame_rate": "30000/1001",
                    "tags": {
                        "handler_name": private_marker,
                        "encoder": private_marker,
                        "comment": private_marker,
                    },
                }
            ],
        }
        completed = subprocess.CompletedProcess(
            args=["ffprobe"],
            returncode=0,
            stdout=json.dumps(ffprobe_payload),
            stderr="",
        )
        with (
            patch("adco_core.ingestion.shutil.which", return_value="/usr/bin/ffprobe"),
            patch("adco_core.ingestion.subprocess.run", return_value=completed),
            patch.object(Path, "read_bytes", side_effect=AssertionError("whole-file read")),
        ):
            sanitized_chunks = parse_file(project, video, "SRC-VIDEO-SANITIZED")
        serialized_video = json.dumps(
            sanitized_chunks[0].metadata,
            ensure_ascii=False,
            sort_keys=True,
        )
        assert private_marker not in serialized_video
        assert str(video) not in serialized_video
        assert "creation_time" not in serialized_video
        assert "handler_name" not in serialized_video
        assert sanitized_chunks[0].metadata["ffprobe"] == {
            "format": {"duration": "12.5"},
            "streams": [
                {
                    "index": 0,
                    "codec_type": "video",
                    "codec_name": "h264",
                    "width": 1920,
                    "height": 1080,
                    "r_frame_rate": "30000/1001",
                }
            ],
        }

        if importlib.util.find_spec("PIL"):
            from PIL import Image

            image_path = project / "private.jpg"
            image = Image.new("RGB", (16, 12), (20, 40, 60))
            exif = Image.Exif()
            exif[271] = "Private Camera Maker"
            exif[272] = "Private Camera Model"
            exif[306] = "2026:07:21 12:34:56"
            exif[274] = 1
            image.save(image_path, exif=exif)
            with patch.object(Path, "read_bytes", side_effect=AssertionError("whole-file read")):
                image_chunks = parse_file(project, image_path, "SRC-IMAGE")
            metadata = image_chunks[0].metadata
            assert metadata.get("exif") == {"Orientation": "1"}, metadata
            serialized = json.dumps(metadata, ensure_ascii=False)
            assert "Private Camera" not in serialized
            assert "2026:07:21" not in serialized


def main() -> int:
    test_all_available_formats_and_long_material()
    test_budget_overflow_is_reported_not_silently_truncated()
    test_media_hashing_is_streaming_and_exif_is_allowlisted()
    print("TEST_MULTIFORMAT_INGESTION=PASS")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
